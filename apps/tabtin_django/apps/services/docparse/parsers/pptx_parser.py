"""
PPTX / PPT 解析器

将 .pptx / .ppt 文件解析为结构化 chunks：
- 每个 Slide 视为一个"逻辑页"
- 标题占位符 → heading chunk
- 文本框 / 正文占位符 → paragraph chunk
- 表格 shape → table chunk (Markdown 表格)
- 图片 shape → image chunk (仅元信息)
- 备注 → note chunk

不复用 pptx_io.read()（TabSlide 渲染级解析，输出 SlideElement JSON，
依赖 image_handler 等外部回调，不适合纯文本提取场景）。
"""

from __future__ import annotations

import hashlib
import logging
import os
import zipfile

from .base import BaseDocumentParser, ChunkResult, PageResult, ParseResult
from .registry import register_parser

logger = logging.getLogger(__name__)

_TITLE_PLACEHOLDER_INDICES = {0, 1}  # idx 0/1 通常是标题或副标题
_MAX_TABLE_ROWS = 500
_MAX_TABLE_COLS = 50
_MAX_SLIDES = 500

MAX_PPTX_FILE_SIZE = 200 * 1024 * 1024  # 200 MB on-disk
MAX_PPTX_UNCOMPRESSED_SIZE = 2 * 1024 * 1024 * 1024  # 2 GB total decompressed
MAX_PPTX_COMPRESSION_RATIO = 100  # flag zip-bombs with ratio > 100:1
MAX_PPTX_ZIP_ENTRIES = 10_000


@register_parser
class PptxParser(BaseDocumentParser):

    def supported_mimes(self) -> list[str]:
        return [
            "application/vnd.openxmlformats-officedocument.presentationml.presentation",
        ]

    def parse(self, file_path: str, **kwargs) -> ParseResult:
        from pptx import Presentation
        from pptx.util import Emu

        if file_path.lower().endswith(".ppt"):
            raise ValueError(
                "不支持旧版 .ppt 格式，请将文件转换为 .pptx 后重新上传"
            )

        _validate_pptx_safe(file_path)

        prs = Presentation(file_path)
        pages: list[PageResult] = []
        doc_title = ""

        slide_width = prs.slide_width or Emu(12192000)
        slide_height = prs.slide_height or Emu(6858000)
        width_pt = slide_width.pt if hasattr(slide_width, "pt") else 0
        height_pt = slide_height.pt if hasattr(slide_height, "pt") else 0

        total_slides = len(prs.slides)
        if total_slides > _MAX_SLIDES:
            logger.warning(
                "PPTX 页数 (%d) 超过上限 (%d)，仅处理前 %d 页: %s",
                total_slides, _MAX_SLIDES, _MAX_SLIDES, file_path,
            )

        for slide_idx, slide in enumerate(prs.slides):
            if slide_idx >= _MAX_SLIDES:
                break
            page_num = slide_idx + 1
            chunks: list[ChunkResult] = []
            seq = 0

            for shape in slide.shapes:
                try:
                    new_chunks = _extract_shape(shape, seq)
                    for c in new_chunks:
                        seq += 1
                        c.sequence = seq
                    chunks.extend(new_chunks)
                except Exception as exc:
                    logger.debug("Slide %d shape 提取失败: %s", page_num, exc)

            notes_text = _extract_notes(slide)
            if notes_text:
                seq += 1
                chunks.append(ChunkResult(
                    chunk_type="note",
                    content=f"[演讲备注]\n{notes_text}",
                    sequence=seq,
                    metadata={"source": "structural", "role": "speaker_notes"},
                ))

            if not doc_title and chunks:
                for c in chunks:
                    if c.chunk_type == "heading" and c.content.strip():
                        doc_title = c.content.strip()[:200]
                        break

            text_content = "\n".join(c.content for c in chunks if c.content)
            pages.append(PageResult(
                page_number=page_num,
                width=width_pt,
                height=height_pt,
                chunks=chunks,
                text_content=text_content,
            ))

        if not pages:
            pages = [PageResult(page_number=1, width=0, height=0, chunks=[], text_content="")]

        return ParseResult(
            pages=pages,
            title=doc_title,
            parse_method="structural",
        )


def _extract_shape(shape, base_seq: int) -> list[ChunkResult]:
    """从单个 shape 提取 chunks，可能返回多个（如一个文本框有多段）。"""
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    chunks: list[ChunkResult] = []

    if shape.has_table:
        md = _table_to_markdown(shape.table)
        if md:
            chunks.append(ChunkResult(
                chunk_type="table", content=md, sequence=0,
                metadata={
                    "source": "structural",
                    "rows": len(shape.table.rows),
                    "cols": len(shape.table.columns),
                },
            ))
        return chunks

    if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
        img_info = _extract_image_info(shape)
        chunks.append(ChunkResult(
            chunk_type="image",
            content=f"[图片: {img_info.get('filename', '未知')}]",
            sequence=0,
            metadata={"source": "structural", **img_info},
        ))
        return chunks

    if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
        for child in shape.shapes:
            try:
                chunks.extend(_extract_shape(child, base_seq + len(chunks)))
            except Exception:
                pass
        return chunks

    if not shape.has_text_frame:
        return chunks

    is_title = _is_title_shape(shape)

    for para in shape.text_frame.paragraphs:
        text = para.text.strip()
        if not text:
            continue

        if is_title:
            level = para.level + 1 if para.level is not None else 1
            chunks.append(ChunkResult(
                chunk_type="heading", content=text, sequence=0,
                heading_level=min(level, 6),
                metadata={"source": "structural"},
            ))
        else:
            chunks.append(ChunkResult(
                chunk_type="paragraph", content=text, sequence=0,
                metadata={"source": "structural"},
            ))

    return chunks


def _is_title_shape(shape) -> bool:
    """判断 shape 是否为标题占位符。"""
    try:
        if hasattr(shape, "placeholder_format") and shape.placeholder_format:
            ph_idx = shape.placeholder_format.idx
            if ph_idx in _TITLE_PLACEHOLDER_INDICES:
                return True
            ph_type = shape.placeholder_format.type
            if ph_type is not None:
                from pptx.enum.shapes import PP_PLACEHOLDER
                if ph_type in (PP_PLACEHOLDER.TITLE, PP_PLACEHOLDER.CENTER_TITLE,
                               PP_PLACEHOLDER.SUBTITLE):
                    return True
    except Exception:
        pass
    return False


def _extract_image_info(shape) -> dict:
    """提取图片 shape 的元信息（不提取二进制内容）。"""
    info: dict = {}
    try:
        image = shape.image
        info["filename"] = image.filename or "image"
        info["content_type"] = image.content_type
        blob = image.blob
        if blob:
            info["hash"] = hashlib.md5(blob).hexdigest()[:12]
            info["size"] = len(blob)
    except Exception:
        info["filename"] = "image"
    return info


_NOTES_PLACEHOLDERS = {
    "单击此处添加备注",
    "Click to add notes",
    "クリックしてノートを追加",
    "Cliquez pour ajouter des notes",
    "Klicken Sie hier, um Notizen hinzuzufügen",
    "Haga clic para agregar notas",
    "클릭하여 노트 추가",
}


def _extract_notes(slide) -> str:
    """提取幻灯片备注文本。"""
    try:
        if slide.has_notes_slide and slide.notes_slide:
            tf = slide.notes_slide.notes_text_frame
            if tf:
                text = tf.text.strip()
                if text and text not in _NOTES_PLACEHOLDERS:
                    return text
    except Exception:
        pass
    return ""


def _table_to_markdown(table) -> str:
    """PPTX 表格 → Markdown 表格。"""
    rows_data: list[list[str]] = []
    for row_idx, row in enumerate(table.rows):
        if row_idx >= _MAX_TABLE_ROWS:
            break
        cells = []
        for col_idx, cell in enumerate(row.cells):
            if col_idx >= _MAX_TABLE_COLS:
                break
            cells.append(cell.text.replace("|", "\\|").replace("\n", " ").strip())
        rows_data.append(cells)

    if not rows_data:
        return ""

    max_cols = max(len(r) for r in rows_data)

    def pad(r: list[str]) -> list[str]:
        return r + [""] * (max_cols - len(r))

    header = pad(rows_data[0])
    lines = ["| " + " | ".join(header) + " |"]
    lines.append("| " + " | ".join("---" for _ in header) + " |")
    for row in rows_data[1:]:
        lines.append("| " + " | ".join(pad(row)) + " |")

    return "\n".join(lines)


def _validate_pptx_safe(file_path: str) -> None:
    """在加载 PPTX 前校验文件大小和 ZIP 解压体积，防止 zip-bomb。"""
    file_size = os.path.getsize(file_path)
    if file_size > MAX_PPTX_FILE_SIZE:
        raise ValueError(
            f"PPTX 文件过大（{file_size / 1024 / 1024:.1f} MB），"
            f"上限 {MAX_PPTX_FILE_SIZE // 1024 // 1024} MB"
        )

    if not zipfile.is_zipfile(file_path):
        raise ValueError("文件不是有效的 PPTX/ZIP 格式")

    with zipfile.ZipFile(file_path, "r") as zf:
        entries = zf.infolist()
        if len(entries) > MAX_PPTX_ZIP_ENTRIES:
            raise ValueError(
                f"PPTX 内部条目数过多（{len(entries)}），"
                f"上限 {MAX_PPTX_ZIP_ENTRIES}"
            )

        total_uncompressed = sum(info.file_size for info in entries)
        if total_uncompressed > MAX_PPTX_UNCOMPRESSED_SIZE:
            raise ValueError(
                f"PPTX 解压后体积过大（{total_uncompressed / 1024 / 1024:.1f} MB），"
                f"上限 {MAX_PPTX_UNCOMPRESSED_SIZE // 1024 // 1024} MB"
            )
        if file_size > 0 and total_uncompressed / file_size > MAX_PPTX_COMPRESSION_RATIO:
            raise ValueError(
                f"PPTX 压缩比异常（{total_uncompressed / file_size:.0f}:1），"
                f"疑似 zip-bomb，上限 {MAX_PPTX_COMPRESSION_RATIO}:1"
            )

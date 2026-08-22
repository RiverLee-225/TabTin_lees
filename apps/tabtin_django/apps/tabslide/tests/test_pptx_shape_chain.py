import importlib.util
from pathlib import Path
from unittest import TestCase

from lxml import etree
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE

# 直接按文件路径加载，避免触发 apps.tabslide.services.__init__ 的 Django 依赖
_MODULE_PATH = Path(__file__).resolve().parents[1] / "services" / "pptx_io.py"
_SPEC = importlib.util.spec_from_file_location("tabslide_pptx_shape_test_module", _MODULE_PATH)
if _SPEC is None or _SPEC.loader is None:
    raise RuntimeError(f"Failed to load module spec from {_MODULE_PATH}")
_PPTX_IO = importlib.util.module_from_spec(_SPEC)
_SPEC.loader.exec_module(_PPTX_IO)

_write_shape_element = _PPTX_IO._write_shape_element
_extract_auto_shape = _PPTX_IO._extract_auto_shape
_extract_shape_pattern = _PPTX_IO._extract_shape_pattern
_extract_shape_fill = _PPTX_IO._extract_shape_fill
_extract_shape_outline = _PPTX_IO._extract_shape_outline
_extract_shadow = _PPTX_IO._extract_shadow
_extract_opacity = _PPTX_IO._extract_opacity
_apply_common_write_props = _PPTX_IO._apply_common_write_props
_map_pptx_shape_to_mso = _PPTX_IO._map_pptx_shape_to_mso


_PNG_DATA_URL = (
    "data:image/png;base64,"
    "iVBORw0KGgoAAAANSUhEUgAAAAEAAAABCAQAAAC1HAwCAAAAC0lEQVR42mP8/x8AAwMCAO3ZfD0AAAAASUVORK5CYII="
)


class _DummyShape:
    def __init__(self, xml: str):
        self._element = etree.fromstring(xml.encode("utf-8"))


class TestPptxShapeChain(TestCase):
    def _new_slide(self):
        prs = Presentation()
        return prs.slides.add_slide(prs.slide_layouts[6])

    def _extract_shape_element(self, shape, shape_id: str = "shape-1"):
        base = {
            "id": shape_id,
            "x": 10,
            "y": 20,
            "width": 200,
            "height": 120,
            "zIndex": 0,
        }
        return _extract_auto_shape(shape, base)

    def test_shape_type_mapping_covers_extended_presets(self):
        mapping_cases = [
            ("round1Rect", "ROUND_1_RECTANGLE", MSO_SHAPE.ROUNDED_RECTANGLE),
            ("snip2DiagRect", "SNIP_2_DIAG_RECTANGLE", MSO_SHAPE.RECTANGLE),
            ("octagon", "OCTAGON", MSO_SHAPE.OCTAGON),
            ("star6", "STAR_6_POINT", MSO_SHAPE.STAR_5_POINT),
            ("callout1", "RECTANGULAR_CALLOUT", MSO_SHAPE.RECTANGLE),
            ("callout2", "ROUNDED_RECTANGULAR_CALLOUT", MSO_SHAPE.RECTANGLE),
            ("cross", "CROSS", MSO_SHAPE.CROSS),
        ]
        for shape_type, attr_name, fallback in mapping_cases:
            with self.subTest(shape_type=shape_type):
                expected = getattr(MSO_SHAPE, attr_name, fallback)
                self.assertEqual(_map_pptx_shape_to_mso(shape_type), expected)

    def test_roundrect_keypoints_roundtrip(self):
        slide = self._new_slide()
        shape = _write_shape_element(
            slide,
            {
                "pptxShapeType": "roundRect",
                "pathFormula": "roundRect",
                "keypoints": [0.24],
                "fill": "#00aa88",
            },
            left=100000,
            top=100000,
            width=2200000,
            height=1400000,
        )
        extracted = self._extract_shape_element(shape, "shape-roundrect-kp")
        props = extracted.get("props", {})
        self.assertEqual(props.get("pathFormula"), "roundRect")
        self.assertIn("keypoints", props)
        self.assertAlmostEqual(props.get("keypoints", [0])[0], 0.24, delta=0.03)

    def test_write_shape_without_outline_shadow_clears_powerpoint_theme_defaults(self):
        slide = self._new_slide()
        shape = _write_shape_element(
            slide,
            {
                "pptxShapeType": "roundRect",
                "pathFormula": "roundRect",
                "keypoints": [0.15, 0.15, 0.15, 0.15],
                "fill": "#FDF8F5",
            },
            left=100000,
            top=100000,
            width=2200000,
            height=600000,
        )

        ns = {
            "p": "http://schemas.openxmlformats.org/presentationml/2006/main",
            "a": "http://schemas.openxmlformats.org/drawingml/2006/main",
        }
        self.assertIsNone(shape._element.find("p:style", ns))

        sp_pr = shape._element.find("p:spPr", ns)
        self.assertIsNotNone(sp_pr)
        assert sp_pr is not None
        line = sp_pr.find("a:ln", ns)
        self.assertIsNotNone(line)
        assert line is not None
        self.assertIsNotNone(line.find("a:noFill", ns))
        self.assertIsNone(sp_pr.find("a:effectLst", ns))

    def test_write_shape_with_theme_outline_keeps_explicit_outline(self):
        slide = self._new_slide()
        shape = _write_shape_element(
            slide,
            {
                "pptxShapeType": "roundRect",
                "pathFormula": "roundRect",
                "keypoints": [0.15, 0.15, 0.15, 0.15],
                "fill": "#FDF8F5",
                "outline": {"width": 2, "color": "#4472C4", "themeKey": "accent1"},
            },
            left=100000,
            top=100000,
            width=2200000,
            height=600000,
        )

        ns = {
            "p": "http://schemas.openxmlformats.org/presentationml/2006/main",
            "a": "http://schemas.openxmlformats.org/drawingml/2006/main",
        }
        self.assertIsNone(shape._element.find("p:style", ns))

        line = shape._element.find("p:spPr/a:ln", ns)
        self.assertIsNotNone(line)
        assert line is not None
        self.assertIsNone(line.find("a:noFill", ns))
        scheme = line.find("a:solidFill/a:schemeClr", ns)
        self.assertIsNotNone(scheme)
        assert scheme is not None
        self.assertEqual(scheme.get("val"), "accent1")

    def test_roundrect_four_corner_keypoints_roundtrip_via_custgeom(self):
        slide = self._new_slide()
        shape = _write_shape_element(
            slide,
            {
                "pptxShapeType": "roundRect",
                "pathFormula": "roundRect",
                "keypoints": [0.26, 0.08, 0.18, 0.12],
                "fill": "#00aa88",
            },
            left=100000,
            top=100000,
            width=2200000,
            height=1400000,
        )

        ns = {
            "p": "http://schemas.openxmlformats.org/presentationml/2006/main",
            "a": "http://schemas.openxmlformats.org/drawingml/2006/main",
        }
        sp_pr = shape._element.find("p:spPr", ns)
        self.assertIsNotNone(sp_pr)
        assert sp_pr is not None
        self.assertIsNotNone(sp_pr.find("a:custGeom", ns))
        self.assertIsNone(sp_pr.find("a:prstGeom", ns))

        extracted = self._extract_shape_element(shape, "shape-roundrect-4corners")
        props = extracted.get("props", {})
        self.assertEqual(props.get("pathFormula"), "roundRect")
        out_kp = props.get("keypoints", [])
        self.assertGreaterEqual(len(out_kp), 4)
        self.assertAlmostEqual(out_kp[0], 0.26, delta=0.05)
        self.assertAlmostEqual(out_kp[1], 0.08, delta=0.05)
        self.assertAlmostEqual(out_kp[2], 0.18, delta=0.05)
        self.assertAlmostEqual(out_kp[3], 0.12, delta=0.05)

    def test_roundrect_uniform_four_corner_keypoints_use_prstgeom(self):
        slide = self._new_slide()
        shape = _write_shape_element(
            slide,
            {
                "pptxShapeType": "roundRect",
                "pathFormula": "roundRect",
                "keypoints": [0.22, 0.22, 0.22, 0.22],
                "fill": "#00aa88",
            },
            left=100000,
            top=100000,
            width=2200000,
            height=1400000,
        )

        ns = {
            "p": "http://schemas.openxmlformats.org/presentationml/2006/main",
            "a": "http://schemas.openxmlformats.org/drawingml/2006/main",
        }
        sp_pr = shape._element.find("p:spPr", ns)
        self.assertIsNotNone(sp_pr)
        assert sp_pr is not None
        prst = sp_pr.find("a:prstGeom", ns)
        self.assertIsNotNone(prst)
        assert prst is not None
        self.assertEqual(prst.get("prst"), "roundRect")
        gd = prst.find("a:avLst/a:gd", ns)
        self.assertIsNotNone(gd)
        assert gd is not None
        self.assertEqual(gd.get("name"), "adj")
        self.assertIn("val 22000", gd.get("fmla", ""))

        extracted = self._extract_shape_element(shape, "shape-roundrect-uniform-4corners")
        props = extracted.get("props", {})
        self.assertEqual(props.get("pathFormula"), "roundRect")
        out_kp = props.get("keypoints", [])
        self.assertGreaterEqual(len(out_kp), 4)
        self.assertAlmostEqual(out_kp[0], 0.22, delta=0.03)
        self.assertAlmostEqual(out_kp[1], 0.22, delta=0.03)
        self.assertAlmostEqual(out_kp[2], 0.22, delta=0.03)
        self.assertAlmostEqual(out_kp[3], 0.22, delta=0.03)

    def test_right_arrow_keypoints_roundtrip(self):
        slide = self._new_slide()
        shape = _write_shape_element(
            slide,
            {
                "pptxShapeType": "rightArrow",
                "pathFormula": "rightArrow",
                "keypoints": [0.62, 0.28],
                "fill": "#ffaa00",
            },
            left=100000,
            top=100000,
            width=2400000,
            height=1200000,
        )
        extracted = self._extract_shape_element(shape, "shape-right-arrow-kp")
        props = extracted.get("props", {})
        self.assertEqual(props.get("pathFormula"), "rightArrow")
        self.assertIn("keypoints", props)
        out_kp = props.get("keypoints", [])
        self.assertGreaterEqual(len(out_kp), 2)
        self.assertAlmostEqual(out_kp[0], 0.62, delta=0.05)
        self.assertAlmostEqual(out_kp[1], 0.28, delta=0.05)

    def test_left_right_arrow_keypoints_roundtrip_keeps_low_adj1(self):
        slide = self._new_slide()
        shape = _write_shape_element(
            slide,
            {
                "pptxShapeType": "leftRightArrow",
                "pathFormula": "leftRightArrow",
                "keypoints": [0.05, 0.45],
                "fill": "#ffaa00",
            },
            left=100000,
            top=100000,
            width=2400000,
            height=1200000,
        )
        extracted = self._extract_shape_element(shape, "shape-left-right-arrow-kp")
        props = extracted.get("props", {})
        self.assertEqual(props.get("pathFormula"), "leftRightArrow")
        out_kp = props.get("keypoints", [])
        self.assertGreaterEqual(len(out_kp), 2)
        self.assertAlmostEqual(out_kp[0], 0.05, delta=0.03)
        self.assertAlmostEqual(out_kp[1], 0.45, delta=0.03)

    def test_up_down_arrow_keypoints_roundtrip_keeps_low_adj1(self):
        slide = self._new_slide()
        shape = _write_shape_element(
            slide,
            {
                "pptxShapeType": "upDownArrow",
                "pathFormula": "upDownArrow",
                "keypoints": [0.04, 0.42],
                "fill": "#33aaee",
            },
            left=100000,
            top=100000,
            width=1400000,
            height=2400000,
        )
        extracted = self._extract_shape_element(shape, "shape-up-down-arrow-kp")
        props = extracted.get("props", {})
        self.assertEqual(props.get("pathFormula"), "upDownArrow")
        out_kp = props.get("keypoints", [])
        self.assertGreaterEqual(len(out_kp), 2)
        self.assertAlmostEqual(out_kp[0], 0.04, delta=0.03)
        self.assertAlmostEqual(out_kp[1], 0.42, delta=0.03)

    def test_pattern_fill_roundtrip_blipfill_and_priority(self):
        slide = self._new_slide()
        shape = _write_shape_element(
            slide,
            {
                "pptxShapeType": "rect",
                "pattern": _PNG_DATA_URL,
                "gradient": {
                    "type": "linear",
                    "rotate": 45,
                    "colors": [{"pos": 0, "color": "#ff0000"}, {"pos": 1, "color": "#0000ff"}],
                },
                "fill": "#00ff00",
            },
            left=100000,
            top=100000,
            width=2000000,
            height=1200000,
        )

        ns = {
            "p": "http://schemas.openxmlformats.org/presentationml/2006/main",
            "a": "http://schemas.openxmlformats.org/drawingml/2006/main",
        }
        sp_pr = shape._element.find("p:spPr", ns)
        self.assertIsNotNone(sp_pr)
        assert sp_pr is not None
        self.assertIsNotNone(sp_pr.find("a:blipFill", ns))
        self.assertIsNone(sp_pr.find("a:gradFill", ns))

        extracted = self._extract_shape_element(shape, "shape-pattern")
        pattern = extracted.get("props", {}).get("pattern")
        self.assertIsInstance(pattern, str)
        self.assertTrue(pattern.startswith("data:image/png;base64,"))

    def test_extract_pattfill_pattern_generates_data_url(self):
        shape = _DummyShape(
            """
            <p:sp xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"
                  xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">
              <p:spPr>
                <a:pattFill prst="cross">
                  <a:fgClr><a:srgbClr val="112233" /></a:fgClr>
                  <a:bgClr><a:srgbClr val="ffffff" /></a:bgClr>
                </a:pattFill>
              </p:spPr>
            </p:sp>
            """
        )
        src = _extract_shape_pattern(shape)
        self.assertIsInstance(src, str)
        assert isinstance(src, str)
        self.assertTrue(src.startswith("data:image/png;base64,"))
        self.assertGreater(len(src), 30)

    def test_radial_gradient_roundtrip(self):
        slide = self._new_slide()
        shape = _write_shape_element(
            slide,
            {
                "pptxShapeType": "ellipse",
                "gradient": {
                    "type": "radial",
                    "rotate": 0,
                    "colors": [
                        {"pos": 0, "color": "#ff0000"},
                        {"pos": 1, "color": "rgba(0,0,255,0.4)"},
                    ],
                },
            },
            left=100000,
            top=100000,
            width=2200000,
            height=1400000,
        )

        extracted = self._extract_shape_element(shape, "shape-radial")
        gradient = extracted.get("props", {}).get("gradient")
        self.assertIsNotNone(gradient)
        assert isinstance(gradient, dict)
        self.assertEqual(gradient.get("type"), "radial")
        self.assertGreaterEqual(len(gradient.get("colors", [])), 2)
        self.assertEqual(gradient["colors"][-1]["color"], "rgba(0,0,255,0.4)")

    def test_outline_roundtrip_supports_dotted_style(self):
        slide = self._new_slide()
        shape = _write_shape_element(
            slide,
            {
                "pptxShapeType": "rect",
                "fill": "#eeeeee",
                "outline": {
                    "color": "rgba(68,85,102,0.5)",
                    "width": 2.5,
                    "style": "dotted",
                },
            },
            left=100000,
            top=100000,
            width=2000000,
            height=1200000,
        )

        extracted = self._extract_shape_element(shape, "shape-outline")
        outline = extracted.get("props", {}).get("outline")
        self.assertIsNotNone(outline)
        assert isinstance(outline, dict)
        self.assertEqual(outline.get("style"), "dotted")
        self.assertGreaterEqual(float(outline.get("width", 0)), 2.0)
        self.assertIn("rgba(", str(outline.get("color", "")))

    def test_shape_fill_and_outline_alpha_roundtrip_keep_independent(self):
        slide = self._new_slide()
        shape = _write_shape_element(
            slide,
            {
                "pptxShapeType": "rect",
                "fill": "rgba(74,109,140,0.35)",
                "outline": {
                    "color": "rgba(74,109,140,0.8)",
                    "width": 1.8,
                    "style": "solid",
                },
            },
            left=100000,
            top=100000,
            width=2000000,
            height=1200000,
        )

        # 没有设置元素整体透明度时，不应回落为 base.opacity
        self.assertIsNone(_extract_opacity(shape))

        extracted = self._extract_shape_element(shape, "shape-fill-outline-alpha")
        props = extracted.get("props", {})
        self.assertEqual(props.get("fill"), "rgba(74,109,140,0.35)")
        outline = props.get("outline")
        self.assertIsNotNone(outline)
        assert isinstance(outline, dict)
        self.assertEqual(outline.get("color"), "rgba(74,109,140,0.8)")
        self.assertNotIn("opacity", extracted)

        ns = {
            "p": "http://schemas.openxmlformats.org/presentationml/2006/main",
            "a": "http://schemas.openxmlformats.org/drawingml/2006/main",
        }
        fill_alpha = shape._element.find("p:spPr/a:solidFill/a:srgbClr/a:alpha", ns)
        line_alpha = shape._element.find("p:spPr/a:ln/a:solidFill/a:srgbClr/a:alpha", ns)
        self.assertIsNotNone(fill_alpha)
        self.assertIsNotNone(line_alpha)
        assert fill_alpha is not None
        assert line_alpha is not None
        self.assertEqual(fill_alpha.get("val"), "35000")
        self.assertEqual(line_alpha.get("val"), "80000")

    def test_extract_shape_alpha_mod_fix_for_fill_and_outline(self):
        shape = _DummyShape(
            """
            <p:sp xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"
                  xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">
              <p:spPr>
                <a:solidFill>
                  <a:srgbClr val="4A6D8C">
                    <a:alphaModFix amt="35000"/>
                  </a:srgbClr>
                </a:solidFill>
                <a:ln w="12700">
                  <a:solidFill>
                    <a:srgbClr val="4A6D8C">
                      <a:alphaModFix amt="80000"/>
                    </a:srgbClr>
                  </a:solidFill>
                </a:ln>
              </p:spPr>
            </p:sp>
            """
        )

        self.assertIsNone(_extract_opacity(shape))
        self.assertEqual(_extract_shape_fill(shape), "rgba(74,109,140,0.35)")
        outline = _extract_shape_outline(shape)
        self.assertIsNotNone(outline)
        assert isinstance(outline, dict)
        self.assertEqual(outline.get("color"), "rgba(74,109,140,0.8)")

    def test_shape_text_roundtrip_keeps_content_and_align(self):
        slide = self._new_slide()
        shape = _write_shape_element(
            slide,
            {
                "pptxShapeType": "roundRect",
                "fill": "#ffffff",
                "text": {
                    "content": "<p><strong>Hello</strong> TabSlide</p>",
                    "align": "right",
                    "verticalAlign": "bottom",
                    "defaultColor": "#112233",
                    "defaultFontSize": 18,
                    "defaultFontFamily": "Arial",
                },
            },
            left=100000,
            top=100000,
            width=2600000,
            height=1600000,
        )

        extracted = self._extract_shape_element(shape, "shape-text")
        text_props = extracted.get("props", {}).get("text")
        self.assertIsNotNone(text_props)
        assert isinstance(text_props, dict)
        self.assertIn("Hello", text_props.get("content", ""))
        self.assertEqual(text_props.get("align"), "right")
        self.assertEqual(text_props.get("verticalAlign"), "bottom")

    def test_apply_common_write_props_writes_flip_and_shadow(self):
        slide = self._new_slide()
        shape = _write_shape_element(
            slide,
            {"pptxShapeType": "rect", "fill": "#ffffff"},
            left=100000,
            top=100000,
            width=2000000,
            height=1200000,
        )

        _apply_common_write_props(
            shape,
            {
                "id": "shape-shadow",
                "type": "shape",
                "props": {
                    "flipH": True,
                    "flipV": True,
                    "shadow": {
                        "h": 3,
                        "v": 4,
                        "blur": 5,
                        "color": "rgba(0,0,0,0.25)",
                    },
                },
            },
        )

        ns = {
            "p": "http://schemas.openxmlformats.org/presentationml/2006/main",
            "a": "http://schemas.openxmlformats.org/drawingml/2006/main",
        }
        xfrm = shape._element.find("p:spPr/a:xfrm", ns)
        self.assertIsNotNone(xfrm)
        assert xfrm is not None
        self.assertEqual(xfrm.get("flipH"), "1")
        self.assertEqual(xfrm.get("flipV"), "1")

        shadow = _extract_shadow(shape)
        self.assertIsNotNone(shadow)
        assert isinstance(shadow, dict)
        self.assertEqual(shadow.get("color"), "#000000")
        self.assertEqual(shadow.get("opacity"), 0.25)

    def test_apply_common_write_props_keeps_gradient_opacity(self):
        slide = self._new_slide()
        shape = _write_shape_element(
            slide,
            {
                "pptxShapeType": "ellipse",
                "gradient": {
                    "type": "linear",
                    "rotate": 60,
                    "colors": [
                        {"pos": 0, "color": "#ff0000"},
                        {"pos": 1, "color": "#0000ff"},
                    ],
                },
            },
            left=100000,
            top=100000,
            width=2200000,
            height=1400000,
        )

        _apply_common_write_props(
            shape,
            {
                "id": "shape-grad-opacity",
                "type": "shape",
                "opacity": 0.42,
                "props": {},
            },
        )

        ns = {
            "p": "http://schemas.openxmlformats.org/presentationml/2006/main",
            "a": "http://schemas.openxmlformats.org/drawingml/2006/main",
        }
        alpha_mod = shape._element.find("p:spPr/a:alphaModFix", ns)
        self.assertIsNotNone(alpha_mod)
        assert alpha_mod is not None
        self.assertEqual(alpha_mod.get("val"), "42000")

        opacity = _extract_opacity(shape)
        self.assertAlmostEqual(opacity or 0, 0.42, places=2)

    def test_cross_keypoint_roundtrip_clamps_to_formula_min(self):
        slide = self._new_slide()
        shape = _write_shape_element(
            slide,
            {
                "pptxShapeType": "plus",
                "pathFormula": "cross",
                "keypoints": [0.01],
                "fill": "#00aa88",
            },
            left=100000,
            top=100000,
            width=2200000,
            height=1400000,
        )
        extracted = self._extract_shape_element(shape, "shape-cross-kp")
        props = extracted.get("props", {})
        self.assertEqual(props.get("pathFormula"), "cross")
        out_kp = props.get("keypoints", [])
        self.assertGreaterEqual(len(out_kp), 1)
        self.assertAlmostEqual(out_kp[0], 0.1, delta=0.03)

    def test_chevron_keypoint_roundtrip_clamps_to_formula_min(self):
        slide = self._new_slide()
        shape = _write_shape_element(
            slide,
            {
                "pptxShapeType": "chevron",
                "pathFormula": "chevron",
                "keypoints": [0.01],
                "fill": "#00aa88",
            },
            left=100000,
            top=100000,
            width=2200000,
            height=1400000,
        )
        extracted = self._extract_shape_element(shape, "shape-chevron-kp")
        props = extracted.get("props", {})
        self.assertEqual(props.get("pathFormula"), "chevron")
        out_kp = props.get("keypoints", [])
        self.assertGreaterEqual(len(out_kp), 1)
        self.assertAlmostEqual(out_kp[0], 0.05, delta=0.03)

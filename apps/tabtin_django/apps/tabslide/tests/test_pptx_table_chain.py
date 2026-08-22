import importlib.util
from pathlib import Path
from unittest import TestCase

from lxml import etree
from pptx import Presentation
from pptx.oxml.ns import qn
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Emu

# 直接按文件路径加载，避免触发 apps.tabslide.services.__init__ 的 Django 依赖
_MODULE_PATH = Path(__file__).resolve().parents[1] / "services" / "pptx_io.py"
_SPEC = importlib.util.spec_from_file_location("tabslide_pptx_io_table_test_module", _MODULE_PATH)
if _SPEC is None or _SPEC.loader is None:
    raise RuntimeError(f"Failed to load module spec from {_MODULE_PATH}")
_PPTX_IO = importlib.util.module_from_spec(_SPEC)
_SPEC.loader.exec_module(_PPTX_IO)

_extract_table_outline = _PPTX_IO._extract_table_outline
_extract_table_borders = _PPTX_IO._extract_table_borders
_extract_table_theme = _PPTX_IO._extract_table_theme
_extract_table_shape = _PPTX_IO._extract_table_shape
_write_table_element = _PPTX_IO._write_table_element
EMU_PER_PT = _PPTX_IO.EMU_PER_PT
EMU_PER_PX = _PPTX_IO.EMU_PER_PX


class _DummyShape:
    def __init__(self, xml: str):
        self._element = etree.fromstring(xml.encode("utf-8"))


class TestPptxTableChain(TestCase):
    def test_extract_table_outline_skips_nofill_side(self):
        shape = _DummyShape(
            """
            <p:graphicFrame xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"
                            xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">
              <a:graphic>
                <a:graphicData>
                  <a:tbl>
                    <a:tblPr>
                      <a:tblBorders>
                        <a:top><a:noFill/></a:top>
                        <a:right w="12700" prstDash="dash">
                          <a:solidFill><a:srgbClr val="FF0000"/></a:solidFill>
                        </a:right>
                      </a:tblBorders>
                    </a:tblPr>
                  </a:tbl>
                </a:graphicData>
              </a:graphic>
            </p:graphicFrame>
            """
        )

        outline = _extract_table_outline(shape)
        self.assertIsNotNone(outline)
        self.assertEqual(outline["style"], "dashed")
        self.assertEqual(outline["width"], 1.0)
        self.assertEqual(outline["color"], "#FF0000")

    def test_extract_table_outline_reads_alpha(self):
        shape = _DummyShape(
            """
            <p:graphicFrame xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"
                            xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">
              <a:graphic>
                <a:graphicData>
                  <a:tbl>
                    <a:tblPr>
                      <a:tblBorders>
                        <a:top w="12700">
                          <a:solidFill>
                            <a:srgbClr val="112233">
                              <a:alpha val="40000"/>
                            </a:srgbClr>
                          </a:solidFill>
                        </a:top>
                      </a:tblBorders>
                    </a:tblPr>
                  </a:tbl>
                </a:graphicData>
              </a:graphic>
            </p:graphicFrame>
            """
        )

        outline = _extract_table_outline(shape)
        self.assertIsNotNone(outline)
        self.assertEqual(outline["color"], "rgba(17,34,51,0.4)")

    def test_extract_table_borders_reads_side_specs(self):
        shape = _DummyShape(
            """
            <p:graphicFrame xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"
                            xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">
              <a:graphic>
                <a:graphicData>
                  <a:tbl>
                    <a:tblPr>
                      <a:tblBorders>
                        <a:top w="25400" prstDash="dot">
                          <a:solidFill><a:srgbClr val="336699"/></a:solidFill>
                        </a:top>
                        <a:insideH>
                          <a:noFill/>
                        </a:insideH>
                      </a:tblBorders>
                    </a:tblPr>
                  </a:tbl>
                </a:graphicData>
              </a:graphic>
            </p:graphicFrame>
            """
        )

        borders = _extract_table_borders(shape)
        self.assertIsNotNone(borders)
        self.assertEqual(borders["top"]["style"], "dotted")
        self.assertEqual(borders["top"]["width"], 2.0)
        self.assertEqual(borders["top"]["color"], "#336699")
        self.assertEqual(borders["insideH"]["width"], 0)

    def test_extract_table_theme_bool_attrs_case_insensitive(self):
        shape = _DummyShape(
            """
            <p:graphicFrame xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"
                            xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">
              <a:graphic>
                <a:graphicData>
                  <a:tbl>
                    <a:tblPr firstRow="True" bandRow="ON" />
                  </a:tbl>
                </a:graphicData>
              </a:graphic>
            </p:graphicFrame>
            """
        )

        theme = _extract_table_theme(shape)
        self.assertIsNotNone(theme)
        self.assertTrue(theme.get("headerRow"))
        self.assertTrue(theme.get("stripedRows"))
        self.assertEqual(theme.get("color"), "#5b9bd5")

    def test_write_table_element_clears_default_theme_and_writes_zero_outline(self):
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        shape = _write_table_element(
            slide=slide,
            props={
                "data": [[{"text": "A1", "colspan": 1, "rowspan": 1}]],
                "outline": {"style": "solid", "width": 0, "color": "#333333"},
            },
            left=Inches(1),
            top=Inches(1),
            width=Inches(4),
            height=Inches(1),
        )
        self.assertIsNotNone(shape)

        tbl_pr = shape.table._tbl.find(qn("a:tblPr"))
        self.assertIsNotNone(tbl_pr)
        # 默认 add_table 会写 firstRow/bandRow；链路应按前端 theme（此处无 theme）清理
        self.assertIsNone(tbl_pr.get("firstRow"))
        self.assertIsNone(tbl_pr.get("bandRow"))

        borders = tbl_pr.find(qn("a:tblBorders"))
        self.assertIsNotNone(borders)
        top_border = borders.find(qn("a:top"))
        self.assertIsNotNone(top_border)
        self.assertIsNotNone(top_border.find(qn("a:noFill")))

    def test_write_table_element_normalizes_col_widths_with_missing_values(self):
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        shape = _write_table_element(
            slide=slide,
            props={
                "data": [
                    [{"text": "A1", "colspan": 1, "rowspan": 1}, {"text": "B1", "colspan": 1, "rowspan": 1}, {"text": "C1", "colspan": 1, "rowspan": 1}],
                ],
                "colWidths": [2, 1],  # 第三列缺失，应该自动补默认权重而不是塌陷
            },
            left=Inches(1),
            top=Inches(1),
            width=Inches(6),
            height=Inches(1),
        )
        self.assertIsNotNone(shape)

        widths = [shape.table.columns[i].width for i in range(3)]
        self.assertGreater(widths[0], widths[1])
        self.assertGreater(widths[1], 1)
        self.assertGreater(widths[2], 1)
        self.assertAlmostEqual(widths[1], widths[2], delta=2000)
        self.assertEqual(sum(widths), int(Inches(6)))

    def test_write_table_element_rich_text_fallback_align_and_color_alpha(self):
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        shape = _write_table_element(
            slide=slide,
            props={
                "data": [[{
                    "text": "A1",
                    "richText": '<p><span style="color: rgba(255, 0, 0, 0.4)">A1</span></p>',
                    "colspan": 1,
                    "rowspan": 1,
                    "style": {"align": "center"},
                }]],
            },
            left=Inches(1),
            top=Inches(1),
            width=Inches(4),
            height=Inches(1),
        )
        self.assertIsNotNone(shape)

        cell = shape.table.cell(0, 0)
        p = cell.text_frame.paragraphs[0]
        self.assertEqual(p.alignment, PP_ALIGN.CENTER)
        run = p.runs[0]
        self.assertEqual(str(run.font.color.rgb), "FF0000")

        nsmap_a = "http://schemas.openxmlformats.org/drawingml/2006/main"
        alpha = run._r.find(f".//{{{nsmap_a}}}rPr/{{{nsmap_a}}}solidFill/{{{nsmap_a}}}srgbClr/{{{nsmap_a}}}alpha")
        self.assertIsNotNone(alpha)
        self.assertEqual(alpha.get("val"), "40000")

    def test_write_table_element_accepts_string_spans_and_cell_min_height(self):
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        shape = _write_table_element(
            slide=slide,
            props={
                "data": [[{
                    "text": "A1",
                    "colspan": "1",
                    "rowspan": "1",
                }]],
                "cellMinHeight": "40",
            },
            left=Inches(1),
            top=Inches(1),
            width=Inches(4),
            height=Inches(1),
        )
        self.assertIsNotNone(shape)
        self.assertEqual(shape.table.rows[0].height, int(40 * EMU_PER_PX))

    def test_write_table_element_accepts_string_outline_width(self):
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        shape = _write_table_element(
            slide=slide,
            props={
                "data": [[{"text": "A1", "colspan": 1, "rowspan": 1}]],
                "outline": {"style": "dashed", "width": "2.5", "color": "#123456"},
            },
            left=Inches(1),
            top=Inches(1),
            width=Inches(4),
            height=Inches(1),
        )
        self.assertIsNotNone(shape)

        tbl_pr = shape.table._tbl.find(qn("a:tblPr"))
        self.assertIsNotNone(tbl_pr)
        borders = tbl_pr.find(qn("a:tblBorders"))
        self.assertIsNotNone(borders)
        top_border = borders.find(qn("a:top"))
        self.assertIsNotNone(top_border)
        self.assertEqual(top_border.get("w"), str(int(2.5 * EMU_PER_PT)))
        self.assertEqual(top_border.get("prstDash"), "dash")

        solid_fill = top_border.find(qn("a:solidFill"))
        self.assertIsNotNone(solid_fill)
        srgb = solid_fill.find(qn("a:srgbClr"))
        self.assertIsNotNone(srgb)
        self.assertEqual(srgb.get("val"), "123456")

    def test_write_table_element_writes_per_side_borders(self):
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        shape = _write_table_element(
            slide=slide,
            props={
                "data": [[{"text": "A1", "colspan": 1, "rowspan": 1}]],
                "outline": {"style": "solid", "width": 1, "color": "#000000"},
                "borders": {
                    "top": {"style": "dashed", "width": 2, "color": "#FF0000"},
                    "insideH": {"style": "solid", "width": 0, "color": "#00FF00"},
                },
            },
            left=Inches(1),
            top=Inches(1),
            width=Inches(4),
            height=Inches(1),
        )
        self.assertIsNotNone(shape)

        tbl_pr = shape.table._tbl.find(qn("a:tblPr"))
        self.assertIsNotNone(tbl_pr)
        borders = tbl_pr.find(qn("a:tblBorders"))
        self.assertIsNotNone(borders)

        top_border = borders.find(qn("a:top"))
        self.assertIsNotNone(top_border)
        self.assertEqual(top_border.get("w"), str(int(2 * EMU_PER_PT)))
        self.assertEqual(top_border.get("prstDash"), "dash")
        top_srgb = top_border.find(qn("a:solidFill")).find(qn("a:srgbClr"))
        self.assertEqual(top_srgb.get("val"), "FF0000")

        inside_h = borders.find(qn("a:insideH"))
        self.assertIsNotNone(inside_h)
        self.assertIsNotNone(inside_h.find(qn("a:noFill")))

    def test_extract_table_shape_preserves_row_heights(self):
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        shape = slide.shapes.add_table(2, 2, Inches(1), Inches(1), Inches(4), Inches(2))
        table = shape.table
        table.rows[0].height = Emu(30 * EMU_PER_PX)
        table.rows[1].height = Emu(50 * EMU_PER_PX)

        extracted = _extract_table_shape(shape, base={})
        props = extracted.get("props", {})
        self.assertEqual(props.get("rowHeights"), [30, 50])
        self.assertEqual(props.get("cellMinHeight"), 30)

    def test_write_table_element_row_heights_prefer_row_heights_over_cell_min_height(self):
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        shape = _write_table_element(
            slide=slide,
            props={
                "data": [
                    [{"text": "A1", "colspan": 1, "rowspan": 1}],
                    [{"text": "A2", "colspan": 1, "rowspan": 1}],
                ],
                "rowHeights": [100, 50],
                "cellMinHeight": 20,
            },
            left=Inches(1),
            top=Inches(1),
            width=Inches(4),
            height=Inches(2),
        )
        self.assertIsNotNone(shape)

        h0 = shape.table.rows[0].height
        h1 = shape.table.rows[1].height
        self.assertEqual(h0 + h1, int(Inches(2)))
        self.assertGreater(h0, h1)
        self.assertAlmostEqual(h0 / h1, 2.0, delta=0.15)


class TestB605RowHeightUnitConversion(TestCase):
    """[B6-05] 回归测试：cellMinHeight/rowHeights 使用 px 单位（EMU_PER_PX=9525），
    而非错误的 pt 单位（EMU_PER_PT=12700），确保行高读写链路无约 33% 偏差。"""

    def test_emu_per_px_constant_is_correct(self):
        self.assertEqual(EMU_PER_PX, 9525)
        self.assertEqual(EMU_PER_PX, 914400 // 96)

    def test_extract_row_height_uses_px_not_pt(self):
        """验证 _extract_table_shape 将 EMU 行高正确转换为 px。"""
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        shape = slide.shapes.add_table(1, 1, Inches(1), Inches(1), Inches(4), Inches(1))

        expected_px = 48
        shape.table.rows[0].height = Emu(expected_px * EMU_PER_PX)

        extracted = _extract_table_shape(shape, base={})
        props = extracted.get("props", {})
        row_heights = props.get("rowHeights", [])
        self.assertEqual(len(row_heights), 1)
        self.assertEqual(row_heights[0], expected_px)

    def test_write_cell_min_height_uses_px(self):
        """验证 _write_table_element 将 cellMinHeight (px) 正确转换为 EMU。"""
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        cell_min_px = 36
        shape = _write_table_element(
            slide=slide,
            props={
                "data": [[{"text": "A", "colspan": 1, "rowspan": 1}]],
                "cellMinHeight": cell_min_px,
            },
            left=Inches(1),
            top=Inches(1),
            width=Inches(4),
            height=Inches(1),
        )
        expected_emu = cell_min_px * EMU_PER_PX
        self.assertEqual(shape.table.rows[0].height, expected_emu)

    def test_row_height_read_write_roundtrip(self):
        """验证行高完整往返链路：write(px→EMU) → read(EMU→px) 无损。"""
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        input_heights_px = [40, 60]
        total_height_emu = sum(h * EMU_PER_PX for h in input_heights_px)
        shape = _write_table_element(
            slide=slide,
            props={
                "data": [
                    [{"text": "R0", "colspan": 1, "rowspan": 1}],
                    [{"text": "R1", "colspan": 1, "rowspan": 1}],
                ],
                "rowHeights": input_heights_px,
                "cellMinHeight": 30,
            },
            left=Inches(1),
            top=Inches(1),
            width=Inches(4),
            height=Emu(total_height_emu),
        )
        extracted = _extract_table_shape(shape, base={})
        props = extracted.get("props", {})
        output_heights = props.get("rowHeights", [])
        self.assertEqual(output_heights, input_heights_px)

    def test_wrong_emu_per_pt_would_cause_33pct_deviation(self):
        """显式验证：若仍使用 EMU_PER_PT(12700)，行高将偏小约 25%。"""
        row_emu = 48 * EMU_PER_PX  # 48px = 457200 EMU
        correct_px = round(row_emu / EMU_PER_PX)
        wrong_pt = round(row_emu / EMU_PER_PT)

        self.assertEqual(correct_px, 48)
        self.assertEqual(wrong_pt, 36)
        self.assertGreater(correct_px, wrong_pt)

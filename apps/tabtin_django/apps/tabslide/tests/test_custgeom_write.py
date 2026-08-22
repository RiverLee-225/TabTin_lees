"""
custGeom 写出回归：当 PPTShapeElement 带 props.path + props.viewBox 时，
pptx_io 写出 OOXML 应输出 <a:custGeom> 真路径而不是降级 prstGeom rect。

覆盖：
  1. SVG path 命令解析（_iter_svg_path_commands）：M / L / 隐式 lineto / Q / C / Z
  2. _apply_svg_path_custom_geometry 写出矩形 path → 含 moveTo + lnTo + close
  3. 写出含 Q 的圆角路径 → 含 quadBezTo
  4. 写出含 C 的复杂路径 → 含 cubicBezTo
  5. 不支持的命令（A 弧）→ 返回 False，不污染 spPr
  6. 端到端 _write_shape_element：
     - pptxShapeType="rect"（命中预设清单）→ 保留 prstGeom
     - 客户端未给 pptxShapeType + 有 path → 走 custGeom
     - pptxShapeType="custom"（不在清单）→ 走 custGeom
"""

from __future__ import annotations

import importlib.util
from pathlib import Path
from unittest import TestCase

from lxml import etree
from pptx import Presentation
from pptx.util import Emu


_MODULE_PATH = Path(__file__).resolve().parents[1] / "services" / "pptx_io.py"
_SPEC = importlib.util.spec_from_file_location("tabslide_pptx_io_custgeom_test", _MODULE_PATH)
if _SPEC is None or _SPEC.loader is None:
    raise RuntimeError(f"Failed to load module spec from {_MODULE_PATH}")
_PPTX_IO = importlib.util.module_from_spec(_SPEC)
_SPEC.loader.exec_module(_PPTX_IO)

_iter_svg_path_commands = _PPTX_IO._iter_svg_path_commands
_apply_svg_path_custom_geometry = _PPTX_IO._apply_svg_path_custom_geometry
_write_shape_element = _PPTX_IO._write_shape_element
_PRESET_GEOM_PREFERRED_TYPES = _PPTX_IO._PRESET_GEOM_PREFERRED_TYPES

NSMAP_A = "http://schemas.openxmlformats.org/drawingml/2006/main"


def _make_blank_slide():
    """构造一张白板 slide，返回 (presentation, slide)。"""
    prs = Presentation()
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)
    return prs, slide


def _get_shape_spPr_xml(shape) -> str:
    """返回 shape 的 spPr 节点 XML 字符串（人眼可读）。"""
    nsmap_p = "http://schemas.openxmlformats.org/presentationml/2006/main"
    sp_pr = shape._element.find(f"{{{nsmap_p}}}spPr")
    if sp_pr is None:
        sp_pr = shape._element.find(f".//{{{NSMAP_A}}}spPr")
    if sp_pr is None:
        return ""
    return etree.tostring(sp_pr, pretty_print=True).decode("utf-8")


# ───────────────────────────────────────────────────────────────
# 1. _iter_svg_path_commands 纯函数行为
# ───────────────────────────────────────────────────────────────


class TestIterSvgPathCommands(TestCase):
    def test_basic_m_l_z(self):
        cmds = list(_iter_svg_path_commands("M 10 20 L 30 40 Z"))
        self.assertEqual(cmds, [("M", [10.0, 20.0]), ("L", [30.0, 40.0]), ("Z", [])])

    def test_implicit_lineto_after_moveto(self):
        # SVG 规范：M 之后多组坐标应当作隐式 L
        cmds = list(_iter_svg_path_commands("M 0 0 100 0 100 100 0 100 Z"))
        kinds = [c[0] for c in cmds]
        self.assertEqual(kinds, ["M", "L", "L", "L", "Z"])
        self.assertEqual(cmds[0][1], [0.0, 0.0])
        self.assertEqual(cmds[1][1], [100.0, 0.0])
        self.assertEqual(cmds[3][1], [0.0, 100.0])

    def test_comma_and_negative_numbers(self):
        cmds = list(_iter_svg_path_commands("M-10,-20 L+5.5,-3e1"))
        self.assertEqual(cmds[0], ("M", [-10.0, -20.0]))
        self.assertEqual(cmds[1], ("L", [5.5, -30.0]))

    def test_quad_and_cubic(self):
        path = "M 0 0 Q 10 10 20 0 C 30 -10 40 10 50 0"
        cmds = list(_iter_svg_path_commands(path))
        kinds = [c[0] for c in cmds]
        self.assertEqual(kinds, ["M", "Q", "C"])
        self.assertEqual(cmds[1][1], [10.0, 10.0, 20.0, 0.0])
        self.assertEqual(cmds[2][1], [30.0, -10.0, 40.0, 10.0, 50.0, 0.0])

    def test_lowercase_relative(self):
        cmds = list(_iter_svg_path_commands("m 5 5 l 10 0 z"))
        kinds = [c[0] for c in cmds]
        self.assertEqual(kinds, ["m", "l", "z"])


# ───────────────────────────────────────────────────────────────
# 2. _apply_svg_path_custom_geometry：写出 custGeom 节点
# ───────────────────────────────────────────────────────────────


class TestApplySvgPathCustomGeometry(TestCase):
    def _add_shape(self, width=Emu(1000000), height=Emu(1000000)):
        from pptx.enum.shapes import MSO_SHAPE
        _, slide = _make_blank_slide()
        return slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Emu(100000), Emu(100000), width, height
        )

    def test_rect_path_writes_moveto_lnto_close(self):
        shape = self._add_shape()
        path_str = "M 0 0 L 100 0 L 100 100 L 0 100 Z"
        ok = _apply_svg_path_custom_geometry(shape, path_str, [100, 100])
        self.assertTrue(ok)

        xml = _get_shape_spPr_xml(shape)
        # prstGeom 应该被替换为 custGeom
        self.assertNotIn("prstGeom", xml)
        self.assertIn("custGeom", xml)
        self.assertIn("moveTo", xml)
        self.assertIn("lnTo", xml)
        self.assertIn("<a:close", xml)

        # 验证 path w/h 等于 shape 的 EMU 尺寸
        from pptx.oxml.ns import qn
        sp_pr = _PPTX_IO._find_sp_pr(shape._element)
        path_el = sp_pr.find(f".//{{{NSMAP_A}}}path")
        self.assertEqual(path_el.get("w"), str(int(shape.width)))
        self.assertEqual(path_el.get("h"), str(int(shape.height)))

    def test_round_rect_path_writes_quadbezto(self):
        shape = self._add_shape()
        # 典型 PPTist 圆角矩形 path
        path_str = (
            "M 10 0 L 90 0 Q 100 0 100 10 "
            "L 100 90 Q 100 100 90 100 "
            "L 10 100 Q 0 100 0 90 "
            "L 0 10 Q 0 0 10 0 Z"
        )
        ok = _apply_svg_path_custom_geometry(shape, path_str, [100, 100])
        self.assertTrue(ok)

        xml = _get_shape_spPr_xml(shape)
        self.assertIn("custGeom", xml)
        self.assertIn("quadBezTo", xml)
        # 至少 4 个 quadBezTo（四个圆角）
        self.assertGreaterEqual(xml.count("quadBezTo"), 4 * 2)  # 开+关 = 8

    def test_complex_cubic_path_writes_cubicbezto(self):
        shape = self._add_shape()
        # 模拟 L 形角标 cubicBezTo
        path_str = (
            "M 114300 0 L 495300 0 "
            "C 558384 0 609600 51216 609600 114300 "
            "L 609600 495300 "
            "C 609600 558384 558384 609600 495300 609600 "
            "L 114300 609600 "
            "C 51216 609600 0 558384 0 495300 "
            "L 0 114300 "
            "C 0 51216 51216 0 114300 0 Z"
        )
        ok = _apply_svg_path_custom_geometry(shape, path_str, [609600, 609600])
        self.assertTrue(ok)

        xml = _get_shape_spPr_xml(shape)
        self.assertIn("custGeom", xml)
        self.assertIn("cubicBezTo", xml)
        # 4 个 cubicBezTo 对应 4 个圆角
        self.assertGreaterEqual(xml.count("<a:cubicBezTo"), 4)

    def test_unsupported_arc_command_returns_false(self):
        shape = self._add_shape()
        # A 命令暂不支持
        path_str = "M 0 0 A 50 50 0 0 1 100 100 Z"
        ok = _apply_svg_path_custom_geometry(shape, path_str, [100, 100])
        self.assertFalse(ok)
        # spPr 不应该被污染——原来的 prstGeom 还在
        xml = _get_shape_spPr_xml(shape)
        self.assertNotIn("custGeom", xml)
        self.assertIn("prstGeom", xml)

    def test_invalid_viewbox_returns_false(self):
        shape = self._add_shape()
        self.assertFalse(
            _apply_svg_path_custom_geometry(shape, "M 0 0 L 100 100", [0, 100])
        )
        self.assertFalse(
            _apply_svg_path_custom_geometry(shape, "M 0 0 L 100 100", None)
        )
        self.assertFalse(
            _apply_svg_path_custom_geometry(shape, "M 0 0 L 100 100", [100])
        )

    def test_empty_path_returns_false(self):
        shape = self._add_shape()
        self.assertFalse(_apply_svg_path_custom_geometry(shape, "", [100, 100]))
        self.assertFalse(_apply_svg_path_custom_geometry(shape, "   ", [100, 100]))

    def test_coordinates_scale_to_emu(self):
        # viewBox 100×100，shape 1,000,000×500,000 EMU
        # SVG 坐标 (50, 50) 应缩放到 (500_000, 250_000)
        shape = self._add_shape(width=Emu(1_000_000), height=Emu(500_000))
        ok = _apply_svg_path_custom_geometry(
            shape, "M 0 0 L 50 50 L 100 100 Z", [100, 100]
        )
        self.assertTrue(ok)

        sp_pr = _PPTX_IO._find_sp_pr(shape._element)
        pts = sp_pr.findall(f".//{{{NSMAP_A}}}pt")
        # 第二个 pt 是 lnTo (50, 50)
        self.assertEqual(pts[1].get("x"), "500000")
        self.assertEqual(pts[1].get("y"), "250000")
        # 第三个 pt 是 lnTo (100, 100) = (1_000_000, 500_000)
        self.assertEqual(pts[2].get("x"), "1000000")
        self.assertEqual(pts[2].get("y"), "500000")


# ───────────────────────────────────────────────────────────────
# 3. _write_shape_element 的分发：simple prstGeom vs custom path
# ───────────────────────────────────────────────────────────────


class TestWriteShapeElementDispatch(TestCase):
    def test_pptx_shape_type_rect_keeps_prstgeom(self):
        """pptxShapeType=rect 在 _PRESET_GEOM_PREFERRED_TYPES 里，
        即使带了 path 也应保留 prstGeom（保留 PowerPoint 原生可编辑性）。"""
        _, slide = _make_blank_slide()
        props = {
            "pptxShapeType": "rect",
            "path": "M 0 0 L 200 0 L 200 200 L 0 200 Z",
            "viewBox": [200, 200],
            "fill": "#FF0000",
        }
        shape = _write_shape_element(
            slide, props, Emu(0), Emu(0), Emu(914400), Emu(914400)
        )
        xml = _get_shape_spPr_xml(shape)
        self.assertIn("prstGeom", xml)
        self.assertNotIn("custGeom", xml)

    def test_no_pptx_shape_type_with_path_uses_custgeom(self):
        """没有 pptxShapeType 但有复杂 path → 走 custGeom。"""
        _, slide = _make_blank_slide()
        props = {
            # 没有 pptxShapeType
            "path": "M 0 0 Q 50 50 100 0 L 100 100 L 0 100 Z",
            "viewBox": [100, 100],
            "fill": "#00FF00",
        }
        shape = _write_shape_element(
            slide, props, Emu(0), Emu(0), Emu(914400), Emu(914400)
        )
        xml = _get_shape_spPr_xml(shape)
        self.assertIn("custGeom", xml)
        self.assertNotIn("prstGeom", xml)
        self.assertIn("quadBezTo", xml)

    def test_unknown_shape_type_with_path_uses_custgeom(self):
        """pptxShapeType 是不在清单里的自定义名 → 走 custGeom。"""
        _, slide = _make_blank_slide()
        props = {
            "pptxShapeType": "customWave",
            "path": "M 0 50 C 25 0 75 100 100 50 Z",
            "viewBox": [100, 100],
            "fill": "#0000FF",
        }
        shape = _write_shape_element(
            slide, props, Emu(0), Emu(0), Emu(914400), Emu(914400)
        )
        xml = _get_shape_spPr_xml(shape)
        self.assertIn("custGeom", xml)
        self.assertIn("cubicBezTo", xml)

    def test_round_rect_with_keypoints_still_uses_round_rect_path(self):
        """带 keypoints 的 roundRect 应继续走 _apply_round_rect_*，
        不被新逻辑覆盖（保护已有圆角矩形通路）。"""
        _, slide = _make_blank_slide()
        props = {
            "pptxShapeType": "roundRect",
            "keypoints": [0.1, 0.1, 0.1, 0.1],
            "fill": "#888888",
        }
        shape = _write_shape_element(
            slide, props, Emu(0), Emu(0), Emu(914400), Emu(914400)
        )
        xml = _get_shape_spPr_xml(shape)
        # 四角一致：走 prstGeom roundRect
        self.assertIn("prstGeom", xml)
        self.assertIn('prst="roundRect"', xml)

    def test_round_rect_with_uneven_keypoints_uses_custgeom(self):
        """四角不一致的 roundRect 通过 _apply_round_rect_custom_geometry
        走 custGeom，新 path 逻辑不应额外干扰。"""
        _, slide = _make_blank_slide()
        props = {
            "pptxShapeType": "roundRect",
            "keypoints": [0.2, 0.05, 0.2, 0.05],
            "fill": "#888888",
        }
        shape = _write_shape_element(
            slide, props, Emu(0), Emu(0), Emu(914400), Emu(914400)
        )
        xml = _get_shape_spPr_xml(shape)
        self.assertIn("custGeom", xml)
        # 不应该有 prstGeom 残留
        self.assertNotIn("prstGeom", xml)

    def test_custgeom_end_to_end_through_save(self):
        """端到端：save 后解压 pptx，验证 slide1.xml 里真的有 custGeom。"""
        import io
        import zipfile

        prs, slide = _make_blank_slide()
        props = {
            "path": "M 0 0 L 100 0 Q 100 50 50 50 L 0 100 Z",
            "viewBox": [100, 100],
            "fill": "#123456",
        }
        _write_shape_element(
            slide, props, Emu(0), Emu(0), Emu(914400), Emu(914400)
        )

        buf = io.BytesIO()
        prs.save(buf)
        buf.seek(0)
        with zipfile.ZipFile(buf, "r") as z:
            xml = z.read("ppt/slides/slide1.xml").decode("utf-8")
        self.assertIn("custGeom", xml)
        self.assertIn("moveTo", xml)
        self.assertIn("quadBezTo", xml)
        self.assertIn("a:close", xml)

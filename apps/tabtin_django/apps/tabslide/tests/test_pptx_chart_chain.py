import importlib.util
import os
import tempfile
import zipfile
from pathlib import Path
from unittest import TestCase

from pptx import Presentation
from pptx.chart.data import CategoryChartData, XyChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.util import Inches
from lxml import etree

# 直接按文件路径加载，避免触发 apps.tabslide.services.__init__ 的 Django 依赖
_MODULE_PATH = Path(__file__).resolve().parents[1] / "services" / "pptx_io.py"
_SPEC = importlib.util.spec_from_file_location("tabslide_pptx_io_chart_test_module", _MODULE_PATH)
if _SPEC is None or _SPEC.loader is None:
    raise RuntimeError(f"Failed to load module spec from {_MODULE_PATH}")
_PPTX_IO = importlib.util.module_from_spec(_SPEC)
_SPEC.loader.exec_module(_PPTX_IO)

read = _PPTX_IO.read
write = _PPTX_IO.write
_chart_has_data_labels = _PPTX_IO._chart_has_data_labels


class _DummyChart:
    def __init__(self, xml: str):
        self.element = etree.fromstring(xml.encode("utf-8"))


class _DummyPlot:
    pass


class TestPptxChartChain(TestCase):
    def _read_single_chart_props(self, pptx_path: str) -> dict:
        pages = read(pptx_path=pptx_path, canvas_width=1920, canvas_height=1080)
        self.assertEqual(len(pages), 1)
        elements = pages[0].get("elements", [])
        charts = [el for el in elements if el.get("type") == "chart"]
        self.assertEqual(len(charts), 1)
        return charts[0].get("props", {})

    @staticmethod
    def _normalize_hex(color: str) -> str:
        if not isinstance(color, str):
            return ""
        return color.strip().lower()

    @staticmethod
    def _labels_as_float(labels):
        out = []
        for item in labels:
            try:
                out.append(float(item))
            except Exception:
                pass
        return out

    @staticmethod
    def _x_series_as_float(x_series):
        out = []
        for row in x_series or []:
            parsed_row = []
            for item in row:
                try:
                    parsed_row.append(float(item))
                except Exception:
                    parsed_row.append(None)
            out.append(parsed_row)
        return out

    @staticmethod
    def _read_chart_xml_root(pptx_path: str, chart_idx: int = 1):
        chart_name = f"ppt/charts/chart{chart_idx}.xml"
        with zipfile.ZipFile(pptx_path, "r") as zf:
            with zf.open(chart_name) as fp:
                return etree.fromstring(fp.read())

    def test_read_maps_all_supported_chart_types_from_native_pptx(self):
        cases = [
            ("bar", XL_CHART_TYPE.COLUMN_CLUSTERED, False),
            ("column", XL_CHART_TYPE.BAR_CLUSTERED, False),
            ("line", XL_CHART_TYPE.LINE, False),
            ("area", XL_CHART_TYPE.AREA, False),
            ("pie", XL_CHART_TYPE.PIE, False),
            ("ring", XL_CHART_TYPE.DOUGHNUT, False),
            ("radar", XL_CHART_TYPE.RADAR, False),
            ("scatter", XL_CHART_TYPE.XY_SCATTER, True),
        ]

        for expected_type, native_type, is_scatter in cases:
            with self.subTest(chart_type=expected_type):
                prs = Presentation()
                slide = prs.slides.add_slide(prs.slide_layouts[6])

                if is_scatter:
                    chart_data = XyChartData()
                    s1 = chart_data.add_series("S1")
                    s2 = chart_data.add_series("S2")
                    for idx, y_val in enumerate((10, 20, 30), start=1):
                        s1.add_data_point(idx, y_val)
                    for idx, y_val in enumerate((15, 25, 35), start=1):
                        s2.add_data_point(idx, y_val)
                else:
                    chart_data = CategoryChartData()
                    chart_data.categories = ["A", "B", "C"]
                    chart_data.add_series("S1", (10, 20, 30))
                    if expected_type not in ("pie", "ring"):
                        chart_data.add_series("S2", (15, 25, 35))

                chart = slide.shapes.add_chart(
                    native_type,
                    Inches(1),
                    Inches(1),
                    Inches(6),
                    Inches(3.5),
                    chart_data,
                ).chart

                chart.has_title = True
                chart.chart_title.text_frame.text = f"{expected_type}-native"
                chart.has_legend = True
                chart.legend.position = XL_LEGEND_POSITION.TOP

                if not is_scatter:
                    plot = chart.plots[0]
                    plot.has_data_labels = True
                    plot.data_labels.show_value = True

                fd, pptx_path = tempfile.mkstemp(suffix=".pptx")
                os.close(fd)
                try:
                    prs.save(pptx_path)
                    props = self._read_single_chart_props(pptx_path)

                    self.assertEqual(props.get("chartType"), expected_type)
                    self.assertEqual(props.get("title"), f"{expected_type}-native")
                    options = props.get("options", {})
                    self.assertTrue(options.get("showLegend"))
                    self.assertEqual(options.get("legendPosition"), "t")

                    data = props.get("data", {})
                    legends = data.get("legends", [])
                    self.assertGreaterEqual(len(legends), 1)
                    self.assertTrue(all(str(item).strip() for item in legends))

                    if is_scatter:
                        self.assertEqual(self._labels_as_float(data.get("labels", []))[:3], [1.0, 2.0, 3.0])
                    else:
                        self.assertEqual(data.get("labels", [])[:3], ["A", "B", "C"])
                        self.assertTrue(options.get("showDataLabel"))
                finally:
                    try:
                        os.remove(pptx_path)
                    except Exception:
                        pass

    def test_write_read_roundtrip_preserves_chart_core_options(self):
        cases = [
            {
                "chartType": "bar",
                "title": "bar-title",
                "labels": ["Q1", "Q2", "Q3", "Q4"],
                "legends": ["A", "B"],
                "series": [[10, 20, 30, 40], [12, 22, 32, 42]],
                "options": {"showLegend": True, "legendPosition": "b", "showDataLabel": True, "stack": True},
                "themeColors": ["#ff0000", "#00aa00"],
                "fill": "#f3f3f3",
                "textColor": "#334455",
                "gridColor": "#aabbcc",
            },
            {
                "chartType": "column",
                "title": "column-title",
                "labels": ["Q1", "Q2", "Q3", "Q4"],
                "legends": ["A", "B"],
                "series": [[8, 18, 28, 38], [9, 19, 29, 39]],
                "options": {"showLegend": True, "legendPosition": "l", "showDataLabel": True, "stack": True},
                "themeColors": ["#4472c4", "#ed7d31"],
            },
            {
                "chartType": "line",
                "title": "line-title",
                "labels": ["Q1", "Q2", "Q3", "Q4"],
                "legends": ["A", "B"],
                "series": [[6, 16, 26, 36], [7, 17, 27, 37]],
                "options": {"showLegend": True, "legendPosition": "t", "showDataLabel": True, "stack": True, "lineSmooth": True},
                "themeColors": ["#70ad47", "#264478"],
            },
            {
                "chartType": "area",
                "title": "area-title",
                "labels": ["Q1", "Q2", "Q3", "Q4"],
                "legends": ["A", "B"],
                "series": [[11, 21, 31, 41], [14, 24, 34, 44]],
                "options": {"showLegend": True, "legendPosition": "r", "showDataLabel": True, "stack": True, "lineSmooth": True},
                "themeColors": ["#5b9bd5", "#ffc000"],
            },
            {
                "chartType": "pie",
                "title": "pie-title",
                "labels": ["A", "B", "C", "D"],
                "legends": ["占比"],
                "series": [[35, 25, 20, 20]],
                "options": {"showLegend": True, "legendPosition": "b", "showDataLabel": True},
                "themeColors": ["#ff6384"],
            },
            {
                "chartType": "ring",
                "title": "ring-title",
                "labels": ["A", "B", "C", "D"],
                "legends": ["占比"],
                "series": [[30, 30, 20, 20]],
                "options": {"showLegend": True, "legendPosition": "t", "showDataLabel": True},
                "themeColors": ["#36a2eb"],
            },
            {
                "chartType": "radar",
                "title": "radar-title",
                "labels": ["维度1", "维度2", "维度3", "维度4", "维度5"],
                "legends": ["A", "B"],
                "series": [[60, 70, 80, 50, 40], [55, 65, 75, 45, 35]],
                "options": {"showLegend": True, "legendPosition": "l", "showDataLabel": True},
                "themeColors": ["#9966ff", "#ff9f40"],
            },
            {
                "chartType": "scatter",
                "title": "scatter-title",
                "labels": ["1", "2", "3", "4"],
                "legends": ["A", "B"],
                "series": [[20, 35, 45, 60], [22, 30, 40, 58]],
                "options": {"showLegend": True, "legendPosition": "r", "showDataLabel": True, "lineSmooth": True},
                "themeColors": ["#4bc0c0", "#ffcd56"],
            },
        ]

        pages = []
        for idx, case in enumerate(cases, start=1):
            props = {
                "chartType": case["chartType"],
                "data": {
                    "labels": case["labels"],
                    "legends": case["legends"],
                    "series": case["series"],
                },
                "options": case["options"],
                "title": case["title"],
                "themeColors": case["themeColors"],
            }
            if case.get("fill"):
                props["fill"] = case["fill"]
            if case.get("textColor"):
                props["textColor"] = case["textColor"]
            if case.get("gridColor"):
                props["gridColor"] = case["gridColor"]

            pages.append(
                {
                    "id": f"page-{idx}",
                    "background": {"type": "color", "value": "#ffffff"},
                    "elements": [
                        {
                            "id": f"chart-{idx}",
                            "type": "chart",
                            "x": 80,
                            "y": 80,
                            "width": 900,
                            "height": 500,
                            "rotate": 0,
                            "zIndex": 0,
                            "props": props,
                        }
                    ],
                }
            )

        fd, pptx_path = tempfile.mkstemp(suffix=".pptx")
        os.close(fd)
        try:
            write(pages=pages, output_path=pptx_path, canvas_width=1920, canvas_height=1080)
            out_pages = read(pptx_path=pptx_path, canvas_width=1920, canvas_height=1080)

            self.assertEqual(len(out_pages), len(cases))

            for idx, case in enumerate(cases):
                with self.subTest(chart_type=case["chartType"]):
                    chart_elements = [el for el in out_pages[idx].get("elements", []) if el.get("type") == "chart"]
                    self.assertEqual(len(chart_elements), 1)
                    props = chart_elements[0].get("props", {})
                    options = props.get("options", {})
                    data = props.get("data", {})

                    self.assertEqual(props.get("chartType"), case["chartType"])
                    self.assertEqual(props.get("title"), case["title"])

                    self.assertTrue(options.get("showLegend"))
                    self.assertEqual(options.get("legendPosition"), case["options"]["legendPosition"])
                    self.assertTrue(options.get("showDataLabel"))

                    if case["chartType"] in ("bar", "column", "line", "area"):
                        self.assertTrue(options.get("stack"))
                    else:
                        self.assertFalse(bool(options.get("stack", False)))

                    if case["chartType"] in ("line", "area", "scatter"):
                        self.assertTrue(options.get("lineSmooth"))
                    else:
                        self.assertFalse(bool(options.get("lineSmooth", False)))

                    self.assertEqual(data.get("legends", [])[: len(case["legends"])], case["legends"])

                    if case["chartType"] == "scatter":
                        self.assertEqual(
                            self._labels_as_float(data.get("labels", []))[: len(case["labels"])],
                            [float(v) for v in case["labels"]],
                        )
                    else:
                        self.assertEqual(
                            data.get("labels", [])[: len(case["labels"])],
                            case["labels"],
                        )

                    out_colors = props.get("themeColors", [])
                    self.assertGreaterEqual(len(out_colors), 1)
                    self.assertEqual(
                        self._normalize_hex(out_colors[0]),
                        self._normalize_hex(case["themeColors"][0]),
                    )

                    if case["chartType"] == "bar":
                        self.assertEqual(self._normalize_hex(props.get("fill", "")), "#f3f3f3")
                        self.assertEqual(self._normalize_hex(props.get("textColor", "")), "#334455")
                        self.assertEqual(self._normalize_hex(props.get("gridColor", "")), "#aabbcc")
        finally:
            try:
                os.remove(pptx_path)
            except Exception:
                pass

    def test_write_read_roundtrip_preserves_pie_slice_theme_colors(self):
        pie_palette = ["#ff6384", "#36a2eb", "#ffce56", "#4bc0c0"]
        pages = [
            {
                "id": "page-pie-colors",
                "background": {"type": "color", "value": "#ffffff"},
                "elements": [
                    {
                        "id": "chart-pie-colors",
                        "type": "chart",
                        "x": 80,
                        "y": 80,
                        "width": 900,
                        "height": 500,
                        "rotate": 0,
                        "zIndex": 0,
                        "props": {
                            "chartType": "pie",
                            "title": "pie-colors",
                            "data": {
                                "labels": ["A", "B", "C", "D"],
                                "legends": ["占比"],
                                "series": [[35, 25, 20, 20]],
                            },
                            "options": {"showLegend": True, "legendPosition": "b", "showDataLabel": True},
                            "themeColors": pie_palette,
                        },
                    }
                ],
            }
        ]

        fd, pptx_path = tempfile.mkstemp(suffix=".pptx")
        os.close(fd)
        try:
            write(pages=pages, output_path=pptx_path, canvas_width=1920, canvas_height=1080)
            out_pages = read(pptx_path=pptx_path, canvas_width=1920, canvas_height=1080)
            self.assertEqual(len(out_pages), 1)

            chart_elements = [el for el in out_pages[0].get("elements", []) if el.get("type") == "chart"]
            self.assertEqual(len(chart_elements), 1)
            props = chart_elements[0].get("props", {})
            out_colors = props.get("themeColors", [])

            self.assertGreaterEqual(len(out_colors), len(pie_palette))
            self.assertEqual(
                [self._normalize_hex(c) for c in out_colors[: len(pie_palette)]],
                [self._normalize_hex(c) for c in pie_palette],
            )
        finally:
            try:
                os.remove(pptx_path)
            except Exception:
                pass

    def test_chart_has_data_labels_detects_show_percent_xml(self):
        chart = _DummyChart(
            """
            <c:chartSpace xmlns:c="http://schemas.openxmlformats.org/drawingml/2006/chart">
              <c:chart>
                <c:plotArea>
                  <c:pieChart>
                    <c:dLbls>
                      <c:showVal val="0" />
                      <c:showPercent val="1" />
                    </c:dLbls>
                  </c:pieChart>
                </c:plotArea>
              </c:chart>
            </c:chartSpace>
            """
        )
        self.assertTrue(_chart_has_data_labels(chart, _DummyPlot()))

    def test_write_read_pie_defaults_align_with_frontend(self):
        pages = [
            {
                "id": "page-pie-defaults",
                "background": {"type": "color", "value": "#ffffff"},
                "elements": [
                    {
                        "id": "chart-pie-defaults",
                        "type": "chart",
                        "x": 80,
                        "y": 80,
                        "width": 900,
                        "height": 500,
                        "rotate": 0,
                        "zIndex": 0,
                        "props": {
                            "chartType": "pie",
                            "data": {
                                "labels": ["A", "B", "C"],
                                "legends": ["占比"],
                                "series": [[35, 25, 40]],
                            },
                            # 不传 options，验证默认行为
                        },
                    }
                ],
            }
        ]

        fd, pptx_path = tempfile.mkstemp(suffix=".pptx")
        os.close(fd)
        try:
            write(pages=pages, output_path=pptx_path, canvas_width=1920, canvas_height=1080)
            out_pages = read(pptx_path=pptx_path, canvas_width=1920, canvas_height=1080)
            chart_elements = [el for el in out_pages[0].get("elements", []) if el.get("type") == "chart"]
            self.assertEqual(len(chart_elements), 1)

            props = chart_elements[0].get("props", {})
            options = props.get("options", {})
            self.assertTrue(options.get("showLegend"))
            self.assertTrue(options.get("showDataLabel"))
        finally:
            try:
                os.remove(pptx_path)
            except Exception:
                pass

    def test_write_read_scatter_preserves_per_series_x_values(self):
        pages = [
            {
                "id": "page-scatter-xseries",
                "background": {"type": "color", "value": "#ffffff"},
                "elements": [
                    {
                        "id": "chart-scatter-xseries",
                        "type": "chart",
                        "x": 80,
                        "y": 80,
                        "width": 900,
                        "height": 500,
                        "rotate": 0,
                        "zIndex": 0,
                        "props": {
                            "chartType": "scatter",
                            "title": "scatter-xseries",
                            "data": {
                                "labels": ["1", "2", "3", "4"],
                                "legends": ["A", "B"],
                                "series": [[20, 35, 45, 60], [22, 30, 40, 58]],
                                "xSeries": [[1, 2, 3, 4], [10, 20, 30, 40]],
                            },
                            "options": {"showLegend": True, "legendPosition": "b", "lineSmooth": True},
                            "themeColors": ["#4bc0c0", "#ffcd56"],
                        },
                    }
                ],
            }
        ]

        fd, pptx_path = tempfile.mkstemp(suffix=".pptx")
        os.close(fd)
        try:
            write(pages=pages, output_path=pptx_path, canvas_width=1920, canvas_height=1080)
            out_pages = read(pptx_path=pptx_path, canvas_width=1920, canvas_height=1080)
            chart_elements = [el for el in out_pages[0].get("elements", []) if el.get("type") == "chart"]
            self.assertEqual(len(chart_elements), 1)

            props = chart_elements[0].get("props", {})
            data = props.get("data", {})
            x_series = self._x_series_as_float(data.get("xSeries", []))
            self.assertEqual(len(x_series), 2)
            self.assertEqual(x_series[0][:4], [1.0, 2.0, 3.0, 4.0])
            self.assertEqual(x_series[1][:4], [10.0, 20.0, 30.0, 40.0])
            self.assertEqual(self._labels_as_float(data.get("labels", []))[:4], [1.0, 2.0, 3.0, 4.0])
        finally:
            try:
                os.remove(pptx_path)
            except Exception:
                pass

    def test_write_exports_series_theme_color_keys_as_scheme_clr(self):
        pages = [
            {
                "id": "page-series-theme-keys",
                "background": {"type": "color", "value": "#ffffff"},
                "elements": [
                    {
                        "id": "chart-series-theme-keys",
                        "type": "chart",
                        "x": 80,
                        "y": 80,
                        "width": 900,
                        "height": 500,
                        "rotate": 0,
                        "zIndex": 0,
                        "props": {
                            "chartType": "bar",
                            "title": "series-theme-keys",
                            "data": {
                                "labels": ["Q1", "Q2", "Q3"],
                                "legends": ["A", "B"],
                                "series": [[10, 20, 30], [12, 22, 32]],
                            },
                            "themeColors": ["#4472C4", "#ED7D31"],
                            "themeColorKeys": ["accent1", "accent2"],
                        },
                    }
                ],
            }
        ]

        fd, pptx_path = tempfile.mkstemp(suffix=".pptx")
        os.close(fd)
        try:
            write(pages=pages, output_path=pptx_path, canvas_width=1920, canvas_height=1080)

            chart_xml = self._read_chart_xml_root(pptx_path, 1)
            ns = {
                "c": "http://schemas.openxmlformats.org/drawingml/2006/chart",
                "a": "http://schemas.openxmlformats.org/drawingml/2006/main",
            }
            series_nodes = chart_xml.findall(".//c:ser", ns)
            self.assertGreaterEqual(len(series_nodes), 2)

            s1 = series_nodes[0].find("./c:spPr/a:solidFill/a:schemeClr", ns)
            s2 = series_nodes[1].find("./c:spPr/a:solidFill/a:schemeClr", ns)
            self.assertIsNotNone(s1)
            self.assertIsNotNone(s2)
            self.assertEqual((s1.get("val") or "").strip(), "accent1")
            self.assertEqual((s2.get("val") or "").strip(), "accent2")

            out_pages = read(pptx_path=pptx_path, canvas_width=1920, canvas_height=1080)
            chart_elements = [el for el in out_pages[0].get("elements", []) if el.get("type") == "chart"]
            self.assertEqual(len(chart_elements), 1)
            props = chart_elements[0].get("props", {})
            out_keys = props.get("themeColorKeys", [])
            self.assertGreaterEqual(len(out_keys), 2)
            self.assertEqual(out_keys[:2], ["accent1", "accent2"])
        finally:
            try:
                os.remove(pptx_path)
            except Exception:
                pass

    def test_write_exports_pie_point_theme_color_keys_as_scheme_clr(self):
        pages = [
            {
                "id": "page-pie-theme-keys",
                "background": {"type": "color", "value": "#ffffff"},
                "elements": [
                    {
                        "id": "chart-pie-theme-keys",
                        "type": "chart",
                        "x": 80,
                        "y": 80,
                        "width": 900,
                        "height": 500,
                        "rotate": 0,
                        "zIndex": 0,
                        "props": {
                            "chartType": "pie",
                            "title": "pie-theme-keys",
                            "data": {
                                "labels": ["A", "B", "C", "D"],
                                "legends": ["占比"],
                                "series": [[35, 25, 20, 20]],
                            },
                            "themeColors": ["#4472C4", "#ED7D31", "#A5A5A5", "#FFC000"],
                            "themeColorKeys": ["accent1", "accent2", "accent3", "accent4"],
                        },
                    }
                ],
            }
        ]

        fd, pptx_path = tempfile.mkstemp(suffix=".pptx")
        os.close(fd)
        try:
            write(pages=pages, output_path=pptx_path, canvas_width=1920, canvas_height=1080)

            chart_xml = self._read_chart_xml_root(pptx_path, 1)
            ns = {
                "c": "http://schemas.openxmlformats.org/drawingml/2006/chart",
                "a": "http://schemas.openxmlformats.org/drawingml/2006/main",
            }
            dpt_scheme_map = {}
            for dpt in chart_xml.findall(".//c:ser/c:dPt", ns):
                idx_el = dpt.find("./c:idx", ns)
                scheme = dpt.find("./c:spPr/a:solidFill/a:schemeClr", ns)
                if idx_el is None or scheme is None:
                    continue
                try:
                    idx = int(idx_el.get("val") or -1)
                except (TypeError, ValueError):
                    continue
                if idx >= 0:
                    dpt_scheme_map[idx] = (scheme.get("val") or "").strip()

            self.assertEqual(dpt_scheme_map.get(0), "accent1")
            self.assertEqual(dpt_scheme_map.get(1), "accent2")
            self.assertEqual(dpt_scheme_map.get(2), "accent3")
            self.assertEqual(dpt_scheme_map.get(3), "accent4")

            out_pages = read(pptx_path=pptx_path, canvas_width=1920, canvas_height=1080)
            chart_elements = [el for el in out_pages[0].get("elements", []) if el.get("type") == "chart"]
            self.assertEqual(len(chart_elements), 1)
            props = chart_elements[0].get("props", {})
            out_keys = props.get("themeColorKeys", [])
            self.assertGreaterEqual(len(out_keys), 4)
            self.assertEqual(out_keys[:4], ["accent1", "accent2", "accent3", "accent4"])
        finally:
            try:
                os.remove(pptx_path)
            except Exception:
                pass

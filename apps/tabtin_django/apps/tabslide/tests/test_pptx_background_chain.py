import importlib.util
from pathlib import Path
from typing import Optional
from unittest import TestCase

from pptx import Presentation
from pptx.oxml import parse_xml
from pptx.oxml.ns import nsdecls

# 直接按文件路径加载，避免触发 apps.tabslide.services.__init__ 的 Django 依赖
_MODULE_PATH = Path(__file__).resolve().parents[1] / "services" / "pptx_io.py"
_SPEC = importlib.util.spec_from_file_location("tabslide_pptx_bg_test_module", _MODULE_PATH)
if _SPEC is None or _SPEC.loader is None:
    raise RuntimeError(f"Failed to load module spec from {_MODULE_PATH}")
_PPTX_IO = importlib.util.module_from_spec(_SPEC)
_SPEC.loader.exec_module(_PPTX_IO)

_set_slide_background_color = _PPTX_IO._set_slide_background_color
_set_slide_background_theme = _PPTX_IO._set_slide_background_theme
_set_slide_background_gradient = _PPTX_IO._set_slide_background_gradient
_set_slide_background_image = _PPTX_IO._set_slide_background_image
_extract_background_from_element = _PPTX_IO._extract_background_from_element

_PNG_DATA_URL = (
    "data:image/png;base64,"
    "iVBORw0KGgoAAAANSUhEUgAAAAEAAAABCAQAAAC1HAwCAAAAC0lEQVR42mP8/x8AAwMCAO3ZfD0AAAAASUVORK5CYII="
)


class TestPptxBackgroundChain(TestCase):
    def _new_slide(self):
        prs = Presentation()
        return prs.slides.add_slide(prs.slide_layouts[6])

    def _inject_bg_ref(self, slide, *, scheme: str, alpha: Optional[float] = None):
        nsmap_p = "http://schemas.openxmlformats.org/presentationml/2006/main"
        c_sld = slide.background._element
        old_bg = c_sld.find(f"{{{nsmap_p}}}bg")
        if old_bg is not None:
            c_sld.remove(old_bg)

        alpha_xml = ""
        if alpha is not None:
            alpha_xml = f'<a:alpha val="{int(alpha * 100000)}"/>'

        bg = parse_xml(
            f'<p:bg {nsdecls("p", "a")}>'
            f'<p:bgRef idx="1001"><a:schemeClr val="{scheme}">{alpha_xml}</a:schemeClr></p:bgRef>'
            f"</p:bg>"
        )

        sp_tree = c_sld.find(f"{{{nsmap_p}}}spTree")
        if sp_tree is not None:
            c_sld.insert(c_sld.index(sp_tree), bg)
        else:
            c_sld.append(bg)

    def _inject_bg_solid_scheme(
        self,
        slide,
        *,
        scheme: str,
        lum_mod: Optional[int] = None,
        lum_off: Optional[int] = None,
        alpha: Optional[float] = None,
    ):
        nsmap_p = "http://schemas.openxmlformats.org/presentationml/2006/main"
        c_sld = slide.background._element
        old_bg = c_sld.find(f"{{{nsmap_p}}}bg")
        if old_bg is not None:
            c_sld.remove(old_bg)

        transforms = ""
        if lum_mod is not None:
            transforms += f'<a:lumMod val="{int(lum_mod)}"/>'
        if lum_off is not None:
            transforms += f'<a:lumOff val="{int(lum_off)}"/>'
        if alpha is not None:
            transforms += f'<a:alpha val="{int(alpha * 100000)}"/>'

        bg = parse_xml(
            f'<p:bg {nsdecls("p", "a")}>'
            f'<p:bgPr><a:solidFill><a:schemeClr val="{scheme}">{transforms}</a:schemeClr></a:solidFill>'
            f"<a:effectLst/></p:bgPr></p:bg>"
        )

        sp_tree = c_sld.find(f"{{{nsmap_p}}}spTree")
        if sp_tree is not None:
            c_sld.insert(c_sld.index(sp_tree), bg)
        else:
            c_sld.append(bg)

    def test_solid_background_roundtrip_with_alpha(self):
        slide = self._new_slide()
        _set_slide_background_color(slide, "rgba(255,0,0,0.5)")

        bg = _extract_background_from_element(slide, None)
        self.assertIsNotNone(bg)
        self.assertEqual(bg["type"], "color")
        self.assertEqual(bg["value"], "rgba(255,0,0,0.5)")

    def test_theme_background_roundtrip_keeps_theme_key(self):
        slide = self._new_slide()
        _set_slide_background_theme(
            slide,
            {"key": "accent1", "color": "rgba(17,34,51,0.4)"},
            fallback_color="#112233",
        )

        bg = _extract_background_from_element(slide, {"accent1": "#4472C4"})
        self.assertIsNotNone(bg)
        self.assertEqual(bg["type"], "color")
        self.assertEqual(bg.get("theme", {}).get("key"), "accent1")
        self.assertIn("rgba(", bg["value"])
        self.assertTrue(bg["value"].endswith(",0.4)"))

    def test_theme_background_roundtrip_supports_bg_tx_alias_keys(self):
        cases = [
            ("bg1", "rgba(255,255,255,0.4)"),
            ("tx1", "rgba(0,0,0,0.4)"),
            ("bg2", "rgba(231,230,230,0.4)"),
            ("tx2", "rgba(68,84,106,0.4)"),
        ]
        for key, color in cases:
            with self.subTest(key=key):
                slide = self._new_slide()
                _set_slide_background_theme(slide, {"key": key, "color": color}, fallback_color=color)
                bg = _extract_background_from_element(slide, None)
                self.assertIsNotNone(bg)
                self.assertEqual(bg["type"], "color")
                self.assertEqual(bg.get("theme", {}).get("key"), key)
                self.assertEqual(bg["value"], color)

    def test_bgref_theme_background_is_extractable(self):
        slide = self._new_slide()
        self._inject_bg_ref(slide, scheme="bg1", alpha=0.55)
        bg = _extract_background_from_element(slide, {"2": "#fefefe"})
        self.assertIsNotNone(bg)
        self.assertEqual(bg["type"], "color")
        self.assertEqual(bg.get("theme", {}).get("key"), "bg1")
        self.assertEqual(bg["value"], "rgba(254,254,254,0.55)")

    def test_solid_theme_background_with_lum_transform_keeps_resolved_color(self):
        slide = self._new_slide()
        self._inject_bg_solid_scheme(
            slide,
            scheme="accent1",
            lum_mod=60000,
            lum_off=20000,
            alpha=0.5,
        )

        bg = _extract_background_from_element(slide, {"5": "#4472C4"})
        self.assertIsNotNone(bg)
        self.assertEqual(bg["type"], "color")
        self.assertEqual(bg.get("theme", {}).get("key"), "accent1")
        self.assertTrue(bg["value"].startswith("rgba("))

        nsmap_a = "http://schemas.openxmlformats.org/drawingml/2006/main"
        scheme_el = slide.background._element.find(f".//{{{nsmap_a}}}schemeClr")
        self.assertIsNotNone(scheme_el)
        expected_hex = _PPTX_IO._apply_color_transforms("#4472C4", scheme_el).lower()
        r = int(expected_hex[1:3], 16)
        g = int(expected_hex[3:5], 16)
        b = int(expected_hex[5:7], 16)
        self.assertEqual(bg["value"], f"rgba({r},{g},{b},0.5)")

    def test_radial_gradient_background_roundtrip(self):
        slide = self._new_slide()
        _set_slide_background_gradient(
            slide,
            {
                "type": "radial",
                "rotate": 0,
                "colors": [
                    {"pos": 0, "color": "#ff0000"},
                    {"pos": 1, "color": "rgba(0,0,255,0.3)"},
                ],
            },
        )

        bg = _extract_background_from_element(slide, None)
        self.assertIsNotNone(bg)
        self.assertEqual(bg["type"], "gradient")
        grad = bg["gradient"]
        self.assertEqual(grad["type"], "radial")
        self.assertGreaterEqual(len(grad["colors"]), 2)
        self.assertEqual(grad["colors"][0]["color"].lower(), "#ff0000")
        self.assertEqual(grad["colors"][-1]["color"], "rgba(0,0,255,0.3)")

    def test_image_background_roundtrip_with_size_modes(self):
        for size in ("cover", "contain", "repeat"):
            with self.subTest(size=size):
                slide = self._new_slide()
                _set_slide_background_image(slide, {"src": _PNG_DATA_URL, "size": size})

                bg = _extract_background_from_element(slide, None)
                self.assertIsNotNone(bg)
                self.assertEqual(bg["type"], "image")
                self.assertEqual(bg["image"]["size"], size)
                self.assertTrue(bg["image"]["src"].startswith("data:image/png;base64,"))

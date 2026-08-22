""": HTML 到 PPTX 的单行、几何与字体契约回归。"""

from __future__ import annotations

import base64
import importlib.util
import os
import sys
import tempfile
import zipfile
from pathlib import Path
from unittest import TestCase, skipUnless

from lxml import etree


_SERVICES = Path(__file__).resolve().parents[1] / "services"
_DJANGO_ROOT = Path(__file__).resolve().parents[3]
if str(_DJANGO_ROOT) not in sys.path:
    sys.path.insert(0, str(_DJANGO_ROOT))


def _load_module(name: str, filename: str):
    spec = importlib.util.spec_from_file_location(name, _SERVICES / filename)
    if spec is None or spec.loader is None:
        raise RuntimeError(f"Failed to load {filename}")
    module = importlib.util.module_from_spec(spec)
    spec.loader.exec_module(module)
    return module


_PPTX_IO = _load_module("tabslide_pptx_text_fidelity_test", "pptx_io.py")
_DOM = _load_module("tabslide_dom_text_fidelity_test", "dom_extractor.py")

A = "http://schemas.openxmlformats.org/drawingml/2006/main"


class FontContractTests(TestCase):
    def test_css_stack_resolves_web_font_to_portable_script_fonts(self):
        latin, east_asian, complex_script = _PPTX_IO._resolve_pptx_typefaces(
            "'Inter', 'Noto Sans SC', 'Microsoft YaHei', sans-serif",
            "上海 36°C",
        )

        self.assertEqual(latin, "Arial")
        self.assertEqual(east_asian, "Microsoft YaHei")
        self.assertEqual(complex_script, "Microsoft YaHei")

    def test_single_exact_font_is_preserved_for_editor_roundtrip(self):
        self.assertEqual(
            _PPTX_IO._resolve_pptx_typefaces("Aptos", "Weather"),
            ("Aptos", "Aptos", "Aptos"),
        )

    def test_latin_only_exact_font_uses_cjk_fallback_for_mixed_text(self):
        self.assertEqual(
            _PPTX_IO._resolve_pptx_typefaces("Arial", "上海 Weather"),
            ("Arial", "Microsoft YaHei", "Microsoft YaHei"),
        )

    def test_custom_primary_font_in_css_stack_is_preserved(self):
        self.assertEqual(
            _PPTX_IO._resolve_pptx_typefaces("'MiSans', Arial, sans-serif", "上海 Weather"),
            ("MiSans", "MiSans", "MiSans"),
        )


class DomTextFlowContractTests(TestCase):
    def test_layout_enrichment_preserves_browser_line_and_flow_geometry(self):
        element = {
            "type": "text",
            "x": 120,
            "y": 100,
            "width": 72,
            "height": 58,
            "content": "<p>36°C</p>",
            "defaultFontName": "Inter",
            "defaultFontSize": 48,
        }
        layout = {
            "x": 120,
            "y": 100,
            "width": 72,
            "height": 58,
            "flowX": 120,
            "flowWidth": 240,
            "sourceLineCount": 1,
            "fontFamily": "Inter, 'Noto Sans SC', 'Microsoft YaHei', sans-serif",
            "runs": [{"text": "36°C", "fontSize": 64}],
        }

        _DOM._enrich_text_element_from_layout(element, layout)

        self.assertEqual(element["x"], 120)
        self.assertEqual(element["width"], 240)
        self.assertEqual(element["sourceLineCount"], 1)
        self.assertFalse(element["wordWrap"])
        self.assertEqual(element["fontFamilyFallbacks"], layout["fontFamily"])

    def test_extractor_measures_rendered_lines_and_available_flow_width(self):
        js = _DOM._JS_TEXT_BOX_RECT
        self.assertIn("getClientRects", js)
        self.assertIn("sourceLineCount", _DOM._PURE_DOM_EXTRACT_SLIDE_JS)
        self.assertIn("flowWidth", _DOM._WALKER_TEXT_EXTRACT_JS)
        self.assertIn("sourceLineCount", _DOM._EXTRACT_SHAPES_JS)
        self.assertIn("wordWrap", _DOM._EXTRACT_SHAPES_JS)
        self.assertIn("defaultFontFamily", _DOM._EXTRACT_SHAPES_JS)

    def test_parent_layout_does_not_overwrite_flex_footer_child(self):
        elements = [
            {
                "id": "left",
                "type": "text",
                "x": 76,
                "y": 641,
                "width": 1005,
                "height": 19,
                "content": "<p>Author / AI Assistant</p>",
                "defaultFontSize": 10.5,
                "defaultColor": "#A8A29E",
            },
            {
                "id": "right",
                "type": "text",
                "x": 1081,
                "y": 641,
                "width": 123,
                "height": 19,
                "content": "<p>Shanghai Weather</p>",
                "defaultFontSize": 10.5,
                "defaultColor": "#A8A29E",
            },
        ]
        layouts = [
            {
                "x": 76,
                "y": 641,
                "width": 1128,
                "height": 19,
                "runs": [
                    {"text": "Author / AI Assistant", "fontSize": 14},
                    {"text": "Shanghai Weather", "fontSize": 14},
                ],
            },
            {
                "x": 76,
                "y": 641,
                "width": 102,
                "height": 19,
                "runs": [{"text": "Author / AI Assistant", "fontSize": 14}],
            },
            {
                "x": 1081,
                "y": 641,
                "width": 123,
                "height": 19,
                "runs": [{"text": "Shanghai Weather", "fontSize": 14}],
            },
        ]

        result = _DOM._postprocess_slide_elements(
            elements,
            text_layout_data=layouts,
        )
        by_id = {element["id"]: element for element in result}

        self.assertIn("Author / AI Assistant", by_id["left"]["content"])
        self.assertNotIn("Shanghai Weather", by_id["left"]["content"])
        self.assertIn("Shanghai Weather", by_id["right"]["content"])


@skipUnless(importlib.util.find_spec("playwright"), "playwright is not installed")
class DomTextFlowBrowserTests(TestCase):
    def test_flex_label_uses_remaining_row_width_without_covering_sibling(self):
        from playwright.sync_api import sync_playwright

        html = """
        <style>
          .ppt-slide { position:relative; width:1280px; height:720px; }
          .row { display:flex; gap:12px; width:300px; align-items:center; }
          .dot { flex:0 0 24px; width:24px; height:24px; border-radius:50%; background:red; }
          .label { font:600 32px Arial; }
        </style>
        <div class="ppt-slide">
          <div class="row"><span class="dot"></span><div class="label">36°C</div></div>
        </div>
        """
        with sync_playwright() as playwright:
            browser = playwright.chromium.launch(headless=True)
            try:
                page = browser.new_page(viewport={"width": 1280, "height": 720})
                page.set_content(html)
                slide = page.query_selector(".ppt-slide")
                self.assertIsNotNone(slide)
                rows = slide.evaluate(
                    _DOM._WALKER_TEXT_EXTRACT_JS,
                    {"canvasWidth": 1280, "canvasHeight": 720},
                )
                label = next(row for row in rows if "36°C" in row.get("content", ""))
                dot_right = page.locator(".dot").bounding_box()["x"] + 24
                self.assertGreaterEqual(label["x"], dot_right)
                self.assertGreater(label["width"], 200)
                self.assertEqual(label["sourceLineCount"], 1)
                self.assertFalse(label["wordWrap"])
            finally:
                browser.close()


class PptxTextFidelityContractTests(TestCase):
    def setUp(self):
        self._tmp = tempfile.TemporaryDirectory()
        self._out = os.path.join(self._tmp.name, "text-fidelity.pptx")

    def tearDown(self):
        self._tmp.cleanup()

    def _write(
        self,
        *,
        source_line_count: int,
        word_wrap: bool,
        font_family: str = "Inter, 'Noto Sans SC', 'Microsoft YaHei', sans-serif",
        font_meta: dict | None = None,
    ) -> None:
        pages = [{
            "id": "p1",
            "width": 1280,
            "height": 720,
            "elements": [{
                "id": "t1",
                "type": "text",
                "x": 80,
                "y": 80,
                "width": 240,
                "height": 58 if source_line_count == 1 else 120,
                "props": {
                    "content": "<p>上海 36°C</p>",
                    "defaultFontName": "Inter",
                    "fontFamilyFallbacks": font_family,
                    "defaultFontSize": 48,
                    "sourceLineCount": source_line_count,
                    "wordWrap": word_wrap,
                },
            }],
        }]
        _PPTX_IO.write(
            pages,
            self._out,
            canvas_width=1280,
            canvas_height=720,
            font_meta=font_meta,
        )

    def _slide_root(self):
        with zipfile.ZipFile(self._out) as archive:
            return etree.fromstring(archive.read("ppt/slides/slide1.xml"))

    def test_browser_single_line_disables_implicit_powerpoint_wrap(self):
        self._write(source_line_count=1, word_wrap=False)
        root = self._slide_root()
        body_pr = next(root.iter(f"{{{A}}}bodyPr"))

        self.assertEqual(body_pr.get("wrap"), "none")

    def test_multiline_text_keeps_word_wrap(self):
        self._write(source_line_count=2, word_wrap=True)
        root = self._slide_root()
        body_pr = next(root.iter(f"{{{A}}}bodyPr"))

        self.assertNotEqual(body_pr.get("wrap"), "none")

    def test_shape_text_single_line_disables_powerpoint_wrap(self):
        pages = [{
            "id": "p1",
            "width": 1280,
            "height": 720,
            "elements": [{
                "id": "badge",
                "type": "shape",
                "x": 80,
                "y": 80,
                "width": 120,
                "height": 36,
                "props": {
                    "pptxShapeType": "roundRect",
                    "fill": "#F97316",
                    "text": {
                        "content": "<p>CONTENTS</p>",
                        "defaultFontSize": 12,
                        "sourceLineCount": 1,
                        "wordWrap": False,
                    },
                },
            }],
        }]
        _PPTX_IO.write(pages, self._out, canvas_width=1280, canvas_height=720)
        root = self._slide_root()
        body_pr = next(root.iter(f"{{{A}}}bodyPr"))

        self.assertEqual(body_pr.get("wrap"), "none")

    def test_single_line_wrap_contract_survives_read_write_roundtrip(self):
        self._write(source_line_count=1, word_wrap=False)
        pages = _PPTX_IO.read(self._out, canvas_width=1280, canvas_height=720)
        text = next(
            element
            for element in pages[0]["elements"]
            if element.get("type") == "text"
        )
        self.assertFalse(text["props"]["wordWrap"])

        roundtrip = os.path.join(self._tmp.name, "roundtrip.pptx")
        _PPTX_IO.write(pages, roundtrip, canvas_width=1280, canvas_height=720)
        with zipfile.ZipFile(roundtrip) as archive:
            root = etree.fromstring(archive.read("ppt/slides/slide1.xml"))
        body_pr = next(root.iter(f"{{{A}}}bodyPr"))
        self.assertEqual(body_pr.get("wrap"), "none")

    def test_writer_uses_portable_latin_and_cjk_typefaces_without_embedding(self):
        self._write(source_line_count=1, word_wrap=False)
        root = self._slide_root()
        run_props = next(root.iter(f"{{{A}}}rPr"))

        self.assertEqual(run_props.find(f"{{{A}}}latin").get("typeface"), "Arial")
        self.assertEqual(run_props.find(f"{{{A}}}ea").get("typeface"), "Microsoft YaHei")
        with zipfile.ZipFile(self._out) as archive:
            self.assertFalse(any(name.startswith("ppt/fonts/") for name in archive.namelist()))

    def test_embedded_web_font_stack_keeps_embedded_primary_typeface(self):
        self._write(
            source_line_count=1,
            word_wrap=False,
            font_family="Inter, Arial, sans-serif",
            font_meta={
                "embedded_fonts": [{
                    "name": "Inter",
                    "style": "normal",
                    "format": "ttf",
                    "data_base64": base64.b64encode(
                        b"\x00\x01\x00\x00" + (b"\x00" * 32)
                    ).decode("ascii"),
                }],
            },
        )
        root = self._slide_root()
        run_props = next(root.iter(f"{{{A}}}rPr"))
        self.assertEqual(run_props.find(f"{{{A}}}latin").get("typeface"), "Inter")
        with zipfile.ZipFile(self._out) as archive:
            self.assertTrue(any(name.startswith("ppt/fonts/") for name in archive.namelist()))

    def test_post_export_lint_accepts_reflow_safe_output(self):
        self._write(source_line_count=1, word_wrap=False)
        self.assertEqual(_PPTX_IO.lint_pptx_text_fidelity(self._out), [])

    def test_post_export_lint_reports_unavailable_font_and_single_line_wrap(self):
        from pptx import Presentation
        from pptx.oxml.ns import qn

        self._write(source_line_count=1, word_wrap=False)
        presentation = Presentation(self._out)
        shape = next(s for s in presentation.slides[0].shapes if s.has_text_frame)
        shape.text_frame.word_wrap = True
        run = shape.text_frame.paragraphs[0].runs[0]
        latin = run._r.get_or_add_rPr().find(qn("a:latin"))
        self.assertIsNotNone(latin)
        latin.set("typeface", "Inter")
        presentation.save(self._out)

        codes = {
            issue["code"]
            for issue in _PPTX_IO.lint_pptx_text_fidelity(self._out)
        }
        self.assertIn("unavailable-web-font", codes)
        self.assertIn("single-line-wrap-enabled", codes)

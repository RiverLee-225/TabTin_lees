"""
gradFill 写出回归（Wave 2）：

PPTElement.props.gradient 来自 dom_extractor 解析 CSS linear-gradient 的产物。
pptx_io 的 _apply_gradient_fill 必须把这个字段稳定地写成 OOXML <a:gradFill>。

覆盖：
  1. linear 渐变 + 纯 RGB（无 alpha）→ 写出 <a:lin ang scaled="0"/> + 2 个 <a:gs>
  2. linear 渐变 + 含 alpha (#RRGGBBAA) → <a:srgbClr> 下挂 <a:alpha>
  3. 多 stop（≥3）渐变 → gsLst 含全部 stop，pos 单调递增
  4. radial 渐变 → 写出 <a:path path="circle"/> + <a:fillToRect>
  5. radial 自定义 center → fillToRect 的 l/t/r/b 反映中心
  6. 端到端 _write_shape_element：props.gradient 优先级高于 props.fill
  7. write() + zipfile 解包：实际 PPTX 内 slide1.xml 含 <a:gradFill>
  8. 旋转角度：rotate=90deg → ang=5400000（OOXML 1/60000 度单位）
"""

from __future__ import annotations

import importlib.util
import os
import tempfile
import zipfile
from pathlib import Path
from unittest import TestCase

from lxml import etree
from pptx import Presentation
from pptx.util import Emu


_MODULE_PATH = Path(__file__).resolve().parents[1] / "services" / "pptx_io.py"
_SPEC = importlib.util.spec_from_file_location("tabslide_pptx_io_gradfill_test", _MODULE_PATH)
if _SPEC is None or _SPEC.loader is None:
    raise RuntimeError(f"Failed to load module spec from {_MODULE_PATH}")
_PPTX_IO = importlib.util.module_from_spec(_SPEC)
_SPEC.loader.exec_module(_PPTX_IO)

_apply_gradient_fill = _PPTX_IO._apply_gradient_fill
_write_shape_element = _PPTX_IO._write_shape_element
_find_sp_pr = _PPTX_IO._find_sp_pr
write = _PPTX_IO.write

NSMAP_A = "http://schemas.openxmlformats.org/drawingml/2006/main"


def _make_blank_slide():
    """构造一张白板 slide，返回 (presentation, slide)。"""
    prs = Presentation()
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)
    return prs, slide


def _add_rect_shape(slide):
    from pptx.enum.shapes import MSO_SHAPE
    return slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Emu(914400), Emu(914400), Emu(2743200), Emu(1828800))


def _grad_fill_of(shape):
    sp_pr = _find_sp_pr(shape._element)
    assert sp_pr is not None
    return sp_pr.find(f"{{{NSMAP_A}}}gradFill")


def _solid_fill_of(shape):
    sp_pr = _find_sp_pr(shape._element)
    assert sp_pr is not None
    return sp_pr.find(f"{{{NSMAP_A}}}solidFill")


class GradientFillUnitTests(TestCase):
    """单元测试：直接调用 _apply_gradient_fill。"""

    def test_linear_two_stop_pure_rgb(self):
        """最常见场景：2 stop 线性渐变，纯 RGB，无 alpha。"""
        _, slide = _make_blank_slide()
        shape = _add_rect_shape(slide)

        gradient = {
            "type": "linear",
            "rotate": 90,  # 上→下
            "colors": [
                {"pos": 0, "color": "#FF0000"},
                {"pos": 1.0, "color": "#0000FF"},
            ],
        }
        _apply_gradient_fill(shape, gradient)

        grad_fill = _grad_fill_of(shape)
        self.assertIsNotNone(grad_fill, "应写出 <a:gradFill>")
        # 2 个 gs
        gs_list = grad_fill.findall(f"{{{NSMAP_A}}}gsLst/{{{NSMAP_A}}}gs")
        self.assertEqual(len(gs_list), 2)
        self.assertEqual(gs_list[0].get("pos"), "0")
        self.assertEqual(gs_list[1].get("pos"), "100000")

        # srgbClr val 大写无 #
        srgb_first = gs_list[0].find(f"{{{NSMAP_A}}}srgbClr")
        self.assertEqual(srgb_first.get("val"), "FF0000")
        srgb_last = gs_list[1].find(f"{{{NSMAP_A}}}srgbClr")
        self.assertEqual(srgb_last.get("val"), "0000FF")
        # 无 alpha 子节点
        self.assertIsNone(srgb_first.find(f"{{{NSMAP_A}}}alpha"))
        self.assertIsNone(srgb_last.find(f"{{{NSMAP_A}}}alpha"))

        # linear 写 <a:lin ang scaled="0"/>，90° → 5400000
        lin = grad_fill.find(f"{{{NSMAP_A}}}lin")
        self.assertIsNotNone(lin)
        self.assertEqual(lin.get("ang"), "5400000")
        self.assertEqual(lin.get("scaled"), "0")

        # 不应残留 solidFill
        self.assertIsNone(_solid_fill_of(shape))

    def test_linear_with_alpha_rrggbbaa(self):
        """含 alpha 的 stop（#RRGGBBAA）→ <a:srgbClr> 下挂 <a:alpha val="..."/>。"""
        _, slide = _make_blank_slide()
        shape = _add_rect_shape(slide)

        gradient = {
            "type": "linear",
            "rotate": 0,
            "colors": [
                {"pos": 0, "color": "#FF0000"},
                # AA = 0xCC = 204/255 ≈ 0.8 → alpha val = 80000
                {"pos": 1.0, "color": "#0000FFCC"},
            ],
        }
        _apply_gradient_fill(shape, gradient)

        grad_fill = _grad_fill_of(shape)
        gs_list = grad_fill.findall(f"{{{NSMAP_A}}}gsLst/{{{NSMAP_A}}}gs")
        self.assertEqual(len(gs_list), 2)

        # 第二个 stop 含 alpha
        srgb = gs_list[1].find(f"{{{NSMAP_A}}}srgbClr")
        alpha_el = srgb.find(f"{{{NSMAP_A}}}alpha")
        self.assertIsNotNone(alpha_el)
        # 0.8 * 100000 = 80000，允许小幅浮点误差（_parse_css_color round 到 3 位）
        alpha_val = int(alpha_el.get("val"))
        self.assertAlmostEqual(alpha_val, 80000, delta=200)

        # 0° 旋转 → ang=0
        lin = grad_fill.find(f"{{{NSMAP_A}}}lin")
        self.assertEqual(lin.get("ang"), "0")

    def test_linear_three_stop(self):
        """3-stop 渐变：颜色全部进入 gsLst，pos 单调。"""
        _, slide = _make_blank_slide()
        shape = _add_rect_shape(slide)

        gradient = {
            "type": "linear",
            "rotate": 45,
            "colors": [
                {"pos": 0, "color": "#FF0000"},
                {"pos": 0.5, "color": "#00FF00"},
                {"pos": 1.0, "color": "#0000FF"},
            ],
        }
        _apply_gradient_fill(shape, gradient)

        grad_fill = _grad_fill_of(shape)
        gs_list = grad_fill.findall(f"{{{NSMAP_A}}}gsLst/{{{NSMAP_A}}}gs")
        self.assertEqual(len(gs_list), 3)
        self.assertEqual(gs_list[0].get("pos"), "0")
        self.assertEqual(gs_list[1].get("pos"), "50000")
        self.assertEqual(gs_list[2].get("pos"), "100000")

        srgb_vals = [gs.find(f"{{{NSMAP_A}}}srgbClr").get("val") for gs in gs_list]
        self.assertEqual(srgb_vals, ["FF0000", "00FF00", "0000FF"])

        # 45° → 45 * 60000 = 2700000
        lin = grad_fill.find(f"{{{NSMAP_A}}}lin")
        self.assertEqual(lin.get("ang"), "2700000")

    def test_radial_default_center(self):
        """radial 渐变：写 <a:path path="circle"/> + <a:fillToRect>，默认中心 50/50。"""
        _, slide = _make_blank_slide()
        shape = _add_rect_shape(slide)

        gradient = {
            "type": "radial",
            "rotate": 0,
            "colors": [
                {"pos": 0, "color": "#FFFFFF"},
                {"pos": 1.0, "color": "#000000"},
            ],
        }
        _apply_gradient_fill(shape, gradient)

        grad_fill = _grad_fill_of(shape)
        # radial 不应有 lin
        self.assertIsNone(grad_fill.find(f"{{{NSMAP_A}}}lin"))

        path_el = grad_fill.find(f"{{{NSMAP_A}}}path")
        self.assertIsNotNone(path_el)
        self.assertEqual(path_el.get("path"), "circle")

        ftr = path_el.find(f"{{{NSMAP_A}}}fillToRect")
        self.assertIsNotNone(ftr)
        self.assertEqual(ftr.get("l"), "50000")
        self.assertEqual(ftr.get("t"), "50000")
        self.assertEqual(ftr.get("r"), "50000")
        self.assertEqual(ftr.get("b"), "50000")

    def test_radial_custom_center(self):
        """radial 渐变 + 自定义中心 → fillToRect l/t/r/b 反映 center。"""
        _, slide = _make_blank_slide()
        shape = _add_rect_shape(slide)

        gradient = {
            "type": "radial",
            "rotate": 0,
            "center": {"x": 0.25, "y": 0.75},
            "colors": [
                {"pos": 0, "color": "#FFAA00"},
                {"pos": 1.0, "color": "#0044AA"},
            ],
        }
        _apply_gradient_fill(shape, gradient)

        grad_fill = _grad_fill_of(shape)
        path_el = grad_fill.find(f"{{{NSMAP_A}}}path")
        ftr = path_el.find(f"{{{NSMAP_A}}}fillToRect")
        # center = (0.25, 0.75) → l=25000, t=75000, r=75000, b=25000
        self.assertEqual(ftr.get("l"), "25000")
        self.assertEqual(ftr.get("t"), "75000")
        self.assertEqual(ftr.get("r"), "75000")
        self.assertEqual(ftr.get("b"), "25000")

    def test_replaces_existing_solid_fill(self):
        """已有 solidFill 被 gradFill 替换；保证 spPr 下只剩一个填充节点。"""
        from pptx.dml.color import RGBColor

        _, slide = _make_blank_slide()
        shape = _add_rect_shape(slide)
        # 先施加一个 solid fill
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor.from_string("AABBCC")
        self.assertIsNotNone(_solid_fill_of(shape))

        gradient = {
            "type": "linear",
            "rotate": 0,
            "colors": [
                {"pos": 0, "color": "#FF0000"},
                {"pos": 1.0, "color": "#0000FF"},
            ],
        }
        _apply_gradient_fill(shape, gradient)

        sp_pr = _find_sp_pr(shape._element)
        # solidFill 应被清掉
        self.assertIsNone(sp_pr.find(f"{{{NSMAP_A}}}solidFill"))
        self.assertIsNotNone(sp_pr.find(f"{{{NSMAP_A}}}gradFill"))


class GradientFillShapeIntegrationTests(TestCase):
    """端到端：通过 _write_shape_element 走 props.gradient 分支。"""

    def test_write_shape_element_picks_gradient_over_fill(self):
        """props.gradient 存在时优先级高于 props.fill。"""
        _, slide = _make_blank_slide()

        props = {
            "viewBox": [800, 600],
            "path": "M 0 0 L 800 0 L 800 600 L 0 600 Z",
            "pptxShapeType": "rect",
            "fill": "#AABBCC",  # 应被 gradient 抑制
            "gradient": {
                "type": "linear",
                "rotate": 90,
                "colors": [
                    {"pos": 0, "color": "#FF0000"},
                    {"pos": 1.0, "color": "#0000FFCC"},
                ],
            },
        }
        shape = _write_shape_element(slide, props, Emu(914400), Emu(914400), Emu(2743200), Emu(1828800))

        self.assertIsNone(_solid_fill_of(shape))
        grad_fill = _grad_fill_of(shape)
        self.assertIsNotNone(grad_fill)

        gs_list = grad_fill.findall(f"{{{NSMAP_A}}}gsLst/{{{NSMAP_A}}}gs")
        self.assertEqual(len(gs_list), 2)

        # 含 alpha 的 stop
        last_srgb = gs_list[-1].find(f"{{{NSMAP_A}}}srgbClr")
        self.assertEqual(last_srgb.get("val"), "0000FF")
        self.assertIsNotNone(last_srgb.find(f"{{{NSMAP_A}}}alpha"))


class GradientFillEndToEndTests(TestCase):
    """端到端：write() 走完整管线，解 PPTX 看 slide1.xml 真实输出。"""

    def test_write_produces_gradfill_in_slide_xml(self):
        """走 write()，解 PPTX，slide1.xml 必须含 <a:gradFill>。"""
        page = {
            "id": "page1",
            "elements": [
                {
                    "id": "shape1",
                    "type": "shape",
                    "x": 100,
                    "y": 100,
                    "width": 800,
                    "height": 600,
                    "rotate": 0,
                    "opacity": 1,
                    "locked": False,
                    "visible": True,
                    "zIndex": 0,
                    "props": {
                        "viewBox": [800, 600],
                        "path": "M 0 0 L 800 0 L 800 600 L 0 600 Z",
                        "pptxShapeType": "rect",
                        "gradient": {
                            "type": "linear",
                            "rotate": 90,
                            "colors": [
                                {"pos": 0, "color": "#FF0000"},
                                {"pos": 1.0, "color": "#0000FFCC"},
                            ],
                        },
                    },
                }
            ],
        }

        with tempfile.NamedTemporaryFile(suffix=".pptx", delete=False) as tmp:
            tmp_path = tmp.name
        try:
            write(pages=[page], output_path=tmp_path)
            with zipfile.ZipFile(tmp_path, "r") as zf:
                slide_xml = zf.read("ppt/slides/slide1.xml")
            tree = etree.fromstring(slide_xml)
            grad_fills = tree.findall(f".//{{{NSMAP_A}}}gradFill")
            self.assertGreaterEqual(len(grad_fills), 1, "slide1.xml 应含至少一个 <a:gradFill>")

            # 至少 1 个 gradFill 含 lin ang=5400000
            lin_angles = [
                lin.get("ang")
                for gf in grad_fills
                for lin in gf.findall(f"{{{NSMAP_A}}}lin")
            ]
            self.assertIn("5400000", lin_angles)

            # 至少 1 个 stop 含 alpha
            alphas = tree.findall(f".//{{{NSMAP_A}}}gradFill//{{{NSMAP_A}}}alpha")
            self.assertGreaterEqual(len(alphas), 1)
        finally:
            try:
                os.unlink(tmp_path)
            except OSError:
                pass


class AIGCMetadataInjectionTests(TestCase):
    """AIGC 溯源：write() 接受 aigc_metadata 时 docProps/custom.xml 要被注入。"""

    def _make_minimal_page(self):
        return {
            "id": "page1",
            "elements": [
                {
                    "id": "shape1",
                    "type": "shape",
                    "x": 100, "y": 100, "width": 200, "height": 100,
                    "rotate": 0, "opacity": 1, "locked": False, "visible": True, "zIndex": 0,
                    "props": {
                        "viewBox": [200, 100],
                        "path": "M 0 0 L 200 0 L 200 100 L 0 100 Z",
                        "pptxShapeType": "rect",
                        "fill": "#AABBCC",
                    },
                }
            ],
        }

    def test_inject_writes_custom_xml(self):
        with tempfile.NamedTemporaryFile(suffix=".pptx", delete=False) as tmp:
            tmp_path = tmp.name
        try:
            write(
                pages=[self._make_minimal_page()],
                output_path=tmp_path,
                aigc_metadata={
                    "projectId": "proj-uuid-001",
                    "organizationId": "wt-uuid-001",
                    "spaceId": "sp-uuid-001",
                    "name": "测试演示文稿",
                },
            )

            with zipfile.ZipFile(tmp_path, "r") as zf:
                names = zf.namelist()
                self.assertIn("docProps/custom.xml", names)
                custom_xml = zf.read("docProps/custom.xml")
                ct_xml = zf.read("[Content_Types].xml")
                rels_xml = zf.read("_rels/.rels")

            # custom.xml 含 TabTin property + JSON 载荷
            self.assertIn(b"TabTin", custom_xml)
            self.assertIn(b"proj-uuid-001", custom_xml)
            self.assertIn(b"wt-uuid-001", custom_xml)
            self.assertIn(b"sp-uuid-001", custom_xml)
            # XML 结构有效
            root = etree.fromstring(custom_xml)
            ns_cp = "http://schemas.openxmlformats.org/officeDocument/2006/custom-properties"
            ns_vt = "http://schemas.openxmlformats.org/officeDocument/2006/docPropsVTypes"
            props = root.findall(f"{{{ns_cp}}}property")
            tabtin_props = [p for p in props if p.get("name") == "TabTin"]
            self.assertEqual(len(tabtin_props), 1)
            self.assertEqual(
                tabtin_props[0].get("fmtid"),
                "{D5CDD505-2E9C-101B-9397-08002B2CF9AE}",
            )
            lpwstr = tabtin_props[0].find(f"{{{ns_vt}}}lpwstr")
            self.assertIsNotNone(lpwstr)
            self.assertIn("TabTin", lpwstr.text)
            self.assertIn("proj-uuid-001", lpwstr.text)

            # Content_Types.xml 含 custom-properties Override
            self.assertIn(b"custom-properties+xml", ct_xml)
            self.assertIn(b"/docProps/custom.xml", ct_xml)

            # _rels/.rels 含 custom-properties relationship
            self.assertIn(b"custom-properties", rels_xml)
            self.assertIn(b"docProps/custom.xml", rels_xml)
        finally:
            try:
                os.unlink(tmp_path)
            except OSError:
                pass

    def test_no_metadata_means_no_injection(self):
        """没传 aigc_metadata 时，PPTX 不应被注入（或至少不报错）。"""
        with tempfile.NamedTemporaryFile(suffix=".pptx", delete=False) as tmp:
            tmp_path = tmp.name
        try:
            write(pages=[self._make_minimal_page()], output_path=tmp_path)
            with zipfile.ZipFile(tmp_path, "r") as zf:
                names = zf.namelist()
                if "docProps/custom.xml" in names:
                    custom_xml = zf.read("docProps/custom.xml")
                    self.assertNotIn(b"TabTin", custom_xml)
        finally:
            try:
                os.unlink(tmp_path)
            except OSError:
                pass

    def test_idempotent_update(self):
        """同一份文件被注入 2 次：name='TabTin' 的 property 仍只有 1 个。"""
        with tempfile.NamedTemporaryFile(suffix=".pptx", delete=False) as tmp:
            tmp_path = tmp.name
        try:
            write(
                pages=[self._make_minimal_page()],
                output_path=tmp_path,
                aigc_metadata={"projectId": "p1", "organizationId": "w1", "spaceId": "s1"},
            )
            # 二次注入（直接调底层函数）
            _PPTX_IO._inject_aigc_metadata(tmp_path, {"projectId": "p2", "organizationId": "w2", "spaceId": "s2"})

            with zipfile.ZipFile(tmp_path, "r") as zf:
                custom_xml = zf.read("docProps/custom.xml")
            root = etree.fromstring(custom_xml)
            ns_cp = "http://schemas.openxmlformats.org/officeDocument/2006/custom-properties"
            tabtin_props = [
                p for p in root.findall(f"{{{ns_cp}}}property") if p.get("name") == "TabTin"
            ]
            self.assertEqual(len(tabtin_props), 1)
            # 内容应是后写入的
            ns_vt = "http://schemas.openxmlformats.org/officeDocument/2006/docPropsVTypes"
            self.assertIn("p2", tabtin_props[0].find(f"{{{ns_vt}}}lpwstr").text)
        finally:
            try:
                os.unlink(tmp_path)
            except OSError:
                pass

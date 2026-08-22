import importlib.util
import os
import tempfile
from pathlib import Path
from unittest import TestCase

from pptx import Presentation

# 直接按文件路径加载，避免触发 apps.tabslide.services.__init__ 的 Django 依赖
_MODULE_PATH = Path(__file__).resolve().parents[1] / "services" / "pptx_io.py"
_SPEC = importlib.util.spec_from_file_location("tabslide_pptx_theme_test_module", _MODULE_PATH)
if _SPEC is None or _SPEC.loader is None:
    raise RuntimeError(f"Failed to load module spec from {_MODULE_PATH}")
_PPTX_IO = importlib.util.module_from_spec(_SPEC)
_SPEC.loader.exec_module(_PPTX_IO)

_build_theme_payload = _PPTX_IO._build_theme_payload
extract_theme_payload = _PPTX_IO.extract_theme_payload
read = _PPTX_IO.read
write = _PPTX_IO.write


class TestPptxThemeChain(TestCase):
    def test_build_theme_payload_maps_colors_and_fonts(self):
        theme_color_map = {
            "1": "#112233",
            "2": "#fefefe",
            "5": "#aa0001",
            "6": "#aa0002",
            "7": "#aa0003",
            "8": "#aa0004",
            "9": "#aa0005",
            "10": "#aa0006",
        }
        theme_fonts = {
            "major_latin": "Aptos Display",
            "minor_ea": "微软雅黑",
        }

        payload = _build_theme_payload(theme_color_map, theme_fonts)

        self.assertEqual(payload.get("backgroundColor"), "#fefefe")
        self.assertEqual(payload.get("fontColor"), "#112233")
        self.assertEqual(
            payload.get("themeColors"),
            ["#aa0001", "#aa0002", "#aa0003", "#aa0004", "#aa0005", "#aa0006"],
        )
        self.assertEqual(payload.get("fontName"), "微软雅黑")
        self.assertEqual(payload.get("headingFontName"), "Aptos Display")

    def test_build_theme_payload_fallbacks_are_stable(self):
        payload = _build_theme_payload({}, {})

        self.assertEqual(str(payload.get("backgroundColor", "")).lower(), "#ffffff")
        self.assertEqual(str(payload.get("fontColor", "")).lower(), "#000000")
        self.assertEqual(len(payload.get("themeColors", [])), 6)
        self.assertEqual(payload.get("fontName"), "Microsoft YaHei")
        self.assertEqual(payload.get("headingFontName"), "Microsoft YaHei")

    def test_extract_theme_payload_from_real_pptx(self):
        fd, pptx_path = tempfile.mkstemp(suffix=".pptx")
        os.close(fd)
        try:
            prs = Presentation()
            prs.slides.add_slide(prs.slide_layouts[6])
            prs.save(pptx_path)

            payload = extract_theme_payload(pptx_path)
            self.assertIn("backgroundColor", payload)
            self.assertIn("fontColor", payload)
            self.assertIn("themeColors", payload)
            self.assertEqual(len(payload.get("themeColors", [])), 6)
            self.assertTrue(payload.get("fontName"))
            self.assertTrue(payload.get("headingFontName"))
        finally:
            try:
                os.remove(pptx_path)
            except Exception:
                pass

    def test_shape_line_table_theme_token_roundtrip(self):
        pages = [
            {
                "id": "page-1",
                "background": {"type": "color", "value": "#ffffff"},
                "notes": "",
                "elements": [
                    {
                        "id": "shape-theme-1",
                        "type": "shape",
                        "x": 80,
                        "y": 80,
                        "width": 260,
                        "height": 140,
                        "zIndex": 0,
                        "props": {
                            "viewBox": [260, 140],
                            "path": "M 0 0 L 260 0 L 260 140 L 0 140 Z",
                            "pptxShapeType": "rect",
                            "fill": "#4472C4",
                            "fillThemeKey": "accent1",
                            "outline": {
                                "style": "solid",
                                "width": 2,
                                "color": "#ED7D31",
                                "themeKey": "accent2",
                            },
                        },
                    },
                    {
                        "id": "line-theme-1",
                        "type": "line",
                        "x": 420,
                        "y": 120,
                        "width": 260,
                        "height": 80,
                        "zIndex": 1,
                        "props": {
                            "start": [0, 0],
                            "end": [260, 80],
                            "style": "solid",
                            "color": "#A5A5A5",
                            "colorThemeKey": "accent3",
                            "lineWidth": 3,
                            "points": ["", ""],
                        },
                    },
                    {
                        "id": "table-theme-1",
                        "type": "table",
                        "x": 80,
                        "y": 280,
                        "width": 360,
                        "height": 120,
                        "zIndex": 2,
                        "props": {
                            "data": [[{
                                "text": "Theme Cell",
                                "colspan": 1,
                                "rowspan": 1,
                                "style": {
                                    "bgColor": "#FFC000",
                                    "bgColorThemeKey": "accent4",
                                    "color": "#70AD47",
                                    "colorThemeKey": "accent6",
                                },
                            }]],
                            "colWidths": [1],
                            "cellMinHeight": 36,
                            "outline": {
                                "style": "solid",
                                "width": 1.5,
                                "color": "#5B9BD5",
                                "themeKey": "accent5",
                            },
                            "theme": {
                                "headerRow": True,
                                "color": "#5B9BD5",
                                "colorThemeKey": "accent5",
                            },
                        },
                    },
                ],
            },
        ]

        fd, pptx_path = tempfile.mkstemp(suffix=".pptx")
        os.close(fd)
        try:
            write(
                pages=pages,
                output_path=pptx_path,
                canvas_width=1920,
                canvas_height=1080,
            )
            out_pages = read(
                pptx_path=pptx_path,
                canvas_width=1920,
                canvas_height=1080,
            )
            self.assertEqual(len(out_pages), 1)
            out_elements = out_pages[0].get("elements", [])

            shape_el = next((el for el in out_elements if el.get("type") == "shape"), None)
            self.assertIsNotNone(shape_el)
            self.assertEqual(shape_el["props"].get("fillThemeKey"), "accent1")
            self.assertEqual(shape_el["props"].get("outline", {}).get("themeKey"), "accent2")

            line_el = next((el for el in out_elements if el.get("type") == "line"), None)
            self.assertIsNotNone(line_el)
            self.assertEqual(line_el["props"].get("colorThemeKey"), "accent3")

            table_el = next((el for el in out_elements if el.get("type") == "table"), None)
            self.assertIsNotNone(table_el)
            self.assertEqual(table_el["props"].get("outline", {}).get("themeKey"), "accent5")
            first_cell = table_el["props"]["data"][0][0]
            style = first_cell.get("style") or {}
            self.assertEqual(style.get("bgColorThemeKey"), "accent4")
            self.assertEqual(style.get("colorThemeKey"), "accent6")
        finally:
            try:
                os.remove(pptx_path)
            except Exception:
                pass

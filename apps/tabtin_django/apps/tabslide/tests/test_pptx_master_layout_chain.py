import importlib.util
import os
import tempfile
from pathlib import Path
from unittest import TestCase
from unittest.mock import patch

from pptx import Presentation

# 直接按文件路径加载，避免触发 apps.tabslide.services.__init__ 的 Django 依赖
_MODULE_PATH = Path(__file__).resolve().parents[1] / "services" / "pptx_io.py"
_SPEC = importlib.util.spec_from_file_location("tabslide_pptx_master_layout_test_module", _MODULE_PATH)
if _SPEC is None or _SPEC.loader is None:
    raise RuntimeError(f"Failed to load module spec from {_MODULE_PATH}")
_PPTX_IO = importlib.util.module_from_spec(_SPEC)
_SPEC.loader.exec_module(_PPTX_IO)

read = _PPTX_IO.read
write = _PPTX_IO.write


class TestPptxMasterLayoutChain(TestCase):
    def _temp_pptx_path(self) -> str:
        fd, path = tempfile.mkstemp(suffix=".pptx")
        os.close(fd)
        return path

    def test_read_extracts_layout_meta_and_placeholder_semantics(self):
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[0])  # Title Slide
        slide.shapes.title.text = "主标题"
        slide.placeholders[1].text = "副标题"

        path = self._temp_pptx_path()
        try:
            prs.save(path)
            pages = read(path, canvas_width=1920, canvas_height=1080)
        finally:
            try:
                os.remove(path)
            except OSError:
                pass

        self.assertEqual(len(pages), 1)
        page = pages[0]
        self.assertIn("layout", page)
        self.assertEqual(page["layout"].get("index"), 0)

        text_elements = [el for el in page.get("elements", []) if el.get("type") == "text"]
        self.assertGreaterEqual(len(text_elements), 2)

        title_el = next(
            (el for el in text_elements if "主标题" in str(el.get("props", {}).get("content", ""))),
            None,
        )
        self.assertIsNotNone(title_el)
        assert title_el is not None
        title_ph = title_el.get("props", {}).get("placeholder", {})
        self.assertIn(title_ph.get("type"), ("title", "ctrTitle"))
        self.assertEqual(title_el.get("props", {}).get("textType"), "title")

        subtitle_el = next(
            (el for el in text_elements if "副标题" in str(el.get("props", {}).get("content", ""))),
            None,
        )
        self.assertIsNotNone(subtitle_el)
        assert subtitle_el is not None
        subtitle_ph = subtitle_el.get("props", {}).get("placeholder", {})
        self.assertIn(subtitle_ph.get("type"), ("subTitle", "body", "obj"))
        self.assertIn(subtitle_el.get("props", {}).get("textType"), ("subtitle", "content"))

    def test_write_uses_layout_and_fills_title_placeholder(self):
        output_path = self._temp_pptx_path()
        try:
            pages = [
                {
                    "id": "page-1",
                    "layout": {"index": 0, "name": "Title Slide"},
                    "background": {"type": "color", "value": "#ffffff"},
                    "elements": [
                        {
                            "id": "text-1",
                            "type": "text",
                            "x": 120,
                            "y": 80,
                            "width": 800,
                            "height": 120,
                            "props": {
                                "content": "<p>导出回填标题</p>",
                                "defaultFontFamily": "Calibri",
                                "defaultColor": "#111111",
                                "placeholder": {"type": "title", "idx": 0},
                                "textType": "title",
                            },
                        }
                    ],
                    "notes": "",
                }
            ]

            write(
                pages=pages,
                output_path=output_path,
                canvas_width=1920,
                canvas_height=1080,
            )

            prs = Presentation(output_path)
            slide = prs.slides[0]
            self.assertEqual(slide.slide_layout.name, prs.slide_layouts[0].name)
            self.assertEqual(slide.shapes.title.text, "导出回填标题")

            shape_names = [str(getattr(sp, "name", "") or "") for sp in slide.shapes]
            self.assertTrue(all("TextBox" not in name for name in shape_names))
        finally:
            try:
                os.remove(output_path)
            except OSError:
                pass

    def test_extract_master_layer_elements_filters_placeholder_and_locks(self):
        class DummyShape:
            def __init__(self, name: str):
                self.name = name

        class DummyMaster:
            def __init__(self, shapes):
                self.shapes = shapes

        class DummyLayout:
            def __init__(self, shapes, master):
                self.shapes = shapes
                self.slide_master = master

        class DummySlide:
            def __init__(self, layout):
                self.slide_layout = layout

        layout_shape = DummyShape("layout-shape")
        placeholder_shape = DummyShape("layout-placeholder")
        master_shape = DummyShape("master-shape")

        slide = DummySlide(
            DummyLayout(
                [layout_shape, placeholder_shape],
                DummyMaster([master_shape]),
            )
        )

        def fake_extract_placeholder_meta(shape):
            if getattr(shape, "name", "") == "layout-placeholder":
                return {"type": "body", "idx": 1}
            return None

        def fake_extract_shape(shape, **kwargs):
            return {
                "id": f"id-{shape.name}",
                "type": "shape",
                "x": 0,
                "y": 0,
                "width": 100,
                "height": 60,
                "rotate": 0,
                "zIndex": kwargs.get("z_index", 0),
                "props": {"fill": "#ffffff"},
            }

        with patch.object(_PPTX_IO, "_extract_placeholder_meta", side_effect=fake_extract_placeholder_meta), patch.object(
            _PPTX_IO,
            "_extract_shape",
            side_effect=fake_extract_shape,
        ):
            elements = _PPTX_IO._extract_master_layer_elements(  # pylint: disable=protected-access
                slide,
                slide_width_emu=9144000,
                slide_height_emu=5143500,
                canvas_width=1920,
                canvas_height=1080,
            )

        self.assertEqual(len(elements), 2)
        self.assertEqual([el.get("id") for el in elements], ["id-master-shape", "id-layout-shape"])
        self.assertTrue(all(el.get("locked") is True for el in elements))
        self.assertEqual([el.get("zIndex") for el in elements], [0, 1])
        self.assertEqual(elements[0].get("props", {}).get("masterSource"), "master")
        self.assertEqual(elements[1].get("props", {}).get("masterSource"), "layout")

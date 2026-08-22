import importlib.util
import os
import tempfile
from pathlib import Path
from unittest import TestCase

from pptx import Presentation
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

# 直接按文件路径加载，避免触发 apps.tabslide.services.__init__ 的 Django 依赖
_MODULE_PATH = Path(__file__).resolve().parents[1] / "services" / "pptx_io.py"
_SPEC = importlib.util.spec_from_file_location("tabslide_pptx_io_text_test_module", _MODULE_PATH)
if _SPEC is None or _SPEC.loader is None:
    raise RuntimeError(f"Failed to load module spec from {_MODULE_PATH}")
_PPTX_IO = importlib.util.module_from_spec(_SPEC)
_SPEC.loader.exec_module(_PPTX_IO)

_parse_inline_style = _PPTX_IO._parse_inline_style
_parse_paragraph_html = _PPTX_IO._parse_paragraph_html
_paragraph_to_html_v2 = _PPTX_IO._paragraph_to_html_v2
_write_text_element = _PPTX_IO._write_text_element
_write_shape_element = _PPTX_IO._write_shape_element
_extract_auto_shape = _PPTX_IO._extract_auto_shape
read = _PPTX_IO.read
write = _PPTX_IO.write


class TestPptxTextChain(TestCase):
    def test_parse_and_write_br_as_soft_line_break(self):
        """HTML <br> 应解析为 break run，并写入 OOXML a:br。"""
        runs = _parse_paragraph_html(
            '<span style="color:">$</span> tabtin code open<br>'
            '<span style="color:">$</span> tabtin space create'
        )
        self.assertTrue(any(r.get("break") or r.get("text") == "\n" for r in runs))
        texts = [r.get("text", "") for r in runs if not r.get("break")]
        joined = "".join(texts)
        self.assertIn("tabtin code open", joined)
        self.assertIn("tabtin space create", joined)

        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        shape = _write_text_element(
            slide=slide,
            props={
                "content": (
                    '<p><span style="color:">$</span> tabtin code open<br>'
                    '<span style="color:">$</span> tabtin space create</p>'
                ),
                "defaultFontFamily": "monospace",
                "defaultFontSize": 12,
                "defaultColor": "#F5A830",
            },
            left=Inches(1),
            top=Inches(1),
            width=Inches(6),
            height=Inches(2),
        )
        nsmap_a = "http://schemas.openxmlformats.org/drawingml/2006/main"
        brs = shape._element.findall(f".//{{{nsmap_a}}}br")
        self.assertGreaterEqual(len(brs), 1, "应写入至少一个 a:br 软换行")

    def test_parse_inline_style_supports_px_pt_and_single_quote_style(self):
        styles = _parse_inline_style(
            "style='font-size:16px;letter-spacing:1.5pt;color:rgba(17,34,51,0.5);font-family:\"Calibri\", \"Segoe UI\", sans-serif'"
        )
        self.assertIn("fontSize", styles)
        self.assertAlmostEqual(styles["fontSize"], 12.0, places=2)  # 16px -> 12pt
        self.assertIn("letterSpacing", styles)
        self.assertAlmostEqual(styles["letterSpacing"], 2.0, places=2)  # 1.5pt -> 2px
        self.assertEqual(styles.get("color"), "rgba(17,34,51,0.5)")
        self.assertEqual(styles.get("fontFamily"), '"Calibri", "Segoe UI", sans-serif')

    def test_parse_inline_style_supports_mixed_quotes_in_font_family(self):
        styles = _parse_inline_style(
            "style=\"font-family:'Segoe UI', PingFang SC, sans-serif; font-size:18px; color:#112233\""
        )
        self.assertEqual(styles.get("fontFamily"), "'Segoe UI', PingFang SC, sans-serif")
        self.assertAlmostEqual(styles.get("fontSize", 0), 13.5, places=2)  # 18px -> 13.5pt
        self.assertEqual(styles.get("color"), "#112233")

    def test_parse_inline_style_supports_theme_color_key_attribute(self):
        styles = _parse_inline_style(
            'data-theme-color-key="accent1" style="color:#4472C4;font-size:16px"'
        )
        self.assertEqual(styles.get("themeColorKey"), "accent1")
        self.assertEqual(styles.get("color"), "#4472C4")

    def test_parse_paragraph_html_restores_outer_span_style_after_nested_override(self):
        runs = _parse_paragraph_html(
            '<span style="color:#FF0000">A<span style="color:#00FF00">B</span>C</span>'
        )
        self.assertEqual(len(runs), 3)
        self.assertEqual(runs[0].get("color"), "#FF0000")
        self.assertEqual(runs[1].get("color"), "#00FF00")
        self.assertEqual(runs[2].get("color"), "#FF0000")

    def test_write_text_element_preserves_justify_decimal_font_and_word_space(self):
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        shape = _write_text_element(
            slide=slide,
            props={
                "content": (
                    '<p style="text-align:justify;line-height:1.5;margin-top:4.5pt;margin-bottom:6.25pt">'
                    '<span style="font-size:18.5pt;color:rgba(17,34,51,0.5)">Hello</span>'
                    "</p>"
                ),
                "defaultFontFamily": "Calibri, 'Segoe UI', Arial, sans-serif",
                "defaultFontSize": 16,
                "defaultColor": "#112233",
                "wordSpace": 2,  # px
            },
            left=Inches(1),
            top=Inches(1),
            width=Inches(6),
            height=Inches(2),
        )

        tf = shape.text_frame
        p = tf.paragraphs[0]
        self.assertEqual(p.alignment, PP_ALIGN.JUSTIFY)
        self.assertAlmostEqual(float(p.line_spacing), 1.5, places=2)
        self.assertIsNotNone(p.space_before)
        self.assertIsNotNone(p.space_after)
        self.assertAlmostEqual(p.space_before.pt, 4.5, places=1)
        self.assertAlmostEqual(p.space_after.pt, 6.25, places=1)

        run = p.runs[0]
        self.assertIsNotNone(run.font.size)
        self.assertAlmostEqual(run.font.size.pt, 18.5, places=2)
        self.assertEqual(str(run.font.color.rgb), "112233")

        nsmap_a = "http://schemas.openxmlformats.org/drawingml/2006/main"
        r_pr = run._r.find(f"{{{nsmap_a}}}rPr")
        self.assertIsNotNone(r_pr)
        assert r_pr is not None
        # 2px -> 1.5pt -> spc=150 (1/100pt)
        self.assertEqual(r_pr.get("spc"), "150")

        alpha = run._r.find(
            f".//{{{nsmap_a}}}solidFill/{{{nsmap_a}}}srgbClr/{{{nsmap_a}}}alpha"
        )
        self.assertIsNotNone(alpha)
        assert alpha is not None
        self.assertEqual(alpha.get("val"), "50000")

    def test_write_read_text_roundtrip_preserves_justify_marks_and_word_space(self):
        pages = [
            {
                "id": "page-1",
                "background": {"type": "color", "value": "#ffffff"},
                "elements": [
                    {
                        "id": "text-1",
                        "type": "text",
                        "x": 100,
                        "y": 120,
                        "width": 800,
                        "height": 240,
                        "rotate": 0,
                        "zIndex": 0,
                        "props": {
                            "content": (
                                '<p style="text-align:justify">'
                                '<span style="font-size:18.5pt">A</span>'
                                "<sup>2</sup><sub>3</sub>"
                                '<mark data-color="#FFF3BF" style="background-color:#FFF3BF">M</mark>'
                                "<s>S</s>"
                                "</p>"
                            ),
                            "defaultFontFamily": "Calibri, 'Segoe UI', Arial, sans-serif",
                            "defaultFontSize": 16,
                            "defaultColor": "#112233",
                            "verticalAlign": "middle",
                            "margin": {"top": 3.6, "right": 7.2, "bottom": 3.6, "left": 7.2},
                            "wordSpace": 2,  # px
                        },
                    }
                ],
            }
        ]

        fd, pptx_path = tempfile.mkstemp(suffix=".pptx")
        os.close(fd)
        try:
            write(pages=pages, output_path=pptx_path, canvas_width=1920, canvas_height=1080)
            out_pages = read(pptx_path=pptx_path, canvas_width=1920, canvas_height=1080)
            self.assertEqual(len(out_pages), 1)

            text_elements = [el for el in out_pages[0].get("elements", []) if el.get("type") == "text"]
            self.assertEqual(len(text_elements), 1)
            props = text_elements[0].get("props", {})
            content = props.get("content", "")

            self.assertIn("text-align:justify", content)
            self.assertIn("<sup>", content)
            self.assertIn("<sub>", content)
            self.assertIn("<mark", content)
            self.assertIn("<s>", content)

            self.assertIn("wordSpace", props)
            self.assertAlmostEqual(float(props.get("wordSpace", 0)), 2.0, delta=0.3)
        finally:
            try:
                os.remove(pptx_path)
            except Exception:
                pass

    def test_shape_text_roundtrip_keeps_sup_sub_mark_and_justify(self):
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        shape = _write_shape_element(
            slide=slide,
            props={
                "pptxShapeType": "rect",
                "fill": "#ffffff",
                "text": {
                    "content": (
                        '<p style="text-align:justify">'
                        '<span style="font-size:16.5pt">A</span>'
                        "<sup>2</sup><sub>3</sub>"
                        '<mark data-color="#FFF3BF" style="background-color:#FFF3BF">M</mark>'
                        "<s>S</s>"
                        "</p>"
                    ),
                    "align": "justify",
                    "verticalAlign": "middle",
                    "defaultColor": "#112233",
                    "defaultFontSize": 16,
                    "defaultFontFamily": "Calibri",
                },
            },
            left=Inches(1),
            top=Inches(1),
            width=Inches(5),
            height=Inches(2),
        )

        extracted = _extract_auto_shape(
            shape,
            {
                "id": "shape-text-rich",
                "x": 10,
                "y": 20,
                "width": 500,
                "height": 200,
                "zIndex": 0,
            },
        )
        text_props = extracted.get("props", {}).get("text", {})
        content = text_props.get("content", "")
        self.assertEqual(text_props.get("align"), "justify")
        self.assertIn("<sup>", content)
        self.assertIn("<sub>", content)
        self.assertIn("<mark", content)
        self.assertIn("<s>", content)

    def test_write_read_text_roundtrip_keeps_decimal_paragraph_space(self):
        pages = [
            {
                "id": "page-ps",
                "background": {"type": "color", "value": "#ffffff"},
                "elements": [
                    {
                        "id": "text-ps",
                        "type": "text",
                        "x": 120,
                        "y": 140,
                        "width": 700,
                        "height": 260,
                        "rotate": 0,
                        "zIndex": 0,
                        "props": {
                            "content": "<p>Line 1</p><p>Line 2</p>",
                            "defaultFontFamily": "Calibri",
                            "defaultFontSize": 16,
                            "defaultColor": "#222222",
                            "paragraphSpace": 6.25,
                        },
                    }
                ],
            }
        ]

        fd, pptx_path = tempfile.mkstemp(suffix=".pptx")
        os.close(fd)
        try:
            write(pages=pages, output_path=pptx_path, canvas_width=1920, canvas_height=1080)
            out_pages = read(pptx_path=pptx_path, canvas_width=1920, canvas_height=1080)
            self.assertEqual(len(out_pages), 1)
            text_elements = [el for el in out_pages[0].get("elements", []) if el.get("type") == "text"]
            self.assertEqual(len(text_elements), 1)
            props = text_elements[0].get("props", {})
            self.assertIn("paragraphSpace", props)
            self.assertAlmostEqual(float(props.get("paragraphSpace", 0)), 6.25, delta=0.2)
        finally:
            try:
                os.remove(pptx_path)
            except Exception:
                pass

    def test_write_read_text_roundtrip_keeps_vertical_margin_and_valign(self):
        pages = [
            {
                "id": "page-v",
                "background": {"type": "color", "value": "#ffffff"},
                "elements": [
                    {
                        "id": "text-v",
                        "type": "text",
                        "x": 100,
                        "y": 120,
                        "width": 500,
                        "height": 260,
                        "rotate": 0,
                        "zIndex": 0,
                        "props": {
                            "content": "<p>Vertical</p>",
                            "defaultFontFamily": "Calibri",
                            "defaultFontSize": 16,
                            "defaultColor": "#333333",
                            "verticalAlign": "bottom",
                            "vertical": True,
                            "margin": {"top": 4.2, "right": 7.8, "bottom": 5.1, "left": 6.4},
                        },
                    }
                ],
            }
        ]

        fd, pptx_path = tempfile.mkstemp(suffix=".pptx")
        os.close(fd)
        try:
            write(pages=pages, output_path=pptx_path, canvas_width=1920, canvas_height=1080)
            out_pages = read(pptx_path=pptx_path, canvas_width=1920, canvas_height=1080)
            self.assertEqual(len(out_pages), 1)
            text_elements = [el for el in out_pages[0].get("elements", []) if el.get("type") == "text"]
            self.assertEqual(len(text_elements), 1)
            props = text_elements[0].get("props", {})

            self.assertEqual(props.get("verticalAlign"), "bottom")
            self.assertTrue(props.get("vertical"))
            margin = props.get("margin", {})
            self.assertAlmostEqual(float(margin.get("top", 0)), 4.2, delta=0.2)
            self.assertAlmostEqual(float(margin.get("right", 0)), 7.8, delta=0.2)
            self.assertAlmostEqual(float(margin.get("bottom", 0)), 5.1, delta=0.2)
            self.assertAlmostEqual(float(margin.get("left", 0)), 6.4, delta=0.2)
        finally:
            try:
                os.remove(pptx_path)
            except Exception:
                pass

    # ──  回归：文本框内边距必须显式写出，避免 PowerPoint 幽灵内边距 ──

    @staticmethod
    def _body_pr_insets(shape):
        ns = {"a": "http://schemas.openxmlformats.org/drawingml/2006/main"}
        body_pr = shape.text_frame._txBody.find("a:bodyPr", ns)
        assert body_pr is not None
        return {attr: body_pr.get(attr) for attr in ("lIns", "tIns", "rIns", "bIns")}

    def test_text_element_without_margin_writes_zero_insets(self):
        """：无 margin 的文本元素必须显式写 inset=0，避免 PPT 套默认 0.1in/0.05in。"""
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        shape = _write_text_element(
            slide=slide,
            props={
                "content": '<p><span style="font-size:16pt">Hello</span></p>',
                "defaultFontSize": 16,
            },
            left=Inches(1),
            top=Inches(1),
            width=Inches(4),
            height=Inches(1),
        )
        insets = self._body_pr_insets(shape)
        self.assertEqual(insets, {"lIns": "0", "tIns": "0", "rIns": "0", "bIns": "0"})

    def test_text_element_with_margin_writes_explicit_insets(self):
        """：带 margin 的文本元素按 margin 写 inset（pt → EMU）。"""
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        shape = _write_text_element(
            slide=slide,
            props={
                "content": '<p><span style="font-size:16pt">Hi</span></p>',
                "defaultFontSize": 16,
                "margin": {"left": 7.2, "right": 7.2, "top": 3.6, "bottom": 3.6},
            },
            left=Inches(1),
            top=Inches(1),
            width=Inches(4),
            height=Inches(1),
        )
        insets = self._body_pr_insets(shape)
        # 7.2pt → 91440 EMU（0.1in）、3.6pt → 45720 EMU（0.05in）
        self.assertEqual(insets["lIns"], "91440")
        self.assertEqual(insets["rIns"], "91440")
        self.assertEqual(insets["tIns"], "45720")
        self.assertEqual(insets["bIns"], "45720")

    def test_extract_text_margins_fills_ooxml_defaults_when_insets_absent(self):
        """：导入侧对缺省 inset 补 OOXML 默认值，保证导入 deck 往返保真。"""
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1))
        box.text_frame.text = "x"
        # 不设置任何 inset —— 模拟依赖 PowerPoint 默认内边距的来源 PPT
        margins = _PPTX_IO._extract_text_margins(box.text_frame)
        self.assertIsNotNone(margins)
        self.assertAlmostEqual(float(margins["left"]), 7.2, delta=0.2)
        self.assertAlmostEqual(float(margins["right"]), 7.2, delta=0.2)
        self.assertAlmostEqual(float(margins["top"]), 3.6, delta=0.2)
        self.assertAlmostEqual(float(margins["bottom"]), 3.6, delta=0.2)

    def test_write_read_text_roundtrip_preserves_nested_list_structure(self):
        pages = [
            {
                "id": "page-list",
                "background": {"type": "color", "value": "#ffffff"},
                "elements": [
                    {
                        "id": "text-list",
                        "type": "text",
                        "x": 120,
                        "y": 120,
                        "width": 820,
                        "height": 300,
                        "rotate": 0,
                        "zIndex": 0,
                        "props": {
                            "content": (
                                "<ul><li><p>A</p><ul><li><p>B</p></li></ul></li><li><p>C</p></li></ul>"
                                '<ol type="a"><li><p>D</p></li></ol>'
                            ),
                            "defaultFontFamily": "Calibri",
                            "defaultFontSize": 16,
                            "defaultColor": "#333333",
                        },
                    }
                ],
            }
        ]

        fd, pptx_path = tempfile.mkstemp(suffix=".pptx")
        os.close(fd)
        try:
            write(pages=pages, output_path=pptx_path, canvas_width=1920, canvas_height=1080)
            out_pages = read(pptx_path=pptx_path, canvas_width=1920, canvas_height=1080)
            self.assertEqual(len(out_pages), 1)
            text_elements = [el for el in out_pages[0].get("elements", []) if el.get("type") == "text"]
            self.assertEqual(len(text_elements), 1)
            content = text_elements[0].get("props", {}).get("content", "")

            self.assertIn("<ul>", content)
            self.assertIn("</ul>", content)
            self.assertIn('<ol type="a">', content)
            self.assertIn(">A<", content)
            self.assertIn(">B<", content)
            self.assertIn(">C<", content)
            self.assertIn(">D<", content)
        finally:
            try:
                os.remove(pptx_path)
            except Exception:
                pass

    def test_write_text_element_theme_color_key_writes_scheme_clr(self):
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        shape = _write_text_element(
            slide=slide,
            props={
                "content": (
                    '<p><span data-theme-color-key="accent1" style="color:#4472C4">ThemeA</span></p>'
                ),
                "defaultFontFamily": "Calibri",
                "defaultFontSize": 16,
                "defaultColor": "#333333",
            },
            left=Inches(1),
            top=Inches(1),
            width=Inches(5),
            height=Inches(2),
        )
        run = shape.text_frame.paragraphs[0].runs[0]
        nsmap_a = "http://schemas.openxmlformats.org/drawingml/2006/main"
        scheme = run._r.find(f".//{{{nsmap_a}}}solidFill/{{{nsmap_a}}}schemeClr")
        self.assertIsNotNone(scheme)
        assert scheme is not None
        self.assertEqual((scheme.get("val") or "").lower(), "accent1")

    def test_write_text_element_maps_inter_to_portable_script_typefaces(self):
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        shape = _write_text_element(
            slide=slide,
            props={
                "content": "<p>Anthropic 打造的下一代人工智能助手</p>",
                "defaultFontFamily": "Inter",
                "defaultFontSize": 21,
                "defaultColor": "#64748B",
            },
            left=Inches(1),
            top=Inches(1),
            width=Inches(5),
            height=Inches(1),
        )

        run = shape.text_frame.paragraphs[0].runs[0]
        nsmap_a = "http://schemas.openxmlformats.org/drawingml/2006/main"
        r_pr = run._r.find(f"{{{nsmap_a}}}rPr")
        self.assertIsNotNone(r_pr)
        assert r_pr is not None

        latin = r_pr.find(f"{{{nsmap_a}}}latin")
        east_asian = r_pr.find(f"{{{nsmap_a}}}ea")
        complex_script = r_pr.find(f"{{{nsmap_a}}}cs")

        self.assertIsNotNone(latin)
        self.assertIsNotNone(east_asian)
        self.assertIsNotNone(complex_script)
        assert latin is not None and east_asian is not None and complex_script is not None
        self.assertEqual(latin.get("typeface"), "Arial")
        self.assertEqual(east_asian.get("typeface"), "Microsoft YaHei")
        self.assertEqual(complex_script.get("typeface"), "Microsoft YaHei")

    def test_write_text_element_maps_unavailable_inter_for_latin_only_runs(self):
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        shape = _write_text_element(
            slide=slide,
            props={
                "content": "<p>Anthropic AI assistant</p>",
                "defaultFontFamily": "Inter",
                "defaultFontSize": 21,
                "defaultColor": "#64748B",
            },
            left=Inches(1),
            top=Inches(1),
            width=Inches(5),
            height=Inches(1),
        )

        run = shape.text_frame.paragraphs[0].runs[0]
        nsmap_a = "http://schemas.openxmlformats.org/drawingml/2006/main"
        r_pr = run._r.find(f"{{{nsmap_a}}}rPr")
        self.assertIsNotNone(r_pr)
        assert r_pr is not None

        latin = r_pr.find(f"{{{nsmap_a}}}latin")
        east_asian = r_pr.find(f"{{{nsmap_a}}}ea")
        complex_script = r_pr.find(f"{{{nsmap_a}}}cs")

        self.assertIsNotNone(latin)
        self.assertIsNotNone(east_asian)
        self.assertIsNotNone(complex_script)
        assert latin is not None and east_asian is not None and complex_script is not None
        self.assertEqual(latin.get("typeface"), "Arial")
        self.assertEqual(east_asian.get("typeface"), "Microsoft YaHei")
        self.assertEqual(complex_script.get("typeface"), "Microsoft YaHei")

    def test_write_read_text_roundtrip_preserves_default_theme_color_key(self):
        pages = [
            {
                "id": "page-theme",
                "background": {"type": "color", "value": "#ffffff"},
                "elements": [
                    {
                        "id": "text-theme",
                        "type": "text",
                        "x": 120,
                        "y": 120,
                        "width": 600,
                        "height": 220,
                        "rotate": 0,
                        "zIndex": 0,
                        "props": {
                            "content": "<p>Theme Default</p>",
                            "defaultFontFamily": "Calibri",
                            "defaultFontSize": 16,
                            "defaultColor": "#ED7D31",
                            "defaultColorThemeKey": "accent2",
                        },
                    }
                ],
            }
        ]

        fd, pptx_path = tempfile.mkstemp(suffix=".pptx")
        os.close(fd)
        try:
            write(pages=pages, output_path=pptx_path, canvas_width=1920, canvas_height=1080)
            out_pages = read(pptx_path=pptx_path, canvas_width=1920, canvas_height=1080)
            self.assertEqual(len(out_pages), 1)
            text_elements = [el for el in out_pages[0].get("elements", []) if el.get("type") == "text"]
            self.assertEqual(len(text_elements), 1)
            props = text_elements[0].get("props", {})
            self.assertEqual(props.get("defaultColorThemeKey"), "accent2")
            content = props.get("content", "")
            self.assertIn('data-theme-color-key="accent2"', content)
        finally:
            try:
                os.remove(pptx_path)
            except Exception:
                pass

    # ── B1-02 回归：段间距不应双重叠加 ──

    def test_paragraph_to_html_v2_no_margin_bottom(self):
        """B1-02 回归：_paragraph_to_html_v2 不应输出 margin-bottom 内联样式"""
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        shape = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(6), Inches(2))
        tf = shape.text_frame
        p = tf.paragraphs[0]
        p.text = "Hello spacing"
        p.space_after = Pt(8)

        html, p_style = _paragraph_to_html_v2(p)
        self.assertNotIn("margin-bottom", p_style)
        self.assertNotIn("margin-bottom", html)

    def test_paragraph_spacing_no_double_stack_roundtrip(self):
        """B1-02 回归：write→read 往返后 paragraphSpace 保留、HTML 无 margin-bottom"""
        pages = [
            {
                "id": "page-b102",
                "background": {"type": "color", "value": "#ffffff"},
                "elements": [
                    {
                        "id": "text-b102",
                        "type": "text",
                        "x": 100,
                        "y": 100,
                        "width": 700,
                        "height": 260,
                        "rotate": 0,
                        "zIndex": 0,
                        "props": {
                            "content": "<p>Para 1</p><p>Para 2</p><p>Para 3</p>",
                            "defaultFontFamily": "Calibri",
                            "defaultFontSize": 16,
                            "defaultColor": "#222222",
                            "paragraphSpace": 8.5,
                        },
                    }
                ],
            }
        ]

        fd, pptx_path = tempfile.mkstemp(suffix=".pptx")
        os.close(fd)
        try:
            write(pages=pages, output_path=pptx_path, canvas_width=1920, canvas_height=1080)
            out_pages = read(pptx_path=pptx_path, canvas_width=1920, canvas_height=1080)
            self.assertEqual(len(out_pages), 1)

            text_els = [e for e in out_pages[0].get("elements", []) if e.get("type") == "text"]
            self.assertEqual(len(text_els), 1)
            props = text_els[0].get("props", {})

            self.assertIn("paragraphSpace", props)
            self.assertAlmostEqual(float(props["paragraphSpace"]), 8.5, delta=0.3)

            content = props.get("content", "")
            self.assertNotIn("margin-bottom", content,
                             "HTML 不应包含 margin-bottom（已由 paragraphSpace 统一控制）")
        finally:
            try:
                os.remove(pptx_path)
            except Exception:
                pass

    def test_legacy_margin_bottom_export_still_writes_space_after(self):
        """B1-02 兼容：存量 HTML 中的 margin-bottom 在导出时仍应写入 space_after"""
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        shape = _write_text_element(
            slide=slide,
            props={
                "content": (
                    '<p style="margin-bottom:9.5pt">'
                    '<span style="font-size:14pt">Legacy</span>'
                    "</p>"
                ),
                "defaultFontFamily": "Calibri",
                "defaultFontSize": 14,
                "defaultColor": "#333333",
            },
            left=Inches(1),
            top=Inches(1),
            width=Inches(6),
            height=Inches(2),
        )
        tf = shape.text_frame
        p = tf.paragraphs[0]
        self.assertIsNotNone(p.space_after)
        self.assertAlmostEqual(p.space_after.pt, 9.5, delta=0.2)


class TestPptxTextFidelityV2(TestCase):
    """#536 后端导出文本保真回归：字号 pt 语义×dpi、autoFit 映射、文本框 fill、元素级行距。"""

    NSMAP_A = "http://schemas.openxmlformats.org/drawingml/2006/main"

    def _write_pages_and_open(self, elements, **write_kwargs):
        """写出单页 pptx 并用 python-pptx 重新打开，返回 (Presentation, 第一页)。"""
        pages = [{
            "id": "page-fidelity",
            "background": {"type": "color", "value": "#ffffff"},
            "elements": elements,
        }]
        fd, pptx_path = tempfile.mkstemp(suffix=".pptx")
        os.close(fd)
        self.addCleanup(lambda: os.path.exists(pptx_path) and os.remove(pptx_path))
        kwargs = {"canvas_width": 1920, "canvas_height": 1080}
        kwargs.update(write_kwargs)
        write(pages=pages, output_path=pptx_path, **kwargs)
        prs = Presentation(pptx_path)
        return prs, prs.slides[0]

    @staticmethod
    def _first_text_run_size_pt(slide):
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    if run.font.size is not None:
                        return run.font.size.pt
        return None

    @staticmethod
    def _first_text_body_pr(slide):
        nsmap = {"a": TestPptxTextFidelityV2.NSMAP_A}
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            return shape.text_frame._txBody.find("a:bodyPr", nsmap)
        return None

    def _text_element(self, **props_override):
        props = {
            "content": "<p>Hello</p>",
            "defaultFontFamily": "Calibri",
            "defaultFontSize": 36,
            "defaultColor": "#222222",
        }
        props.update(props_override)
        return {
            "id": "text-fidelity",
            "type": "text",
            "x": 100, "y": 120, "width": 800, "height": 240,
            "rotate": 0, "zIndex": 0,
            "props": props,
        }

    # ── P0-1：pt 语义字号按 dpi_scale 缩放，不再被当 px 多缩小 ──

    def test_default_font_size_scaled_by_dpi_not_shrunk_to_px(self):
        """1920 画布(13.333in slide)：defaultFontSize=36(pt) 应写出 ≈24pt，而非旧 bug 的 18pt。"""
        prs, slide = self._write_pages_and_open([self._text_element(defaultFontSize=36)])
        size_pt = self._first_text_run_size_pt(slide)
        self.assertIsNotNone(size_pt)
        # 36 × (96/144) = 24pt；明确不是把 pt 当 px 又 ×0.75 得到的 18pt
        self.assertAlmostEqual(size_pt, 24.0, delta=0.6)
        self.assertGreater(size_pt, 21.0)

    def test_default_font_size_pt_semantics_on_native_canvas(self):
        """slide 物理尺寸=画布(1px=1/96in, dpi_scale≈1)：defaultFontSize=36 应原样写出 36pt。"""
        prs, slide = self._write_pages_and_open(
            [self._text_element(defaultFontSize=36)],
            source_slide_width_emu=1920 * 9525,   # 20in
            source_slide_height_emu=1080 * 9525,  # 11.25in
        )
        size_pt = self._first_text_run_size_pt(slide)
        self.assertIsNotNone(size_pt)
        self.assertAlmostEqual(size_pt, 36.0, delta=0.6)

    def test_inline_pt_font_size_scaled_by_dpi(self):
        """内联 font-size:30pt 在 13.333in slide 上应缩放为 ≈20pt。"""
        prs, slide = self._write_pages_and_open([
            self._text_element(content='<p><span style="font-size:30pt">Big</span></p>')
        ])
        size_pt = self._first_text_run_size_pt(slide)
        self.assertIsNotNone(size_pt)
        self.assertAlmostEqual(size_pt, 20.0, delta=0.6)  # 30 × 0.6667

    # ── P0-2：autoFit 按语义映射，不再无条件 normAutofit ──

    def test_autofit_shrink_maps_to_norm_autofit(self):
        prs, slide = self._write_pages_and_open([self._text_element(autoFit="shrink")])
        body_pr = self._first_text_body_pr(slide)
        self.assertIsNotNone(body_pr)
        self.assertIsNotNone(body_pr.find(f"{{{self.NSMAP_A}}}normAutofit"))
        self.assertIsNone(body_pr.get("autoFit"))  # 不再写无效属性

    def test_autofit_resize_maps_to_sp_autofit(self):
        prs, slide = self._write_pages_and_open([self._text_element(autoFit="resize")])
        body_pr = self._first_text_body_pr(slide)
        self.assertIsNotNone(body_pr)
        self.assertIsNotNone(body_pr.find(f"{{{self.NSMAP_A}}}spAutoFit"))
        self.assertIsNone(body_pr.find(f"{{{self.NSMAP_A}}}normAutofit"))

    def test_autofit_unset_defaults_to_shrink_to_fit(self):
        # ：不内嵌字体后，缺字机器回退更宽字体会换行撑破框跑版；未设置 autoFit
        # 默认 shrink-to-fit（normAutofit），让打开方自动缩小字号塞回框，保文字可编辑。
        prs, slide = self._write_pages_and_open([self._text_element()])  # 无 autoFit
        body_pr = self._first_text_body_pr(slide)
        self.assertIsNotNone(body_pr)
        self.assertIsNotNone(body_pr.find(f"{{{self.NSMAP_A}}}normAutofit"))
        self.assertIsNone(body_pr.find(f"{{{self.NSMAP_A}}}noAutofit"))
        self.assertIsNone(body_pr.find(f"{{{self.NSMAP_A}}}spAutoFit"))

    # ── P0-3：文本框 fill 背景写出 ──

    def test_text_box_solid_fill_written(self):
        prs, slide = self._write_pages_and_open([self._text_element(fill="#FFAA00")])
        text_shape = next((s for s in slide.shapes if s.has_text_frame), None)
        self.assertIsNotNone(text_shape)
        # SOLID fill 应被写出且颜色匹配
        self.assertEqual(str(text_shape.fill.fore_color.rgb), "FFAA00")

    def test_text_box_transparent_fill_stays_background(self):
        prs, slide = self._write_pages_and_open([self._text_element(fill="transparent")])
        text_shape = next((s for s in slide.shapes if s.has_text_frame), None)
        self.assertIsNotNone(text_shape)
        # transparent → 不写 solid fill，保持无填充（背景）
        from pptx.enum.dml import MSO_FILL
        self.assertEqual(text_shape.fill.type, MSO_FILL.BACKGROUND)

    # ── P0-4：元素级 lineHeight 默认行距 ──

    def test_element_line_height_applied_as_default(self):
        prs, slide = self._write_pages_and_open([
            self._text_element(content="<p>Default LH</p>", lineHeight=1.8)
        ])
        text_shape = next((s for s in slide.shapes if s.has_text_frame), None)
        self.assertIsNotNone(text_shape)
        p = text_shape.text_frame.paragraphs[0]
        self.assertIsNotNone(p.line_spacing)
        self.assertAlmostEqual(float(p.line_spacing), 1.8, places=2)

    def test_inline_line_height_overrides_element_default(self):
        prs, slide = self._write_pages_and_open([
            self._text_element(
                content='<p style="line-height:1.2">Inline LH</p>',
                lineHeight=1.8,
            )
        ])
        text_shape = next((s for s in slide.shapes if s.has_text_frame), None)
        self.assertIsNotNone(text_shape)
        p = text_shape.text_frame.paragraphs[0]
        self.assertIsNotNone(p.line_spacing)
        self.assertAlmostEqual(float(p.line_spacing), 1.2, places=2)


class TestSlideBackgroundColorKey(TestCase):
    """normalize_background_for_api 产出 solid+color 时，导出必须写出页背景。"""

    def test_write_solid_background_with_color_key(self):
        pages = [{
            "id": "page-1",
            "background": {"type": "solid", "color": "#F5A830"},
            "elements": [],
        }]
        fd, pptx_path = tempfile.mkstemp(suffix=".pptx")
        os.close(fd)
        try:
            write(pages=pages, output_path=pptx_path, canvas_width=1280, canvas_height=720)
            prs = Presentation(pptx_path)
            fill = prs.slides[0].background.fill
            from pptx.enum.dml import MSO_FILL
            self.assertEqual(fill.type, MSO_FILL.SOLID)
            self.assertEqual(str(fill.fore_color.rgb), "F5A830")
        finally:
            os.unlink(pptx_path)

    def test_write_solid_background_with_value_key(self):
        pages = [{
            "id": "page-1",
            "background": {"type": "solid", "value": "#F9F8F6"},
            "elements": [],
        }]
        fd, pptx_path = tempfile.mkstemp(suffix=".pptx")
        os.close(fd)
        try:
            write(pages=pages, output_path=pptx_path, canvas_width=1280, canvas_height=720)
            prs = Presentation(pptx_path)
            fill = prs.slides[0].background.fill
            from pptx.enum.dml import MSO_FILL
            self.assertEqual(fill.type, MSO_FILL.SOLID)
            self.assertEqual(str(fill.fore_color.rgb), "F9F8F6")
        finally:
            os.unlink(pptx_path)

import importlib.util
import os
import tempfile
import zipfile
from pathlib import Path
from unittest import TestCase

from lxml import etree
from pptx import Presentation
from pptx.enum.shapes import MSO_CONNECTOR
from pptx.util import Inches

# 直接按文件路径加载，避免触发 apps.tabslide.services.__init__ 的 Django 依赖
_MODULE_PATH = Path(__file__).resolve().parents[1] / "services" / "pptx_io.py"
_SPEC = importlib.util.spec_from_file_location("tabslide_pptx_io_group_test_module", _MODULE_PATH)
if _SPEC is None or _SPEC.loader is None:
    raise RuntimeError(f"Failed to load module spec from {_MODULE_PATH}")
_PPTX_IO = importlib.util.module_from_spec(_SPEC)
_SPEC.loader.exec_module(_PPTX_IO)

read = _PPTX_IO.read
write = _PPTX_IO.write

_NS_P = "http://schemas.openxmlformats.org/presentationml/2006/main"
_PNG_DATA_URL = (
    "data:image/png;base64,"
    "iVBORw0KGgoAAAANSUhEUgAAAAEAAAABCAQAAAC1HAwCAAAAC0lEQVR42mP8/x8AAwMCAO3ZfD0AAAAASUVORK5CYII="
)


class TestPptxGroupChain(TestCase):
    def _new_tmp_pptx_path(self) -> str:
        fd, pptx_path = tempfile.mkstemp(suffix=".pptx")
        os.close(fd)
        return pptx_path

    def _count_group_nodes(self, pptx_path: str, slide_no: int = 1) -> int:
        with zipfile.ZipFile(pptx_path, "r") as zf:
            xml_bytes = zf.read(f"ppt/slides/slide{slide_no}.xml")
        doc = etree.fromstring(xml_bytes)
        groups = doc.findall(f".//{{{_NS_P}}}grpSp")
        return len(groups)

    def test_read_nested_group_recursively_flattens_members(self):
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])

        # 先创建一个内层组合：文本 + 连接器
        tb_inner = slide.shapes.add_textbox(Inches(1.0), Inches(1.0), Inches(2.0), Inches(0.8))
        tb_inner.text_frame.text = "Inner Text"
        conn = slide.shapes.add_connector(
            MSO_CONNECTOR.STRAIGHT,
            Inches(1.0),
            Inches(2.0),
            Inches(3.5),
            Inches(2.6),
        )
        inner_group = slide.shapes.add_group_shape([tb_inner, conn])

        # 再创建外层组合：内层组合 + 另一个文本框
        tb_outer = slide.shapes.add_textbox(Inches(4.0), Inches(1.2), Inches(2.0), Inches(0.8))
        tb_outer.text_frame.text = "Outer Text"
        outer_group = slide.shapes.add_group_shape([inner_group, tb_outer])
        outer_group.rotation = 18

        pptx_path = self._new_tmp_pptx_path()
        try:
            prs.save(pptx_path)
            pages = read(pptx_path=pptx_path, canvas_width=1920, canvas_height=1080)
            self.assertEqual(len(pages), 1)

            elements = pages[0]["elements"]
            # 外层组内应至少有 3 个可见子元素（2 文本 + 1 线条）
            self.assertGreaterEqual(len(elements), 3)

            grouped = [e for e in elements if e.get("groupId")]
            self.assertGreaterEqual(len(grouped), 3)
            gid_set = {e.get("groupId") for e in grouped}
            # 嵌套组应在读取后扁平为同一 groupId（父组语义）
            self.assertEqual(len(gid_set), 1)

            line_el = next((e for e in grouped if e.get("type") == "line"), None)
            self.assertIsNotNone(line_el)
            self.assertIn("start", line_el.get("props", {}))
            self.assertIn("end", line_el.get("props", {}))
            # 连接器几何已由 start/end 表达，不应再携带重复 flip 标记
            self.assertNotIn("flipH", line_el)
            self.assertNotIn("flipV", line_el)
        finally:
            try:
                os.remove(pptx_path)
            except Exception:
                pass

    def test_write_read_group_roundtrip_keeps_group_and_connector(self):
        pages = [
            {
                "id": "page-1",
                "elements": [
                    {
                        "id": "txt-g-1",
                        "type": "text",
                        "x": 100,
                        "y": 120,
                        "width": 320,
                        "height": 80,
                        "rotate": 12,
                        "zIndex": 0,
                        "groupId": "g-1",
                        "props": {
                            "content": "<p>Group Text</p>",
                            "defaultFontFamily": "Microsoft YaHei",
                            "defaultFontSize": 28,
                            "defaultColor": "#222222",
                        },
                    },
                    {
                        "id": "line-g-1",
                        "type": "line",
                        "x": 180,
                        "y": 220,
                        "width": 360,
                        "height": 140,
                        "rotate": 22,
                        "zIndex": 1,
                        "groupId": "g-1",
                        "flipH": True,
                        "props": {
                            "start": [0, 0],
                            "end": [360, 140],
                            "style": "dashed",
                            "color": "#0A84FF",
                            "lineWidth": 3,
                            "points": ["", "arrow"],
                            "broken": [140, 50],
                        },
                    },
                    {
                        "id": "txt-free",
                        "type": "text",
                        "x": 620,
                        "y": 120,
                        "width": 280,
                        "height": 80,
                        "rotate": 0,
                        "zIndex": 2,
                        "props": {
                            "content": "<p>Ungrouped</p>",
                            "defaultFontFamily": "Microsoft YaHei",
                            "defaultFontSize": 24,
                            "defaultColor": "#444444",
                        },
                    },
                ],
                "background": {"type": "color", "value": "#ffffff"},
                "notes": "",
            }
        ]

        pptx_path = self._new_tmp_pptx_path()
        try:
            write(pages=pages, output_path=pptx_path, canvas_width=1920, canvas_height=1080)
            # 连续 groupId 应写出真正 grpSp
            self.assertGreaterEqual(self._count_group_nodes(pptx_path), 1)

            out_pages = read(pptx_path=pptx_path, canvas_width=1920, canvas_height=1080)
            self.assertEqual(len(out_pages), 1)
            out_elements = out_pages[0]["elements"]

            grouped = [e for e in out_elements if e.get("groupId")]
            self.assertGreaterEqual(len(grouped), 2)
            self.assertEqual(len({e.get("groupId") for e in grouped}), 1)

            grouped_types = {e.get("type") for e in grouped}
            self.assertIn("text", grouped_types)
            self.assertIn("line", grouped_types)

            line_out = next(e for e in grouped if e.get("type") == "line")
            self.assertIn("start", line_out.get("props", {}))
            self.assertIn("end", line_out.get("props", {}))
            self.assertNotIn("flipH", line_out)
            self.assertNotIn("flipV", line_out)
        finally:
            try:
                os.remove(pptx_path)
            except Exception:
                pass

    def test_write_read_group_image_colormask_keeps_overlay_semantics(self):
        pages = [
            {
                "id": "page-1",
                "elements": [
                    {
                        "id": "img-g-mask",
                        "type": "image",
                        "x": 160,
                        "y": 180,
                        "width": 300,
                        "height": 180,
                        "rotate": 9,
                        "zIndex": 0,
                        "groupId": "g-mask",
                        "props": {
                            "src": _PNG_DATA_URL,
                            "radius": 14,
                            "colorMask": "rgba(12,34,56,0.45)",
                        },
                    },
                    {
                        "id": "txt-g-mask",
                        "type": "text",
                        "x": 520,
                        "y": 220,
                        "width": 240,
                        "height": 90,
                        "rotate": 0,
                        "zIndex": 1,
                        "groupId": "g-mask",
                        "props": {
                            "content": "<p>Mask Group</p>",
                            "defaultFontFamily": "Microsoft YaHei",
                            "defaultFontSize": 24,
                            "defaultColor": "#222222",
                        },
                    },
                ],
                "background": {"type": "color", "value": "#ffffff"},
                "notes": "",
            }
        ]

        pptx_path = self._new_tmp_pptx_path()
        try:
            write(pages=pages, output_path=pptx_path, canvas_width=1920, canvas_height=1080)
            self.assertGreaterEqual(self._count_group_nodes(pptx_path), 1)

            out_pages = read(pptx_path=pptx_path, canvas_width=1920, canvas_height=1080)
            out_elements = out_pages[0]["elements"]

            grouped = [e for e in out_elements if e.get("groupId")]
            self.assertEqual(len(grouped), 2)
            self.assertEqual(len({e.get("groupId") for e in grouped}), 1)
            self.assertEqual({e.get("type") for e in grouped}, {"image", "text"})

            image_out = next(e for e in grouped if e.get("type") == "image")
            image_props = image_out.get("props", {})
            self.assertEqual(image_props.get("colorMask"), "rgba(12,34,56,0.45)")
            self.assertAlmostEqual(float(image_props.get("radius", 0)), 14, delta=2.0)
        finally:
            try:
                os.remove(pptx_path)
            except Exception:
                pass

    def test_write_non_contiguous_group_falls_back_to_flat_order(self):
        pages = [
            {
                "id": "page-1",
                "elements": [
                    {
                        "id": "txt-g-1",
                        "type": "text",
                        "x": 80,
                        "y": 100,
                        "width": 260,
                        "height": 70,
                        "rotate": 0,
                        "zIndex": 0,
                        "groupId": "g-2",
                        "props": {
                            "content": "<p>G2-A</p>",
                            "defaultFontFamily": "Microsoft YaHei",
                            "defaultFontSize": 22,
                            "defaultColor": "#222222",
                        },
                    },
                    {
                        "id": "txt-mid",
                        "type": "text",
                        "x": 400,
                        "y": 100,
                        "width": 260,
                        "height": 70,
                        "rotate": 0,
                        "zIndex": 1,
                        "props": {
                            "content": "<p>Middle</p>",
                            "defaultFontFamily": "Microsoft YaHei",
                            "defaultFontSize": 22,
                            "defaultColor": "#222222",
                        },
                    },
                    {
                        "id": "line-g-2",
                        "type": "line",
                        "x": 120,
                        "y": 220,
                        "width": 300,
                        "height": 120,
                        "rotate": 10,
                        "zIndex": 2,
                        "groupId": "g-2",
                        "props": {
                            "start": [0, 0],
                            "end": [300, 120],
                            "style": "solid",
                            "color": "#16A34A",
                            "lineWidth": 2,
                            "points": ["", ""],
                        },
                    },
                ],
                "background": {"type": "color", "value": "#ffffff"},
                "notes": "",
            }
        ]

        pptx_path = self._new_tmp_pptx_path()
        try:
            write(pages=pages, output_path=pptx_path, canvas_width=1920, canvas_height=1080)
            # 非连续 groupId 必须降级平铺，避免错误成组破坏层级
            self.assertEqual(self._count_group_nodes(pptx_path), 0)

            out_pages = read(pptx_path=pptx_path, canvas_width=1920, canvas_height=1080)
            out_elements = out_pages[0]["elements"]
            self.assertEqual(len(out_elements), 3)
            self.assertTrue(all(not el.get("groupId") for el in out_elements))
        finally:
            try:
                os.remove(pptx_path)
            except Exception:
                pass

    def test_read_write_group_flip_line_keeps_canonical_line_geometry(self):
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])

        tb = slide.shapes.add_textbox(Inches(1.0), Inches(1.0), Inches(2.0), Inches(0.8))
        tb.text_frame.text = "Group"
        conn = slide.shapes.add_connector(
            MSO_CONNECTOR.STRAIGHT,
            Inches(1.0),
            Inches(2.0),
            Inches(3.2),
            Inches(2.7),
        )
        group = slide.shapes.add_group_shape([tb, conn])

        # 通过 XML 设置组合整体变换（flipH + rotation）
        ns_a = "http://schemas.openxmlformats.org/drawingml/2006/main"
        grp_sp_pr = group._element.find(f"{{{_NS_P}}}grpSpPr")
        self.assertIsNotNone(grp_sp_pr)
        xfrm = grp_sp_pr.find(f"{{{ns_a}}}xfrm")
        self.assertIsNotNone(xfrm)
        xfrm.set("flipH", "1")
        xfrm.set("rot", str(int(22 * 60000)))

        src_path = self._new_tmp_pptx_path()
        out_path = self._new_tmp_pptx_path()
        try:
            prs.save(src_path)
            pages = read(pptx_path=src_path, canvas_width=1920, canvas_height=1080)
            self.assertEqual(len(pages), 1)

            in_line = next(e for e in pages[0]["elements"] if e.get("type") == "line")
            self.assertIn("start", in_line.get("props", {}))
            self.assertIn("end", in_line.get("props", {}))
            # 线条在组翻转后应规范为点位几何，不再额外携带 flip 标记
            self.assertNotIn("flipH", in_line)
            self.assertNotIn("flipV", in_line)

            write(pages=pages, output_path=out_path, canvas_width=1920, canvas_height=1080)
            out_pages = read(pptx_path=out_path, canvas_width=1920, canvas_height=1080)
            out_line = next(e for e in out_pages[0]["elements"] if e.get("type") == "line")
            self.assertNotIn("flipH", out_line)
            self.assertNotIn("flipV", out_line)

            self.assertAlmostEqual(in_line.get("x", 0), out_line.get("x", 0), places=3)
            self.assertAlmostEqual(in_line.get("y", 0), out_line.get("y", 0), places=3)
            self.assertAlmostEqual(in_line.get("width", 0), out_line.get("width", 0), places=3)
            self.assertAlmostEqual(in_line.get("height", 0), out_line.get("height", 0), places=3)
            self.assertAlmostEqual(in_line.get("rotate", 0), out_line.get("rotate", 0), places=2)

            in_start = in_line.get("props", {}).get("start")
            out_start = out_line.get("props", {}).get("start")
            in_end = in_line.get("props", {}).get("end")
            out_end = out_line.get("props", {}).get("end")
            self.assertIsNotNone(in_start)
            self.assertIsNotNone(out_start)
            self.assertIsNotNone(in_end)
            self.assertIsNotNone(out_end)
            self.assertAlmostEqual(float(in_start[0]), float(out_start[0]), places=3)
            self.assertAlmostEqual(float(in_start[1]), float(out_start[1]), places=3)
            self.assertAlmostEqual(float(in_end[0]), float(out_end[0]), places=3)
            self.assertAlmostEqual(float(in_end[1]), float(out_end[1]), places=3)
        finally:
            try:
                os.remove(src_path)
            except Exception:
                pass
            try:
                os.remove(out_path)
            except Exception:
                pass

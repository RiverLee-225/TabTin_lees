#!/usr/bin/env python3
"""
slides_generator - HTML 演示文稿转 PPTX 转换工具

核心功能：
  1. 解析 HTML 中的 .ppt-slide 容器
  2. 使用 Playwright + Chromium 逐张截图（@2x 高清）
  3. 使用 python-pptx 拼装为标准 PPTX 文件

使用方式：
  # 基础用法
  python slides_generator.py input.pptx.html -o output.pptx

  # 指定标题
  python slides_generator.py input.pptx.html -o output.pptx --title "我的演示文稿"

  # 高清模式（3x 缩放）
  python slides_generator.py input.pptx.html -o output.pptx --scale 3

  # 作为 Python 模块调用
  from slides_generator import generate_pptx
  generate_pptx("input.html", "output.pptx", title="演示文稿")

技术规格：
  - 幻灯片尺寸: 1280×720px (16:9)
  - PPTX 尺寸: 12192000×6858000 EMU
  - 默认 @2x 渲染（2560×1440 实际像素）
  - 支持 ECharts / Plotly / Chart.js 图表渲染
  - 支持 Tailwind CSS / 自定义字体
"""

import argparse
import asyncio
import logging
import os
import re
import sys
import tempfile
from pathlib import Path
from typing import Optional

logger = logging.getLogger(__name__)

# ============================================================================
# 常量定义
# ============================================================================

# 幻灯片尺寸（像素）- 与 .ppt-slide 一致
SLIDE_WIDTH_PX = 1280
SLIDE_HEIGHT_PX = 720

# PPTX 尺寸（EMU 单位）
# 1 EMU = 1/914400 英寸, 1 像素 = 9525 EMU (at 96 DPI)
EMU_PER_PX = 9525
SLIDE_WIDTH_EMU = SLIDE_WIDTH_PX * EMU_PER_PX   # 12192000
SLIDE_HEIGHT_EMU = SLIDE_HEIGHT_PX * EMU_PER_PX  # 6858000

# 默认渲染缩放（@2x 高清）
DEFAULT_SCALE_FACTOR = 2

# 等待图表渲染的超时时间（毫秒）
CHART_RENDER_TIMEOUT = 5000

# 幻灯片选择器
SLIDE_SELECTOR = ".ppt-slide"


# ============================================================================
# 核心转换逻辑
# ============================================================================

async def render_slides_async(
    html_path: str,
    output_dir: str,
    scale_factor: int = DEFAULT_SCALE_FACTOR,
    wait_for_charts: bool = True,
) -> list[str]:
    """
    使用 Playwright 渲染 HTML 中的每张幻灯片为 PNG 截图。

    Args:
        html_path: HTML 文件路径
        output_dir: 截图输出目录
        scale_factor: 渲染缩放倍数（2 = @2x 高清）
        wait_for_charts: 是否等待图表库渲染完成

    Returns:
        PNG 截图文件路径列表（按幻灯片顺序）
    """
    from playwright.async_api import async_playwright

    html_path = os.path.abspath(html_path)
    file_url = f"file://{html_path}"

    screenshots = []

    async with async_playwright() as p:
        browser = await p.chromium.launch(
            headless=True,
            args=[
                "--no-sandbox",
                "--disable-gpu",
                "--disable-dev-shm-usage",
                "--font-render-hinting=none",  # 字体渲染一致性
            ],
        )

        page = await browser.new_page(
            viewport={"width": SLIDE_WIDTH_PX, "height": SLIDE_HEIGHT_PX},
            device_scale_factor=scale_factor,
        )

        # 加载 HTML
        logger.info(f"加载 HTML: {html_path}")
        await page.goto(file_url, wait_until="networkidle")

        # 等待图表库渲染完成
        if wait_for_charts:
            await _wait_for_charts(page)

        # 获取所有幻灯片元素
        slides = await page.query_selector_all(SLIDE_SELECTOR)

        if not slides:
            logger.warning(f"未找到 '{SLIDE_SELECTOR}' 元素，尝试备用选择器...")
            # 备用：尝试 section, .slide 等
            for fallback_selector in ["section.slide", ".slide", "section", "[data-slide]"]:
                slides = await page.query_selector_all(fallback_selector)
                if slides:
                    logger.info(f"使用备用选择器 '{fallback_selector}' 找到 {len(slides)} 张幻灯片")
                    break

        if not slides:
            raise ValueError("HTML 中未找到任何幻灯片元素。请确保使用 .ppt-slide 类名包裹每张幻灯片。")

        logger.info(f"找到 {len(slides)} 张幻灯片，开始截图...")

        for i, slide in enumerate(slides):
            screenshot_path = os.path.join(output_dir, f"slide_{i + 1:03d}.png")

            # 滚动到当前幻灯片位置
            await slide.scroll_into_view_if_needed()
            await page.wait_for_timeout(200)  # 短暂等待渲染稳定

            # 截图单个幻灯片元素
            await slide.screenshot(
                path=screenshot_path,
                type="png",
            )

            screenshots.append(screenshot_path)

            slide_type = await slide.get_attribute("type") or "unknown"
            logger.info(f"  [{i + 1}/{len(slides)}] 截图完成: {slide_type} -> {screenshot_path}")

        await browser.close()

    return screenshots


async def _wait_for_charts(page) -> None:
    """
    等待页面中的图表库渲染完成。
    支持 ECharts、Plotly、Chart.js。
    """
    # 等待 ECharts 渲染完成
    try:
        await page.evaluate("""
            () => {
                return new Promise((resolve) => {
                    // 检查是否有 ECharts 实例
                    if (typeof echarts !== 'undefined') {
                        const instances = echarts.getInstanceByDom ?
                            document.querySelectorAll('[_echarts_instance_]') : [];
                        if (instances.length > 0) {
                            let rendered = 0;
                            instances.forEach(el => {
                                const chart = echarts.getInstanceByDom(el);
                                if (chart) {
                                    chart.on('finished', () => {
                                        rendered++;
                                        if (rendered >= instances.length) resolve();
                                    });
                                }
                            });
                            // 超时保护
                            setTimeout(resolve, 3000);
                            return;
                        }
                    }
                    // 没有 ECharts，直接返回
                    resolve();
                });
            }
        """)
    except Exception:
        pass  # ECharts 不存在，跳过

    # 等待 Plotly 渲染完成
    try:
        await page.evaluate("""
            () => {
                return new Promise((resolve) => {
                    if (typeof Plotly !== 'undefined') {
                        // Plotly 图表在 DOM 中有 .plotly 类
                        const plots = document.querySelectorAll('.js-plotly-plot');
                        if (plots.length > 0) {
                            let rendered = 0;
                            plots.forEach(el => {
                                el.on('plotly_afterplot', () => {
                                    rendered++;
                                    if (rendered >= plots.length) resolve();
                                });
                            });
                            setTimeout(resolve, 3000);
                            return;
                        }
                    }
                    resolve();
                });
            }
        """)
    except Exception:
        pass

    # 通用等待：确保所有图片加载完成
    try:
        await page.evaluate("""
            () => {
                return Promise.all(
                    Array.from(document.images)
                        .filter(img => !img.complete)
                        .map(img => new Promise(resolve => {
                            img.onload = img.onerror = resolve;
                        }))
                );
            }
        """)
    except Exception:
        pass

    # 最终等待一小段时间确保渲染稳定
    await page.wait_for_timeout(500)


def assemble_pptx(
    screenshot_paths: list[str],
    output_path: str,
    title: str = "Presentation",
) -> str:
    """
    将截图拼装为 PPTX 文件。

    Args:
        screenshot_paths: 截图文件路径列表（按幻灯片顺序）
        output_path: 输出 PPTX 文件路径
        title: 演示文稿标题

    Returns:
        输出文件路径
    """
    from pptx import Presentation
    from pptx.util import Emu
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    prs = Presentation()

    # 设置幻灯片尺寸为 16:9 (1280x720)
    prs.slide_width = Emu(SLIDE_WIDTH_EMU)
    prs.slide_height = Emu(SLIDE_HEIGHT_EMU)

    # 设置演示文稿标题
    prs.core_properties.title = title

    # 使用空白布局
    blank_layout = prs.slide_layouts[6]  # 通常索引 6 是空白布局

    for i, screenshot_path in enumerate(screenshot_paths):
        logger.info(f"  添加幻灯片 {i + 1}/{len(screenshot_paths)}")

        slide = prs.slides.add_slide(blank_layout)

        # 将截图作为全幻灯片背景图片插入
        # 位置: 左上角 (0, 0)，尺寸: 全幻灯片
        slide.shapes.add_picture(
            screenshot_path,
            left=Emu(0),
            top=Emu(0),
            width=Emu(SLIDE_WIDTH_EMU),
            height=Emu(SLIDE_HEIGHT_EMU),
        )

    # 保存
    os.makedirs(os.path.dirname(os.path.abspath(output_path)), exist_ok=True)
    prs.save(output_path)
    logger.info(f"PPTX 已保存: {output_path}")

    return output_path


# ============================================================================
# 高级模式：混合转换（文本可编辑）
# ============================================================================

async def render_slides_hybrid_async(
    html_path: str,
    output_path: str,
    title: str = "Presentation",
    scale_factor: int = DEFAULT_SCALE_FACTOR,
) -> str:
    """
    混合模式转换：提取文本为可编辑文本框 + 背景截图。

    混合模式实现：
    - 背景/形状 → 截图作为幻灯片背景
    - 文本内容 → 提取为 python-pptx 文本框（可编辑）

    注意：这是高级模式，需要更复杂的 DOM 解析。
    基础模式（全截图）已经能覆盖绝大多数需求。

    Args:
        html_path: HTML 文件路径
        output_path: 输出 PPTX 路径
        title: 演示文稿标题
        scale_factor: 渲染缩放

    Returns:
        输出文件路径
    """
    from playwright.async_api import async_playwright
    from pptx import Presentation
    from pptx.util import Emu, Pt
    from pptx.dml.color import RGBColor

    html_path = os.path.abspath(html_path)
    file_url = f"file://{html_path}"

    prs = Presentation()
    prs.slide_width = Emu(SLIDE_WIDTH_EMU)
    prs.slide_height = Emu(SLIDE_HEIGHT_EMU)
    prs.core_properties.title = title

    blank_layout = prs.slide_layouts[6]

    async with async_playwright() as p:
        browser = await p.chromium.launch(
            headless=True,
            args=["--no-sandbox", "--disable-gpu"],
        )

        page = await browser.new_page(
            viewport={"width": SLIDE_WIDTH_PX, "height": SLIDE_HEIGHT_PX},
            device_scale_factor=scale_factor,
        )

        await page.goto(file_url, wait_until="networkidle")
        await _wait_for_charts(page)

        slides_elements = await page.query_selector_all(SLIDE_SELECTOR)

        if not slides_elements:
            raise ValueError("未找到幻灯片元素")

        with tempfile.TemporaryDirectory(prefix="slides_gen_") as tmp_dir:
            for i, slide_el in enumerate(slides_elements):
                logger.info(f"  处理幻灯片 {i + 1}/{len(slides_elements)} (混合模式)")

                slide = prs.slides.add_slide(blank_layout)

                # --- Step 1: 深度提取所有文本元素的信息 ---
                text_data = await slide_el.evaluate("""
                    (slideEl) => {
                        const slideRect = slideEl.getBoundingClientRect();
                        const results = [];
                        const processedTexts = new Set();  // 去重

                        // 递归遍历所有 DOM 节点，找到包含文本的叶子节点
                        function walkDOM(node) {
                            // 跳过 script, style, svg 内部文本
                            if (['SCRIPT', 'STYLE', 'SVG', 'CANVAS', 'IMG', 'BR', 'HR'].includes(node.tagName)) return;

                            // 检查该节点是否是"文本叶子"——即它自己有直接文本内容
                            // 或者它的所有子节点都是文本/内联元素
                            const childElements = Array.from(node.children || []).filter(
                                c => !['SCRIPT', 'STYLE', 'BR', 'HR'].includes(c.tagName)
                            );

                            const directText = Array.from(node.childNodes)
                                .filter(n => n.nodeType === 3)  // TEXT_NODE
                                .map(n => n.textContent.trim())
                                .join('');

                            const isLeafText = childElements.length === 0 && directText.length > 0;
                            const isInlineContainer = childElements.length > 0 &&
                                childElements.every(c => {
                                    const display = window.getComputedStyle(c).display;
                                    return display === 'inline' || display === 'inline-block';
                                });

                            const innerText = node.innerText ? node.innerText.trim() : '';

                            if ((isLeafText || isInlineContainer) && innerText.length > 0) {
                                const rect = node.getBoundingClientRect();
                                const style = window.getComputedStyle(node);

                                // 跳过不可见元素
                                if (style.visibility === 'hidden' || style.display === 'none' ||
                                    parseFloat(style.opacity) === 0) return;
                                // 跳过零尺寸
                                if (rect.width < 2 || rect.height < 2) return;
                                // 跳过超出幻灯片范围的元素
                                const relX = rect.left - slideRect.left;
                                const relY = rect.top - slideRect.top;
                                if (relX > 1280 || relY > 720 || relX + rect.width < 0 || relY + rect.height < 0) return;

                                // 去重：同一位置同一文本不重复提取
                                const key = `${Math.round(relX)}_${Math.round(relY)}_${innerText.substring(0, 20)}`;
                                if (processedTexts.has(key)) return;
                                processedTexts.add(key);

                                results.push({
                                    text: innerText,
                                    x: Math.max(0, relX),
                                    y: Math.max(0, relY),
                                    width: Math.min(rect.width, 1280 - Math.max(0, relX)),
                                    height: Math.min(rect.height, 720 - Math.max(0, relY)),
                                    fontSize: parseFloat(style.fontSize),
                                    fontWeight: style.fontWeight,
                                    fontStyle: style.fontStyle,
                                    color: style.color,
                                    textAlign: style.textAlign,
                                    lineHeight: style.lineHeight,
                                    letterSpacing: style.letterSpacing,
                                    tag: node.tagName.toLowerCase(),
                                });
                                return;  // 不再递归子节点
                            }

                            // 递归子元素
                            for (const child of childElements) {
                                walkDOM(child);
                            }
                        }

                        walkDOM(slideEl);
                        return results;
                    }
                """)

                # --- Step 2: 隐藏所有文本，截图纯背景 ---
                await slide_el.evaluate("""
                    (slideEl) => {
                        // 隐藏所有文本节点（通过设置颜色为透明）
                        const walker = document.createTreeWalker(
                            slideEl, NodeFilter.SHOW_ELEMENT, null, false
                        );
                        const textEls = [];
                        while (walker.nextNode()) {
                            const el = walker.currentNode;
                            if (['SCRIPT', 'STYLE', 'SVG', 'CANVAS', 'IMG'].includes(el.tagName)) continue;
                            const hasDirectText = Array.from(el.childNodes)
                                .some(n => n.nodeType === 3 && n.textContent.trim().length > 0);
                            if (hasDirectText || (el.children.length === 0 && el.innerText && el.innerText.trim())) {
                                el.dataset.origColor = el.style.color;
                                el.style.color = 'transparent';
                                // 也隐藏 Font Awesome 图标
                                if (el.classList.contains('fa') || el.classList.contains('fas') ||
                                    el.classList.contains('far') || el.classList.contains('fab')) {
                                    el.dataset.origVis = el.style.visibility;
                                    el.style.visibility = 'hidden';
                                }
                            }
                        }
                    }
                """)

                await slide_el.scroll_into_view_if_needed()
                await page.wait_for_timeout(150)

                bg_path = os.path.join(tmp_dir, f"bg_{i + 1:03d}.png")
                await slide_el.screenshot(path=bg_path, type="png")

                # 恢复文本可见性
                await slide_el.evaluate("""
                    (slideEl) => {
                        const walker = document.createTreeWalker(
                            slideEl, NodeFilter.SHOW_ELEMENT, null, false
                        );
                        while (walker.nextNode()) {
                            const el = walker.currentNode;
                            if (el.dataset.origColor !== undefined) {
                                el.style.color = el.dataset.origColor;
                                delete el.dataset.origColor;
                            }
                            if (el.dataset.origVis !== undefined) {
                                el.style.visibility = el.dataset.origVis;
                                delete el.dataset.origVis;
                            }
                        }
                    }
                """)

                # --- Step 3: 插入背景图片 ---
                slide.shapes.add_picture(
                    bg_path,
                    left=Emu(0),
                    top=Emu(0),
                    width=Emu(SLIDE_WIDTH_EMU),
                    height=Emu(SLIDE_HEIGHT_EMU),
                )

                # 添加可编辑文本框
                for td in text_data:
                    _add_text_box(slide, td)

        await browser.close()

    prs.save(output_path)
    logger.info(f"PPTX（混合模式）已保存: {output_path}")
    return output_path


def _add_text_box(slide, text_data: dict) -> None:
    """在幻灯片上添加可编辑文本框，精确还原 CSS 样式。"""
    from pptx.util import Emu, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

    x = int(text_data["x"] * EMU_PER_PX)
    y = int(text_data["y"] * EMU_PER_PX)
    w = int(text_data["width"] * EMU_PER_PX)
    h = int(text_data["height"] * EMU_PER_PX)

    # 安全边界检查
    if w <= 0 or h <= 0:
        return

    # 给文本框留出一点内边距余量
    margin_emu = int(2 * EMU_PER_PX)  # 2px 边距

    txBox = slide.shapes.add_textbox(Emu(x), Emu(y), Emu(w + margin_emu), Emu(h + margin_emu))
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.auto_size = None  # 不自动调整大小

    # 设置文本框内边距为 0（让文字紧贴定位）
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)

    # 处理多行文本
    lines = text_data["text"].split("\n")

    for line_idx, line in enumerate(lines):
        if line_idx == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()

        run = p.add_run()
        run.text = line

        # 字号：CSS px → pt（1px ≈ 0.75pt）
        font_size = text_data.get("fontSize", 16)
        run.font.size = Pt(font_size * 0.75)

        # 粗体
        font_weight = str(text_data.get("fontWeight", "400"))
        run.font.bold = font_weight in ("bold", "700", "800", "900")

        # 斜体
        if text_data.get("fontStyle") == "italic":
            run.font.italic = True

        # 颜色
        color_str = text_data.get("color", "rgb(0, 0, 0)")
        rgb = _parse_css_color(color_str)
        if rgb:
            run.font.color.rgb = RGBColor(*rgb)

        # 字间距
        letter_spacing = text_data.get("letterSpacing", "normal")
        if letter_spacing and letter_spacing != "normal":
            try:
                sp_val = float(letter_spacing.replace("px", "").replace("em", ""))
                if "em" in str(letter_spacing):
                    sp_val = sp_val * font_size
                run.font._element.attrib[
                    "{http://schemas.openxmlformats.org/drawingml/2006/main}spc"
                ] = str(int(sp_val * 100))
            except (ValueError, TypeError):
                pass

        # 对齐
        align_map = {
            "center": PP_ALIGN.CENTER,
            "right": PP_ALIGN.RIGHT,
            "justify": PP_ALIGN.JUSTIFY,
            "left": PP_ALIGN.LEFT,
            "start": PP_ALIGN.LEFT,
            "end": PP_ALIGN.RIGHT,
        }
        p.alignment = align_map.get(text_data.get("textAlign", "left"), PP_ALIGN.LEFT)

        # 行高
        line_height = text_data.get("lineHeight", "normal")
        if line_height and line_height != "normal":
            try:
                lh_val = float(line_height.replace("px", ""))
                if lh_val > 5:  # px 值
                    p.line_spacing = Pt(lh_val * 0.75)
                else:  # 倍数值
                    p.line_spacing = lh_val
            except (ValueError, TypeError):
                pass


def _parse_css_color(color_str: str) -> Optional[tuple]:
    """解析 CSS 颜色值为 (R, G, B) 元组。"""
    # rgb(r, g, b)
    m = re.match(r"rgb\(\s*(\d+)\s*,\s*(\d+)\s*,\s*(\d+)\s*\)", color_str)
    if m:
        return (int(m.group(1)), int(m.group(2)), int(m.group(3)))

    # rgba(r, g, b, a)
    m = re.match(r"rgba\(\s*(\d+)\s*,\s*(\d+)\s*,\s*(\d+)\s*,\s*[\d.]+\s*\)", color_str)
    if m:
        return (int(m.group(1)), int(m.group(2)), int(m.group(3)))

    # #RRGGBB
    m = re.match(r"#([0-9a-fA-F]{6})", color_str)
    if m:
        hex_str = m.group(1)
        return (int(hex_str[0:2], 16), int(hex_str[2:4], 16), int(hex_str[4:6], 16))

    # #RGB
    m = re.match(r"#([0-9a-fA-F]{3})$", color_str)
    if m:
        hex_str = m.group(1)
        return (int(hex_str[0] * 2, 16), int(hex_str[1] * 2, 16), int(hex_str[2] * 2, 16))

    return None


# ============================================================================
# 统一入口
# ============================================================================

def generate_pptx(
    html_path: str,
    output_path: str,
    title: str = "Presentation",
    scale_factor: int = DEFAULT_SCALE_FACTOR,
    mode: str = "screenshot",
) -> str:
    """
    将 HTML 演示文稿转换为 PPTX。

    Args:
        html_path: HTML 文件路径（包含 .ppt-slide 元素）
        output_path: 输出 PPTX 文件路径
        title: 演示文稿标题
        scale_factor: 渲染缩放倍数（2=@2x, 3=@3x）
        mode: 转换模式
            - "screenshot": 全截图模式（默认，最稳定）
            - "hybrid": 混合模式（文本可编辑）

    Returns:
        输出文件路径
    """
    if not os.path.exists(html_path):
        raise FileNotFoundError(f"HTML 文件不存在: {html_path}")

    if mode == "hybrid":
        return asyncio.run(render_slides_hybrid_async(
            html_path, output_path, title, scale_factor
        ))

    # 默认：全截图模式
    with tempfile.TemporaryDirectory(prefix="slides_gen_") as tmp_dir:
        # Step 1: Playwright 截图
        screenshots = asyncio.run(render_slides_async(
            html_path, tmp_dir, scale_factor
        ))

        if not screenshots:
            raise RuntimeError("截图失败：未生成任何幻灯片图片")

        # Step 2: 拼装 PPTX
        return assemble_pptx(screenshots, output_path, title)


# ============================================================================
# 辅助：从 HTML 字符串生成（content 参数）
# ============================================================================

def generate_pptx_from_content(
    content: str,
    file_path: str,
    slide_title: str = "Presentation",
    scale_factor: int = DEFAULT_SCALE_FACTOR,
    mode: str = "screenshot",
) -> str:
    """
    从 HTML 字符串生成 PPTX。

    Args:
        content: HTML 格式的演示文稿内容（完整 HTML 字符串）
        file_path: 输出路径（.pptx.html 或 .pptx）
        slide_title: 演示文稿标题
        scale_factor: 渲染缩放
        mode: 转换模式

    Returns:
        PPTX 输出文件路径
    """
    # 确定输出路径
    if file_path.endswith(".pptx.html"):
        html_path = file_path
        pptx_path = file_path.replace(".pptx.html", ".pptx")
    elif file_path.endswith(".pptx"):
        pptx_path = file_path
        html_path = file_path.replace(".pptx", ".pptx.html")
    else:
        html_path = file_path + ".html"
        pptx_path = file_path + ".pptx"

    # 保存 HTML 源文件
    os.makedirs(os.path.dirname(os.path.abspath(html_path)), exist_ok=True)
    with open(html_path, "w", encoding="utf-8") as f:
        f.write(content)
    logger.info(f"HTML 源文件已保存: {html_path}")

    # 转换
    return generate_pptx(html_path, pptx_path, slide_title, scale_factor, mode)


# ============================================================================
# CLI 入口
# ============================================================================

def main():
    parser = argparse.ArgumentParser(
        description="slides_generator - HTML 演示文稿转 PPTX 转换工具",
        formatter_class=argparse.RawDescriptionHelpFormatter,
        epilog="""
示例:
  # 基础转换
  %(prog)s presentation.html -o presentation.pptx

  # 指定标题 + 高清模式
  %(prog)s presentation.html -o output.pptx --title "季度报告" --scale 3

  # 混合模式（文本可编辑）
  %(prog)s presentation.html -o output.pptx --mode hybrid

  # 从标准输入读取 HTML
  cat slides.html | %(prog)s - -o output.pptx
        """,
    )

    parser.add_argument(
        "input",
        help="HTML 文件路径（使用 '-' 从标准输入读取）",
    )
    parser.add_argument(
        "-o", "--output",
        required=True,
        help="输出 PPTX 文件路径",
    )
    parser.add_argument(
        "--title",
        default="Presentation",
        help="演示文稿标题（默认: Presentation）",
    )
    parser.add_argument(
        "--scale",
        type=int,
        default=DEFAULT_SCALE_FACTOR,
        choices=[1, 2, 3],
        help=f"渲染缩放倍数（默认: {DEFAULT_SCALE_FACTOR}）",
    )
    parser.add_argument(
        "--mode",
        choices=["screenshot", "hybrid"],
        default="screenshot",
        help="转换模式: screenshot=全截图（默认）, hybrid=混合（文本可编辑）",
    )
    parser.add_argument(
        "-v", "--verbose",
        action="store_true",
        help="显示详细日志",
    )

    args = parser.parse_args()

    # 配置日志
    logging.basicConfig(
        level=logging.DEBUG if args.verbose else logging.INFO,
        format="%(message)s",
    )

    # 处理标准输入
    if args.input == "-":
        content = sys.stdin.read()
        with tempfile.NamedTemporaryFile(
            mode="w", suffix=".html", delete=False, encoding="utf-8"
        ) as tmp:
            tmp.write(content)
            html_path = tmp.name
    else:
        html_path = args.input

    try:
        output = generate_pptx(
            html_path=html_path,
            output_path=args.output,
            title=args.title,
            scale_factor=args.scale,
            mode=args.mode,
        )
        print(f"✓ 转换完成: {output}")
    except Exception as e:
        logger.error(f"✗ 转换失败: {e}")
        sys.exit(1)
    finally:
        # 清理临时文件
        if args.input == "-" and os.path.exists(html_path):
            os.unlink(html_path)


if __name__ == "__main__":
    main()

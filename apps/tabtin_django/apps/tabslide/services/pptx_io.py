"""
pptx_io — PPTX ↔ SlideElement[] 双向转换（核心模块）

纯函数实现，不依赖 Django ORM。
可被 Django API 调用，未来也可被 CLI / 子进程调用。

职责：
  - read():  PPTX → List[SlidePage]（SlideElement JSON）
  - write(): List[SlidePage] → PPTX
  - EMU ↔ px 换算
  - 样式映射（字体/颜色/粗细/行高等）
"""

import functools
import logging
import math
import re
import uuid
from contextvars import ContextVar
from pathlib import Path
from typing import Any, Dict, List, Optional, Set, Tuple

logger = logging.getLogger(__name__)

# ============================================================================
# 常量
# ============================================================================

# PPT 标准 EMU（English Metric Units）
# 1 inch = 914400 EMU
# 注意：
# - 元素几何坐标（x/y/width/height）统一通过 slide_emu:canvas_px 的比例换算，
#   见 emu_to_px()/px_to_emu()，保证导入导出链路可逆。
# - EMU_PER_PT 仅用于线宽、段落间距等绝对样式量（1pt = 12700 EMU）。
EMU_PER_PT = 12700
EMU_PER_PX = 9525  # 914400 / 96；用于行高等需要 EMU↔px 转换的场景
EMU_PER_INCH = 914400
COORD_DECIMALS = 3
ROTATE_DECIMALS = 2
PX_PER_PT = 96.0 / 72.0
PT_PER_PX = 72.0 / 96.0

# PowerPoint 对未显式声明 inset 的文本框套用的默认内边距（EMU）：左右 0.1in、上下 0.05in。
# 导入侧用它补齐缺省的 inset，保证导入 deck 往返导出时保留有效内边距；
# 而浏览器抽取的文本元素不带 margin，导出侧按 0 写，使 PPT 文本区贴合测量框。
_DEFAULT_TEXT_INSET_EMU = {
    "left": 91440,
    "right": 91440,
    "top": 45720,
    "bottom": 45720,
}

# base64 编码膨胀约 4/3，反推单个字体原始字节上限（与 slide_service 的 10MB base64 对齐）
MAX_SINGLE_FONT_RAW_SIZE = 10 * 1024 * 1024 * 3 // 4  # ~7.5 MB

# ── 字体 fallback 映射 ──
# PPT 嵌入字体在 Web 端不可用，映射到最接近的 web-safe 字体链
_FONT_FALLBACK_MAP: Dict[str, str] = {
    # 英文无衬线
    "Calibri": "Calibri, 'Segoe UI', 'Helvetica Neue', Arial, sans-serif",
    "Calibri Light": "'Calibri Light', 'Segoe UI Light', 'Helvetica Neue Light', Arial, sans-serif",
    "Arial": "Arial, 'Helvetica Neue', Helvetica, sans-serif",
    "Arial Black": "'Arial Black', 'Helvetica Neue', Impact, sans-serif",
    "Helvetica": "Helvetica, 'Helvetica Neue', Arial, sans-serif",
    "Segoe UI": "'Segoe UI', 'Helvetica Neue', Arial, sans-serif",
    "Verdana": "Verdana, Geneva, sans-serif",
    "Tahoma": "Tahoma, Geneva, sans-serif",
    "Trebuchet MS": "'Trebuchet MS', 'Segoe UI', sans-serif",
    "Century Gothic": "'Century Gothic', 'Helvetica Neue', sans-serif",
    "Franklin Gothic Medium": "'Franklin Gothic Medium', 'Arial Narrow', Arial, sans-serif",
    "Gill Sans MT": "'Gill Sans MT', 'Gill Sans', 'Helvetica Neue', sans-serif",
    "Bahnschrift": "Bahnschrift, 'Segoe UI', sans-serif",
    "Impact": "Impact, 'Arial Black', sans-serif",
    # 英文衬线
    "Cambria": "Cambria, Georgia, 'Times New Roman', serif",
    "Georgia": "Georgia, Cambria, 'Times New Roman', serif",
    "Garamond": "Garamond, Georgia, 'Times New Roman', serif",
    "Book Antiqua": "'Book Antiqua', Palatino, Georgia, serif",
    "Palatino Linotype": "'Palatino Linotype', Palatino, Georgia, serif",
    "Times New Roman": "'Times New Roman', Times, serif",
    # 英文等宽
    "Consolas": "Consolas, Monaco, 'Courier New', monospace",
    "Courier New": "'Courier New', Courier, monospace",
    # 中文
    "微软雅黑": "'Microsoft YaHei', '微软雅黑', 'PingFang SC', 'Hiragino Sans GB', sans-serif",
    "Microsoft YaHei": "'Microsoft YaHei', '微软雅黑', 'PingFang SC', 'Hiragino Sans GB', sans-serif",
    "Microsoft YaHei UI": "'Microsoft YaHei UI', '微软雅黑', 'PingFang SC', sans-serif",
    "等线": "'DengXian', '等线', 'Microsoft YaHei', 'PingFang SC', sans-serif",
    "DengXian": "'DengXian', '等线', 'Microsoft YaHei', 'PingFang SC', sans-serif",
    "DengXian Light": "'DengXian Light', '等线', 'Microsoft YaHei', 'PingFang SC', sans-serif",
    "宋体": "'SimSun', '宋体', 'STSong', serif",
    "SimSun": "'SimSun', '宋体', 'STSong', serif",
    "新宋体": "'NSimSun', '新宋体', 'SimSun', serif",
    "NSimSun": "'NSimSun', '新宋体', 'SimSun', serif",
    "黑体": "'SimHei', '黑体', 'STHeiti', 'Microsoft YaHei', sans-serif",
    "SimHei": "'SimHei', '黑体', 'STHeiti', 'Microsoft YaHei', sans-serif",
    "楷体": "'KaiTi', '楷体', 'STKaiti', serif",
    "KaiTi": "'KaiTi', '楷体', 'STKaiti', serif",
    "仿宋": "'FangSong', '仿宋', 'STFangsong', serif",
    "FangSong": "'FangSong', '仿宋', 'STFangsong', serif",
    "幼圆": "'YouYuan', '幼圆', 'Microsoft YaHei', sans-serif",
    "YouYuan": "'YouYuan', '幼圆', 'Microsoft YaHei', sans-serif",
    "华文细黑": "'STXihei', '华文细黑', 'PingFang SC', sans-serif",
    "STXihei": "'STXihei', '华文细黑', 'PingFang SC', sans-serif",
    # macOS 中文
    "PingFang SC": "'PingFang SC', 'Microsoft YaHei', 'Hiragino Sans GB', sans-serif",
    "苹方": "'PingFang SC', 'Microsoft YaHei', 'Hiragino Sans GB', sans-serif",
    # 日文
    "游ゴシック": "'Yu Gothic', 'Hiragino Kaku Gothic Pro', 'MS PGothic', sans-serif",
    "Yu Gothic": "'Yu Gothic', 'Hiragino Kaku Gothic Pro', 'MS PGothic', sans-serif",
    "MS PGothic": "'MS PGothic', 'Yu Gothic', sans-serif",
    "Meiryo": "Meiryo, 'Yu Gothic', sans-serif",
}

_GENERIC_FONT_FAMILY_KEYWORDS = {
    "inherit",
    "initial",
    "unset",
    "revert",
    "revert-layer",
    "serif",
    "sans-serif",
    "monospace",
    "cursive",
    "fantasy",
    "system-ui",
    "ui-serif",
    "ui-sans-serif",
    "ui-monospace",
    "ui-rounded",
    "emoji",
    "math",
    "fangsong",
}

_DEFAULT_CJK_WEB_FALLBACK = "'Microsoft YaHei', 'PingFang SC', 'Hiragino Sans GB', 'Noto Sans SC', 'Source Han Sans SC', sans-serif"
_DEFAULT_PPTX_CJK_TYPEFACE = "Microsoft YaHei"
_LATIN_ONLY_TYPEFACE_CJK_FALLBACKS = {
    "arial",
    "calibri",
    "helvetica",
    "inter",
    "roboto",
    "segoe ui",
}

# 浏览器可通过 @font-face 使用这些 Web 字体，但导出的 PPTX 在目标电脑上通常找不到它们。
# PPTX 的 latin/ea/cs 每个脚本只能声明一个字体，不能照搬 CSS fallback 列表；因此只在
# “明确是 CSS 字体栈或常见 Web 字体”时收敛到跨 Office 更稳定的脚本字体。单一、明确的
# 编辑器字体（如 Aptos/自定义企业字体）仍原样保留，避免破坏 PPT 导入再导出的语义。
_WEB_ONLY_OR_NON_SYSTEM_LATIN_FONTS = {
    "inter",
    "noto sans",
    "noto sans sc",
    "open sans",
    "roboto",
    "source sans pro",
    "source han sans sc",
}
_PORTABLE_LATIN_TYPEFACES = {
    "arial": "Arial",
    "calibri": "Calibri",
    "aptos": "Aptos",
    "courier new": "Courier New",
    "georgia": "Georgia",
    "tahoma": "Tahoma",
    "times new roman": "Times New Roman",
    "trebuchet ms": "Trebuchet MS",
    "verdana": "Verdana",
}
_PREFERRED_CJK_TYPEFACES = (
    "Microsoft YaHei",
    "PingFang SC",
    "Hiragino Sans GB",
    "Noto Sans SC",
    "Source Han Sans SC",
)

# write() 期间按调用隔离已声明嵌入的字体；ContextVar 避免并发导出互相污染。
_ACTIVE_EMBEDDED_TYPEFACES: ContextVar[frozenset[str]] = ContextVar(
    "pptx_io_active_embedded_typefaces",
    default=frozenset(),
)


def _embedded_typefaces_from_font_meta(font_meta: Any) -> frozenset[str]:
    if not isinstance(font_meta, dict):
        return frozenset()
    entries = font_meta.get("embedded_fonts")
    if not isinstance(entries, list):
        return frozenset()
    return frozenset(
        str(entry.get("name") or "").strip().casefold()
        for entry in entries
        if isinstance(entry, dict)
        and str(entry.get("name") or "").strip()
        and (entry.get("data_base64") or entry.get("oss_url"))
    )


def _split_font_family_tokens(font_family: Any) -> List[str]:
    """把 CSS font-family 拆成有序字体名；忽略泛型与 var()/calc() 等表达式。"""
    if not isinstance(font_family, str) or not font_family.strip():
        return []
    tokens: List[str] = []
    quote: Optional[str] = None
    depth = 0
    start = 0
    for i, ch in enumerate(font_family):
        if quote:
            if ch == quote and (i == 0 or font_family[i - 1] != "\\"):
                quote = None
            continue
        if ch in ('"', "'"):
            quote = ch
        elif ch == "(":
            depth += 1
        elif ch == ")" and depth > 0:
            depth -= 1
        elif ch == "," and depth == 0:
            token = font_family[start:i].strip().strip("'\"")
            if token:
                tokens.append(token)
            start = i + 1
    token = font_family[start:].strip().strip("'\"")
    if token:
        tokens.append(token)
    return [
        token
        for token in tokens
        if token.casefold() not in _GENERIC_FONT_FAMILY_KEYWORDS
        and "(" not in token
        and ")" not in token
    ]


def _resolve_pptx_typefaces(font_family: Any, text: Any = "") -> Tuple[str, str, str]:
    """将 CSS 字体栈解析为 PPTX 的 latin/ea/cs 三脚本字体。

    CSS 可以逐字形回退，DrawingML 却只能为每个脚本指定一个 typeface。若首选是 Web 字体
    或输入本来就是 fallback 栈，选择稳定的脚本字体，避免目标机缺字体后发生不可控替换；
    单一明确字体继续原样往返。
    """
    raw = str(font_family or "").strip()
    tokens = _split_font_family_tokens(raw)
    primary = tokens[0] if tokens else (_strip_font_fallback(raw) or "Arial")
    primary_is_web_font = primary.casefold() in _WEB_ONLY_OR_NON_SYSTEM_LATIN_FONTS
    primary_is_embedded = primary.casefold() in _ACTIVE_EMBEDDED_TYPEFACES.get()
    # 计算样式通常返回完整 fallback 栈；逗号本身不代表首选字体不可用。
    # 未知首选可能是本机安装或 font_meta 嵌入的品牌字体，必须继续由 run 引用它，
    # 否则嵌入成功也不会生效。仅对明确已知的 Web-only 字体收敛到 Office 稳定字体。
    if not primary_is_web_font or primary_is_embedded:
        if (
            _contains_cjk_text(text)
            and primary.casefold() in _LATIN_ONLY_TYPEFACE_CJK_FALLBACKS
        ):
            return primary, _DEFAULT_PPTX_CJK_TYPEFACE, _DEFAULT_PPTX_CJK_TYPEFACE
        return primary, primary, primary

    latin = next(
        (
            _PORTABLE_LATIN_TYPEFACES[token.casefold()]
            for token in tokens
            if token.casefold() in _PORTABLE_LATIN_TYPEFACES
        ),
        "Arial",
    )
    token_names = {token.casefold(): token for token in tokens}
    east_asian = next(
        (
            canonical
            for canonical in _PREFERRED_CJK_TYPEFACES
            if canonical.casefold() in token_names
        ),
        _DEFAULT_PPTX_CJK_TYPEFACE,
    )
    return latin, east_asian, east_asian


def _split_first_font_family_token(font_family: str) -> str:
    quote: Optional[str] = None
    depth = 0
    i = 0
    while i < len(font_family):
        ch = font_family[i]
        if quote:
            if ch == "\\":
                i += 2
                continue
            if ch == quote:
                quote = None
            i += 1
            continue
        if ch in ('"', "'"):
            prev = font_family[i - 1] if i > 0 else ""
            if (not prev) or prev.isspace() or prev in ",(":
                quote = ch
                i += 1
                continue
        if ch == "(":
            depth += 1
            i += 1
            continue
        if ch == ")":
            depth = max(0, depth - 1)
            i += 1
            continue
        if ch == "," and depth == 0:
            return font_family[:i]
        i += 1
    return font_family


def _normalize_primary_font_name(font_family: Optional[str]) -> Optional[str]:
    if not font_family:
        return None
    raw = str(font_family).strip()
    if not raw:
        return None

    first = _split_first_font_family_token(raw).strip().strip("'\"")
    if not first:
        return None
    lower = first.lower()
    if lower.startswith("var("):
        return None
    if lower in _GENERIC_FONT_FAMILY_KEYWORDS:
        return None
    return first


def _strip_font_fallback(font_family: Optional[str]) -> Optional[str]:
    """从 CSS font-family 链中提取第一个字体名（写回 PPTX 时使用）"""
    return _normalize_primary_font_name(font_family)


def _resolve_theme_font_reference(
    font_name: Optional[str],
    theme_fonts: Optional[Dict[str, str]] = None,
) -> Optional[str]:
    """
    解析 PPTX 中的主题字体引用。

    PPTX XML 中的字体名可能是主题引用：
    - "+mj-lt" → majorFont latin
    - "+mn-lt" → minorFont latin
    - "+mj-ea" → majorFont east asian
    - "+mn-ea" → minorFont east asian

    如果是主题引用，替换为实际字体名；否则原样返回。
    """
    if not font_name:
        return None
    raw_name = str(font_name).strip()
    if not raw_name:
        return None
    if not raw_name.startswith("+"):
        return raw_name
    if not theme_fonts:
        return None

    ref_map = {
        "+mj-lt": "major_latin",
        "+mn-lt": "minor_latin",
        "+mj-ea": "major_ea",
        "+mn-ea": "minor_ea",
        "+mj-cs": "major_cs",
        "+mn-cs": "minor_cs",
    }
    key = ref_map.get(raw_name.lower())
    if key:
        resolved = theme_fonts.get(key)
        if isinstance(resolved, str) and resolved.strip():
            return resolved.strip()
    return None


def _font_with_fallback(font_name: Optional[str]) -> Optional[str]:
    """为字体名添加 web-safe fallback 链"""
    primary = _normalize_primary_font_name(font_name)
    if not primary:
        return None
    # 跳过未解析的主题引用
    if primary.startswith("+"):
        return None
    # 精确匹配 + 忽略大小写匹配
    fb = _FONT_FALLBACK_MAP.get(primary)
    if not fb:
        lower = primary.lower()
        for key, value in _FONT_FALLBACK_MAP.items():
            if key.lower() == lower:
                fb = value
                break
    if fb:
        return fb
    # 不在映射中的字体，保留主字体名 + 通用 CJK fallback
    escaped = primary.replace("\\", "\\\\").replace("'", "\\'")
    return f"'{escaped}', {_DEFAULT_CJK_WEB_FALLBACK}"


# 默认幻灯片尺寸（像素）
#  canvas 统一：默认 1280×720（= DEFAULT_SLIDE_WIDTH_EMU / 9525，
# 与 html-spec 及 PPTX 页面 1:1）。调用方一般显式传项目 canvas，此处仅兜底。
DEFAULT_SLIDE_WIDTH_PX = 1280
DEFAULT_SLIDE_HEIGHT_PX = 720

# 默认 PPTX 尺寸（EMU）: 标准 16:9 宽屏
DEFAULT_SLIDE_WIDTH_EMU = 12192000  # 13.333 inches
DEFAULT_SLIDE_HEIGHT_EMU = 6858000  # 7.5 inches

# LaTeX 语义元数据（写入图片 altText）
LATEX_META_PREFIX = "TABSLIDE_LATEX_V1:"
MAX_LATEX_META_LENGTH = 24_000

# Media 语义元数据（写入媒体 shape altText，保留前端专有语义）
MEDIA_META_PREFIX = "TABSLIDE_MEDIA_V1:"
MAX_MEDIA_META_LENGTH = 8_000

# PPTX 导出时单张图片下载的最大字节数（50 MB），超出则跳过以防 OOM
MAX_EXPORT_IMAGE_BYTES = 50 * 1024 * 1024

MEDIA_MIME_BY_EXT: Dict[str, str] = {
    "mp4": "video/mp4",
    "m4v": "video/mp4",
    "mov": "video/quicktime",
    "webm": "video/webm",
    "ogv": "video/ogg",
    "avi": "video/x-msvideo",
    "wmv": "video/x-ms-wmv",
    "mpeg": "video/mpeg",
    "mpg": "video/mpeg",
    "mp3": "audio/mpeg",
    "m4a": "audio/mp4",
    "aac": "audio/aac",
    "wav": "audio/wav",
    "ogg": "audio/ogg",
    "oga": "audio/ogg",
    "flac": "audio/flac",
    "wma": "audio/x-ms-wma",
}

MEDIA_EXT_BY_MIME: Dict[str, str] = {v: k for k, v in MEDIA_MIME_BY_EXT.items()}
MEDIA_AUDIO_EXTS = {"mp3", "m4a", "aac", "wav", "ogg", "oga", "flac", "wma"}
MEDIA_VIDEO_EXTS = {"mp4", "m4v", "mov", "webm", "ogv", "avi", "wmv", "mpeg", "mpg"}


# ============================================================================
# 单位换算
# ============================================================================


def emu_to_px(emu: int, slide_emu: int, slide_px: int) -> float:
    """EMU → px，按比例换算"""
    if slide_emu == 0:
        return 0
    return round(emu * slide_px / slide_emu, COORD_DECIMALS)


def px_to_emu(px: float, slide_px: int, slide_emu: int) -> int:
    """px → EMU，按比例换算（四舍五入减少往返精度损失）"""
    if slide_px == 0:
        return 0
    return round(px * slide_emu / slide_px)


def _find_sp_pr(sp_elem):
    """在 shape XML 中定位 spPr，兼容 p:spPr / a:spPr 及不同层级。"""
    if sp_elem is None:
        return None
    nsmap_p = "http://schemas.openxmlformats.org/presentationml/2006/main"
    nsmap_a = "http://schemas.openxmlformats.org/drawingml/2006/main"
    sp_pr = sp_elem.find(f"{{{nsmap_p}}}spPr")
    if sp_pr is None:
        sp_pr = sp_elem.find(f"{{{nsmap_a}}}spPr")
    if sp_pr is None:
        sp_pr = sp_elem.find(f".//{{{nsmap_p}}}spPr")
    if sp_pr is None:
        sp_pr = sp_elem.find(f".//{{{nsmap_a}}}spPr")
    return sp_pr


def _remove_shape_theme_style(shape) -> None:
    """清理 add_shape() 自动带入的主题 style，避免无显式样式的元素被 PowerPoint 套默认线/阴影。"""
    try:
        nsmap_p = "http://schemas.openxmlformats.org/presentationml/2006/main"
        style = shape._element.find(f"{{{nsmap_p}}}style")
        if style is not None:
            shape._element.remove(style)
    except Exception:
        pass


def _apply_no_line_fill(shape) -> None:
    """显式写入无边框，覆盖 PowerPoint 对 auto shape 的默认主题线条。"""
    try:
        from lxml import etree
        from pptx.oxml.ns import qn

        nsmap_a = "http://schemas.openxmlformats.org/drawingml/2006/main"
        sp_pr = _find_sp_pr(shape._element)
        if sp_pr is None:
            return

        ln = sp_pr.find(f"{{{nsmap_a}}}ln")
        if ln is None:
            ln = etree.SubElement(sp_pr, qn("a:ln"))
        ln.attrib.clear()
        for child in list(ln):
            ln.remove(child)
        etree.SubElement(ln, qn("a:noFill"))
    except Exception:
        pass


def _export_dpi_scale(slide_width_emu: int, canvas_width: int) -> float:
    slide_width_inches = slide_width_emu / EMU_PER_INCH
    actual_dpi = canvas_width / slide_width_inches if slide_width_inches > 0 else 96
    if actual_dpi <= 0:
        return 1.0
    return 96.0 / actual_dpi


def _find_transform_xfrm(sp_elem):
    """
    定位当前 shape 自身的 xfrm，避免误命中后代节点（例如 extLst 内部的 xfrm）。
    """
    if sp_elem is None:
        return None

    nsmap_p = "http://schemas.openxmlformats.org/presentationml/2006/main"
    nsmap_a = "http://schemas.openxmlformats.org/drawingml/2006/main"

    # 1) 优先取根节点直属 spPr（最准确）
    sp_pr = sp_elem.find(f"{{{nsmap_p}}}spPr")
    if sp_pr is None:
        sp_pr = sp_elem.find(f"{{{nsmap_a}}}spPr")
    if sp_pr is not None:
        xfrm = sp_pr.find(f"{{{nsmap_a}}}xfrm")
        if xfrm is not None:
            return xfrm

    # 2) 兜底：兼容少数结构差异（但仍只取 spPr 直属 xfrm）
    sp_pr = _find_sp_pr(sp_elem)
    if sp_pr is not None:
        xfrm = sp_pr.find(f"{{{nsmap_a}}}xfrm")
        if xfrm is not None:
            return xfrm

    return None


def _coerce_bool_flag(raw: Any) -> Optional[bool]:
    """将多种输入规范化为 bool，无法识别时返回 None。"""
    if isinstance(raw, bool):
        return raw
    if isinstance(raw, (int, float)):
        if raw == 1:
            return True
        if raw == 0:
            return False
        return None
    if isinstance(raw, str):
        normalized = raw.strip().lower()
        if normalized in ("1", "true", "yes", "on"):
            return True
        if normalized in ("0", "false", "no", "off"):
            return False
    return None


def _resolve_flip_flag(props: Dict[str, Any], element: Dict[str, Any], key: str) -> Optional[bool]:
    """
    解析 flipH/flipV：props 优先于 element，并保留显式 false 语义。
    """
    if isinstance(props, dict):
        prop_val = _coerce_bool_flag(props.get(key))
        if prop_val is not None:
            return prop_val
    return _coerce_bool_flag(element.get(key) if isinstance(element, dict) else None)


def _parse_z_index(raw: Any) -> Optional[int]:
    """解析 zIndex 为非负整数；无法解析返回 None。"""
    try:
        value = float(raw)
    except (TypeError, ValueError):
        return None
    if not math.isfinite(value):
        return None
    return max(0, int(value))


def _normalize_z_index(raw: Any, fallback: int = 0) -> int:
    parsed = _parse_z_index(raw)
    if parsed is None:
        return max(0, int(fallback))
    return parsed


def _sort_elements_by_z_index(elements: Any) -> List[Dict[str, Any]]:
    """
    对元素按 zIndex 排序：
    - 有效 zIndex：按数值升序（同值保持输入顺序）
    - 无效 zIndex：保留输入顺序并放在有效元素之后
    """
    if not isinstance(elements, list):
        return []

    keyed: List[Tuple[Tuple[int, int, int], Dict[str, Any]]] = []
    for idx, el in enumerate(elements):
        if not isinstance(el, dict):
            continue
        parsed = _parse_z_index(el.get("zIndex"))
        if parsed is None:
            key = (1, idx, idx)
        else:
            key = (0, parsed, idx)
        keyed.append((key, el))

    keyed.sort(key=lambda item: item[0])
    return [el for _, el in keyed]


def _encode_latex_alt_text(payload: Dict[str, Any]) -> Optional[str]:
    """
    将 latex 元素语义编码为 altText，确保 image 降级后仍可还原。
    """
    if not isinstance(payload, dict):
        return None

    latex = payload.get("latex")
    if not isinstance(latex, str):
        return None
    latex = latex.strip()
    if not latex:
        return None

    safe_payload: Dict[str, Any] = {"latex": latex}
    if isinstance(payload.get("svg"), str) and payload.get("svg"):
        safe_payload["svg"] = payload["svg"]
    if isinstance(payload.get("path"), str) and payload.get("path"):
        safe_payload["path"] = payload["path"]
    if isinstance(payload.get("viewBox"), (list, tuple)) and len(payload["viewBox"]) == 2:
        safe_payload["viewBox"] = [payload["viewBox"][0], payload["viewBox"][1]]
    if isinstance(payload.get("color"), str) and payload.get("color"):
        safe_payload["color"] = payload["color"]
    if payload.get("strokeWidth") is not None:
        safe_payload["strokeWidth"] = payload["strokeWidth"]
    if payload.get("fixedRatio") is not None:
        safe_payload["fixedRatio"] = payload["fixedRatio"]

    import base64 as _b64
    import json as _json

    def _encode(obj: Dict[str, Any]) -> str:
        raw = _json.dumps(obj, ensure_ascii=False, separators=(",", ":")).encode("utf-8")
        return LATEX_META_PREFIX + _b64.b64encode(raw).decode("ascii")

    encoded = _encode(safe_payload)
    if len(encoded) <= MAX_LATEX_META_LENGTH:
        return encoded

    if "svg" in safe_payload:
        safe_payload.pop("svg", None)
        encoded = _encode(safe_payload)
        if len(encoded) <= MAX_LATEX_META_LENGTH:
            return encoded

    if "path" in safe_payload:
        safe_payload.pop("path", None)
        encoded = _encode(safe_payload)
        if len(encoded) <= MAX_LATEX_META_LENGTH:
            return encoded

    minimal_payload: Dict[str, Any] = {
        "latex": f"{latex[:2048]}..." if len(latex) > 2048 else latex,
    }
    if "color" in safe_payload:
        minimal_payload["color"] = safe_payload["color"]
    if "strokeWidth" in safe_payload:
        minimal_payload["strokeWidth"] = safe_payload["strokeWidth"]
    if "fixedRatio" in safe_payload:
        minimal_payload["fixedRatio"] = safe_payload["fixedRatio"]

    return _encode(minimal_payload)


def _escape_xml_text(raw: str) -> str:
    return (
        raw.replace("&", "&amp;")
        .replace("<", "&lt;")
        .replace(">", "&gt;")
        .replace('"', "&quot;")
        .replace("'", "&apos;")
    )


def _build_latex_placeholder_svg_data_url(
    latex: str,
    color: str = "#111111",
    width_px: Optional[float] = None,
    height_px: Optional[float] = None,
) -> Optional[str]:
    """
    构建 LaTeX 导出兜底 SVG（文本占位），避免公式在缺失 svg/path/raster 时直接丢失。
    """
    if not isinstance(latex, str):
        return None
    source = " ".join(latex.strip().split())
    if not source:
        return None

    try:
        w = int(round(float(width_px))) if width_px is not None else 320
    except (TypeError, ValueError):
        w = 320
    try:
        h = int(round(float(height_px))) if height_px is not None else 96
    except (TypeError, ValueError):
        h = 96

    safe_w = max(120, min(2400, w))
    safe_h = max(40, min(900, h))
    safe_source = source[:237] + "..." if len(source) > 240 else source
    safe_text = _escape_xml_text(safe_source)
    safe_color = _escape_xml_text((color or "#111111").strip() or "#111111")
    font_size = max(12, min(28, int(round(safe_h * 0.28))))
    y = int(round(safe_h / 2))

    svg = (
        f'<svg xmlns="http://www.w3.org/2000/svg" '
        f'viewBox="0 0 {safe_w} {safe_h}" width="{safe_w}" height="{safe_h}" '
        f'preserveAspectRatio="xMidYMid meet">'
        f'<rect x="0" y="0" width="{safe_w}" height="{safe_h}" fill="transparent"/>'
        f'<text x="12" y="{y}" fill="{safe_color}" dominant-baseline="middle" text-anchor="start" '
        f'font-family="Times New Roman, serif" font-size="{font_size}">{safe_text}</text>'
        f'</svg>'
    )

    import base64 as _b64

    svg_b64 = _b64.b64encode(svg.encode("utf-8")).decode("ascii")
    return f"data:image/svg+xml;base64,{svg_b64}"


def _parse_ooxml_percentage_value(raw: Any, allow_negative: bool = False) -> Optional[float]:
    """
    解析 OOXML 百分比值（单位 1/100000）到 0~1 浮点数。

    例如：
    - 50000  -> 0.5
    - 100000 -> 1.0
    """
    if raw is None:
        return None
    try:
        parsed = float(str(raw).strip())
    except (TypeError, ValueError):
        return None

    if not allow_negative:
        parsed = max(0.0, parsed)
    return parsed / 100000.0


def _extract_alpha_from_color_transforms(color_el, decimals: int = 4) -> Optional[float]:
    """
    从颜色节点（srgbClr/schemeClr 等）提取有效 alpha。

    支持常见透明度变换：
    - <a:alpha val="..."/>      绝对 alpha
    - <a:alphaMod val="..."/>   alpha 乘法
    - <a:alphaModFix amt="..."/> alpha 固定乘法
    - <a:alphaOff val="..."/>   alpha 偏移
    """
    if color_el is None:
        return None

    alpha_value = 1.0
    has_alpha_transform = False
    try:
        for child in list(color_el):
            tag = child.tag.split("}")[-1] if isinstance(child.tag, str) and "}" in child.tag else str(child.tag)
            if tag == "alpha":
                parsed = _parse_ooxml_percentage_value(child.get("val") or child.get("amt"))
                if parsed is None:
                    continue
                alpha_value = parsed
                has_alpha_transform = True
            elif tag in ("alphaMod", "alphaModFix"):
                parsed = _parse_ooxml_percentage_value(child.get("val") or child.get("amt"))
                if parsed is None:
                    continue
                alpha_value *= parsed
                has_alpha_transform = True
            elif tag == "alphaOff":
                parsed = _parse_ooxml_percentage_value(
                    child.get("val") or child.get("amt"),
                    allow_negative=True,
                )
                if parsed is None:
                    continue
                alpha_value += parsed
                has_alpha_transform = True
    except Exception:
        return None

    if not has_alpha_transform:
        return None
    return round(max(0.0, min(1.0, alpha_value)), decimals)


def _extract_solid_fill_alpha(solid_fill_el, decimals: int = 4) -> Optional[float]:
    """从 <a:solidFill> 节点中提取颜色 alpha（支持 alpha/alphaMod/alphaModFix）。"""
    if solid_fill_el is None:
        return None
    try:
        for color_el in list(solid_fill_el):
            alpha = _extract_alpha_from_color_transforms(color_el, decimals=decimals)
            if alpha is not None:
                return alpha
    except Exception:
        return None
    return None


def _extract_color_alpha(color_elem) -> Optional[float]:
    """从颜色对象提取 alpha（兼容 alpha/alphaMod/alphaModFix/alphaOff）。"""
    if color_elem is None:
        return None
    try:
        # 尝试获取底层 XML 元素
        xml_el = None
        if hasattr(color_elem, "_color"):
            # python-pptx ColorFormat → 内部 _color 对象
            inner = color_elem._color
            if hasattr(inner, "_xFill"):
                # solidFill 下的颜色元素
                xfill = inner._xFill
                if xfill is not None:
                    for child in xfill:
                        alpha = _extract_alpha_from_color_transforms(child, decimals=2)
                        if alpha is not None:
                            return alpha
        # 直接尝试 element 属性
        if hasattr(color_elem, "element"):
            xml_el = color_elem.element
        elif hasattr(color_elem, "_element"):
            xml_el = color_elem._element
        if xml_el is not None:
            tag = xml_el.tag.split("}")[-1] if isinstance(xml_el.tag, str) and "}" in xml_el.tag else str(xml_el.tag)
            if tag == "solidFill":
                return _extract_solid_fill_alpha(xml_el, decimals=2)
            return _extract_alpha_from_color_transforms(xml_el, decimals=2)
    except Exception:
        pass
    return None


def _alignment_to_str(alignment, default: str = "left") -> str:
    """将 PPTX 段落对齐枚举安全映射到前端字符串。"""
    if alignment is None:
        return default

    # 优先用枚举名，规避不同 python-pptx 版本的 int 值差异
    try:
        align_name = getattr(alignment, "name", None)
        if align_name:
            name_map = {
                "LEFT": "left",
                "CENTER": "center",
                "RIGHT": "right",
                "JUSTIFY": "justify",
                "JUSTIFY_LOW": "justify",
                "DISTRIBUTE": "justify",
                "THAI_DISTRIBUTE": "justify",
            }
            mapped = name_map.get(str(align_name).upper())
            if mapped:
                return mapped
    except Exception:
        pass

    # 兜底整型值（兼容历史数据）
    try:
        align_val = int(alignment)
    except Exception:
        return default

    align_map = {
        0: "left",
        1: "left",
        2: "center",
        3: "right",
        4: "justify",
        5: "justify",
        6: "justify",
    }
    return align_map.get(align_val, default)


def _hex_to_rgba(hex_color: str, alpha: float) -> str:
    """#RRGGBB + alpha(0~1) → rgba(R,G,B,A)"""
    h = hex_color.lstrip("#")
    r, g, b = int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16)
    return f"rgba({r},{g},{b},{alpha})"


def _parse_css_color(color_str: str) -> Tuple[str, Optional[float]]:
    """
    解析 CSS 颜色字符串，统一返回 (hex_6, alpha)。

    支持格式：
    - "#4A6D8C" → ("#4A6D8C", None)
    - "rgba(74,109,140,0.2)" → ("#4A6D8C", 0.2)
    - "#4A6D8C" → ("#4A6D8C", None)

    写回 PPTX 时使用：hex → RGBColor, alpha → <a:alpha>
    """
    if not color_str:
        return ("#000000", None)

    color_str = color_str.strip()

    if color_str.startswith("rgba"):
        import re as _re
        m = _re.match(r"rgba?\((\d+),\s*(\d+),\s*(\d+),?\s*([\d.]*)\)", color_str)
        if m:
            r, g, b = int(m.group(1)), int(m.group(2)), int(m.group(3))
            a = float(m.group(4)) if m.group(4) else 1.0
            return (f"#{r:02x}{g:02x}{b:02x}", a if a < 1.0 else None)

    if color_str.startswith("rgb("):
        import re as _re
        m = _re.match(r"rgb\((\d+),\s*(\d+),\s*(\d+)\)", color_str)
        if m:
            r, g, b = int(m.group(1)), int(m.group(2)), int(m.group(3))
            return (f"#{r:02x}{g:02x}{b:02x}", None)

    # #RGB / #RGBA / #RRGGBB / #RRGGBBAA
    if color_str.startswith("#"):
        h = color_str[1:]
        if len(h) == 3:
            return (f"#{h[0]}{h[0]}{h[1]}{h[1]}{h[2]}{h[2]}", None)
        if len(h) == 4:
            alpha = int(h[3] + h[3], 16) / 255
            return (
                f"#{h[0]}{h[0]}{h[1]}{h[1]}{h[2]}{h[2]}",
                round(alpha, 3) if alpha < 1.0 else None,
            )
        if len(h) == 6:
            return (f"#{h}", None)
        if len(h) == 8:
            alpha = int(h[6:8], 16) / 255
            return (f"#{h[:6]}", round(alpha, 3) if alpha < 1.0 else None)

    # hsl()/hsla()
    if color_str.startswith("hsl"):
        import re as _re
        m = _re.match(
            r"hsla?\(\s*([\d.]+)\s*,\s*([\d.]+)%\s*,\s*([\d.]+)%\s*(?:,\s*([\d.]+)\s*)?\)",
            color_str,
        )
        if m:
            h_val = float(m.group(1))
            s_val = float(m.group(2)) / 100
            l_val = float(m.group(3))  / 100
            alpha = float(m.group(4)) if m.group(4) else 1.0
            a = s_val * min(l_val, 1 - l_val)
            def _f(n: float) -> int:
                k = (n + h_val / 30) % 12
                return max(0, min(255, round(255 * (l_val - a * max(min(k - 3, 9 - k, 1), -1)))))
            r, g, b = _f(0), _f(8), _f(4)
            return (f"#{r:02x}{g:02x}{b:02x}", alpha if alpha < 1.0 else None)

    # 假定 hex 格式
    return (color_str, None)


def _format_numeric(value: float, decimals: int = 2) -> str:
    """格式化数值，尽量保留精度并去掉无意义尾零。"""
    try:
        rounded = round(float(value), decimals)
    except (TypeError, ValueError):
        return str(value)
    if abs(rounded - round(rounded)) < 1e-9:
        return str(int(round(rounded)))
    return f"{rounded:.{decimals}f}".rstrip("0").rstrip(".")


def _parse_css_length_value(raw: Any, default_unit: str = "pt") -> Optional[Tuple[float, str]]:
    """
    解析 CSS 长度（支持 pt/px/无单位）。

    返回 (数值, 单位)：
    - 单位为 "pt" / "px"
    - 无单位默认按 "pt" 处理（保持历史兼容）
    """
    import re

    if raw is None:
        return None
    if isinstance(raw, (int, float)):
        return (float(raw), default_unit)

    s = str(raw).strip().lower()
    if not s:
        return None

    m = re.match(r"^(-?\d+(?:\.\d+)?)\s*(pt|px)?$", s)
    if not m:
        return None

    val = float(m.group(1))
    unit = m.group(2) or default_unit
    return (val, unit)


def _css_length_to_pt(raw: Any) -> Optional[float]:
    """CSS 长度 → pt（支持 px/pt）。"""
    parsed = _parse_css_length_value(raw, default_unit="pt")
    if not parsed:
        return None
    val, unit = parsed
    if unit == "px":
        return val * PT_PER_PX
    return val


def _css_length_to_px(raw: Any) -> Optional[float]:
    """CSS 长度 → px（支持 px/pt）。"""
    parsed = _parse_css_length_value(raw, default_unit="px")
    if not parsed:
        return None
    val, unit = parsed
    if unit == "pt":
        return val * PX_PER_PT
    return val


def _parse_css_line_height_ratio(raw: Any) -> Optional[float]:
    """解析 CSS line-height 为倍数值（如 1.5，150% -> 1.5）。"""
    if raw is None:
        return None
    if isinstance(raw, (int, float)):
        try:
            v = float(raw)
            return v if v > 0 else None
        except (TypeError, ValueError):
            return None

    s = str(raw).strip().lower()
    if not s:
        return None

    if s.endswith("%"):
        try:
            pct = float(s[:-1])
            return round(pct / 100.0, 4) if pct > 0 else None
        except (TypeError, ValueError):
            return None

    try:
        v = float(s)
        return v if v > 0 else None
    except (TypeError, ValueError):
        return None


def _apply_color_transforms(base_hex: str, color_element) -> str:
    """
    对基础色应用 OOXML 颜色修饰符（lumMod/lumOff/tint/shade/alpha 等）。

    PPTX 中主题色经常带有修饰符：
    - lumMod (亮度倍率): 50000 = 50%
    - lumOff (亮度偏移): 40000 = +40%
    - tint (淡化): 60000 = 向白色混合 60%
    - shade (加深): 75000 = 向黑色混合 75%
    - satMod (饱和度调整)

    公式参考 ECMA-376 ST_SchemeColor 规范
    """
    nsmap_a = "http://schemas.openxmlformats.org/drawingml/2006/main"

    if color_element is None:
        return base_hex

    # 提取修饰符
    lum_mod = None
    lum_off = None
    tint_val = None
    shade_val = None
    sat_mod = None

    for child in color_element:
        tag = child.tag.split("}")[-1] if "}" in child.tag else child.tag
        val_str = child.get("val")
        if not val_str:
            continue
        try:
            val = int(val_str)
        except (ValueError, TypeError):
            continue

        if tag == "lumMod":
            lum_mod = val / 100000.0
        elif tag == "lumOff":
            lum_off = val / 100000.0
        elif tag == "tint":
            tint_val = val / 100000.0
        elif tag == "shade":
            shade_val = val / 100000.0
        elif tag == "satMod":
            sat_mod = val / 100000.0

    # 无修饰符，返回原色
    if lum_mod is None and lum_off is None and tint_val is None and shade_val is None and sat_mod is None:
        return base_hex

    # 解析 base_hex → RGB
    hex_clean = base_hex.lstrip("#")[:6]
    try:
        r = int(hex_clean[0:2], 16)
        g = int(hex_clean[2:4], 16)
        b = int(hex_clean[4:6], 16)
    except (ValueError, IndexError):
        return base_hex

    # 转 HSL（0-1 范围）
    import colorsys
    r_f, g_f, b_f = r / 255.0, g / 255.0, b / 255.0
    h, l, s = colorsys.rgb_to_hls(r_f, g_f, b_f)  # 注意 colorsys 返回 (h, l, s)

    # 应用 tint（向白色混合）
    if tint_val is not None:
        l = l * tint_val + (1.0 - tint_val)

    # 应用 shade（向黑色混合）
    if shade_val is not None:
        l = l * shade_val

    # 应用 lumMod + lumOff
    if lum_mod is not None:
        l = l * lum_mod
    if lum_off is not None:
        l = l + lum_off

    l = max(0.0, min(1.0, l))

    # 应用 satMod（饱和度倍率）
    if sat_mod is not None:
        s = max(0.0, min(1.0, s * sat_mod))

    r_out, g_out, b_out = colorsys.hls_to_rgb(h, l, s)
    return f"#{int(r_out * 255):02x}{int(g_out * 255):02x}{int(b_out * 255):02x}"


def _extract_color_with_transforms(
    color_xml_parent,
    theme_color_map: Optional[Dict[str, str]] = None,
) -> Optional[str]:
    """
    从 XML 颜色元素（<a:solidFill> 的子元素）提取颜色值，包含修饰符计算。
    同时处理 srgbClr 和 schemeClr 两种情况。
    """
    nsmap_a = "http://schemas.openxmlformats.org/drawingml/2006/main"
    if color_xml_parent is None:
        return None

    # 先查 srgbClr（固定颜色）
    srgb = color_xml_parent.find(f"{{{nsmap_a}}}srgbClr")
    if srgb is not None:
        base_hex = f"#{srgb.get('val', '000000')}"
        return _apply_color_transforms(base_hex, srgb)

    # 再查 schemeClr（主题颜色）
    scheme = color_xml_parent.find(f"{{{nsmap_a}}}schemeClr")
    if scheme is not None:
        scheme_val = (scheme.get("val", "") or "").strip()
        base_hex = _resolve_scheme_color(scheme_val, theme_color_map)
        if base_hex:
            return _apply_color_transforms(base_hex, scheme)

    return None


def color_to_hex(color, theme_color_map: Optional[Dict[str, str]] = None) -> Optional[str]:
    """python-pptx 颜色对象 → #RRGGBB 字符串"""
    if color is None:
        return None
    try:
        if hasattr(color, "rgb") and color.rgb is not None:
            return f"#{color.rgb}"
        if hasattr(color, "theme_color") and color.theme_color is not None:
            # python-pptx theme_color 返回 MSO_THEME_COLOR 枚举，
            # str(枚举) 输出类似 "ACCENT_1 (5)"，不是纯数字，
            # 所以必须用 int() 取值再转 str 做 key lookup。
            try:
                tc_int_key = str(int(color.theme_color))
            except (ValueError, TypeError):
                tc_int_key = None

            if tc_int_key and theme_color_map:
                canonical_key = _normalize_theme_color_key(tc_int_key)
                mapped = theme_color_map.get(tc_int_key)
                if not mapped and canonical_key:
                    mapped = theme_color_map.get(canonical_key)
                if mapped:
                    return mapped
            # 降级：常用主题色默认值
            return _default_theme_color(color.theme_color)
    except Exception:
        pass
    return None


# Office 默认主题色映射（Office 2013+ 默认主题 "Office"）
_DEFAULT_THEME_COLORS: Dict[int, str] = {
    1: "#000000",   # DARK_1 (dk1) — 通常是黑色
    2: "#FFFFFF",   # LIGHT_1 (lt1) — 通常是白色
    3: "#44546A",   # DARK_2 (dk2)
    4: "#E7E6E6",   # LIGHT_2 (lt2)
    5: "#4472C4",   # ACCENT_1
    6: "#ED7D31",   # ACCENT_2
    7: "#A5A5A5",   # ACCENT_3
    8: "#FFC000",   # ACCENT_4
    9: "#5B9BD5",   # ACCENT_5
    10: "#70AD47",  # ACCENT_6
    11: "#0563C1",  # HYPERLINK
    12: "#954F72",  # FOLLOWED_HYPERLINK
}

_THEME_SCHEME_TO_KEY: Dict[str, str] = {
    "dk1": "1",
    "lt1": "2",
    "dk2": "3",
    "lt2": "4",
    "accent1": "5",
    "accent2": "6",
    "accent3": "7",
    "accent4": "8",
    "accent5": "9",
    "accent6": "10",
    "hlink": "11",
    "folhlink": "12",
    # alias（OOXML 在某些节点会出现 bg/tx 语义键）
    "tx1": "1",
    "bg1": "2",
    "tx2": "3",
    "bg2": "4",
}

_THEME_KEY_ALIASES: Dict[str, str] = {
    "13": "1",  # TEXT_1 -> DARK_1
    "14": "2",  # BACKGROUND_1 -> LIGHT_1
    "15": "3",  # TEXT_2 -> DARK_2
    "16": "4",  # BACKGROUND_2 -> LIGHT_2
}

_THEME_KEY_ALIAS_REVERSE: Dict[str, Tuple[str, ...]] = {
    "1": ("13",),
    "2": ("14",),
    "3": ("15",),
    "4": ("16",),
}


def _normalize_theme_color_key(key: Optional[str]) -> Optional[str]:
    if key is None:
        return None
    key_str = str(key).strip()
    if not key_str:
        return None
    return _THEME_KEY_ALIASES.get(key_str, key_str)


def _resolve_theme_color_by_key(
    key: Optional[str],
    theme_color_map: Optional[Dict[str, str]] = None,
) -> Optional[str]:
    """按主题色 key（字符串数字）解析十六进制颜色。"""
    if key is None:
        return None
    raw_key = str(key).strip()
    if not raw_key:
        return None

    canonical_key = _normalize_theme_color_key(raw_key)

    if theme_color_map:
        mapped = theme_color_map.get(raw_key)
        if not mapped and canonical_key:
            mapped = theme_color_map.get(canonical_key)
        if mapped:
            return mapped

    if canonical_key and canonical_key.isdigit():
        return _DEFAULT_THEME_COLORS.get(int(canonical_key))
    return None


def _resolve_scheme_color(
    scheme_val: str,
    theme_color_map: Optional[Dict[str, str]] = None,
) -> Optional[str]:
    """按 schemeClr val（如 accent1/bg1/tx1）解析颜色。"""
    key = _THEME_SCHEME_TO_KEY.get((scheme_val or "").strip().lower())
    return _resolve_theme_color_by_key(key, theme_color_map)


_TEXT_THEME_KEY_MAP: Dict[str, str] = {
    "1": "tx1",
    "13": "tx1",
    "dk1": "tx1",
    "dark1": "tx1",
    "dark_1": "tx1",
    "tx1": "tx1",
    "text1": "tx1",
    "text_1": "tx1",
    "2": "bg1",
    "14": "bg1",
    "lt1": "bg1",
    "light1": "bg1",
    "light_1": "bg1",
    "bg1": "bg1",
    "background1": "bg1",
    "background_1": "bg1",
    "3": "tx2",
    "15": "tx2",
    "dk2": "tx2",
    "dark2": "tx2",
    "dark_2": "tx2",
    "tx2": "tx2",
    "text2": "tx2",
    "text_2": "tx2",
    "4": "bg2",
    "16": "bg2",
    "lt2": "bg2",
    "light2": "bg2",
    "light_2": "bg2",
    "bg2": "bg2",
    "background2": "bg2",
    "background_2": "bg2",
    "accent1": "accent1",
    "accent_1": "accent1",
    "accent2": "accent2",
    "accent_2": "accent2",
    "accent3": "accent3",
    "accent_3": "accent3",
    "accent4": "accent4",
    "accent_4": "accent4",
    "accent5": "accent5",
    "accent_5": "accent5",
    "accent6": "accent6",
    "accent_6": "accent6",
    "11": "hlink",
    "hlink": "hlink",
    "hyperlink": "hlink",
    "12": "folhlink",
    "folhlink": "folhlink",
    "followed_hyperlink": "folhlink",
}

_TEXT_THEME_KEY_TO_MSO: Dict[str, str] = {
    "tx1": "DARK_1",
    "bg1": "LIGHT_1",
    "tx2": "DARK_2",
    "bg2": "LIGHT_2",
    "accent1": "ACCENT_1",
    "accent2": "ACCENT_2",
    "accent3": "ACCENT_3",
    "accent4": "ACCENT_4",
    "accent5": "ACCENT_5",
    "accent6": "ACCENT_6",
    "hlink": "HYPERLINK",
    "folhlink": "FOLLOWED_HYPERLINK",
}


def _normalize_text_theme_key(raw_key: Optional[Any]) -> Optional[str]:
    if raw_key is None:
        return None
    raw = str(raw_key).strip().lower()
    if not raw:
        return None
    return _TEXT_THEME_KEY_MAP.get(raw)


def _theme_enum_to_text_theme_key(theme_color_enum: Optional[Any]) -> Optional[str]:
    """MSO_THEME_COLOR 枚举值 -> 文本主题色 key（tx1/bg1/accent1...）"""
    if theme_color_enum is None:
        return None
    try:
        enum_int = int(theme_color_enum)
    except (TypeError, ValueError):
        return None
    canonical = _normalize_theme_color_key(str(enum_int)) or str(enum_int)
    return _normalize_text_theme_key(canonical)


def _extract_theme_key_from_solid_fill(solid_fill_el) -> Optional[str]:
    """
    从 <a:solidFill> 提取主题色 key。

    仅在 schemeClr 无亮度/色调变换时返回 key，避免将 tint/shade 误简化为纯主题色。
    """
    nsmap_a = "http://schemas.openxmlformats.org/drawingml/2006/main"
    if solid_fill_el is None:
        return None
    scheme = solid_fill_el.find(f"{{{nsmap_a}}}schemeClr")
    if scheme is None:
        return None

    # 仅允许 alpha 变换，其它变换会改变实际颜色，不应简化为纯 theme key
    for child in list(scheme):
        tag = child.tag.split("}")[-1] if isinstance(child.tag, str) and "}" in child.tag else str(child.tag)
        if tag != "alpha":
            return None

    return _normalize_text_theme_key(scheme.get("val"))


def _extract_theme_key_from_color_obj(color_obj) -> Optional[str]:
    """从 python-pptx ColorFormat 对象提取主题色 key（tx1/bg1/accent1...）。"""
    if color_obj is None:
        return None
    try:
        if hasattr(color_obj, "theme_color") and color_obj.theme_color is not None:
            return _theme_enum_to_text_theme_key(color_obj.theme_color)
    except Exception:
        return None
    return None


def _apply_theme_color_to_solid_parent(
    parent_el,
    theme_key: Any,
    color_str_for_alpha: Optional[Any] = None,
    transforms: Optional[Dict[str, float]] = None,
) -> bool:
    """
    在任意包含 <a:solidFill> 的父节点上写入 schemeClr 主题色。

    parent_el 可为：
    - spPr（形状填充）
    - ln（线条/轮廓）
    - tcPr（表格单元格填充）
    - table border side（a:top/a:right/...）
    """
    if parent_el is None:
        return False
    try:
        from lxml import etree
        from pptx.oxml.ns import qn

        nsmap_a = "http://schemas.openxmlformats.org/drawingml/2006/main"
        normalized_key = _normalize_text_theme_key(theme_key)
        if not normalized_key:
            return False

        solid_fill = parent_el.find(f"{{{nsmap_a}}}solidFill")
        if solid_fill is None:
            solid_fill = etree.SubElement(parent_el, qn("a:solidFill"))
        for child in list(solid_fill):
            solid_fill.remove(child)

        scheme_el = etree.SubElement(solid_fill, qn("a:schemeClr"))
        scheme_el.set("val", normalized_key)

        # 写入颜色变换子元素（tint/shade/lumMod/lumOff）
        if transforms and isinstance(transforms, dict):
            for t_name, t_val in transforms.items():
                if t_name == "alpha":
                    continue
                try:
                    t_el = etree.SubElement(scheme_el, qn(f"a:{t_name}"))
                    t_el.set("val", str(int(max(0.0, min(1.0, float(t_val))) * 100000)))
                except Exception:
                    pass

        if color_str_for_alpha is not None:
            try:
                _, alpha = _parse_css_color(str(color_str_for_alpha))
            except Exception:
                alpha = None
            if alpha is not None and alpha < 1.0:
                alpha_el = etree.SubElement(scheme_el, qn("a:alpha"))
                alpha_el.set("val", str(int(max(0.0, min(1.0, float(alpha))) * 100000)))

        return True
    except Exception:
        return False


def _extract_run_theme_color_key(run_obj) -> Optional[str]:
    """从 run XML 或 run.font.color 提取主题色 key（tx1/bg1/accent1...）。"""
    try:
        nsmap_a = "http://schemas.openxmlformats.org/drawingml/2006/main"
        r_pr = run_obj._r.find(f"{{{nsmap_a}}}rPr")
        if r_pr is not None:
            solid_fill = r_pr.find(f"{{{nsmap_a}}}solidFill")
            if solid_fill is not None:
                theme_key = _extract_theme_key_from_solid_fill(solid_fill)
                if theme_key:
                    return theme_key
                # 存在 schemeClr 但带有变换（如 tint/shade）时，不应简化为纯 theme key
                if solid_fill.find(f"{{{nsmap_a}}}schemeClr") is not None:
                    return None
    except Exception:
        pass

    try:
        font_color = run_obj.font.color if hasattr(run_obj, "font") else None
        if font_color is not None and hasattr(font_color, "theme_color") and font_color.theme_color is not None:
            return _theme_enum_to_text_theme_key(font_color.theme_color)
    except Exception:
        pass
    return None


def _default_theme_color(theme_color_enum) -> Optional[str]:
    """根据 MSO_THEME_COLOR 枚举值返回默认颜色"""
    try:
        val = int(theme_color_enum)
        canonical = _normalize_theme_color_key(str(val))
        if canonical and canonical.isdigit():
            val = int(canonical)
        color = _DEFAULT_THEME_COLORS.get(val)
        if color:
            return color
    except (ValueError, TypeError):
        pass
    return None


def _extract_theme_color_map(prs) -> Dict[str, str]:
    """从演示文稿 theme 中提取主题色映射表"""
    theme_map: Dict[str, str] = {}
    try:
        # 获取第一个 slideMaster 的主题
        if prs.slide_masters and len(prs.slide_masters) > 0:
            master = prs.slide_masters[0]
            nsmap = {
                "a": "http://schemas.openxmlformats.org/drawingml/2006/main",
                "p": "http://schemas.openxmlformats.org/presentationml/2006/main",
            }

            for rel in master.part.rels.values():
                if "theme" in rel.reltype:
                    theme_part = rel.target_part
                    theme_xml = theme_part.element
                    clr_scheme = theme_xml.find(".//a:clrScheme", nsmap)
                    if clr_scheme is not None:
                        color_name_map = {
                            "dk1": "1",  # DARK_1
                            "lt1": "2",  # LIGHT_1
                            "dk2": "3",  # DARK_2
                            "lt2": "4",  # LIGHT_2
                            "accent1": "5",
                            "accent2": "6",
                            "accent3": "7",
                            "accent4": "8",
                            "accent5": "9",
                            "accent6": "10",
                            "hlink": "11",
                            "folhlink": "12",
                        }
                        for child in clr_scheme:
                            tag = child.tag.split("}")[-1] if "}" in child.tag else child.tag
                            tc_key = color_name_map.get(str(tag).lower())
                            if tc_key:
                                srgb = child.find("a:srgbClr", nsmap)
                                if srgb is not None:
                                    color_hex = f"#{srgb.get('val', '000000')}"
                                    theme_map[tc_key] = color_hex
                                    for alias in _THEME_KEY_ALIAS_REVERSE.get(tc_key, ()):
                                        theme_map[alias] = color_hex
                                else:
                                    sys_clr = child.find("a:sysClr", nsmap)
                                    if sys_clr is not None:
                                        last_clr = sys_clr.get("lastClr")
                                        if last_clr:
                                            color_hex = f"#{last_clr}"
                                            theme_map[tc_key] = color_hex
                                            for alias in _THEME_KEY_ALIAS_REVERSE.get(tc_key, ()):
                                                theme_map[alias] = color_hex
                    break
    except Exception as e:
        logger.debug(f"Failed to extract theme color map: {e}")

    return theme_map


def _extract_theme_color_map_for_master(master) -> Dict[str, str]:
    """从单个 slideMaster 的 theme 中提取主题色映射表。"""
    theme_map: Dict[str, str] = {}
    try:
        nsmap = {
            "a": "http://schemas.openxmlformats.org/drawingml/2006/main",
        }
        for rel in master.part.rels.values():
            if "theme" in rel.reltype:
                theme_part = rel.target_part
                theme_xml = theme_part.element
                clr_scheme = theme_xml.find(".//a:clrScheme", nsmap)
                if clr_scheme is not None:
                    color_name_map = {
                        "dk1": "1",  # DARK_1
                        "lt1": "2",  # LIGHT_1
                        "dk2": "3",  # DARK_2
                        "lt2": "4",  # LIGHT_2
                        "accent1": "5",
                        "accent2": "6",
                        "accent3": "7",
                        "accent4": "8",
                        "accent5": "9",
                        "accent6": "10",
                        "hlink": "11",
                        "folhlink": "12",
                    }
                    for child in clr_scheme:
                        tag = child.tag.split("}")[-1] if "}" in child.tag else child.tag
                        tc_key = color_name_map.get(str(tag).lower())
                        if tc_key:
                            srgb = child.find("a:srgbClr", nsmap)
                            if srgb is not None:
                                color_hex = f"#{srgb.get('val', '000000')}"
                                theme_map[tc_key] = color_hex
                                for alias in _THEME_KEY_ALIAS_REVERSE.get(tc_key, ()):
                                    theme_map[alias] = color_hex
                            else:
                                sys_clr = child.find("a:sysClr", nsmap)
                                if sys_clr is not None:
                                    last_clr = sys_clr.get("lastClr")
                                    if last_clr:
                                        color_hex = f"#{last_clr}"
                                        theme_map[tc_key] = color_hex
                                        for alias in _THEME_KEY_ALIAS_REVERSE.get(tc_key, ()):
                                            theme_map[alias] = color_hex
                break
    except Exception as e:
        logger.debug(f"Failed to extract theme color map for master: {e}")
    return theme_map


def _extract_theme_fonts(prs) -> Dict[str, str]:
    """
    从 theme XML 提取主题字体。

    返回 {"major_latin": "Calibri Light", "major_ea": "等线 Light",
           "minor_latin": "Calibri", "minor_ea": "等线"}

    PPTX 字体继承规则：
    - majorFont → 标题类 placeholder (title, ctrTitle)
    - minorFont → 正文类 placeholder (body, subTitle) + 普通文本框
    """
    fonts: Dict[str, str] = {}
    try:
        if prs.slide_masters and len(prs.slide_masters) > 0:
            master = prs.slide_masters[0]
            nsmap = {"a": "http://schemas.openxmlformats.org/drawingml/2006/main"}

            for rel in master.part.rels.values():
                if "theme" in rel.reltype:
                    theme_xml = rel.target_part.element
                    font_scheme = theme_xml.find(".//a:fontScheme", nsmap)
                    if font_scheme is not None:
                        for font_type in ("majorFont", "minorFont"):
                            prefix = "major" if font_type == "majorFont" else "minor"
                            font_el = font_scheme.find(f"a:{font_type}", nsmap)
                            if font_el is not None:
                                latin = font_el.find("a:latin", nsmap)
                                ea = font_el.find("a:ea", nsmap)
                                cs = font_el.find("a:cs", nsmap)
                                if latin is not None and latin.get("typeface"):
                                    fonts[f"{prefix}_latin"] = latin.get("typeface")
                                if ea is not None and ea.get("typeface"):
                                    fonts[f"{prefix}_ea"] = ea.get("typeface")
                                if cs is not None and cs.get("typeface"):
                                    fonts[f"{prefix}_cs"] = cs.get("typeface")
                    break
    except Exception as e:
        logger.debug(f"Failed to extract theme fonts: {e}")

    return fonts


def _pick_theme_font_family(theme_fonts: Optional[Dict[str, str]], prefix: str) -> Optional[str]:
    """按优先级从主题字体集中选择字体族（优先东亚字体）。"""
    if not isinstance(theme_fonts, dict):
        return None
    for suffix in ("ea", "latin", "cs"):
        key = f"{prefix}_{suffix}"
        raw = theme_fonts.get(key)
        if isinstance(raw, str) and raw.strip():
            return raw.strip()
    return None


def _build_theme_payload(
    theme_color_map: Optional[Dict[str, str]] = None,
    theme_fonts: Optional[Dict[str, str]] = None,
) -> Dict[str, Any]:
    """
    构建前端可消费的主题 JSON。

    输出结构与前端 SlideTheme 对齐：
    {
      backgroundColor, fontColor, themeColors[accent1..6],
      bg2Color, tx2Color, hlinkColor, folHlinkColor,
      fontName, headingFontName
    }
    """
    bg_color = _resolve_theme_color_by_key("2", theme_color_map) or "#FFFFFF"
    font_color = _resolve_theme_color_by_key("1", theme_color_map) or "#000000"

    theme_colors: List[str] = []
    for key in ("5", "6", "7", "8", "9", "10"):
        default_hex = _DEFAULT_THEME_COLORS.get(int(key), "#4472C4")
        color = _resolve_theme_color_by_key(key, theme_color_map) or default_hex
        theme_colors.append(color)

    # 扩展主题色：bg2(lt2)/tx2(dk2)/hlink/folHlink
    bg2_color = _resolve_theme_color_by_key("4", theme_color_map) or "#E7E6E6"
    tx2_color = _resolve_theme_color_by_key("3", theme_color_map) or "#44546A"
    hlink_color = _resolve_theme_color_by_key("11", theme_color_map) or "#0563C1"
    fol_hlink_color = _resolve_theme_color_by_key("12", theme_color_map) or "#954F72"

    font_name = _pick_theme_font_family(theme_fonts, "minor") or "Microsoft YaHei"
    heading_font_name = _pick_theme_font_family(theme_fonts, "major") or font_name

    return {
        "backgroundColor": bg_color,
        "fontColor": font_color,
        "themeColors": theme_colors,
        "bg2Color": bg2_color,
        "tx2Color": tx2_color,
        "hlinkColor": hlink_color,
        "folHlinkColor": fol_hlink_color,
        "fontName": font_name,
        "headingFontName": heading_font_name,
    }


def extract_theme_payload(pptx_path: str) -> Dict[str, Any]:
    """
    公共接口：从 PPTX 提取主题配色与主题字体（前端 SlideTheme）。
    """
    from pptx import Presentation

    prs = Presentation(pptx_path)
    theme_color_map = _extract_theme_color_map(prs)
    theme_fonts = _extract_theme_fonts(prs)
    return _build_theme_payload(theme_color_map, theme_fonts)


def _extract_master_text_styles(prs) -> Dict[str, Dict[str, Any]]:
    """
    从 slideMaster 提取默认文本样式 (txStyles)。

    返回 {"title": {"fontSize": 44, "color": "#000000", ...},
           "body": {"fontSize": 32, "color": "#404040", ...},
           "other": {"fontSize": 18, "color": "#404040", ...}}

    这些是 placeholder 元素在没有显式样式时的默认值。
    """
    styles: Dict[str, Dict[str, Any]] = {}
    try:
        if prs.slide_masters and len(prs.slide_masters) > 0:
            master = prs.slide_masters[0]
            nsmap = {
                "a": "http://schemas.openxmlformats.org/drawingml/2006/main",
                "p": "http://schemas.openxmlformats.org/presentationml/2006/main",
            }

            theme_color_map = _extract_theme_color_map(prs)
            tx_styles = master.element.find(".//p:txStyles", nsmap)
            if tx_styles is not None:
                for style_name, xml_name in [("title", "p:titleStyle"), ("body", "p:bodyStyle"), ("other", "p:otherStyle")]:
                    style_el = tx_styles.find(xml_name, nsmap)
                    if style_el is not None:
                        # 取第一级段落属性
                        lvl1 = style_el.find("a:lvl1pPr", nsmap)
                        if lvl1 is not None:
                            st: Dict[str, Any] = {}
                            def_rPr = lvl1.find("a:defRPr", nsmap)
                            if def_rPr is not None:
                                sz = def_rPr.get("sz")
                                if sz:
                                    st["fontSize"] = int(sz) // 100  # hundredths of pt → pt

                                b = def_rPr.get("b")
                                if b == "1":
                                    st["bold"] = True

                                # 颜色
                                solid = def_rPr.find("a:solidFill", nsmap)
                                if solid is not None:
                                    srgb = solid.find("a:srgbClr", nsmap)
                                    if srgb is not None:
                                        st["color"] = f"#{srgb.get('val', '000000')}"
                                    else:
                                        scheme = solid.find("a:schemeClr", nsmap)
                                        if scheme is not None:
                                            scheme_val = scheme.get("val", "")
                                            scheme_to_key = {
                                                "dk1": "1", "lt1": "2", "dk2": "3", "lt2": "4",
                                                "accent1": "5", "accent2": "6", "accent3": "7",
                                                "accent4": "8", "accent5": "9", "accent6": "10",
                                                "hlink": "11", "folhlink": "12",
                                            }
                                            tc_key = scheme_to_key.get(scheme_val)
                                            if tc_key and tc_key in theme_color_map:
                                                st["color"] = theme_color_map[tc_key]

                            if st:
                                styles[style_name] = st

                        # 提取 lvl2-lvl9 的段落属性
                        for lvl_num in range(2, 10):
                            lvl_key = f"a:lvl{lvl_num}pPr"
                            lvl_el = style_el.find(lvl_key, nsmap)
                            if lvl_el is not None:
                                lvl_st: Dict[str, Any] = {}
                                lvl_defRPr = lvl_el.find("a:defRPr", nsmap)
                                if lvl_defRPr is not None:
                                    sz = lvl_defRPr.get("sz")
                                    if sz:
                                        lvl_st["fontSize"] = int(sz) // 100
                                    b = lvl_defRPr.get("b")
                                    if b == "1":
                                        lvl_st["bold"] = True
                                    solid = lvl_defRPr.find("a:solidFill", nsmap)
                                    if solid is not None:
                                        srgb = solid.find("a:srgbClr", nsmap)
                                        if srgb is not None:
                                            lvl_st["color"] = f"#{srgb.get('val', '000000')}"
                                        else:
                                            scheme = solid.find("a:schemeClr", nsmap)
                                            if scheme is not None:
                                                scheme_val = scheme.get("val", "")
                                                scheme_to_key = {
                                                    "dk1": "1", "lt1": "2", "dk2": "3", "lt2": "4",
                                                    "accent1": "5", "accent2": "6", "accent3": "7",
                                                    "accent4": "8", "accent5": "9", "accent6": "10",
                                                    "hlink": "11", "folhlink": "12",
                                                }
                                                tc_key = scheme_to_key.get(scheme_val)
                                                if tc_key and tc_key in theme_color_map:
                                                    lvl_st["color"] = theme_color_map[tc_key]
                                if lvl_st:
                                    styles[f"{style_name}_lvl{lvl_num}"] = lvl_st
    except Exception as e:
        logger.debug(f"Failed to extract master text styles: {e}")

    return styles


def _extract_placeholder_meta(shape) -> Optional[Dict[str, Any]]:
    """
    提取 shape 的 placeholder 元数据。

    返回示例：
    {
        "type": "title",
        "idx": 0,
        "orient": "horz",
        "sz": "full",
    }
    """
    try:
        nsmap = {"p": "http://schemas.openxmlformats.org/presentationml/2006/main"}
        nv_sp_pr = shape._element.find(".//p:nvSpPr", nsmap)
        if nv_sp_pr is None:
            return None
        nv_pr = nv_sp_pr.find("p:nvPr", nsmap)
        if nv_pr is None:
            return None
        ph = nv_pr.find("p:ph", nsmap)
        if ph is None:
            return None

        meta: Dict[str, Any] = {}
        ph_type = ph.get("type")
        if ph_type:
            meta["type"] = str(ph_type)
        ph_idx = ph.get("idx")
        if ph_idx is not None:
            try:
                meta["idx"] = int(str(ph_idx).strip())
            except Exception:
                meta["idx"] = str(ph_idx)
        orient = ph.get("orient")
        if orient:
            meta["orient"] = str(orient)
        sz = ph.get("sz")
        if sz:
            meta["sz"] = str(sz)
        return meta
    except Exception:
        return None


def _detect_placeholder_type(shape) -> Optional[str]:
    """
    检测 shape 是否是 placeholder，返回 placeholder type 字符串。

    PPTX placeholder types:
    - "title", "ctrTitle" → 标题（使用 majorFont）
    - "subTitle", "body" → 正文（使用 minorFont）
    - "dt", "ftr", "sldNum" → 辅助信息
    - None → 不是 placeholder 或类型未指定（视为 body/other）
    """
    meta = _extract_placeholder_meta(shape)
    if not meta:
        return None
    ph_type = meta.get("type")
    return str(ph_type) if ph_type else None


def _map_placeholder_type_to_text_type(placeholder_type: Optional[str]) -> Optional[str]:
    """placeholder type → 前端 textType 语义映射。"""
    if not placeholder_type:
        return None
    if placeholder_type in ("title", "ctrTitle"):
        return "title"
    if placeholder_type == "subTitle":
        return "subtitle"
    if placeholder_type in ("body", "obj"):
        return "content"
    return None


def _resolve_theme_font(
    theme_fonts: Dict[str, str],
    placeholder_type: Optional[str] = None,
) -> Optional[str]:
    """
    根据 placeholder 类型选择合适的主题字体。

    标题类 (title, ctrTitle) → majorFont
    其他 (body, subTitle, None, etc.) → minorFont

    优先返回 East Asian 字体（适配中文 PPT），如果没有则返回 Latin 字体。
    """
    if not theme_fonts:
        return None

    is_title = placeholder_type in ("title", "ctrTitle")
    prefix = "major" if is_title else "minor"

    # 优先 EA（中文环境更常用）→ Latin → CS
    for suffix in ("ea", "latin", "cs"):
        font = theme_fonts.get(f"{prefix}_{suffix}")
        if font:
            return font

    return None


def _resolve_master_style(
    master_styles: Dict[str, Dict[str, Any]],
    placeholder_type: Optional[str] = None,
) -> Dict[str, Any]:
    """根据 placeholder 类型返回对应的 master 默认样式"""
    if not master_styles:
        return {}

    if placeholder_type in ("title", "ctrTitle"):
        return master_styles.get("title", {})
    elif placeholder_type in ("body", "subTitle", "obj"):
        return master_styles.get("body", {})
    else:
        return master_styles.get("other", {})


# ============================================================================
# PPTX 文件格式校验
# ============================================================================


class InvalidPptxError(ValueError):
    """上传的文件不是有效的 PPTX 格式。"""
    pass


MAX_DECOMPRESSED_SIZE = 200 * 1024 * 1024  # 200 MB


def validate_pptx_file(pptx_path: str) -> None:
    """
    校验文件是否为合法 PPTX：ZIP magic bytes + 包含 [Content_Types].xml。
    同时检测 ZIP bomb（解压后总大小超过 200MB 则拒绝）。
    不合法时抛出 InvalidPptxError。
    """
    import zipfile

    if not zipfile.is_zipfile(pptx_path):
        raise InvalidPptxError("文件不是有效的 ZIP/PPTX 格式")

    try:
        with zipfile.ZipFile(pptx_path, "r") as zf:
            if "[Content_Types].xml" not in zf.namelist():
                raise InvalidPptxError("ZIP 中缺少 [Content_Types].xml，不是有效的 PPTX 文件")
            total_uncompressed = sum(info.file_size for info in zf.infolist())
            if total_uncompressed > MAX_DECOMPRESSED_SIZE:
                raise InvalidPptxError(
                    f"PPTX 解压后大小 ({total_uncompressed / 1024 / 1024:.1f} MB) "
                    f"超过限制 ({MAX_DECOMPRESSED_SIZE / 1024 / 1024:.0f} MB)，疑似 ZIP bomb"
                )
    except zipfile.BadZipFile:
        raise InvalidPptxError("文件 ZIP 结构损坏，无法解析")


# ============================================================================
# READ: PPTX → SlideElement[]
# ============================================================================


def _read_pages_from_prs(
    prs,
    pptx_path: str,
    canvas_width: int = DEFAULT_SLIDE_WIDTH_PX,
    canvas_height: int = DEFAULT_SLIDE_HEIGHT_PX,
    image_handler: Optional[Any] = None,
) -> List[Dict[str, Any]]:
    """内部函数：从已加载的 Presentation 对象提取所有页面。"""
    slide_width_emu = prs.slide_width or DEFAULT_SLIDE_WIDTH_EMU
    slide_height_emu = prs.slide_height or DEFAULT_SLIDE_HEIGHT_EMU

    theme_color_map = _extract_theme_color_map(prs)
    theme_fonts = _extract_theme_fonts(prs)
    master_styles = _extract_master_text_styles(prs)

    if theme_fonts:
        logger.info(f"Theme fonts: {theme_fonts}")
    if master_styles:
        logger.info(f"Master styles: {list(master_styles.keys())}")

    master_to_idx: Dict[int, int] = {}
    for idx, master in enumerate(prs.slide_masters):
        master_to_idx[id(master)] = idx
    per_master_theme_cache: Dict[int, Dict[str, str]] = {0: theme_color_map}

    pages = []
    for slide_idx, slide in enumerate(prs.slides):
        try:
            slide_master = slide.slide_layout.slide_master
            master_idx = master_to_idx.get(id(slide_master), 0)
        except Exception:
            master_idx = 0
        if master_idx not in per_master_theme_cache:
            try:
                per_master_theme_cache[master_idx] = _extract_theme_color_map_for_master(prs.slide_masters[master_idx])
            except Exception:
                per_master_theme_cache[master_idx] = theme_color_map
        slide_theme_color_map = per_master_theme_cache.get(master_idx, theme_color_map)

        page = _extract_slide(
            slide,
            slide_idx=slide_idx,
            slide_width_emu=slide_width_emu,
            slide_height_emu=slide_height_emu,
            canvas_width=canvas_width,
            canvas_height=canvas_height,
            theme_color_map=slide_theme_color_map,
            theme_fonts=theme_fonts,
            master_styles=master_styles,
            image_handler=image_handler,
        )
        pages.append(page)

    logger.info(f"pptx_io.read: {pptx_path} → {len(pages)} pages")
    return pages


def read(
    pptx_path: str,
    canvas_width: int = DEFAULT_SLIDE_WIDTH_PX,
    canvas_height: int = DEFAULT_SLIDE_HEIGHT_PX,
    image_handler: Optional[Any] = None,
) -> List[Dict[str, Any]]:
    """
    读取 PPTX 文件，返回 SlidePage[] JSON。

    Args:
        pptx_path: PPTX 文件路径
        canvas_width: 编辑器画布宽度（px），用于 EMU→px 换算
        canvas_height: 编辑器画布高度（px）
        image_handler: 可选的图片处理回调 (bytes, content_type, idx) → src_url。
                       传入时图片上传到外部存储返回 URL，否则使用 base64 内联。

    Returns:
        List[SlidePage]，每个 SlidePage 包含:
        {
            "id": "page-uuid",
            "elements": [SlideElement, ...],
            "background": { "type": "color", "value": "#ffffff" },
            "notes": "演讲备注"
        }
    """
    from pptx import Presentation

    prs = Presentation(pptx_path)
    return _read_pages_from_prs(prs, pptx_path, canvas_width, canvas_height, image_handler)


def read_all(
    pptx_path: str,
    canvas_width: int = DEFAULT_SLIDE_WIDTH_PX,
    canvas_height: int = DEFAULT_SLIDE_HEIGHT_PX,
    image_handler: Optional[Any] = None,
) -> Dict[str, Any]:
    """
    单次 Presentation() 初始化，同时提取 pages + theme_fonts + theme_payload。
    合并了 read() / extract_theme_fonts() / extract_theme_payload() 三次独立解析。

    extract_embedded_fonts() 直接使用 zipfile 提取，不需要 Presentation 对象，
    仍需单独调用。

    Returns:
        {"pages": [...], "theme_fonts": {...}, "theme_payload": {...}}
    """
    from pptx import Presentation

    prs = Presentation(pptx_path)
    pages = _read_pages_from_prs(prs, pptx_path, canvas_width, canvas_height, image_handler)
    theme_color_map = _extract_theme_color_map(prs)
    theme_fonts = _extract_theme_fonts(prs)
    theme_payload = _build_theme_payload(theme_color_map, theme_fonts)

    return {
        "pages": pages,
        "theme_fonts": theme_fonts,
        "theme_payload": theme_payload,
    }


def extract_theme_fonts(pptx_path: str) -> Dict[str, str]:
    """
    公共接口：从 PPTX 提取主题字体信息。

    返回如 {"major_latin": "Calibri Light", "major_ea": "等线 Light",
            "minor_latin": "Calibri", "minor_ea": "等线"} 的字典。
    """
    from pptx import Presentation
    prs = Presentation(pptx_path)
    return _extract_theme_fonts(prs)


def extract_embedded_fonts(pptx_path: str) -> List[Dict[str, Any]]:
    """
    从 PPTX 文件中提取嵌入字体。

    PPTX (OPC ZIP) 结构：
    - ppt/presentation.xml 包含 <p:embeddedFontLst>，声明嵌入字体
    - ppt/_rels/presentation.xml.rels 映射 rId → fonts/*.fntdata
    - ppt/fonts/*.fntdata 是字体文件（可能经 OOXML 混淆处理）

    Returns:
        List[{
            "name": "Calibri",          # 字体名称
            "style": "normal",          # normal | bold | italic | bolditalic
            "format": "truetype",       # truetype | opentype
            "data_base64": "AAEAA..."   # base64 编码的字体二进制
        }]
    """
    import base64
    import zipfile

    from lxml import etree

    fonts: List[Dict[str, Any]] = []
    seen_variants: set[str] = set()

    try:
        with zipfile.ZipFile(pptx_path, "r") as zf:
            namelist = zf.namelist()

            # 1. 读取 presentation.xml
            pres_path = "ppt/presentation.xml"
            if pres_path not in namelist:
                return fonts

            pres_xml = etree.parse(zf.open(pres_path))
            nsmap = {
                "p": "http://schemas.openxmlformats.org/presentationml/2006/main",
                "a": "http://schemas.openxmlformats.org/drawingml/2006/main",
                "r": "http://schemas.openxmlformats.org/officeDocument/2006/relationships",
            }

            # 2. 查找嵌入字体声明
            emb_font_list = pres_xml.findall(".//p:embeddedFontLst/p:embeddedFont", nsmap)
            if not emb_font_list:
                logger.debug("No embedded fonts found in PPTX")
                return fonts

            # 3. 读取 relationships（rId → Target 映射）
            rels_path = "ppt/_rels/presentation.xml.rels"
            rels_map: Dict[str, str] = {}
            if rels_path in namelist:
                rels_xml = etree.parse(zf.open(rels_path))
                for rel in rels_xml.iter():
                    if rel.tag.endswith("Relationship"):
                        rid = rel.get("Id", "")
                        target = rel.get("Target", "")
                        if rid and target:
                            rels_map[rid] = target

            # 4. 提取每个嵌入字体
            for emb_font in emb_font_list:
                font_elem = emb_font.find("p:font", nsmap)
                if font_elem is None:
                    continue
                font_name = font_elem.get("typeface", "")
                if not font_name:
                    continue

                # 检查每种样式变体
                variants = [
                    ("p:regular", "normal"),
                    ("p:bold", "bold"),
                    ("p:italic", "italic"),
                    ("p:boldItalic", "bolditalic"),
                ]
                for variant_tag, style_name in variants:
                    var_elem = emb_font.find(variant_tag, nsmap)
                    if var_elem is None:
                        continue

                    r_embed = var_elem.get(
                        "{http://schemas.openxmlformats.org/officeDocument/2006/relationships}embed"
                    )
                    if not r_embed or r_embed not in rels_map:
                        continue

                    font_target = rels_map[r_embed]
                    # Target 可以是相对路径如 "fonts/font1.fntdata"
                    font_part_path = f"ppt/{font_target}" if not font_target.startswith("ppt/") else font_target
                    if font_part_path not in namelist:
                        # 也尝试不带 ppt/ 前缀
                        if font_target in namelist:
                            font_part_path = font_target
                        else:
                            continue

                    try:
                        info = zf.getinfo(font_part_path)
                    except KeyError:
                        continue
                    if info.file_size > MAX_SINGLE_FONT_RAW_SIZE:
                        logger.warning(
                            "Embedded font skipped before read: %s (%s) too large (%d bytes)",
                            font_name, font_part_path, info.file_size,
                        )
                        continue

                    raw_data = zf.read(font_part_path)
                    if not raw_data or len(raw_data) < 32:
                        continue

                    # 5. 反混淆处理
                    font_data = _deobfuscate_font_data(raw_data, font_target)

                    # 6. 检测字体格式
                    fmt = _detect_font_format(font_data)
                    if not fmt:
                        logger.debug(f"Cannot determine font format for {font_name} ({font_target})")
                        continue

                    dedupe_key = f"{font_name.strip().lower()}::{style_name}"
                    if dedupe_key in seen_variants:
                        continue
                    seen_variants.add(dedupe_key)

                    fonts.append({
                        "name": font_name,
                        "style": style_name,
                        "format": fmt,
                        "data_base64": base64.b64encode(font_data).decode("ascii"),
                    })
                    logger.info(f"Extracted embedded font: {font_name} ({style_name})")

    except Exception as e:
        logger.warning(f"Failed to extract embedded fonts: {e}", exc_info=True)

    return fonts


def _deobfuscate_font_data(raw_data: bytes, part_target: str) -> bytes:
    """
    OOXML 字体反混淆（ECMA-376 Part 2, §13.2.1）。

    算法：
    1. 从字体 part 路径中提取 GUID
    2. GUID → 16 字节 key，重复为 32 字节
    3. XOR 前 32 字节

    如果字体已经是有效格式（未混淆），直接返回原数据。
    """
    import re

    # 先检查是否已经是有效字体（未混淆）
    if _detect_font_format(raw_data):
        return raw_data

    # 尝试从 part 路径中提取 GUID
    # 模式: {XXXXXXXX-XXXX-XXXX-XXXX-XXXXXXXXXXXX}.fntdata
    guid_match = re.search(
        r"\{?([0-9a-fA-F]{8}-[0-9a-fA-F]{4}-[0-9a-fA-F]{4}-[0-9a-fA-F]{4}-[0-9a-fA-F]{12})\}?",
        part_target,
    )

    if not guid_match:
        # 无 GUID → 尝试暴力检测：部分生成器在文件头存储了 GUID
        # 或者尝试所有可能的方式
        return raw_data

    guid_str = guid_match.group(1)

    try:
        key = _guid_to_key(guid_str)
        key32 = key * 2  # 16 bytes → 32 bytes

        result = bytearray(raw_data)
        for i in range(min(32, len(result))):
            result[i] ^= key32[i]

        deobfuscated = bytes(result)
        if _detect_font_format(deobfuscated):
            return deobfuscated
    except Exception:
        pass

    return raw_data


def _guid_to_key(guid_str: str) -> bytes:
    """
    GUID 字符串转 16 字节反混淆 key。

    OOXML 规定 GUID 字节序：前三段 little-endian，后两段 big-endian。
    例如: "B77A0B10-B005-45EF-843B-D81C4C2B0EC6"
    → bytes(10, 0B, 7A, B7, 05, B0, EF, 45, 84, 3B, D8, 1C, 4C, 2B, 0E, C6)
    """
    import struct

    parts = guid_str.split("-")
    if len(parts) != 5:
        raise ValueError(f"Invalid GUID: {guid_str}")

    # Part 1-3: little-endian
    p1 = struct.pack("<I", int(parts[0], 16))
    p2 = struct.pack("<H", int(parts[1], 16))
    p3 = struct.pack("<H", int(parts[2], 16))
    # Part 4-5: big-endian (raw hex bytes)
    p4 = bytes.fromhex(parts[3])
    p5 = bytes.fromhex(parts[4])

    return p1 + p2 + p3 + p4 + p5


def _detect_font_format(data: bytes) -> Optional[str]:
    """检测字体二进制的格式，返回 CSS @font-face format 名称"""
    if len(data) < 4:
        return None

    # TrueType: 0x00010000 或 'true'
    if data[:4] == b"\x00\x01\x00\x00" or data[:4] == b"true":
        return "truetype"
    # OpenType (CFF): 'OTTO'
    if data[:4] == b"OTTO":
        return "opentype"
    # TrueType Collection: 'ttcf'
    if data[:4] == b"ttcf":
        return "truetype"
    # WOFF
    if data[:4] == b"wOFF":
        return "woff"
    # WOFF2
    if data[:4] == b"wOF2":
        return "woff2"

    return None


def _extract_slide_layout_meta(slide) -> Optional[Dict[str, Any]]:
    """
    提取 slide 对应的 layout/master 元数据，供前端保存后导出时还原。
    """
    try:
        layout = getattr(slide, "slide_layout", None)
        if layout is None:
            return None

        meta: Dict[str, Any] = {}
        layout_name = getattr(layout, "name", None)
        if layout_name:
            meta["name"] = str(layout_name)
        try:
            part = getattr(layout, "part", None)
            partname = getattr(part, "partname", None) if part is not None else None
            if partname:
                meta["partName"] = str(partname)
        except Exception:
            pass

        layout_idx = None
        try:
            prs = slide.part.package.presentation_part.presentation  # type: ignore[attr-defined]
            for idx, candidate in enumerate(prs.slide_layouts):
                candidate_part = getattr(candidate, "part", None)
                candidate_partname = getattr(candidate_part, "partname", None) if candidate_part is not None else None
                if candidate is layout or (
                    candidate_partname
                    and meta.get("partName")
                    and str(candidate_partname) == str(meta.get("partName"))
                ):
                    layout_idx = idx
                    break
        except Exception:
            pass
        if layout_idx is not None:
            meta["index"] = int(layout_idx)

        master = getattr(layout, "slide_master", None)
        if master is not None:
            master_name = getattr(master, "name", None)
            if master_name:
                meta["masterName"] = str(master_name)
            try:
                master_part = getattr(master, "part", None)
                master_partname = (
                    getattr(master_part, "partname", None) if master_part is not None else None
                )
                if master_partname:
                    meta["masterPartName"] = str(master_partname)
            except Exception:
                pass

        return meta or None
    except Exception:
        return None


def _iter_layout_master_shapes(slide) -> List[Tuple[str, Any]]:
    """
    枚举当前 slide 可继承的版式元素（master -> layout）。

    返回值为 [(source, shape), ...]，source in {"master", "layout"}。
    """
    collected: List[Tuple[str, Any]] = []
    try:
        layout = getattr(slide, "slide_layout", None)
    except Exception:
        layout = None
    if layout is None:
        return collected

    try:
        master = getattr(layout, "slide_master", None)
        if master is not None:
            for shape in getattr(master, "shapes", []):
                collected.append(("master", shape))
    except Exception:
        pass

    try:
        for shape in getattr(layout, "shapes", []):
            collected.append(("layout", shape))
    except Exception:
        pass

    return collected


def _extract_master_layer_elements(
    slide,
    slide_width_emu: int,
    slide_height_emu: int,
    canvas_width: int,
    canvas_height: int,
    theme_color_map: Optional[Dict[str, str]] = None,
    theme_fonts: Optional[Dict[str, str]] = None,
    master_styles: Optional[Dict[str, Dict[str, Any]]] = None,
    image_handler: Optional[Any] = None,
) -> List[Dict[str, Any]]:
    """
    提取 layout/master 的静态元素，作为前端只读层渲染数据。

    规则：
    - 过滤 placeholder，避免与可编辑占位符文本重复
    - 导出顺序按 master -> layout，layout 覆盖 master 的视觉层
    - 元素统一 locked=true，防止误进入编辑链路
    """
    collected: List[Dict[str, Any]] = []
    z_index = 0
    for source, shape in _iter_layout_master_shapes(slide):
        try:
            if _extract_placeholder_meta(shape):
                continue
        except Exception:
            pass

        extracted = _extract_shape(
            shape,
            z_index=z_index,
            slide_width_emu=slide_width_emu,
            slide_height_emu=slide_height_emu,
            canvas_width=canvas_width,
            canvas_height=canvas_height,
            theme_color_map=theme_color_map,
            theme_fonts=theme_fonts,
            master_styles=master_styles,
            image_handler=image_handler,
        )
        z_index += 1
        if not extracted:
            continue

        extracted_items = extracted if isinstance(extracted, list) else [extracted]
        for item in extracted_items:
            if not isinstance(item, dict):
                continue
            item["locked"] = True
            props = item.get("props")
            if isinstance(props, dict):
                props = dict(props)
            else:
                props = {}
            props["masterSource"] = source
            item["props"] = props
            collected.append(item)

    if not collected:
        return []

    # 与普通元素一致，先合并颜色蒙版，再重排 zIndex，避免显示断层。
    collected = _merge_image_color_mask_overlays(collected)
    for idx, el in enumerate(collected):
        try:
            el["zIndex"] = idx
        except Exception:
            pass
    return collected


def _extract_slide(
    slide,
    slide_idx: int,
    slide_width_emu: int,
    slide_height_emu: int,
    canvas_width: int,
    canvas_height: int,
    theme_color_map: Optional[Dict[str, str]] = None,
    theme_fonts: Optional[Dict[str, str]] = None,
    master_styles: Optional[Dict[str, Dict[str, Any]]] = None,
    image_handler: Optional[Any] = None,
) -> Dict[str, Any]:
    """提取单张幻灯片的所有元素"""

    layout_meta = _extract_slide_layout_meta(slide)
    master_elements = _extract_master_layer_elements(
        slide,
        slide_width_emu=slide_width_emu,
        slide_height_emu=slide_height_emu,
        canvas_width=canvas_width,
        canvas_height=canvas_height,
        theme_color_map=theme_color_map,
        theme_fonts=theme_fonts,
        master_styles=master_styles,
        image_handler=image_handler,
    )

    elements = []
    for z_index, shape in enumerate(slide.shapes):
        extracted = _extract_shape(
            shape,
            z_index=z_index,
            slide_width_emu=slide_width_emu,
            slide_height_emu=slide_height_emu,
            canvas_width=canvas_width,
            canvas_height=canvas_height,
            theme_color_map=theme_color_map,
            theme_fonts=theme_fonts,
            master_styles=master_styles,
            image_handler=image_handler,
        )
        if extracted:
            if isinstance(extracted, list):
                elements.extend(extracted)
            else:
                elements.append(extracted)

    # 颜色蒙版写回采用“图片 + 顶层半透明矩形”方案，这里在读取时合并回 image.props.colorMask，
    # 避免前端出现一张图片+一个独立矩形的体验断层。
    elements = _merge_image_color_mask_overlays(elements)

    # 统一重建 zIndex，确保分组展开后仍严格保持视觉层级顺序。
    # 注意：group 子元素在提取时可能会使用临时 zIndex（如 z*100+idx），
    # 这里按最终扁平顺序重排为 0..N-1，避免后续写回时层级错乱。
    for idx, el in enumerate(elements):
        try:
            el["zIndex"] = idx
        except Exception:
            pass

    # 背景（支持从 layout / master 继承）
    # 优先从 slide 自身提取，如果 slide 没有显式背景再查 layout → master
    slide_bg = _extract_background_from_element(slide, theme_color_map, image_handler=image_handler)
    if slide_bg:
        background = slide_bg
    else:
        layout_bg = _extract_layout_or_master_background(slide, theme_color_map, image_handler=image_handler)
        background = layout_bg or {"type": "color", "value": "#ffffff"}
        # 标记背景来自 layout/master 继承，写回时跳过（让模板自身定义）
        background["inherited"] = True

    # 备注
    notes = ""
    if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
        raw_notes = slide.notes_slide.notes_text_frame.text or ""
        notes = raw_notes.replace("\r\n", "\n")

    result = {
        "id": f"page-{uuid.uuid4().hex[:12]}",
        "elements": elements,
        "background": background,
        "notes": notes,
    }
    if layout_meta:
        result["layout"] = layout_meta
    if master_elements:
        result["masterElements"] = master_elements
    return result


def _is_close_num(a: Any, b: Any, tol: float = 1.0) -> bool:
    try:
        return abs(float(a) - float(b)) <= tol
    except Exception:
        return False


def _to_rgba_color_with_alpha(color: str, fallback_opacity: Optional[float] = None) -> Optional[str]:
    if not color or not isinstance(color, str):
        return None

    hex_color, alpha = _parse_css_color(color)
    if alpha is None and fallback_opacity is not None:
        try:
            op = float(fallback_opacity)
            if 0 < op < 1:
                alpha = op
        except Exception:
            pass

    if alpha is None or alpha <= 0:
        return None
    alpha = max(0.0, min(1.0, float(alpha)))
    return _hex_to_rgba(hex_color, round(alpha, 3))


def _extract_color_mask_from_overlay(image_el: Dict[str, Any], shape_el: Dict[str, Any]) -> Optional[str]:
    """检测 shape 是否是 image 上方的颜色蒙版矩形，命中则返回 rgba 颜色字符串。"""
    if image_el.get("type") != "image" or shape_el.get("type") != "shape":
        return None

    # 需与图片几何几乎完全一致，且紧邻上一层（zIndex + 1）
    for key in ("x", "y", "width", "height"):
        if not _is_close_num(image_el.get(key, 0), shape_el.get(key, 0), tol=1.0):
            return None
    if not _is_close_num(image_el.get("rotate", 0), shape_el.get("rotate", 0), tol=0.6):
        return None

    img_z = image_el.get("zIndex")
    shp_z = shape_el.get("zIndex")
    try:
        # 兼容两类 zIndex：
        # 1) 常规读写后的整数层级（如 3 -> 4）
        # 2) 组合展开阶段的小数偏移层级（如 0.0001 -> 0.0002）
        # 这里仅要求“蒙版在图片上方”，其余由“相邻元素 + 几何一致”约束兜底。
        if img_z is not None and shp_z is not None and float(shp_z) <= float(img_z):
            return None
    except Exception:
        pass

    props = shape_el.get("props") or {}
    if not isinstance(props, dict):
        return None

    shape_type = props.get("pptxShapeType")
    if shape_type and shape_type not in (
        "rect", "ellipse", "roundRect", "round1Rect", "round2SameRect", "round2DiagRect", "snipRndRect",
    ):
        return None

    # 蒙版形状应“只有填充、无文本、无边框/渐变/图案/阴影”
    if props.get("text") or props.get("outline") or props.get("gradient") or props.get("pattern"):
        return None
    if shape_el.get("shadow"):
        return None

    fill = props.get("fill")
    if not isinstance(fill, str) or fill in ("none", "transparent", ""):
        return None

    return _to_rgba_color_with_alpha(fill, shape_el.get("opacity"))


def _merge_image_color_mask_overlays(elements: List[Dict[str, Any]]) -> List[Dict[str, Any]]:
    if len(elements) < 2:
        return elements

    merged: List[Dict[str, Any]] = []
    idx = 0
    while idx < len(elements):
        current = elements[idx]

        if current.get("type") == "image" and idx + 1 < len(elements):
            nxt = elements[idx + 1]
            color_mask = _extract_color_mask_from_overlay(current, nxt)
            if color_mask:
                img_props = current.setdefault("props", {})
                if isinstance(img_props, dict):
                    img_props["colorMask"] = color_mask

                    # 写回时蒙版矩形会继承图片圆角，这里也回填 radius，保障前端一致。
                    nxt_props = nxt.get("props") or {}
                    if isinstance(nxt_props, dict) and "radius" in nxt_props and "radius" not in img_props:
                        try:
                            r = float(nxt_props["radius"])
                            if r > 0:
                                img_props["radius"] = round(r, 1)
                        except Exception:
                            pass

                merged.append(current)
                idx += 2
                continue

        merged.append(current)
        idx += 1

    return merged


def _normalize_element_link_payload(raw_link: Any) -> Optional[Dict[str, str]]:
    """规范化元素 link 字段（支持 web/slide）。"""
    if not isinstance(raw_link, dict):
        return None

    raw_target = raw_link.get("target")
    if raw_target is None:
        return None
    target = str(raw_target).strip()
    if not target:
        return None

    raw_type = str(raw_link.get("type", "")).strip().lower()
    if raw_type == "slide":
        page_id = _link_target_to_page_id(target)
        if not page_id:
            return None
        return {"type": "slide", "target": page_id}

    if raw_type == "web":
        web_target = _normalize_web_hyperlink_target(target)
        if not web_target:
            return None
        return {"type": "web", "target": web_target}

    # 兼容历史/脏数据：按 target 形态推断
    page_id = _link_target_to_page_id(target)
    if page_id:
        return {"type": "slide", "target": page_id}

    web_target = _normalize_web_hyperlink_target(target)
    if web_target:
        return {"type": "web", "target": web_target}

    return None


def _link_target_to_page_id(target: str) -> Optional[str]:
    """将 slide 链接目标标准化为 page-N。"""
    s = str(target or "").strip()
    if not s:
        return None

    m_page = re.match(r"^page-(\d+)$", s, flags=re.IGNORECASE)
    if m_page:
        idx = int(m_page.group(1))
        if idx > 0:
            return f"page-{idx}"

    m_slide_xml = re.match(r"^slide(\d+)\.xml$", s, flags=re.IGNORECASE)
    if m_slide_xml:
        idx = int(m_slide_xml.group(1))
        if idx > 0:
            return f"page-{idx}"

    if s.isdigit():
        idx = int(s)
        if idx > 0:
            return f"page-{idx}"

    return None


_SAFE_WEB_LINK_SCHEMES = {"http", "https", "mailto", "tel"}


def _normalize_web_hyperlink_target(raw_target: Any) -> Optional[str]:
    """规范化外部超链接目标，仅允许安全协议。"""
    target = str(raw_target or "").strip()
    if not target:
        return None

    if target.startswith("//"):
        return f"https:{target}"

    m_scheme = re.match(r"^([a-zA-Z][a-zA-Z\d+.-]*):", target)
    if m_scheme:
        scheme = m_scheme.group(1).lower()
        if scheme not in _SAFE_WEB_LINK_SCHEMES:
            return None
        return target

    if re.search(r"\s", target):
        return None

    return f"https://{target}"


def _normalize_rich_text_hyperlink_payload(raw_target: Any) -> Optional[Dict[str, str]]:
    """
    规范化富文本 run 超链接（支持 web/slide）。

    约定：
    - 页内跳转：#page-N / page-N / slideN.xml / N
    - 外部链接：同 _normalize_web_hyperlink_target
    """
    target = str(raw_target or "").strip()
    if not target:
        return None

    slide_candidate = target[1:].strip() if target.startswith("#") else target
    page_id = _link_target_to_page_id(slide_candidate)
    if page_id:
        return {"type": "slide", "target": page_id}

    web_target = _normalize_web_hyperlink_target(target)
    if web_target:
        return {"type": "web", "target": web_target}

    return None


def _page_id_to_slide_index(page_id: str) -> Optional[int]:
    m = re.match(r"^page-(\d+)$", str(page_id or "").strip(), flags=re.IGNORECASE)
    if not m:
        return None
    slide_idx = int(m.group(1))
    if slide_idx <= 0:
        return None
    return slide_idx


def _resolve_target_slide_by_part_owner(part_owner, target: str):
    """根据 target（page-N / slideN.xml / N）解析目标 slide。"""
    page_id = _link_target_to_page_id(target)
    if not page_id:
        return None

    slide_idx = _page_id_to_slide_index(page_id)
    if not slide_idx:
        return None

    try:
        prs = part_owner.part.package.presentation_part.presentation
        if slide_idx > len(prs.slides):
            return None
        return prs.slides[slide_idx - 1]
    except Exception:
        return None


def _slide_part_to_page_id(part_owner, target_part) -> Optional[str]:
    """
    将 slide part（relationship target_part）映射为 page-N。
    """
    if target_part is None:
        return None

    target_partname = str(getattr(target_part, "partname", "") or "")
    if not target_partname:
        return None

    try:
        prs = part_owner.part.package.presentation_part.presentation
    except Exception:
        return None

    for idx, slide in enumerate(prs.slides, start=1):
        try:
            slide_part = getattr(slide, "part", None)
            slide_partname = str(getattr(slide_part, "partname", "") or "")
            if slide_partname and slide_partname == target_partname:
                return f"page-{idx}"
        except Exception:
            continue
    return None


def _shape_target_slide_to_page_id(shape, target_slide) -> Optional[str]:
    """将 click_action.target_slide 映射为 page-N。"""
    if target_slide is None:
        return None

    try:
        target_part = getattr(target_slide, "part", None)
        page_id = _slide_part_to_page_id(shape, target_part)
        if page_id:
            return page_id
    except Exception:
        pass

    # 兜底：按对象/slide_id 对比
    try:
        prs = shape.part.package.presentation_part.presentation
        target_slide_id = getattr(target_slide, "slide_id", None)
        for idx, slide in enumerate(prs.slides, start=1):
            if slide is target_slide:
                return f"page-{idx}"
            if target_slide_id is not None and getattr(slide, "slide_id", None) == target_slide_id:
                return f"page-{idx}"
    except Exception:
        return None
    return None


def _extract_shape_hyperlink(shape) -> Optional[Dict[str, str]]:
    """提取元素整体超链接（shape/image/table 等 clickAction）。"""
    try:
        if not hasattr(shape, "click_action"):
            return None

        click_action = shape.click_action
        if click_action is None or not hasattr(click_action, "hyperlink"):
            return None

        hyperlink = click_action.hyperlink
        target_slide = getattr(click_action, "target_slide", None)
        target_page_id = _shape_target_slide_to_page_id(shape, target_slide)
        if target_page_id:
            return {"type": "slide", "target": target_page_id}

        if hyperlink is None:
            return None

        address = getattr(hyperlink, "address", None)
        if not address:
            return None

        target = str(address).strip()
        if not target:
            return None

        # python-pptx 对内部跳转通常返回 slideN.xml
        page_id = _link_target_to_page_id(target)
        if page_id:
            return {"type": "slide", "target": page_id}

        web_target = _normalize_web_hyperlink_target(target)
        if web_target:
            return {"type": "web", "target": web_target}
        return None
    except Exception:
        return None


def _resolve_target_slide(shape, target: str):
    """
    根据 link.target（page-N/slideN.xml/N）解析目标 slide（用于写回内部跳转）。
    """
    return _resolve_target_slide_by_part_owner(shape, target)


def _apply_shape_hyperlink_write(shape, raw_link: Any) -> None:
    """将元素整体超链接写回到 shape.click_action。"""
    try:
        if not hasattr(shape, "click_action"):
            return

        click_action = shape.click_action
        if click_action is None or not hasattr(click_action, "hyperlink"):
            return

        link = _normalize_element_link_payload(raw_link)
        if not link:
            # 仅在显式传入 link 字段时清理旧链接
            if isinstance(raw_link, dict) or raw_link is None:
                try:
                    click_action.target_slide = None
                except Exception:
                    pass
                try:
                    click_action.hyperlink.address = None
                except Exception:
                    pass
            return

        link_type = link.get("type")
        target = link.get("target", "")
        if link_type == "slide":
            target_slide = _resolve_target_slide(shape, target)
            if target_slide is not None:
                try:
                    click_action.hyperlink.address = None
                except Exception:
                    pass
                click_action.target_slide = target_slide
            else:
                # 回退：保底写为外链，避免导出静默丢失
                click_action.hyperlink.address = target
            return

        try:
            click_action.target_slide = None
        except Exception:
            pass
        click_action.hyperlink.address = target
    except Exception:
        pass


def _extract_run_hyperlink(run) -> Optional[Dict[str, str]]:
    """
    提取 run 级超链接（支持外链与页内跳转）。
    """
    nsmap_a = "http://schemas.openxmlformats.org/drawingml/2006/main"
    nsmap_r = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"

    # 1) 优先解析内部页跳转：a:rPr/a:hlinkClick + slide relationship
    try:
        r_pr = run._r.find(f"{{{nsmap_a}}}rPr")
        if r_pr is not None:
            hlink_click = r_pr.find(f"{{{nsmap_a}}}hlinkClick")
            if hlink_click is not None:
                rel_id = hlink_click.get(f"{{{nsmap_r}}}id") or hlink_click.get("id")
                if rel_id:
                    rel = None
                    rels = getattr(run.part, "rels", None)
                    if rels is not None:
                        try:
                            rel = rels[rel_id]
                        except Exception:
                            try:
                                rel = rels.get(rel_id)
                            except Exception:
                                rel = None
                    target_part = getattr(rel, "target_part", None) if rel is not None else None
                    page_id = _slide_part_to_page_id(run, target_part)
                    if page_id:
                        return {"type": "slide", "target": page_id}
    except Exception:
        pass

    # 2) 外链（及部分生成器写成 slideN.xml 的内部跳转）
    try:
        address = getattr(getattr(run, "hyperlink", None), "address", None)
        if not address:
            return None
        target = str(address).strip()
        if not target:
            return None

        page_id = _link_target_to_page_id(target)
        if page_id:
            return {"type": "slide", "target": page_id}

        web_target = _normalize_web_hyperlink_target(target)
        if web_target:
            return {"type": "web", "target": web_target}
    except Exception:
        pass

    return None


def _clear_run_hyperlink(run_obj) -> None:
    nsmap_a = "http://schemas.openxmlformats.org/drawingml/2006/main"
    try:
        run_obj.hyperlink.address = None
    except Exception:
        pass
    try:
        r_pr = _ensure_run_rpr(run_obj)
        for old in list(r_pr.findall(f"{{{nsmap_a}}}hlinkClick")):
            r_pr.remove(old)
    except Exception:
        pass


def _apply_run_hyperlink_write(run_obj, raw_target: Any) -> None:
    """将 run_data.hyperlink 写回到 run（支持 web/slide）。"""
    link = _normalize_rich_text_hyperlink_payload(raw_target)

    # 先清理旧链接，避免 web/slide 切换残留
    _clear_run_hyperlink(run_obj)
    if not link:
        return

    link_type = link.get("type")
    target = link.get("target", "")

    if link_type == "web":
        try:
            run_obj.hyperlink.address = target
        except Exception:
            pass
        return

    if link_type != "slide":
        return

    target_slide = _resolve_target_slide_by_part_owner(run_obj, target)
    if target_slide is None:
        return

    try:
        from lxml import etree
        from pptx.opc.constants import RELATIONSHIP_TYPE as RT

        nsmap_a = "http://schemas.openxmlformats.org/drawingml/2006/main"
        nsmap_r = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"

        rid = run_obj.part.relate_to(target_slide.part, RT.SLIDE)
        r_pr = _ensure_run_rpr(run_obj)
        hlink_click = etree.SubElement(r_pr, f"{{{nsmap_a}}}hlinkClick")
        hlink_click.set(f"{{{nsmap_r}}}id", rid)
        hlink_click.set("action", "ppaction://hlinksldjump")
    except Exception:
        pass


def _extract_shape(
    shape,
    z_index: int,
    slide_width_emu: int,
    slide_height_emu: int,
    canvas_width: int,
    canvas_height: int,
    theme_color_map: Optional[Dict[str, str]] = None,
    theme_fonts: Optional[Dict[str, str]] = None,
    master_styles: Optional[Dict[str, Dict[str, Any]]] = None,
    group_id: Optional[str] = None,
    image_handler: Optional[Any] = None,
) -> Optional[Any]:
    """
    提取单个 shape 为 SlideElement。

    对于 GROUP 类型，递归展开所有子元素，返回 List[Dict]。
    其他类型返回单个 Dict 或 None。
    """
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    media_shape_types = [MSO_SHAPE_TYPE.MEDIA]
    web_video_type = getattr(MSO_SHAPE_TYPE, "WEB_VIDEO", None)
    if web_video_type is not None:
        media_shape_types.append(web_video_type)

    # 基础位置信息（使用 is not None 判断避免 0 值被误处理）
    base = {
        "id": str(uuid.uuid4())[:8],
        "x": emu_to_px(shape.left if shape.left is not None else 0, slide_width_emu, canvas_width),
        "y": emu_to_px(shape.top if shape.top is not None else 0, slide_height_emu, canvas_height),
        "width": emu_to_px(shape.width if shape.width is not None else 0, slide_width_emu, canvas_width),
        "height": emu_to_px(shape.height if shape.height is not None else 0, slide_height_emu, canvas_height),
        "rotate": round(float(shape.rotation) % 360, ROTATE_DECIMALS) if hasattr(shape, "rotation") and shape.rotation else 0,
        "zIndex": z_index,
    }

    if shape.name:
        base["name"] = shape.name

    if group_id:
        base["groupId"] = group_id

    link = _extract_shape_hyperlink(shape)
    if link:
        base["link"] = link

    opacity = _extract_opacity(shape)
    if opacity is not None and opacity < 1.0:
        base["opacity"] = opacity

    shadow = _extract_shadow(shape)
    if shadow:
        base["shadow"] = shadow

    try:
        xfrm = _find_transform_xfrm(shape._element)
        if xfrm is not None:
            if xfrm.get("flipH") == "1":
                base["flipH"] = True
            if xfrm.get("flipV") == "1":
                base["flipV"] = True
    except Exception:
        pass

    # 锁定状态提取（noMove/noResize/noRot → locked: true）
    try:
        nsmap_p_lock = "http://schemas.openxmlformats.org/presentationml/2006/main"
        nsmap_a_lock = "http://schemas.openxmlformats.org/drawingml/2006/main"
        lock_containers = [
            shape._element.find(f".//{{{nsmap_p_lock}}}cNvSpPr"),
            shape._element.find(f".//{{{nsmap_p_lock}}}cNvPicPr"),
            shape._element.find(f".//{{{nsmap_p_lock}}}cNvCxnSpPr"),
            shape._element.find(f".//{{{nsmap_p_lock}}}cNvGrpSpPr"),
        ]
        lock_tags = ("spLocks", "picLocks", "cxnSpLocks", "grpSpLocks")
        for container in lock_containers:
            if container is None:
                continue
            for tag in lock_tags:
                lock_node = container.find(f"{{{nsmap_a_lock}}}{tag}")
                if lock_node is None:
                    continue
                if (
                    lock_node.get("noMove") == "1"
                    or lock_node.get("noResize") == "1"
                    or lock_node.get("noRot") == "1"
                ):
                    base["locked"] = True
                    break
            if base.get("locked"):
                break
    except Exception:
        pass

    # 可见性提取（cNvPr@hidden="1" → visible: false）
    try:
        nsmap_p = "http://schemas.openxmlformats.org/presentationml/2006/main"
        c_nv_pr = shape._element.find(f".//{{{nsmap_p}}}cNvPr")
        if c_nv_pr is not None:
            hidden_attr = str(c_nv_pr.get("hidden", "")).strip().lower()
            if hidden_attr in {"1", "true", "t"}:
                base["visible"] = False
    except Exception:
        pass

    # 检测 placeholder（用于主题字体/样式继承 + 版式占位符语义保留）
    ph_meta = _extract_placeholder_meta(shape)
    ph_type = str(ph_meta.get("type")) if isinstance(ph_meta, dict) and ph_meta.get("type") else None

    if _is_connector_shape(shape):
        return _extract_connector_shape(
            shape, base, slide_width_emu, slide_height_emu,
            canvas_width, canvas_height, theme_color_map,
        )

    if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
        return _extract_group_shape(
            shape,
            z_index=z_index,
            slide_width_emu=slide_width_emu,
            slide_height_emu=slide_height_emu,
            canvas_width=canvas_width,
            canvas_height=canvas_height,
            theme_color_map=theme_color_map,
            theme_fonts=theme_fonts,
            master_styles=master_styles,
            image_handler=image_handler,
        )
    elif shape.shape_type == MSO_SHAPE_TYPE.TABLE:
        return _extract_table_shape(shape, base, theme_color_map, theme_fonts=theme_fonts)
    elif shape.shape_type in tuple(media_shape_types):
        media_el = _extract_media_shape(shape, base, image_handler=image_handler)
        if media_el is not None:
            return media_el
        # 媒体结构解析失败时回退到图片提取，避免元素直接丢失。
        if shape.shape_type == MSO_SHAPE_TYPE.MEDIA and hasattr(shape, "image"):
            logger.debug("Failed to extract media shape as media element, fallback to picture extraction")
            try:
                return _extract_image_shape(shape, base, image_handler=image_handler)
            except Exception:
                pass
        return None
    elif shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
        return _extract_image_shape(shape, base, image_handler=image_handler)
    elif shape.shape_type == MSO_SHAPE_TYPE.CHART:
        return _extract_chart_shape(shape, base, theme_color_map=theme_color_map)
    elif shape.shape_type in (MSO_SHAPE_TYPE.AUTO_SHAPE, MSO_SHAPE_TYPE.FREEFORM):
        # 版式占位符（标题/副标题/正文等）优先按文本元素提取，保留占位符语义链路。
        if ph_type in ("title", "ctrTitle", "subTitle", "body", "obj") and shape.has_text_frame:
            return _extract_text_shape(
                shape,
                base,
                theme_color_map,
                theme_fonts,
                master_styles,
                ph_type,
                ph_meta,
            )
        # 检测是否是伪装成 AutoShape 的文本框（txBox="1"）
        # PPTX XML: <p:cNvSpPr txBox="1"/> 表示这是一个纯文本框
        if _is_text_box(shape):
            return _extract_text_shape(
                shape,
                base,
                theme_color_map,
                theme_fonts,
                master_styles,
                ph_type,
                ph_meta,
            )
        # 启发式检测：透明无边框的基础矩形 + 有文本 → 视为文本元素
        # 很多第三方工具生成的 PPTX 不设 txBox="1"，但实际上是文本容器
        if _is_text_dominant_shape(shape):
            return _extract_text_shape(
                shape,
                base,
                theme_color_map,
                theme_fonts,
                master_styles,
                ph_type,
                ph_meta,
            )
        # 优先识别 custom geometry 的线条（客户端导出的非直线会走该路径）
        custom_line = _extract_custom_geom_line_shape(shape, base, theme_color_map)
        if custom_line is not None:
            return custom_line
        return _extract_auto_shape(
            shape,
            base,
            theme_color_map,
            theme_fonts,
            master_styles,
            ph_type,
            ph_meta,
            image_handler=image_handler,
        )
    elif shape.has_text_frame:
        return _extract_text_shape(
            shape,
            base,
            theme_color_map,
            theme_fonts,
            master_styles,
            ph_type,
            ph_meta,
        )
    else:
        logger.debug(f"Unsupported shape type: {shape.shape_type} ({shape.name})")
        return None


def _is_text_box(shape) -> bool:
    """
    判断 AUTO_SHAPE 是否实际上是文本框。

    PPTX 中"插入文本框"创建的元素在 XML 中的标志：
    <p:nvSpPr><p:cNvSpPr txBox="1"/></p:nvSpPr>

    这类元素虽然 shape_type == AUTO_SHAPE，但应当作为可编辑文本处理。
    """
    try:
        nsmap = {"p": "http://schemas.openxmlformats.org/presentationml/2006/main"}
        nv_sp_pr = shape._element.find("p:nvSpPr", nsmap)
        if nv_sp_pr is not None:
            c_nv_sp_pr = nv_sp_pr.find("p:cNvSpPr", nsmap)
            if c_nv_sp_pr is not None:
                return c_nv_sp_pr.get("txBox") == "1"
    except Exception:
        pass
    return False


def _is_text_dominant_shape(shape) -> bool:
    """
    启发式判断：AUTO_SHAPE 是否本质上是一个文本容器。

    很多 PPTX 文件（尤其第三方工具或模板生成的）使用透明无边框的矩形作为
    文本容器，但不设 txBox="1"。这些形状在视觉上等同于文本框。

    判断条件（必须全部满足）：
    1. 形状有文本内容（非空白）
    2. 形状几何类型是基础矩形（rect / RECTANGLE / 无自定义 preset）
    3. 没有可见的填充（noFill 或无填充节点）
    4. 没有可见的描边（noFill 或宽度为 0）
    """
    from pptx.oxml.ns import qn

    try:
        if not shape.has_text_frame:
            return False
        text = shape.text_frame.text.strip()
        if not text:
            return False

        # 只对基础矩形应用（不影响圆形、三角形、箭头等有视觉意义的形状）
        _RECT_TYPES = {"rect", "rectangle", "round_rectangle", "snip1Rect", "snip2SameRect"}
        try:
            from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
            auto_type = shape.auto_shape_type
            type_name = auto_type.name.lower() if auto_type else ""
        except Exception:
            type_name = ""
        # 也检查 XML 中的 prstGeom
        prst_geom = ""
        try:
            sp_pr = shape._element.find(qn("a:spPr")) or shape._element.find(
                ".//{http://schemas.openxmlformats.org/drawingml/2006/main}spPr"
            )
            if sp_pr is not None:
                prst = sp_pr.find(qn("a:prstGeom"))
                if prst is not None:
                    prst_geom = (prst.get("prst") or "").lower()
        except Exception:
            pass

        is_rect = type_name in _RECT_TYPES or prst_geom in ("rect", "")
        if not is_rect:
            return False

        # 检查填充：noFill 或无填充 → 透明
        has_visible_fill = False
        try:
            sp_pr = shape._element.find(qn("a:spPr")) or shape._element.find(
                ".//{http://schemas.openxmlformats.org/drawingml/2006/main}spPr"
            )
            if sp_pr is not None:
                no_fill = sp_pr.find(qn("a:noFill"))
                solid_fill = sp_pr.find(qn("a:solidFill"))
                grad_fill = sp_pr.find(qn("a:gradFill"))
                patt_fill = sp_pr.find(qn("a:pattFill"))
                blip_fill = sp_pr.find(qn("a:blipFill"))
                if solid_fill is not None or grad_fill is not None or patt_fill is not None or blip_fill is not None:
                    has_visible_fill = True
                # 没有任何填充节点也视为 noFill
        except Exception:
            pass

        # 检查描边：noFill 或 width=0 → 无可见描边
        has_visible_outline = False
        try:
            sp_pr = shape._element.find(qn("a:spPr")) or shape._element.find(
                ".//{http://schemas.openxmlformats.org/drawingml/2006/main}spPr"
            )
            if sp_pr is not None:
                ln = sp_pr.find(qn("a:ln"))
                if ln is not None:
                    ln_no_fill = ln.find(qn("a:noFill"))
                    if ln_no_fill is None:
                        # 有 <a:ln> 但没有 noFill → 检查宽度
                        w = ln.get("w")
                        if w is not None and int(w) > 0:
                            has_visible_outline = True
                        elif w is None:
                            # 没有显式宽度，检查是否有颜色填充
                            ln_solid = ln.find(qn("a:solidFill"))
                            if ln_solid is not None:
                                has_visible_outline = True
        except Exception:
            pass

        return not has_visible_fill and not has_visible_outline
    except Exception:
        return False


def _is_connector_shape(shape) -> bool:
    """判断 shape 是否是连接器/线条（OOXML <p:cxnSp> 元素）"""
    try:
        tag = shape._element.tag
        if tag and "cxnSp" in tag:
            return True
        # 某些 python-pptx 版本中，连接器有 begin_x / end_x 属性
        if hasattr(shape, "begin_x") or hasattr(shape, "end_x"):
            return True
    except Exception:
        pass
    return False


def _map_ooxml_arrow_to_line_point(arrow_type: str) -> str:
    """OOXML 端点类型 → 前端 line point"""
    if not arrow_type or arrow_type == "none":
        return ""
    if arrow_type == "oval":
        return "dot"
    if arrow_type in ("diamond", "stealth", "triangle", "arrow"):
        return arrow_type
    # 兜底：未知端点按箭头处理，避免丢语义
    return "arrow"


def _map_line_point_to_ooxml_arrow(line_point: str) -> str:
    """前端 line point → OOXML 端点类型"""
    if not line_point or line_point == "none":
        return "none"
    if line_point == "dot":
        return "oval"
    if line_point in ("diamond", "stealth", "triangle", "arrow"):
        return line_point
    return "arrow"


def _map_ooxml_dash_to_line_style(prst_val: Optional[str]) -> Optional[str]:
    """OOXML prstDash 值 → 前端 line style。

    保留更多线型以提升 round-trip 保真度：
    solid / dashed / dotted / dashDot / longDash / longDashDot
    """
    dash_map = {
        "solid": "solid",
        "dash": "dashed",
        "dot": "dotted",
        "lgDash": "longDash",
        "dashDot": "dashDot",
        "sysDash": "dashed",
        "sysDot": "dotted",
        "lgDashDot": "longDashDot",
        "lgDashDotDot": "longDashDot",  # 最接近的近似
    }
    if not prst_val:
        return None
    return dash_map.get(prst_val)


def _extract_connector_shape(
    shape, base: dict,
    slide_width_emu: int, slide_height_emu: int,
    canvas_width: int, canvas_height: int,
    theme_color_map: Optional[Dict[str, str]] = None,
) -> Dict[str, Any]:
    """提取连接器/线条 → LineElement"""

    # 提取线条样式
    line_color = "#333333"
    line_width = 2
    line_style = "solid"
    line_alpha = 1.0
    line_theme_key = None
    start_arrow = ""
    end_arrow = ""

    try:
        ln = shape.line if hasattr(shape, "line") else None
        if ln:
            c = color_to_hex(ln.color, theme_color_map) if hasattr(ln, "color") and ln.color else None
            if c:
                line_color = c
            tk = _extract_theme_key_from_color_obj(ln.color) if hasattr(ln, "color") else None
            if tk:
                line_theme_key = tk
            # 检查线条颜色 alpha（单独记录到 line_alpha，避免写入 rgba 造成双重透明）
            try:
                nsmap_a = "http://schemas.openxmlformats.org/drawingml/2006/main"
                ln_el = shape._element.find(f".//{{{nsmap_a}}}ln")
                if ln_el is not None:
                    solid = ln_el.find(f"{{{nsmap_a}}}solidFill")
                    if solid is not None and len(solid) > 0:
                        theme_key = _extract_theme_key_from_solid_fill(solid)
                        if theme_key:
                            line_theme_key = theme_key
                        elif solid.find(f"{{{nsmap_a}}}schemeClr") is not None:
                            line_theme_key = None
                        alpha_float = _extract_solid_fill_alpha(solid, decimals=4)
                        if alpha_float is not None:
                            line_alpha = max(0.0, min(1.0, alpha_float))
            except Exception:
                pass
            if hasattr(ln, "width") and ln.width:
                line_width = max(0.1, round(ln.width / EMU_PER_PT, 2))
            if hasattr(ln, "dash_style") and ln.dash_style:
                dash_val = int(ln.dash_style) if ln.dash_style else 1
                dash_map = {
                    1: "solid",      # SOLID
                    2: "dashed",     # DASH
                    3: "dotted",     # ROUND_DOT
                    4: "dotted",     # SQUARE_DOT
                    5: "dashDot",    # DASH_DOT
                    6: "longDash",   # LONG_DASH
                    7: "longDashDot",  # LONG_DASH_DOT
                    8: "longDashDot",  # LONG_DASH_DOT_DOT (closest approximation)
                }
                line_style = dash_map.get(dash_val, "solid")
    except Exception:
        pass

    # 兜底：直接从 OOXML 读取 prstDash（python-pptx 某些连接器场景下 dash_style 不可靠）
    try:
        nsmap_a_dash = "http://schemas.openxmlformats.org/drawingml/2006/main"
        ln_xml = shape._element.find(f".//{{{nsmap_a_dash}}}ln")
        if ln_xml is not None:
            prst_dash = ln_xml.find(f"{{{nsmap_a_dash}}}prstDash")
            if prst_dash is not None:
                mapped = _map_ooxml_dash_to_line_style(prst_dash.get("val"))
                if mapped:
                    line_style = mapped
    except Exception:
        pass

    # 将线条颜色 alpha 折叠到元素 opacity，保持前端语义一致（color 尽量保持纯色）。
    if line_alpha < 1.0:
        try:
            base_opacity = float(base.get("opacity", 1.0))
        except (TypeError, ValueError):
            base_opacity = 1.0
        combined_opacity = max(0.0, min(1.0, round(base_opacity * line_alpha, 4)))
        if combined_opacity < 1.0:
            base["opacity"] = round(combined_opacity, 4)

    # 提取箭头端点（含尺寸 w/len）
    start_arrow_size: Dict[str, str] = {}
    end_arrow_size: Dict[str, str] = {}
    try:
        from lxml import etree
        nsmap = {"a": "http://schemas.openxmlformats.org/drawingml/2006/main"}
        el = shape._element

        head_end = el.find(".//a:ln/a:headEnd", nsmap)
        if head_end is not None:
            htype = head_end.get("type", "")
            start_arrow = _map_ooxml_arrow_to_line_point(htype)
            hw = head_end.get("w", "")
            hlen = head_end.get("len", "")
            if hw and hw != "med":
                start_arrow_size["w"] = hw
            if hlen and hlen != "med":
                start_arrow_size["len"] = hlen

        tail_end = el.find(".//a:ln/a:tailEnd", nsmap)
        if tail_end is not None:
            ttype = tail_end.get("type", "")
            end_arrow = _map_ooxml_arrow_to_line_point(ttype)
            tw = tail_end.get("w", "")
            tlen = tail_end.get("len", "")
            if tw and tw != "med":
                end_arrow_size["w"] = tw
            if tlen and tlen != "med":
                end_arrow_size["len"] = tlen
    except Exception:
        pass

    # 计算起点终点（相对于绘制区域）
    # 连接器的绘制区域就是 bounding box（x, y, width, height）
    # 线条起点和终点可以通过 flipH / flipV 推断
    x = base.get("x", 0)
    y = base.get("y", 0)
    w = base.get("width", 100)
    h = base.get("height", 0)

    flip_h = False
    flip_v = False
    try:
        xfrm = _find_transform_xfrm(shape._element)
        if xfrm is not None:
            flip_h = xfrm.get("flipH") == "1"
            flip_v = xfrm.get("flipV") == "1"
    except Exception:
        pass

    # 线条起终点
    if flip_h and flip_v:
        start = [w, h]
        end = [0, 0]
    elif flip_h:
        start = [w, 0]
        end = [0, h]
    elif flip_v:
        start = [0, h]
        end = [w, 0]
    else:
        start = [0, 0]
        end = [w, h]

    line_props: Dict[str, Any] = {
        "start": start,
        "end": end,
        "style": line_style,
        "color": line_color,
        "lineWidth": line_width,
        "points": [start_arrow, end_arrow],
    }
    if start_arrow_size or end_arrow_size:
        line_props["pointSizes"] = [start_arrow_size, end_arrow_size]
    if line_theme_key:
        line_props["colorThemeKey"] = line_theme_key

    # 检测连接器类型：折线/曲线
    try:
        nsmap = {
            "a": "http://schemas.openxmlformats.org/drawingml/2006/main",
            "p": "http://schemas.openxmlformats.org/presentationml/2006/main",
        }
        sp_pr = shape._element.find(".//p:spPr", nsmap)
        if sp_pr is None:
            sp_pr = shape._element.find(".//a:spPr", nsmap)
        if sp_pr is not None:
            prst_geom = sp_pr.find("a:prstGeom", nsmap)
            if prst_geom is not None:
                prst = prst_geom.get("prst", "")
                # 读取调整值（连接器 4/5 通常会有多个 gd）
                av_lst = prst_geom.find("a:avLst", nsmap)
                adj_vals: List[float] = []
                if av_lst is not None:
                    for gd in av_lst.findall("a:gd", nsmap):
                        fmla = gd.get("fmla", "")
                        if "val" not in fmla:
                            continue
                        try:
                            adj_vals.append(int(fmla.split()[-1]) / 100000)
                        except (ValueError, IndexError):
                            continue
                if not adj_vals:
                    adj_vals = [0.5]

                # 连接器编号（如 bentConnector3 / curvedConnector5）
                conn_num = 3
                try:
                    import re as _re
                    m = _re.search(r"(\d+)$", prst or "")
                    if m:
                        conn_num = int(m.group(1))
                except Exception:
                    conn_num = 3

                if "bentConnector" in prst:
                    if conn_num >= 4:
                        # 双折线：优先使用前两个调整值，缺失时回退到同一比例
                        ax = adj_vals[0]
                        ay = adj_vals[1] if len(adj_vals) > 1 else adj_vals[0]
                        line_props["broken2"] = [w * ax, h * ay]
                    else:
                        # 单折线
                        ax = adj_vals[0]
                        ay = adj_vals[1] if len(adj_vals) > 1 else adj_vals[0]
                        line_props["broken"] = [w * ax, h * ay]
                elif "curvedConnector" in prst:
                    if conn_num >= 4:
                        # 三次贝塞尔：curvedConnector4 可有多个调整值。
                        # 标准情况下 adj1/adj2 分别控制两个控制点的 X 比率，
                        # 若有 adj3/adj4 则分别控制 Y 比率；否则 Y 回退到与 X 相同比率。
                        a1 = adj_vals[0]
                        a2 = adj_vals[1] if len(adj_vals) > 1 else 1 - a1
                        a3 = adj_vals[2] if len(adj_vals) > 2 else a1
                        a4 = adj_vals[3] if len(adj_vals) > 3 else a2
                        line_props["cubic"] = [[w * a1, h * a3], [w * a2, h * a4]]
                    else:
                        # 二次贝塞尔
                        a1 = adj_vals[0]
                        a2 = adj_vals[1] if len(adj_vals) > 1 else a1
                        line_props["curve"] = [w * a1, h * a2]
    except Exception:
        pass

    # 连接器的 start/end 已根据 xfrm flipH/flipV 还原到前端几何语义，
    # 若继续保留 base.flipH/base.flipV，会在前端渲染时重复翻转。
    base.pop("flipH", None)
    base.pop("flipV", None)

    return {
        **base,
        "type": "line",
        "props": line_props,
    }


def _extract_custom_geom_line_shape(
    shape,
    base: Dict[str, Any],
    theme_color_map: Optional[Dict[str, str]] = None,
) -> Optional[Dict[str, Any]]:
    """
    提取 custom geometry 线条（主要用于客户端 custGeom 导出的折线/贝塞尔）。

    识别规则：
    - 必须存在 <a:custGeom>
    - 路径为开放路径（不含 <a:close>）
    - 仅由 moveTo/lnTo/quadBezTo/cubicBezTo 组成
    """
    try:
        nsmap = {"a": "http://schemas.openxmlformats.org/drawingml/2006/main"}
        nsmap_p = {"p": "http://schemas.openxmlformats.org/presentationml/2006/main"}
        sp_pr = shape._element.find(".//p:spPr", nsmap_p)
        if sp_pr is None:
            sp_pr = shape._element.find(".//a:spPr", nsmap)
        if sp_pr is None:
            return None

        cust_geom = sp_pr.find("a:custGeom", nsmap)
        if cust_geom is None:
            return None

        path_lst = cust_geom.find("a:pathLst", nsmap)
        if path_lst is None:
            return None
        path_el = path_lst.find("a:path", nsmap)
        if path_el is None:
            return None

        path_w = int(path_el.get("w", "0")) or 1
        path_h = int(path_el.get("h", "0")) or 1
        w = float(base.get("width", 0) or 0)
        h = float(base.get("height", 0) or 0)
        sx = w / path_w if path_w else 1.0
        sy = h / path_h if path_h else 1.0

        start_pt: Optional[List[float]] = None
        end_pt: Optional[List[float]] = None
        line_pts: List[List[float]] = []
        quad_ctrl: Optional[List[float]] = None
        cubic_ctrl: Optional[List[List[float]]] = None
        seen_curve = False

        def _read_pt(pt_el):
            if pt_el is None:
                return None
            x = round(int(pt_el.get("x", "0")) * sx, 1)
            y = round(int(pt_el.get("y", "0")) * sy, 1)
            return [x, y]

        for cmd in path_el:
            tag = cmd.tag.split("}")[-1] if "}" in cmd.tag else cmd.tag
            if tag == "moveTo":
                pt = _read_pt(cmd.find("a:pt", nsmap))
                if pt is None:
                    return None
                start_pt = pt
                end_pt = pt
            elif tag == "lnTo":
                if seen_curve:
                    # 不支持混合路径，避免误识别复杂形状
                    return None
                pt = _read_pt(cmd.find("a:pt", nsmap))
                if pt is None:
                    return None
                line_pts.append(pt)
                end_pt = pt
            elif tag == "quadBezTo":
                if line_pts or quad_ctrl is not None or cubic_ctrl is not None:
                    return None
                pts = cmd.findall("a:pt", nsmap)
                if len(pts) < 2:
                    return None
                cp = _read_pt(pts[0])
                ep = _read_pt(pts[1])
                if cp is None or ep is None:
                    return None
                quad_ctrl = cp
                end_pt = ep
                seen_curve = True
            elif tag == "cubicBezTo":
                if line_pts or quad_ctrl is not None or cubic_ctrl is not None:
                    return None
                pts = cmd.findall("a:pt", nsmap)
                if len(pts) < 3:
                    return None
                cp1 = _read_pt(pts[0])
                cp2 = _read_pt(pts[1])
                ep = _read_pt(pts[2])
                if cp1 is None or cp2 is None or ep is None:
                    return None
                cubic_ctrl = [cp1, cp2]
                end_pt = ep
                seen_curve = True
            elif tag == "close":
                return None
            else:
                # arcTo 等复杂命令不作为 line 处理
                return None

        if start_pt is None or end_pt is None:
            return None

        # 线条样式提取（与连接器保持一致）
        line_color = "#333333"
        line_width = 2
        line_style = "solid"
        line_alpha = 1.0
        line_theme_key = None
        start_arrow = ""
        start_arrow_size: Dict[str, str] = {}
        end_arrow_size: Dict[str, str] = {}
        end_arrow = ""
        try:
            ln = shape.line if hasattr(shape, "line") else None
            if ln:
                c = color_to_hex(ln.color, theme_color_map) if hasattr(ln, "color") and ln.color else None
                if c:
                    line_color = c
                tk = _extract_theme_key_from_color_obj(ln.color) if hasattr(ln, "color") else None
                if tk:
                    line_theme_key = tk
                if hasattr(ln, "width") and ln.width:
                    line_width = max(0.1, round(ln.width / EMU_PER_PT, 2))
                if hasattr(ln, "dash_style") and ln.dash_style:
                    dash_val = int(ln.dash_style) if ln.dash_style else 1
                    line_style = {
                        1: "solid",      # SOLID
                        2: "dashed",     # DASH
                        3: "dotted",     # ROUND_DOT
                        4: "dotted",     # SQUARE_DOT
                        5: "dashDot",    # DASH_DOT
                        6: "longDash",   # LONG_DASH
                        7: "longDashDot",  # LONG_DASH_DOT
                        8: "longDashDot",  # LONG_DASH_DOT_DOT
                    }.get(dash_val, "solid")

            ln_el = shape._element.find(".//a:ln", nsmap)
            if ln_el is not None:
                # 线条颜色 alpha
                solid = ln_el.find("a:solidFill", nsmap)
                if solid is not None and len(solid) > 0:
                    theme_key = _extract_theme_key_from_solid_fill(solid)
                    if theme_key:
                        line_theme_key = theme_key
                    elif solid.find("a:schemeClr", nsmap) is not None:
                        line_theme_key = None
                    alpha_float = _extract_solid_fill_alpha(solid, decimals=4)
                    if alpha_float is not None:
                        line_alpha = max(0.0, min(1.0, alpha_float))

                head_end = ln_el.find("a:headEnd", nsmap)
                tail_end = ln_el.find("a:tailEnd", nsmap)
                if head_end is not None:
                    start_arrow = _map_ooxml_arrow_to_line_point(head_end.get("type", ""))
                    hw = head_end.get("w", "")
                    hlen = head_end.get("len", "")
                    if hw and hw != "med":
                        start_arrow_size["w"] = hw
                    if hlen and hlen != "med":
                        start_arrow_size["len"] = hlen
                if tail_end is not None:
                    end_arrow = _map_ooxml_arrow_to_line_point(tail_end.get("type", ""))
                    tw = tail_end.get("w", "")
                    tlen = tail_end.get("len", "")
                    if tw and tw != "med":
                        end_arrow_size["w"] = tw
                    if tlen and tlen != "med":
                        end_arrow_size["len"] = tlen

                # 兜底：直接读取 OOXML prstDash（某些导入源下 python-pptx 的 dash_style 为空）。
                prst_dash = ln_el.find("a:prstDash", nsmap)
                if prst_dash is not None:
                    mapped = _map_ooxml_dash_to_line_style(prst_dash.get("val"))
                    if mapped:
                        line_style = mapped
        except Exception:
            pass

        # 将线条颜色 alpha 折叠到元素 opacity，保持 color 为纯色并与 connector 路径一致。
        if line_alpha < 1.0:
            try:
                base_opacity = float(base.get("opacity", 1.0))
            except (TypeError, ValueError):
                base_opacity = 1.0
            combined_opacity = max(0.0, min(1.0, round(base_opacity * line_alpha, 4)))
            if combined_opacity < 1.0:
                base["opacity"] = round(combined_opacity, 4)

        line_props: Dict[str, Any] = {
            "start": start_pt,
            "end": end_pt,
            "style": line_style,
            "color": line_color,
            "lineWidth": line_width,
            "points": [start_arrow, end_arrow],
        }
        if start_arrow_size or end_arrow_size:
            line_props["pointSizes"] = [start_arrow_size, end_arrow_size]
        if line_theme_key:
            line_props["colorThemeKey"] = line_theme_key

        if cubic_ctrl is not None:
            line_props["cubic"] = cubic_ctrl
        elif quad_ctrl is not None:
            line_props["curve"] = quad_ctrl
        elif len(line_pts) >= 3:
            # custom geometry 导出的双折线：第一折点作为 broken2 控制点
            line_props["broken2"] = line_pts[0]
        elif len(line_pts) == 2:
            line_props["broken"] = line_pts[0]

        return {
            **base,
            "type": "line",
            "props": line_props,
        }
    except Exception:
        return None


def _extract_group_shape(
    shape,
    z_index: int,
    slide_width_emu: int,
    slide_height_emu: int,
    canvas_width: int,
    canvas_height: int,
    theme_color_map: Optional[Dict[str, str]] = None,
    theme_fonts: Optional[Dict[str, str]] = None,
    master_styles: Optional[Dict[str, Dict[str, Any]]] = None,
    image_handler: Optional[Any] = None,
) -> List[Dict[str, Any]]:
    """递归提取组合元素内的所有子元素，标记相同的 groupId。

    关键处理：
    1. 组合自身可能有 rotation / flipH / flipV
    2. python-pptx 已将子元素坐标从组内坐标空间转换为幻灯片坐标空间
       （通过 chOff/chExt → off/ext 映射），但 **不处理旋转/翻转**
    3. 当组合有旋转时，需要将子元素绕组合中心旋转
    4. 当组合有翻转时，需要将子元素在组合中心镜像
    """
    import math

    group_id = str(uuid.uuid4())[:8]
    group_name = ""
    elements = []

    # ── 读取组合自身的变换 ──
    group_rotation = 0.0
    group_flip_h = False
    group_flip_v = False
    group_cx = 0.0  # 组合中心 x (px)
    group_cy = 0.0  # 组合中心 y (px)

    try:
        group_rotation = shape.rotation if hasattr(shape, "rotation") else 0.0
    except Exception as e:
        logger.warning(f"Failed to read group rotation: {e}")

    try:
        raw_group_name = shape.name if hasattr(shape, "name") else ""
        if isinstance(raw_group_name, str):
            group_name = raw_group_name.strip()
    except Exception:
        group_name = ""

    try:
        nsmap_a = "http://schemas.openxmlformats.org/drawingml/2006/main"
        nsmap_p = "http://schemas.openxmlformats.org/presentationml/2006/main"
        # python-pptx GroupShape._element tag 是 p:grpSp，子元素 grpSpPr 也在 p: 命名空间
        grp_sp_pr = shape._element.find(f"{{{nsmap_p}}}grpSpPr")
        if grp_sp_pr is None:
            # fallback: 无命名空间或其他变体
            grp_sp_pr = shape._element.find("grpSpPr")
        if grp_sp_pr is not None:
            xfrm = grp_sp_pr.find(f"{{{nsmap_a}}}xfrm")
            if xfrm is not None:
                if xfrm.get("flipH") == "1":
                    group_flip_h = True
                if xfrm.get("flipV") == "1":
                    group_flip_v = True
    except Exception as e:
        logger.warning(f"Failed to read group flipH/flipV: {e}")

    # 组合中心 (px) — 用于旋转/翻转子元素坐标
    try:
        gx = emu_to_px(shape.left if shape.left is not None else 0, slide_width_emu, canvas_width)
        gy = emu_to_px(shape.top if shape.top is not None else 0, slide_height_emu, canvas_height)
        gw = emu_to_px(shape.width if shape.width is not None else 0, slide_width_emu, canvas_width)
        gh = emu_to_px(shape.height if shape.height is not None else 0, slide_height_emu, canvas_height)
        group_cx = gx + gw / 2
        group_cy = gy + gh / 2
    except Exception as e:
        logger.warning(f"Failed to compute group center: {e}")

    has_group_transform = (
        abs(group_rotation) > 0.01 or group_flip_h or group_flip_v
    )

    try:
        for child_idx, child_shape in enumerate(shape.shapes):
            # 将组内元素 zIndex 压缩在当前组层级附近，避免与后续顶层元素层级穿插。
            # 使用小数偏移可稳定保留组内相对顺序，同时不破坏组与组外元素的前后关系。
            child_z_index = z_index + (child_idx + 1) * 0.0001
            child = _extract_shape(
                child_shape,
                z_index=child_z_index,
                slide_width_emu=slide_width_emu,
                slide_height_emu=slide_height_emu,
                canvas_width=canvas_width,
                canvas_height=canvas_height,
                theme_color_map=theme_color_map,
                theme_fonts=theme_fonts,
                master_styles=master_styles,
                group_id=group_id,
                image_handler=image_handler,
            )
            if not child:
                continue

            child_list = child if isinstance(child, list) else [child]
            for c in child_list:
                c["groupId"] = group_id
                if group_name:
                    c["groupName"] = group_name

                # ── 应用组合的旋转/翻转到子元素 ──
                if has_group_transform:
                    _apply_group_transform_to_child(
                        c,
                        group_cx, group_cy,
                        group_rotation,
                        group_flip_h, group_flip_v,
                    )

            elements.extend(child_list)
    except Exception as e:
        logger.warning(f"Failed to extract group shapes: {e}")

    return elements


def _apply_group_transform_to_child(
    child: Dict[str, Any],
    group_cx: float,
    group_cy: float,
    group_rotation: float,
    group_flip_h: bool,
    group_flip_v: bool,
) -> None:
    """将组合的旋转/翻转应用到子元素的坐标和属性上。

    变换顺序（与 OOXML 规范一致）：
    1. 先翻转（镜像子元素中心相对于组合中心）
    2. 再旋转（绕组合中心旋转子元素中心坐标，并叠加旋转角度）
    """
    import math

    cx = child.get("x", 0) + child.get("width", 0) / 2
    cy = child.get("y", 0) + child.get("height", 0) / 2

    def _mirror_line_points(horizontal: bool = False, vertical: bool = False) -> None:
        """将 line 几何点位直接镜像，避免对连接器引入额外 flip 标记。"""
        if child.get("type") != "line":
            return
        props = child.get("props")
        if not isinstance(props, dict):
            return
        try:
            w = float(child.get("width", 0) or 0)
            h = float(child.get("height", 0) or 0)
        except (TypeError, ValueError):
            w, h = 0.0, 0.0

        def _mirror_pt(raw_pt):
            if not isinstance(raw_pt, (list, tuple)) or len(raw_pt) < 2:
                return raw_pt
            try:
                px = float(raw_pt[0])
                py = float(raw_pt[1])
            except (TypeError, ValueError):
                return raw_pt
            if horizontal:
                px = w - px
            if vertical:
                py = h - py
            return [round(px, COORD_DECIMALS), round(py, COORD_DECIMALS)]

        for key in ("start", "end", "broken", "broken2", "curve"):
            if key in props:
                props[key] = _mirror_pt(props.get(key))

        cubic = props.get("cubic")
        if isinstance(cubic, (list, tuple)) and len(cubic) == 2:
            props["cubic"] = [_mirror_pt(cubic[0]), _mirror_pt(cubic[1])]

    # 1. 翻转：沿组合中心轴镜像
    if group_flip_h:
        cx = 2 * group_cx - cx
        child_rotate = child.get("rotate", 0)
        child["rotate"] = -child_rotate
        if child.get("type") == "line":
            _mirror_line_points(horizontal=True)
            # 线条若本身携带 flip 标记（历史数据），与组翻转做布尔合成。
            if "flipH" in child:
                child["flipH"] = not bool(child.get("flipH"))
        else:
            # 子元素自身翻转取反
            child["flipH"] = not child.get("flipH", False)
    if group_flip_v:
        cy = 2 * group_cy - cy
        child_rotate = child.get("rotate", 0)
        child["rotate"] = -child_rotate
        if child.get("type") == "line":
            _mirror_line_points(vertical=True)
            if "flipV" in child:
                child["flipV"] = not bool(child.get("flipV"))
        else:
            child["flipV"] = not child.get("flipV", False)

    # 2. 旋转：绕组合中心旋转子元素中心坐标
    if abs(group_rotation) > 0.01:
        rad = math.radians(group_rotation)
        cos_a = math.cos(rad)
        sin_a = math.sin(rad)
        dx = cx - group_cx
        dy = cy - group_cy
        cx = group_cx + dx * cos_a - dy * sin_a
        cy = group_cy + dx * sin_a + dy * cos_a
        child["rotate"] = child.get("rotate", 0) + group_rotation

    # 写回坐标（从中心坐标还原为左上角坐标）
    child["x"] = round(cx - child.get("width", 0) / 2, COORD_DECIMALS)
    child["y"] = round(cy - child.get("height", 0) / 2, COORD_DECIMALS)

    # 标准化旋转角度到 [0, 360)
    rot = child.get("rotate", 0)
    if rot:
        child["rotate"] = round(rot % 360, ROTATE_DECIMALS)

    # 清理 False 的翻转标记
    if child.get("flipH") is False:
        child.pop("flipH", None)
    if child.get("flipV") is False:
        child.pop("flipV", None)


def _extract_text_shape(
    shape, base: dict,
    theme_color_map: Optional[Dict[str, str]] = None,
    theme_fonts: Optional[Dict[str, str]] = None,
    master_styles: Optional[Dict[str, Dict[str, Any]]] = None,
    placeholder_type: Optional[str] = None,
    placeholder_meta: Optional[Dict[str, Any]] = None,
) -> Dict[str, Any]:
    """提取文本框 → TextElement（支持主题字体和 Master 样式继承）"""
    tf = shape.text_frame

    # 获取 master 默认样式（根据 placeholder 类型）
    m_style = _resolve_master_style(master_styles or {}, placeholder_type)

    # 1. 扫描默认样式（取第一个非空 run 的样式，作为全局默认值）
    default_font_family = None
    default_font_size = None
    default_color = None
    default_color_theme_key = None
    for para in tf.paragraphs:
        for run in para.runs:
            font = run.font
            if font.name and not default_font_family:
                # 解析可能的主题字体引用
                resolved = _resolve_theme_font_reference(font.name, theme_fonts) if font.name.startswith("+") else font.name
                if resolved:
                    default_font_family = resolved
            if font.size and not default_font_size:
                default_font_size = round(font.size / 12700, 2)  # EMU → pt
            if not default_color:
                c = color_to_hex(font.color, theme_color_map) if font.color else None
                if c:
                    default_color = c
                    default_color_theme_key = _extract_run_theme_color_key(run)
            if default_font_family and default_font_size and default_color:
                break
        if default_font_family and default_font_size and default_color:
            break

    # 主题字体 fallback：如果没有显式字体，使用主题字体
    if not default_font_family and theme_fonts:
        theme_font = _resolve_theme_font(theme_fonts, placeholder_type)
        if theme_font:
            default_font_family = theme_font
            logger.debug(f"Using theme font '{theme_font}' for placeholder type '{placeholder_type}'")

    # Master 样式 fallback：如果没有显式字号/颜色，使用 master 默认值
    if not default_font_size and m_style.get("fontSize"):
        default_font_size = m_style["fontSize"]
    if not default_color and m_style.get("color"):
        default_color = m_style["color"]

    # 2. 分析段落：检测项目符号/编号 + 提取 HTML
    para_infos = []
    line_spacings = []
    for para in tf.paragraphs:
        para_html, p_style = _paragraph_to_html_v2(
            para, theme_color_map,
            default_font_size=default_font_size,
            default_font_family=default_font_family,
            default_color=default_color,
            theme_fonts=theme_fonts,
        )
        bullet_type = _detect_bullet_type(para)
        number_format = _detect_number_format(para) if bullet_type == "number" else None
        bullet_char = _detect_bullet_char(para) if bullet_type == "bullet" else None
        bullet_style = _detect_bullet_style(para, theme_color_map) if bullet_type else None
        level = para.level if hasattr(para, "level") and para.level else 0
        para_infos.append({
            "html": para_html,
            "p_style": p_style,
            "bullet": bullet_type,
            "numberFormat": number_format,
            "bulletChar": bullet_char,
            "bulletStyle": bullet_style,
            "level": level,
        })

        ls = _extract_line_spacing(para)
        if ls is not None:
            line_spacings.append(ls)

    # 3. 生成 HTML，将连续的 bullet/number 段落分组为 <ul>/<ol>
    content = _build_html_with_lists(para_infos) if para_infos else "<p></p>"

    props: Dict[str, Any] = {"content": content}
    if isinstance(placeholder_meta, dict) and placeholder_meta:
        props["placeholder"] = dict(placeholder_meta)
    text_type = _map_placeholder_type_to_text_type(placeholder_type)
    if text_type:
        props["textType"] = text_type
    if default_font_family:
        props["defaultFontFamily"] = _font_with_fallback(default_font_family)
    if default_font_size:
        props["defaultFontSize"] = default_font_size
    if default_color:
        props["defaultColor"] = default_color
    if default_color_theme_key:
        props["defaultColorThemeKey"] = default_color_theme_key

    # 行高（取最常见的行距值）
    if line_spacings:
        avg_ls = round(sum(line_spacings) / len(line_spacings), 2)
        if avg_ls > 0 and avg_ls != 1.0:
            props["lineHeight"] = avg_ls

    # 段间距
    para_space = _extract_paragraph_spacing(tf)
    if para_space:
        props["paragraphSpace"] = para_space

    # 垂直对齐
    vert_align = _extract_text_vertical_align(tf)
    if vert_align:
        props["verticalAlign"] = vert_align

    # 内边距（文本框 margins）
    margins = _extract_text_margins(tf)
    if margins:
        props["margin"] = margins

    # 文本框背景填充色
    try:
        if hasattr(shape, "fill") and shape.fill and shape.fill.type is not None:
            from pptx.enum.dml import MSO_THEME_COLOR
            fill_type = shape.fill.type
            if fill_type == 1:  # SOLID
                fg = shape.fill.fore_color
                if fg and fg.rgb:
                    props["fill"] = f"#{fg.rgb}"
                elif fg and hasattr(fg, "theme_color") and fg.theme_color is not None and theme_color_map:
                    tc_key = str(int(fg.theme_color))
                    if tc_key in theme_color_map:
                        props["fill"] = theme_color_map[tc_key]
    except Exception:
        pass

    # 文本框边框
    outline = _extract_shape_outline(shape, theme_color_map)
    if outline:
        props["outline"] = outline

    # 竖排文字 & 自动适应（bodyPr 属性）
    try:
        nsmap = {"a": "http://schemas.openxmlformats.org/drawingml/2006/main"}
        body_pr = tf._txBody.find("a:bodyPr", nsmap)
        if body_pr is not None:
            vert_attr = body_pr.get("vert")
            if vert_attr and vert_attr in ("vert", "vert270", "eaVert", "wordArtVert"):
                props["vertical"] = True
            wrap_attr = body_pr.get("wrap")
            if wrap_attr == "none":
                props["wordWrap"] = False
            elif wrap_attr == "square":
                props["wordWrap"] = True
            # 自动适应：normAutofit（缩小字体以适应）
            if body_pr.find("a:normAutofit", nsmap) is not None:
                props["autoFit"] = "shrink"
            # 自动适应：spAutoFit（调整形状大小以适应文字）
            elif body_pr.find("a:spAutoFit", nsmap) is not None:
                props["autoFit"] = "resize"
    except Exception:
        pass

    # 字间距（取第一个非零 a:spc 值）
    try:
        nsmap_a = "http://schemas.openxmlformats.org/drawingml/2006/main"
        for para in tf.paragraphs:
            for run in para.runs:
                rPr = run._r.find(f"{{{nsmap_a}}}rPr")
                if rPr is not None:
                    spc = rPr.get("spc")
                    if spc and int(spc) != 0:
                        # OOXML spc 单位是 1/100pt，前端统一使用 px
                        props["wordSpace"] = round(int(spc) / 100 * PX_PER_PT, 2)
                        break
            if "wordSpace" in props:
                break
    except Exception:
        pass

    return {**base, "type": "text", "props": props}


def _paragraph_to_html(
    para,
    theme_color_map: Optional[Dict[str, str]] = None,
    theme_fonts: Optional[Dict[str, str]] = None,
) -> str:
    """段落 → HTML 片段（旧接口，保持向后兼容）"""
    html, _ = _paragraph_to_html_v2(para, theme_color_map, theme_fonts=theme_fonts)
    return html


def _paragraph_to_html_v2(
    para,
    theme_color_map: Optional[Dict[str, str]] = None,
    default_font_size: Optional[float] = None,
    default_font_family: Optional[str] = None,
    default_color: Optional[str] = None,
    theme_fonts: Optional[Dict[str, str]] = None,
) -> Tuple[str, str]:
    """
    段落 → (HTML 内容, p 标签 style 字符串)

    增强版：
    - 提取 defRPr（段落级默认 run 属性）作为中间 fallback
    - 对齐方式（居中/右对齐/两端对齐）→ p style
    - 对无显式字号的 run 使用 defRPr > default_font_size 填充
    - 对无显式字体的 run 使用 defRPr > default_font_family 填充
    """
    # 提取段落级 defRPr（PPTX 样式继承链的关键一环）
    para_def_font_family = None
    para_def_font_size = None
    para_def_color = None
    try:
        nsmap_a = "http://schemas.openxmlformats.org/drawingml/2006/main"
        pPr = para._p.find(f"{{{nsmap_a}}}pPr")
        if pPr is not None:
            defRPr = pPr.find(f"{{{nsmap_a}}}defRPr")
            if defRPr is not None:
                # 字号 (hundredths of pt)
                sz = defRPr.get("sz")
                if sz:
                    para_def_font_size = int(sz) / 100.0

                # 字体（支持主题引用解析 +mj-lt / +mn-lt 等）
                latin = defRPr.find(f"{{{nsmap_a}}}latin")
                if latin is not None and latin.get("typeface"):
                    raw_font = latin.get("typeface")
                    resolved = _resolve_theme_font_reference(raw_font, theme_fonts)
                    if resolved:
                        para_def_font_family = resolved
                if not para_def_font_family:
                    ea = defRPr.find(f"{{{nsmap_a}}}ea")
                    if ea is not None and ea.get("typeface"):
                        raw_font = ea.get("typeface")
                        resolved = _resolve_theme_font_reference(raw_font, theme_fonts)
                        if resolved:
                            para_def_font_family = resolved

                # 颜色（含 lumMod/tint/shade 修饰符计算）
                solid = defRPr.find(f"{{{nsmap_a}}}solidFill")
                if solid is not None:
                    resolved_color = _extract_color_with_transforms(solid, theme_color_map)
                    if resolved_color:
                        para_def_color = resolved_color
    except Exception:
        pass

    # 构建 fallback 链: run 显式 > defRPr > shape default
    eff_font_size = para_def_font_size or default_font_size
    eff_font_family = para_def_font_family or default_font_family
    eff_color = para_def_color or default_color

    parts = []

    from html import escape as html_escape

    for run in para.runs:
        text = run.text
        if not text:
            continue
        font = run.font
        styles = []
        tags_open = []
        tags_close = []

        # 超链接（外链 + 页内跳转）
        hyperlink_payload = _extract_run_hyperlink(run)

        if font.bold:
            tags_open.append("<b>")
            tags_close.insert(0, "</b>")
        if font.italic:
            tags_open.append("<i>")
            tags_close.insert(0, "</i>")
        if font.underline:
            # 保留 OOXML 下划线样式（sng/dbl/wavy/dotted 等 18 种）
            u_style = None
            try:
                nsmap_a_u = "http://schemas.openxmlformats.org/drawingml/2006/main"
                rPr_u = run._r.find(f"{{{nsmap_a_u}}}rPr")
                if rPr_u is not None:
                    u_style = rPr_u.get("u")
            except Exception:
                pass
            if u_style and u_style not in ("none", "sng"):
                tags_open.append(f'<u data-underline-style="{u_style}">')
            else:
                tags_open.append("<u>")
            tags_close.insert(0, "</u>")

        # 删除线（python-pptx 未直接暴露 strike 属性，从 XML 读取）
        try:
            nsmap_a_strike = "http://schemas.openxmlformats.org/drawingml/2006/main"
            rPr = run._r.find(f"{{{nsmap_a_strike}}}rPr")
            if rPr is not None:
                strike_val = rPr.get("strike")
                if strike_val and strike_val != "noStrike":
                    tags_open.append("<s>")
                    tags_close.insert(0, "</s>")
        except Exception:
            pass

        # 高亮背景色（<a:highlight>）
        try:
            nsmap_a_hl = "http://schemas.openxmlformats.org/drawingml/2006/main"
            rPr_hl = run._r.find(f"{{{nsmap_a_hl}}}rPr")
            if rPr_hl is not None:
                highlight_el = rPr_hl.find(f"{{{nsmap_a_hl}}}highlight")
                if highlight_el is not None:
                    hl_color = _extract_color_with_transforms(highlight_el, theme_color_map)
                    if hl_color:
                        tags_open.insert(0, f'<mark data-color="{hl_color}" style="background-color:{hl_color}">')
                        tags_close.append("</mark>")
        except Exception:
            pass

        # 上标 / 下标（baseline 属性：正值=上标，负值=下标）
        try:
            nsmap_a_bl = "http://schemas.openxmlformats.org/drawingml/2006/main"
            rPr_bl = run._r.find(f"{{{nsmap_a_bl}}}rPr")
            if rPr_bl is not None:
                baseline = rPr_bl.get("baseline")
                if baseline:
                    baseline_val = int(baseline)
                    if baseline_val > 0:
                        tags_open.append("<sup>")
                        tags_close.insert(0, "</sup>")
                    elif baseline_val < 0:
                        tags_open.append("<sub>")
                        tags_close.insert(0, "</sub>")
        except Exception:
            pass

        # 颜色：优先从 XML 直接读取（含 lumMod/tint/shade 修饰符），降级到 python-pptx API
        font_hex = None
        run_theme_color_key = None
        try:
            nsmap_a_clr = "http://schemas.openxmlformats.org/drawingml/2006/main"
            rPr_clr = run._r.find(f"{{{nsmap_a_clr}}}rPr")
            if rPr_clr is not None:
                solid_fill = rPr_clr.find(f"{{{nsmap_a_clr}}}solidFill")
                if solid_fill is not None:
                    font_hex = _extract_color_with_transforms(solid_fill, theme_color_map)
                    run_theme_color_key = _extract_theme_key_from_solid_fill(solid_fill)
        except Exception:
            pass
        if not run_theme_color_key:
            run_theme_color_key = _extract_run_theme_color_key(run)
        if not font_hex:
            font_hex = color_to_hex(font.color, theme_color_map) if font.color else None
        if font_hex:
            run_alpha = _extract_run_color_alpha(run)
            if run_alpha is not None and run_alpha < 1.0:
                font_hex = _hex_to_rgba(font_hex, run_alpha)
            styles.append(f"color:{font_hex}")
        elif eff_color:
            styles.append(f"color:{eff_color}")

        # 字号：run 显式 > defRPr > shape default
        if font.size:
            size_pt = round(font.size / 12700, 2)
            styles.append(f"font-size:{_format_numeric(size_pt)}pt")
        elif eff_font_size:
            styles.append(f"font-size:{_format_numeric(float(eff_font_size))}pt")

        # 字体：run 显式 > defRPr > shape default (附带 web-safe fallback 链)
        # 注意：font.name 可能是主题引用 (+mj-lt 等)，需先解析
        run_font_name = font.name
        if run_font_name and run_font_name.startswith("+"):
            run_font_name = _resolve_theme_font_reference(run_font_name, theme_fonts)
        if run_font_name:
            styles.append(f"font-family:{_font_with_fallback(run_font_name)}")
        elif eff_font_family:
            styles.append(f"font-family:{_font_with_fallback(eff_font_family)}")

        # 字符间距（run 级别 a:rPr spc 属性，单位 1/100pt → px）
        try:
            nsmap_a_spc = "http://schemas.openxmlformats.org/drawingml/2006/main"
            rPr_spc = run._r.find(f"{{{nsmap_a_spc}}}rPr")
            if rPr_spc is not None:
                spc_val = rPr_spc.get("spc")
                if spc_val and int(spc_val) != 0:
                    spc_px = round(int(spc_val) / 100 * PX_PER_PT, 2)
                    styles.append(f"letter-spacing:{_format_numeric(spc_px)}px")
        except Exception:
            pass

        span_attrs = []
        if run_theme_color_key:
            span_attrs.append(f'data-theme-color-key="{html_escape(run_theme_color_key)}"')
        if styles:
            span_attrs.append(f'style="{";".join(styles)}"')
        if span_attrs:
            tags_open.insert(0, f"<span {' '.join(span_attrs)}>")
            tags_close.append("</span>")

        # 包裹超链接
        if hyperlink_payload:
            if hyperlink_payload.get("type") == "slide":
                href = f'#{hyperlink_payload.get("target", "")}'
                tags_open.insert(0, f'<a href="{html_escape(href)}">')
            else:
                hyperlink_url = hyperlink_payload.get("target", "")
                tags_open.insert(0, f'<a href="{html_escape(hyperlink_url)}" target="_blank" rel="noopener">')
            tags_close.append("</a>")

        # 对文本做 HTML 转义，防止 < > & " 等字符破坏 HTML 结构
        safe_text = html_escape(text)
        parts.append("".join(tags_open) + safe_text + "".join(tags_close))

    html = "".join(parts)

    # ── 段落级别 style（text-align）──
    # 始终显式输出 text-align（包括 left），防止被父容器样式覆盖
    p_styles = []
    try:
        alignment = para.alignment
        if alignment is not None:
            align_str = _alignment_to_str(alignment, "left")
            p_styles.append(f"text-align:{align_str}")
        else:
            # 未显式设置对齐 → PPT 默认为 left
            p_styles.append("text-align:left")
    except Exception:
        pass

    # 段落行距（内联到 p style，比例行距输出纯数字，绝对行距输出 pt 值）
    try:
        ls_css = _extract_line_spacing_css(para)
        if ls_css:
            p_styles.append(f"line-height:{ls_css}")
    except Exception:
        pass

    # 段前间距（内联到 p style，保留每段独立的间距信息）
    try:
        if para.space_before is not None and para.space_before > 0:
            spc_pt = max(0.01, round(para.space_before / 12700, 2))
            p_styles.append(f"margin-top:{_format_numeric(spc_pt)}pt")
    except Exception:
        pass

    # 段后间距：不再内联 margin-bottom，统一由元素级 paragraphSpace 控制，
    # 避免与前端 CSS 叠加导致间距翻倍（B1-02）。

    # 左缩进（来自 para.left_indent，如果 pPr 有 marL 属性）
    try:
        nsmap_a_indent = "http://schemas.openxmlformats.org/drawingml/2006/main"
        pPr_indent = para._p.find(f"{{{nsmap_a_indent}}}pPr")
        if pPr_indent is not None:
            mar_l = pPr_indent.get("marL")
            if mar_l and int(mar_l) > 0:
                indent_px = round(int(mar_l) / 12700 * PX_PER_PT, 2)  # EMU → pt → px
                p_styles.append(f"padding-left:{_format_numeric(indent_px)}px")
            # 首行缩进（可正可负，负值为悬挂缩进）
            first_indent = pPr_indent.get("indent")
            if first_indent:
                indent_val = int(first_indent)
                if indent_val != 0:
                    indent_pt = round(indent_val / 12700, 2)
                    p_styles.append(f"text-indent:{_format_numeric(indent_pt)}pt")
    except Exception:
        pass

    return html, ";".join(p_styles)


def _extract_line_spacing(para) -> Optional[float]:
    """
    提取段落行距，返回 CSS line-height 值。

    python-pptx 的 line_spacing 属性：
    - float: 比例行距（如 1.5 = 1.5倍行距）→ 直接返回
    - int (EMU): 绝对行距（固定值，单位 EMU。如 228600 = 18pt）→ 返回 None
      绝对行距通过 _extract_line_spacing_css() 单独处理

    返回 CSS line-height 倍数（纯数字，如 1.5），或 None。
    """
    try:
        ls = para.line_spacing
        if ls is None:
            return None
        if isinstance(ls, float):
            return ls
        # int (EMU) 由 _extract_line_spacing_css 处理
    except Exception:
        pass
    return None


def _extract_line_spacing_css(para) -> Optional[str]:
    """
    提取段落行距的 CSS 值。

    - float: 比例行距 → 返回纯数字字符串（如 "1.5"）
    - int (EMU): 绝对行距 → 返回带 pt 单位的字符串（如 "18pt"），
      保留精确值，避免转为 ratio 时的精度损失。

    返回 CSS line-height 字符串，或 None。
    """
    try:
        ls = para.line_spacing
        if ls is None:
            return None
        if isinstance(ls, float):
            if ls > 0 and ls != 1.0:
                return _format_numeric(ls, 3)
            return None
        if isinstance(ls, int):
            # int 表示 EMU 绝对值 → 输出为 pt
            abs_pt = round(ls / 12700, 2)
            if abs_pt > 0:
                return f"{_format_numeric(abs_pt, 2)}pt"
    except Exception:
        pass
    return None


def _extract_paragraph_spacing(tf) -> Optional[float]:
    """提取文本框的段间距（仅取 space_after，space_before 已在逐段 HTML margin-top 中处理）"""
    try:
        for para in tf.paragraphs:
            if para.space_after is not None and para.space_after > 0:
                return max(0.01, round(para.space_after / 12700, 2))
    except Exception:
        pass
    return None


def _extract_text_vertical_align(tf) -> Optional[str]:
    """提取文本框垂直对齐方式"""
    try:
        from lxml import etree
        nsmap = {"a": "http://schemas.openxmlformats.org/drawingml/2006/main"}
        # python-pptx TextFrame 没有直接暴露 anchor 属性
        # 通过 XML 查找 <a:bodyPr anchor="ctr|b|t">
        body_pr = tf._txBody.find("a:bodyPr", nsmap)
        if body_pr is not None:
            anchor = body_pr.get("anchor")
            anchor_map = {
                "t": "top",
                "ctr": "middle",
                "b": "bottom",
            }
            if anchor and anchor in anchor_map:
                return anchor_map[anchor]
    except Exception:
        pass
    return None


def _extract_text_margins(tf) -> Optional[Dict[str, float]]:
    """提取文本框内边距（pt）。

    缺省的 inset 用 OOXML 默认值补齐（左右 0.1in、上下 0.05in），保证导入 deck
    往返导出时保留 PowerPoint 的有效内边距，而非被导出侧当作 0（会让文字换行漂移）。
    """
    try:
        nsmap = {"a": "http://schemas.openxmlformats.org/drawingml/2006/main"}
        body_pr = tf._txBody.find("a:bodyPr", nsmap)
        if body_pr is None:
            return None
        margins: Dict[str, float] = {}
        # bodyPr 属性：lIns, tIns, rIns, bIns（EMU）
        for attr, key in [("lIns", "left"), ("tIns", "top"), ("rIns", "right"), ("bIns", "bottom")]:
            val = body_pr.get(attr)
            emu = int(val) if val is not None else _DEFAULT_TEXT_INSET_EMU[key]
            margins[key] = round(emu / EMU_PER_PT, 1)
        return margins
    except Exception:
        pass
    return None


def _detect_bullet_type(para) -> Optional[str]:
    """
    检测段落的项目符号/编号类型。

    返回:
        "bullet" - 无序列表（圆点、方点、破折号等字符符号）
        "number" - 有序列表（1. 2. 或 a. b. 等自动编号）
        None     - 无项目符号
    """
    try:
        nsmap = {"a": "http://schemas.openxmlformats.org/drawingml/2006/main"}
        pPr = para._p.find("a:pPr", nsmap)
        if pPr is None:
            return None

        # 显式声明无项目符号
        if pPr.find("a:buNone", nsmap) is not None:
            return None

        # 字符型项目符号（• ■ ● ➢ 等）
        if pPr.find("a:buChar", nsmap) is not None:
            return "bullet"

        # 自动编号（1. 2. 或 a) b) 等）
        if pPr.find("a:buAutoNum", nsmap) is not None:
            return "number"

        # 图片型项目符号（少见但存在）
        if pPr.find("a:buBlip", nsmap) is not None:
            return "bullet"

        # 有 level > 0 但无 buNone：可能继承了母版的项目符号
        # 这里保守处理：只有显式声明才算列表
    except Exception:
        pass
    return None


def _detect_bullet_char(para) -> Optional[str]:
    """
    提取 <a:buChar> 的自定义符号字符。
    返回字符本身（如 '■', '●', '➢'），非默认 '•' 时才有意义。
    """
    try:
        nsmap = {"a": "http://schemas.openxmlformats.org/drawingml/2006/main"}
        pPr = para._p.find("a:pPr", nsmap)
        if pPr is not None:
            bu_char = pPr.find("a:buChar", nsmap)
            if bu_char is not None:
                char_val = bu_char.get("char")
                if char_val and char_val != "\u2022":
                    return char_val
    except Exception:
        pass
    return None


def _detect_bullet_style(para, theme_color_map=None) -> Optional[Dict[str, str]]:
    """
    提取项目符号的颜色、字号、字体。

    返回 {"color": "#FF0000", "fontSize": "80%", "fontFamily": "Wingdings"} 中存在的子集，
    或 None（无特殊格式）。
    """
    try:
        nsmap = {"a": "http://schemas.openxmlformats.org/drawingml/2006/main"}
        pPr = para._p.find("a:pPr", nsmap)
        if pPr is None:
            return None
        result: Dict[str, str] = {}
        # buClr — 项目符号颜色
        bu_clr = pPr.find("a:buClr", nsmap)
        if bu_clr is not None:
            color = _extract_color_with_transforms(bu_clr, theme_color_map)
            if color:
                result["color"] = color
        # buSzPct — 项目符号大小（百分比）
        bu_sz_pct = pPr.find("a:buSzPct", nsmap)
        if bu_sz_pct is not None:
            val = bu_sz_pct.get("val")
            if val:
                result["fontSize"] = f"{int(val) // 1000}%"
        # buSzPts — 项目符号大小（绝对点数）
        bu_sz_pts = pPr.find("a:buSzPts", nsmap)
        if bu_sz_pts is not None:
            val = bu_sz_pts.get("val")
            if val:
                result["fontSize"] = f"{int(val) // 100}pt"
        # buFont — 项目符号字体
        bu_font = pPr.find("a:buFont", nsmap)
        if bu_font is not None:
            typeface = bu_font.get("typeface")
            if typeface:
                result["fontFamily"] = typeface
        return result if result else None
    except Exception:
        pass
    return None


def _detect_number_format(para) -> Optional[str]:
    """
    检测有序列表的编号格式（buAutoNum 的 type 属性）。

    返回 OOXML buAutoNum type 字符串，如:
        "arabicPeriod"    → 1. 2. 3.
        "arabicParenR"    → 1) 2) 3)
        "alphaLcPeriod"   → a. b. c.
        "alphaUcPeriod"   → A. B. C.
        "romanLcPeriod"   → i. ii. iii.
        "romanUcPeriod"   → I. II. III.
        None              → 非有序列表或无法检测
    """
    try:
        nsmap = {"a": "http://schemas.openxmlformats.org/drawingml/2006/main"}
        pPr = para._p.find("a:pPr", nsmap)
        if pPr is not None:
            bu_auto = pPr.find("a:buAutoNum", nsmap)
            if bu_auto is not None:
                return bu_auto.get("type")
    except Exception:
        pass
    return None


# OOXML buAutoNum type → HTML <ol type> 映射
_NUM_FMT_TO_OL_TYPE = {
    "alphaLcParenBoth": "a", "alphaLcParenR": "a", "alphaLcPeriod": "a",
    "alphaUcParenBoth": "A", "alphaUcParenR": "A", "alphaUcPeriod": "A",
    "romanLcParenBoth": "i", "romanLcParenR": "i", "romanLcPeriod": "i",
    "romanUcParenBoth": "I", "romanUcParenR": "I", "romanUcPeriod": "I",
}

# HTML <ol type> → OOXML buAutoNum type 映射（逆向，用于写回）
_OL_TYPE_TO_NUM_FMT = {
    "a": "alphaLcPeriod", "A": "alphaUcPeriod",
    "i": "romanLcPeriod", "I": "romanUcPeriod",
    "1": "arabicPeriod",
}


def _number_format_to_ol_type(num_fmt: Optional[str]) -> Optional[str]:
    """OOXML buAutoNum type → HTML <ol type> 属性值"""
    if not num_fmt:
        return None
    return _NUM_FMT_TO_OL_TYPE.get(num_fmt)


def _build_nested_list(items: List[Dict[str, Any]], list_tag: str, close_tag: Optional[str] = None) -> str:
    """
    将扁平列表项（带 level 属性）构建为 Tiptap 兼容的嵌套 HTML 列表。

    Tiptap/ProseMirror 表示多级列表的方式是嵌套 <ul>/<ol>（放在父 <li> 内部），
    而不是 <li style="margin-left:..."> 扁平结构。

    list_tag: 开标签（可含属性，如 'ol type="a"'）
    close_tag: 关闭标签名（如 'ol'），不传则从 list_tag 提取
    """
    if not items:
        return ""

    ct = close_tag or list_tag.split()[0]  # 'ol type="a"' → 'ol'
    parts: List[str] = []
    parts.append(f"<{list_tag}>")
    depth = 0

    for idx, item in enumerate(items):
        target = item.get("level", 0)
        style_attr = f' style="{item["p_style"]}"' if item.get("p_style") else ""

        # ── 层级增加：在当前 <li> 内部嵌套（子级用基础标签名，不带属性） ──
        if target > depth:
            if idx == 0:
                while depth < target:
                    parts.append(f"<li><p></p><{ct}>")
                    depth += 1
            else:
                parts.append(f"<{ct}>")
                depth += 1
                while depth < target:
                    parts.append(f"<li><p></p><{ct}>")
                    depth += 1

        # ── 层级减少：关闭嵌套列表 ──
        elif target < depth:
            while depth > target:
                parts.append(f"</{ct}>")
                parts.append("</li>")
                depth -= 1

        # ── 添加当前列表项 ──
        parts.append(f"<li><p{style_attr}>{item['html']}</p>")

        # 判断是否关闭当前 <li>：如果下一项更深，保持开启让其成为父节点
        next_level = items[idx + 1].get("level", 0) if idx + 1 < len(items) else -1
        if next_level <= target:
            parts.append("</li>")

    # 关闭所有剩余的嵌套层级
    while depth > 0:
        parts.append(f"</{ct}>")
        parts.append("</li>")
        depth -= 1

    parts.append(f"</{ct}>")
    return "".join(parts)


def _escape_html_attr(val: str) -> str:
    """转义 HTML 属性值中的特殊字符。"""
    return val.replace('&', '&amp;').replace('"', '&quot;').replace('<', '&lt;').replace('>', '&gt;')


def _build_html_with_lists(para_infos: List[Dict[str, Any]]) -> str:
    """
    将段落信息列表转换为包含 <ul>/<ol> 的 HTML。

    Tiptap/ProseMirror 兼容格式：
    - <ul><li><p>text</p></li></ul> 无序列表（多级嵌套）
    - <ol><li><p>text</p></li></ol> 有序列表（多级嵌套）
    - <p>text</p> 普通段落
    """
    html_parts = []
    i = 0

    while i < len(para_infos):
        info = para_infos[i]
        bullet = info["bullet"]

        if bullet == "bullet":
            # 收集连续的无序列表项，交由嵌套构建器处理
            items = []
            while i < len(para_infos) and para_infos[i]["bullet"] == "bullet":
                items.append(para_infos[i])
                i += 1
            # 从首项提取自定义符号 + 样式 → data-bullet-* 属性
            bc = items[0].get("bulletChar") if items else None
            bs = items[0].get("bulletStyle") if items else None
            ul_attrs = []
            if bc:
                ul_attrs.append(f'data-bullet-char="{_escape_html_attr(bc)}"')
            if bs:
                if bs.get("color"):
                    ul_attrs.append(f'data-bullet-color="{_escape_html_attr(bs["color"])}"')
                if bs.get("fontSize"):
                    ul_attrs.append(f'data-bullet-font-size="{_escape_html_attr(bs["fontSize"])}"')
                if bs.get("fontFamily"):
                    ul_attrs.append(f'data-bullet-font="{_escape_html_attr(bs["fontFamily"])}"')
            ul_tag = 'ul ' + ' '.join(ul_attrs) if ul_attrs else 'ul'
            html_parts.append(_build_nested_list(items, ul_tag, "ul"))

        elif bullet == "number":
            # 收集连续的有序列表项，交由嵌套构建器处理
            items = []
            while i < len(para_infos) and para_infos[i]["bullet"] == "number":
                items.append(para_infos[i])
                i += 1
            # 从首项推断编号格式 → HTML <ol type>
            num_fmt = items[0].get("numberFormat") if items else None
            ol_type = _number_format_to_ol_type(num_fmt)
            ol_tag = f'ol type="{ol_type}"' if ol_type else "ol"
            html_parts.append(_build_nested_list(items, ol_tag, "ol"))

        else:
            # 普通段落（含缩进支持：para.level > 0 时添加 padding-left）
            p_style = info["p_style"]
            level = info.get("level", 0)
            if level > 0:
                indent_css = f"padding-left:{level * 24}px"
                p_style = f"{indent_css};{p_style}" if p_style else indent_css
            style_attr = f' style="{p_style}"' if p_style else ""
            html_parts.append(f"<p{style_attr}>{info['html']}</p>")
            i += 1

    return "".join(html_parts) if html_parts else "<p></p>"


def _extract_table_shape(
    shape,
    base: dict,
    theme_color_map: Optional[Dict[str, str]] = None,
    theme_fonts: Optional[Dict[str, str]] = None,
) -> Dict[str, Any]:
    """提取表格 → TableElement（包含单元格样式、合并单元格占位、行高）"""
    table = shape.table

    def _is_merged_placeholder(cell_obj) -> bool:
        """判断当前 cell 是否为被合并覆盖的占位单元格。"""
        try:
            if hasattr(cell_obj, "is_spanned") and cell_obj.is_spanned:
                return True
        except Exception:
            pass

        tc_el = cell_obj._tc if hasattr(cell_obj, "_tc") else None
        if tc_el is None:
            return False

        for attr in ("hMerge", "vMerge"):
            val = tc_el.get(attr)
            if not val:
                continue
            val_l = str(val).strip().lower()
            if val_l in ("1", "true", "continue", "cont"):
                return True
        return False

    rows = []
    row_heights = []
    for row in table.rows:
        # 提取行高（EMU → px，与前端坐标系一致）
        try:
            rh = row.height
            row_heights.append(round(rh / EMU_PER_PX) if rh and rh > 0 else 0)
        except Exception:
            row_heights.append(0)

        cells = []
        for cell in row.cells:
            if _is_merged_placeholder(cell):
                # 被合并覆盖的占位单元格 → 标记为 0
                cells.append({"text": "", "colspan": 0, "rowspan": 0})
                continue

            cell_data: Dict[str, Any] = {"text": cell.text or ""}

            # 如果单元格有多段落或混合格式，生成 richText HTML
            try:
                if cell.text_frame and len(cell.text_frame.paragraphs) > 0:
                    has_rich_content = False
                    para_count = len(cell.text_frame.paragraphs)
                    total_runs = sum(len(p.runs) for p in cell.text_frame.paragraphs)
                    # 检测是否有混合格式（多段落、多 run、或 run 有格式）
                    if para_count > 1 or total_runs > 1:
                        has_rich_content = True
                    elif total_runs == 1:
                        r = cell.text_frame.paragraphs[0].runs[0]
                        if r.font.bold or r.font.italic or r.font.underline:
                            has_rich_content = True
                        elif _extract_run_hyperlink(r):
                            has_rich_content = True
                    if has_rich_content:
                        html_parts = []
                        for para in cell.text_frame.paragraphs:
                            para_html, p_style = _paragraph_to_html_v2(
                                para,
                                theme_color_map,
                                theme_fonts=theme_fonts,
                            )
                            style_attr = f' style="{p_style}"' if p_style else ""
                            html_parts.append(f"<p{style_attr}>{para_html}</p>")
                        rich_html = "".join(html_parts)
                        if rich_html:
                            cell_data["richText"] = rich_html
            except Exception:
                pass

            if cell.span_height > 1:
                cell_data["rowspan"] = cell.span_height
            if cell.span_width > 1:
                cell_data["colspan"] = cell.span_width

            # 提取单元格样式（取第一个 run 的基本格式）
            cell_style = _extract_cell_style(cell, theme_color_map, theme_fonts=theme_fonts)
            if cell_style:
                # 兼容历史扁平字段，同时补充 style 对象供新前端直接消费
                cell_data.update(cell_style)
                cell_data["style"] = dict(cell_style)

            cells.append(cell_data)
        rows.append(cells)

    # 列宽比例
    total_width = sum(max(col.width, 0) for col in table.columns)
    col_widths: List[float] = []
    if total_width > 0:
        raw_ratios = [max(col.width, 0) / total_width for col in table.columns]
        ratio_sum = sum(raw_ratios)
        if ratio_sum > 0:
            col_widths = [round(r / ratio_sum, 6) for r in raw_ratios]
            if col_widths:
                drift = round(1 - sum(col_widths), 6)
                col_widths[-1] = round(col_widths[-1] + drift, 6)

    # 计算单元格最小高度（取所有行高的最小值或平均值）
    valid_heights = [h for h in row_heights if h > 0]
    cell_min_height = min(valid_heights) if valid_heights else 36
    normalized_row_heights = [h if h > 0 else cell_min_height for h in row_heights]

    props: Dict[str, Any] = {
        "data": rows,
        "colWidths": col_widths,
        "cellMinHeight": cell_min_height,
    }
    if normalized_row_heights:
        props["rowHeights"] = normalized_row_heights

    # 提取表格边框（六向 + 整体兜底）
    borders = _extract_table_borders(shape, theme_color_map)
    if borders:
        props["borders"] = borders
    outline = _extract_table_outline(shape, theme_color_map, pre_extracted_borders=borders)
    if not outline:
        outline = _extract_shape_outline(shape, theme_color_map)
    if outline:
        props["outline"] = outline

    # 提取表格主题属性（交替行/列、首行首列高亮等）
    theme = _extract_table_theme(shape, theme_color_map)
    if theme:
        props["theme"] = theme

    return {
        **base,
        "type": "table",
        "props": props,
    }


def _extract_table_borders(shape, theme_color_map: Optional[Dict[str, str]] = None) -> Optional[Dict[str, Any]]:
    """从 <a:tblPr><a:tblBorders> 提取六向边框（top/right/bottom/left/insideH/insideV）。"""
    try:
        nsmap = {"a": "http://schemas.openxmlformats.org/drawingml/2006/main"}
        tbl_pr = shape._element.find(".//a:tbl/a:tblPr", nsmap) if hasattr(shape, "_element") else None
        if tbl_pr is None:
            return None

        borders_el = tbl_pr.find("a:tblBorders", nsmap)
        if borders_el is None:
            return None

        dash_map = {
            "solid": "solid",
            "dash": "dashed",
            "dot": "dotted",
            "lgDash": "longDash",
            "dashDot": "dashDot",
            "sysDash": "dashed",
            "sysDot": "dotted",
            "lgDashDot": "longDashDot",
            "lgDashDotDot": "longDashDot",
        }
        scheme_to_key = {
            "dk1": "1", "lt1": "2", "dk2": "3", "lt2": "4",
            "accent1": "5", "accent2": "6", "accent3": "7",
            "accent4": "8", "accent5": "9", "accent6": "10",
        }

        sides = ("top", "right", "bottom", "left", "insideH", "insideV")
        result: Dict[str, Any] = {}

        for side in sides:
            border_el = borders_el.find(f"a:{side}", nsmap)
            if border_el is None:
                continue

            side_data: Dict[str, Any] = {"style": "solid", "width": 0, "color": "#333333"}
            if border_el.find("a:noFill", nsmap) is not None:
                result[side] = side_data
                continue

            width_px = 1.0
            w_attr = border_el.get("w")
            if w_attr:
                try:
                    width_px = round(int(w_attr) / EMU_PER_PT, 1)
                except Exception:
                    pass

            dash_val = border_el.get("prstDash")
            if not dash_val:
                dash_el = border_el.find("a:prstDash", nsmap)
                if dash_el is not None:
                    dash_val = dash_el.get("val")
            style = dash_map.get(dash_val or "solid", "solid")

            color = None
            theme_key = None
            solid_fill = border_el.find("a:solidFill", nsmap)
            if solid_fill is not None:
                theme_key = _extract_theme_key_from_solid_fill(solid_fill)
                srgb = solid_fill.find("a:srgbClr", nsmap)
                if srgb is not None and srgb.get("val"):
                    color = f"#{srgb.get('val')}"
                    alpha_el = srgb.find("a:alpha", nsmap)
                    if alpha_el is not None and alpha_el.get("val"):
                        try:
                            a_float = round(int(alpha_el.get("val")) / 100000, 2)
                            if a_float < 1.0:
                                color = _hex_to_rgba(color, a_float)
                        except (TypeError, ValueError):
                            pass
                if not color:
                    scheme = solid_fill.find("a:schemeClr", nsmap)
                    if scheme is not None and theme_color_map:
                        tc_key = scheme_to_key.get(scheme.get("val", ""))
                        if tc_key and tc_key in theme_color_map:
                            color = theme_color_map[tc_key]
                            alpha_el = scheme.find("a:alpha", nsmap)
                            if alpha_el is not None and alpha_el.get("val"):
                                try:
                                    a_float = round(int(alpha_el.get("val")) / 100000, 2)
                                    if a_float < 1.0:
                                        color = _hex_to_rgba(color, a_float)
                                except (TypeError, ValueError):
                                    pass

            result[side] = {
                "style": style,
                "width": max(width_px, 0),
                "color": color or "#333333",
                **({"themeKey": theme_key} if theme_key else {}),
            }

        if result:
            return result

    except Exception:
        pass
    return None


def _extract_table_outline(
    shape,
    theme_color_map: Optional[Dict[str, str]] = None,
    pre_extracted_borders: Optional[Dict[str, Any]] = None,
) -> Optional[Dict[str, Any]]:
    """从六向边框推导表格整体边框兜底样式。"""
    try:
        borders = pre_extracted_borders or _extract_table_borders(shape, theme_color_map)
        if not isinstance(borders, dict) or not borders:
            return None

        side_priority = {"top": 6, "right": 5, "bottom": 4, "left": 3, "insideH": 2, "insideV": 1}
        candidates: List[Dict[str, Any]] = []
        for side, spec in borders.items():
            if not isinstance(spec, dict):
                continue
            width = spec.get("width", 0)
            try:
                width_num = max(float(width), 0.0)
            except (TypeError, ValueError):
                width_num = 0.0
            candidates.append({
                "sidePriority": side_priority.get(side, 0),
                "style": spec.get("style", "solid"),
                "width": width_num,
                "color": spec.get("color", "#333333") or "#333333",
                "themeKey": spec.get("themeKey"),
            })

        if candidates:
            chosen = max(
                candidates,
                key=lambda it: (1 if it.get("width", 0) > 0 else 0, it.get("sidePriority", 0), it.get("width", 0)),
            )
            chosen.pop("sidePriority", None)
            return chosen
    except Exception:
        pass
    return None


def _extract_table_theme(shape, theme_color_map: Optional[Dict[str, str]] = None) -> Optional[Dict[str, Any]]:
    """
    从 <a:tblPr> 提取表格主题属性：交替行/列颜色、首行/首列/末行/末列高亮。

    PPTX XML 示例：
      <a:tblPr firstRow="1" bandRow="1" lastCol="0">
        <a:tblStyleId>{5C22544A-7EE6-4342-B048-85BDC9FD1C3A}</a:tblStyleId>
      </a:tblPr>
    """
    try:
        nsmap_a = "http://schemas.openxmlformats.org/drawingml/2006/main"
        tbl_el = shape._element.find(f".//{{{nsmap_a}}}tbl") if hasattr(shape, "_element") else None
        if tbl_el is None:
            return None

        tbl_pr = tbl_el.find(f"{{{nsmap_a}}}tblPr")
        if tbl_pr is None:
            return None

        # 读取布尔属性（兼容 "1"/"true"/"on"/"yes"）
        def _bool_attr(name: str) -> bool:
            val = tbl_pr.get(name)
            return str(val).strip().lower() in ("1", "true", "on", "yes")

        header_row = _bool_attr("firstRow")
        header_col = _bool_attr("firstCol")
        footer_row = _bool_attr("lastRow")
        last_col = _bool_attr("lastCol")
        striped_rows = _bool_attr("bandRow")
        striped_cols = _bool_attr("bandCol")

        # 如果没有任何主题属性被设置，返回 None
        if not any([header_row, header_col, footer_row, last_col, striped_rows, striped_cols]):
            return None

        # 尝试从表格首行第一个单元格的背景色推断主题色
        theme_color = None
        theme_color_key = None
        if header_row:
            try:
                first_cell = shape.table.cell(0, 0)
                if first_cell.fill and first_cell.fill.type == 1:  # SOLID
                    theme_color = color_to_hex(first_cell.fill.fore_color, theme_color_map)
                    theme_color_key = _extract_theme_key_from_color_obj(first_cell.fill.fore_color)
                    try:
                        nsmap = {"a": "http://schemas.openxmlformats.org/drawingml/2006/main"}
                        tc_pr = first_cell._tc.find("a:tcPr", nsmap) if hasattr(first_cell, "_tc") else None
                        if tc_pr is not None:
                            solid_fill = tc_pr.find("a:solidFill", nsmap)
                            if solid_fill is not None:
                                tk = _extract_theme_key_from_solid_fill(solid_fill)
                                if tk:
                                    theme_color_key = tk
                                elif solid_fill.find("a:schemeClr", nsmap) is not None:
                                    # schemeClr 有变换（tint/shade）时，提取基础 scheme val 作为参考
                                    _sc_el = solid_fill.find("a:schemeClr", nsmap)
                                    _sc_val = _sc_el.get("val") if _sc_el is not None else None
                                    if _sc_val:
                                        theme_color_key = _normalize_text_theme_key(_sc_val)
                                    else:
                                        theme_color_key = None
                    except Exception:
                        pass
            except Exception:
                pass

        # 如果无法推断，使用默认主题色
        if not theme_color:
            theme_color = "#5b9bd5"

        theme: Dict[str, Any] = {"color": theme_color}
        if theme_color_key:
            theme["colorThemeKey"] = theme_color_key
        if header_row:
            theme["headerRow"] = True
        if header_col:
            theme["headerCol"] = True
        if footer_row:
            theme["footerRow"] = True
        if last_col:
            theme["lastCol"] = True
        if striped_rows:
            theme["stripedRows"] = True
        if striped_cols:
            theme["stripedCols"] = True

        return theme

    except Exception:
        return None


def _extract_cell_style(
    cell,
    theme_color_map: Optional[Dict[str, str]] = None,
    theme_fonts: Optional[Dict[str, str]] = None,
) -> Dict[str, Any]:
    """提取表格单元格的样式信息（含背景色、字体、对齐、垂直对齐等）"""
    style: Dict[str, Any] = {}

    try:
        # 背景色（支持 alpha）
        fill = cell.fill
        if fill and fill.type is not None and fill.type == 1:  # SOLID
            bg_color = color_to_hex(fill.fore_color, theme_color_map)
            bg_theme_key = _extract_theme_key_from_color_obj(fill.fore_color)
            if bg_color:
                # 检查 cell fill 的 alpha
                try:
                    nsmap_a = "http://schemas.openxmlformats.org/drawingml/2006/main"
                    tcPr = cell._tc.find(f"{{{nsmap_a}}}tcPr") if hasattr(cell, "_tc") else None
                    if tcPr is not None:
                        solid = tcPr.find(f"{{{nsmap_a}}}solidFill")
                        if solid is not None and len(solid) > 0:
                            tk = _extract_theme_key_from_solid_fill(solid)
                            if tk:
                                bg_theme_key = tk
                            elif solid.find(f"{{{nsmap_a}}}schemeClr") is not None:
                                bg_theme_key = None
                            alpha_el = solid[0].find(f"{{{nsmap_a}}}alpha")
                            if alpha_el is not None:
                                a_val = alpha_el.get("val")
                                if a_val:
                                    a_float = round(int(a_val) / 100000, 2)
                                    if a_float < 1.0:
                                        bg_color = _hex_to_rgba(bg_color, a_float)
                except Exception:
                    pass
                style["bgColor"] = bg_color
                if bg_theme_key:
                    style["bgColorThemeKey"] = bg_theme_key
    except Exception:
        pass

    # 垂直对齐 + 单元格内边距：从 <a:tcPr> 提取
    try:
        nsmap_a = "http://schemas.openxmlformats.org/drawingml/2006/main"
        tcPr = cell._tc.find(f"{{{nsmap_a}}}tcPr") if hasattr(cell, "_tc") else None
        if tcPr is not None:
            anchor = tcPr.get("anchor")
            if anchor:
                v_align_map = {"t": "top", "ctr": "middle", "b": "bottom"}
                v_align = v_align_map.get(anchor)
                if v_align and v_align != "top":  # top 是默认值，可省略
                    style["verticalAlign"] = v_align

            # 单元格内边距 (EMU → px)，OOXML 默认值: marL=91440, marR=91440, marT=45720, marB=45720
            default_margins = {"marL": 91440, "marR": 91440, "marT": 45720, "marB": 45720}
            margin_keys = {"marL": "paddingLeft", "marR": "paddingRight", "marT": "paddingTop", "marB": "paddingBottom"}
            margins = {}
            has_non_default = False
            for xml_attr, fe_key in margin_keys.items():
                raw = tcPr.get(xml_attr)
                if raw is not None:
                    try:
                        emu_val = int(raw)
                        px_val = round(emu_val / 12700, 1)
                        margins[fe_key] = px_val
                        if emu_val != default_margins.get(xml_attr, 0):
                            has_non_default = True
                    except (ValueError, TypeError):
                        pass
            if has_non_default and margins:
                style["padding"] = margins

            # ── 单元格级别边框 <a:tcBorders>（per-cell borders） ──
            _tc_bdr_el = tcPr.find(f"{{{nsmap_a}}}tcBorders")
            if _tc_bdr_el is not None:
                _cb_dash = {
                    "solid": "solid", "dash": "dashed", "dot": "dotted",
                    "lgDash": "longDash", "dashDot": "dashDot", "sysDash": "dashed",
                    "sysDot": "dotted", "lgDashDot": "longDashDot", "lgDashDotDot": "longDashDot",
                }
                _cb_scheme = {
                    "dk1": "1", "lt1": "2", "dk2": "3", "lt2": "4",
                    "accent1": "5", "accent2": "6", "accent3": "7",
                    "accent4": "8", "accent5": "9", "accent6": "10",
                }
                _cb_result: Dict[str, Any] = {}
                for _cb_side in ("top", "right", "bottom", "left"):
                    _cb_el = _tc_bdr_el.find(f"{{{nsmap_a}}}{_cb_side}")
                    if _cb_el is None:
                        continue
                    if _cb_el.find(f"{{{nsmap_a}}}noFill") is not None:
                        _cb_result[_cb_side] = {"style": "solid", "width": 0, "color": "#000000"}
                        continue
                    _cb_w = _cb_el.get("w")
                    if not _cb_w:
                        continue
                    _cb_wpx = round(int(_cb_w) / 12700, 1)
                    if _cb_wpx <= 0:
                        continue
                    _cb_color = "#000000"
                    _cb_tk = None
                    _cb_sf = _cb_el.find(f"{{{nsmap_a}}}solidFill")
                    if _cb_sf is not None:
                        _cb_srgb = _cb_sf.find(f"{{{nsmap_a}}}srgbClr")
                        if _cb_srgb is not None:
                            _cb_color = "#" + _cb_srgb.get("val", "000000")
                        _cb_sc = _cb_sf.find(f"{{{nsmap_a}}}schemeClr")
                        if _cb_sc is not None:
                            _cb_sv = _cb_sc.get("val", "")
                            _cb_tk = _cb_scheme.get(_cb_sv)
                            if theme_color_map and _cb_sv in theme_color_map:
                                _cb_color = "#" + theme_color_map[_cb_sv]
                    _cb_dash_val = _cb_el.get("prstDash", "solid")
                    _cb_style = _cb_dash.get(_cb_dash_val, "solid")
                    _cb_spec: Dict[str, Any] = {"style": _cb_style, "width": _cb_wpx, "color": _cb_color}
                    if _cb_tk:
                        _cb_spec["themeKey"] = _cb_tk
                    _cb_result[_cb_side] = _cb_spec
                if _cb_result:
                    style["cellBorders"] = _cb_result
    except Exception:
        pass

    try:
        # 文本样式（取第一个非空 run）
        if cell.text_frame:
            for para in cell.text_frame.paragraphs:
                # 段落对齐
                try:
                    if para.alignment is not None:
                        align_str = _alignment_to_str(para.alignment, "left")
                        if align_str:
                            style["align"] = align_str
                except Exception:
                    pass

                for run in para.runs:
                    font = run.font
                    if font.bold and "bold" not in style:
                        style["bold"] = True
                    if font.italic and "italic" not in style:
                        style["italic"] = True
                    if font.underline and "underline" not in style:
                        style["underline"] = True
                    if font.size and "fontSize" not in style:
                        style["fontSize"] = round(font.size / 12700)
                    if font.name and "fontFamily" not in style:
                        resolved_font = _resolve_theme_font_reference(font.name, theme_fonts) if font.name.startswith("+") else font.name
                        font_with_fallback = _font_with_fallback(resolved_font)
                        if font_with_fallback:
                            style["fontFamily"] = font_with_fallback
                    if font.color and "color" not in style:
                        c = color_to_hex(font.color, theme_color_map)
                        if c:
                            # 检查 alpha
                            run_alpha = _extract_run_color_alpha(run)
                            if run_alpha is not None and run_alpha < 1.0:
                                c = _hex_to_rgba(c, run_alpha)
                            style["color"] = c
                        color_theme_key = _extract_run_theme_color_key(run)
                        if color_theme_key:
                            style["colorThemeKey"] = color_theme_key
                    break
                if style:
                    break
    except Exception:
        pass

    return style


def _extract_shape_alt_text(shape) -> Optional[str]:
    """
    读取 shape 替代文本（Alt Text）。

    优先返回 cNvPr@descr；为空时回退 cNvPr@title。
    """
    try:
        nsmap_p = "http://schemas.openxmlformats.org/presentationml/2006/main"
        c_nv_pr = shape._element.find(f".//{{{nsmap_p}}}cNvPr")
        if c_nv_pr is None:
            return None
        descr = (c_nv_pr.get("descr") or "").strip()
        if descr:
            return descr
        title = (c_nv_pr.get("title") or "").strip()
        if title:
            return title
    except Exception:
        pass
    return None


def _normalize_media_ext(raw: Any) -> Optional[str]:
    if raw is None:
        return None
    s = str(raw).strip().lower()
    if not s:
        return None
    if s.startswith("."):
        s = s[1:]
    s = s.split("?", 1)[0].split("#", 1)[0]
    return s or None


def _extract_file_ext_from_target(raw_target: Any) -> Optional[str]:
    if raw_target is None:
        return None
    target = str(raw_target).strip()
    if not target:
        return None
    clean = target.split("#", 1)[0].split("?", 1)[0]
    file_name = clean.rsplit("/", 1)[-1]
    if "." not in file_name:
        return None
    return _normalize_media_ext(file_name.rsplit(".", 1)[-1])


def _infer_media_mime(ext: Optional[str], fallback: Optional[str] = None) -> Optional[str]:
    normalized_ext = _normalize_media_ext(ext)
    if normalized_ext and normalized_ext in MEDIA_MIME_BY_EXT:
        return MEDIA_MIME_BY_EXT[normalized_ext]
    if isinstance(fallback, str) and "/" in fallback:
        return fallback.strip().lower()
    return None


def _guess_media_kind(content_type: Optional[str], ext: Optional[str]) -> Optional[str]:
    ct = (content_type or "").strip().lower()
    normalized_ext = _normalize_media_ext(ext)
    if ct.startswith("audio/"):
        return "audio"
    if ct.startswith("video/"):
        return "video"
    if normalized_ext in MEDIA_AUDIO_EXTS:
        return "audio"
    if normalized_ext in MEDIA_VIDEO_EXTS:
        return "video"
    return None


def _resolve_shape_rel_source(shape, rel_id: str, asset_handler: Optional[Any] = None) -> Optional[Dict[str, Any]]:
    import base64 as b64

    if not rel_id:
        return None
    rels = getattr(getattr(shape, "part", None), "rels", None)
    if rels is None:
        return None

    try:
        rel = rels[rel_id]
    except Exception:
        return None

    target_ref = str(getattr(rel, "target_ref", "") or "")
    rel_type = str(getattr(rel, "reltype", "") or "")
    target_part = getattr(rel, "target_part", None)
    content_type = getattr(target_part, "content_type", None) if target_part is not None else None
    blob = getattr(target_part, "blob", None) if target_part is not None else None
    ext = _extract_file_ext_from_target(target_ref)

    src = ""
    if blob:
        ct = _infer_media_mime(ext, content_type) or "application/octet-stream"
        if asset_handler and callable(asset_handler):
            try:
                uploaded_url = asset_handler(blob, ct)
                if uploaded_url:
                    src = uploaded_url
            except Exception as e:
                logger.warning(f"asset_handler failed for rel={rel_id}, falling back to base64: {e}")
        if not src:
            src = f"data:{ct};base64,{b64.b64encode(blob).decode('ascii')}"
        content_type = ct
    else:
        src = target_ref
        if not content_type:
            content_type = _infer_media_mime(ext)

    if not src:
        return None

    return {
        "src": src,
        "content_type": content_type,
        "ext": ext,
        "target_ref": target_ref,
        "rel_type": rel_type,
    }


def _extract_media_rel_candidates(shape) -> List[str]:
    ns = {
        "p": "http://schemas.openxmlformats.org/presentationml/2006/main",
        "a": "http://schemas.openxmlformats.org/drawingml/2006/main",
        "r": "http://schemas.openxmlformats.org/officeDocument/2006/relationships",
        "p14": "http://schemas.microsoft.com/office/powerpoint/2010/main",
    }

    candidates: List[str] = []
    seen: Set[str] = set()

    def _push(val: Optional[str]) -> None:
        if not val:
            return
        rid = str(val).strip()
        if not rid or rid in seen:
            return
        seen.add(rid)
        candidates.append(rid)

    try:
        nv_pr = shape._element.find(".//p:nvPicPr/p:nvPr", ns)
        if nv_pr is None:
            nv_pr = shape._element.find(".//p:nvPr", ns)
        if nv_pr is None:
            return candidates

        audio_file = nv_pr.find("a:audioFile", ns)
        if audio_file is not None:
            _push(audio_file.get(f"{{{ns['r']}}}embed"))
            _push(audio_file.get(f"{{{ns['r']}}}link"))

        video_file = nv_pr.find("a:videoFile", ns)
        if video_file is not None:
            _push(video_file.get(f"{{{ns['r']}}}embed"))
            _push(video_file.get(f"{{{ns['r']}}}link"))

        media_node = nv_pr.find(".//p14:media", ns)
        if media_node is not None:
            _push(media_node.get(f"{{{ns['r']}}}embed"))
            _push(media_node.get(f"{{{ns['r']}}}link"))
    except Exception:
        pass

    return candidates


def _encode_media_alt_text(media_type: str, props: Dict[str, Any]) -> Optional[str]:
    if media_type not in ("video", "audio") or not isinstance(props, dict):
        return None

    payload: Dict[str, Any] = {"type": media_type}
    autoplay = _coerce_bool_flag(props.get("autoplay"))
    if autoplay is not None:
        payload["autoplay"] = autoplay

    ext = _normalize_media_ext(props.get("ext"))
    if ext:
        payload["ext"] = ext

    if media_type == "audio":
        loop = _coerce_bool_flag(props.get("loop"))
        if loop is not None:
            payload["loop"] = loop
        fixed_ratio = _coerce_bool_flag(props.get("fixedRatio"))
        if fixed_ratio is not None:
            payload["fixedRatio"] = fixed_ratio
        color = props.get("color")
        if isinstance(color, str) and color.strip():
            payload["color"] = color.strip()

    import base64 as _b64
    import json as _json

    try:
        raw = _json.dumps(payload, ensure_ascii=False, separators=(",", ":")).encode("utf-8")
        encoded = MEDIA_META_PREFIX + _b64.b64encode(raw).decode("ascii")
        if len(encoded) <= MAX_MEDIA_META_LENGTH:
            return encoded
    except Exception:
        pass
    return None


def _decode_media_alt_text(raw_alt_text: Optional[str]) -> Optional[Dict[str, Any]]:
    import base64 as _b64
    import json as _json

    if not isinstance(raw_alt_text, str):
        return None
    text = raw_alt_text.strip()
    if not text.startswith(MEDIA_META_PREFIX):
        return None
    payload_b64 = text[len(MEDIA_META_PREFIX):].strip()
    if not payload_b64:
        return None

    try:
        decoded = _b64.b64decode(payload_b64).decode("utf-8")
        data = _json.loads(decoded)
    except Exception:
        return None

    if not isinstance(data, dict):
        return None
    media_type = data.get("type")
    if media_type not in ("video", "audio"):
        return None
    return data


def _extract_media_shape(
    shape,
    base: Dict[str, Any],
    image_handler: Optional[Any] = None,
) -> Optional[Dict[str, Any]]:
    """
    提取媒体 shape（MSO_SHAPE_TYPE.MEDIA / WEB_VIDEO）为 video/audio 元素。
    """
    ns = {
        "p": "http://schemas.openxmlformats.org/presentationml/2006/main",
        "a": "http://schemas.openxmlformats.org/drawingml/2006/main",
        "r": "http://schemas.openxmlformats.org/officeDocument/2006/relationships",
    }

    media_info: Optional[Dict[str, Any]] = None
    media_kind: Optional[str] = None
    rel_ids = _extract_media_rel_candidates(shape)

    for rid in rel_ids:
        info = _resolve_shape_rel_source(shape, rid, asset_handler=image_handler)
        if not info:
            continue
        kind = _guess_media_kind(info.get("content_type"), info.get("ext"))
        if kind is None:
            rel_type = str(info.get("rel_type") or "").lower()
            if "audio" in rel_type:
                kind = "audio"
            elif "video" in rel_type or "media" in rel_type:
                kind = "video"
        if kind:
            media_info = info
            media_kind = kind
            break
        if media_info is None:
            media_info = info

    if media_kind is None and media_info is not None:
        # 关系类型无法判断时，按扩展名兜底；仍无法判断则默认按 video 处理。
        media_kind = _guess_media_kind(media_info.get("content_type"), media_info.get("ext")) or "video"

    if media_info is None or media_kind is None:
        return None

    # 读取封面图（poster）：来自 blipFill 的 image rel。
    poster_src: Optional[str] = None
    try:
        blip = shape._element.find(".//p:blipFill/a:blip", ns)
        if blip is None:
            blip = shape._element.find(".//a:blip", ns)
        if blip is not None:
            poster_rid = blip.get(f"{{{ns['r']}}}embed")
            if poster_rid:
                poster_info = _resolve_shape_rel_source(shape, poster_rid, asset_handler=image_handler)
                if poster_info and str(poster_info.get("content_type") or "").startswith("image/"):
                    poster_src = poster_info.get("src")
    except Exception:
        poster_src = None

    meta = _decode_media_alt_text(_extract_shape_alt_text(shape))
    meta_autoplay = _coerce_bool_flag(meta.get("autoplay")) if isinstance(meta, dict) else None
    meta_ext = _normalize_media_ext(meta.get("ext")) if isinstance(meta, dict) else None
    rel_ext = _normalize_media_ext(media_info.get("ext"))
    ext = rel_ext
    known_exts = MEDIA_AUDIO_EXTS | MEDIA_VIDEO_EXTS
    if meta_ext and (not ext or ext not in known_exts):
        ext = meta_ext

    props: Dict[str, Any] = {
        "src": media_info.get("src", ""),
        "autoplay": bool(meta_autoplay) if meta_autoplay is not None else False,
    }
    if ext:
        props["ext"] = ext

    if media_kind == "video":
        if poster_src:
            props["poster"] = poster_src
        return {**base, "type": "video", "props": props}

    # audio
    loop_flag = _coerce_bool_flag(meta.get("loop")) if isinstance(meta, dict) else None
    fixed_ratio = _coerce_bool_flag(meta.get("fixedRatio")) if isinstance(meta, dict) else None
    color = meta.get("color") if isinstance(meta, dict) else None
    props["loop"] = bool(loop_flag) if loop_flag is not None else False
    props["fixedRatio"] = bool(fixed_ratio) if fixed_ratio is not None else True
    props["color"] = color.strip() if isinstance(color, str) and color.strip() else "#666666"
    return {**base, "type": "audio", "props": props}


def _extract_image_shape(shape, base: dict, image_handler: Optional[Any] = None) -> Dict[str, Any]:
    """提取图片 → ImageElement（支持 OSS 上传或 base64 内联）"""
    import base64 as b64

    image = shape.image
    content_type = image.content_type if hasattr(image, "content_type") else "image/png"

    src = ""
    try:
        if hasattr(image, "blob") and image.blob:
            blob_data = image.blob

            # 如果提供了 image_handler（如 OSS 上传），优先使用
            if image_handler and callable(image_handler):
                try:
                    uploaded_url = image_handler(blob_data, content_type)
                    if uploaded_url:
                        src = uploaded_url
                except Exception as e:
                    logger.warning(f"image_handler failed, falling back to base64: {e}")

            # 降级或无 handler 时使用 base64
            if not src:
                encoded = b64.b64encode(blob_data).decode("ascii")
                src = f"data:{content_type};base64,{encoded}"
    except Exception as e:
        logger.warning(f"Failed to extract image blob: {e}")

    props: Dict[str, Any] = {
        "src": src,
    }

    alt_text = _extract_shape_alt_text(shape)
    if alt_text:
        props["altText"] = alt_text

    nsmap_a = "http://schemas.openxmlformats.org/drawingml/2006/main"
    nsmap_p = "http://schemas.openxmlformats.org/presentationml/2006/main"

    # ── 1. 图片裁剪 (srcRect) ──
    try:
        blip_fill = shape._element.find(f".//{{{nsmap_p}}}blipFill")
        if blip_fill is None:
            blip_fill = shape._element.find(f".//{{{nsmap_a}}}blipFill")
        if blip_fill is not None:
            src_rect = blip_fill.find(f"{{{nsmap_a}}}srcRect")
            if src_rect is not None:
                t = int(src_rect.get("t", "0")) / 100000
                r = int(src_rect.get("r", "0")) / 100000
                b_val = int(src_rect.get("b", "0")) / 100000
                l_val = int(src_rect.get("l", "0")) / 100000
                if any(v > 0.001 for v in [t, r, b_val, l_val]):
                    props["clip"] = {
                        "shape": "rect",
                        "range": [
                            [l_val, t],
                            [1 - r, t],
                            [1 - r, 1 - b_val],
                            [l_val, 1 - b_val],
                        ],
                    }

            # ── 2. Object-fit：stretch/fillRect ──
            stretch = blip_fill.find(f"{{{nsmap_a}}}stretch")
            if stretch is not None:
                fill_rect = stretch.find(f"{{{nsmap_a}}}fillRect")
                if fill_rect is not None:
                    fr_t = int(fill_rect.get("t", "0"))
                    fr_r = int(fill_rect.get("r", "0"))
                    fr_b = int(fill_rect.get("b", "0"))
                    fr_l = int(fill_rect.get("l", "0"))
                    # fillRect 有值 → 图片不完全填充，类似 contain
                    if any(v != 0 for v in [fr_t, fr_r, fr_b, fr_l]):
                        props["objectFit"] = "contain"
                    # fillRect 全 0 或不存在 → cover
                # 无 fillRect → 默认 stretch = cover

            # tile = 平铺模式（少见，暂不映射）
    except Exception:
        pass

    # ── 3. 图片边框 ──
    outline = _extract_shape_outline(shape, None)
    if outline:
        props["outline"] = outline

    # ── 4. 几何形状（圆角 / 圆形 / 多角变体） ──
    try:
        nsmap_p = "http://schemas.openxmlformats.org/presentationml/2006/main"
        spPr = shape._element.find(f".//{{{nsmap_p}}}spPr")
        if spPr is None:
            spPr = shape._element.find(f".//{{{nsmap_a}}}spPr")
        if spPr is not None:
            prstGeom = spPr.find(f"{{{nsmap_a}}}prstGeom")
            if prstGeom is not None:
                prst = prstGeom.get("prst", "")

                if prst == "ellipse":
                    # 圆形/椭圆遮罩
                    props["clip"] = {"shape": "ellipse", "range": []}

                elif prst in ("roundRect", "round1Rect", "round2SameRect", "round2DiagRect", "snipRndRect"):
                    # 圆角矩形（含多角变体）
                    avLst = prstGeom.find(f"{{{nsmap_a}}}avLst")
                    if avLst is not None:
                        gd = avLst.find(f"{{{nsmap_a}}}gd")
                        if gd is not None:
                            fmla = gd.get("fmla", "")
                            if "val" in fmla:
                                val = int(fmla.split()[-1]) if fmla.split() else 16667
                                radius_px = round(val / 100000 * min(base.get("width", 100), base.get("height", 100)))
                                if radius_px > 0:
                                    props["radius"] = radius_px
                        else:
                            # roundRect 无 avLst/gd → 使用默认 16.67%
                            default_radius = round(16667 / 100000 * min(base.get("width", 100), base.get("height", 100)))
                            if default_radius > 0:
                                props["radius"] = default_radius
                    else:
                        default_radius = round(16667 / 100000 * min(base.get("width", 100), base.get("height", 100)))
                        if default_radius > 0:
                            props["radius"] = default_radius
    except Exception:
        pass

    # ── 5. 图片级透明度 (blip alphaModFix) ──
    try:
        blip = shape._element.find(f".//{{{nsmap_a}}}blip")
        if blip is not None:
            alpha_mod = blip.find(f"{{{nsmap_a}}}alphaModFix")
            if alpha_mod is not None:
                amt = alpha_mod.get("amt")
                if amt:
                    img_opacity = round(int(amt) / 100000, 2)
                    if img_opacity < 1.0:
                        # 将 blip 级透明度合并到 base.opacity
                        base_opacity = base.get("opacity", 1.0)
                        base["opacity"] = round(base_opacity * img_opacity, 4)
    except Exception:
        pass

    # ── 6. 颜色滤镜 (blip 子元素) ──
    filters: Dict[str, Any] = {}
    try:
        blip = shape._element.find(f".//{{{nsmap_a}}}blip")
        if blip is not None:
            # 灰度
            if blip.find(f"{{{nsmap_a}}}grayscl") is not None:
                filters["grayscale"] = 1

            # 模糊: <a:blur rad="12700"/>
            blur_el = blip.find(f"{{{nsmap_a}}}blur")
            if blur_el is not None:
                blur_rad = blur_el.get("rad")
                if blur_rad:
                    try:
                        blur_px = round(max(0, int(blur_rad)) / EMU_PER_PT, 2)
                        if blur_px > 0:
                            filters["blur"] = blur_px
                    except Exception:
                        pass

            # 反相
            if blip.find(f"{{{nsmap_a}}}inv") is not None:
                filters["invert"] = 1

            # 亮度/对比度: <a:lum bright="20000" contrast="-10000"/>
            lum = blip.find(f"{{{nsmap_a}}}lum")
            if lum is not None:
                bright = lum.get("bright")
                if bright:
                    # PPTX bright: -100000 ~ 100000 → CSS brightness: 0 ~ 2 (1=原始)
                    filters["brightness"] = round(1 + int(bright) / 100000, 2)
                contrast = lum.get("contrast")
                if contrast:
                    # PPTX contrast: -100000 ~ 100000 → CSS contrast: 0 ~ 2 (1=原始)
                    filters["contrast"] = round(1 + int(contrast) / 100000, 2)

            # 色相/饱和度/亮度: <a:hsl hue="0" sat="100000" lum="0"/>
            hsl = blip.find(f"{{{nsmap_a}}}hsl")
            if hsl is not None:
                hue = hsl.get("hue")
                if hue:
                    # PPTX hue: 0 ~ 21600000 (60000 = 1度) → CSS hue-rotate(deg)
                    filters["hueRotate"] = round(int(hue) / 60000)
                sat = hsl.get("sat")
                if sat:
                    # PPTX sat: -100000 ~ 100000 → CSS saturate: 0 ~ 2
                    filters["saturate"] = round(1 + int(sat) / 100000, 2)

            # 双色调效果 → 近似灰度 + 棕褐色
            if blip.find(f"{{{nsmap_a}}}duotone") is not None:
                filters.setdefault("sepia", 0.8)
                filters.setdefault("saturate", 1.5)

            # 二值化 → 近似高对比度
            bi_level = blip.find(f"{{{nsmap_a}}}biLevel")
            if bi_level is not None:
                filters["contrast"] = 2.0
                filters["grayscale"] = 1
    except Exception:
        pass

    if filters:
        props["filters"] = filters

    # 锁定宽高比（noChangeAspect）
    try:
        nsmap_a_lock = "http://schemas.openxmlformats.org/drawingml/2006/main"
        cNvPicPr = shape._element.find(f".//{{{nsmap_a_lock}}}cNvPicPr")
        if cNvPicPr is not None:
            pic_locks = cNvPicPr.find(f"{{{nsmap_a_lock}}}picLocks")
            if pic_locks is not None and pic_locks.get("noChangeAspect") == "1":
                props["fixedRatio"] = True
    except Exception:
        pass

    return {**base, "type": "image", "props": props}


def _extract_first_drawingml_color(
    node,
    theme_color_map: Optional[Dict[str, str]] = None,
) -> Optional[str]:
    """从任意 DrawingML 节点提取第一个颜色（srgbClr / schemeClr）。"""
    color, _ = _extract_first_drawingml_color_and_theme_key(node, theme_color_map)
    return color


def _extract_first_drawingml_color_and_theme_key(
    node,
    theme_color_map: Optional[Dict[str, str]] = None,
) -> Tuple[Optional[str], Optional[str]]:
    """
    从任意 DrawingML 节点提取第一个颜色与主题色 key。

    返回：(color, theme_key)
    - color: #RRGGBB
    - theme_key: tx1/bg1/accent1...（仅当 schemeClr 无 tint/shade 等变换时返回）
    """
    nsmap_a = "http://schemas.openxmlformats.org/drawingml/2006/main"
    if node is None:
        return None, None

    # 优先显式 RGB
    try:
        for srgb in node.iter(f"{{{nsmap_a}}}srgbClr"):
            val = srgb.get("val")
            if val and len(val) == 6:
                return f"#{val}", None
    except Exception:
        pass

    # 其次主题色
    try:
        for scheme in node.iter(f"{{{nsmap_a}}}schemeClr"):
            scheme_name = (scheme.get("val") or "").strip()
            if not scheme_name:
                continue

            # 颜色始终尝试按 scheme 解析（优先文档主题，失败回退 Office 默认）
            color = _resolve_scheme_color(scheme_name, theme_color_map)

            # key 仅在纯 themeClr(+alpha) 语义时保留
            theme_key: Optional[str] = None
            has_non_alpha_transform = False
            for child in list(scheme):
                tag = child.tag.split("}")[-1] if isinstance(child.tag, str) and "}" in child.tag else str(child.tag)
                if tag != "alpha":
                    has_non_alpha_transform = True
                    break
            if not has_non_alpha_transform:
                theme_key = _normalize_text_theme_key(scheme_name)

            if color:
                return color, theme_key
            if theme_key:
                return None, theme_key
    except Exception:
        pass

    return None, None


def _extract_series_color(
    series, theme_color_map: Optional[Dict[str, str]] = None,
) -> Optional[str]:
    """从图表系列中提取颜色（支持显式 RGB / 主题色 / XML 级降级）。"""
    color, _ = _extract_series_color_and_theme_key(series, theme_color_map)
    return color


def _extract_series_color_and_theme_key(
    series, theme_color_map: Optional[Dict[str, str]] = None,
) -> Tuple[Optional[str], Optional[str]]:
    """从图表系列中提取颜色与主题色 key。"""
    # 策略 1: python-pptx API — 显式填充色
    try:
        fmt = series.format
        if fmt and fmt.fill and fmt.fill.type is not None:
            fore_color = fmt.fill.fore_color
            rgb = fore_color.rgb
            theme_key = _extract_theme_key_from_color_obj(fore_color)
            if rgb:
                return "#" + str(rgb), theme_key
            if theme_key:
                resolved = _resolve_scheme_color(theme_key, theme_color_map)
                if resolved:
                    return resolved, theme_key
                return None, theme_key
    except Exception:
        pass

    # 策略 2: python-pptx API — 线条色（折线图等无填充的类型）
    try:
        fmt = series.format
        if fmt and fmt.line and fmt.line.color and fmt.line.color.type is not None:
            line_color = fmt.line.color
            rgb = line_color.rgb
            theme_key = _extract_theme_key_from_color_obj(line_color)
            if rgb:
                return "#" + str(rgb), theme_key
            if theme_key:
                resolved = _resolve_scheme_color(theme_key, theme_color_map)
                if resolved:
                    return resolved, theme_key
                return None, theme_key
    except Exception:
        pass

    # 策略 3: XML 级提取 — 直接从 series XML 中解析 srgbClr / schemeClr
    try:
        ser_el = series._element
        extracted_color, extracted_key = _extract_first_drawingml_color_and_theme_key(ser_el, theme_color_map)
        if extracted_color or extracted_key:
            return extracted_color, extracted_key
    except Exception:
        pass

    return None, None


def _extract_pie_point_colors(
    series, theme_color_map: Optional[Dict[str, str]] = None,
) -> List[str]:
    """
    提取饼图/环形图每个扇区颜色，优先点级颜色（dPt），回退系列色。
    """
    colors, _ = _extract_pie_point_colors_and_theme_keys(series, theme_color_map)
    return colors


def _extract_pie_point_colors_and_theme_keys(
    series, theme_color_map: Optional[Dict[str, str]] = None,
) -> Tuple[List[str], List[Optional[str]]]:
    """
    提取饼图/环形图每个扇区颜色与主题色 key，优先点级（dPt），回退系列级。
    """
    nsmap_c = "http://schemas.openxmlformats.org/drawingml/2006/chart"

    colors_by_idx: Dict[int, str] = {}
    keys_by_idx: Dict[int, Optional[str]] = {}
    series_fallback, series_fallback_key = _extract_series_color_and_theme_key(series, theme_color_map)

    # 策略 1：python-pptx 点级 API
    try:
        for idx, point in enumerate(getattr(series, "points", []) or []):
            color = None
            theme_key = None
            try:
                fmt = point.format
                if fmt and fmt.fill and fmt.fill.type is not None:
                    fore_color = fmt.fill.fore_color
                    rgb = fore_color.rgb
                    theme_key = _extract_theme_key_from_color_obj(fore_color)
                    if rgb:
                        color = "#" + str(rgb)
                    elif theme_key:
                        color = _resolve_scheme_color(theme_key, theme_color_map)
            except Exception:
                pass
            if not color:
                try:
                    fmt = point.format
                    if fmt and fmt.line and fmt.line.color and fmt.line.color.type is not None:
                        line_color = fmt.line.color
                        rgb = line_color.rgb
                        if not theme_key:
                            theme_key = _extract_theme_key_from_color_obj(line_color)
                        if rgb:
                            color = "#" + str(rgb)
                        elif theme_key:
                            color = _resolve_scheme_color(theme_key, theme_color_map)
                except Exception:
                    pass
            if color or theme_key:
                if theme_key:
                    keys_by_idx[idx] = theme_key
            if color:
                colors_by_idx[idx] = color
    except Exception:
        pass

    # 策略 2：XML 点级 dPt
    try:
        ser_el = series._element
        for dpt in ser_el.findall(f"{{{nsmap_c}}}dPt"):
            idx_el = dpt.find(f"{{{nsmap_c}}}idx")
            if idx_el is None:
                continue
            try:
                idx = int(idx_el.get("val") or -1)
            except (TypeError, ValueError):
                continue
            if idx < 0:
                continue
            extracted_color, extracted_key = _extract_first_drawingml_color_and_theme_key(dpt, theme_color_map)
            if idx not in colors_by_idx and extracted_color:
                colors_by_idx[idx] = extracted_color
            if extracted_key and (idx not in keys_by_idx or not keys_by_idx.get(idx)):
                keys_by_idx[idx] = extracted_key
    except Exception:
        pass

    point_count = 0
    try:
        values = list(series.values) if series.values else []
        point_count = len(values)
    except Exception:
        point_count = 0
    if point_count <= 0 and colors_by_idx:
        point_count = max(colors_by_idx.keys()) + 1

    # 未获取到点级颜色时，保留系列色语义（单色）
    if not colors_by_idx and not keys_by_idx:
        if series_fallback:
            return [series_fallback], [series_fallback_key]
        return [], []

    out_colors: List[str] = []
    out_keys: List[Optional[str]] = []
    for idx in range(point_count):
        color = colors_by_idx.get(idx) or series_fallback
        theme_key = keys_by_idx.get(idx)
        if not theme_key:
            theme_key = series_fallback_key
        if color:
            out_colors.append(color)
            out_keys.append(theme_key)

    return out_colors, out_keys


def _extract_chart_series_name(series, default_name: str) -> str:
    """提取图表系列名，优先 python-pptx API，失败时回退 XML 解析。"""
    nsmap_c = "http://schemas.openxmlformats.org/drawingml/2006/chart"

    try:
        name = getattr(series, "name", None)
        if name is not None:
            name_text = str(name).strip()
            if name_text:
                return name_text
    except Exception:
        pass

    try:
        ser_el = series._element
        tx_el = ser_el.find(f"{{{nsmap_c}}}tx")
        if tx_el is not None:
            # 常见结构：tx/strRef/strCache/pt/v
            for cache_tag in ("strCache", "numCache"):
                cache_el = tx_el.find(f".//{{{nsmap_c}}}{cache_tag}")
                if cache_el is None:
                    continue
                pts = sorted(
                    cache_el.findall(f"{{{nsmap_c}}}pt"),
                    key=lambda node: int(node.get("idx") or 0),
                )
                for pt in pts:
                    v_el = pt.find(f"{{{nsmap_c}}}v")
                    if v_el is not None and v_el.text:
                        v_text = v_el.text.strip()
                        if v_text:
                            return v_text

            # 兜底：tx/v
            v_el = tx_el.find(f"{{{nsmap_c}}}v")
            if v_el is not None and v_el.text:
                v_text = v_el.text.strip()
                if v_text:
                    return v_text
    except Exception:
        pass

    return default_name


def _extract_scatter_labels_from_series(series) -> List[str]:
    """从散点图 series XML 提取 xVal，作为前端 labels。"""
    nsmap_c = "http://schemas.openxmlformats.org/drawingml/2006/chart"
    try:
        ser_el = series._element
        x_val_el = ser_el.find(f"{{{nsmap_c}}}xVal")
        if x_val_el is None:
            return []

        cache = x_val_el.find(f".//{{{nsmap_c}}}numCache")
        if cache is None:
            cache = x_val_el.find(f".//{{{nsmap_c}}}strCache")
        if cache is None:
            cache = x_val_el.find(f"{{{nsmap_c}}}numLit")
        if cache is None:
            cache = x_val_el.find(f"{{{nsmap_c}}}strLit")
        if cache is None:
            return []

        pts = sorted(
            cache.findall(f"{{{nsmap_c}}}pt"),
            key=lambda node: int(node.get("idx") or 0),
        )
        labels: List[str] = []
        for pt in pts:
            v_el = pt.find(f"{{{nsmap_c}}}v")
            if v_el is not None and v_el.text is not None:
                labels.append(str(v_el.text))
        return labels
    except Exception:
        return []


def _chart_has_data_labels(chart, plot) -> bool:
    """判断图表是否开启数据标签（含 scatter 的 XML 兜底）。"""
    nsmap_c = "http://schemas.openxmlformats.org/drawingml/2006/chart"

    try:
        if hasattr(plot, "has_data_labels") and bool(plot.has_data_labels):
            return True
    except Exception:
        pass

    # 对 scatter 等 python-pptx API 不支持 has_data_labels 的图表，走 XML 兜底
    try:
        chart_xml = chart.element
        for d_lbls in chart_xml.iter(f"{{{nsmap_c}}}dLbls"):
            visible_toggles = (
                "showVal",
                "showPercent",
                "showCatName",
                "showSerName",
                "showBubbleSize",
                "showLegendKey",
            )
            has_explicit_toggle = False
            for toggle in visible_toggles:
                node = d_lbls.find(f"{{{nsmap_c}}}{toggle}")
                if node is None:
                    continue
                has_explicit_toggle = True
                raw = str(node.get("val") or "1").strip().lower()
                if raw not in ("0", "false"):
                    return True
            # 没有显式开关时按“开启标签”处理，贴合 Office 默认行为。
            if not has_explicit_toggle:
                return True
    except Exception:
        pass

    return False


def _extract_chart_area_fill_color(
    chart, theme_color_map: Optional[Dict[str, str]] = None,
) -> str:
    """提取 chartSpace 级背景色，避免误读系列颜色。"""
    nsmap_a = "http://schemas.openxmlformats.org/drawingml/2006/main"

    try:
        chart_xml = chart.element
        sp_pr = chart_xml.find(f"{{{nsmap_a}}}spPr")
        if sp_pr is None:
            return ""

        solid = sp_pr.find(f"{{{nsmap_a}}}solidFill")
        if solid is None:
            return ""

        srgb = solid.find(f"{{{nsmap_a}}}srgbClr")
        if srgb is not None:
            val = srgb.get("val")
            if val and len(val) == 6:
                return f"#{val}"

        scheme = solid.find(f"{{{nsmap_a}}}schemeClr")
        if scheme is not None and theme_color_map:
            scheme_to_key = {
                "dk1": "1",
                "lt1": "2",
                "dk2": "3",
                "lt2": "4",
                "accent1": "5",
                "accent2": "6",
                "accent3": "7",
                "accent4": "8",
                "accent5": "9",
                "accent6": "10",
            }
            tc_key = scheme_to_key.get((scheme.get("val") or "").strip())
            if tc_key and tc_key in theme_color_map:
                return theme_color_map[tc_key]
    except Exception:
        pass

    return ""


def _extract_chart_shape(
    shape, base: dict,
    theme_color_map: Optional[Dict[str, str]] = None,
) -> Dict[str, Any]:
    """提取图表 → ChartElement（完整类型映射 + 颜色/标题/选项提取）"""
    from pptx.enum.chart import XL_CHART_TYPE

    chart = shape.chart

    # ── 完整的图表类型映射 ──
    chart_type_map = {
        # 柱状图（垂直柱）
        XL_CHART_TYPE.COLUMN_CLUSTERED: "bar",
        XL_CHART_TYPE.COLUMN_STACKED: "bar",
        XL_CHART_TYPE.COLUMN_STACKED_100: "bar",
        XL_CHART_TYPE.THREE_D_COLUMN: "bar",
        XL_CHART_TYPE.THREE_D_COLUMN_CLUSTERED: "bar",
        XL_CHART_TYPE.THREE_D_COLUMN_STACKED: "bar",
        XL_CHART_TYPE.THREE_D_COLUMN_STACKED_100: "bar",
        # 3D 变体柱状图（锥/柱/棱锥 → 统一映射 bar）
        XL_CHART_TYPE.CONE_COL: "bar",
        XL_CHART_TYPE.CONE_COL_CLUSTERED: "bar",
        XL_CHART_TYPE.CONE_COL_STACKED: "bar",
        XL_CHART_TYPE.CONE_COL_STACKED_100: "bar",
        XL_CHART_TYPE.CYLINDER_COL: "bar",
        XL_CHART_TYPE.CYLINDER_COL_CLUSTERED: "bar",
        XL_CHART_TYPE.CYLINDER_COL_STACKED: "bar",
        XL_CHART_TYPE.CYLINDER_COL_STACKED_100: "bar",
        XL_CHART_TYPE.PYRAMID_COL: "bar",
        XL_CHART_TYPE.PYRAMID_COL_CLUSTERED: "bar",
        XL_CHART_TYPE.PYRAMID_COL_STACKED: "bar",
        XL_CHART_TYPE.PYRAMID_COL_STACKED_100: "bar",

        # 条形图（水平条）
        XL_CHART_TYPE.BAR_CLUSTERED: "column",
        XL_CHART_TYPE.BAR_STACKED: "column",
        XL_CHART_TYPE.BAR_STACKED_100: "column",
        XL_CHART_TYPE.THREE_D_BAR_CLUSTERED: "column",
        XL_CHART_TYPE.THREE_D_BAR_STACKED: "column",
        XL_CHART_TYPE.THREE_D_BAR_STACKED_100: "column",
        XL_CHART_TYPE.CONE_BAR_CLUSTERED: "column",
        XL_CHART_TYPE.CONE_BAR_STACKED: "column",
        XL_CHART_TYPE.CONE_BAR_STACKED_100: "column",
        XL_CHART_TYPE.CYLINDER_BAR_CLUSTERED: "column",
        XL_CHART_TYPE.CYLINDER_BAR_STACKED: "column",
        XL_CHART_TYPE.CYLINDER_BAR_STACKED_100: "column",
        XL_CHART_TYPE.PYRAMID_BAR_CLUSTERED: "column",
        XL_CHART_TYPE.PYRAMID_BAR_STACKED: "column",
        XL_CHART_TYPE.PYRAMID_BAR_STACKED_100: "column",

        # 折线图
        XL_CHART_TYPE.LINE: "line",
        XL_CHART_TYPE.LINE_MARKERS: "line",
        XL_CHART_TYPE.LINE_STACKED: "line",
        XL_CHART_TYPE.LINE_STACKED_100: "line",
        XL_CHART_TYPE.LINE_MARKERS_STACKED: "line",
        XL_CHART_TYPE.LINE_MARKERS_STACKED_100: "line",
        XL_CHART_TYPE.THREE_D_LINE: "line",

        # 面积图
        XL_CHART_TYPE.AREA: "area",
        XL_CHART_TYPE.AREA_STACKED: "area",
        XL_CHART_TYPE.AREA_STACKED_100: "area",
        XL_CHART_TYPE.THREE_D_AREA: "area",
        XL_CHART_TYPE.THREE_D_AREA_STACKED: "area",
        XL_CHART_TYPE.THREE_D_AREA_STACKED_100: "area",

        # 饼图
        XL_CHART_TYPE.PIE: "pie",
        XL_CHART_TYPE.PIE_EXPLODED: "pie",
        XL_CHART_TYPE.THREE_D_PIE: "pie",
        XL_CHART_TYPE.THREE_D_PIE_EXPLODED: "pie",
        XL_CHART_TYPE.PIE_OF_PIE: "pie",
        XL_CHART_TYPE.BAR_OF_PIE: "pie",

        # 环形图
        XL_CHART_TYPE.DOUGHNUT: "ring",
        XL_CHART_TYPE.DOUGHNUT_EXPLODED: "ring",

        # 雷达图
        XL_CHART_TYPE.RADAR: "radar",
        XL_CHART_TYPE.RADAR_FILLED: "radar",
        XL_CHART_TYPE.RADAR_MARKERS: "radar",

        # 散点图
        XL_CHART_TYPE.XY_SCATTER: "scatter",
        XL_CHART_TYPE.XY_SCATTER_LINES: "scatter",
        XL_CHART_TYPE.XY_SCATTER_LINES_NO_MARKERS: "scatter",
        XL_CHART_TYPE.XY_SCATTER_SMOOTH: "scatter",
        XL_CHART_TYPE.XY_SCATTER_SMOOTH_NO_MARKERS: "scatter",

        # 气泡图 → 降级为散点图
        XL_CHART_TYPE.BUBBLE: "scatter",
        XL_CHART_TYPE.BUBBLE_THREE_D_EFFECT: "scatter",
    }
    raw_chart_type = chart.chart_type
    chart_type = chart_type_map.get(raw_chart_type, "bar")
    if raw_chart_type not in chart_type_map:
        logger.info(f"Unmapped chart type {raw_chart_type} → fallback to 'bar'")

    # ── 检测是否为堆叠/平滑类型 ──
    _stacked_types = {
        XL_CHART_TYPE.COLUMN_STACKED, XL_CHART_TYPE.COLUMN_STACKED_100,
        XL_CHART_TYPE.BAR_STACKED, XL_CHART_TYPE.BAR_STACKED_100,
        XL_CHART_TYPE.THREE_D_COLUMN_STACKED, XL_CHART_TYPE.THREE_D_COLUMN_STACKED_100,
        XL_CHART_TYPE.THREE_D_BAR_STACKED, XL_CHART_TYPE.THREE_D_BAR_STACKED_100,
        XL_CHART_TYPE.LINE_STACKED, XL_CHART_TYPE.LINE_STACKED_100,
        XL_CHART_TYPE.LINE_MARKERS_STACKED, XL_CHART_TYPE.LINE_MARKERS_STACKED_100,
        XL_CHART_TYPE.AREA_STACKED, XL_CHART_TYPE.AREA_STACKED_100,
        XL_CHART_TYPE.THREE_D_AREA_STACKED, XL_CHART_TYPE.THREE_D_AREA_STACKED_100,
        # 锥/柱/棱锥堆叠变体
        XL_CHART_TYPE.CONE_COL_STACKED, XL_CHART_TYPE.CONE_COL_STACKED_100,
        XL_CHART_TYPE.CONE_BAR_STACKED, XL_CHART_TYPE.CONE_BAR_STACKED_100,
        XL_CHART_TYPE.CYLINDER_COL_STACKED, XL_CHART_TYPE.CYLINDER_COL_STACKED_100,
        XL_CHART_TYPE.CYLINDER_BAR_STACKED, XL_CHART_TYPE.CYLINDER_BAR_STACKED_100,
        XL_CHART_TYPE.PYRAMID_COL_STACKED, XL_CHART_TYPE.PYRAMID_COL_STACKED_100,
        XL_CHART_TYPE.PYRAMID_BAR_STACKED, XL_CHART_TYPE.PYRAMID_BAR_STACKED_100,
    }
    _smooth_types = {
        XL_CHART_TYPE.XY_SCATTER_SMOOTH, XL_CHART_TYPE.XY_SCATTER_SMOOTH_NO_MARKERS,
    }
    is_stacked = chart.chart_type in _stacked_types
    is_smooth = chart.chart_type in _smooth_types
    is_radar_filled = chart.chart_type == XL_CHART_TYPE.RADAR_FILLED
    if not is_smooth and chart_type in ("line", "area"):
        # 折线图平滑通常是 series.smooth 属性；面积图 AreaSeries 无此 API，需检查 XML
        try:
            plot = chart.plots[0]
            for series in plot.series:
                if hasattr(series, "smooth") and bool(series.smooth):
                    is_smooth = True
                    break
            if not is_smooth and chart_type == "area":
                _ns_c_smooth = "http://schemas.openxmlformats.org/drawingml/2006/chart"
                for ser_el in chart.element.iter(f"{{{_ns_c_smooth}}}ser"):
                    smooth_el = ser_el.find(f"{{{_ns_c_smooth}}}smooth")
                    if smooth_el is not None and smooth_el.get("val") in ("1", "true"):
                        is_smooth = True
                        break
        except Exception:
            pass

    # ── 提取数据 ──
    data: Dict[str, Any] = {"labels": [], "legends": [], "series": []}
    is_scatter_type = chart_type in ("scatter",)
    try:
        plot = chart.plots[0]
        scatter_x_series: List[List[str]] = []
        for idx, series in enumerate(plot.series):
            legend = _extract_chart_series_name(series, default_name=f"Series {idx + 1}")
            data["legends"].append(legend)
            raw_values = list(series.values) if series.values else []
            values = [v if v is not None and v == v else 0 for v in raw_values]
            data["series"].append(values)
            if is_scatter_type:
                x_vals = _extract_scatter_labels_from_series(series)
                scatter_x_series.append(x_vals if x_vals else [])
        # 散点图：从 series.xVal 提取真实 X 坐标作为 labels
        if is_scatter_type:
            normalized_x_series: List[List[str]] = []
            for idx, row in enumerate(scatter_x_series):
                series_len = len(data["series"][idx]) if idx < len(data["series"]) else len(row)
                normalized_row: List[str] = []
                for j in range(series_len):
                    raw_x = row[j] if j < len(row) else None
                    if raw_x is None or str(raw_x).strip() == "":
                        normalized_row.append(str(j + 1))
                    else:
                        normalized_row.append(str(raw_x))
                normalized_x_series.append(normalized_row)

            if normalized_x_series:
                data["xSeries"] = normalized_x_series
                data["labels"] = normalized_x_series[0]
        # 非散点图或散点图未取到 X 值时，使用 categories
        if not data["labels"]:
            if hasattr(plot, "categories") and plot.categories:
                data["labels"] = [str(c) for c in plot.categories]
        if is_scatter_type and not data["labels"]:
            max_len = max((len(row) for row in data["series"]), default=0)
            data["labels"] = [str(i + 1) for i in range(max_len)]
    except Exception as e:
        logger.warning(f"Chart data extraction failed: {e}")

    # ── 提取图表标题 ──
    title = ""
    try:
        if hasattr(chart, "has_title") and chart.has_title and chart.chart_title:
            title = chart.chart_title.text_frame.text
    except Exception:
        pass

    # ── 提取系列颜色（themeColors）+ 主题 key（themeColorKeys） ──
    theme_colors: List[str] = []
    theme_color_keys: List[Optional[str]] = []
    try:
        plot = chart.plots[0]
        if chart_type in ("pie", "ring"):
            # 饼图/环形图优先提取点级颜色，保证主题色映射按扇区保真。
            if len(plot.series) > 0:
                point_colors, point_theme_keys = _extract_pie_point_colors_and_theme_keys(plot.series[0], theme_color_map)
                if point_colors:
                    theme_colors.extend(point_colors)
                    if point_theme_keys:
                        theme_color_keys.extend(point_theme_keys[: len(point_colors)])
        else:
            for series in plot.series:
                color, theme_key = _extract_series_color_and_theme_key(series, theme_color_map)
                if color:
                    theme_colors.append(color)
                    theme_color_keys.append(theme_key)
    except Exception:
        pass

    # ── 构建 options ──
    options: Dict[str, Any] = {}
    if is_stacked:
        options["stack"] = True
    if is_smooth:
        options["lineSmooth"] = True
    if is_radar_filled:
        options["radarFilled"] = True

    # 图例位置和显示状态
    try:
        if hasattr(chart, "has_legend") and chart.has_legend:
            options["showLegend"] = True
            legend_obj = chart.legend
            if legend_obj is not None:
                # python-pptx Legend.position → XL_LEGEND_POSITION enum
                from pptx.enum.chart import XL_LEGEND_POSITION
                pos_map = {
                    XL_LEGEND_POSITION.BOTTOM: "b",
                    XL_LEGEND_POSITION.LEFT: "l",
                    XL_LEGEND_POSITION.RIGHT: "r",
                    XL_LEGEND_POSITION.TOP: "t",
                }
                try:
                    pos = legend_obj.position
                    if pos in pos_map:
                        options["legendPosition"] = pos_map[pos]
                except Exception:
                    pass
        else:
            options["showLegend"] = False
    except Exception:
        pass

    # 数据标签（plot API + XML 兜底）
    try:
        plot = chart.plots[0]
        show_data_label = _chart_has_data_labels(chart, plot)
        options["showDataLabel"] = show_data_label
    except Exception:
        pass

    # ── 提取样式属性（fill / textColor / gridColor） ──
    fill_color = ""
    text_color = ""
    grid_color = ""
    try:
        # 图表区域背景色（chartSpace/spPr）
        fill_color = _extract_chart_area_fill_color(chart, theme_color_map=theme_color_map)
    except Exception:
        pass

    _chart_scheme_to_key = {
        "dk1": "1", "lt1": "2", "dk2": "3", "lt2": "4",
        "accent1": "5", "accent2": "6", "accent3": "7",
        "accent4": "8", "accent5": "9", "accent6": "10",
        "tx1": "1", "bg1": "2", "tx2": "3", "bg2": "4",
    }

    try:
        # 坐标轴文字颜色（从第一个坐标轴的 txPr 提取，srgbClr 优先，schemeClr 兜底）
        ns_a = "http://schemas.openxmlformats.org/drawingml/2006/main"
        ns_c = "http://schemas.openxmlformats.org/drawingml/2006/chart"
        chart_xml = chart.element
        for ax in chart_xml.iter(f"{{{ns_c}}}catAx", f"{{{ns_c}}}valAx"):
            tx_pr = ax.find(f"{{{ns_c}}}txPr")
            if tx_pr is not None:
                for srgb in tx_pr.iter(f"{{{ns_a}}}srgbClr"):
                    val = srgb.get("val")
                    if val:
                        text_color = f"#{val}"
                        break
                if not text_color:
                    for _sc in tx_pr.iter(f"{{{ns_a}}}schemeClr"):
                        _sc_val = (_sc.get("val") or "").strip()
                        _tc_key = _chart_scheme_to_key.get(_sc_val)
                        if _tc_key and theme_color_map and _tc_key in theme_color_map:
                            text_color = theme_color_map[_tc_key]
                            break
            if text_color:
                break
    except Exception:
        pass

    try:
        # 网格线颜色（从第一个坐标轴的 majorGridlines 提取，srgbClr 优先，schemeClr 兜底）
        ns_a = "http://schemas.openxmlformats.org/drawingml/2006/main"
        ns_c = "http://schemas.openxmlformats.org/drawingml/2006/chart"
        chart_xml = chart.element
        for ax in chart_xml.iter(f"{{{ns_c}}}catAx", f"{{{ns_c}}}valAx"):
            grid = ax.find(f"{{{ns_c}}}majorGridlines")
            if grid is not None:
                sp_pr = grid.find(f"{{{ns_c}}}spPr")
                if sp_pr is not None:
                    for srgb in sp_pr.iter(f"{{{ns_a}}}srgbClr"):
                        val = srgb.get("val")
                        if val:
                            grid_color = f"#{val}"
                            break
                    if not grid_color:
                        for _sc in sp_pr.iter(f"{{{ns_a}}}schemeClr"):
                            _sc_val = (_sc.get("val") or "").strip()
                            _tc_key = _chart_scheme_to_key.get(_sc_val)
                            if _tc_key and theme_color_map and _tc_key in theme_color_map:
                                grid_color = theme_color_map[_tc_key]
                                break
            if grid_color:
                break
    except Exception:
        pass

    props: Dict[str, Any] = {
        "chartType": chart_type,
        "data": data,
    }
    if title:
        props["title"] = title
    if theme_colors:
        props["themeColors"] = theme_colors
    if theme_colors and theme_color_keys and any(key for key in theme_color_keys):
        aligned_theme_keys: List[Optional[str]] = []
        for idx in range(len(theme_colors)):
            aligned_theme_keys.append(theme_color_keys[idx] if idx < len(theme_color_keys) else None)
        props["themeColorKeys"] = aligned_theme_keys
    if options:
        props["options"] = options
    if fill_color:
        props["fill"] = fill_color
    if text_color:
        props["textColor"] = text_color
    if grid_color:
        props["gridColor"] = grid_color

    return {
        **base,
        "type": "chart",
        "props": props,
    }


def _extract_auto_shape(
    shape, base: dict,
    theme_color_map: Optional[Dict[str, str]] = None,
    theme_fonts: Optional[Dict[str, str]] = None,
    master_styles: Optional[Dict[str, Dict[str, Any]]] = None,
    placeholder_type: Optional[str] = None,
    placeholder_meta: Optional[Dict[str, Any]] = None,
    image_handler: Optional[Any] = None,
) -> Dict[str, Any]:
    """提取自动形状 → ShapeElement 或带形状的 TextElement（支持主题字体/样式继承）"""
    from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE

    pptx_shape_type = _get_auto_shape_type_name(shape)
    shape_formula = _map_shape_type_to_formula(pptx_shape_type)
    shape_keypoints = _extract_shape_keypoints(shape, pptx_shape_type, shape_formula)

    fill_color = _extract_shape_fill(shape, theme_color_map)
    fill_theme_key = _extract_shape_fill_theme_key(shape)
    gradient = _extract_shape_gradient(shape, theme_color_map)
    pattern = _extract_shape_pattern(shape, image_handler=image_handler, theme_color_map=theme_color_map)
    outline = _extract_shape_outline(shape, theme_color_map)

    # 获取 master 默认样式（根据 placeholder 类型）
    m_style = _resolve_master_style(master_styles or {}, placeholder_type)

    if shape.has_text_frame and shape.text_frame.text.strip():
        tf = shape.text_frame

        shape_text_font_family = None
        shape_text_font_size = None
        shape_text_color = None
        shape_text_color_theme_key = None
        for para in tf.paragraphs:
            for run in para.runs:
                font = run.font
                if font.name and not shape_text_font_family:
                    resolved = _resolve_theme_font_reference(font.name, theme_fonts) if font.name.startswith("+") else font.name
                    if resolved:
                        shape_text_font_family = resolved
                if font.size and not shape_text_font_size:
                    shape_text_font_size = round(font.size / 12700, 2)
                if not shape_text_color:
                    c = color_to_hex(font.color, theme_color_map) if font.color else None
                    if c:
                        shape_text_color = c
                        shape_text_color_theme_key = _extract_run_theme_color_key(run)
                if shape_text_font_family and shape_text_font_size and shape_text_color:
                    break
            if shape_text_font_family and shape_text_font_size and shape_text_color:
                break

        # 主题字体 fallback
        if not shape_text_font_family and theme_fonts:
            theme_font = _resolve_theme_font(theme_fonts, placeholder_type)
            if theme_font:
                shape_text_font_family = theme_font

        # Master 样式 fallback
        if not shape_text_font_size and m_style.get("fontSize"):
            shape_text_font_size = m_style["fontSize"]
        if not shape_text_color and m_style.get("color"):
            shape_text_color = m_style["color"]

        html_parts = []
        for para in tf.paragraphs:
            para_html, p_style = _paragraph_to_html_v2(
                para, theme_color_map,
                default_font_size=shape_text_font_size,
                default_font_family=shape_text_font_family,
                default_color=shape_text_color,
                theme_fonts=theme_fonts,
            )
            style_attr = f' style="{p_style}"' if p_style else ""
            html_parts.append(f"<p{style_attr}>{para_html}</p>")
        shape_text_content = "".join(html_parts) or tf.text.strip()

        # 提取段落对齐方式（扫描全部段落，取最常见的）
        shape_text_align = "left"
        try:
            align_counts: Dict[str, int] = {}
            align_order: List[str] = []

            def _has_visible_text(para) -> bool:
                try:
                    if para.text and para.text.strip():
                        return True
                except Exception:
                    pass
                try:
                    for run in para.runs:
                        if run.text and run.text.strip():
                            return True
                except Exception:
                    pass
                return False

            for p in tf.paragraphs:
                # 忽略空段落，避免被 PPT 默认段落对齐干扰
                if not _has_visible_text(p):
                    continue
                if p.alignment is not None:
                    a = _alignment_to_str(p.alignment, "left")
                else:
                    a = "left"
                if a not in align_counts:
                    align_order.append(a)
                align_counts[a] = align_counts.get(a, 0) + 1
            if align_counts:
                best_align = align_order[0]
                best_count = align_counts[best_align]
                for a in align_order[1:]:
                    count = align_counts[a]
                    if count > best_count:
                        best_align = a
                        best_count = count
                shape_text_align = best_align
            elif tf.paragraphs and tf.paragraphs[0].alignment is not None:
                shape_text_align = _alignment_to_str(tf.paragraphs[0].alignment, "left")
        except Exception:
            pass

        # 垂直对齐
        shape_vert_align = _extract_text_vertical_align(tf) or "top"

        text_props: Dict[str, Any] = {
            "content": shape_text_content,
            "align": shape_text_align,
            "verticalAlign": shape_vert_align,
        }
        if shape_text_color:
            text_props["defaultColor"] = shape_text_color
        if shape_text_color_theme_key:
            text_props["defaultColorThemeKey"] = shape_text_color_theme_key
        if shape_text_font_size:
            text_props["defaultFontSize"] = shape_text_font_size
        if shape_text_font_family:
            text_props["defaultFontFamily"] = _font_with_fallback(shape_text_font_family)

        _w = base.get("width", 100)
        _h = base.get("height", 100)
        svg_path_text, custgeom_no_fill = _extract_custgeom_path(shape, _w, _h)
        # 如果 custGeom 无路径，用预设形状生成兜底
        if not svg_path_text:
            svg_path_text = _generate_preset_shape_path(pptx_shape_type, _w, _h)
        if svg_path_text and not shape_keypoints:
            inferred_round_rect = _infer_round_rect_keypoints_from_svg_path(svg_path_text, float(_w), float(_h))
            if inferred_round_rect:
                shape_formula = "roundRect"
                shape_keypoints = inferred_round_rect
                if not pptx_shape_type:
                    pptx_shape_type = "roundRect"

        props: Dict[str, Any] = {
            "viewBox": [_w, _h],
            "path": svg_path_text or "",
            "text": text_props,
        }
        # 填充：无填充时显式写 "none"，区别于"主题色"
        # custGeom 子路径全部标记 fill="none" 时，也视为无填充（线条图标场景）
        if custgeom_no_fill:
            props["fill"] = "none"
        elif fill_color:
            props["fill"] = fill_color
            if fill_theme_key:
                props["fillThemeKey"] = fill_theme_key
                fill_theme_transforms = _extract_fill_theme_transforms(shape)
                if fill_theme_transforms:
                    props["fillThemeTransforms"] = fill_theme_transforms
        else:
            props["fill"] = "none"
        if pptx_shape_type:
            props["pptxShapeType"] = pptx_shape_type
        if shape_formula:
            props["pathFormula"] = shape_formula
        if shape_keypoints:
            props["keypoints"] = shape_keypoints
        if gradient:
            props["gradient"] = gradient
        if pattern:
            props["pattern"] = pattern
        if outline:
            props["outline"] = outline
        return {**base, "type": "shape", "props": props}

    # 纯形状
    _w = base.get("width", 100)
    _h = base.get("height", 100)
    svg_path, custgeom_no_fill = _extract_custgeom_path(shape, _w, _h)
    # 如果 custGeom 无路径，用预设形状生成兜底
    if not svg_path:
        svg_path = _generate_preset_shape_path(pptx_shape_type, _w, _h)
    if svg_path and not shape_keypoints:
        inferred_round_rect = _infer_round_rect_keypoints_from_svg_path(svg_path, float(_w), float(_h))
        if inferred_round_rect:
            shape_formula = "roundRect"
            shape_keypoints = inferred_round_rect
            if not pptx_shape_type:
                pptx_shape_type = "roundRect"

    props: Dict[str, Any] = {
        "viewBox": [_w, _h],
        "path": svg_path or "",
    }
    # custGeom 子路径全部标记 fill="none" 时，也视为无填充（线条图标场景）
    if custgeom_no_fill:
        props["fill"] = "none"
    elif fill_color:
        props["fill"] = fill_color
        if fill_theme_key:
            props["fillThemeKey"] = fill_theme_key
            fill_theme_transforms = _extract_fill_theme_transforms(shape)
            if fill_theme_transforms:
                props["fillThemeTransforms"] = fill_theme_transforms
    else:
        props["fill"] = "none"
    if pptx_shape_type:
        props["pptxShapeType"] = pptx_shape_type
    if shape_formula:
        props["pathFormula"] = shape_formula
    if shape_keypoints:
        props["keypoints"] = shape_keypoints
    if gradient:
        props["gradient"] = gradient
    if pattern:
        props["pattern"] = pattern
    if outline:
        props["outline"] = outline

    return {
        **base,
        "type": "shape",
        "props": props,
    }


def _get_auto_shape_type_name(shape) -> Optional[str]:
    """获取 auto shape 类型名称，用于前端 pathFormula 映射"""
    try:
        if hasattr(shape, "auto_shape_type") and shape.auto_shape_type is not None:
            ast = shape.auto_shape_type

            enum_name = None
            if hasattr(ast, "name") and ast.name:
                enum_name = str(ast.name).upper()
            elif isinstance(ast, str):
                enum_name = ast.upper()

            # 优先走枚举名映射，避免不同 python-pptx 版本中数值漂移导致错配
            if enum_name:
                name_map = {
                    # 矩形类
                    "RECTANGLE": "rect",
                    "ROUNDED_RECTANGLE": "roundRect",
                    "ROUND_1_RECTANGLE": "round1Rect",
                    "ROUND_2_SAME_RECTANGLE": "round2SameRect",
                    "ROUND_2_DIAG_RECTANGLE": "round2DiagRect",
                    "SNIP_2_DIAG_RECTANGLE": "snip2DiagRect",
                    "SNIP_ROUND_RECTANGLE": "snipRoundRect",
                    # 基础几何
                    "OVAL": "ellipse",
                    "DIAMOND": "diamond",
                    "ISOSCELES_TRIANGLE": "triangle",
                    "RIGHT_TRIANGLE": "rtTriangle",
                    "PARALLELOGRAM": "parallelogram",
                    "TRAPEZOID": "trapezoid",
                    "PENTAGON": "pentagon",
                    "HEXAGON": "hexagon",
                    "OCTAGON": "octagon",
                    # 星形
                    "STAR_4_POINT": "star4",
                    "STAR_5_POINT": "star5",
                    "STAR_6_POINT": "star6",
                    # 箭头
                    "RIGHT_ARROW": "rightArrow",
                    "LEFT_ARROW": "leftArrow",
                    "UP_ARROW": "upArrow",
                    "DOWN_ARROW": "downArrow",
                    "LEFT_RIGHT_ARROW": "leftRightArrow",
                    "UP_DOWN_ARROW": "upDownArrow",
                    "NOTCHED_RIGHT_ARROW": "notchedRightArrow",
                    # 其他
                    "HEART": "heart",
                    "LIGHTNING_BOLT": "lightningBolt",
                    "CROSS": "plus",
                    "CHEVRON": "chevron",
                    "CLOUD": "cloud",
                    "RECTANGULAR_CALLOUT": "callout1",
                    "ROUNDED_RECTANGULAR_CALLOUT": "callout2",
                }
                if enum_name in name_map:
                    return name_map[enum_name]
                return _snake_to_camel(enum_name.lower())
    except Exception as e:
        logger.debug(f"Failed to get auto_shape_type: {e}")
    return None


def _snake_to_camel(name: str) -> str:
    """snake_case → camelCase"""
    parts = name.split("_")
    return parts[0] + "".join(p.capitalize() for p in parts[1:])


_SHAPE_TYPE_TO_FORMULA: Dict[str, str] = {
    "rect": "rect",
    "roundRect": "roundRect",
    "round1Rect": "roundRect",
    "round2SameRect": "roundRect",
    "round2DiagRect": "roundRect",
    "snipRoundRect": "cutRect",
    "snip2DiagRect": "cutRect",
    "ellipse": "ellipse",
    "triangle": "triangle",
    "rightTriangle": "rtTriangle",
    "rtTriangle": "rtTriangle",
    "diamond": "diamond",
    "parallelogram": "parallelogram",
    "trapezoid": "trapezoid",
    "pentagon": "pentagon",
    "hexagon": "hexagon",
    "octagon": "octagon",
    "star4": "star4",
    "star5": "star5",
    "star6": "star6",
    "star6Point": "star6",
    "rightArrow": "rightArrow",
    "leftArrow": "leftArrow",
    "upArrow": "upArrow",
    "downArrow": "downArrow",
    "leftRightArrow": "leftRightArrow",
    "upDownArrow": "upDownArrow",
    "notchedRightArrow": "notchedRightArrow",
    "heart": "heart",
    "lightningBolt": "lightningBolt",
    "cloud": "cloud",
    "chevron": "chevron",
    "callout1": "callout1",
    "callout2": "callout2",
    "plus": "cross",
    "cross": "cross",
}

_FORMULA_KEYPOINT_META: Dict[str, Dict[str, Any]] = {
    "roundRect": {"names": ["adj"], "ranges": [(0.0, 0.5)]},
    "roundRectSingle": {"names": ["adj"], "ranges": [(0.0, 0.5)]},
    "cutRect": {"names": ["adj"], "ranges": [(0.0, 0.5)]},
    "triangle": {"names": ["adj"], "ranges": [(0.0, 1.0)]},
    "parallelogram": {"names": ["adj"], "ranges": [(0.0, 0.5)]},
    "trapezoid": {"names": ["adj"], "ranges": [(0.0, 0.5)]},
    "star4": {"names": ["adj"], "ranges": [(0.1, 0.9)]},
    "star5": {"names": ["adj"], "ranges": [(0.1, 0.9)]},
    "star6": {"names": ["adj"], "ranges": [(0.1, 0.9)]},
    # 与前端 ShapePathFormulas.range 保持一致，避免回读到前端出现“当前值低于可编辑下限”。
    "cross": {"names": ["adj"], "ranges": [(0.1, 0.5)]},
    "chevron": {"names": ["adj"], "ranges": [(0.05, 0.5)]},
    "rightArrow": {"names": ["adj1", "adj2"], "ranges": [(0.0, 1.0), (0.0, 0.5)]},
    "leftArrow": {"names": ["adj1", "adj2"], "ranges": [(0.0, 1.0), (0.0, 0.5)]},
    "upArrow": {"names": ["adj1", "adj2"], "ranges": [(0.0, 1.0), (0.0, 0.5)]},
    "downArrow": {"names": ["adj1", "adj2"], "ranges": [(0.0, 1.0), (0.0, 0.5)]},
    "notchedRightArrow": {"names": ["adj1", "adj2"], "ranges": [(0.0, 1.0), (0.0, 0.5)]},
    # 与前端 ShapePathFormulas 保持一致：adj1 允许 0~0.5，避免编辑后回读被后端窄范围裁剪。
    "leftRightArrow": {"names": ["adj1", "adj2"], "ranges": [(0.0, 0.5), (0.0, 0.5)]},
    "upDownArrow": {"names": ["adj1", "adj2"], "ranges": [(0.0, 0.5), (0.0, 0.5)]},
}


def _map_shape_type_to_formula(shape_type: Optional[str]) -> Optional[str]:
    if not shape_type:
        return None
    raw = str(shape_type).strip()
    return _SHAPE_TYPE_TO_FORMULA.get(raw)


_ROUND_RECT_PRESET_TYPES: Set[str] = {"roundRect", "round1Rect", "round2SameRect", "round2DiagRect"}


def _clamp_round_rect_ratio(value: Any, fallback: float = 0.1) -> float:
    try:
        val = float(value)
    except Exception:
        val = float(fallback)
    if not math.isfinite(val):
        val = float(fallback)
    return max(0.0, min(0.5, val))


def _normalize_round_rect_keypoints(raw: Any, fallback: float = 0.1) -> Optional[List[float]]:
    if not isinstance(raw, list) or not raw:
        return None

    fb = _clamp_round_rect_ratio(fallback, 0.1)
    first = _clamp_round_rect_ratio(raw[0], fb)
    if len(raw) == 1:
        return [first, first, first, first]

    second = _clamp_round_rect_ratio(raw[1], first)
    if len(raw) == 2:
        return [first, second, first, second]

    third = _clamp_round_rect_ratio(raw[2], first)
    if len(raw) == 3:
        return [first, second, third, second]

    fourth = _clamp_round_rect_ratio(raw[3], first)
    return [first, second, third, fourth]


def _round_rect_keypoints_uniform(keypoints: List[float], tol: float = 1e-4) -> bool:
    if len(keypoints) < 4:
        return False
    first = keypoints[0]
    return all(abs(v - first) <= tol for v in keypoints[1:4])


def _extract_round_rect_keypoints(shape, shape_type: Optional[str]) -> Optional[List[float]]:
    """提取 round 系列预设形状的四角圆角比例（左上/右上/右下/左下）。"""
    try:
        nsmap_a = "http://schemas.openxmlformats.org/drawingml/2006/main"
        sp_pr = _find_sp_pr(shape._element)
        if sp_pr is None:
            return None
        prst_geom = sp_pr.find(f"{{{nsmap_a}}}prstGeom")
        if prst_geom is None:
            return None

        prst = (prst_geom.get("prst") or (shape_type or "")).strip()
        if prst not in _ROUND_RECT_PRESET_TYPES:
            return None

        ratio = 16667 / 100000.0  # PowerPoint roundRect 默认 16.667%
        av_lst = prst_geom.find(f"{{{nsmap_a}}}avLst")
        if av_lst is not None:
            fallback_val: Optional[int] = None
            for gd in av_lst.findall(f"{{{nsmap_a}}}gd"):
                fmla = gd.get("fmla") or ""
                parts = fmla.split()
                if not parts:
                    continue
                try:
                    parsed = int(parts[-1])
                except Exception:
                    continue
                if fallback_val is None:
                    fallback_val = parsed
                name = (gd.get("name") or "").strip().lower()
                if name in {"adj", "adj1"}:
                    ratio = parsed / 100000.0
                    break
            else:
                if fallback_val is not None:
                    ratio = fallback_val / 100000.0

        ratio = _clamp_round_rect_ratio(ratio, 16667 / 100000.0)
        if prst == "roundRect":
            return [round(ratio, 4)] * 4
        if prst == "round1Rect":
            return [round(ratio, 4), 0.0, 0.0, 0.0]
        if prst == "round2DiagRect":
            return [round(ratio, 4), 0.0, round(ratio, 4), 0.0]
        if prst == "round2SameRect":
            return [round(ratio, 4), 0.0, 0.0, round(ratio, 4)]
    except Exception:
        return None
    return None


def _extract_shape_keypoints(shape, shape_type: Optional[str], formula: Optional[str] = None) -> Optional[List[float]]:
    """从 prstGeom/avLst/gd 提取可调关键点，统一为 0~1（按公式范围裁剪）。"""
    formula_name = formula or _map_shape_type_to_formula(shape_type)
    if not formula_name:
        return None
    if formula_name == "roundRect":
        round_rect_kp = _extract_round_rect_keypoints(shape, shape_type)
        if round_rect_kp:
            return round_rect_kp
    meta = _FORMULA_KEYPOINT_META.get(formula_name)
    if not meta:
        return None

    try:
        nsmap_a = "http://schemas.openxmlformats.org/drawingml/2006/main"
        sp_pr = _find_sp_pr(shape._element)
        if sp_pr is None:
            return None
        prst_geom = sp_pr.find(f"{{{nsmap_a}}}prstGeom")
        if prst_geom is None:
            return None
        av_lst = prst_geom.find(f"{{{nsmap_a}}}avLst")
        if av_lst is None:
            return None

        gd_nodes = av_lst.findall(f"{{{nsmap_a}}}gd")
        if not gd_nodes:
            return None

        gd_values: Dict[str, int] = {}
        ordered_vals: List[int] = []
        for gd in gd_nodes:
            name = gd.get("name") or ""
            fmla = gd.get("fmla") or ""
            raw_val: Optional[int] = None
            parts = fmla.split()
            if parts:
                try:
                    raw_val = int(parts[-1])
                except Exception:
                    raw_val = None
            if raw_val is None:
                continue
            gd_values[name] = raw_val
            ordered_vals.append(raw_val)

        if not ordered_vals:
            return None

        keypoints: List[float] = []
        fallback_idx = 0
        names: List[str] = meta.get("names", [])
        ranges: List[Tuple[float, float]] = meta.get("ranges", [])

        for idx, gd_name in enumerate(names):
            raw_val = gd_values.get(gd_name)
            if raw_val is None and gd_name == "adj1":
                raw_val = gd_values.get("adj")
            if raw_val is None and fallback_idx < len(ordered_vals):
                raw_val = ordered_vals[fallback_idx]
                fallback_idx += 1
            if raw_val is None:
                break

            ratio = raw_val / 100000.0
            min_v, max_v = ranges[idx] if idx < len(ranges) else (0.0, 1.0)
            ratio = max(min_v, min(max_v, ratio))
            keypoints.append(round(ratio, 4))

        return keypoints or None
    except Exception:
        return None


def _generate_preset_shape_path(shape_type: Optional[str], w: float, h: float) -> Optional[str]:
    """
    根据预设形状类型生成 SVG path，作为前端渲染的兜底保障。

    当 shape 使用 <a:prstGeom>（预设几何）时，不会有 <a:custGeom> 节点，
    此函数确保后端总是返回有效的 path，即使前端 pathFormula 缺失也能正确显示。
    """
    if not shape_type:
        return None

    import math

    generators = {
        "rect": lambda: f"M 0 0 L {w} 0 L {w} {h} L 0 {h} Z",
        "roundRect": lambda: _gen_round_rect_path(w, h, 0.1),
        "round1Rect": lambda: _gen_round_1_rect_path(w, h, 0.2),
        "round2SameRect": lambda: _gen_round_rect_4_path(w, h, [0.1, 0.0, 0.0, 0.1]),
        "round2DiagRect": lambda: _gen_round_rect_4_path(w, h, [0.1, 0.0, 0.1, 0.0]),
        "snip2DiagRect": lambda: _gen_snip_2_diag_rect_path(w, h, 0.1),
        "snipRoundRect": lambda: _gen_snip_2_diag_rect_path(w, h, 0.1),
        "ellipse": lambda: _gen_ellipse_path(w, h),
        "diamond": lambda: f"M {w/2} 0 L {w} {h/2} L {w/2} {h} L 0 {h/2} Z",
        "triangle": lambda: f"M {w/2} 0 L {w} {h} L 0 {h} Z",
        "rtTriangle": lambda: f"M 0 {h} L {w} {h} L 0 0 Z",
        "rightTriangle": lambda: f"M 0 {h} L {w} {h} L 0 0 Z",
        "parallelogram": lambda: f"M {w*0.25} 0 L {w} 0 L {w*0.75} {h} L 0 {h} Z",
        "trapezoid": lambda: f"M {w*0.2} 0 L {w*0.8} 0 L {w} {h} L 0 {h} Z",
        "pentagon": lambda: _gen_regular_polygon_path(w, h, 5),
        "hexagon": lambda: _gen_regular_polygon_path(w, h, 6),
        "octagon": lambda: _gen_octagon_path(w, h, 0.3),
        "star5": lambda: _gen_star_path(w, h, 5, 0.38),
        "star4": lambda: _gen_star_path(w, h, 4, 0.38),
        "star6": lambda: _gen_star_path(w, h, 6, 0.4),
        "star6Point": lambda: _gen_star_path(w, h, 6, 0.4),
        "rightArrow": lambda: f"M 0 {h*0.3} L {w*0.5} {h*0.3} L {w*0.5} 0 L {w} {h/2} L {w*0.5} {h} L {w*0.5} {h*0.7} L 0 {h*0.7} Z",
        "leftArrow": lambda: f"M {w} {h*0.3} L {w*0.5} {h*0.3} L {w*0.5} 0 L 0 {h/2} L {w*0.5} {h} L {w*0.5} {h*0.7} L {w} {h*0.7} Z",
        "upArrow": lambda: f"M {w/2} 0 L {w} {h*0.5} L {w*0.7} {h*0.5} L {w*0.7} {h} L {w*0.3} {h} L {w*0.3} {h*0.5} L 0 {h*0.5} Z",
        "downArrow": lambda: f"M {w*0.3} 0 L {w*0.7} 0 L {w*0.7} {h*0.5} L {w} {h*0.5} L {w/2} {h} L 0 {h*0.5} L {w*0.3} {h*0.5} Z",
        "leftRightArrow": lambda: (
            f"M 0 {h/2} L {w*0.2} 0 L {w*0.2} {h*0.3} L {w*0.8} {h*0.3} L {w*0.8} 0 "
            f"L {w} {h/2} L {w*0.8} {h} L {w*0.8} {h*0.7} L {w*0.2} {h*0.7} L {w*0.2} {h} Z"
        ),
        "upDownArrow": lambda: (
            f"M {w/2} 0 L {w} {h*0.2} L {w*0.7} {h*0.2} L {w*0.7} {h*0.8} L {w} {h*0.8} "
            f"L {w/2} {h} L 0 {h*0.8} L {w*0.3} {h*0.8} L {w*0.3} {h*0.2} L 0 {h*0.2} Z"
        ),
        "notchedRightArrow": lambda: (
            f"M 0 {h*0.3} L {w*0.5} {h*0.3} L {w*0.5} 0 L {w} {h/2} "
            f"L {w*0.5} {h} L {w*0.5} {h*0.7} L 0 {h*0.7} L {h*0.12} {h/2} Z"
        ),
        "heart": lambda: (
            f"M {w/2} {h*0.25} C {w/2} {h*0.1} {w*0.25} 0 {w*0.1} 0 "
            f"C 0 0 0 {h*0.15} 0 {h*0.3} C 0 {h*0.55} {w/2} {h*0.7} {w/2} {h} "
            f"C {w/2} {h*0.7} {w} {h*0.55} {w} {h*0.3} "
            f"C {w} {h*0.15} {w} 0 {w*0.9} 0 C {w*0.75} 0 {w/2} {h*0.1} {w/2} {h*0.25} Z"
        ),
        "lightningBolt": lambda: (
            f"M {w*0.37} 0 L {w*0.63} 0 L {w*0.52} {h*0.35} L {w*0.75} {h*0.35} "
            f"L {w*0.3} {h} L {w*0.42} {h*0.5} L {w*0.2} {h*0.5} Z"
        ),
        "plus": lambda: _gen_cross_path(w, h, 0.3),
        "cloud": lambda: (
            f"M {w*0.25} {h*0.7} C {w*0.05} {h*0.7} 0 {h*0.55} {w*0.08} {h*0.4} "
            f"C {w*0.02} {h*0.25} {w*0.15} {h*0.1} {w*0.3} {h*0.15} "
            f"C {w*0.35} {h*0.02} {w*0.5} 0 {w*0.6} {h*0.08} "
            f"C {w*0.7} 0 {w*0.85} {h*0.05} {w*0.88} {h*0.22} "
            f"C {w} {h*0.25} {w} {h*0.45} {w*0.9} {h*0.55} "
            f"C {w*0.98} {h*0.65} {w*0.9} {h*0.8} {w*0.75} {h*0.8} "
            f"L {w*0.25} {h*0.8} C {w*0.1} {h*0.8} {w*0.05} {h*0.78} {w*0.25} {h*0.7} Z"
        ),
        "chevron": lambda: (
            f"M 0 0 L {w*0.75} 0 L {w} {h/2} L {w*0.75} {h} "
            f"L 0 {h} L {w*0.25} {h/2} Z"
        ),
        "callout1": lambda: (
            f"M 0 0 L {w} 0 L {w} {h*0.7} L 0 {h*0.7} Z "
            f"M {w*0.2} {h*0.7} L {w*0.15} {h} L {w*0.35} {h*0.7}"
        ),
        "callout2": lambda: (
            f"M {w*0.05} 0 L {w*0.95} 0 Q {w} 0 {w} {h*0.05} L {w} {h*0.65} "
            f"Q {w} {h*0.7} {w*0.95} {h*0.7} L {w*0.35} {h*0.7} L {w*0.15} {h} "
            f"L {w*0.2} {h*0.7} L {w*0.05} {h*0.7} Q 0 {h*0.7} 0 {h*0.65} "
            f"L 0 {h*0.05} Q 0 0 {w*0.05} 0 Z"
        ),
    }

    gen = generators.get(shape_type)
    if gen:
        try:
            return gen()
        except Exception:
            return None
    return None


def _gen_ellipse_path(w: float, h: float) -> str:
    """生成椭圆 SVG path（四段三次贝塞尔曲线近似）"""
    rx = w / 2
    ry = h / 2
    cx = rx
    cy = ry
    k = 0.5522847498
    kx = rx * k
    ky = ry * k
    return (
        f"M {cx} 0 "
        f"C {cx + kx} 0 {w} {cy - ky} {w} {cy} "
        f"C {w} {cy + ky} {cx + kx} {h} {cx} {h} "
        f"C {cx - kx} {h} 0 {cy + ky} 0 {cy} "
        f"C 0 {cy - ky} {cx - kx} 0 {cx} 0 Z"
    )


def _gen_round_rect_path(w: float, h: float, ratio: float = 0.1) -> str:
    """生成圆角矩形 SVG path"""
    r = _clamp_round_rect_ratio(ratio, 0.1)
    return _gen_round_rect_4_path(w, h, [r, r, r, r])


def _gen_round_rect_4_path(w: float, h: float, ratios: List[float]) -> str:
    """按四角比例生成圆角矩形 SVG path（左上/右上/右下/左下）。"""
    short = max(1e-6, min(w, h))
    kps = _normalize_round_rect_keypoints(ratios, 0.1) or [0.1, 0.1, 0.1, 0.1]
    rtl = short * kps[0]
    rtr = short * kps[1]
    rbr = short * kps[2]
    rbl = short * kps[3]
    return (
        f"M {rtl} 0 L {w - rtr} 0 Q {w} 0 {w} {rtr} "
        f"L {w} {h - rbr} Q {w} {h} {w - rbr} {h} "
        f"L {rbl} {h} Q 0 {h} 0 {h - rbl} "
        f"L 0 {rtl} Q 0 0 {rtl} 0 Z"
    )


def _gen_round_1_rect_path(w: float, h: float, ratio: float = 0.2) -> str:
    """生成单圆角矩形（左上圆角）"""
    r = min(w, h) * ratio
    return f"M {r} 0 L {w} 0 L {w} {h} L 0 {h} L 0 {r} Q 0 0 {r} 0 Z"


def _gen_snip_2_diag_rect_path(w: float, h: float, ratio: float = 0.1) -> str:
    """生成双对角剪角矩形（左上 + 右下）"""
    c = min(w, h) * ratio
    return f"M {c} 0 L {w} 0 L {w} {h - c} L {w - c} {h} L 0 {h} L 0 {c} Z"


def _gen_octagon_path(w: float, h: float, ratio: float = 0.3) -> str:
    """生成八边形"""
    c = min(w, h) * ratio
    return (
        f"M {c} 0 L {w-c} 0 L {w} {c} L {w} {h-c} "
        f"L {w-c} {h} L {c} {h} L 0 {h-c} L 0 {c} Z"
    )


def _gen_regular_polygon_path(w: float, h: float, sides: int) -> str:
    """生成正多边形 SVG path"""
    import math
    cx = w / 2
    cy = h / 2
    r = min(w, h) / 2
    points = []
    for i in range(sides):
        angle = (i * 360 / sides - 90) * math.pi / 180
        points.append(f"{round(cx + r * math.cos(angle), 1)} {round(cy + r * math.sin(angle), 1)}")
    return f"M {' L '.join(points)} Z"


def _gen_star_path(w: float, h: float, points_count: int, inner_ratio: float) -> str:
    """生成星形 SVG path"""
    import math
    cx = w / 2
    cy = h / 2
    outer_r = min(w, h) / 2
    inner_r = outer_r * inner_ratio
    pts = []
    angle_step = 360 / points_count
    for i in range(points_count):
        outer_angle = (i * angle_step - 90) * math.pi / 180
        pts.append(f"{round(cx + outer_r * math.cos(outer_angle), 1)} {round(cy + outer_r * math.sin(outer_angle), 1)}")
        inner_angle = ((i * angle_step + angle_step / 2) - 90) * math.pi / 180
        pts.append(f"{round(cx + inner_r * math.cos(inner_angle), 1)} {round(cy + inner_r * math.sin(inner_angle), 1)}")
    return f"M {' L '.join(pts)} Z"


def _gen_cross_path(w: float, h: float, ratio: float = 0.3) -> str:
    """生成十字形 SVG path"""
    arm = min(w, h) * ratio
    cx = w / 2
    cy = h / 2
    return (
        f"M {cx - arm} 0 L {cx + arm} 0 L {cx + arm} {cy - arm} L {w} {cy - arm} "
        f"L {w} {cy + arm} L {cx + arm} {cy + arm} L {cx + arm} {h} L {cx - arm} {h} "
        f"L {cx - arm} {cy + arm} L 0 {cy + arm} L 0 {cy - arm} L {cx - arm} {cy - arm} Z"
    )


def _extract_custgeom_path(shape, width_px: float, height_px: float) -> Tuple[Optional[str], bool]:
    """
    从 shape 的 <a:custGeom> 中提取自定义几何路径，转为 SVG path d 属性。

    返回 (path_d_string, all_paths_no_fill):
      - path_d_string: SVG path d 属性字符串，失败时为 None
      - all_paths_no_fill: 当所有 <a:path> 都标记 fill="none" 时为 True，
        表示该形状路径仅用于描边（线条图标常见场景）

    OOXML custGeom path 命令：
      <a:moveTo>   → M x y
      <a:lnTo>     → L x y
      <a:cubicBezTo> → C x1 y1 x2 y2 x y
      <a:quadBezTo>  → Q x1 y1 x y
      <a:arcTo>    → SVG arc 命令
      <a:close>    → Z
    """
    try:
        import math
        nsmap = {"a": "http://schemas.openxmlformats.org/drawingml/2006/main"}
        el = shape._element

        cust_geom = el.find(".//a:custGeom", nsmap)
        if cust_geom is None:
            return None, False

        path_lst = cust_geom.find("a:pathLst", nsmap)
        if path_lst is None:
            return None, False

        paths = path_lst.findall("a:path", nsmap)
        if not paths:
            return None, False

        all_d = []
        has_any_filled_path = False
        for path_el in paths:
            # 检查子路径的 fill 属性：fill="none" 表示仅描边
            path_fill_attr = (path_el.get("fill") or "").strip().lower()
            if path_fill_attr != "none":
                has_any_filled_path = True

            # custGeom path 有自己的坐标空间（w/h 属性定义）
            try:
                path_w = int(float(path_el.get("w", "0"))) or 1
                path_h = int(float(path_el.get("h", "0"))) or 1
            except (ValueError, TypeError):
                path_w = 1
                path_h = 1

            # 缩放因子：将 custGeom 坐标映射到实际 px 尺寸
            sx = width_px / path_w if path_w else 1
            sy = height_px / path_h if path_h else 1

            # 维护当前点，用于 arcTo 端点计算
            cur_x, cur_y = 0.0, 0.0

            d_parts = []
            for cmd in path_el:
                tag = cmd.tag.split("}")[-1] if "}" in cmd.tag else cmd.tag

                if tag == "moveTo":
                    pt = cmd.find("a:pt", nsmap)
                    if pt is not None:
                        x = round(float(pt.get("x", "0")) * sx, 1)
                        y = round(float(pt.get("y", "0")) * sy, 1)
                        d_parts.append(f"M {x} {y}")
                        cur_x, cur_y = x, y

                elif tag == "lnTo":
                    pt = cmd.find("a:pt", nsmap)
                    if pt is not None:
                        x = round(float(pt.get("x", "0")) * sx, 1)
                        y = round(float(pt.get("y", "0")) * sy, 1)
                        d_parts.append(f"L {x} {y}")
                        cur_x, cur_y = x, y

                elif tag == "cubicBezTo":
                    pts = cmd.findall("a:pt", nsmap)
                    if len(pts) >= 3:
                        coords = []
                        for pt in pts[:3]:
                            coords.append(round(float(pt.get("x", "0")) * sx, 1))
                            coords.append(round(float(pt.get("y", "0")) * sy, 1))
                        d_parts.append(f"C {coords[0]} {coords[1]} {coords[2]} {coords[3]} {coords[4]} {coords[5]}")
                        cur_x, cur_y = coords[4], coords[5]

                elif tag == "quadBezTo":
                    pts = cmd.findall("a:pt", nsmap)
                    if len(pts) >= 2:
                        coords = []
                        for pt in pts[:2]:
                            coords.append(round(float(pt.get("x", "0")) * sx, 1))
                            coords.append(round(float(pt.get("y", "0")) * sy, 1))
                        d_parts.append(f"Q {coords[0]} {coords[1]} {coords[2]} {coords[3]}")
                        cur_x, cur_y = coords[2], coords[3]

                elif tag == "arcTo":
                    try:
                        wr = float(cmd.get("wR", "0")) * sx
                        hr = float(cmd.get("hR", "0")) * sy
                        st_ang = float(cmd.get("stAng", "0")) / 60000  # OOXML 60000ths of a degree
                        sw_ang = float(cmd.get("swAng", "0")) / 60000
                        if wr > 0 and hr > 0 and abs(sw_ang) > 0:
                            st_ang_rad = math.radians(st_ang)
                            end_ang_rad = math.radians(st_ang + sw_ang)
                            # OOXML arcTo: 弧心 = 当前点 - (wr*cos(stAng), hr*sin(stAng))
                            cx = cur_x - wr * math.cos(st_ang_rad)
                            cy = cur_y - hr * math.sin(st_ang_rad)
                            # 终点 = 弧心 + (wr*cos(endAng), hr*sin(endAng))
                            end_x = round(cx + wr * math.cos(end_ang_rad), 1)
                            end_y = round(cy + hr * math.sin(end_ang_rad), 1)
                            large_arc = 1 if abs(sw_ang) > 180 else 0
                            sweep = 1 if sw_ang > 0 else 0
                            d_parts.append(
                                f"A {round(wr, 1)} {round(hr, 1)} 0 {large_arc} {sweep} {end_x} {end_y}"
                            )
                            cur_x, cur_y = end_x, end_y
                    except (ValueError, TypeError):
                        pass

                elif tag == "close":
                    d_parts.append("Z")

            if d_parts:
                all_d.append(" ".join(d_parts))

        all_paths_no_fill = bool(all_d) and not has_any_filled_path
        path_str = " ".join(all_d) if all_d else None
        return path_str, all_paths_no_fill

    except Exception as e:
        logger.warning(f"Failed to extract custGeom path: {e}")
        return None, False


def _parse_simple_svg_path(path_d: str) -> Optional[List[Tuple[str, List[float]]]]:
    """解析简化 SVG path（仅 M/L/Q/Z 绝对命令）。"""
    if not path_d:
        return None
    tokens = re.findall(r"[A-Za-z]|-?\d+(?:\.\d+)?", path_d)
    if not tokens:
        return None

    result: List[Tuple[str, List[float]]] = []
    idx = 0
    while idx < len(tokens):
        token = tokens[idx]
        if not re.match(r"^[A-Za-z]$", token):
            return None
        cmd = token.upper()
        idx += 1
        if cmd in {"M", "L"}:
            if idx + 1 >= len(tokens):
                return None
            try:
                x = float(tokens[idx])
                y = float(tokens[idx + 1])
            except Exception:
                return None
            idx += 2
            result.append((cmd, [x, y]))
            continue
        if cmd == "Q":
            if idx + 3 >= len(tokens):
                return None
            try:
                vals = [float(tokens[idx]), float(tokens[idx + 1]), float(tokens[idx + 2]), float(tokens[idx + 3])]
            except Exception:
                return None
            idx += 4
            result.append((cmd, vals))
            continue
        if cmd == "Z":
            result.append((cmd, []))
            continue
        return None
    return result


def _infer_round_rect_keypoints_from_svg_path(path_d: Optional[str], width_px: float, height_px: float) -> Optional[List[float]]:
    """
    从 custGeom 回读得到的 SVG path 推断四角圆角比例（左上/右上/右下/左下）。
    仅识别本模块写出的 roundRect 四角路径模板。
    """
    if not path_d or width_px <= 0 or height_px <= 0:
        return None

    commands = _parse_simple_svg_path(path_d)
    if not commands:
        return None

    expected = ["M", "L", "Q", "L", "Q", "L", "Q", "L", "Q", "Z"]
    if len(commands) != len(expected):
        return None
    if [item[0] for item in commands] != expected:
        return None

    tol = max(1.5, min(width_px, height_px) * 0.02)
    near = lambda a, b: abs(a - b) <= tol  # noqa: E731

    m_x, m_y = commands[0][1]
    l1_x, l1_y = commands[1][1]
    q1_cx, q1_cy, q1_x, q1_y = commands[2][1]
    l2_x, l2_y = commands[3][1]
    q2_cx, q2_cy, q2_x, q2_y = commands[4][1]
    l3_x, l3_y = commands[5][1]
    q3_cx, q3_cy, q3_x, q3_y = commands[6][1]
    l4_x, l4_y = commands[7][1]
    q4_cx, q4_cy, q4_x, q4_y = commands[8][1]

    if not (
        near(m_y, 0)
        and near(l1_y, 0)
        and near(q1_cx, width_px)
        and near(q1_cy, 0)
        and near(q1_x, width_px)
        and near(l2_x, width_px)
        and near(q2_cx, width_px)
        and near(q2_cy, height_px)
        and near(q2_y, height_px)
        and near(l3_y, height_px)
        and near(q3_cx, 0)
        and near(q3_cy, height_px)
        and near(q3_x, 0)
        and near(l4_x, 0)
        and near(q4_cx, 0)
        and near(q4_cy, 0)
        and near(q4_y, 0)
    ):
        return None

    r_tl_candidates = [m_x, q4_x, l4_y]
    if max(r_tl_candidates) - min(r_tl_candidates) > tol:
        return None
    r_tl = sum(r_tl_candidates) / len(r_tl_candidates)

    r_tr_candidates = [q1_y, width_px - l1_x]
    if max(r_tr_candidates) - min(r_tr_candidates) > tol:
        return None
    r_tr = sum(r_tr_candidates) / len(r_tr_candidates)

    r_br_candidates = [height_px - l2_y, width_px - q2_x]
    if max(r_br_candidates) - min(r_br_candidates) > tol:
        return None
    r_br = sum(r_br_candidates) / len(r_br_candidates)

    r_bl_candidates = [l3_x, height_px - q3_y]
    if max(r_bl_candidates) - min(r_bl_candidates) > tol:
        return None
    r_bl = sum(r_bl_candidates) / len(r_bl_candidates)

    short_side = max(1e-6, min(width_px, height_px))
    radii = [r_tl, r_tr, r_br, r_bl]
    ratios: List[float] = []
    for radius in radii:
        if radius < tol:
            ratios.append(0.0)
        else:
            ratios.append(round(_clamp_round_rect_ratio(radius / short_side, 0.0), 4))

    if not any(v > 0 for v in ratios):
        return None
    return ratios


def _extract_shape_fill_theme_key(shape) -> Optional[str]:
    """提取形状纯色填充的主题色 key（仅 schemeClr + 可选 alpha）。"""
    try:
        nsmap_a = "http://schemas.openxmlformats.org/drawingml/2006/main"
        sp_pr = _find_sp_pr(shape._element)
        if sp_pr is not None:
            solid_fill = sp_pr.find(f"{{{nsmap_a}}}solidFill")
            if solid_fill is not None:
                key = _extract_theme_key_from_solid_fill(solid_fill)
                if key:
                    return key
                # 存在非纯主题变换（如 tint/shade）时仍返回 key（变换信息由 _extract_fill_theme_transforms 提取）
                scheme_el = solid_fill.find(f"{{{nsmap_a}}}schemeClr")
                if scheme_el is not None:
                    return _normalize_text_theme_key(scheme_el.get("val"))
    except Exception:
        pass

    try:
        fill = getattr(shape, "fill", None)
        if fill and fill.type is not None and fill.type == 1:  # SOLID
            return _extract_theme_key_from_color_obj(fill.fore_color)
    except Exception:
        pass
    return None


def _extract_fill_theme_transforms(shape) -> Optional[Dict[str, float]]:
    """
    Extract color transforms (tint/shade/lumMod/lumOff) from shape's solid fill.
    Returns a dict of transform names → values (0-1 scale), or None if no transforms.
    """
    try:
        nsmap_a = "http://schemas.openxmlformats.org/drawingml/2006/main"
        sp_pr = _find_sp_pr(shape._element)
        if sp_pr is None:
            return None
        solid_fill = sp_pr.find(f"{{{nsmap_a}}}solidFill")
        if solid_fill is None:
            return None
        scheme = solid_fill.find(f"{{{nsmap_a}}}schemeClr")
        if scheme is None:
            return None

        transforms = {}
        for child in list(scheme):
            tag = child.tag.split("}")[-1] if isinstance(child.tag, str) and "}" in child.tag else str(child.tag)
            if tag == "alpha":
                continue  # alpha handled separately
            val_str = child.get("val")
            if val_str:
                try:
                    transforms[tag] = int(val_str) / 100000.0
                except (ValueError, TypeError):
                    pass

        return transforms if transforms else None
    except Exception:
        return None


def _extract_shape_fill(shape, theme_color_map: Optional[Dict[str, str]] = None) -> Optional[str]:
    """
    提取形状的填充颜色。
    透明度语义：
    - 填充透明度保留在 fill 颜色本身（rgba）；
    - 元素整体透明度由 _extract_opacity（spPr 直属 alphaModFix）单独承载。
    """
    try:
        # 先检查 XML 级别的 <a:noFill>，优先级高于 python-pptx 的主题继承
        nsmap_a = "http://schemas.openxmlformats.org/drawingml/2006/main"
        sp_pr = _find_sp_pr(shape._element)
        if sp_pr is not None:
            no_fill = sp_pr.find(f"{{{nsmap_a}}}noFill")
            if no_fill is not None:
                return None  # 显式 noFill → 无填充

        fill = getattr(shape, "fill", None)
        if fill and fill.type is not None:
            if fill.type == 1:  # SOLID
                hex_color = color_to_hex(fill.fore_color, theme_color_map)
                if hex_color:
                    fill_alpha = None
                    if sp_pr is not None:
                        solid_fill = sp_pr.find(f"{{{nsmap_a}}}solidFill")
                        fill_alpha = _extract_solid_fill_alpha(solid_fill, decimals=2)
                    if fill_alpha is not None and fill_alpha < 1.0:
                        return _hex_to_rgba(hex_color, fill_alpha)
                    return hex_color

        # 兜底：直接解析 XML（兼容 DummyShape / API 无法读取 fill 的场景）
        if sp_pr is not None:
            solid_fill = sp_pr.find(f"{{{nsmap_a}}}solidFill")
            if solid_fill is not None:
                hex_color = _extract_color_with_transforms(solid_fill, theme_color_map)
                if hex_color:
                    fill_alpha = _extract_solid_fill_alpha(solid_fill, decimals=2)
                    if fill_alpha is not None and fill_alpha < 1.0:
                        return _hex_to_rgba(hex_color, fill_alpha)
                    return hex_color
    except Exception:
        pass
    return None


def _extract_fill_alpha(shape) -> Optional[float]:
    """从形状的 solidFill 颜色中提取 alpha 值"""
    try:
        nsmap_a = "http://schemas.openxmlformats.org/drawingml/2006/main"
        spPr = _find_sp_pr(shape._element)
        if spPr is None:
            return None
        solid_fill = spPr.find(f"{{{nsmap_a}}}solidFill")
        if solid_fill is not None:
            return _extract_solid_fill_alpha(solid_fill, decimals=2)
    except Exception:
        pass
    return None


def _extract_shape_gradient(shape, theme_color_map: Optional[Dict[str, str]] = None) -> Optional[Dict[str, Any]]:
    """提取形状的渐变填充"""
    try:
        fill = shape.fill
        fill_type = getattr(fill, "type", None) if fill else None

        is_gradient_fill = False
        if fill_type is not None:
            # 兼容 python-pptx 版本差异：GRADIENT 可能是 2 或 3，或枚举对象
            try:
                fill_type_int = int(fill_type)
                is_gradient_fill = fill_type_int in (2, 3)
            except Exception:
                pass

            if not is_gradient_fill:
                fill_type_name = ""
                try:
                    fill_type_name = getattr(fill_type, "name", "") or str(fill_type)
                except Exception:
                    fill_type_name = str(fill_type)
                is_gradient_fill = "GRADIENT" in fill_type_name.upper()

        if fill and is_gradient_fill:
            # 检测渐变类型（线性 vs 径向）
            grad_type = "linear"
            _shape_grad_center = None
            try:
                nsmap_a_gt = "http://schemas.openxmlformats.org/drawingml/2006/main"
                sp_pr = _find_sp_pr(shape._element)
                if sp_pr is not None:
                    grad_fill = sp_pr.find(f"{{{nsmap_a_gt}}}gradFill")
                    if grad_fill is not None:
                        # 检查是否有 <a:path> 节点（径向/矩形渐变）
                        path_node = grad_fill.find(f"{{{nsmap_a_gt}}}path")
                        if path_node is not None:
                            grad_type = "radial"
                            # 提取径向渐变中心点
                            _ft = path_node.find(f"{{{nsmap_a_gt}}}fillToRect")
                            if _ft is not None:
                                try:
                                    _shape_cx = int(_ft.get("l", "50000")) / 100000
                                    _shape_cy = int(_ft.get("t", "50000")) / 100000
                                    _shape_grad_center = {"x": round(_shape_cx, 4), "y": round(_shape_cy, 4)}
                                except (ValueError, TypeError):
                                    _shape_grad_center = None
                            else:
                                _shape_grad_center = None
            except Exception:
                pass

            gradient: Dict[str, Any] = {
                "type": grad_type,
                "rotate": 0,
                "colors": [],
            }
            # 径向渐变中心点（从上方 path_node 解析）
            if grad_type == "radial" and _shape_grad_center is not None:
                gradient["center"] = _shape_grad_center
            # 尝试获取渐变角度（仅线性有效）
            try:
                if hasattr(fill, "gradient_angle") and fill.gradient_angle is not None:
                    gradient["rotate"] = fill.gradient_angle
            except Exception:
                pass

            # 获取渐变停靠点（支持 per-stop alpha）
            try:
                nsmap_a = "http://schemas.openxmlformats.org/drawingml/2006/main"
                for stop in fill.gradient_stops:
                    stop_color = color_to_hex(stop.color, theme_color_map) or "#000000"
                    stop_pos = stop.position if hasattr(stop, "position") else 0
                    # 从 XML 提取 stop 颜色的 alpha
                    try:
                        gs_el = stop._gs if hasattr(stop, "_gs") else None
                        if gs_el is not None:
                            for child in gs_el:
                                alpha_el = child.find(f"{{{nsmap_a}}}alpha")
                                if alpha_el is not None:
                                    a_val = alpha_el.get("val")
                                    if a_val:
                                        a_float = round(int(a_val) / 100000, 2)
                                        if a_float < 1.0:
                                            stop_color = _hex_to_rgba(stop_color, a_float)
                                    break
                    except Exception:
                        pass
                    gradient["colors"].append({
                        "pos": round(stop_pos, 3),
                        "color": stop_color,
                    })
            except Exception:
                pass

            if gradient["colors"]:
                return gradient
    except Exception:
        pass
    return None


def _extract_shape_pattern(
    shape,
    image_handler: Optional[Any] = None,
    theme_color_map: Optional[Dict[str, str]] = None,
) -> Optional[str]:
    """提取形状填充图案（图片填充 blipFill / 预设图案 pattFill）。"""
    import base64 as b64

    try:
        nsmap_a = "http://schemas.openxmlformats.org/drawingml/2006/main"
        nsmap_p = "http://schemas.openxmlformats.org/presentationml/2006/main"
        nsmap_r = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"

        sp_pr = _find_sp_pr(shape._element)
        if sp_pr is None:
            return None

        # 1) 图片填充
        blip_fill = sp_pr.find(f"{{{nsmap_a}}}blipFill")
        if blip_fill is not None:
            blip = blip_fill.find(f"{{{nsmap_a}}}blip")
            if blip is not None:
                r_embed = blip.get(f"{{{nsmap_r}}}embed")
                if not r_embed:
                    # 链接型图片填充（少见）
                    r_link = blip.get(f"{{{nsmap_r}}}link")
                    if r_link:
                        return r_link
                else:
                    image_part = None
                    if hasattr(shape, "part") and shape.part is not None:
                        part = shape.part

                        # 路径 1：python-pptx 通用 API
                        if hasattr(part, "related_part"):
                            try:
                                image_part = part.related_part(r_embed)
                            except Exception:
                                image_part = None

                        # 路径 2：关系表兜底（兼容部分版本对象结构）
                        if image_part is None and hasattr(part, "rels"):
                            try:
                                rel = part.rels.get(r_embed) if hasattr(part.rels, "get") else part.rels[r_embed]
                                if rel is not None:
                                    image_part = getattr(rel, "target_part", None) or getattr(rel, "_target", None)
                            except Exception:
                                image_part = None

                        # 路径 3：按 relationship id 取关联 part
                        if image_part is None and hasattr(part, "part_related_by"):
                            try:
                                image_part = part.part_related_by(r_embed)
                            except Exception:
                                image_part = None

                    if image_part is not None and hasattr(image_part, "blob"):
                        blob_data = image_part.blob
                        content_type = getattr(image_part, "content_type", "image/png")
                        src = ""

                        if image_handler and callable(image_handler):
                            try:
                                uploaded_url = image_handler(blob_data, content_type)
                                if uploaded_url:
                                    src = uploaded_url
                            except Exception as e:
                                logger.warning(f"shape pattern image_handler failed, fallback to base64: {e}")

                        if not src:
                            encoded = b64.b64encode(blob_data).decode("ascii")
                            src = f"data:{content_type};base64,{encoded}"
                        if src:
                            return src

        # 2) 预设图案填充（pattFill）→ 生成 SVG data URL 近似还原
        patt_fill = sp_pr.find(f"{{{nsmap_a}}}pattFill")
        if patt_fill is not None:
            prst = patt_fill.get("prst", "pct10")
            fg = _extract_pattfill_color(patt_fill.find(f"{{{nsmap_a}}}fgClr"), theme_color_map) or "#666666"
            bg = _extract_pattfill_color(patt_fill.find(f"{{{nsmap_a}}}bgClr"), theme_color_map) or "#ffffff"
            return _build_pattfill_data_url(prst, fg, bg)
    except Exception:
        return None
    return None


def _extract_pattfill_color(
    color_node,
    theme_color_map: Optional[Dict[str, str]] = None,
) -> Optional[str]:
    """从 pattFill 的 fgClr/bgClr 节点解析颜色。"""
    if color_node is None:
        return None

    nsmap_a = "http://schemas.openxmlformats.org/drawingml/2006/main"
    try:
        srgb = color_node.find(f"{{{nsmap_a}}}srgbClr")
        if srgb is not None and srgb.get("val"):
            return f"#{srgb.get('val')}"

        scheme = color_node.find(f"{{{nsmap_a}}}schemeClr")
        if scheme is not None and theme_color_map:
            scheme_name = scheme.get("val", "")
            scheme_to_key = {
                "dk1": "1", "lt1": "2", "dk2": "3", "lt2": "4",
                "accent1": "5", "accent2": "6", "accent3": "7",
                "accent4": "8", "accent5": "9", "accent6": "10",
            }
            tc_key = scheme_to_key.get(scheme_name)
            if tc_key and tc_key in theme_color_map:
                return theme_color_map[tc_key]
    except Exception:
        return None
    return None


def _build_pattfill_data_url(prst: str, fg_color: str, bg_color: str) -> str:
    """将 OOXML pattFill 近似转换成可渲染的 PNG data URL。"""
    import base64 as b64
    import io

    try:
        from PIL import Image, ImageDraw, ImageColor
    except Exception:
        # PIL 不可用时回退纯色 png（透明）
        return "data:image/png;base64,"

    prst_l = (prst or "").lower()
    tile = 12

    def _safe_color(val: str, fallback: str) -> tuple:
        try:
            return ImageColor.getrgb(val)
        except Exception:
            return ImageColor.getrgb(fallback)

    fg = _safe_color(fg_color, "#666666")
    bg = _safe_color(bg_color, "#ffffff")

    img = Image.new("RGB", (tile, tile), bg)
    draw = ImageDraw.Draw(img)

    if "diag" in prst_l:
        draw.line([(0, tile - 1), (tile - 1, 0)], fill=fg, width=1)
        draw.line([(-tile // 3, tile - 1), (tile - 1, tile // 3)], fill=fg, width=1)
        draw.line([(tile // 3, tile - 1), (tile - 1, -tile // 3)], fill=fg, width=1)
    elif "cross" in prst_l:
        draw.line([(0, tile // 2), (tile - 1, tile // 2)], fill=fg, width=1)
        draw.line([(tile // 2, 0), (tile // 2, tile - 1)], fill=fg, width=1)
    else:
        # pct*/dot* 等统一做点阵近似
        dot_r = 1
        draw.ellipse((2 - dot_r, 2 - dot_r, 2 + dot_r, 2 + dot_r), fill=fg)
        draw.ellipse((tile - 3 - dot_r, tile - 3 - dot_r, tile - 3 + dot_r, tile - 3 + dot_r), fill=fg)

    buf = io.BytesIO()
    img.save(buf, format="PNG")
    encoded = b64.b64encode(buf.getvalue()).decode("ascii")
    return f"data:image/png;base64,{encoded}"


def _extract_run_color_alpha(run) -> Optional[float]:
    """
    从文本 run 的 <a:rPr><a:solidFill>/<a:srgbClr>/<a:alpha> 提取 alpha。
    适用于 run.font.color 有 alpha 的场景。
    """
    try:
        nsmap_a = "http://schemas.openxmlformats.org/drawingml/2006/main"
        rPr = run._r.find(f"{{{nsmap_a}}}rPr") if hasattr(run, "_r") else None
        if rPr is None:
            return None
        solid = rPr.find(f"{{{nsmap_a}}}solidFill")
        if solid is None:
            return None
        for color_el in solid:
            alpha = _extract_alpha_from_color_transforms(color_el, decimals=2)
            if alpha is not None:
                return alpha
    except Exception:
        pass
    return None


def _extract_line_color_alpha(shape) -> Optional[float]:
    """从 <a:ln><a:solidFill> 下的颜色元素中提取 alpha"""
    try:
        nsmap_a = "http://schemas.openxmlformats.org/drawingml/2006/main"
        spPr = _find_sp_pr(shape._element)
        if spPr is None:
            return None
        ln = spPr.find(f"{{{nsmap_a}}}ln")
        if ln is None:
            return None
        solid = ln.find(f"{{{nsmap_a}}}solidFill")
        if solid is None:
            return None
        return _extract_solid_fill_alpha(solid, decimals=2)
    except Exception:
        pass
    return None


def _extract_shape_outline(shape, theme_color_map: Optional[Dict[str, str]] = None) -> Optional[Dict[str, Any]]:
    """
    提取形状的轮廓/边框。

    使用双重策略：
    1. 先尝试 python-pptx 高级 API
    2. 失败时直接解析 XML（<a:ln>）— 覆盖主题色引用等 API 未暴露的场景
    """
    outline: Dict[str, Any] = {}

    # ── 策略 1: python-pptx 高级 API ──
    try:
        if hasattr(shape, "line") and shape.line:
            line = shape.line
            try:
                if line.width and line.width > 0:
                    outline["width"] = round(line.width / EMU_PER_PT, 1)
            except Exception:
                pass

            try:
                if line.color and line.color.rgb:
                    hex_c = f"#{line.color.rgb}"
                    # 尝试从 XML 读取 line color 的 alpha
                    ln_alpha = _extract_line_color_alpha(shape)
                    if ln_alpha is not None and ln_alpha < 1.0:
                        outline["color"] = _hex_to_rgba(hex_c, ln_alpha)
                    else:
                        outline["color"] = hex_c
            except Exception:
                pass

            try:
                if line.color and line.color.theme_color is not None:
                    theme_key = _theme_enum_to_text_theme_key(line.color.theme_color)
                    if theme_key:
                        outline["themeKey"] = theme_key
            except Exception:
                pass

            try:
                if not outline.get("color") and line.color and line.color.theme_color is not None:
                    tc_key = str(int(line.color.theme_color))
                    if theme_color_map and tc_key in theme_color_map:
                        hex_c = theme_color_map[tc_key]
                        ln_alpha = _extract_line_color_alpha(shape)
                        if ln_alpha is not None and ln_alpha < 1.0:
                            outline["color"] = _hex_to_rgba(hex_c, ln_alpha)
                        else:
                            outline["color"] = hex_c
            except Exception:
                pass

            try:
                if line.dash_style:
                    dash_val = int(line.dash_style)
                    dash_map = {1: "solid", 2: "dashed", 3: "dotted", 4: "dotted", 5: "dashDot", 6: "longDash", 7: "longDashDot", 8: "longDashDot"}
                    outline["style"] = dash_map.get(dash_val, "solid")
            except Exception:
                pass
    except Exception:
        pass

    # ── 策略 2: 直接解析 XML <a:ln>（兜底 + dash 样式校正）──
    try:
        nsmap = {"a": "http://schemas.openxmlformats.org/drawingml/2006/main"}
        spPr = _find_sp_pr(shape._element)
        if spPr is not None:
            ln = spPr.find("a:ln", nsmap)
            if ln is not None:
                # 宽度
                w_attr = ln.get("w")
                if w_attr and not outline.get("width"):
                    outline["width"] = round(int(w_attr) / EMU_PER_PT, 1)

                solid = ln.find("a:solidFill", nsmap)
                if solid is not None:
                    theme_key = _extract_theme_key_from_solid_fill(solid)
                    if theme_key:
                        outline["themeKey"] = theme_key
                    elif solid.find("a:schemeClr", nsmap) is not None:
                        # 存在 tint/shade 等主题变换时不应保留纯 key 语义
                        outline.pop("themeKey", None)

                # 颜色: <a:solidFill><a:srgbClr val="..."/>
                # 同时提取颜色的独立 alpha
                if not outline.get("color"):
                    if solid is not None:
                        srgb = solid.find("a:srgbClr", nsmap)
                        if srgb is not None:
                            hex_val = f"#{srgb.get('val', '000000')}"
                            alpha_float = _extract_alpha_from_color_transforms(srgb, decimals=2)
                            if alpha_float is not None and alpha_float < 1.0:
                                outline["color"] = _hex_to_rgba(hex_val, alpha_float)
                            else:
                                outline["color"] = hex_val
                        else:
                            # 主题色 <a:schemeClr val="dk1"/>
                            scheme = solid.find("a:schemeClr", nsmap)
                            if scheme is not None and theme_color_map:
                                scheme_name = scheme.get("val", "")
                                scheme_to_key = {
                                    "dk1": "1", "lt1": "2", "dk2": "3", "lt2": "4",
                                    "accent1": "5", "accent2": "6", "accent3": "7",
                                    "accent4": "8", "accent5": "9", "accent6": "10",
                                }
                                tc_key = scheme_to_key.get(scheme_name)
                                if tc_key and tc_key in theme_color_map:
                                    hex_val = theme_color_map[tc_key]
                                    alpha_float = _extract_alpha_from_color_transforms(scheme, decimals=2)
                                    if alpha_float is not None and alpha_float < 1.0:
                                        outline["color"] = _hex_to_rgba(hex_val, alpha_float)
                                    else:
                                        outline["color"] = hex_val

                # 虚线样式（优先 XML，覆盖 API 误判为 solid 的场景）
                dash_val_str = ln.get("prstDash")
                if not dash_val_str:
                    prst_dash = ln.find("a:prstDash", nsmap)
                    if prst_dash is not None:
                        dash_val_str = prst_dash.get("val", "solid")
                if dash_val_str:
                    dash_str_map = {
                        "solid": "solid", "dash": "dashed", "dot": "dotted",
                        "lgDash": "longDash", "dashDot": "dashDot",
                        "sysDash": "dashed", "sysDot": "dotted",
                        "lgDashDot": "longDashDot", "lgDashDotDot": "longDashDot",
                    }
                    xml_style = dash_str_map.get(dash_val_str, "solid")
                    if xml_style != "solid" or not outline.get("style"):
                        outline["style"] = xml_style

                # 如果 <a:ln> 存在但有 <a:noFill>，说明显式无描边
                no_fill = ln.find("a:noFill", nsmap)
                if no_fill is not None:
                    return None
    except Exception:
        pass

    if outline.get("width") or outline.get("color"):
        outline.setdefault("style", "solid")
        outline.setdefault("width", 1)
        outline.setdefault("color", "#333333")
        return outline
    return None


def _extract_shadow(shape) -> Optional[Dict[str, Any]]:
    """提取形状的阴影效果（兼容 sp / cxnSp / pic / graphicFrame）"""
    try:
        # python-pptx ShadowFormat 可能对连接器/其他类型不可用，先尝试高级 API
        try:
            shadow = shape.shadow
            if shadow is not None and hasattr(shadow, "inherit") and shadow.inherit:
                return None
        except (AttributeError, TypeError):
            pass

        result: Dict[str, Any] = {}

        # 通过 XML 直接提取（更可靠，兼容所有元素类型）
        from lxml import etree
        nsmap = {"a": "http://schemas.openxmlformats.org/drawingml/2006/main"}
        el = shape._element

        # 查找 <a:effectLst><a:outerShdw ...>
        outer_shdw = el.find(".//a:effectLst/a:outerShdw", nsmap)
        if outer_shdw is None:
            return None

        # 模糊半径（EMU → px）
        blur_rad = outer_shdw.get("blurRad")
        if blur_rad:
            result["blur"] = round(int(blur_rad) / EMU_PER_PT, 1)

        # 偏移距离和角度
        dist = outer_shdw.get("dist")
        dir_val = outer_shdw.get("dir")  # 方向，单位是 60000 分之一度
        if dist:
            dist_px = round(int(dist) / EMU_PER_PT, 1)
            if dir_val:
                import math
                angle_deg = int(dir_val) / 60000
                angle_rad = math.radians(angle_deg)
                result["h"] = round(dist_px * math.cos(angle_rad), 1)
                result["v"] = round(dist_px * math.sin(angle_rad), 1)
            else:
                result["h"] = dist_px
                result["v"] = dist_px

        # 阴影颜色
        srgb = outer_shdw.find("a:srgbClr", nsmap)
        if srgb is not None:
            result["color"] = f"#{srgb.get('val', '000000')}"
            alpha = srgb.find("a:alpha", nsmap)
            if alpha is not None:
                result["opacity"] = round(int(alpha.get("val", "100000")) / 100000, 2)
        else:
            result["color"] = "#000000"

        if result:
            result.setdefault("blur", 4)
            result.setdefault("h", 2)
            result.setdefault("v", 2)
            result.setdefault("color", "#000000")
            result.setdefault("opacity", 0.3)
            return result

    except Exception:
        pass
    return None


def _extract_opacity(shape) -> Optional[float]:
    """
    提取元素整体不透明度（0.0 ~ 1.0）。

    仅读取 spPr 直属 alphaModFix，避免把 fill/outline 自身透明度误当成整体透明度。
    """
    try:
        nsmap = {
            "a": "http://schemas.openxmlformats.org/drawingml/2006/main",
            "p": "http://schemas.openxmlformats.org/presentationml/2006/main",
        }
        # 只在 spPr 下搜索，避免误取阴影/渐变的 alpha
        spPr = shape._element.find(".//p:spPr", nsmap)
        if spPr is None:
            spPr = shape._element.find(".//a:spPr", nsmap)
        if spPr is None:
            return None

        # 检查 spPr 直属 alphaModFix（形状整体透明度）
        # 注意：不要用 ".//a:alphaModFix"，否则会误命中图片 blipFill 下的 alpha，
        # 导致在 _extract_image_shape() 中再次叠加，出现透明度平方问题。
        alpha_mod = spPr.find("a:alphaModFix", nsmap)
        if alpha_mod is not None:
            parsed_alpha = _parse_ooxml_percentage_value(alpha_mod.get("val") or alpha_mod.get("amt"))
            if parsed_alpha is not None:
                return round(max(0.0, min(1.0, parsed_alpha)), 4)
    except Exception:
        pass
    return None


def _extract_background_from_element(bg_owner, theme_color_map: Optional[Dict[str, str]] = None, image_handler: Optional[Any] = None) -> Optional[Dict[str, Any]]:
    """
    从一个具有 .background 属性的对象（slide / slide_layout / slide_master）提取背景。
    返回 None 表示没有显式背景。
    """
    import base64 as b64

    try:
        try:
            from pptx.enum.dml import MSO_FILL  # type: ignore
            fill_solid = MSO_FILL.SOLID
            fill_gradient = MSO_FILL.GRADIENT
            fill_picture = MSO_FILL.PICTURE
        except Exception:
            # 兜底值，兼容 python-pptx 版本差异
            fill_solid = 1
            fill_gradient = 3
            fill_picture = 6

        bg = bg_owner.background

        # theme 引用背景（p:bgRef），常见于模板页/母版背景继承场景
        # 注意：必须在访问 bg.fill 之前读取 XML；某些 python-pptx 访问路径会将 bgRef 规范化为 bgPr/noFill。
        try:
            nsmap_p = "http://schemas.openxmlformats.org/presentationml/2006/main"
            nsmap_a = "http://schemas.openxmlformats.org/drawingml/2006/main"
            bg_elem = bg_owner.background._element if hasattr(bg_owner.background, "_element") else None
            if bg_elem is not None:
                bg_ref = bg_elem.find(f".//{{{nsmap_p}}}bgRef")
                if bg_ref is not None:
                    theme_key = None
                    color = None
                    scheme_el = bg_ref.find(f"{{{nsmap_a}}}schemeClr")
                    if scheme_el is not None:
                        theme_key = (scheme_el.get("val") or "").strip() or None
                        if theme_key:
                            color = _resolve_scheme_color(theme_key, theme_color_map)
                    if not color:
                        color = _extract_color_with_transforms(bg_ref, theme_color_map)
                    if color:
                        if scheme_el is not None:
                            alpha_el = scheme_el.find(f"{{{nsmap_a}}}alpha")
                            if alpha_el is not None and alpha_el.get("val"):
                                alpha = round(int(alpha_el.get("val")) / 100000, 2)
                                if alpha < 1.0:
                                    color = _hex_to_rgba(color, alpha)
                        result = {"type": "color", "value": color}
                        if theme_key:
                            theme_obj_ref: Dict[str, Any] = {"key": theme_key, "color": color}
                            # 保留 bgRef 中 schemeClr 上的变换参数
                            if scheme_el is not None:
                                transforms_ref: Dict[str, int] = {}
                                for tag_name in ("lumMod", "lumOff", "tint", "shade", "satMod"):
                                    tr_el = scheme_el.find(f"{{{nsmap_a}}}{tag_name}")
                                    if tr_el is not None and tr_el.get("val"):
                                        try:
                                            transforms_ref[tag_name] = int(tr_el.get("val"))
                                        except (ValueError, TypeError):
                                            pass
                                if transforms_ref:
                                    theme_obj_ref["transforms"] = transforms_ref
                            result["theme"] = theme_obj_ref
                        return result
        except Exception as e:
            logger.debug(f"Failed to extract background bgRef: {e}")

        if bg.fill and bg.fill.type is not None:
            fill = bg.fill

            if fill.type == fill_solid:  # SOLID
                color = color_to_hex(fill.fore_color, theme_color_map)
                theme_key = None
                # 检查背景 fill 的 alpha
                try:
                    nsmap_a_bg = "http://schemas.openxmlformats.org/drawingml/2006/main"
                    bg_elem = bg_owner.background._element if hasattr(bg_owner.background, "_element") else None
                    if bg_elem is not None:
                        solid_fill_el = bg_elem.find(f".//{{{nsmap_a_bg}}}solidFill")
                        if solid_fill_el is not None and len(solid_fill_el) > 0:
                            transformed_color = _extract_color_with_transforms(solid_fill_el, theme_color_map)
                            if transformed_color:
                                color = transformed_color
                            scheme_el = solid_fill_el.find(f"{{{nsmap_a_bg}}}schemeClr")
                            if scheme_el is not None and scheme_el.get("val"):
                                theme_key = scheme_el.get("val")
                                if not color:
                                    color = _resolve_scheme_color(theme_key, theme_color_map)
                            alpha_el = solid_fill_el[0].find(f"{{{nsmap_a_bg}}}alpha")
                            if alpha_el is not None:
                                a_val = alpha_el.get("val")
                                if a_val:
                                    a_float = round(int(a_val) / 100000, 2)
                                    if a_float < 1.0 and color:
                                        color = _hex_to_rgba(color, a_float)
                except Exception:
                    pass

                if color:
                    result = {"type": "color", "value": color}
                    if theme_key:
                        theme_obj: Dict[str, Any] = {"key": theme_key, "color": color}
                        # 保留 lumMod/lumOff/tint/shade 变换参数，便于导出时回写
                        try:
                            nsmap_a_tr = "http://schemas.openxmlformats.org/drawingml/2006/main"
                            if bg_elem is not None:
                                solid_fill_tr = bg_elem.find(f".//{{{nsmap_a_tr}}}solidFill")
                                if solid_fill_tr is not None:
                                    scheme_tr = solid_fill_tr.find(f"{{{nsmap_a_tr}}}schemeClr")
                                    if scheme_tr is not None:
                                        transforms: Dict[str, int] = {}
                                        for tag_name in ("lumMod", "lumOff", "tint", "shade", "satMod"):
                                            tr_el = scheme_tr.find(f"{{{nsmap_a_tr}}}{tag_name}")
                                            if tr_el is not None and tr_el.get("val"):
                                                try:
                                                    transforms[tag_name] = int(tr_el.get("val"))
                                                except (ValueError, TypeError):
                                                    pass
                                        if transforms:
                                            theme_obj["transforms"] = transforms
                        except Exception:
                            pass
                        result["theme"] = theme_obj
                    return result

            elif fill.type == fill_gradient:  # GRADIENT
                gradient: Dict[str, Any] = {
                    "type": "linear",
                    "rotate": 0,
                    "colors": [],
                }
                try:
                    nsmap_a = "http://schemas.openxmlformats.org/drawingml/2006/main"
                    bg_elem = bg_owner.background._element if hasattr(bg_owner.background, "_element") else None
                    if bg_elem is not None:
                        grad_fill = bg_elem.find(f".//{{{nsmap_a}}}gradFill")
                        if grad_fill is not None:
                            path_node = grad_fill.find(f"{{{nsmap_a}}}path")
                            if path_node is not None:
                                gradient["type"] = "radial"
                                # 提取径向渐变中心点（fillToRect 的 l/t 值）
                                _ft = path_node.find(f"{{{nsmap_a}}}fillToRect")
                                if _ft is not None:
                                    try:
                                        _cx = int(_ft.get("l", "50000")) / 100000
                                        _cy = int(_ft.get("t", "50000")) / 100000
                                        if abs(_cx - 0.5) > 0.001 or abs(_cy - 0.5) > 0.001:
                                            gradient["center"] = {"x": round(_cx, 3), "y": round(_cy, 3)}
                                    except (ValueError, TypeError):
                                        pass
                except Exception:
                    pass
                try:
                    if hasattr(fill, "gradient_angle") and fill.gradient_angle is not None:
                        gradient["rotate"] = fill.gradient_angle
                except Exception:
                    pass

                try:
                    nsmap_a = "http://schemas.openxmlformats.org/drawingml/2006/main"
                    for stop in fill.gradient_stops:
                        stop_color = color_to_hex(stop.color, theme_color_map) or "#000000"
                        stop_pos = stop.position if hasattr(stop, "position") else 0
                        # per-stop alpha
                        try:
                            gs_el = stop._gs if hasattr(stop, "_gs") else None
                            if gs_el is not None:
                                transformed_stop_color = _extract_color_with_transforms(gs_el, theme_color_map)
                                if transformed_stop_color:
                                    stop_color = transformed_stop_color
                                for child in gs_el:
                                    alpha_el = child.find(f"{{{nsmap_a}}}alpha")
                                    if alpha_el is not None:
                                        a_val = alpha_el.get("val")
                                        if a_val:
                                            a_float = round(int(a_val) / 100000, 2)
                                            if a_float < 1.0:
                                                stop_color = _hex_to_rgba(stop_color, a_float)
                                        break
                        except Exception:
                            pass
                        gradient["colors"].append({
                            "pos": round(stop_pos, 3),
                            "color": stop_color,
                        })
                except Exception:
                    pass

                if gradient["colors"]:
                    return {"type": "gradient", "gradient": gradient}

            elif fill.type in (fill_picture, 5):  # PICTURE / BACKGROUND(部分版本)
                try:
                    nsmap = {
                        "a": "http://schemas.openxmlformats.org/drawingml/2006/main",
                        "r": "http://schemas.openxmlformats.org/officeDocument/2006/relationships",
                    }
                    bg_elem = bg._element
                    blip = bg_elem.find(".//a:blip", nsmap)
                    if blip is not None:
                        embed = blip.get(f"{{{nsmap['r']}}}embed")
                        if embed:
                            rel = bg_owner.part.rels[embed]
                            image_part = rel.target_part
                            content_type = image_part.content_type
                            blob = image_part.blob

                            # 优先 OSS 上传，降级 base64
                            src = ""
                            if image_handler and callable(image_handler):
                                try:
                                    uploaded_url = image_handler(blob, content_type)
                                    if uploaded_url:
                                        src = uploaded_url
                                except Exception as e:
                                    logger.warning(f"Background image_handler failed, falling back to base64: {e}")

                            if not src:
                                encoded = b64.b64encode(blob).decode("ascii")
                                src = f"data:{content_type};base64,{encoded}"

                            size = "cover"
                            try:
                                blip_fill = bg_elem.find(".//a:blipFill", nsmap)
                                if blip_fill is not None:
                                    tile = blip_fill.find("a:tile", nsmap)
                                    if tile is not None:
                                        size = "repeat"
                                    else:
                                        stretch = blip_fill.find("a:stretch", nsmap)
                                        fill_rect = stretch.find("a:fillRect", nsmap) if stretch is not None else None
                                        if fill_rect is not None:
                                            vals = [fill_rect.get(k, "0") for k in ("t", "r", "b", "l")]
                                            if any(
                                                str(v).strip() not in ("", "0")
                                                and int(str(v).strip() or "0") != 0
                                                for v in vals
                                            ):
                                                size = "contain"
                            except Exception:
                                pass

                            return {
                                "type": "image",
                                "image": {
                                    "src": src,
                                    "size": size,
                                },
                            }
                except Exception as e:
                    logger.debug(f"Failed to extract background image: {e}")

    except Exception:
        pass

    return None


def _extract_background(slide, theme_color_map: Optional[Dict[str, str]] = None, image_handler: Optional[Any] = None) -> Dict[str, Any]:
    """提取幻灯片背景（纯色 / 渐变 / 图片）"""
    result = _extract_background_from_element(slide, theme_color_map, image_handler=image_handler)
    if result:
        return result
    return {"type": "color", "value": "#ffffff"}


def _extract_layout_or_master_background(slide, theme_color_map: Optional[Dict[str, str]] = None, image_handler: Optional[Any] = None) -> Optional[Dict[str, Any]]:
    """
    从 slide_layout → slide_master 逐级查找背景。

    PPTX 背景继承链：slide → slide_layout → slide_master
    如果 slide 没有显式背景，应该从 layout/master 继承。
    """
    try:
        layout = slide.slide_layout
        if layout:
            bg = _extract_background_from_element(layout, theme_color_map, image_handler=image_handler)
            if bg:
                return bg

            # layout 也没有，查 master
            master = layout.slide_master
            if master:
                bg = _extract_background_from_element(master, theme_color_map, image_handler=image_handler)
                if bg:
                    return bg
    except Exception as e:
        logger.debug(f"Failed to extract layout/master background: {e}")

    return None


# ============================================================================
# WRITE: SlideElement[] → PPTX
# ============================================================================


def _clear_template_slides(prs) -> None:
    """清除模板中的已有幻灯片，递归清理子关系（图片/媒体）避免孤立文件。

    使用 python-pptx 私有 API（_sldIdLst / drop_rel），调用方需 try-except 防护。
    """
    _ns_r = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"
    while len(prs.slides) > 0:
        sldId = prs.slides._sldIdLst[0]
        rId = sldId.get(f"{{{_ns_r}}}id") or sldId.attrib.get("r:id")
        if rId:
            try:
                slide_part = prs.part.related_parts.get(rId)
                if slide_part is not None and hasattr(slide_part, "rels"):
                    for sub_rId in list(slide_part.rels.keys()):
                        slide_part.drop_rel(sub_rId)
            except (KeyError, AttributeError):
                pass
            prs.part.drop_rel(rId)
        del prs.slides._sldIdLst[0]


def write(
    pages: List[Dict[str, Any]],
    output_path: str,
    canvas_width: int = DEFAULT_SLIDE_WIDTH_PX,
    canvas_height: int = DEFAULT_SLIDE_HEIGHT_PX,
    template_path: Optional[str] = None,
    source_slide_width_emu: Optional[int] = None,
    source_slide_height_emu: Optional[int] = None,
    font_meta: Optional[Dict[str, Any]] = None,
    aigc_metadata: Optional[Dict[str, Any]] = None,
) -> str:
    """
    将 SlidePage[] JSON 写入 PPTX 文件。

    Args:
        pages: List[SlidePage] JSON
        output_path: 输出 PPTX 路径
        canvas_width: 编辑器画布宽度（px）
        canvas_height: 编辑器画布高度（px）
        template_path: 可选的 PPTX 模板路径
        source_slide_width_emu: 导入源 PPT 的原始宽度（EMU，可选）
        source_slide_height_emu: 导入源 PPT 的原始高度（EMU，可选）
        font_meta: 字体元数据 {embedded_fonts: [...], theme_fonts: {...}}，
                   嵌入字体条目需包含 name/style/format 及 data_base64 或 oss_url
        aigc_metadata: AIGC 溯源元数据，写入 docProps/custom.xml，
                       例：{"projectId": "...", "organizationId": "...", "spaceId": "...", "name": "..."}

    Returns:
        输出文件路径
    """
    from pptx import Presentation
    from pptx.util import Emu

    if template_path:
        prs = Presentation(template_path)
        try:
            _clear_template_slides(prs)
        except Exception:
            import logging
            logging.getLogger(__name__).warning(
                "模板幻灯片清理失败（python-pptx 私有 API 可能已变更），回退到空白演示文稿"
            )
            prs = Presentation()
    else:
        prs = Presentation()

    slide_width_emu, slide_height_emu = _resolve_slide_emu_for_write(
        canvas_width=canvas_width,
        canvas_height=canvas_height,
        source_slide_width_emu=source_slide_width_emu,
        source_slide_height_emu=source_slide_height_emu,
    )
    prs.slide_width = slide_width_emu
    prs.slide_height = slide_height_emu

    # 动态查找空白布局（名称包含 "Blank"/"blank"/"空白"）
    blank_layout = None
    for layout in prs.slide_layouts:
        name = (layout.name or "").lower()
        if "blank" in name or "空白" in name:
            blank_layout = layout
            break
    if blank_layout is None:
        # fallback: 用最后一个布局（通常是空白），或如果只有少数布局用第一个
        blank_layout = prs.slide_layouts[-1] if len(prs.slide_layouts) > 0 else prs.slide_layouts[0]

    # 两阶段写回：
    # 1) 先创建全部 slide（确保页内跳转可解析到任意目标页，包括前向引用）
    # 2) 再按顺序写入每页内容
    slide_page_pairs = []
    for page in pages:
        target_layout = _resolve_slide_layout_for_page(prs, page, blank_layout)
        slide = prs.slides.add_slide(target_layout)
        slide_page_pairs.append((slide, page))

    embedded_typefaces_token = _ACTIVE_EMBEDDED_TYPEFACES.set(
        _embedded_typefaces_from_font_meta(font_meta)
    )
    try:
        for slide, page in slide_page_pairs:
            _write_slide(
                slide,
                page,
                slide_width_emu=slide_width_emu,
                slide_height_emu=slide_height_emu,
                canvas_width=canvas_width,
                canvas_height=canvas_height,
            )
    finally:
        _ACTIVE_EMBEDDED_TYPEFACES.reset(embedded_typefaces_token)

    prs.save(output_path)

    embedded_fonts = (font_meta or {}).get("embedded_fonts")
    if embedded_fonts:
        _embed_fonts_into_pptx(output_path, embedded_fonts)

    # 导出后处理（一次读写完成 AIGC 溯源注入 + autofit fontScale 预算 + Keynote 兼容清洗）：
    # - AIGC：把 TabTin / projectId / organizationId / spaceId 写入 docProps/custom.xml。
    # - fontScale：不内嵌字体后缺字回退更宽字体会换行，WPS/PPT 打开不重算 autofit，
    #   故用字体度量把缩放比算好写死。
    # - Keynote 清洗：python-pptx 三处 OOXML 结构令 Keynote 拒绝导入，统一修正。
    # 各步内部独立 try，单步失败不影响其余，也绝不让导出整体失败。
    _postprocess_pptx(output_path, aigc_metadata=aigc_metadata)

    # 导出后再从最终 OOXML 回流检查一次。这里检查的是用户实际拿到的文件，而不是写入前
    # 的内存对象；发现风险只告警、不阻断下载，便于线上诊断包直接定位字体/换行契约。
    text_fidelity_issues = lint_pptx_text_fidelity(output_path)
    if text_fidelity_issues:
        logger.warning(
            "pptx text fidelity lint: %d issue(s) in %s: %s",
            len(text_fidelity_issues),
            output_path,
            text_fidelity_issues[:8],
        )

    logger.info(f"pptx_io.write: {len(pages)} pages → {output_path}")
    return output_path


# ── autofit fontScale 预算 ──
# 缺字机器把 Noto Sans SC / Inter 替换成回退字体（如微软雅黑/Arial），回退字通常更宽。
# 我们服务端只有 Noto/Inter 度量、拿不到目标机回退字度量，故用一个有依据的加宽余量估算：
# 经验上中英混排回退字宽度膨胀约 5~10%（CJK 汉字多为全角、膨胀小，拉丁字差异较大）。取 8%。
_AUTOFIT_FALLBACK_WIDTH_MARGIN = 1.08
# fontScale 下限，避免极端情况把文字缩到不可读。
_AUTOFIT_MIN_FONT_SCALE = 0.5
_EMU_PER_PT_AUTOFIT = 12700


@functools.lru_cache(maxsize=1)
def _autofit_metric_font():
    """加载用于度量的内置字体（Noto Sans SC 含 CJK + 拉丁），返回 (cmap, hmtx, unitsPerEm)。

    复用 dom_extractor 量框用的同一份 woff2 资产，避免额外内嵌用 TTF。
    """
    from fontTools.ttLib import TTFont

    fpath = (
        Path(__file__).resolve().parent.parent
        / "assets" / "fonts" / "NotoSansSC-sc-400.woff2"
    )
    f = TTFont(str(fpath), lazy=True)
    cmap = f.getBestCmap()
    hmtx = f["hmtx"]
    upm = f["head"].unitsPerEm
    adv = {gname: hmtx.metrics[gname][0] for gname in hmtx.metrics}
    return cmap, adv, upm


def _measure_text_width_pt(text: str, size_pt: float) -> float:
    """用内置字体度量一段文字（不换行）的自然宽度（pt）。未知字形按全角 em 估。"""
    try:
        cmap, adv, upm = _autofit_metric_font()
    except Exception:
        return len(text) * size_pt  # 兜底：按全角估
    total_units = 0
    for ch in text:
        gname = cmap.get(ord(ch))
        if gname is not None and gname in adv:
            total_units += adv[gname]
        else:
            total_units += upm  # 未知字符按全角 em
    return (total_units / upm) * size_pt


def lint_pptx_text_fidelity(pptx_path: str) -> List[Dict[str, Any]]:
    """检查最终 PPTX 中可导致字体替换或意外换行的高置信风险。

    返回结构化问题，供导出日志、回归测试和后续诊断包复用。该检查不修改文件，也不把
    不同 Office 套件的排版差异伪装成导出失败。
    """
    import zipfile
    from lxml import etree

    issues: List[Dict[str, Any]] = []
    A = "http://schemas.openxmlformats.org/drawingml/2006/main"
    P = "http://schemas.openxmlformats.org/presentationml/2006/main"

    try:
        with zipfile.ZipFile(pptx_path, "r") as archive:
            names = set(archive.namelist())
            embedded_typefaces: Set[str] = set()
            if "ppt/presentation.xml" in names:
                presentation = etree.fromstring(archive.read("ppt/presentation.xml"))
                for font_el in presentation.iter(f"{{{P}}}font"):
                    typeface = (font_el.get("typeface") or "").strip().casefold()
                    if typeface:
                        embedded_typefaces.add(typeface)

            slide_names = sorted(
                name
                for name in names
                if name.startswith("ppt/slides/slide") and name.endswith(".xml")
            )
            for slide_name in slide_names:
                root = etree.fromstring(archive.read(slide_name))
                for shape_index, shape in enumerate(root.iter(f"{{{P}}}sp"), start=1):
                    txbody = shape.find(f"{{{P}}}txBody")
                    if txbody is None:
                        continue
                    text = "".join(
                        node.text or ""
                        for node in txbody.iter(f"{{{A}}}t")
                    ).strip()
                    if not text:
                        continue
                    shape_ref = {
                        "slide": slide_name,
                        "shapeIndex": shape_index,
                        "text": text[:80],
                    }

                    for font_el in txbody.iter(f"{{{A}}}latin"):
                        typeface = (font_el.get("typeface") or "").strip()
                        if (
                            typeface.casefold() in _WEB_ONLY_OR_NON_SYSTEM_LATIN_FONTS
                            and typeface.casefold() not in embedded_typefaces
                        ):
                            issues.append({
                                **shape_ref,
                                "code": "unavailable-web-font",
                                "font": typeface,
                            })
                            break

                    body_pr = txbody.find(f"{{{A}}}bodyPr")
                    ext = shape.find(f"{{{P}}}spPr/{{{A}}}xfrm/{{{A}}}ext")
                    if body_pr is None or ext is None or body_pr.get("wrap") == "none":
                        continue
                    # 无显式换行、单段、框高接近单行字号时，它在 HTML 中高度疑似单行；
                    # 若仍允许 PPT 自动换行，就是可观测的契约缺口。
                    paragraphs = txbody.findall(f"{{{A}}}p")
                    has_explicit_break = any(
                        paragraph.find(f".//{{{A}}}br") is not None
                        for paragraph in paragraphs
                    )
                    sizes: List[float] = []
                    for rpr in txbody.iter(f"{{{A}}}rPr"):
                        try:
                            sizes.append(int(rpr.get("sz", "0")) / 100.0)
                        except ValueError:
                            pass
                    try:
                        box_height_pt = int(ext.get("cy", "0")) / _EMU_PER_PT_AUTOFIT
                    except ValueError:
                        box_height_pt = 0
                    max_size = max(sizes) if sizes else 0
                    if (
                        len(paragraphs) == 1
                        and not has_explicit_break
                        and max_size > 0
                        and 0 < box_height_pt <= max_size * 1.65
                    ):
                        issues.append({
                            **shape_ref,
                            "code": "single-line-wrap-enabled",
                            "boxHeightPt": round(box_height_pt, 2),
                            "fontSizePt": round(max_size, 2),
                        })
    except Exception as exc:
        return [{"code": "text-fidelity-lint-failed", "error": str(exc)}]
    return issues


def _autofit_fontscale_entries(entries: Dict[str, bytes]) -> None:
    """为带 normAutofit 的文本框预算 fontScale 并写死（就地修改 entries）。

    思路：文本框尺寸是编辑器按内置字体量出的、贴合文字。缺字机器回退更宽字体后同样的字
    需要更多横向空间→多出行→撑破框高。这里用内置字体度量 + 回退加宽余量估算回退字下的
    换行行数与总高，若超过框内高，则算出使其塞回框的 fontScale（≤100%）写入 normAutofit，
    使 WPS/PowerPoint 无需重算即按缩放渲染，避免换行跑版。仅缩不放，且有下限保护。
    """
    from lxml import etree

    A = "http://schemas.openxmlformats.org/drawingml/2006/main"
    P = "http://schemas.openxmlformats.org/presentationml/2006/main"

    slide_names = [
        n for n in entries
        if n.startswith("ppt/slides/slide") and n.endswith(".xml")
    ]
    if not slide_names:
        return

    for sname in slide_names:
        root = etree.fromstring(entries[sname])
        sp_changed = False
        for sp in root.iter(f"{{{P}}}sp"):
            txbody = sp.find(f"{{{P}}}txBody")
            if txbody is None:
                continue
            body_pr = txbody.find(f"{{{A}}}bodyPr")
            if body_pr is None:
                continue
            norm = body_pr.find(f"{{{A}}}normAutofit")
            if norm is None:
                continue
            ext = sp.find(f"{{{P}}}spPr/{{{A}}}xfrm/{{{A}}}ext")
            if ext is None:
                continue
            try:
                box_w = int(ext.get("cx", "0"))
                box_h = int(ext.get("cy", "0"))
            except ValueError:
                continue
            if box_w <= 0 or box_h <= 0:
                continue

            def _ins(attr):
                v = body_pr.get(attr)
                try:
                    return int(v) if v is not None else 0
                except ValueError:
                    return 0

            inner_w = box_w - _ins("lIns") - _ins("rIns")
            inner_h = box_h - _ins("tIns") - _ins("bIns")
            if inner_w <= 0 or inner_h <= 0:
                continue
            inner_w_pt = inner_w / _EMU_PER_PT_AUTOFIT

            # 只处理「单行文字」这一最常见、且可仅凭宽度安全判定的跑版场景（标题/标签）：
            # 某段在内置字体下本来一行放得下（nat_w ≤ 内宽），但缺字回退加宽后会超一行
            # （nat_w×余量 > 内宽）→ 按 内宽/(nat_w×余量) 温和缩放，让它继续保持一行。
            # 本来就多行的正文（nat_w > 内宽）不参与约束——后端无法可靠复现其换行排版，
            # 强行按高度缩放会误缩，宁可不动、信任编辑器量框。
            scale = 1.0
            for para in txbody.findall(f"{{{A}}}p"):
                sizes = []
                text_parts = []
                for r in para.findall(f"{{{A}}}r"):
                    rpr = r.find(f"{{{A}}}rPr")
                    sz = rpr.get("sz") if rpr is not None else None
                    if sz:
                        try:
                            sizes.append(int(sz) / 100.0)
                        except ValueError:
                            pass
                    t = r.find(f"{{{A}}}t")
                    if t is not None and t.text:
                        text_parts.append(t.text)
                text = "".join(text_parts)
                if not text.strip():
                    continue
                size_pt = max(sizes) if sizes else 18.0
                nat_w_pt = _measure_text_width_pt(text, size_pt)
                # 本来就不止一行 → 跳过（不误缩多行正文）
                if nat_w_pt > inner_w_pt:
                    continue
                needed = nat_w_pt * _AUTOFIT_FALLBACK_WIDTH_MARGIN
                if needed > inner_w_pt:
                    scale = min(scale, inner_w_pt / needed)

            scale = max(_AUTOFIT_MIN_FONT_SCALE, scale)
            if scale < 0.999:
                norm.set("fontScale", str(int(round(scale * 100000))))
                sp_changed = True

        if sp_changed:
            entries[sname] = etree.tostring(
                root, xml_declaration=True, encoding="UTF-8", standalone=True
            )


def _sanitize_entries(entries: Dict[str, bytes]) -> None:
    """清洗 python-pptx 输出使其能被 Apple Keynote / macOS 导入（就地修改 entries，）。

    python-pptx 的输出在 PowerPoint/PowerPoint Online/LibreOffice 都能打开，但 Keynote
    校验更严，有三处结构会被直接拒绝（"文件格式无效/无法导入"）：

    1. sldSz 尺寸与 type 矛盾：非 4:3 尺寸却仍标 type="screen4x3"（python-pptx 改尺寸时
       不更新 type）。Keynote 比对尺寸与 type 发现矛盾即拒绝 → 去掉 type（可选属性，
       打开方按 cx/cy 推导），近 16:9/4:3 标准值时顺带吸附到规范 EMU。
    2. 缺 notesMasterIdLst：建了 notesMaster 关系却没在 presentation.xml 登记 master
       （违反 ECMA-376 §19.2.1.30）→ 按关系补写 notesMasterIdLst。
    3. Windows 打印机设置桩：ppt/printerSettings/*.bin 及其 content-type，macOS 无对应
       驱动，Keynote 拒绝声明该类型的文件 → 删除该 part、关系与 content-type 声明。

    幂等、可安全重复执行。参考：github.com/anthropics/skills#1167。
    """
    import re

    # ── #3 移除 Windows 打印机设置桩 ──
    printer_parts = [n for n in entries if n.startswith("ppt/printerSettings/")]
    for n in printer_parts:
        del entries[n]
    if printer_parts:
        rels_path = "ppt/_rels/presentation.xml.rels"
        if rels_path in entries:
            rels = entries[rels_path].decode("utf-8")
            entries[rels_path] = re.sub(
                r"<Relationship[^>]*printerSettings[^>]*/>", "", rels
            ).encode("utf-8")
        ct_path = "[Content_Types].xml"
        if ct_path in entries:
            ct = entries[ct_path].decode("utf-8")
            entries[ct_path] = re.sub(
                r'<Default Extension="bin"[^>]*ContentType="[^"]*printerSettings"[^>]*/>',
                "",
                ct,
            ).encode("utf-8")

    # ── #1 / #2 修正 presentation.xml ──
    pres_path = "ppt/presentation.xml"
    rels_path = "ppt/_rels/presentation.xml.rels"
    if pres_path in entries:
        pres = entries[pres_path].decode("utf-8")

        # #1：sldSz 去矛盾 type
        m = re.search(r"<p:sldSz\s+([^/]+?)/>", pres)
        if m:
            attrs = m.group(1)
            cx_m = re.search(r'cx="(\d+)"', attrs)
            cy_m = re.search(r'cy="(\d+)"', attrs)
            if cx_m and cy_m:
                cx, cy = int(cx_m.group(1)), int(cy_m.group(1))
                if abs(cx - 12192000) <= 60000 and abs(cy - 6858000) <= 60000:
                    new_el = '<p:sldSz cx="12192000" cy="6858000"/>'
                elif abs(cx - 9144000) <= 60000 and abs(cy - 6858000) <= 60000:
                    new_el = '<p:sldSz cx="9144000" cy="6858000" type="screen4x3"/>'
                else:
                    new_attrs = re.sub(r'\s+type="[^"]+"', "", attrs).strip()
                    new_el = f"<p:sldSz {new_attrs}/>"
                pres = pres.replace(m.group(0), new_el)

        # #2：补 notesMasterIdLst
        if "<p:notesMasterIdLst" not in pres and rels_path in entries:
            rels = entries[rels_path].decode("utf-8")
            rel_m = re.search(
                r'<Relationship Id="(rId\d+)"[^>]*?notesMaster[^>]*?/>', rels
            )
            if rel_m:
                rid = rel_m.group(1)
                block = (
                    f'<p:notesMasterIdLst><p:notesMasterId r:id="{rid}"/>'
                    f"</p:notesMasterIdLst>"
                )
                pres = re.sub(
                    r"</p:sldIdLst>", "</p:sldIdLst>" + block, pres, count=1
                )

        entries[pres_path] = pres.encode("utf-8")


def _sanitize_pptx_for_keynote(pptx_path: str) -> None:
    """path 版薄封装（供直接调用/测试）；生产导出走 _postprocess_pptx 单次读写。"""
    import zipfile

    with zipfile.ZipFile(pptx_path, "r") as zin:
        entries = {n: zin.read(n) for n in zin.namelist()}
    _sanitize_entries(entries)
    with zipfile.ZipFile(pptx_path, "w", zipfile.ZIP_DEFLATED) as zout:
        for name, data in entries.items():
            zout.writestr(name, data)


def _postprocess_pptx(pptx_path: str, aigc_metadata: Optional[Dict[str, Any]] = None) -> None:
    """导出后处理：一次读入 zip，依次应用 AIGC 注入 / autofit fontScale 预算 / Keynote 清洗，
    一次写回。合并原先三次全量 zip 读写为一次。各步独立 try，互不影响。
    """
    import zipfile

    with zipfile.ZipFile(pptx_path, "r") as zin:
        entries = {n: zin.read(n) for n in zin.namelist()}

    if aigc_metadata:
        try:
            _inject_aigc_entries(entries, aigc_metadata)
        except Exception as exc:
            logger.warning("AIGC metadata injection failed for %s: %s", pptx_path, exc)
    try:
        _autofit_fontscale_entries(entries)
    except Exception as exc:
        logger.warning("autofit fontScale precompute failed for %s: %s", pptx_path, exc)
    try:
        _sanitize_entries(entries)
    except Exception as exc:
        logger.warning("Keynote sanitize failed for %s: %s", pptx_path, exc)

    with zipfile.ZipFile(pptx_path, "w", zipfile.ZIP_DEFLATED) as zout:
        for name, data in entries.items():
            zout.writestr(name, data)


# ============================================================================
# AIGC 溯源：docProps/custom.xml 注入
# ============================================================================

# Microsoft 标准 custom property fmtid（所有 OOXML custom property 都用同一个 GUID）
_AIGC_CUSTOM_PROP_FMTID = "{D5CDD505-2E9C-101B-9397-08002B2CF9AE}"
_AIGC_CUSTOM_XML_NS = "http://schemas.openxmlformats.org/officeDocument/2006/custom-properties"
_AIGC_CUSTOM_XML_NS_VT = "http://schemas.openxmlformats.org/officeDocument/2006/docPropsVTypes"
_AIGC_CONTENT_TYPE = "application/vnd.openxmlformats-officedocument.custom-properties+xml"
_AIGC_REL_TYPE = (
    "http://schemas.openxmlformats.org/officeDocument/2006/relationships/custom-properties"
)
_AIGC_DEFAULT_VERSION = "2026.05"


def _build_aigc_payload(metadata: Dict[str, Any]) -> str:
    """构造写入 <vt:lpwstr> 的 JSON 载荷，含固定标识 + 用户传入的 projectId/organizationId/spaceId。"""
    import json as _json
    from datetime import datetime, timezone

    payload: Dict[str, Any] = {
        "app": "TabTin",
        "version": str(metadata.get("version") or _AIGC_DEFAULT_VERSION),
        "generatedAt": datetime.now(timezone.utc).strftime("%Y-%m-%dT%H:%M:%SZ"),
    }
    for key in ("projectId", "organizationId", "spaceId", "name"):
        value = metadata.get(key)
        if value is not None and value != "":
            payload[key] = str(value)
    # 透传额外字段（如 agentId / userId 等扩展），保持向前兼容
    for key, value in metadata.items():
        if key in payload or key in ("version",):
            continue
        if value is None or value == "":
            continue
        payload[key] = value if isinstance(value, (int, float, bool)) else str(value)
    return _json.dumps(payload, ensure_ascii=False, separators=(",", ":"))


def _inject_aigc_metadata(pptx_path: str, metadata: Dict[str, Any]) -> None:
    """
    在已生成的 PPTX 文件里注入 AIGC 溯源信息（docProps/custom.xml）。

    必要时会同步更新 [Content_Types].xml 与 _rels/.rels，使 PowerPoint 能正常识别。
    幂等：重复调用会更新（而非重复添加）名为 "TabTin" 的 property。

    任何异常都会向上抛，由调用方决定是否吞掉（write() 会包 try/except）。
    """
    import os as _os
    import zipfile as _zipfile

    if not _os.path.exists(pptx_path):
        raise FileNotFoundError(pptx_path)

    with _zipfile.ZipFile(pptx_path, "r") as zf:
        entries = {n: zf.read(n) for n in zf.namelist()}
    _inject_aigc_entries(entries, metadata)
    with _zipfile.ZipFile(pptx_path, "w", _zipfile.ZIP_DEFLATED) as zf_out:
        for name, data in entries.items():
            zf_out.writestr(name, data)


def _inject_aigc_entries(entries: Dict[str, bytes], metadata: Dict[str, Any]) -> None:
    """AIGC 溯源注入（就地修改 entries）：写/更新 docProps/custom.xml 及其 content-type、rels。"""
    payload_json = _build_aigc_payload(metadata)

    new_custom_xml = _build_or_update_custom_xml(entries.get("docProps/custom.xml"), payload_json)
    new_ct_xml = _ensure_content_type_for_custom_xml(entries.get("[Content_Types].xml"))
    new_rels_xml = _ensure_rels_for_custom_xml(entries.get("_rels/.rels"))

    entries["docProps/custom.xml"] = new_custom_xml
    if new_ct_xml is not None:
        entries["[Content_Types].xml"] = new_ct_xml
    if new_rels_xml is not None:
        entries["_rels/.rels"] = new_rels_xml


def _build_or_update_custom_xml(existing_xml: Optional[bytes], payload_json: str) -> bytes:
    """生成或更新 docProps/custom.xml；幂等地写 name='TabTin' 的 property。"""
    from lxml import etree as _etree

    nsmap = {None: _AIGC_CUSTOM_XML_NS, "vt": _AIGC_CUSTOM_XML_NS_VT}

    if existing_xml:
        try:
            root = _etree.fromstring(existing_xml)
        except _etree.XMLSyntaxError:
            root = _etree.Element(f"{{{_AIGC_CUSTOM_XML_NS}}}Properties", nsmap=nsmap)
    else:
        root = _etree.Element(f"{{{_AIGC_CUSTOM_XML_NS}}}Properties", nsmap=nsmap)

    # 移除已有的 name="TabTin"
    target = None
    for prop in list(root):
        if prop.tag.endswith("}property") and prop.get("name") == "TabTin":
            target = prop
            break
    if target is not None:
        root.remove(target)

    # 计算下一个 pid（property id），从 2 开始，避开已有的占用值
    used_pids = set()
    for prop in root:
        if prop.tag.endswith("}property"):
            try:
                used_pids.add(int(prop.get("pid", "0") or "0"))
            except ValueError:
                pass
    next_pid = 2
    while next_pid in used_pids:
        next_pid += 1

    prop_el = _etree.SubElement(root, f"{{{_AIGC_CUSTOM_XML_NS}}}property")
    prop_el.set("fmtid", _AIGC_CUSTOM_PROP_FMTID)
    prop_el.set("pid", str(next_pid))
    prop_el.set("name", "TabTin")
    lpwstr = _etree.SubElement(prop_el, f"{{{_AIGC_CUSTOM_XML_NS_VT}}}lpwstr")
    lpwstr.text = payload_json

    return _etree.tostring(
        root,
        xml_declaration=True,
        encoding="UTF-8",
        standalone=True,
        pretty_print=False,
    )


def _ensure_content_type_for_custom_xml(existing_xml: Optional[bytes]) -> Optional[bytes]:
    """确保 [Content_Types].xml 含 custom-properties 的 Override 节点；已含则返回 None（无需改写）。"""
    from lxml import etree as _etree

    if not existing_xml:
        return None
    ns_ct = "http://schemas.openxmlformats.org/package/2006/content-types"
    try:
        root = _etree.fromstring(existing_xml)
    except _etree.XMLSyntaxError:
        return None
    for child in root:
        if child.tag.endswith("}Override") and child.get("PartName") == "/docProps/custom.xml":
            return None  # 已注册
    override = _etree.SubElement(root, f"{{{ns_ct}}}Override")
    override.set("PartName", "/docProps/custom.xml")
    override.set("ContentType", _AIGC_CONTENT_TYPE)
    return _etree.tostring(
        root,
        xml_declaration=True,
        encoding="UTF-8",
        standalone=True,
        pretty_print=False,
    )


def _ensure_rels_for_custom_xml(existing_xml: Optional[bytes]) -> Optional[bytes]:
    """确保 _rels/.rels 含 custom-properties relationship；已含则返回 None。"""
    from lxml import etree as _etree

    if not existing_xml:
        return None
    ns_rels = "http://schemas.openxmlformats.org/package/2006/relationships"
    try:
        root = _etree.fromstring(existing_xml)
    except _etree.XMLSyntaxError:
        return None
    used_ids = set()
    for child in root:
        if child.tag.endswith("}Relationship"):
            if child.get("Type") == _AIGC_REL_TYPE and child.get("Target") in (
                "docProps/custom.xml", "/docProps/custom.xml"
            ):
                return None  # 已注册
            rid = child.get("Id", "")
            if rid.startswith("rId"):
                try:
                    used_ids.add(int(rid[3:]))
                except ValueError:
                    pass
    next_id = 1
    while next_id in used_ids:
        next_id += 1
    rel = _etree.SubElement(root, f"{{{ns_rels}}}Relationship")
    rel.set("Id", f"rId{next_id}")
    rel.set("Type", _AIGC_REL_TYPE)
    rel.set("Target", "docProps/custom.xml")
    return _etree.tostring(
        root,
        xml_declaration=True,
        encoding="UTF-8",
        standalone=True,
        pretty_print=False,
    )


def _parse_oss_url_to_object_key(url: str) -> Optional[str]:
    """如果 URL 指向我们配置的 OSS bucket，返回其中的 object_key；否则 None。

    设计原因：
      在某些本地开发环境下（用户跑透明代理/ClashX/Surge 等），所有公网域名
      会被 DNS 劫持到 198.18.x.x 假 IP（RFC 5735 reserved range），
      ssrf_safe_urlopen 看到私有网段会拒绝（这是正确的 SSRF 行为）。

      其他业务模块（TabData / TabDoc / OSS 上传等）都通过 OSS SDK 直接传
      object_key 工作，不受影响。只有 pptx_io export 是从 PPTElement.src 里
      拿 HTTPS URL 走 ssrf_safe_urlopen，是唯一受 198.18 影响的下载路径。

      这个函数把 OSS URL → object_key，让 pptx_io 也能走 SDK 下载，
      跟其他模块行为一致。非我们 OSS 的 URL（外部图片）继续走 SSRF 安全的 HTTPS 下载。

    支持的 URL 形态：
      - https://{bucket}.{endpoint}/{object_key}
      - https://{cdn_domain}/{object_key}
      - https://{bucket}.{endpoint}/{object_key}?{signature}  (带签名直接也吃)
    """
    if not url or not isinstance(url, str):
        return None
    if not (url.startswith("http://") or url.startswith("https://")):
        return None
    try:
        from django.conf import settings
        bucket = (getattr(settings, "ALIYUN_OSS_BUCKET_NAME", "") or "").strip().lower()
        endpoint = (getattr(settings, "ALIYUN_OSS_ENDPOINT", "") or "").strip().lower()
        # CDN 域名（如果配置了）
        cdn_domain = (getattr(settings, "ALIYUN_OSS_CDN_DOMAIN", "") or "").strip().lower()
        if not bucket or not endpoint:
            return None
    except Exception:
        return None

    try:
        from urllib.parse import urlparse
        parsed = urlparse(url)
    except Exception:
        return None
    host = (parsed.hostname or "").lower()
    if not host:
        return None

    # 匹配 {bucket}.{endpoint}（标准 OSS URL）
    expected_oss_host = f"{bucket}.{endpoint}"
    expected_oss_internal_host = f"{bucket}.{endpoint.replace('.aliyuncs.com', '-internal.aliyuncs.com')}"

    matched = (
        host == expected_oss_host
        or host == expected_oss_internal_host
        or (cdn_domain and host == cdn_domain)
    )
    if not matched:
        return None

    # 提取 path 作为 object_key（去掉前导 /）
    object_key = (parsed.path or "").lstrip("/")
    if not object_key:
        return None
    return object_key


def _download_oss_via_sdk(object_key: str) -> Optional[bytes]:
    """通过 OSS SDK 下载（绕过 DNS 解析，跟项目其他模块行为一致）。"""
    try:
        from apps.services.oss.services.factory import get_oss_service
        svc = get_oss_service()
        if not svc:
            return None
        result = svc.download_file(object_key)
        if not result.get("success"):
            logger.warning("OSS SDK download failed: %s", result.get("message"))
            return None
        return result.get("data", {}).get("content")
    except Exception as exc:
        logger.warning("OSS SDK download error: %s", exc)
        return None


def _download_image_smart(
    src: str,
    *,
    max_bytes: Optional[int] = None,
    timeout: int = 15,
) -> Optional[bytes]:
    """统一图片下载入口，自动选择 OSS SDK 或 HTTPS。

    优先级：
      1. data: 协议 → 直接 base64 解码
      2. 命中我们 OSS bucket 的 URL → 走 OSS SDK（避开 DNS 劫持 / 198.18 假 IP）
      3. 其他 http(s) → 走 ssrf_safe_urlopen（保留对外部 URL 的 SSRF 防护）

    返回 None 表示下载失败，调用方应优雅处理（log + 跳过该图片，不阻断整个 export）。
    """
    if not src or not isinstance(src, str):
        return None

    # data URL
    if src.startswith("data:"):
        try:
            import base64 as b64
            _, b64_data = src.split(",", 1)
            data = b64.b64decode(b64_data)
            if max_bytes and len(data) > max_bytes:
                logger.warning("data URL image too large (%d > %d), skipping", len(data), max_bytes)
                return None
            return data
        except Exception as e:
            logger.warning("Failed to decode base64 image: %s", e)
            return None

    # 我们自己的 OSS → 走 SDK
    object_key = _parse_oss_url_to_object_key(src)
    if object_key:
        data = _download_oss_via_sdk(object_key)
        if data is not None:
            if max_bytes and len(data) > max_bytes:
                logger.warning("OSS image too large (%d > %d), skipping: %s", len(data), max_bytes, object_key)
                return None
            return data
        # SDK 失败时，fallback 到 HTTPS（兼容 OSS bucket 配置变化等情况）
        logger.info("OSS SDK download failed, falling back to HTTPS: %s", object_key)

    # 外部 HTTP/HTTPS → 走 SSRF 安全的下载
    if src.startswith("http://") or src.startswith("https://"):
        try:
            from apps.services.common.url_security import ssrf_safe_urlopen
            data = ssrf_safe_urlopen(src, timeout=timeout, max_read_bytes=(max_bytes + 1) if max_bytes else None)
            if max_bytes and data and len(data) > max_bytes:
                logger.warning("Image too large (%d > %d), skipping: %s", len(data), max_bytes, src[:120])
                return None
            return data
        except ValueError as e:
            logger.warning("SSRF blocked image URL: %s reason=%s", src[:120], e)
            return None
        except Exception as e:
            logger.warning("Failed to download image from URL: %s", e)
            return None

    logger.warning("[SECURITY] Blocked unsupported/local image src: %s", src[:200])
    return None


def _download_font_from_oss(oss_url: str) -> Optional[bytes]:
    """下载单个 OSS 字体（用于并发线程池）。

    优先走 OSS SDK（绕过 DNS 劫持），fallback HTTPS。
    """
    object_key = _parse_oss_url_to_object_key(oss_url)
    if object_key:
        data = _download_oss_via_sdk(object_key)
        if data is not None:
            return data
        logger.info("OSS SDK font download failed, falling back to HTTPS: %s", object_key)
    try:
        from apps.services.common.url_security import ssrf_safe_urlopen
        return ssrf_safe_urlopen(oss_url, timeout=15)
    except ValueError as exc:
        logger.warning("SSRF blocked font URL: %s reason=%s", oss_url[:120], exc)
        return None
    except Exception as exc:
        logger.warning("Failed to download font from %s: %s", oss_url, exc)
        return None


def _prefetch_oss_fonts(font_entries: List[Dict[str, Any]]) -> Dict[str, bytes]:
    """并发预下载所有需要 OSS 获取的字体，避免串行阻塞 Django worker。"""
    from concurrent.futures import ThreadPoolExecutor, as_completed

    oss_tasks: Dict[str, str] = {}
    for entry in font_entries:
        url = entry.get("oss_url", "")
        b64 = entry.get("data_base64", "")
        if url and not b64:
            key = f"{entry.get('name', '')}:{entry.get('style', '')}"
            oss_tasks[key] = url

    if not oss_tasks:
        return {}

    results: Dict[str, bytes] = {}
    with ThreadPoolExecutor(max_workers=min(len(oss_tasks), 4)) as pool:
        futures = {pool.submit(_download_font_from_oss, url): key for key, url in oss_tasks.items()}
        for future in as_completed(futures):
            key = futures[future]
            data = future.result()
            if data:
                results[key] = data
    return results


def _resolve_font_data_for_embed(
    font_entry: Dict[str, Any],
    prefetched: Optional[Dict[str, bytes]] = None,
) -> Optional[bytes]:
    """从 font_meta 条目获取字体二进制数据（base64 解码、预取缓存或 OSS 下载）。"""
    import base64

    data_b64 = font_entry.get("data_base64", "")
    if data_b64:
        try:
            return base64.b64decode(data_b64)
        except Exception as exc:
            logger.warning("Failed to decode font base64 for %s: %s", font_entry.get("name"), exc)
            return None

    oss_url = font_entry.get("oss_url", "")
    if oss_url:
        cache_key = f"{font_entry.get('name', '')}:{font_entry.get('style', '')}"
        if prefetched and cache_key in prefetched:
            return prefetched[cache_key]
        return _download_font_from_oss(oss_url)

    return None


_PRES_CHILD_ORDER = [
    "sldMasterIdLst", "notesMasterIdLst", "handoutMasterIdLst",
    "sldIdLst", "sldSz", "notesSz", "smartTags", "embeddedFontLst",
]


def _embed_fonts_into_pptx(
    pptx_path: str,
    embedded_fonts: List[Dict[str, Any]],
) -> None:
    """
    在已保存的 PPTX 中注入嵌入字体。

    操作 OOXML ZIP 结构：
    - ppt/fonts/{GUID}.fntdata — 混淆后的字体二进制
    - ppt/presentation.xml    — <p:embeddedFontLst> 声明
    - ppt/_rels/presentation.xml.rels — relationship 条目
    - [Content_Types].xml     — fntdata 内容类型
    """
    import os
    import shutil
    import tempfile
    import zipfile

    from lxml import etree

    font_groups: Dict[str, Dict[str, Dict[str, Any]]] = {}
    for font in embedded_fonts:
        if not isinstance(font, dict):
            continue
        name = (font.get("name") or "").strip()
        if not name:
            continue
        style = (font.get("style") or "normal").strip().lower()
        if style not in ("normal", "bold", "italic", "bolditalic"):
            style = "normal"
        if not font.get("data_base64") and not font.get("oss_url"):
            continue
        font_groups.setdefault(name, {})[style] = font

    if not font_groups:
        return

    p_ns = "http://schemas.openxmlformats.org/presentationml/2006/main"
    r_ns = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"
    rels_pkg_ns = "http://schemas.openxmlformats.org/package/2006/relationships"
    font_rel_type = "http://schemas.openxmlformats.org/officeDocument/2006/relationships/font"
    ct_ns = "http://schemas.openxmlformats.org/package/2006/content-types"

    style_to_tag = {
        "normal": "regular",
        "bold": "bold",
        "italic": "italic",
        "bolditalic": "boldItalic",
    }

    tmp_fd, tmp_path = tempfile.mkstemp(suffix=".pptx")
    os.close(tmp_fd)

    try:
        with zipfile.ZipFile(pptx_path, "r") as zin:
            pres_bytes = zin.read("ppt/presentation.xml")
            rels_bytes = zin.read("ppt/_rels/presentation.xml.rels")
            ct_bytes = zin.read("[Content_Types].xml")

            pres_root = etree.fromstring(pres_bytes)
            rels_root = etree.fromstring(rels_bytes)
            ct_root = etree.fromstring(ct_bytes)

            # 确保 fntdata 的 Content-Type 已注册
            has_fntdata_ct = any(
                el.get("Extension") == "fntdata"
                for el in ct_root
                if el.tag == f"{{{ct_ns}}}Default"
            )
            if not has_fntdata_ct:
                etree.SubElement(
                    ct_root,
                    f"{{{ct_ns}}}Default",
                    Extension="fntdata",
                    ContentType="application/x-fontdata",
                )

            # 计算下一个可用 rId
            existing_rids: Set[int] = set()
            for rel in rels_root:
                rid = rel.get("Id", "")
                if rid.startswith("rId"):
                    try:
                        existing_rids.add(int(rid[3:]))
                    except ValueError:
                        pass
            next_rid = max(existing_rids, default=0) + 1

            # 找到或创建 embeddedFontLst
            emb_lst = pres_root.find(f"{{{p_ns}}}embeddedFontLst")
            if emb_lst is None:
                emb_lst = etree.Element(f"{{{p_ns}}}embeddedFontLst")
                insert_idx = 0
                for idx, child in enumerate(pres_root):
                    local = etree.QName(child).localname
                    if local in _PRES_CHILD_ORDER:
                        pos = _PRES_CHILD_ORDER.index(local)
                        target_pos = _PRES_CHILD_ORDER.index("embeddedFontLst")
                        if pos < target_pos:
                            insert_idx = idx + 1
                pres_root.insert(insert_idx, emb_lst)

            # 并发预下载 OSS 字体，避免串行阻塞
            all_font_entries = [
                entry for styles in font_groups.values() for entry in styles.values()
            ]
            prefetched = _prefetch_oss_fonts(all_font_entries)

            _VALID_FONT_MAGIC = (
                b"\x00\x01\x00\x00",  # TrueType (TTF)
                b"OTTO",              # OpenType (OTF)
                b"true",              # TrueType (Apple variant)
                b"typ1",              # PostScript in SFNT
            )
            font_parts: List[Tuple[str, bytes]] = []
            for font_name, styles in font_groups.items():
                variant_infos: List[Tuple[str, bytes, str]] = []
                for style, font_entry in styles.items():
                    font_data = _resolve_font_data_for_embed(font_entry, prefetched=prefetched)
                    if not font_data or len(font_data) < 32:
                        continue
                    if not font_data[:4] in _VALID_FONT_MAGIC:
                        logger.warning(
                            "Skipping non-TTF/OTF font for embed: %s/%s (magic: %r)",
                            font_name, style, font_data[:4],
                        )
                        continue

                    # PowerPoint 的 ppt/fonts/*.fntdata（content type application/x-fontdata）
                    # 存放的是原始 TTF/OTF 字节，读取方（PowerPoint/Keynote/WPS）会直接按
                    # sfnt 解析。绝不能做 OOXML/Word 式 XOR 混淆——那会破坏 sfnt 头 32 字节，
                    # 令所有消费方解析内嵌字体失败进而崩溃。
                    tag_name = style_to_tag.get(style, "regular")
                    font_guid = str(uuid.uuid4()).upper()
                    variant_infos.append((tag_name, font_data, font_guid))

                if not variant_infos:
                    continue

                emb_font_el = etree.SubElement(emb_lst, f"{{{p_ns}}}embeddedFont")
                font_el = etree.SubElement(emb_font_el, f"{{{p_ns}}}font")
                font_el.set("typeface", font_name)

                for tag_name, font_bytes, font_guid in variant_infos:
                    part_rel_path = f"fonts/{{{font_guid}}}.fntdata"
                    zip_path = f"ppt/{part_rel_path}"
                    font_parts.append((zip_path, font_bytes))

                    rid = f"rId{next_rid}"
                    next_rid += 1
                    rel_el = etree.SubElement(rels_root, f"{{{rels_pkg_ns}}}Relationship")
                    rel_el.set("Id", rid)
                    rel_el.set("Type", font_rel_type)
                    rel_el.set("Target", part_rel_path)

                    variant_el = etree.SubElement(emb_font_el, f"{{{p_ns}}}{tag_name}")
                    variant_el.set(f"{{{r_ns}}}embed", rid)

            if not font_parts:
                try:
                    os.unlink(tmp_path)
                except OSError:
                    pass
                return

            # 标记「已内嵌 TrueType 字体」。embedTrueTypeFonts 默认 false（ECMA-376），
            # 若写了 embeddedFontLst + 字体 part 却不置位，包会自相矛盾（字体已嵌但标志说未嵌），
            # PowerPoint/Keynote/WPS 打开时判定损坏甚至崩溃。
            # 我们嵌的是完整字体（非子集），saveSubsetFonts 置 0，避免读取方按子集处理。
            pres_root.set("embedTrueTypeFonts", "1")
            pres_root.set("saveSubsetFonts", "0")

            # 写入新的 ZIP
            with zipfile.ZipFile(tmp_path, "w", zipfile.ZIP_DEFLATED) as zout:
                skip = {"ppt/presentation.xml", "ppt/_rels/presentation.xml.rels", "[Content_Types].xml"}
                for item in zin.infolist():
                    if item.filename not in skip:
                        zout.writestr(item, zin.read(item.filename))

                zout.writestr("ppt/presentation.xml", etree.tostring(
                    pres_root, xml_declaration=True, encoding="UTF-8", standalone=True,
                ))
                zout.writestr("ppt/_rels/presentation.xml.rels", etree.tostring(
                    rels_root, xml_declaration=True, encoding="UTF-8", standalone=True,
                ))
                zout.writestr("[Content_Types].xml", etree.tostring(
                    ct_root, xml_declaration=True, encoding="UTF-8", standalone=True,
                ))
                for zip_path, data in font_parts:
                    zout.writestr(zip_path, data)

        shutil.move(tmp_path, pptx_path)
        logger.info("Embedded %d font variant(s) into PPTX", len(font_parts))

    except Exception:
        logger.warning("Failed to embed fonts into PPTX: %s", pptx_path, exc_info=True)
        try:
            os.unlink(tmp_path)
        except OSError:
            pass


def _resolve_slide_emu_for_write(
    canvas_width: int,
    canvas_height: int,
    source_slide_width_emu: Optional[int] = None,
    source_slide_height_emu: Optional[int] = None,
) -> Tuple[int, int]:
    """
    计算写回时的幻灯片 EMU 尺寸。

    优先使用导入源 PPT 的原始 EMU（保证 EMU ↔ px 往返可逆）；
    若未提供，则按历史逻辑使用默认宽度 + 画布比例推导高度。
    """
    if (
        isinstance(source_slide_width_emu, int)
        and isinstance(source_slide_height_emu, int)
        and source_slide_width_emu > 0
        and source_slide_height_emu > 0
    ):
        return int(source_slide_width_emu), int(source_slide_height_emu)

    canvas_ratio = canvas_height / canvas_width if canvas_width > 0 else (9 / 16)
    slide_width_emu = DEFAULT_SLIDE_WIDTH_EMU
    slide_height_emu = int(slide_width_emu * canvas_ratio)
    return slide_width_emu, slide_height_emu


def _normalize_layout_name(raw: Any) -> str:
    return str(raw or "").strip().lower()


def _resolve_slide_layout_for_page(prs, page: Dict[str, Any], fallback_layout):
    """
    根据 page.layout 元数据优先恢复原始布局，失败时回退到 blank layout。
    """
    layout_meta = page.get("layout")
    if not isinstance(layout_meta, dict):
        return fallback_layout

    raw_index = layout_meta.get("index")
    idx_val: Optional[int] = None
    if isinstance(raw_index, int):
        idx_val = raw_index
    elif isinstance(raw_index, str):
        stripped = raw_index.strip()
        if stripped.lstrip("-").isdigit():
            try:
                idx_val = int(stripped)
            except Exception:
                idx_val = None
    if idx_val is not None and 0 <= idx_val < len(prs.slide_layouts):
        return prs.slide_layouts[idx_val]

    part_name = str(layout_meta.get("partName") or "").strip()
    if part_name:
        for layout in prs.slide_layouts:
            try:
                candidate = str(getattr(getattr(layout, "part", None), "partname", "") or "")
                if not candidate:
                    continue
                if candidate == part_name or candidate.endswith(part_name):
                    return layout
            except Exception:
                continue

    layout_name = _normalize_layout_name(layout_meta.get("name"))
    if layout_name:
        for layout in prs.slide_layouts:
            if _normalize_layout_name(getattr(layout, "name", "")) == layout_name:
                return layout

    return fallback_layout


def _find_matching_placeholder_shape(
    slide,
    placeholder_meta: Dict[str, Any],
    used_shape_ids: Optional[Set[int]] = None,
):
    """
    按 idx 优先、type 兜底匹配当前 slide 上的 placeholder shape。
    """
    if not isinstance(placeholder_meta, dict):
        return None

    target_idx = placeholder_meta.get("idx")
    target_type = str(placeholder_meta.get("type") or "").strip()
    if target_type in ("title", "ctrTitle"):
        candidate_types = {"title", "ctrTitle"}
    elif target_type in ("subTitle", "body", "obj"):
        candidate_types = {"subTitle", "body", "obj"}
    else:
        candidate_types = {target_type} if target_type else set()
    used = used_shape_ids or set()

    try:
        candidates = list(getattr(slide, "placeholders", []))
    except Exception:
        candidates = []

    if not candidates:
        try:
            candidates = [sp for sp in slide.shapes if hasattr(sp, "text_frame")]
        except Exception:
            candidates = []

    if target_idx is not None:
        for shape in candidates:
            try:
                shape_id = int(getattr(shape, "shape_id"))
                if shape_id in used:
                    continue
                meta = _extract_placeholder_meta(shape)
                if not meta:
                    continue
                if meta.get("idx") == target_idx:
                    return shape
            except Exception:
                continue

    if candidate_types:
        for shape in candidates:
            try:
                shape_id = int(getattr(shape, "shape_id"))
                if shape_id in used:
                    continue
                meta = _extract_placeholder_meta(shape)
                if not meta:
                    continue
                if str(meta.get("type") or "") in candidate_types:
                    return shape
            except Exception:
                continue

    return None


def _fill_text_element_into_placeholder(
    slide,
    element: Dict[str, Any],
    used_shape_ids: Set[int],
) -> bool:
    """
    将 text element 内容写入现有 placeholder，命中返回 True。
    """
    if element.get("type") != "text":
        return False
    props = element.get("props")
    if not isinstance(props, dict):
        return False

    placeholder_meta = props.get("placeholder")
    if not isinstance(placeholder_meta, dict):
        text_type = str(props.get("textType") or "").strip().lower()
        type_by_text_type = {
            "title": "title",
            "subtitle": "subTitle",
            "content": "body",
            "item": "body",
        }
        fallback_type = type_by_text_type.get(text_type)
        if not fallback_type:
            return False
        placeholder_meta = {"type": fallback_type}

    target_shape = _find_matching_placeholder_shape(slide, placeholder_meta, used_shape_ids)
    if target_shape is None or not hasattr(target_shape, "text_frame"):
        return False

    try:
        _apply_text_props_to_shape(target_shape, props, apply_outline=False)
        used_shape_ids.add(int(getattr(target_shape, "shape_id")))
        return True
    except Exception as e:
        logger.debug(f"Failed to fill placeholder text: {e}")
        return False


def _build_write_batches(elements: List[Dict[str, Any]]) -> List[Dict[str, Any]]:
    """
    基于 zIndex 构建写入批次，保证分组元素与非分组元素保持全局层级顺序。

    返回格式：
    - {"kind": "single", "element": {...}}
    - {"kind": "group", "groupId": "...", "elements": [{...}, ...]}
    """
    groups: Dict[str, List[Dict[str, Any]]] = {}
    group_indices: Dict[str, List[int]] = {}
    for idx, el in enumerate(elements):
        gid = el.get("groupId")
        if gid:
            groups.setdefault(gid, []).append(el)
            group_indices.setdefault(gid, []).append(idx)

    non_contiguous_groups = set()
    for gid, idx_list in group_indices.items():
        if not idx_list:
            continue
        if idx_list[-1] - idx_list[0] + 1 != len(idx_list):
            non_contiguous_groups.add(gid)

    emitted_groups = set()
    batches: List[Dict[str, Any]] = []
    for el in elements:
        gid = el.get("groupId")
        if not gid:
            batches.append({"kind": "single", "element": el})
            continue

        if gid in non_contiguous_groups:
            # 非连续 groupId 说明层级顺序已被打散，按单元素写入以保住视觉层级。
            batches.append({"kind": "single", "element": el})
            continue

        if gid in emitted_groups:
            continue
        emitted_groups.add(gid)

        members = groups.get(gid, [])
        if len(members) < 2:
            # 单元素“伪分组”按普通元素写入，避免空 grpSp
            for member in (members or [el]):
                batches.append({"kind": "single", "element": member})
            continue

        batches.append({"kind": "group", "groupId": gid, "elements": members})

    return batches


def _next_shape_id(sp_tree) -> int:
    """计算当前 slide spTree 下可用的下一个 p:cNvPr@id（max + 1）。"""
    nsmap_p = "http://schemas.openxmlformats.org/presentationml/2006/main"
    max_id = 0
    try:
        for node in sp_tree.findall(f".//{{{nsmap_p}}}cNvPr"):
            raw = node.get("id")
            if raw is None:
                continue
            try:
                max_id = max(max_id, int(raw))
            except (TypeError, ValueError):
                continue
    except Exception:
        pass
    return max_id + 1


def _write_slide(
    slide,
    page: Dict[str, Any],
    slide_width_emu: int,
    slide_height_emu: int,
    canvas_width: int,
    canvas_height: int,
) -> None:
    """写入单张幻灯片"""
    # 背景（支持纯色、渐变、图片）
    # 如果背景是从 layout/master 继承来的，跳过设置（让模板自身的 layout/master 定义背景）
    bg = page.get("background", {})
    if not bg.get("inherited"):
        bg_type = bg.get("type", "")
        # 兼容两套字段：抽取/DB 常用 value；normalize_background_for_api 产出 color
        bg_color = bg.get("value") or bg.get("color")
        if bg_type in ("color", "solid"):
            theme_meta = bg.get("theme")
            if isinstance(theme_meta, dict) and theme_meta.get("key"):
                _set_slide_background_theme(slide, theme_meta, fallback_color=bg_color)
            elif bg_color:
                _set_slide_background_color(slide, bg_color)
        elif bg_type == "theme":
            _set_slide_background_theme(slide, bg.get("theme", {}), fallback_color=bg_color)
        elif bg_type == "gradient" and bg.get("gradient"):
            _set_slide_background_gradient(slide, bg["gradient"])
        elif bg_type == "image" and bg.get("image"):
            _set_slide_background_image(slide, bg["image"])

    consumed_placeholder_element_ids: Set[str] = set()
    used_placeholder_shape_ids: Set[int] = set()

    # 先尝试将占位符文本回填到当前 layout 的 placeholder 中，避免导出时丢失版式语义。
    for element in page.get("elements", []):
        if not isinstance(element, dict):
            continue
        if _fill_text_element_into_placeholder(slide, element, used_placeholder_shape_ids):
            el_id = element.get("id")
            if el_id:
                consumed_placeholder_element_ids.add(str(el_id))

    # ── 按 zIndex 排序，保证层级顺序正确（PPTX 中后写入的元素在上层） ──
    raw_elements = page.get("elements", [])
    if consumed_placeholder_element_ids:
        raw_elements = [
            el for el in raw_elements
            if str(el.get("id", "")) not in consumed_placeholder_element_ids
        ]
    elements = _sort_elements_by_z_index(raw_elements)
    batches = _build_write_batches(elements)

    for batch in batches:
        kind = batch.get("kind")
        if kind == "single":
            element = batch.get("element", {})
            try:
                _write_element(
                    slide,
                    element,
                    slide_width_emu=slide_width_emu,
                    slide_height_emu=slide_height_emu,
                    canvas_width=canvas_width,
                    canvas_height=canvas_height,
                )
            except Exception as e:
                logger.warning(
                    f"Failed to write element {element.get('id')}/{element.get('type')}: {e}"
                )
            continue

        # group
        gid = batch.get("groupId")
        group_elements = batch.get("elements", [])
        try:
            _write_group_element(
                slide,
                group_elements,
                slide_width_emu=slide_width_emu,
                slide_height_emu=slide_height_emu,
                canvas_width=canvas_width,
                canvas_height=canvas_height,
            )
        except Exception as e:
            logger.warning(f"Failed to write group {gid}: {e}, falling back to flat write")
            for element in group_elements:
                try:
                    _write_element(
                        slide,
                        element,
                        slide_width_emu=slide_width_emu,
                        slide_height_emu=slide_height_emu,
                        canvas_width=canvas_width,
                        canvas_height=canvas_height,
                    )
                except Exception:
                    pass

    # 备注
    notes_text = page.get("notes", "")
    if notes_text:
        notes_slide = slide.notes_slide
        notes_slide.notes_text_frame.text = notes_text


def _write_group_element(
    slide,
    group_elements: List[Dict[str, Any]],
    slide_width_emu: int,
    slide_height_emu: int,
    canvas_width: int,
    canvas_height: int,
) -> None:
    """将一组相同 groupId 的元素写入 PPTX 的 grpSp（组合形状）。

    OOXML 组合结构：
    <p:grpSp>
      <p:nvGrpSpPr>...</p:nvGrpSpPr>
      <p:grpSpPr>
        <a:xfrm>
          <a:off x="组合左上角" y="组合左上角"/>
          <a:ext cx="组合宽" cy="组合高"/>
          <a:chOff x="子坐标原点x" y="子坐标原点y"/>
          <a:chExt cx="子坐标范围宽" cy="子坐标范围高"/>
        </a:xfrm>
      </p:grpSpPr>
      <p:sp>子元素1...</p:sp>
      <p:sp>子元素2...</p:sp>
    </p:grpSp>

    采用简化策略：chOff = off, chExt = ext（子坐标空间 = 幻灯片坐标空间）
    这样子元素的 EMU 坐标可以直接使用，无需额外转换。
    """
    from lxml import etree
    from pptx.oxml.ns import qn

    nsmap_p = "http://schemas.openxmlformats.org/presentationml/2006/main"
    nsmap_a = "http://schemas.openxmlformats.org/drawingml/2006/main"
    nsmap_r = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"

    group_name = ""
    for el in group_elements:
        raw_name = el.get("groupName")
        if isinstance(raw_name, str) and raw_name.strip():
            group_name = raw_name.strip()
            break
    if not group_name:
        group_name = "Group"

    # 计算组合的 bounding box (px → EMU)
    # 考虑子元素旋转后的实际占用区域（旋转矩形的 AABB）
    import math

    def _rotated_bbox(x, y, w, h, rot_deg):
        """计算旋转矩形的轴对齐包围盒"""
        if abs(rot_deg) < 0.01:
            return x, y, x + w, y + h
        rad = math.radians(rot_deg)
        cos_a = abs(math.cos(rad))
        sin_a = abs(math.sin(rad))
        # 旋转后的宽高
        rw = w * cos_a + h * sin_a
        rh = w * sin_a + h * cos_a
        # 以中心为基准计算新的 AABB
        cx, cy = x + w / 2, y + h / 2
        return cx - rw / 2, cy - rh / 2, cx + rw / 2, cy + rh / 2

    all_boxes = [
        _rotated_bbox(el.get("x", 0), el.get("y", 0), el.get("width", 0), el.get("height", 0), el.get("rotate", 0))
        for el in group_elements
    ]
    min_x = min(b[0] for b in all_boxes)
    min_y = min(b[1] for b in all_boxes)
    max_x = max(b[2] for b in all_boxes)
    max_y = max(b[3] for b in all_boxes)

    grp_left = px_to_emu(min_x, canvas_width, slide_width_emu)
    grp_top = px_to_emu(min_y, canvas_height, slide_height_emu)
    grp_width = px_to_emu(max_x - min_x, canvas_width, slide_width_emu)
    grp_height = px_to_emu(max_y - min_y, canvas_height, slide_height_emu)

    # 构建 grpSp XML
    grpSp = etree.SubElement(slide.shapes._spTree, qn("p:grpSp"))

    # nvGrpSpPr
    nvGrpSpPr = etree.SubElement(grpSp, qn("p:nvGrpSpPr"))
    cNvPr = etree.SubElement(nvGrpSpPr, qn("p:cNvPr"))
    cNvPr.set("id", str(_next_shape_id(slide.shapes._spTree)))
    cNvPr.set("name", group_name)
    etree.SubElement(nvGrpSpPr, qn("p:cNvGrpSpPr"))
    etree.SubElement(nvGrpSpPr, qn("p:nvPr"))

    # grpSpPr
    grpSpPr = etree.SubElement(grpSp, qn("p:grpSpPr"))
    xfrm = etree.SubElement(grpSpPr, qn("a:xfrm"))
    off = etree.SubElement(xfrm, qn("a:off"))
    off.set("x", str(grp_left))
    off.set("y", str(grp_top))
    ext = etree.SubElement(xfrm, qn("a:ext"))
    ext.set("cx", str(max(grp_width, 1)))
    ext.set("cy", str(max(grp_height, 1)))
    # 子坐标空间等于组合坐标空间（简化方案）
    chOff = etree.SubElement(xfrm, qn("a:chOff"))
    chOff.set("x", str(grp_left))
    chOff.set("y", str(grp_top))
    chExt = etree.SubElement(xfrm, qn("a:chExt"))
    chExt.set("cx", str(max(grp_width, 1)))
    chExt.set("cy", str(max(grp_height, 1)))

    # 将子元素写入到 slide 的 spTree，然后把新增的 shape XML 节点移到 grpSp 下。
    # python-pptx 的 add_textbox / add_picture 等 API 总是追加到 spTree，
    # 我们利用 count 差值检测新增节点，并移动所有新增节点（处理一个元素可能产生多个节点的情况）。
    for element in group_elements:
        try:
            sp_tree = slide.shapes._spTree
            count_before = len(sp_tree)
            _write_element(
                slide, element,
                slide_width_emu=slide_width_emu,
                slide_height_emu=slide_height_emu,
                canvas_width=canvas_width,
                canvas_height=canvas_height,
            )
            count_after = len(sp_tree)
            # 将所有新增的 shape 节点按原始顺序移到 grpSp 下。
            # 注意：某些元素（如 image + colorMask）会一次写出多个节点，
            # 若倒序迁移会破坏层级关系，导致读取侧无法合并回 image.props.colorMask。
            if count_after > count_before:
                new_nodes = list(sp_tree)[count_before:count_after]
                for new_node in new_nodes:
                    sp_tree.remove(new_node)
                    grpSp.append(new_node)
        except Exception as e:
            logger.warning(f"Failed to write group child {element.get('id')}/{element.get('type')}: {e}")

    # 如果 grpSp 中没有任何子形状，移除空组合（避免导出损坏的 PPTX）
    child_shape_count = len([c for c in grpSp if c.tag != qn("p:nvGrpSpPr") and c.tag != qn("p:grpSpPr")])
    if child_shape_count == 0:
        logger.warning("Empty group shape removed (no children written successfully)")
        slide.shapes._spTree.remove(grpSp)


def _write_element(
    slide,
    element: Dict[str, Any],
    slide_width_emu: int,
    slide_height_emu: int,
    canvas_width: int,
    canvas_height: int,
) -> None:
    """写入单个元素到幻灯片"""
    el_type = element.get("type")
    props = element.get("props", {})

    left = px_to_emu(element.get("x", 0), canvas_width, slide_width_emu)
    top = px_to_emu(element.get("y", 0), canvas_height, slide_height_emu)
    width = px_to_emu(element.get("width", 100), canvas_width, slide_width_emu)
    height = px_to_emu(element.get("height", 50), canvas_height, slide_height_emu)

    from pptx.util import Emu

    # 编辑器画布按 CSS 1px=1/96in 渲染（1920px 画布≈20in），字号统一以 pt 语义存储；
    # 而后端导出的 slide 物理宽度通常是 13.333in 标准宽屏。dpi_scale 即「编辑器画布 →
    # 导出 slide」的等比缩放系数（13.333/20≈0.667），用于在更小的 slide 上保持与编辑器
    # 一致的「字号 ∶ 幅面」视觉比例。无此缩放时字号会相对 slide 偏大、布局溢出。
    dpi_scale = _export_dpi_scale(slide_width_emu, canvas_width)

    shape = None
    if el_type == "text":
        if dpi_scale < 0.95 or dpi_scale > 1.05:
            props = dict(props)
            # 元素级 defaultFontSize 是 pt 语义（编辑器以 `${defaultFontSize}pt` 渲染），
            # 仅按 dpi_scale 等比缩放并保留 pt 单位，切勿当成 px 再做一次 px→pt 换算。
            raw_fs = props.get("defaultFontSize")
            if isinstance(raw_fs, (int, float)) and raw_fs > 0:
                props["defaultFontSize"] = f"{raw_fs * dpi_scale}pt"
            content = props.get("content", "")
            if content:
                import re
                # 内联 font-size 同样按 dpi_scale 等比缩放并保留原单位（编辑器内联字号为 pt，
                # 历史富文本可能残留 px）。px→pt 的换算交给后续 _css_length_to_pt 完成。
                def _scale_font(m: re.Match) -> str:
                    val = float(m.group(1))
                    unit = m.group(2) or 'px'
                    return f"font-size:{val * dpi_scale:.1f}{unit}"
                props["content"] = re.sub(
                    r'font-size:\s*([\d.]+)(px|pt)?',
                    _scale_font,
                    content,
                )
        shape = _write_text_element(slide, props, Emu(left), Emu(top), Emu(width), Emu(height))
    elif el_type == "table":
        shape = _write_table_element(slide, props, Emu(left), Emu(top), Emu(width), Emu(height))
    elif el_type == "image":
        shape = _write_image_element(
            slide,
            props,
            Emu(left),
            Emu(top),
            Emu(width),
            Emu(height),
            width_px=element.get("width"),
            height_px=element.get("height"),
        )
    elif el_type == "video":
        shape = _write_media_element(
            slide,
            media_type="video",
            props=props if isinstance(props, dict) else {},
            left=Emu(left),
            top=Emu(top),
            width=Emu(width),
            height=Emu(height),
        )
    elif el_type == "audio":
        shape = _write_media_element(
            slide,
            media_type="audio",
            props=props if isinstance(props, dict) else {},
            left=Emu(left),
            top=Emu(top),
            width=Emu(width),
            height=Emu(height),
        )
    elif el_type == "latex":
        latex_props = dict(props) if isinstance(props, dict) else {}
        # LaTeX 导出优先走 SVG，避免使用历史 rasterSrc 放大后模糊。
        svg_markup = latex_props.get("svg")
        if not (isinstance(svg_markup, str) and svg_markup.strip()):
            path_val = latex_props.get("path")
            view_box_val = latex_props.get("viewBox")
            if (
                isinstance(path_val, str)
                and path_val.strip()
                and isinstance(view_box_val, (list, tuple))
                and len(view_box_val) == 2
            ):
                try:
                    vb_w = float(view_box_val[0])
                    vb_h = float(view_box_val[1])
                    if vb_w > 0 and vb_h > 0:
                        color = str(latex_props.get("color", "#111111") or "#111111")
                        stroke_width = latex_props.get("strokeWidth", 0)
                        stroke_w = float(stroke_width) if stroke_width is not None else 0.0
                        safe_path = _escape_xml_text(path_val)
                        safe_color = _escape_xml_text(color)
                        svg_markup = (
                            f'<svg xmlns="http://www.w3.org/2000/svg" '
                            f'viewBox="0 0 {vb_w} {vb_h}" width="{vb_w}" height="{vb_h}" '
                            f'preserveAspectRatio="xMidYMid meet" style="color:{safe_color}">'
                            f'<path d="{safe_path}" fill="currentColor" stroke="currentColor" '
                            f'stroke-width="{stroke_w}"/></svg>'
                        )
                        latex_props["svg"] = svg_markup
                except Exception:
                    pass
        if isinstance(svg_markup, str) and svg_markup.strip():
            try:
                import base64 as _b64
                svg_b64 = _b64.b64encode(svg_markup.encode("utf-8")).decode("ascii")
                latex_props["src"] = f"data:image/svg+xml;base64,{svg_b64}"
            except Exception:
                pass
        elif not latex_props.get("src"):
            raster_src = latex_props.get("rasterSrc")
            if isinstance(raster_src, str) and raster_src:
                latex_props["src"] = raster_src
        if not latex_props.get("src"):
            latex_source = latex_props.get("latex")
            if isinstance(latex_source, str) and latex_source.strip():
                placeholder_src = _build_latex_placeholder_svg_data_url(
                    latex_source,
                    color=str(latex_props.get("color", "#111111") or "#111111"),
                    width_px=element.get("width"),
                    height_px=element.get("height"),
                )
                if placeholder_src:
                    latex_props["src"] = placeholder_src
        raw_alt_text = latex_props.get("altText")
        has_latex_meta = isinstance(raw_alt_text, str) and raw_alt_text.strip().startswith(LATEX_META_PREFIX)
        if not has_latex_meta:
            auto_alt_text = _encode_latex_alt_text(latex_props)
            if auto_alt_text:
                latex_props["altText"] = auto_alt_text
        shape = _write_image_element(
            slide,
            latex_props,
            Emu(left),
            Emu(top),
            Emu(width),
            Emu(height),
            width_px=element.get("width"),
            height_px=element.get("height"),
            svg_scale=3.0,
        )
    elif el_type == "chart":
        shape = _write_chart_element(slide, props, Emu(left), Emu(top), Emu(width), Emu(height))
    elif el_type == "shape":
        shape = _write_shape_element(slide, props, Emu(left), Emu(top), Emu(width), Emu(height))
    elif el_type == "line":
        _write_line_element(slide, element, slide_width_emu, slide_height_emu, canvas_width, canvas_height)
    else:
        logger.warning(f"Unknown element type: {el_type}")

    if shape is not None:
        _apply_common_write_props(shape, element)

    # 图片颜色蒙版 → 在图片上方叠加一个同尺寸半透明矩形
    if el_type == "image" and shape is not None:
        color_mask = props.get("colorMask")
        if color_mask and isinstance(color_mask, str):
            try:
                from pptx.util import Emu as _Emu
                from pptx.enum.shapes import MSO_SHAPE
                from pptx.dml.color import RGBColor

                hex_cm, cm_alpha = _parse_css_color(color_mask)
                hex_6 = hex_cm.lstrip("#")[:6]
                if len(hex_6) == 6 and cm_alpha is not None and cm_alpha > 0:
                    mask_shape = slide.shapes.add_shape(
                        MSO_SHAPE.RECTANGLE,
                        _Emu(left), _Emu(top), _Emu(width), _Emu(height),
                    )
                    mask_shape.fill.solid()
                    mask_shape.fill.fore_color.rgb = RGBColor.from_string(hex_6)
                    mask_shape.line.fill.background()

                    # 设置填充色的 alpha
                    _apply_color_alpha(mask_shape, "fill", cm_alpha)

                    # 如果图片有圆角，叠层也需要匹配
                    radius = props.get("radius")
                    clip = props.get("clip")
                    is_ellipse_clip = isinstance(clip, dict) and clip.get("shape") == "ellipse"

                    if is_ellipse_clip:
                        try:
                            from lxml import etree
                            from pptx.oxml.ns import qn as _qn
                            nsmap_a = "http://schemas.openxmlformats.org/drawingml/2006/main"
                            spPr = _find_sp_pr(mask_shape._element)
                            if spPr is not None:
                                for old_geom in spPr.findall(f"{{{nsmap_a}}}prstGeom"):
                                    spPr.remove(old_geom)
                                prst_geom = etree.SubElement(spPr, _qn("a:prstGeom"))
                                prst_geom.set("prst", "ellipse")
                                etree.SubElement(prst_geom, _qn("a:avLst"))
                        except Exception:
                            pass
                    elif radius and float(radius) > 0:
                        try:
                            from lxml import etree
                            from pptx.oxml.ns import qn as _qn
                            nsmap_a = "http://schemas.openxmlformats.org/drawingml/2006/main"
                            spPr = _find_sp_pr(mask_shape._element)
                            if spPr is not None:
                                for old_geom in spPr.findall(f"{{{nsmap_a}}}prstGeom"):
                                    spPr.remove(old_geom)
                                prst_geom = etree.SubElement(spPr, _qn("a:prstGeom"))
                                prst_geom.set("prst", "roundRect")
                                av_lst = etree.SubElement(prst_geom, _qn("a:avLst"))
                                img_w = element.get("width", 0)
                                img_h = element.get("height", 0)
                                shorter_px = min(float(img_w or 0), float(img_h or 0))
                                if shorter_px > 0:
                                    adj_val = min(int(max(0.0, float(radius) / shorter_px) * 100000), 50000)
                                else:
                                    shorter_side = min(int(mask_shape.width), int(mask_shape.height))
                                    adj_val = min(int(float(radius) * EMU_PER_PT / max(shorter_side, 1) * 100000), 50000)
                                if adj_val > 0:
                                    gd = etree.SubElement(av_lst, _qn("a:gd"))
                                    gd.set("name", "adj")
                                    gd.set("fmla", f"val {adj_val}")
                        except Exception:
                            pass

                    # 匹配图片旋转
                    rotate = element.get("rotate", 0)
                    if rotate:
                        mask_shape.rotation = float(rotate) % 360
            except Exception as e:
                logger.debug(f"Failed to write image colorMask overlay: {e}")


def _ensure_run_rpr(run_obj):
    """确保 run 有 a:rPr 节点，便于写入底层文本样式。"""
    from lxml import etree

    nsmap_a = "http://schemas.openxmlformats.org/drawingml/2006/main"
    r_pr = run_obj._r.find(f"{{{nsmap_a}}}rPr")
    if r_pr is None:
        r_pr = etree.SubElement(run_obj._r, f"{{{nsmap_a}}}rPr")
        run_obj._r.insert(0, r_pr)
    return r_pr


def _apply_run_color_with_alpha(run_obj, color_str: Any) -> None:
    """写入 run 字体颜色（支持 rgba alpha）。"""
    from lxml import etree
    from pptx.dml.color import RGBColor

    if not color_str:
        return

    hex_fc, alpha_fc = _parse_css_color(str(color_str))
    hex_6 = hex_fc.lstrip("#")[:6]
    if len(hex_6) != 6:
        return

    run_obj.font.color.rgb = RGBColor.from_string(hex_6)
    if alpha_fc is None or alpha_fc >= 1.0:
        return

    nsmap_a = "http://schemas.openxmlformats.org/drawingml/2006/main"
    r_pr = _ensure_run_rpr(run_obj)
    solid_fill = r_pr.find(f"{{{nsmap_a}}}solidFill")
    if solid_fill is None:
        solid_fill = etree.SubElement(r_pr, f"{{{nsmap_a}}}solidFill")
    color_el = solid_fill.find(f"{{{nsmap_a}}}srgbClr")
    if color_el is None:
        color_el = etree.SubElement(solid_fill, f"{{{nsmap_a}}}srgbClr")
    color_el.set("val", hex_6)

    for old_alpha in color_el.findall(f"{{{nsmap_a}}}alpha"):
        color_el.remove(old_alpha)
    alpha_el = etree.SubElement(color_el, f"{{{nsmap_a}}}alpha")
    alpha_el.set("val", str(int(max(0.0, min(1.0, float(alpha_fc))) * 100000)))


def _apply_run_theme_color(run_obj, theme_key: Any, color_str_for_alpha: Optional[Any] = None) -> bool:
    """按主题色 key 写入 run 字体颜色（schemeClr），可附加 alpha。"""
    from lxml import etree
    from pptx.enum.dml import MSO_THEME_COLOR

    normalized_key = _normalize_text_theme_key(theme_key)
    if not normalized_key:
        return False

    enum_name = _TEXT_THEME_KEY_TO_MSO.get(normalized_key)
    if not enum_name:
        return False

    try:
        run_obj.font.color.theme_color = getattr(MSO_THEME_COLOR, enum_name)
    except Exception:
        return False

    if color_str_for_alpha is None:
        return True

    _, alpha_fc = _parse_css_color(str(color_str_for_alpha))
    if alpha_fc is None or alpha_fc >= 1.0:
        return True

    try:
        nsmap_a = "http://schemas.openxmlformats.org/drawingml/2006/main"
        r_pr = _ensure_run_rpr(run_obj)
        solid_fill = r_pr.find(f"{{{nsmap_a}}}solidFill")
        if solid_fill is None:
            solid_fill = etree.SubElement(r_pr, f"{{{nsmap_a}}}solidFill")
        scheme_el = solid_fill.find(f"{{{nsmap_a}}}schemeClr")
        if scheme_el is None:
            scheme_el = etree.SubElement(solid_fill, f"{{{nsmap_a}}}schemeClr")
        scheme_el.set("val", normalized_key)
        for old_alpha in scheme_el.findall(f"{{{nsmap_a}}}alpha"):
            scheme_el.remove(old_alpha)
        alpha_el = etree.SubElement(scheme_el, f"{{{nsmap_a}}}alpha")
        alpha_el.set("val", str(int(max(0.0, min(1.0, float(alpha_fc))) * 100000)))
    except Exception:
        pass
    return True


def _apply_run_style(
    run_obj,
    run_data: Dict[str, Any],
    default_font: Optional[str] = None,
    default_size: Optional[Any] = None,
    default_color: Optional[str] = None,
    default_color_theme_key: Optional[str] = None,
) -> None:
    """将 run_data 样式写入 python-pptx run。"""
    from pptx.util import Pt

    font = run_obj.font
    font_name = run_data.get("fontFamily") or default_font
    font_size = run_data.get("fontSize") if run_data.get("fontSize") is not None else default_size
    font_color = run_data.get("color") or default_color
    run_theme_key = _normalize_text_theme_key(run_data.get("themeColorKey"))
    default_theme_key = _normalize_text_theme_key(default_color_theme_key)

    if font_name:
        latin_typeface, _, _ = _resolve_pptx_typefaces(
            str(font_name),
            getattr(run_obj, "text", ""),
        )
        font.name = latin_typeface
        _apply_run_typeface_scripts(run_obj, str(font_name))

    font_size_pt = _css_length_to_pt(font_size)
    if font_size_pt is not None and font_size_pt > 0:
        font.size = Pt(round(font_size_pt, 3))

    applied_theme_color = False
    if run_theme_key:
        applied_theme_color = _apply_run_theme_color(run_obj, run_theme_key, font_color)
    elif default_theme_key and run_data.get("color") is None:
        # 默认主题色仅在 run 未显式设置 color 时生效
        applied_theme_color = _apply_run_theme_color(run_obj, default_theme_key, font_color)

    if not applied_theme_color and font_color:
        try:
            _apply_run_color_with_alpha(run_obj, font_color)
        except Exception:
            pass

    if run_data.get("bold"):
        font.bold = True
    if run_data.get("italic"):
        font.italic = True
    if run_data.get("underline"):
        u_style = run_data.get("underlineStyle")
        if u_style and u_style != "sng":
            try:
                r_pr = _ensure_run_rpr(run_obj)
                r_pr.set("u", u_style)
            except Exception:
                font.underline = True
        else:
            font.underline = True

    nsmap_a = "http://schemas.openxmlformats.org/drawingml/2006/main"

    # 删除线
    if run_data.get("strikethrough"):
        try:
            r_pr = _ensure_run_rpr(run_obj)
            r_pr.set("strike", "sngStrike")
        except Exception:
            pass

    # 上标 / 下标
    baseline_type = run_data.get("baseline")
    if baseline_type:
        try:
            r_pr = _ensure_run_rpr(run_obj)
            if baseline_type == "super":
                r_pr.set("baseline", "30000")
            elif baseline_type == "sub":
                r_pr.set("baseline", "-25000")
        except Exception:
            pass

    # 高亮
    hl_color = run_data.get("highlight")
    if hl_color:
        try:
            from lxml import etree
            r_pr = _ensure_run_rpr(run_obj)
            for old_hl in r_pr.findall(f"{{{nsmap_a}}}highlight"):
                r_pr.remove(old_hl)
            highlight_node = etree.SubElement(r_pr, f"{{{nsmap_a}}}highlight")
            hex_hl, _ = _parse_css_color(str(hl_color))
            srgb_node = etree.SubElement(highlight_node, f"{{{nsmap_a}}}srgbClr")
            srgb_node.set("val", hex_hl.lstrip("#")[:6])
        except Exception:
            pass

    # 超链接
    _apply_run_hyperlink_write(run_obj, run_data.get("hyperlink"))

    # 字间距（px -> spc(1/100pt)）
    run_ls = run_data.get("letterSpacing")
    ls_px = _css_length_to_px(run_ls)
    if ls_px is not None and abs(ls_px) > 1e-9:
        try:
            r_pr = _ensure_run_rpr(run_obj)
            spc_val = int(round(ls_px * PT_PER_PX * 100))
            r_pr.set("spc", str(spc_val))
        except Exception:
            pass


def _insert_paragraph_line_break(paragraph) -> None:
    """在段落末尾插入 OOXML 软换行 `<a:br/>`（对应 HTML `<br>`）。"""
    from lxml import etree
    from pptx.oxml.ns import qn

    etree.SubElement(paragraph._p, qn("a:br"))


def _append_text_run_with_breaks(
    paragraph,
    run_data: Dict[str, Any],
    *,
    default_font: Any = None,
    default_size: Any = None,
    default_color: Any = None,
    default_color_theme_key: Any = None,
) -> None:
    """写入单个 run；支持 `break` 标记与文本内 `\\n` 软换行。"""
    if not isinstance(run_data, dict):
        return
    if run_data.get("break") and (run_data.get("text") in (None, "", "\n")):
        _insert_paragraph_line_break(paragraph)
        return

    text = run_data.get("text", "")
    if text is None:
        text = ""
    text = str(text)
    parts = text.split("\n")
    for i, part in enumerate(parts):
        if part:
            run = paragraph.add_run()
            run.text = part
            _apply_run_style(
                run,
                run_data,
                default_font=default_font,
                default_size=default_size,
                default_color=default_color,
                default_color_theme_key=default_color_theme_key,
            )
        if i < len(parts) - 1:
            _insert_paragraph_line_break(paragraph)


def _contains_cjk_text(text: Any) -> bool:
    if not isinstance(text, str):
        return False
    return any(
        "\u3400" <= ch <= "\u9fff"
        or "\uf900" <= ch <= "\ufaff"
        or "\u3040" <= ch <= "\u30ff"
        or "\uac00" <= ch <= "\ud7af"
        for ch in text
    )


def _apply_run_typeface_scripts(run_obj, font_family: str) -> None:
    """按 CSS fallback 契约写 latin/ea/cs，避免打开方做不可控字体替换。"""
    try:
        from lxml import etree
        from pptx.oxml.ns import qn

        r_pr = _ensure_run_rpr(run_obj)
        latin_typeface, east_asian_typeface, complex_script_typeface = _resolve_pptx_typefaces(
            font_family,
            getattr(run_obj, "text", ""),
        )

        for tag_name, typeface in (
            ("latin", latin_typeface),
            ("ea", east_asian_typeface),
            ("cs", complex_script_typeface),
        ):
            font_el = r_pr.find(qn(f"a:{tag_name}"))
            if font_el is None:
                font_el = etree.SubElement(r_pr, qn(f"a:{tag_name}"))
            font_el.set("typeface", typeface)
    except Exception:
        pass


def _apply_text_box_fill(shape, props: dict) -> None:
    """
    将文本元素的 fill 写为文本框 solid 背景（编辑器以 `background: element.fill` 渲染）。

    后端此前从不写出导致有底色文本框被导成透明。fill 为空 / 'transparent' / 'none'
    时保持默认无填充（不覆盖）。PPTTextElement 只有纯色 fill（无主题色 / 渐变，主题色
    在读取时已拍平为 hex），故只处理 solid + alpha。
    """
    fill_color = props.get("fill")
    if not isinstance(fill_color, str) or fill_color.strip().lower() in ("", "transparent", "none"):
        return
    try:
        from pptx.dml.color import RGBColor
        hex_color, fill_alpha = _parse_css_color(fill_color)
        hex_6 = hex_color.lstrip("#")
        if len(hex_6) != 6:
            return
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor.from_string(hex_6)
        if fill_alpha is not None and fill_alpha < 1.0:
            _apply_color_alpha(shape, "fill", fill_alpha)
    except Exception:
        pass


def _apply_text_frame_insets(text_frame, margin: Any) -> None:
    """显式写入文本框四向内边距（bodyPr lIns/tIns/rIns/bIns，单位 EMU）。

    编辑器/浏览器测量的文本框紧贴文字、无内边距；若不显式写 0，PowerPoint 会套用
    默认 0.1in(左右)/0.05in(上下)，导致文字提前换行、纵向撑出固定框（noAutofit）。
    有 margin 时按 margin 写，缺省的边写 0，使 PPT 文本区与来源测量框一致。
    """
    try:
        nsmap = {"a": "http://schemas.openxmlformats.org/drawingml/2006/main"}
        body_pr = text_frame._txBody.find("a:bodyPr", nsmap)
        if body_pr is None:
            return
        margin_map = margin if isinstance(margin, dict) else {}
        for attr, key in (("lIns", "left"), ("tIns", "top"), ("rIns", "right"), ("bIns", "bottom")):
            raw = margin_map.get(key)
            pt = _css_length_to_pt(raw) if raw is not None else None
            emu = int(round(pt * EMU_PER_PT)) if pt is not None else 0
            body_pr.set(attr, str(emu))
    except Exception:
        pass


def _apply_text_word_wrap(text_frame, props: dict) -> None:
    """把浏览器/导入链路的换行契约统一写入 DrawingML。"""
    explicit_word_wrap = props.get("wordWrap")
    if isinstance(explicit_word_wrap, bool):
        text_frame.word_wrap = explicit_word_wrap
        return
    try:
        text_frame.word_wrap = int(props.get("sourceLineCount")) != 1
    except (TypeError, ValueError):
        # 旧数据没有浏览器行数时保持历史行为。
        text_frame.word_wrap = True


def _apply_text_props_to_shape(shape, props: dict, apply_outline: bool = True) -> None:
    """
    将 TextElement props 写入任意可编辑文本 shape（文本框或 placeholder）。
    """
    from pptx.util import Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN

    if not hasattr(shape, "text_frame"):
        return

    tf = shape.text_frame
    tf.clear()
    # HTML 抽取链路记录浏览器最终行数；固定画布中已是一行的文字不能交给
    # PowerPoint 因字体替换再次换行。
    _apply_text_word_wrap(tf, props)

    # 垂直对齐
    v_align = props.get("verticalAlign", "top")
    try:
        nsmap = {"a": "http://schemas.openxmlformats.org/drawingml/2006/main"}
        body_pr = tf._txBody.find("a:bodyPr", nsmap)
        if body_pr is not None:
            anchor_xml = {"top": "t", "middle": "ctr", "bottom": "b"}
            body_pr.set("anchor", anchor_xml.get(v_align, "t"))
    except Exception:
        pass

    # 内边距：始终显式写出，避免 PowerPoint 套用默认 0.1in/0.05in 幽灵内边距。
    _apply_text_frame_insets(tf, props.get("margin"))

    content = props.get("content", "")
    default_font = (
        props.get("fontFamilyFallbacks")
        or props.get("defaultFontFamily")
        or props.get("defaultFontName")
    )
    # defaultFontSize 是 pt 语义（编辑器以 `${defaultFontSize}pt` 渲染、旧前端导出按 pt 透传），
    # 数值缺单位时按 pt 解释，不能当成 px 否则会被 _css_length_to_pt 再缩小 0.75。
    raw_default_size = props.get("defaultFontSize")
    default_size = f"{raw_default_size}pt" if isinstance(raw_default_size, (int, float)) and raw_default_size > 0 else raw_default_size
    default_color = props.get("defaultColor")
    default_color_theme_key = props.get("defaultColorThemeKey")

    paragraphs = _parse_html_to_paragraphs(content)
    if not paragraphs:
        from html import unescape as _html_unescape
        plain = _html_unescape(str(content).replace("&nbsp;", " "))
        paragraphs = [{"runs": [{"text": plain, "bold": False, "italic": False, "underline": False}]}]

    align_map = {
        "left": PP_ALIGN.LEFT,
        "center": PP_ALIGN.CENTER,
        "right": PP_ALIGN.RIGHT,
        "justify": PP_ALIGN.JUSTIFY,
    }

    default_text_align = props.get("defaultTextAlign")
    # 元素级默认行距（编辑器 lineHeight 为无单位倍数），段落未显式声明行距时回退使用
    default_line_height = _parse_css_line_height_ratio(props.get("lineHeight"))

    for p_idx, para_data in enumerate(paragraphs):
        if p_idx == 0:
            p = tf.paragraphs[0] if tf.paragraphs else tf.add_paragraph()
        else:
            p = tf.add_paragraph()

        # 段落对齐：HTML 内联样式优先，否则用 defaultTextAlign
        p_align = para_data.get("align")
        if p_align and p_align in align_map:
            p.alignment = align_map[p_align]
        elif default_text_align and default_text_align in align_map:
            p.alignment = align_map[default_text_align]

        # 段落行距：HTML 内联 line-height 优先，否则回退到元素级默认行距
        p_line_height = para_data.get("lineHeight")
        if p_line_height:
            lh_val = _parse_css_line_height_ratio(p_line_height)
            if lh_val is not None and lh_val > 0:
                p.line_spacing = lh_val
        elif default_line_height is not None and default_line_height > 0:
            p.line_spacing = default_line_height

        # 段前间距（来自 HTML <p> 的 margin-top，如 "6pt"）
        p_margin_top = para_data.get("marginTop")
        if p_margin_top:
            mt_pt = _css_length_to_pt(p_margin_top)
            if mt_pt is not None and mt_pt > 0:
                p.space_before = Pt(round(mt_pt, 3))

        # 段后间距（来自 HTML <p> 的 margin-bottom，如 "6pt"）
        p_margin_bottom = para_data.get("marginBottom")
        if p_margin_bottom:
            mb_pt = _css_length_to_pt(p_margin_bottom)
            if mb_pt is not None and mb_pt > 0:
                p.space_after = Pt(round(mb_pt, 3))

        # 段落缩进（padding-left → marL, text-indent → indent）
        p_padding_left = para_data.get("paddingLeft")
        p_text_indent = para_data.get("textIndent")
        if p_padding_left or p_text_indent:
            try:
                nsmap_a_pi = "http://schemas.openxmlformats.org/drawingml/2006/main"
                pPr_pi = p._p.find(f"{{{nsmap_a_pi}}}pPr")
                if pPr_pi is None:
                    from lxml import etree
                    pPr_pi = etree.SubElement(p._p, f"{{{nsmap_a_pi}}}pPr")
                    p._p.insert(0, pPr_pi)
                if p_padding_left:
                    pl_px = _css_length_to_px(p_padding_left)
                    if pl_px is not None:
                        mar_l_emu = int(round(pl_px * PT_PER_PX * 12700))
                        pPr_pi.set("marL", str(mar_l_emu))
                if p_text_indent:
                    ti_pt = _css_length_to_pt(p_text_indent)
                    if ti_pt is not None:
                        indent_emu = int(round(ti_pt * 12700))
                        pPr_pi.set("indent", str(indent_emu))
            except Exception:
                pass

        runs = para_data.get("runs", [])
        if not runs:
            run = p.add_run()
            run.text = ""
            continue

        # 列表项目符号
        bullet_type = para_data.get("bullet")
        if bullet_type:
            try:
                from lxml import etree
                nsmap_a = "http://schemas.openxmlformats.org/drawingml/2006/main"
                pPr = p._p.find(f"{{{nsmap_a}}}pPr")
                if pPr is None:
                    pPr = etree.SubElement(p._p, f"{{{nsmap_a}}}pPr")
                    # 确保 pPr 在 runs 之前
                    p._p.insert(0, pPr)
                # 设置缩进级别
                level = para_data.get("level", 0)
                if level > 0:
                    pPr.set("lvl", str(level))
                    pPr.set("marL", str(level * 457200))  # 每级 0.5 inch
                    pPr.set("indent", str(-228600))  # 悬挂缩进
                # 项目符号样式（颜色/大小/字体）
                _bu_sty = para_data.get("bulletStyle")
                if _bu_sty:
                    _bu_clr = _bu_sty.get("color")
                    if _bu_clr:
                        _buClrEl = etree.SubElement(pPr, f"{{{nsmap_a}}}buClr")
                        _srgbEl = etree.SubElement(_buClrEl, f"{{{nsmap_a}}}srgbClr")
                        _srgbEl.set("val", _bu_clr.lstrip("#")[:6].upper())
                    _bu_fsz = _bu_sty.get("fontSize")
                    if _bu_fsz:
                        if _bu_fsz.endswith("%"):
                            _pct = int(_bu_fsz.rstrip("%")) * 1000
                            _szEl = etree.SubElement(pPr, f"{{{nsmap_a}}}buSzPct")
                            _szEl.set("val", str(_pct))
                        elif _bu_fsz.endswith("pt"):
                            _pts = int(float(_bu_fsz.rstrip("pt")) * 100)
                            _szEl = etree.SubElement(pPr, f"{{{nsmap_a}}}buSzPts")
                            _szEl.set("val", str(_pts))
                    _bu_fnt = _bu_sty.get("fontFamily")
                    if _bu_fnt:
                        _buFntEl = etree.SubElement(pPr, f"{{{nsmap_a}}}buFont")
                        _buFntEl.set("typeface", _bu_fnt)
                if bullet_type == "bullet":
                    bu_char = etree.SubElement(pPr, f"{{{nsmap_a}}}buChar")
                    bu_char.set("char", para_data.get("bulletChar") or "\u2022")
                elif bullet_type == "number":
                    bu_auto = etree.SubElement(pPr, f"{{{nsmap_a}}}buAutoNum")
                    num_fmt = para_data.get("numberFormat") or "arabicPeriod"
                    bu_auto.set("type", num_fmt)
            except Exception:
                pass

        for run_data in runs:
            _append_text_run_with_breaks(
                p,
                run_data,
                default_font=default_font,
                default_size=default_size,
                default_color=default_color,
                default_color_theme_key=default_color_theme_key,
            )

    # 段间距：统一由元素级 paragraphSpace 控制（B1-02 修复后不再有内联 margin-bottom）
    paragraph_space = props.get("paragraphSpace")
    paragraph_space_pt = _css_length_to_pt(paragraph_space)
    if paragraph_space_pt is not None and paragraph_space_pt > 0:
        for p_ps in tf.paragraphs:
            try:
                # 兼容旧数据：如果该段已有独立段后间距（存量 HTML margin-bottom），则跳过
                if p_ps.space_after is not None and p_ps.space_after > 0:
                    continue
                p_ps.space_after = Pt(round(paragraph_space_pt, 3))
            except Exception:
                pass

    # 竖排文字
    if props.get("vertical"):
        try:
            nsmap = {"a": "http://schemas.openxmlformats.org/drawingml/2006/main"}
            body_pr = tf._txBody.find("a:bodyPr", nsmap)
            if body_pr is not None:
                body_pr.set("vert", "eaVert")
        except Exception:
            pass

    # 字间距（应用到所有 run）
    word_space = props.get("wordSpace")
    word_space_px = _css_length_to_px(word_space)
    if word_space_px is not None and abs(word_space_px) > 1e-9:
        try:
            nsmap_a_ws = "http://schemas.openxmlformats.org/drawingml/2006/main"
            spc_val = str(int(round(word_space_px * PT_PER_PX * 100)))
            for p in tf.paragraphs:
                for run in p.runs:
                    rPr = run._r.find(f"{{{nsmap_a_ws}}}rPr")
                    if rPr is None:
                        from lxml import etree
                        rPr = etree.SubElement(run._r, f"{{{nsmap_a_ws}}}rPr")
                        run._r.insert(0, rPr)
                    rPr.set("spc", spc_val)
        except Exception:
            pass

    # 文本框背景填充与边框（placeholder 回填时不覆盖其底色/边框）
    if apply_outline:
        _apply_text_box_fill(shape, props)
        outline = props.get("outline")
        if outline and isinstance(outline, dict):
            try:
                ln = shape.line
                if outline.get("width"):
                    ln.width = Pt(float(outline["width"]))
                if outline.get("color"):
                    hex_out, _ = _parse_css_color(outline["color"])
                    ln.color.rgb = RGBColor.from_string(hex_out.lstrip("#")[:6])
                _dash_enum = {"dashed": 2, "dotted": 3, "dashDot": 5, "longDash": 6, "longDashDot": 7}
                if outline.get("style") in _dash_enum:
                    ln.dash_style = _dash_enum[outline["style"]]
            except Exception:
                pass


def _write_text_element(slide, props: dict, left, top, width, height):
    """写入文本元素（精确还原 HTML 富文本格式，含段落对齐）。返回 shape 对象。"""
    tx_box = slide.shapes.add_textbox(left, top, width, height)
    _apply_text_props_to_shape(tx_box, props, apply_outline=True)
    tf = tx_box.text_frame
    _apply_text_autofit(tf, props.get("autoFit"))
    return tx_box


# 编辑器 autoFit 语义 → OOXML bodyPr autofit 子元素：
# - 'shrink'：溢出时缩小字体  → normAutofit
# - 'resize'：调整文本框适应文字 → spAutoFit
# - 未设置：默认 shrink-to-fit → normAutofit
#
# 默认用 normAutofit 而非 noAutofit：我们在服务端用 Noto Sans SC 量框，目标机器缺该字体时
# 回退到更宽的字体（如微软雅黑）→ 同一行文字超出框宽 → 换行/撑破框 → 跑版。
# normAutofit 让打开的应用在文字溢出时自动等比缩小字号塞回框内，既不换行跑版、又保留文字
# 可编辑（不同于转曲线）；字体存在时不触发缩放，渲染与编辑器一致。
_AUTOFIT_TAG_MAP = {
    "shrink": "normAutofit",
    "resize": "spAutoFit",
}
_DEFAULT_AUTOFIT_TAG = "normAutofit"


def _apply_text_autofit(text_frame, auto_fit: Any) -> None:
    """按编辑器 autoFit 语义写入 bodyPr 的 autofit 子元素；未设置时默认 shrink-to-fit。"""
    try:
        from lxml import etree
        nsmap_a = 'http://schemas.openxmlformats.org/drawingml/2006/main'
        body_pr = text_frame._txBody.find(f'{{{nsmap_a}}}bodyPr')
        if body_pr is None:
            return
        # 清掉历史/默认的 autofit 子元素与无效属性后重新写入
        body_pr.attrib.pop('autoFit', None)
        for child in list(body_pr):
            tag = etree.QName(child.tag).localname if isinstance(child.tag, str) else ''
            if tag in ('noAutofit', 'normAutofit', 'spAutoFit'):
                body_pr.remove(child)
        fit_tag = _AUTOFIT_TAG_MAP.get(str(auto_fit or '').strip(), _DEFAULT_AUTOFIT_TAG)
        etree.SubElement(body_pr, f'{{{nsmap_a}}}{fit_tag}')
    except Exception:
        pass


def _parse_html_to_paragraphs(html: str) -> List[Dict[str, Any]]:
    """
    解析 HTML 富文本为段落+运行结构（含段落对齐、行距、间距、列表类型）。

    使用 lxml 解析，正确处理嵌套 <ul>/<ol> 列表结构。
    支持的结构：
    - <p style="text-align:center; margin-top:6pt">...</p> → 普通段落
    - <ul><li><p>...</p><ul><li>嵌套...</li></ul></li></ul> → 嵌套无序列表
    - <ol><li><p>...</p></li></ol> → 有序列表
    - <a href="...">text</a> → 超链接（在 run 中标记 hyperlink）
    """
    import re
    from html import unescape

    if not html or not html.strip():
        return []

    paragraphs: List[Dict[str, Any]] = []

    def _extract_attr(attrs_str: str, attr_name: str) -> Optional[str]:
        """提取 HTML 属性值，支持单/双引号与无引号写法。"""
        attr_pat = re.compile(
            rf"{re.escape(attr_name)}\s*=\s*(['\"])(.*?)\1",
            re.IGNORECASE | re.DOTALL,
        )
        m = attr_pat.search(attrs_str or "")
        if m:
            return m.group(2).strip()

        # 无引号兜底：style=color:red
        attr_pat_unquoted = re.compile(
            rf"{re.escape(attr_name)}\s*=\s*([^\s>]+)",
            re.IGNORECASE,
        )
        m_unquoted = attr_pat_unquoted.search(attrs_str or "")
        if m_unquoted:
            return m_unquoted.group(1).strip()
        return None

    def _parse_p_style(style_str: str) -> Dict[str, Any]:
        """从 <p> 的 style 属性解析段落级样式"""
        result: Dict[str, Any] = {}
        if not style_str:
            return result
        for prop in style_str.split(";"):
            prop = prop.strip()
            if ":" not in prop:
                continue
            key, val = prop.split(":", 1)
            key = key.strip().lower()
            val = val.strip()
            if key == "text-align":
                result["align"] = val
            elif key == "line-height":
                result["lineHeight"] = val
            elif key == "margin-top":
                result["marginTop"] = val
            elif key == "margin-bottom":
                result["marginBottom"] = val
            elif key == "padding-left":
                result["paddingLeft"] = val
            elif key == "text-indent":
                result["textIndent"] = val
        return result

    def _p_element_to_para(
        p_html_inner: str,
        p_style_str: str,
        bullet_type: Optional[str] = None,
        level: int = 0,
        number_format: Optional[str] = None,
        bullet_char: Optional[str] = None,
        bullet_style: Optional[Dict[str, str]] = None,
    ) -> Dict[str, Any]:
        """将 <p> 标签内容和样式转换为段落数据"""
        style_data = _parse_p_style(p_style_str)
        runs = _parse_paragraph_html(p_html_inner)
        result: Dict[str, Any] = {
            "runs": runs,
            "align": style_data.get("align"),
            "lineHeight": style_data.get("lineHeight"),
        }
        # 段前/段后间距
        if style_data.get("marginTop"):
            result["marginTop"] = style_data["marginTop"]
        if style_data.get("marginBottom"):
            result["marginBottom"] = style_data["marginBottom"]
        # 缩进
        if style_data.get("paddingLeft"):
            result["paddingLeft"] = style_data["paddingLeft"]
        if style_data.get("textIndent"):
            result["textIndent"] = style_data["textIndent"]
        # 列表类型和级别
        if bullet_type:
            result["bullet"] = bullet_type
            result["level"] = level
        if number_format:
            result["numberFormat"] = number_format
        if bullet_char:
            result["bulletChar"] = bullet_char
        if bullet_style:
            result["bulletStyle"] = bullet_style
        return result

    # ── 使用 lxml 解析，正确处理嵌套列表 ──
    try:
        from lxml import etree

        # 包裹一层 <root> 确保是合法 XML
        root = etree.fromstring(f"<root>{html}</root>")

        def _collect(node, bullet_type=None, level=0, number_format=None, bullet_char=None, bullet_style=None):
            """递归遍历 XML 树，收集段落数据"""
            for child in node:
                tag = child.tag.lower() if isinstance(child.tag, str) else ""
                if tag == "p":
                    style_str = child.get("style", "")
                    inner_html = etree.tostring(child, method="html", encoding="unicode")
                    inner_html = re.sub(r"^<p[^>]*>", "", inner_html, count=1)
                    inner_html = re.sub(r"</p>\s*$", "", inner_html, count=1)
                    paragraphs.append(_p_element_to_para(inner_html, style_str, bullet_type, level, number_format, bullet_char, bullet_style))
                elif tag == "ul":
                    bc = child.get("data-bullet-char") or bullet_char
                    # 提取项目符号样式
                    _bs = {}
                    _bc_color = child.get("data-bullet-color")
                    _bc_fsize = child.get("data-bullet-font-size")
                    _bc_font = child.get("data-bullet-font")
                    if _bc_color:
                        _bs["color"] = _bc_color
                    if _bc_fsize:
                        _bs["fontSize"] = _bc_fsize
                    if _bc_font:
                        _bs["fontFamily"] = _bc_font
                    _bstyle = _bs if _bs else bullet_style
                    _collect(child, bullet_type="bullet", level=level, bullet_char=bc, bullet_style=_bstyle)
                elif tag == "ol":
                    ol_type = child.get("type")
                    ol_num_fmt = _OL_TYPE_TO_NUM_FMT.get(ol_type, "arabicPeriod") if ol_type else number_format
                    _collect(child, bullet_type="number", level=level, number_format=ol_num_fmt)
                elif tag == "li":
                    has_block_child = False
                    for li_child in child:
                        li_tag = li_child.tag.lower() if isinstance(li_child.tag, str) else ""
                        if li_tag == "p":
                            has_block_child = True
                            style_str = li_child.get("style", "")
                            inner_html = etree.tostring(li_child, method="html", encoding="unicode")
                            inner_html = re.sub(r"^<p[^>]*>", "", inner_html, count=1)
                            inner_html = re.sub(r"</p>\s*$", "", inner_html, count=1)
                            paragraphs.append(_p_element_to_para(inner_html, style_str, bullet_type, level, number_format, bullet_char, bullet_style))
                        elif li_tag in ("ul", "ol"):
                            has_block_child = True
                            inner_type = "bullet" if li_tag == "ul" else "number"
                            inner_fmt = None
                            inner_bc = None
                            inner_bstyle = None
                            if li_tag == "ol":
                                inner_ol_type = li_child.get("type")
                                inner_fmt = _OL_TYPE_TO_NUM_FMT.get(inner_ol_type, "arabicPeriod") if inner_ol_type else number_format
                            elif li_tag == "ul":
                                inner_bc = li_child.get("data-bullet-char") or bullet_char
                                # 提取嵌套层的项目符号样式
                                _nbs = {}
                                _nbc_color = li_child.get("data-bullet-color")
                                _nbc_fsize = li_child.get("data-bullet-font-size")
                                _nbc_font = li_child.get("data-bullet-font")
                                if _nbc_color:
                                    _nbs["color"] = _nbc_color
                                if _nbc_fsize:
                                    _nbs["fontSize"] = _nbc_fsize
                                if _nbc_font:
                                    _nbs["fontFamily"] = _nbc_font
                                inner_bstyle = _nbs if _nbs else bullet_style
                            _collect(li_child, bullet_type=inner_type, level=level + 1, number_format=inner_fmt, bullet_char=inner_bc, bullet_style=inner_bstyle)
                    if not has_block_child:
                        li_inner = etree.tostring(child, method="html", encoding="unicode")
                        li_inner = re.sub(r"^<li[^>]*>", "", li_inner, count=1)
                        li_inner = re.sub(r"</li>\s*$", "", li_inner, count=1)
                        if li_inner.strip():
                            paragraphs.append(_p_element_to_para(li_inner.strip(), "", bullet_type, level, number_format, bullet_char, bullet_style))

        _collect(root)

        if paragraphs:
            return paragraphs
    except Exception as e:
        logger.debug(f"lxml HTML parse failed, falling back to regex: {e}")

    # ── Fallback: 简单 regex（处理非嵌套情况） ──
    p_pattern = re.compile(r"<p([^>]*)>(.*?)</p>", re.DOTALL | re.IGNORECASE)
    for m in p_pattern.finditer(html):
        style_str = _extract_attr(m.group(1), "style") or ""
        paragraphs.append(_p_element_to_para(m.group(2).strip(), style_str))

    if not paragraphs:
        runs = _parse_paragraph_html(html)
        paragraphs.append({"runs": runs, "align": None, "lineHeight": None})

    return paragraphs


def _parse_paragraph_html(html: str) -> List[Dict[str, Any]]:
    """
    解析单个段落内的 HTML 为 run 列表。

    处理嵌套标签：<span style="..."><b><i>text</i></b></span>
    支持超链接：<a href="...">text</a> → run 中包含 hyperlink 字段
    """
    import re
    from html import unescape

    if not html:
        return [{"text": ""}]

    runs: List[Dict[str, Any]] = []

    # 使用简单的状态机解析（不用 html.parser 避免过度复杂）
    # 策略：找到所有文本节点，并记录其周围的标签状态
    pos = 0
    current_styles: Dict[str, Any] = {}
    _missing_style = object()
    # 记录 span 进入前被覆盖的样式，关闭 span 时可精确恢复外层样式。
    span_restore_stack: List[Dict[str, Any]] = []
    hyperlink_stack: List[Optional[str]] = []  # 超链接 URL 栈

    while pos < len(html):
        # 找下一个标签
        tag_start = html.find("<", pos)
        if tag_start == -1:
            # 剩余全是文本
            text = unescape(html[pos:])
            if text:
                runs.append(_make_run(text, current_styles))
            break

        # 标签前的文本
        if tag_start > pos:
            text = unescape(html[pos:tag_start])
            if text:
                runs.append(_make_run(text, current_styles))

        tag_end = html.find(">", tag_start)
        if tag_end == -1:
            break

        tag_content = html[tag_start + 1:tag_end]
        pos = tag_end + 1

        # 解析标签
        is_closing = tag_content.startswith("/")
        if is_closing:
            tag_name = tag_content[1:].strip().lower().split()[0] if tag_content[1:].strip() else ""
            # 弹出样式
            if tag_name in ("b", "strong"):
                current_styles.pop("bold", None)
            elif tag_name in ("i", "em"):
                current_styles.pop("italic", None)
            elif tag_name == "u":
                current_styles.pop("underline", None)
                current_styles.pop("underlineStyle", None)
            elif tag_name in ("s", "del", "strike"):
                current_styles.pop("strikethrough", None)
            elif tag_name == "sup":
                current_styles.pop("baseline", None)
            elif tag_name == "sub":
                current_styles.pop("baseline", None)
            elif tag_name == "mark":
                current_styles.pop("highlight", None)
            elif tag_name == "span":
                if span_restore_stack:
                    restore_map = span_restore_stack.pop()
                    for k, prev_val in restore_map.items():
                        if prev_val is _missing_style:
                            current_styles.pop(k, None)
                        else:
                            current_styles[k] = prev_val
            elif tag_name == "a":
                if hyperlink_stack:
                    hyperlink_stack.pop()
                current_styles.pop("hyperlink", None)
        else:
            tag_parts = tag_content.strip().rstrip("/").split(None, 1)
            tag_name = tag_parts[0].lower() if tag_parts else ""
            tag_attrs = tag_parts[1] if len(tag_parts) > 1 else ""

            if tag_name in ("b", "strong"):
                current_styles["bold"] = True
            elif tag_name in ("i", "em"):
                current_styles["italic"] = True
            elif tag_name == "u":
                current_styles["underline"] = True
                # 保留 OOXML 下划线样式
                us_match = re.search(
                    r'data-underline-style\s*=\s*["\']([^"\']+)["\']',
                    tag_attrs,
                    re.IGNORECASE,
                )
                if us_match:
                    current_styles["underlineStyle"] = us_match.group(1).strip()
            elif tag_name in ("s", "del", "strike"):
                current_styles["strikethrough"] = True
            elif tag_name == "sup":
                current_styles["baseline"] = "super"
            elif tag_name == "sub":
                current_styles["baseline"] = "sub"
            elif tag_name == "mark":
                # 提取 data-color 或 style 中的 background-color
                mark_color = None
                data_color_match = re.search(
                    r"data-color\s*=\s*(['\"])(.*?)\1",
                    tag_attrs,
                    re.IGNORECASE | re.DOTALL,
                )
                if data_color_match:
                    mark_color = data_color_match.group(2).strip()
                if not mark_color:
                    bg_match = re.search(r'background-color:\s*([^;"]+)', tag_attrs, re.IGNORECASE)
                    if bg_match:
                        mark_color = bg_match.group(1).strip()
                if mark_color:
                    current_styles["highlight"] = mark_color
            elif tag_name == "span":
                span_styles = _parse_inline_style(tag_attrs)
                restore_map: Dict[str, Any] = {}
                for k, v in span_styles.items():
                    restore_map[k] = current_styles.get(k, _missing_style)
                    current_styles[k] = v
                span_restore_stack.append(restore_map)
            elif tag_name == "a":
                # 提取 href 属性
                href_match = re.search(
                    r"href\s*=\s*(['\"])(.*?)\1",
                    tag_attrs,
                    re.IGNORECASE | re.DOTALL,
                )
                href = unescape(href_match.group(2).strip()) if href_match else None
                hyperlink_stack.append(href)
                if href:
                    current_styles["hyperlink"] = href
            elif tag_name == "br":
                # 软换行：写入特殊 run，_write_text_element 再落 a:br
                runs.append({"text": "\n", "break": True})

    return runs if runs else [{"text": ""}]


def _make_run(text: str, styles: Dict[str, Any]) -> Dict[str, Any]:
    """创建一个 run 数据"""
    run: Dict[str, Any] = {"text": text}
    if styles.get("bold"):
        run["bold"] = True
    if styles.get("italic"):
        run["italic"] = True
    if styles.get("underline"):
        run["underline"] = True
        if styles.get("underlineStyle"):
            run["underlineStyle"] = styles["underlineStyle"]
    if styles.get("strikethrough"):
        run["strikethrough"] = True
    if styles.get("color"):
        run["color"] = styles["color"]
    if styles.get("fontSize"):
        run["fontSize"] = styles["fontSize"]
    if styles.get("fontFamily"):
        run["fontFamily"] = styles["fontFamily"]
    if styles.get("hyperlink"):
        run["hyperlink"] = styles["hyperlink"]
    if styles.get("baseline"):
        run["baseline"] = styles["baseline"]
    if styles.get("highlight"):
        run["highlight"] = styles["highlight"]
    if styles.get("themeColorKey"):
        run["themeColorKey"] = styles["themeColorKey"]
    if styles.get("letterSpacing") is not None:
        run["letterSpacing"] = styles["letterSpacing"]
    return run


def _parse_inline_style(attrs_str: str) -> Dict[str, Any]:
    """解析 HTML 标签的 style 属性，提取字体样式信息"""
    import re

    result: Dict[str, Any] = {}
    # 提取 style="..." 或 style='...'（支持值内包含另一种引号）
    style_match = re.search(
        r"style\s*=\s*(['\"])(.*?)\1",
        attrs_str,
        re.IGNORECASE | re.DOTALL,
    )
    if style_match:
        style_str = style_match.group(2).strip()

        for prop in style_str.split(";"):
            prop = prop.strip()
            if not prop:
                continue
            if ":" not in prop:
                continue
            key, val = prop.split(":", 1)
            key = key.strip().lower()
            val = val.strip()

            if key == "color":
                result["color"] = val
            elif key == "font-size":
                # 统一转换为 pt，避免 px/pt 混用导致字号漂移
                size_pt = _css_length_to_pt(val)
                if size_pt is not None and size_pt > 0:
                    result["fontSize"] = round(size_pt, 3)
            elif key == "font-family":
                result["fontFamily"] = val
            elif key == "letter-spacing":
                # 前端约定 letter-spacing 为 px
                ls_px = _css_length_to_px(val)
                if ls_px is not None:
                    result["letterSpacing"] = round(ls_px, 3)
            elif key == "font-weight":
                # 与 dom_extractor 一致：≥600 视为粗体（OOXML 无 800 档，只能 b=1）
                if val in ("bold", "bolder") or (val.isdigit() and int(val) >= 600):
                    result["bold"] = True
            elif key == "font-style":
                if val == "italic":
                    result["italic"] = True
            elif key == "text-decoration":
                if "underline" in val:
                    result["underline"] = True
                if "line-through" in val:
                    result["strikethrough"] = True

    # 提取 data-theme-color-key（支持单双引号）
    theme_key_match = re.search(
        r"data-theme-color-key\s*=\s*(['\"])(.*?)\1",
        attrs_str,
        re.IGNORECASE | re.DOTALL,
    )
    if theme_key_match:
        theme_key = _normalize_text_theme_key(theme_key_match.group(2))
        if theme_key:
            result["themeColorKey"] = theme_key

    return result


def _write_cell_rich_text(cell, rich_text_html: str, cell_data: dict) -> None:
    """
    将 richText HTML 写入表格单元格（多段落、多 run 格式）。

    复用 _parse_html_to_paragraphs 解析 HTML，
    然后逐段逐 run 写入 cell.text_frame。
    """
    from pptx.util import Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    from pptx.oxml.ns import qn
    from lxml import etree

    paragraphs = _parse_html_to_paragraphs(rich_text_html)
    if not paragraphs:
        cell.text = cell_data.get("text", "")
        return

    # 清空单元格默认文本
    cell.text = ""
    tf = cell.text_frame

    align_map = {
        "left": PP_ALIGN.LEFT,
        "center": PP_ALIGN.CENTER,
        "right": PP_ALIGN.RIGHT,
        "justify": PP_ALIGN.JUSTIFY,
    }

    # 从 cell_data 取默认样式（作为 run 样式的 fallback）
    style = cell_data.get("style") or {}
    default_font = style.get("fontName") or style.get("fontFamily") or cell_data.get("fontName") or cell_data.get("fontFamily")
    default_size = style.get("fontSize") or cell_data.get("fontSize")
    default_color = style.get("color") or cell_data.get("color")
    default_color_theme_key = style.get("colorThemeKey") or cell_data.get("colorThemeKey")
    default_align = style.get("align") or cell_data.get("align")

    def _apply_run_color(run_obj, color_str: Any, theme_key: Optional[Any] = None) -> None:
        normalized_theme_key = _normalize_text_theme_key(theme_key)
        if normalized_theme_key and _apply_run_theme_color(run_obj, normalized_theme_key, color_str):
            return
        if not color_str:
            return
        hex_fc, alpha_fc = _parse_css_color(str(color_str))
        hex_6 = hex_fc.lstrip("#")[:6]
        if len(hex_6) != 6:
            return

        # 基础色通过高层 API 写入
        run_obj.font.color.rgb = RGBColor.from_string(hex_6)

        # alpha 通过底层 XML 写入
        if alpha_fc is None or alpha_fc >= 1.0:
            return
        nsmap_a = "http://schemas.openxmlformats.org/drawingml/2006/main"
        rPr = run_obj._r.find(f"{{{nsmap_a}}}rPr")
        if rPr is None:
            rPr = etree.SubElement(run_obj._r, qn("a:rPr"))
            run_obj._r.insert(0, rPr)

        solid_fill = rPr.find(f"{{{nsmap_a}}}solidFill")
        if solid_fill is None:
            solid_fill = etree.SubElement(rPr, qn("a:solidFill"))

        color_el = solid_fill.find(f"{{{nsmap_a}}}srgbClr")
        if color_el is None:
            color_el = etree.SubElement(solid_fill, qn("a:srgbClr"))
        color_el.set("val", hex_6)

        for old_alpha in color_el.findall(f"{{{nsmap_a}}}alpha"):
            color_el.remove(old_alpha)
        alpha_el = etree.SubElement(color_el, qn("a:alpha"))
        alpha_el.set("val", str(int(max(0.0, min(1.0, float(alpha_fc))) * 100000)))

    for p_idx, para_data in enumerate(paragraphs):
        if p_idx == 0:
            p = tf.paragraphs[0] if tf.paragraphs else tf.add_paragraph()
        else:
            p = tf.add_paragraph()

        # 段落对齐
        p_align = para_data.get("align")
        if p_align and p_align in align_map:
            p.alignment = align_map[p_align]
        elif default_align and default_align in align_map:
            p.alignment = align_map[default_align]

        # 段落行距
        p_line_height = para_data.get("lineHeight")
        if p_line_height:
            try:
                lh_val = float(p_line_height)
                if lh_val > 0:
                    p.line_spacing = lh_val
            except (ValueError, TypeError):
                pass

        # 列表项目符号（与文本框写入逻辑对齐）
        bullet_type = para_data.get("bullet")
        if bullet_type:
            try:
                nsmap_a_bu = "http://schemas.openxmlformats.org/drawingml/2006/main"
                pPr_bu = p._p.find(f"{{{nsmap_a_bu}}}pPr")
                if pPr_bu is None:
                    pPr_bu = etree.SubElement(p._p, f"{{{nsmap_a_bu}}}pPr")
                    p._p.insert(0, pPr_bu)
                level = para_data.get("level", 0)
                if level > 0:
                    pPr_bu.set("lvl", str(level))
                    pPr_bu.set("marL", str(level * 457200))
                    pPr_bu.set("indent", str(-228600))
                _bu_sty = para_data.get("bulletStyle")
                if _bu_sty:
                    _bu_clr = _bu_sty.get("color")
                    if _bu_clr:
                        _buClrEl = etree.SubElement(pPr_bu, f"{{{nsmap_a_bu}}}buClr")
                        _srgbEl = etree.SubElement(_buClrEl, f"{{{nsmap_a_bu}}}srgbClr")
                        _srgbEl.set("val", _bu_clr.lstrip("#")[:6].upper())
                    _bu_fsz = _bu_sty.get("fontSize")
                    if _bu_fsz:
                        if _bu_fsz.endswith("%"):
                            _pct = int(_bu_fsz.rstrip("%")) * 1000
                            _szEl = etree.SubElement(pPr_bu, f"{{{nsmap_a_bu}}}buSzPct")
                            _szEl.set("val", str(_pct))
                        elif _bu_fsz.endswith("pt"):
                            _pts = int(float(_bu_fsz.rstrip("pt")) * 100)
                            _szEl = etree.SubElement(pPr_bu, f"{{{nsmap_a_bu}}}buSzPts")
                            _szEl.set("val", str(_pts))
                    _bu_fnt = _bu_sty.get("fontFamily")
                    if _bu_fnt:
                        _buFntEl = etree.SubElement(pPr_bu, f"{{{nsmap_a_bu}}}buFont")
                        _buFntEl.set("typeface", _bu_fnt)
                if bullet_type == "bullet":
                    bu_char_el = etree.SubElement(pPr_bu, f"{{{nsmap_a_bu}}}buChar")
                    bu_char_el.set("char", para_data.get("bulletChar") or "\u2022")
                elif bullet_type == "number":
                    bu_auto_el = etree.SubElement(pPr_bu, f"{{{nsmap_a_bu}}}buAutoNum")
                    num_fmt = para_data.get("numberFormat") or "arabicPeriod"
                    bu_auto_el.set("type", num_fmt)
            except Exception:
                pass

        runs = para_data.get("runs", [])
        if not runs:
            run = p.add_run()
            run.text = ""
            continue

        for run_data in runs:
            run = p.add_run()
            run.text = run_data.get("text", "")

            font = run.font
            font_name = run_data.get("fontFamily") or default_font
            font_size = run_data.get("fontSize") or default_size
            font_color = run_data.get("color") or default_color
            run_theme_key = run_data.get("themeColorKey") or default_color_theme_key

            if font_name:
                primary_font_name = _strip_font_fallback(str(font_name))
                if primary_font_name:
                    font.name = primary_font_name
            if font_size:
                font_size_pt = _css_length_to_pt(font_size)
                if font_size_pt is not None and font_size_pt > 0:
                    font.size = Pt(round(font_size_pt, 3))
            if font_color or run_theme_key:
                _apply_run_color(run, font_color, run_theme_key)

            if run_data.get("bold"):
                font.bold = True
            if run_data.get("italic"):
                font.italic = True
            if run_data.get("underline"):
                font.underline = True

            # 删除线
            if run_data.get("strikethrough"):
                try:
                    nsmap_a_s = "http://schemas.openxmlformats.org/drawingml/2006/main"
                    rPr = run._r.find(f"{{{nsmap_a_s}}}rPr")
                    if rPr is None:
                        from lxml import etree
                        rPr = etree.SubElement(run._r, f"{{{nsmap_a_s}}}rPr")
                        run._r.insert(0, rPr)
                    rPr.set("strike", "sngStrike")
                except Exception:
                    pass

            # 超链接
            _apply_run_hyperlink_write(run, run_data.get("hyperlink"))


def _write_table_element(slide, props: dict, left, top, width, height):
    """写入表格元素（含单元格样式、列宽、合并单元格、边框、行高）。返回 shape 对象。"""
    from pptx.util import Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    from pptx.oxml.ns import qn
    from lxml import etree

    data = props.get("data", [])
    if not data:
        return None

    rows = len(data)
    cols = max(len(row) for row in data) if data else 0
    if rows == 0 or cols == 0:
        return None

    def _to_non_negative_float(raw: Any, default: float) -> float:
        try:
            val = float(raw)
            if math.isfinite(val) and val >= 0:
                return val
        except (TypeError, ValueError):
            pass
        return default

    def _to_span_int(raw: Any, default: int = 1) -> int:
        try:
            val = int(float(raw))
            if val >= 0:
                return val
        except (TypeError, ValueError):
            pass
        return default

    table_shape = slide.shapes.add_table(rows, cols, left, top, width, height)
    table = table_shape.table

    # 列宽
    col_widths = props.get("colWidths")
    if isinstance(col_widths, list) and cols > 0:
        safe: List[float] = []
        for i in range(cols):
            raw = col_widths[i] if i < len(col_widths) else 1
            try:
                v = float(raw)
                safe.append(v if v >= 0 and v == v else 0.0)
            except (TypeError, ValueError):
                safe.append(0.0)

        ratio_sum = sum(safe)
        if ratio_sum <= 0:
            safe = [1.0 / cols for _ in range(cols)]
        else:
            safe = [v / ratio_sum for v in safe]

        total_width_emu = int(width)
        consumed = 0
        for i in range(cols):
            if i == cols - 1:
                col_w = max(total_width_emu - consumed, 1)
            else:
                col_w = int(round(total_width_emu * safe[i]))
                col_w = max(col_w, 1)
                consumed += col_w
            table.columns[i].width = Emu(col_w)

    # 行高：优先使用 rowHeights（逐行高度），否则回退到 cellMinHeight
    applied_row_heights = False
    row_heights = props.get("rowHeights")
    if isinstance(row_heights, list) and rows > 0:
        safe_heights: List[float] = []
        for i in range(rows):
            raw = row_heights[i] if i < len(row_heights) else None
            safe_heights.append(_to_non_negative_float(raw, 0.0))

        if any(h > 0 for h in safe_heights):
            cell_min_height = _to_non_negative_float(props.get("cellMinHeight"), 0.0)
            fallback_height = cell_min_height if cell_min_height > 0 else (float(height) / EMU_PER_PX / rows)
            safe_heights = [h if h > 0 else fallback_height for h in safe_heights]

            target_total_px = max(float(height) / EMU_PER_PX, 0.0)
            current_total_px = sum(safe_heights)
            if target_total_px > 0 and current_total_px > 0:
                safe_heights = [h * target_total_px / current_total_px for h in safe_heights]

            target_total_emu = int(height)
            consumed = 0
            for i in range(rows):
                if i == rows - 1:
                    row_h_emu = max(target_total_emu - consumed, 1)
                else:
                    row_h_emu = max(int(round(safe_heights[i] * EMU_PER_PX)), 1)
                    consumed += row_h_emu
                try:
                    table.rows[i].height = Emu(row_h_emu)
                except Exception:
                    pass
            applied_row_heights = True

    if not applied_row_heights:
        cell_min_height = _to_non_negative_float(props.get("cellMinHeight"), 0.0)
        if cell_min_height > 0:
            min_height_emu = Emu(int(cell_min_height * EMU_PER_PX))
            for r_idx in range(rows):
                try:
                    table.rows[r_idx].height = min_height_emu
                except Exception:
                    pass

    # 表格主题/边框写入（统一入口，避免重复与默认属性污染）
    try:
        from lxml import etree
        from pptx.oxml.ns import qn

        nsmap_a = "http://schemas.openxmlformats.org/drawingml/2006/main"
        tbl = table._tbl

        def _ensure_tbl_pr():
            tbl_pr = tbl.find(f"{{{nsmap_a}}}tblPr")
            if tbl_pr is None:
                tbl_pr = etree.SubElement(tbl, qn("a:tblPr"))
                tbl.insert(0, tbl_pr)
            return tbl_pr

        # 先清理默认 firstRow/bandRow，再按前端 theme 精确回写
        tblPr = _ensure_tbl_pr()
        bool_attrs = {
            "headerRow": "firstRow",
            "headerCol": "firstCol",
            "footerRow": "lastRow",
            "lastCol": "lastCol",
            "stripedRows": "bandRow",
            "stripedCols": "bandCol",
        }
        for pptx_attr in bool_attrs.values():
            if pptx_attr in tblPr.attrib:
                tblPr.attrib.pop(pptx_attr, None)

        theme = props.get("theme")
        if theme and isinstance(theme, dict):
            for fe_key, pptx_attr in bool_attrs.items():
                if theme.get(fe_key):
                    tblPr.set(pptx_attr, "1")

        # 表格边框：优先使用六向 borders，缺失项回退 outline
        outline = props.get("outline") if isinstance(props.get("outline"), dict) else {}
        o_color_str = str(outline.get("color", "#000000") or "#000000")
        o_theme_key = _normalize_text_theme_key(outline.get("themeKey"))
        o_width = _to_non_negative_float(outline.get("width", 1), 1.0)
        o_style_raw = str(outline.get("style", "solid") or "solid").lower()
        o_style = o_style_raw if o_style_raw in ("solid", "dashed", "dotted") else "solid"

        def _normalize_border_side(raw_side: Any) -> Dict[str, Any]:
            side_dict = raw_side if isinstance(raw_side, dict) else {}
            side_color_str = str(side_dict.get("color", o_color_str) or o_color_str)
            side_width = _to_non_negative_float(side_dict.get("width", o_width), o_width)
            side_style_raw = str(side_dict.get("style", o_style) or o_style).lower()
            side_style = side_style_raw if side_style_raw in ("solid", "dashed", "dotted") else o_style
            side_hex, _ = _parse_css_color(side_color_str)
            side_theme_key = _normalize_text_theme_key(side_dict.get("themeKey")) or o_theme_key
            return {
                "style": side_style,
                "width_emu": int(side_width * EMU_PER_PT),
                "hex": side_hex.lstrip("#")[:6] if side_hex else "000000",
                "colorStr": side_color_str,
                "themeKey": side_theme_key,
            }

        sides = ["top", "right", "bottom", "left", "insideH", "insideV"]
        borders_prop = props.get("borders")
        side_specs: Dict[str, Dict[str, Any]] = {}
        if isinstance(borders_prop, dict):
            for side in sides:
                if side in borders_prop:
                    side_specs[side] = _normalize_border_side(borders_prop.get(side))
                elif outline:
                    side_specs[side] = _normalize_border_side(None)
        elif outline:
            for side in sides:
                side_specs[side] = _normalize_border_side(None)

        if side_specs:
            borders_el = tblPr.find(qn("a:tblBorders"))
            if borders_el is None:
                borders_el = etree.SubElement(tblPr, qn("a:tblBorders"))
            else:
                for child in list(borders_el):
                    borders_el.remove(child)

            dash_map = {"solid": "solid", "dashed": "dash", "dotted": "dot", "dashDot": "dashDot", "longDash": "lgDash", "longDashDot": "lgDashDot"}
            for border_side in sides:
                spec = side_specs.get(border_side)
                if not spec:
                    continue
                border_el = etree.SubElement(borders_el, qn(f"a:{border_side}"))
                if spec["width_emu"] <= 0:
                    etree.SubElement(border_el, qn("a:noFill"))
                    continue

                border_el.set("w", str(spec["width_emu"]))
                border_el.set("cmpd", "sng")
                dash_val = dash_map.get(spec["style"], "solid")
                if dash_val != "solid":
                    border_el.set("prstDash", dash_val)
                applied_theme = False
                if spec.get("themeKey"):
                    applied_theme = _apply_theme_color_to_solid_parent(
                        border_el,
                        spec.get("themeKey"),
                        spec.get("colorStr"),
                    )
                if not applied_theme:
                    solid_fill = etree.SubElement(border_el, qn("a:solidFill"))
                    srgb = etree.SubElement(solid_fill, qn("a:srgbClr"))
                    srgb.set("val", spec["hex"])
                    _, border_alpha = _parse_css_color(spec.get("colorStr", ""))
                    if border_alpha is not None and border_alpha < 1.0:
                        alpha_el = etree.SubElement(srgb, qn("a:alpha"))
                        alpha_el.set("val", str(int(max(0.0, min(1.0, float(border_alpha))) * 100000)))
    except Exception as e:
        logger.debug(f"Failed to write table theme/outline: {e}")

    align_map = {
        "left": PP_ALIGN.LEFT,
        "center": PP_ALIGN.CENTER,
        "right": PP_ALIGN.RIGHT,
        "justify": PP_ALIGN.JUSTIFY,
    }

    def _apply_run_font_color(run_obj, color_str: Any, theme_key: Optional[Any] = None) -> None:
        normalized_theme_key = _normalize_text_theme_key(theme_key)
        if normalized_theme_key and _apply_run_theme_color(run_obj, normalized_theme_key, color_str):
            return
        if not color_str:
            return
        hex_fc, alpha_fc = _parse_css_color(str(color_str))
        hex_6 = hex_fc.lstrip("#")[:6]
        if len(hex_6) != 6:
            return

        run_obj.font.color.rgb = RGBColor.from_string(hex_6)
        if alpha_fc is None or alpha_fc >= 1.0:
            return

        nsmap_a = "http://schemas.openxmlformats.org/drawingml/2006/main"
        rPr = run_obj._r.find(f"{{{nsmap_a}}}rPr")
        if rPr is None:
            rPr = etree.SubElement(run_obj._r, qn("a:rPr"))
            run_obj._r.insert(0, rPr)

        solid_fill = rPr.find(f"{{{nsmap_a}}}solidFill")
        if solid_fill is None:
            solid_fill = etree.SubElement(rPr, qn("a:solidFill"))
        color_el = solid_fill.find(f"{{{nsmap_a}}}srgbClr")
        if color_el is None:
            color_el = etree.SubElement(solid_fill, qn("a:srgbClr"))
        color_el.set("val", hex_6)
        for old_alpha in color_el.findall(f"{{{nsmap_a}}}alpha"):
            color_el.remove(old_alpha)
        alpha_el = etree.SubElement(color_el, qn("a:alpha"))
        alpha_el.set("val", str(int(max(0.0, min(1.0, float(alpha_fc))) * 100000)))

    for r, row_data in enumerate(data):
        for c, cell_data in enumerate(row_data):
            if c >= cols or r >= rows:
                continue

            cell = table.cell(r, c)

            if isinstance(cell_data, dict):
                # 被合并的单元格跳过
                cs = _to_span_int(cell_data.get("colspan", 1), 1)
                rs = _to_span_int(cell_data.get("rowspan", 1), 1)
                if cs == 0 or rs == 0:
                    continue

                # 合并单元格（先合并再写入内容，确保写到合并后的主单元格）
                if cs > 1 or rs > 1:
                    try:
                        end_r = min(r + rs - 1, rows - 1)
                        end_c = min(c + cs - 1, cols - 1)
                        if end_r > r or end_c > c:
                            cell.merge(table.cell(end_r, end_c))
                            cell = table.cell(r, c)
                    except Exception:
                        pass

                # 写入单元格内容：优先使用 richText（HTML），否则使用纯文本
                rich_text = cell_data.get("richText")
                if rich_text:
                    _write_cell_rich_text(cell, rich_text, cell_data)
                else:
                    cell.text = cell_data.get("text", "")

                # 单元格样式
                style = cell_data.get("style") or {}

                # 背景色
                bg_color = style.get("bgColor") or cell_data.get("bgColor")
                bg_theme_key = style.get("bgColorThemeKey") or cell_data.get("bgColorThemeKey")
                if (bg_color and isinstance(bg_color, str)) or bg_theme_key:
                    try:
                        safe_bg_color = bg_color if isinstance(bg_color, str) and bg_color else "#000000"
                        hex_c, c_alpha = _parse_css_color(safe_bg_color)
                        hex_6 = hex_c.lstrip("#")[:6]
                        if len(hex_6) == 6:
                            cell.fill.solid()
                            applied_theme_bg = False
                            if bg_theme_key:
                                nsmap_a = "http://schemas.openxmlformats.org/drawingml/2006/main"
                                tcPr = cell._tc.find(f"{{{nsmap_a}}}tcPr") if hasattr(cell, "_tc") else None
                                if tcPr is not None:
                                    applied_theme_bg = _apply_theme_color_to_solid_parent(
                                        tcPr,
                                        bg_theme_key,
                                        safe_bg_color,
                                    )
                            if not applied_theme_bg:
                                cell.fill.fore_color.rgb = RGBColor.from_string(hex_6)
                                if c_alpha is not None and c_alpha < 1.0:
                                    nsmap_a = "http://schemas.openxmlformats.org/drawingml/2006/main"
                                    tcPr = cell._tc.find(f"{{{nsmap_a}}}tcPr") if hasattr(cell, "_tc") else None
                                    if tcPr is not None:
                                        solid = tcPr.find(f"{{{nsmap_a}}}solidFill")
                                        if solid is not None and len(solid) > 0:
                                            from lxml import etree
                                            from pptx.oxml.ns import qn
                                            color_el = solid[0]
                                            for old_alpha in color_el.findall(f"{{{nsmap_a}}}alpha"):
                                                color_el.remove(old_alpha)
                                            alpha_el = etree.SubElement(color_el, qn("a:alpha"))
                                            alpha_el.set("val", str(int(c_alpha * 100000)))
                    except Exception:
                        pass

                # 垂直对齐 + 单元格内边距
                v_align = style.get("verticalAlign") or cell_data.get("verticalAlign")
                padding = style.get("padding") or cell_data.get("padding")
                if v_align or padding:
                    try:
                        nsmap_a = "http://schemas.openxmlformats.org/drawingml/2006/main"
                        tcPr = cell._tc.find(f"{{{nsmap_a}}}tcPr") if hasattr(cell, "_tc") else None
                        if tcPr is None and hasattr(cell, "_tc"):
                            from lxml import etree
                            from pptx.oxml.ns import qn
                            tcPr = etree.SubElement(cell._tc, qn("a:tcPr"))
                        if tcPr is not None:
                            if v_align:
                                v_map = {"top": "t", "middle": "ctr", "bottom": "b"}
                                anchor = v_map.get(v_align, "ctr")
                                tcPr.set("anchor", anchor)
                            if padding and isinstance(padding, dict):
                                pad_map = {"paddingLeft": "marL", "paddingRight": "marR",
                                           "paddingTop": "marT", "paddingBottom": "marB"}
                                for fe_key, xml_attr in pad_map.items():
                                    val = padding.get(fe_key)
                                    if val is not None:
                                        try:
                                            emu_val = int(float(val) * 12700)
                                            tcPr.set(xml_attr, str(max(emu_val, 0)))
                                        except (ValueError, TypeError):
                                            pass
                    except Exception:
                        pass

                # 单元格级别边框（per-cell borders）
                cell_borders = style.get("cellBorders") or cell_data.get("cellBorders")
                if cell_borders and isinstance(cell_borders, dict):
                    try:
                        nsmap_a = "http://schemas.openxmlformats.org/drawingml/2006/main"
                        tcPr = cell._tc.find(f"{{{nsmap_a}}}tcPr") if hasattr(cell, "_tc") else None
                        if tcPr is None and hasattr(cell, "_tc"):
                            tcPr = etree.SubElement(cell._tc, qn("a:tcPr"))
                        if tcPr is not None:
                            tc_bdr = tcPr.find(qn("a:tcBorders"))
                            if tc_bdr is None:
                                tc_bdr = etree.SubElement(tcPr, qn("a:tcBorders"))
                            else:
                                for child in list(tc_bdr):
                                    tc_bdr.remove(child)
                            _cb_dash_map = {"solid": "solid", "dashed": "dash", "dotted": "dot", "dashDot": "dashDot", "longDash": "lgDash", "longDashDot": "lgDashDot"}
                            for cb_side in ("top", "right", "bottom", "left"):
                                cb_spec = cell_borders.get(cb_side)
                                if not cb_spec or not isinstance(cb_spec, dict):
                                    continue
                                cb_el = etree.SubElement(tc_bdr, qn(f"a:{cb_side}"))
                                cb_w = _to_non_negative_float(cb_spec.get("width", 0), 0.0)
                                if cb_w <= 0:
                                    etree.SubElement(cb_el, qn("a:noFill"))
                                    continue
                                cb_el.set("w", str(int(cb_w * EMU_PER_PT)))
                                cb_el.set("cmpd", "sng")
                                cb_dash = _cb_dash_map.get(str(cb_spec.get("style", "solid")), "solid")
                                if cb_dash != "solid":
                                    cb_el.set("prstDash", cb_dash)
                                cb_applied_theme = False
                                cb_tk = cb_spec.get("themeKey")
                                if cb_tk:
                                    cb_applied_theme = _apply_theme_color_to_solid_parent(
                                        cb_el, cb_tk, cb_spec.get("color"),
                                    )
                                if not cb_applied_theme:
                                    cb_sf = etree.SubElement(cb_el, qn("a:solidFill"))
                                    cb_srgb = etree.SubElement(cb_sf, qn("a:srgbClr"))
                                    cb_hex, cb_alpha = _parse_css_color(str(cb_spec.get("color", "#000000") or "#000000"))
                                    cb_srgb.set("val", cb_hex.lstrip("#")[:6])
                                    if cb_alpha is not None and cb_alpha < 1.0:
                                        cb_a_el = etree.SubElement(cb_srgb, qn("a:alpha"))
                                        cb_a_el.set("val", str(int(max(0.0, min(1.0, float(cb_alpha))) * 100000)))
                    except Exception:
                        pass

                # 文本样式（仅在非 richText 时应用，richText 已包含内联格式）
                if not rich_text and cell.text_frame.paragraphs:
                    p = cell.text_frame.paragraphs[0]

                    # 对齐
                    cell_align = style.get("align") or cell_data.get("align")
                    if cell_align and cell_align in align_map:
                        p.alignment = align_map[cell_align]

                    # 字体样式
                    for run in p.runs:
                        if style.get("bold") or cell_data.get("bold"):
                            run.font.bold = True
                        if style.get("italic") or cell_data.get("italic"):
                            run.font.italic = True
                        if style.get("underline") or cell_data.get("underline"):
                            run.font.underline = True

                        font_size = style.get("fontSize") or cell_data.get("fontSize")
                        if font_size:
                            font_size_pt = _css_length_to_pt(font_size)
                            if font_size_pt is not None and font_size_pt > 0:
                                run.font.size = Pt(round(font_size_pt, 3))

                        font_family = style.get("fontName") or style.get("fontFamily") or cell_data.get("fontName") or cell_data.get("fontFamily")
                        if font_family:
                            primary_font_name = _strip_font_fallback(str(font_family))
                            if primary_font_name:
                                run.font.name = primary_font_name

                        font_color = style.get("color") or cell_data.get("color")
                        font_color_theme_key = style.get("colorThemeKey") or cell_data.get("colorThemeKey")
                        if (font_color and isinstance(font_color, str)) or font_color_theme_key:
                            try:
                                _apply_run_font_color(run, font_color, font_color_theme_key)
                            except Exception:
                                pass
            else:
                cell.text = str(cell_data)

    return table_shape


def _write_chart_element(slide, props: dict, left, top, width, height):
    """写入图表元素到 PPTX。使用 python-pptx 的 chart API 创建原生图表。"""
    from pptx.chart.data import CategoryChartData
    from pptx.enum.chart import XL_CHART_TYPE

    chart_type_str = props.get("chartType", "bar")
    data = props.get("data", {})
    labels = data.get("labels", [])
    legends = data.get("legends", [])
    series_data = data.get("series", [])

    if not series_data:
        logger.warning("Chart write: no series data, skipping")
        return None

    options = props.get("options", {})
    if not isinstance(options, dict):
        options = {}
    is_stacked = bool(options.get("stack", False))

    # ── 散点图：使用 XyChartData ──
    if chart_type_str == "scatter":
        return _write_scatter_chart(slide, props, left, top, width, height)

    # ── 饼图/环形图：只用第一个 series ──
    is_pie = chart_type_str in ("pie", "ring")

    # 前端类型 → python-pptx XL_CHART_TYPE 映射
    _type_map = {
        "bar": XL_CHART_TYPE.COLUMN_CLUSTERED,
        "column": XL_CHART_TYPE.BAR_CLUSTERED,
        "line": XL_CHART_TYPE.LINE,
        "area": XL_CHART_TYPE.AREA,
        "pie": XL_CHART_TYPE.PIE,
        "ring": XL_CHART_TYPE.DOUGHNUT,
        "radar": XL_CHART_TYPE.RADAR,
    }

    # 堆叠变体
    if is_stacked and not is_pie:
        _stacked_map = {
            "bar": XL_CHART_TYPE.COLUMN_STACKED,
            "column": XL_CHART_TYPE.BAR_STACKED,
            "line": XL_CHART_TYPE.LINE_STACKED,
            "area": XL_CHART_TYPE.AREA_STACKED,
        }
        xl_chart_type = _stacked_map.get(chart_type_str, _type_map.get(chart_type_str, XL_CHART_TYPE.COLUMN_CLUSTERED))
    else:
        xl_chart_type = _type_map.get(chart_type_str, XL_CHART_TYPE.COLUMN_CLUSTERED)

    # 雷达图填充变体
    if chart_type_str == "radar" and options.get("radarFilled"):
        xl_chart_type = XL_CHART_TYPE.RADAR_FILLED

    chart_data = CategoryChartData()
    chart_data.categories = labels if labels else [f"Cat {i+1}" for i in range(len(series_data[0]) if series_data else 0)]

    if is_pie:
        # 饼图/环形图只用第一个系列
        name = legends[0] if legends else "Series 1"
        safe_values = [v if v is not None and v == v else 0 for v in series_data[0]]
        chart_data.add_series(name, safe_values)
    else:
        for i, values in enumerate(series_data):
            name = legends[i] if i < len(legends) else f"Series {i + 1}"
            safe_values = [v if v is not None and v == v else 0 for v in values]
            chart_data.add_series(name, safe_values)

    try:
        graphic_frame = slide.shapes.add_chart(
            xl_chart_type, left, top, width, height, chart_data
        )
        chart = graphic_frame.chart

        # 图例显示与位置
        show_legend = options.get("showLegend")
        if show_legend is not None:
            chart.has_legend = bool(show_legend)
        else:
            # 与前端渲染/导出默认行为对齐：
            # - pie/ring 默认显示图例
            # - 其他图表默认仅在多系列时显示
            chart.has_legend = is_pie or len(series_data) > 1

        if chart.has_legend:
            legend_pos = options.get("legendPosition", "b")
            try:
                from pptx.enum.chart import XL_LEGEND_POSITION
                _pos_map = {
                    "b": XL_LEGEND_POSITION.BOTTOM,
                    "t": XL_LEGEND_POSITION.TOP,
                    "l": XL_LEGEND_POSITION.LEFT,
                    "r": XL_LEGEND_POSITION.RIGHT,
                }
                if legend_pos in _pos_map:
                    chart.legend.position = _pos_map[legend_pos]
            except Exception:
                pass

        # 数据标签
        show_data_label = options.get("showDataLabel")
        if show_data_label is None:
            # 与前端默认对齐：pie/ring 默认开启数据标签
            show_data_label = is_pie
        if show_data_label:
            try:
                plot = chart.plots[0]
                plot.has_data_labels = True
                data_labels = plot.data_labels
                if is_pie:
                    # 饼图/环形图优先显示百分比，贴近前端 ECharts 标签语义。
                    try:
                        data_labels.show_percentage = True
                    except Exception:
                        pass
                    try:
                        data_labels.show_value = False
                    except Exception:
                        pass
                else:
                    data_labels.show_value = True
            except Exception:
                pass

        # 写回图表标题
        title = props.get("title", "") or props.get("chartTitle", "")
        if title:
            chart.has_title = True
            chart.chart_title.has_text_frame = True
            chart.chart_title.text_frame.text = title

        # 写回平滑选项（折线图/面积图）
        if options.get("lineSmooth") and chart_type_str in ("line", "area"):
            try:
                if chart_type_str == "line":
                    for series in chart.plots[0].series:
                        series.smooth = True
                else:
                    # AreaSeries 无 smooth 属性，需通过 XML 注入
                    from lxml import etree as _et
                    from pptx.oxml.ns import qn as _qn
                    _ns_c = "http://schemas.openxmlformats.org/drawingml/2006/chart"
                    for ser_el in chart.element.iter(f"{{{_ns_c}}}ser"):
                        old = ser_el.find(f"{{{_ns_c}}}smooth")
                        if old is not None:
                            ser_el.remove(old)
                        smooth_el = _et.SubElement(ser_el, _qn("c:smooth"))
                        smooth_el.set("val", "1")
            except Exception:
                pass

        _apply_chart_colors(chart, props.get("themeColors", []), props.get("themeColorKeys", []))
        _apply_chart_style(chart, props)
        return graphic_frame
    except Exception as e:
        logger.warning(f"Chart write failed for type '{chart_type_str}': {e}")
        return None


def _set_scatter_data_labels(chart) -> None:
    """为 scatter chart 注入 dLbls（python-pptx 对 scatter 未暴露 has_data_labels API）。"""
    from lxml import etree
    from pptx.oxml.ns import qn

    nsmap_c = "http://schemas.openxmlformats.org/drawingml/2006/chart"
    chart_xml = chart.element
    scatter_chart = chart_xml.find(f".//{{{nsmap_c}}}scatterChart")
    if scatter_chart is None:
        return

    old = scatter_chart.find(f"{{{nsmap_c}}}dLbls")
    if old is not None:
        scatter_chart.remove(old)

    d_lbls = etree.Element(qn("c:dLbls"))
    for tag, val in (
        ("showLegendKey", "0"),
        ("showVal", "1"),
        ("showCatName", "0"),
        ("showSerName", "0"),
        ("showPercent", "0"),
        ("showBubbleSize", "0"),
        ("showLeaderLines", "1"),
    ):
        node = etree.SubElement(d_lbls, qn(f"c:{tag}"))
        node.set("val", val)

    first_ax = scatter_chart.find(f"{{{nsmap_c}}}axId")
    if first_ax is not None:
        scatter_chart.insert(scatter_chart.index(first_ax), d_lbls)
    else:
        scatter_chart.append(d_lbls)


def _write_scatter_chart(slide, props: dict, left, top, width, height):
    """写入散点图（使用 XyChartData，保留 X-Y 关系）。"""
    from pptx.chart.data import XyChartData
    from pptx.enum.chart import XL_CHART_TYPE

    data = props.get("data", {})
    labels = data.get("labels", [])
    legends = data.get("legends", [])
    series_data = data.get("series", [])
    x_series_data = data.get("xSeries", [])

    options = props.get("options", {})
    if not isinstance(options, dict):
        options = {}

    chart_data = XyChartData()
    for i, values in enumerate(series_data):
        name = legends[i] if i < len(legends) else f"Series {i + 1}"
        s = chart_data.add_series(name)
        x_row = []
        if isinstance(x_series_data, list) and i < len(x_series_data) and isinstance(x_series_data[i], list):
            x_row = x_series_data[i]
        for j, y_val in enumerate(values):
            x_val = float(j + 1)
            if j < len(x_row):
                try:
                    x_val = float(x_row[j])
                except (ValueError, TypeError):
                    x_val = float(j + 1)
            elif j < len(labels):
                try:
                    x_val = float(labels[j])
                except (ValueError, TypeError):
                    x_val = float(j + 1)
            y = y_val if y_val is not None and y_val == y_val else 0
            s.add_data_point(x_val, y)

    # 保留平滑散点图类型（XY_SCATTER_SMOOTH）
    chart_kind = XL_CHART_TYPE.XY_SCATTER_SMOOTH if options.get("lineSmooth") else XL_CHART_TYPE.XY_SCATTER

    try:
        graphic_frame = slide.shapes.add_chart(
            chart_kind, left, top, width, height, chart_data
        )
        chart = graphic_frame.chart

        # 图例显示 + 位置
        show_legend = options.get("showLegend")
        chart.has_legend = bool(show_legend) if show_legend is not None else (len(series_data) > 1)
        if chart.has_legend:
            legend_pos = options.get("legendPosition", "b")
            try:
                from pptx.enum.chart import XL_LEGEND_POSITION
                _pos_map = {
                    "b": XL_LEGEND_POSITION.BOTTOM,
                    "t": XL_LEGEND_POSITION.TOP,
                    "l": XL_LEGEND_POSITION.LEFT,
                    "r": XL_LEGEND_POSITION.RIGHT,
                }
                if legend_pos in _pos_map:
                    chart.legend.position = _pos_map[legend_pos]
            except Exception:
                pass

        # 数据标签
        if options.get("showDataLabel"):
            try:
                plot = chart.plots[0]
                plot.has_data_labels = True
                plot.data_labels.show_value = True
            except Exception:
                # scatter 图表 python-pptx 无 has_data_labels API，走 XML 注入兜底
                try:
                    _set_scatter_data_labels(chart)
                except Exception:
                    pass

        # 标题
        title = props.get("title", "") or props.get("chartTitle", "")
        if title:
            chart.has_title = True
            chart.chart_title.has_text_frame = True
            chart.chart_title.text_frame.text = title

        _apply_chart_colors(chart, props.get("themeColors", []), props.get("themeColorKeys", []))
        _apply_chart_style(chart, props)
        return graphic_frame
    except Exception as e:
        logger.warning(f"Scatter chart write failed: {e}")
        return None


def _apply_chart_theme_color_to_series(
    series,
    theme_key: Any,
    color_str_for_alpha: Optional[Any] = None,
) -> bool:
    """对图表系列写入 schemeClr（填充/线条/marker）。"""
    try:
        from lxml import etree
        from pptx.oxml.ns import qn
    except Exception:
        return False

    try:
        nsmap_c = "http://schemas.openxmlformats.org/drawingml/2006/chart"
        nsmap_a = "http://schemas.openxmlformats.org/drawingml/2006/main"
        ser_el = series._element

        sp_pr = ser_el.find(f"{{{nsmap_c}}}spPr")
        if sp_pr is None:
            sp_pr = etree.SubElement(ser_el, qn("c:spPr"))

        applied_fill = _apply_theme_color_to_solid_parent(sp_pr, theme_key, color_str_for_alpha)

        ln = sp_pr.find(f"{{{nsmap_a}}}ln")
        if ln is None:
            ln = etree.SubElement(sp_pr, qn("a:ln"))
        applied_line = _apply_theme_color_to_solid_parent(ln, theme_key, color_str_for_alpha)

        applied_marker = False
        marker = ser_el.find(f"{{{nsmap_c}}}marker")
        if marker is not None:
            marker_sp_pr = marker.find(f"{{{nsmap_c}}}spPr")
            if marker_sp_pr is not None:
                applied_marker = _apply_theme_color_to_solid_parent(
                    marker_sp_pr,
                    theme_key,
                    color_str_for_alpha,
                )

        return bool(applied_fill or applied_line or applied_marker)
    except Exception:
        return False


def _apply_chart_theme_color_to_point(
    series,
    point_idx: int,
    theme_key: Any,
    color_str_for_alpha: Optional[Any] = None,
) -> bool:
    """对饼/环图扇区写入 schemeClr（dPt/spPr）。"""
    if point_idx < 0:
        return False

    try:
        from lxml import etree
        from pptx.oxml.ns import qn
    except Exception:
        return False

    try:
        nsmap_c = "http://schemas.openxmlformats.org/drawingml/2006/chart"
        nsmap_a = "http://schemas.openxmlformats.org/drawingml/2006/main"
        ser_el = series._element

        dpt_node = None
        for dpt in ser_el.findall(f"{{{nsmap_c}}}dPt"):
            idx_el = dpt.find(f"{{{nsmap_c}}}idx")
            if idx_el is None:
                continue
            try:
                cur_idx = int(idx_el.get("val") or -1)
            except (TypeError, ValueError):
                continue
            if cur_idx == point_idx:
                dpt_node = dpt
                break

        if dpt_node is None:
            dpt_node = etree.SubElement(ser_el, qn("c:dPt"))
            idx_el = etree.SubElement(dpt_node, qn("c:idx"))
            idx_el.set("val", str(point_idx))

        sp_pr = dpt_node.find(f"{{{nsmap_c}}}spPr")
        if sp_pr is None:
            sp_pr = etree.SubElement(dpt_node, qn("c:spPr"))

        applied_fill = _apply_theme_color_to_solid_parent(sp_pr, theme_key, color_str_for_alpha)

        ln = sp_pr.find(f"{{{nsmap_a}}}ln")
        if ln is None:
            ln = etree.SubElement(sp_pr, qn("a:ln"))
        applied_line = _apply_theme_color_to_solid_parent(ln, theme_key, color_str_for_alpha)

        return bool(applied_fill or applied_line)
    except Exception:
        return False


def _apply_chart_colors(chart, theme_colors: list, theme_color_keys: Optional[list] = None):
    """
    将图表颜色写回：
    1) 若有 themeColorKeys，优先写 schemeClr；
    2) 无 key 时回退 srgbClr。
    """
    from pptx.dml.color import RGBColor
    from pptx.enum.chart import XL_CHART_TYPE

    if not isinstance(theme_colors, list):
        theme_colors = []
    if not isinstance(theme_color_keys, list):
        theme_color_keys = []
    if not theme_colors and not theme_color_keys:
        return

    def _normalize_color_item(raw: Any) -> Optional[str]:
        if not isinstance(raw, str):
            return None
        hex_color, _ = _parse_css_color(raw)
        if not hex_color:
            return None
        hex_6 = hex_color.lstrip("#")[:6]
        if len(hex_6) != 6:
            return None
        return hex_6

    def _normalize_key_item(raw: Any) -> Optional[str]:
        return _normalize_text_theme_key(raw)

    def _apply_rgb_to_series(series, color_raw: Any) -> bool:
        color_hex = _normalize_color_item(color_raw)
        if not color_hex:
            return False
        try:
            rgb = RGBColor.from_string(color_hex)
            try:
                series.format.fill.solid()
                series.format.fill.fore_color.rgb = rgb
            except Exception:
                pass
            try:
                series.format.line.color.rgb = rgb
            except Exception:
                pass
            try:
                series.marker.format.fill.solid()
                series.marker.format.fill.fore_color.rgb = rgb
            except Exception:
                pass
            return True
        except Exception:
            return False

    def _apply_rgb_to_point(point, color_raw: Any) -> bool:
        color_hex = _normalize_color_item(color_raw)
        if not color_hex:
            return False
        try:
            rgb = RGBColor.from_string(color_hex)
            try:
                point.format.fill.solid()
                point.format.fill.fore_color.rgb = rgb
            except Exception:
                pass
            try:
                point.format.line.color.rgb = rgb
            except Exception:
                pass
            return True
        except Exception:
            return False

    try:
        plot = chart.plots[0]
        color_count = len(theme_colors)
        key_count = len(theme_color_keys)
        item_count = max(color_count, key_count)

        pie_types = {
            XL_CHART_TYPE.PIE,
            XL_CHART_TYPE.PIE_EXPLODED,
            XL_CHART_TYPE.THREE_D_PIE,
            XL_CHART_TYPE.THREE_D_PIE_EXPLODED,
            XL_CHART_TYPE.PIE_OF_PIE,
            XL_CHART_TYPE.BAR_OF_PIE,
            XL_CHART_TYPE.DOUGHNUT,
            XL_CHART_TYPE.DOUGHNUT_EXPLODED,
        }

        if chart.chart_type in pie_types and len(plot.series) > 0:
            series = plot.series[0]
            for idx, point in enumerate(series.points):
                if idx >= item_count:
                    break
                color_raw = theme_colors[idx] if idx < color_count else None
                key = _normalize_key_item(theme_color_keys[idx] if idx < key_count else None)
                applied = False
                if key:
                    applied = _apply_chart_theme_color_to_point(
                        series,
                        idx,
                        key,
                        color_raw,
                    )
                if not applied and color_raw is not None:
                    _apply_rgb_to_point(point, color_raw)
            return

        for idx, series in enumerate(plot.series):
            if idx >= item_count:
                break
            color_raw = theme_colors[idx] if idx < color_count else None
            key = _normalize_key_item(theme_color_keys[idx] if idx < key_count else None)
            applied = False
            if key:
                applied = _apply_chart_theme_color_to_series(series, key, color_raw)
            if not applied and color_raw is not None:
                _apply_rgb_to_series(series, color_raw)
    except Exception as e:
        logger.debug(f"Chart color application failed: {e}")


def _apply_chart_style(chart, props: dict) -> None:
    """将 fill / textColor / gridColor 写回图表 XML（覆盖已有节点，避免重复堆叠）。"""

    fill_color = props.get("fill", "")
    text_color = props.get("textColor", "")
    grid_color = props.get("gridColor", "")

    if not (fill_color or text_color or grid_color):
        return

    try:
        from lxml import etree
        from pptx.oxml.ns import qn

        def _ensure_child(parent, tag):
            node = parent.find(tag)
            if node is None:
                node = etree.SubElement(parent, tag)
            return node

        def _drop_children(parent, tag):
            for child in list(parent):
                if child.tag == tag:
                    parent.remove(child)

        def _set_solid_fill(parent, hex_6: str):
            _drop_children(parent, qn("a:solidFill"))
            solid = etree.SubElement(parent, qn("a:solidFill"))
            srgb = etree.SubElement(solid, qn("a:srgbClr"))
            srgb.set("val", hex_6)

        def _set_line_color(parent, hex_6: str):
            ln = _ensure_child(parent, qn("a:ln"))
            _drop_children(ln, qn("a:solidFill"))
            solid = etree.SubElement(ln, qn("a:solidFill"))
            srgb = etree.SubElement(solid, qn("a:srgbClr"))
            srgb.set("val", hex_6)

        # 图表区域背景色
        if fill_color and isinstance(fill_color, str):
            hex_f, _ = _parse_css_color(fill_color)
            hex_6 = hex_f.lstrip("#")[:6]
            if len(hex_6) == 6:
                chart_space = chart.element
                sp_pr = _ensure_child(chart_space, qn("a:spPr"))
                _set_solid_fill(sp_pr, hex_6)

        # 坐标轴文字颜色 + 网格线颜色
        if text_color or grid_color:
            hex_tc = ""
            hex_gc = ""
            if text_color and isinstance(text_color, str):
                hex_tc_raw, _ = _parse_css_color(text_color)
                hex_tc = hex_tc_raw.lstrip("#")[:6] if hex_tc_raw else ""
            if grid_color and isinstance(grid_color, str):
                hex_gc_raw, _ = _parse_css_color(grid_color)
                hex_gc = hex_gc_raw.lstrip("#")[:6] if hex_gc_raw else ""

            if len(hex_tc) == 6 or len(hex_gc) == 6:
                ns_c = "http://schemas.openxmlformats.org/drawingml/2006/chart"
                chart_xml = chart.element

                for ax_tag in (f"{{{ns_c}}}catAx", f"{{{ns_c}}}valAx"):
                    for ax in chart_xml.iter(ax_tag):
                        if len(hex_tc) == 6:
                            tx_pr = _ensure_child(ax, qn("c:txPr"))
                            _ensure_child(tx_pr, qn("a:bodyPr"))
                            _ensure_child(tx_pr, qn("a:lstStyle"))
                            paragraph = _ensure_child(tx_pr, qn("a:p"))
                            paragraph_pr = _ensure_child(paragraph, qn("a:pPr"))
                            def_r_pr = _ensure_child(paragraph_pr, qn("a:defRPr"))
                            _set_solid_fill(def_r_pr, hex_tc)

                        if len(hex_gc) == 6:
                            grid_el = _ensure_child(ax, qn("c:majorGridlines"))
                            sp_pr = _ensure_child(grid_el, qn("c:spPr"))
                            _set_line_color(sp_pr, hex_gc)

    except Exception as e:
        logger.debug(f"Chart style application failed: {e}")


def _guess_image_format(image_bytes: bytes, src_hint: Optional[str] = None) -> Optional[str]:
    """尽力推断图片格式（扩展名或 MIME 后缀）"""
    try:
        if src_hint:
            import re as _re
            if src_hint.startswith("data:image/"):
                m = _re.match(r"data:image/([a-zA-Z0-9.+-]+)", src_hint)
                if m:
                    ext = m.group(1).lower()
                    if ext == "svg+xml":
                        return "svg"
                    if ext in ("jpeg", "jpg"):
                        return "jpg"
                    return ext
            path_hint = src_hint.split("?", 1)[0].split("#", 1)[0].lower()
            if "." in path_hint:
                ext = path_hint.rsplit(".", 1)[-1]
                if ext:
                    if ext == "jpeg":
                        return "jpg"
                    return ext
    except Exception:
        pass

    head = image_bytes[:512].lstrip().lower()
    if head.startswith(b"\x89png\r\n\x1a\n"):
        return "png"
    if head.startswith(b"\xff\xd8\xff"):
        return "jpg"
    if head.startswith(b"gif87a") or head.startswith(b"gif89a"):
        return "gif"
    if head.startswith(b"bm"):
        return "bmp"
    if image_bytes[:4] == b"RIFF" and image_bytes[8:12] == b"WEBP":
        return "webp"
    if head.startswith(b"ii*\x00") or head.startswith(b"mm\x00*"):
        return "tiff"
    if head.startswith(b"<svg") or b"<svg" in head:
        return "svg"
    return None


def _extract_svg_canvas_size(svg_bytes: bytes) -> Tuple[int, int]:
    """
    从 SVG 内容中推断渲染尺寸（px）。
    优先 width/height，其次 viewBox，最后降级到 1024x1024。
    """
    import re as _re

    default_size = (1024, 1024)
    try:
        from lxml import etree

        root = etree.fromstring(svg_bytes)
        if root is None:
            return default_size

        def _parse_len(raw_val: Optional[str]) -> Optional[float]:
            if not raw_val:
                return None
            m = _re.match(r"^\s*([0-9]+(?:\.[0-9]+)?)", str(raw_val))
            if not m:
                return None
            num = float(m.group(1))
            return num if num > 0 else None

        width = _parse_len(root.get("width"))
        height = _parse_len(root.get("height"))

        if (width is None or height is None) and root.get("viewBox"):
            parts = [p for p in _re.split(r"[,\s]+", root.get("viewBox", "").strip()) if p]
            if len(parts) == 4:
                try:
                    vb_w = float(parts[2])
                    vb_h = float(parts[3])
                    if vb_w > 0 and width is None:
                        width = vb_w
                    if vb_h > 0 and height is None:
                        height = vb_h
                except Exception:
                    pass

        if width is None or height is None:
            return default_size

        # 防止超大 SVG 占用过多内存
        max_side = max(width, height)
        if max_side > 4096:
            scale = 4096.0 / max_side
            width *= scale
            height *= scale

        return (max(1, int(round(width))), max(1, int(round(height))))
    except Exception:
        return default_size


def _resolve_svg_render_size(
    svg_bytes: bytes,
    target_width_px: Optional[float] = None,
    target_height_px: Optional[float] = None,
    supersample: float = 1.0,
) -> Tuple[int, int]:
    """
    计算 SVG 栅格化输出尺寸。

    规则：
    - 优先使用目标渲染尺寸（元素在画布上的宽高），避免 SVG 内部 width/height 过小导致模糊。
    - 使用 supersample 做超采样提升清晰度。
    - 对超大尺寸做上限裁剪，防止内存/导出体积失控。
    """
    base_w, base_h = _extract_svg_canvas_size(svg_bytes)

    safe_scale = 1.0
    try:
        safe_scale = float(supersample)
    except (TypeError, ValueError):
        safe_scale = 1.0
    if safe_scale <= 0:
        safe_scale = 1.0

    ratio = (base_h / base_w) if base_w > 0 else 1.0
    if ratio <= 0:
        ratio = 1.0

    w_target: Optional[float] = None
    h_target: Optional[float] = None
    try:
        if target_width_px is not None:
            w_val = float(target_width_px)
            if w_val > 0:
                w_target = w_val
    except (TypeError, ValueError):
        w_target = None
    try:
        if target_height_px is not None:
            h_val = float(target_height_px)
            if h_val > 0:
                h_target = h_val
    except (TypeError, ValueError):
        h_target = None

    if w_target is None and h_target is None:
        width = max(1, int(round(base_w * safe_scale)))
        height = max(1, int(round(base_h * safe_scale)))
    else:
        if w_target is None and h_target is not None:
            w_target = h_target / ratio if ratio > 0 else float(base_w)
        if h_target is None and w_target is not None:
            h_target = w_target * ratio

        width = max(1, int(round((w_target or float(base_w)) * safe_scale)))
        height = max(1, int(round((h_target or float(base_h)) * safe_scale)))

    max_side = 4096
    if max(width, height) > max_side:
        scale_down = max_side / float(max(width, height))
        width = max(1, int(width * scale_down))
        height = max(1, int(height * scale_down))

    max_pixels = 16_000_000
    if width * height > max_pixels:
        import math as _math
        scale_down = _math.sqrt(max_pixels / float(width * height))
        width = max(1, int(width * scale_down))
        height = max(1, int(height * scale_down))

    return width, height


_PLAYWRIGHT_SVG_TIMEOUT_SECONDS = 30


def _png_is_useless_solid(png_bytes: bytes) -> bool:
    """检测「全透明」或「全不透明白」PNG——stroke-only SVG 坏栅格的典型产物。"""
    try:
        import io
        from PIL import Image

        img = Image.open(io.BytesIO(png_bytes)).convert("RGBA")
        extrema = img.getextrema()  # ((rmin,rmax), (gmin,gmax), (bmin,bmax), (amin,amax))
        if not extrema or len(extrema) < 4:
            return False
        (rmin, rmax), (gmin, gmax), (bmin, bmax), (amin, amax) = extrema
        # 全透明
        if amax == 0:
            return True
        # 全不透明且 RGB 恒为白（旧 img 截图路径对 stroke-only 常出此结果）
        if amin == 255 and amax == 255 and rmin == 255 and rmax == 255 and gmin == 255 and gmax == 255 and bmin == 255 and bmax == 255:
            return True
        return False
    except Exception:
        return False


def _prepare_inline_svg_markup(svg_bytes: bytes, width: int, height: int) -> str:
    """把 SVG bytes 变成可内联进 HTML 的 markup（补 xmlns / 尺寸）。"""
    try:
        text = svg_bytes.decode("utf-8")
    except UnicodeDecodeError:
        text = svg_bytes.decode("utf-8", errors="replace")
    text = text.strip()
    if not text.lower().startswith("<svg"):
        # 容错：包一层
        text = f'<svg xmlns="http://www.w3.org/2000/svg">{text}</svg>'
    if "xmlns=" not in text[:200]:
        text = text.replace("<svg", '<svg xmlns="http://www.w3.org/2000/svg"', 1)
    # 强制渲染尺寸，避免 viewBox-only 时落成 0×0 / 浏览器默认
    import re as _re

    if _re.search(r'\bwidth\s*=', text[:300], flags=_re.IGNORECASE):
        text = _re.sub(r'\bwidth\s*=\s*"[^"]*"', f'width="{width}"', text, count=1, flags=_re.IGNORECASE)
    else:
        text = text.replace("<svg", f'<svg width="{width}"', 1)
    if _re.search(r'\bheight\s*=', text[:300], flags=_re.IGNORECASE):
        text = _re.sub(r'\bheight\s*=\s*"[^"]*"', f'height="{height}"', text, count=1, flags=_re.IGNORECASE)
    else:
        text = text.replace("<svg", f'<svg height="{height}"', 1)
    return text


def _render_svg_to_png_with_playwright(
    svg_bytes: bytes,
    target_width_px: Optional[float] = None,
    target_height_px: Optional[float] = None,
    supersample: float = 1.0,
) -> Optional[bytes]:
    """
    在缺少 CairoSVG 时，使用 Playwright(Chromium) 渲染 SVG 为 PNG。
    整体限时 _PLAYWRIGHT_SVG_TIMEOUT_SECONDS 秒，避免 Chromium 卡住无限阻塞。

    使用内联 <svg> + omit_background，避免 data: URL <img> 对 stroke-only
    图标产出全白不透明方块。
    """
    from concurrent.futures import ThreadPoolExecutor, TimeoutError as FuturesTimeoutError

    def _do_render() -> Optional[bytes]:
        from playwright.sync_api import sync_playwright

        width, height = _resolve_svg_render_size(
            svg_bytes,
            target_width_px=target_width_px,
            target_height_px=target_height_px,
            supersample=supersample,
        )
        svg_markup = _prepare_inline_svg_markup(svg_bytes, width, height)
        html = (
            "<!doctype html><html><head><style>"
            "html,body{margin:0;padding:0;background:transparent;}"
            f"svg{{display:block;width:{width}px;height:{height}px;}}"
            "</style></head><body>"
            f"{svg_markup}"
            "</body></html>"
        )

        with sync_playwright() as p:
            browser = p.chromium.launch(headless=True)
            try:
                page = browser.new_page(
                    viewport={"width": max(width, 1), "height": max(height, 1)},
                    device_scale_factor=1,
                )
                page.set_content(html, wait_until="load")
                locator = page.locator("svg").first
                locator.wait_for(state="visible", timeout=3000)
                png = locator.screenshot(type="png", omit_background=True)
                if _png_is_useless_solid(png):
                    logger.warning(
                        "Playwright SVG render produced useless solid PNG (%dx%d); discard",
                        width, height,
                    )
                    return None
                return png
            finally:
                browser.close()

    pool = ThreadPoolExecutor(max_workers=1, thread_name_prefix="pw-svg")
    try:
        future = pool.submit(_do_render)
        return future.result(timeout=_PLAYWRIGHT_SVG_TIMEOUT_SECONDS)
    except FuturesTimeoutError:
        logger.warning("Playwright SVG render timed out after %ds", _PLAYWRIGHT_SVG_TIMEOUT_SECONDS)
        future.cancel()
        return None
    except Exception as e:
        logger.warning("Failed to render SVG with Playwright: %s", e)
        return None
    finally:
        pool.shutdown(wait=False, cancel_futures=True)


def _normalize_image_bytes_for_pptx(
    image_bytes: bytes,
    src_hint: Optional[str] = None,
    target_width_px: Optional[float] = None,
    target_height_px: Optional[float] = None,
    supersample: float = 1.0,
) -> bytes:
    """
    将 python-pptx 可能不支持的图片（如 webp/svg）尽量转为 PNG。
    转换失败时返回原始 bytes，交由后续写入逻辑兜底。
    """
    fmt = (_guess_image_format(image_bytes, src_hint) or "").lower()
    native_ok = {"png", "jpg", "jpeg", "gif", "bmp", "tif", "tiff", "wmf", "emf"}
    if fmt in native_ok:
        return image_bytes

    # SVG 优先尝试 CairoSVG 转 PNG
    if fmt == "svg":
        out_w, out_h = _resolve_svg_render_size(
            image_bytes,
            target_width_px=target_width_px,
            target_height_px=target_height_px,
            supersample=supersample,
        )
        try:
            import cairosvg
            png_bytes = cairosvg.svg2png(
                bytestring=image_bytes,
                output_width=out_w,
                output_height=out_h,
            )
            if png_bytes and not _png_is_useless_solid(png_bytes):
                return png_bytes
            logger.warning("CairoSVG produced useless solid PNG; trying Playwright fallback")
            pw_png = _render_svg_to_png_with_playwright(
                image_bytes,
                target_width_px=target_width_px,
                target_height_px=target_height_px,
                supersample=supersample,
            )
            if pw_png:
                return pw_png
            return png_bytes or image_bytes
        except ImportError:
            png_bytes = _render_svg_to_png_with_playwright(
                image_bytes,
                target_width_px=target_width_px,
                target_height_px=target_height_px,
                supersample=supersample,
            )
            if png_bytes:
                return png_bytes
            logger.warning("SVG image detected but cairosvg/playwright unavailable, keep original bytes")
            return image_bytes
        except Exception as e:
            logger.warning(f"Failed to convert SVG to PNG: {e}")
            png_bytes = _render_svg_to_png_with_playwright(
                image_bytes,
                target_width_px=target_width_px,
                target_height_px=target_height_px,
                supersample=supersample,
            )
            if png_bytes:
                return png_bytes
            return image_bytes

    # 其他非常见格式（如 webp）尝试 Pillow 转 PNG
    try:
        import io
        from PIL import Image
        img = Image.open(io.BytesIO(image_bytes))
        if img.mode in ("RGBA", "LA") or (img.mode == "P" and "transparency" in img.info):
            img = img.convert("RGBA")
        else:
            img = img.convert("RGB")
        buf = io.BytesIO()
        img.save(buf, format="PNG")
        return buf.getvalue()
    except ImportError:
        logger.warning(f"Image format '{fmt or 'unknown'}' may be unsupported and Pillow is unavailable")
    except Exception as e:
        logger.warning(f"Failed to normalize image format '{fmt or 'unknown'}': {e}")

    return image_bytes


def _read_binary_source_for_write(src: str) -> Optional[Tuple[bytes, Optional[str]]]:
    """
    读取通用二进制资源（data URL / http(s)），返回 (bytes, mime?)。
    安全策略：不允许本地文件路径，防止任意文件读取漏洞。
    """
    import base64 as b64

    if not isinstance(src, str):
        return None
    source = src.strip()
    if not source:
        return None

    if source.startswith("data:"):
        try:
            header, b64_data = source.split(",", 1)
            mime = None
            if header.lower().startswith("data:"):
                mime = header[5:].split(";", 1)[0].strip().lower() or None
            return b64.b64decode(b64_data), mime
        except Exception as exc:
            logger.warning("Failed to decode data URL for write: %s (src=%s…)", exc, source[:80])
            return None

    if source.startswith("http://") or source.startswith("https://"):
        try:
            from apps.services.common.url_security import ssrf_safe_request
            resp = ssrf_safe_request("GET", source, timeout=20, allow_redirects=True)
            resp.raise_for_status()
            content_type = (resp.headers.get("Content-Type") or "").split(";", 1)[0].strip().lower()
            payload = resp.content
            if len(payload) > MAX_EXPORT_IMAGE_BYTES:
                logger.warning("Resource too large (%d bytes), skipping: %s…", len(payload), source[:120])
                return None
            return payload, content_type or None
        except ValueError as exc:
            logger.warning("SSRF blocked URL for write: %s reason=%s", source[:120], exc)
            return None
        except Exception as exc:
            logger.warning("Failed to download resource for write: %s (url=%s…)", exc, source[:120])
            return None

    logger.warning("Unsupported source scheme for write (local paths disallowed): %s…", source[:120])
    return None


def _write_shape_alt_text(shape, alt_text: Optional[str]) -> None:
    if not isinstance(alt_text, str) or not alt_text.strip():
        return

    nsmap_p = "http://schemas.openxmlformats.org/presentationml/2006/main"
    try:
        c_nv_pr = shape._element.find(f".//{{{nsmap_p}}}cNvPr")
        if c_nv_pr is not None:
            c_nv_pr.set("descr", alt_text.strip())
    except Exception:
        pass


def _write_media_element(
    slide,
    media_type: str,
    props: dict,
    left,
    top,
    width,
    height,
):
    """
    写入媒体元素（video/audio），返回 shape 对象。
    """
    import io
    import os as _os
    import tempfile as _tempfile
    from pptx.util import Emu

    src = props.get("src", "")
    if not isinstance(src, str) or not src.strip():
        logger.warning(f"{media_type} element has no src, skipping")
        return None

    source_payload = _read_binary_source_for_write(src)
    if not source_payload:
        logger.warning(f"{media_type} src not resolvable: {str(src)[:120]}")
        return None
    media_bytes, source_mime = source_payload
    if not media_bytes:
        logger.warning(f"{media_type} source is empty")
        return None

    ext_hint = _normalize_media_ext(props.get("ext")) or _extract_file_ext_from_target(src)
    mime_type = _infer_media_mime(ext_hint, source_mime)
    expected_prefix = "audio/" if media_type == "audio" else "video/"
    if not (isinstance(mime_type, str) and mime_type.startswith(expected_prefix)):
        mime_type = "audio/mpeg" if media_type == "audio" else "video/mp4"

    safe_width = width if int(width) > 0 else Emu(1)
    safe_height = height if int(height) > 0 else Emu(1)

    kwargs: Dict[str, Any] = {"mime_type": mime_type}
    if media_type == "video":
        poster = props.get("poster")
        if isinstance(poster, str) and poster.strip():
            poster_payload = _read_binary_source_for_write(poster)
            if poster_payload and poster_payload[0]:
                poster_bytes, poster_mime = poster_payload
                if not poster_mime or poster_mime.startswith("image/"):
                    kwargs["poster_frame_image"] = io.BytesIO(poster_bytes)

    preferred_ext = _normalize_media_ext(ext_hint) or MEDIA_EXT_BY_MIME.get(mime_type)
    if not preferred_ext:
        preferred_ext = "mp3" if media_type == "audio" else "mp4"

    tmp_file_path = None
    try:
        tmp_file = _tempfile.NamedTemporaryFile(suffix=f".{preferred_ext}", delete=False)
        tmp_file.write(media_bytes)
        tmp_file.flush()
        tmp_file_path = tmp_file.name
        tmp_file.close()

        shape = slide.shapes.add_movie(
            tmp_file_path,
            left,
            top,
            safe_width,
            safe_height,
            **kwargs,
        )
    finally:
        if tmp_file_path and _os.path.exists(tmp_file_path):
            try:
                _os.remove(tmp_file_path)
            except Exception:
                pass

    meta_alt_text = _encode_media_alt_text(media_type, props if isinstance(props, dict) else {})
    if meta_alt_text:
        _write_shape_alt_text(shape, meta_alt_text)

    return shape


def _write_image_element(
    slide,
    props: dict,
    left,
    top,
    width,
    height,
    width_px: Optional[float] = None,
    height_px: Optional[float] = None,
    svg_scale: float = 2.0,
):
    """写入图片元素（支持 base64/URL、裁剪、边框、圆角、滤镜；本地路径已禁用）。返回 shape 对象。"""
    import base64 as b64
    import io
    from pptx.util import Pt
    from pptx.dml.color import RGBColor

    src = props.get("src", "")
    if not src:
        logger.warning("Image element has no src, skipping")
        return None

    # ── 0. 加载图片字节（OSS 走 SDK，其他 URL 走 SSRF 安全 HTTPS）──
    image_bytes = _download_image_smart(src, max_bytes=MAX_EXPORT_IMAGE_BYTES)
    if not image_bytes:
        return None

    # ── 0a. 规范化图片格式（webp/svg 等转 PNG，失败则保留原始字节） ──
    image_bytes = _normalize_image_bytes_for_pptx(
        image_bytes,
        src_hint=src,
        target_width_px=width_px,
        target_height_px=height_px,
        supersample=svg_scale,
    )

    image_stream = io.BytesIO(image_bytes)
    shape = slide.shapes.add_picture(image_stream, left, top, width, height)

    if shape is None:
        return None

    # ── 0b. 替代文本（altText）写回 cNvPr@descr ──
    alt_text = props.get("altText")
    nsmap_p = "http://schemas.openxmlformats.org/presentationml/2006/main"
    if isinstance(alt_text, str):
        alt_text = alt_text.strip()
        if alt_text:
            try:
                c_nv_pr = shape._element.find(f".//{{{nsmap_p}}}cNvPr")
                if c_nv_pr is not None:
                    c_nv_pr.set("descr", alt_text)
            except Exception as e:
                logger.debug(f"Failed to write image altText: {e}")

    nsmap_a = "http://schemas.openxmlformats.org/drawingml/2006/main"

    # ── 1. 图片裁剪 (clip → srcRect) ──
    clip = props.get("clip")
    is_ellipse_clip = bool(isinstance(clip, dict) and clip.get("shape") == "ellipse")
    if clip and isinstance(clip, dict):
        try:
            clip_range = clip.get("range", [])
            if clip_range and len(clip_range) >= 4:
                from lxml import etree
                from pptx.oxml.ns import qn

                # clip.range = [[left,top],[right,top],[right,bottom],[left,bottom]]
                l_pct = clip_range[0][0] if clip_range[0] else 0
                t_pct = clip_range[0][1] if clip_range[0] else 0
                r_pct = 1 - (clip_range[1][0] if clip_range[1] else 1)
                b_pct = 1 - (clip_range[2][1] if clip_range[2] else 1)

                # PPTX srcRect 值为千分比（0-100000）
                blip_fill = shape._element.find(f".//{{{nsmap_p}}}blipFill")
                if blip_fill is None:
                    blip_fill = shape._element.find(f".//{{{nsmap_a}}}blipFill")
                if blip_fill is not None:
                    src_rect = blip_fill.find(f"{{{nsmap_a}}}srcRect")
                    if src_rect is None:
                        src_rect = etree.SubElement(blip_fill, qn("a:srcRect"))
                    src_rect.set("l", str(int(l_pct * 100000)))
                    src_rect.set("t", str(int(t_pct * 100000)))
                    src_rect.set("r", str(int(r_pct * 100000)))
                    src_rect.set("b", str(int(b_pct * 100000)))
        except Exception as e:
            logger.debug(f"Failed to write image clip: {e}")

    # ── 1a. 圆形/椭圆裁剪（clip.shape=ellipse → prstGeom=ellipse） ──
    if is_ellipse_clip:
        try:
            from lxml import etree
            from pptx.oxml.ns import qn
            spPr = _find_sp_pr(shape._element)
            if spPr is not None:
                for old_geom in spPr.findall(f"{{{nsmap_a}}}prstGeom"):
                    spPr.remove(old_geom)
                prst_geom = etree.SubElement(spPr, qn("a:prstGeom"))
                prst_geom.set("prst", "ellipse")
                etree.SubElement(prst_geom, qn("a:avLst"))
        except Exception as e:
            logger.debug(f"Failed to apply ellipse image clip: {e}")

    # ── 1b. objectFit → stretch/fillRect ──
    object_fit = props.get("objectFit")
    if object_fit:
        try:
            from lxml import etree
            from pptx.oxml.ns import qn

            blip_fill = shape._element.find(f".//{{{nsmap_p}}}blipFill")
            if blip_fill is None:
                blip_fill = shape._element.find(f".//{{{nsmap_a}}}blipFill")
            if blip_fill is not None:
                stretch = blip_fill.find(f"{{{nsmap_a}}}stretch")
                if stretch is None:
                    stretch = etree.SubElement(blip_fill, qn("a:stretch"))
                fill_rect = stretch.find(f"{{{nsmap_a}}}fillRect")
                if fill_rect is None:
                    fill_rect = etree.SubElement(stretch, qn("a:fillRect"))
                if object_fit == "contain":
                    fill_rect.set("t", "10000")
                    fill_rect.set("r", "10000")
                    fill_rect.set("b", "10000")
                    fill_rect.set("l", "10000")
                else:
                    for attr in ("t", "r", "b", "l"):
                        if fill_rect.get(attr):
                            del fill_rect.attrib[attr]
        except Exception as e:
            logger.debug(f"Failed to apply objectFit: {e}")

    # ── 2. 图片边框 (outline) ──
    outline = props.get("outline")
    if outline and isinstance(outline, dict):
        try:
            ln = shape.line
            if outline.get("width"):
                ln.width = Pt(float(outline["width"]))
            if outline.get("color"):
                hex_c, ln_alpha = _parse_css_color(outline["color"])
                hex_6 = hex_c.lstrip("#")
                if len(hex_6) == 6:
                    ln.color.rgb = RGBColor.from_string(hex_6)
                if ln_alpha is not None and ln_alpha < 1.0:
                    _apply_color_alpha(shape, "outline", ln_alpha)
            dash_style = outline.get("style", "solid")
            _outline_dash_to_ooxml = {
                "dashed": "dash", "dotted": "dot",
                "dashDot": "dashDot", "longDash": "lgDash", "longDashDot": "lgDashDot",
            }
            _outline_dash_val = _outline_dash_to_ooxml.get(dash_style)
            if _outline_dash_val:
                from pptx.oxml.ns import qn
                from lxml import etree
                ln_elem = shape._element.find(f".//{{{nsmap_a}}}ln")
                if ln_elem is not None:
                    prstDash = etree.SubElement(ln_elem, qn("a:prstDash"))
                    prstDash.set("val", _outline_dash_val)
        except Exception as e:
            logger.debug(f"Failed to apply image outline: {e}")

    # ── 3. 圆角 (radius → prstGeom roundRect + avLst) ──
    radius = props.get("radius")
    if not is_ellipse_clip and radius and float(radius) > 0:
        try:
            from pptx.oxml.ns import qn
            from lxml import etree

            spPr = _find_sp_pr(shape._element)
            if spPr is not None:
                # 移除默认的矩形几何
                for old_geom in spPr.findall(f"{{{nsmap_a}}}prstGeom"):
                    spPr.remove(old_geom)
                # 添加圆角矩形几何
                prst_geom = etree.SubElement(spPr, qn("a:prstGeom"))
                prst_geom.set("prst", "roundRect")
                av_lst = etree.SubElement(prst_geom, qn("a:avLst"))
                # adj 使用相对比例（radius / min(width,height)），避免依赖固定 EMU_PER_PT
                shorter_side_px = None
                try:
                    if width_px is not None and height_px is not None:
                        w_px = float(width_px)
                        h_px = float(height_px)
                        if w_px > 0 and h_px > 0:
                            shorter_side_px = min(w_px, h_px)
                except (TypeError, ValueError):
                    shorter_side_px = None

                if shorter_side_px is not None:
                    ratio = float(radius) / shorter_side_px
                    adj_val = min(int(max(0.0, ratio) * 100000), 50000)
                else:
                    # 回退：使用 shape EMU 近似计算
                    shorter_side = min(int(shape.width), int(shape.height))
                    adj_val = min(int(float(radius) * EMU_PER_PT / max(shorter_side, 1) * 100000), 50000)

                if adj_val > 0:
                    gd = etree.SubElement(av_lst, qn("a:gd"))
                    gd.set("name", "adj")
                    gd.set("fmla", f"val {adj_val}")
        except Exception as e:
            logger.debug(f"Failed to apply image radius: {e}")

    # ── 4. 图片滤镜 (filters → blip 子元素) ──
    filters = props.get("filters")
    if filters and isinstance(filters, dict):
        try:
            from lxml import etree
            from pptx.oxml.ns import qn

            blip_fill = shape._element.find(f".//{{{nsmap_p}}}blipFill")
            if blip_fill is None:
                blip_fill = shape._element.find(f".//{{{nsmap_a}}}blipFill")
            blip = blip_fill.find(f"{{{nsmap_a}}}blip") if blip_fill is not None else None
            if blip is not None:
                # 模糊
                blur = filters.get("blur")
                if blur is not None:
                    try:
                        blur_val = max(0.0, float(blur))
                        if blur_val > 1e-3:
                            blur_el = etree.SubElement(blip, qn("a:blur"))
                            blur_el.set("rad", str(int(blur_val * EMU_PER_PT)))
                    except Exception:
                        pass

                # 反相
                invert = filters.get("invert")
                if invert is not None:
                    try:
                        if float(invert) > 1e-3:
                            etree.SubElement(blip, qn("a:inv"))
                    except Exception:
                        pass

                # 灰度
                if filters.get("grayscale"):
                    etree.SubElement(blip, qn("a:grayscl"))

                # 亮度和对比度 → <a:lum>
                brightness = filters.get("brightness")
                contrast = filters.get("contrast")
                if brightness is not None or contrast is not None:
                    lum = etree.SubElement(blip, qn("a:lum"))
                    if brightness is not None:
                        bright_val = int((float(brightness) - 1.0) * 100000)
                        lum.set("bright", str(bright_val))
                    if contrast is not None:
                        contrast_val = int((float(contrast) - 1.0) * 100000)
                        lum.set("contrast", str(contrast_val))

                # 色相旋转和饱和度 → <a:hsl>
                hue_rotate = filters.get("hueRotate")
                saturate = filters.get("saturate")
                if hue_rotate is not None or (saturate is not None and saturate != 1):
                    hsl = etree.SubElement(blip, qn("a:hsl"))
                    if hue_rotate is not None:
                        hsl.set("hue", str(int(float(hue_rotate) * 60000)))
                    if saturate is not None:
                        sat_val = int((float(saturate) - 1.0) * 100000)
                        hsl.set("sat", str(sat_val))

                # 棕褐色调 → sepia（使用 <a:duotone> 双色调近似，读取端会将其映射回 sepia）
                sepia = filters.get("sepia")
                if sepia and float(sepia) > 0:
                    duotone = etree.SubElement(blip, qn("a:duotone"))
                    # 经典棕褐色调：暗部 #3D2B1F（深棕），亮部 #FFF5E1（暖象牙）
                    dark = etree.SubElement(duotone, qn("a:srgbClr"))
                    dark.set("val", "3D2B1F")
                    light = etree.SubElement(duotone, qn("a:srgbClr"))
                    light.set("val", "FFF5E1")
        except Exception as e:
            logger.debug(f"Failed to apply image filters: {e}")

    # ── 5. 锁定宽高比 (fixedRatio → picLocks.noChangeAspect) ──
    if props.get("fixedRatio"):
        try:
            from lxml import etree
            from pptx.oxml.ns import qn
            nsmap_a_lock = "http://schemas.openxmlformats.org/drawingml/2006/main"
            cNvPicPr = shape._element.find(f".//{{{nsmap_a_lock}}}cNvPicPr")
            if cNvPicPr is not None:
                pic_locks = cNvPicPr.find(f"{{{nsmap_a_lock}}}picLocks")
                if pic_locks is None:
                    pic_locks = etree.SubElement(cNvPicPr, qn("a:picLocks"))
                pic_locks.set("noChangeAspect", "1")
        except Exception as e:
            logger.debug(f"Failed to write fixedRatio: {e}")

    return shape


def _apply_shape_keypoints(shape, shape_type: Optional[str], keypoints: List[float], path_formula: Optional[str] = None) -> None:
    """将前端 keypoints 写回到 prstGeom/avLst/gd。"""
    if not keypoints:
        return

    formula_name = path_formula or _map_shape_type_to_formula(shape_type)
    if not formula_name:
        return
    meta = _FORMULA_KEYPOINT_META.get(formula_name)
    if not meta:
        return

    try:
        from lxml import etree
        from pptx.oxml.ns import qn

        nsmap_a = "http://schemas.openxmlformats.org/drawingml/2006/main"
        sp_pr = _find_sp_pr(shape._element)
        if sp_pr is None:
            return

        prst_geom = sp_pr.find(f"{{{nsmap_a}}}prstGeom")
        if prst_geom is None:
            return

        av_lst = prst_geom.find(f"{{{nsmap_a}}}avLst")
        if av_lst is None:
            av_lst = etree.SubElement(prst_geom, qn("a:avLst"))

        for old in av_lst.findall(f"{{{nsmap_a}}}gd"):
            av_lst.remove(old)

        names: List[str] = meta.get("names", [])
        ranges: List[Tuple[float, float]] = meta.get("ranges", [])
        for idx, name in enumerate(names):
            if idx >= len(keypoints):
                break
            try:
                value = float(keypoints[idx])
            except Exception:
                continue
            min_v, max_v = ranges[idx] if idx < len(ranges) else (0.0, 1.0)
            value = max(min_v, min(max_v, value))
            gd = etree.SubElement(av_lst, qn("a:gd"))
            gd.set("name", name)
            gd.set("fmla", f"val {int(round(value * 100000))}")
    except Exception:
        pass


def _apply_round_rect_preset_geometry(shape, ratio: float) -> None:
    """写入 roundRect 预设几何（四角一致）。"""
    try:
        from lxml import etree
        from pptx.oxml.ns import qn

        nsmap_a = "http://schemas.openxmlformats.org/drawingml/2006/main"
        sp_pr = _find_sp_pr(shape._element)
        if sp_pr is None:
            return

        for old_geom in sp_pr.findall(f"{{{nsmap_a}}}prstGeom"):
            sp_pr.remove(old_geom)
        for old_geom in sp_pr.findall(f"{{{nsmap_a}}}custGeom"):
            sp_pr.remove(old_geom)

        prst_geom = etree.SubElement(sp_pr, qn("a:prstGeom"))
        prst_geom.set("prst", "roundRect")
        av_lst = etree.SubElement(prst_geom, qn("a:avLst"))
        gd = etree.SubElement(av_lst, qn("a:gd"))
        gd.set("name", "adj")
        gd.set("fmla", f"val {int(round(_clamp_round_rect_ratio(ratio) * 100000))}")
    except Exception:
        pass


def _apply_round_rect_custom_geometry(shape, keypoints: List[float]) -> None:
    """写入四角可独立配置的 custGeom 圆角矩形。"""
    try:
        from lxml import etree
        from pptx.oxml.ns import qn

        nsmap_a = "http://schemas.openxmlformats.org/drawingml/2006/main"
        sp_pr = _find_sp_pr(shape._element)
        if sp_pr is None:
            return

        for old_geom in sp_pr.findall(f"{{{nsmap_a}}}prstGeom"):
            sp_pr.remove(old_geom)
        for old_geom in sp_pr.findall(f"{{{nsmap_a}}}custGeom"):
            sp_pr.remove(old_geom)

        kps = _normalize_round_rect_keypoints(keypoints, 0.1) or [0.1, 0.1, 0.1, 0.1]
        w = max(1, int(shape.width))
        h = max(1, int(shape.height))
        short = min(w, h)
        rtl = int(round(short * kps[0]))
        rtr = int(round(short * kps[1]))
        rbr = int(round(short * kps[2]))
        rbl = int(round(short * kps[3]))

        cust_geom = etree.SubElement(sp_pr, qn("a:custGeom"))
        etree.SubElement(cust_geom, qn("a:avLst"))
        etree.SubElement(cust_geom, qn("a:gdLst"))
        etree.SubElement(cust_geom, qn("a:ahLst"))
        etree.SubElement(cust_geom, qn("a:cxnLst"))
        rect = etree.SubElement(cust_geom, qn("a:rect"))
        rect.set("l", "l")
        rect.set("t", "t")
        rect.set("r", "r")
        rect.set("b", "b")

        path_lst = etree.SubElement(cust_geom, qn("a:pathLst"))
        path = etree.SubElement(path_lst, qn("a:path"))
        path.set("w", str(w))
        path.set("h", str(h))

        def _append_pt(parent, x: int, y: int) -> None:
            pt = etree.SubElement(parent, qn("a:pt"))
            pt.set("x", str(int(x)))
            pt.set("y", str(int(y)))

        move_to = etree.SubElement(path, qn("a:moveTo"))
        _append_pt(move_to, rtl, 0)

        ln_1 = etree.SubElement(path, qn("a:lnTo"))
        _append_pt(ln_1, w - rtr, 0)

        q_1 = etree.SubElement(path, qn("a:quadBezTo"))
        _append_pt(q_1, w, 0)
        _append_pt(q_1, w, rtr)

        ln_2 = etree.SubElement(path, qn("a:lnTo"))
        _append_pt(ln_2, w, h - rbr)

        q_2 = etree.SubElement(path, qn("a:quadBezTo"))
        _append_pt(q_2, w, h)
        _append_pt(q_2, w - rbr, h)

        ln_3 = etree.SubElement(path, qn("a:lnTo"))
        _append_pt(ln_3, rbl, h)

        q_3 = etree.SubElement(path, qn("a:quadBezTo"))
        _append_pt(q_3, 0, h)
        _append_pt(q_3, 0, h - rbl)

        ln_4 = etree.SubElement(path, qn("a:lnTo"))
        _append_pt(ln_4, 0, rtl)

        q_4 = etree.SubElement(path, qn("a:quadBezTo"))
        _append_pt(q_4, 0, 0)
        _append_pt(q_4, rtl, 0)

        etree.SubElement(path, qn("a:close"))
    except Exception:
        pass


# ── SVG path → OOXML custGeom ─────────────────────────────────────
#
# 用于把 PPTElement.shape 上的 props.path + props.viewBox（来自 PPTist /
# crawlspace / 模型直生成）转成 OOXML <a:custGeom>，让 PowerPoint 真按 path
# 渲染装饰元素（波纹/连接线/L 形角标/复杂圆角等），不再降级成裸 prstGeom rect。
#
# 设计取舍：
#   - 仅支持 M/L/H/V/Q/C/Z（绝对 + 相对），覆盖 PPTist 内置 shapes.ts 与
#     模型常用的命令；S/T/A 暂不支持，回退到 prstGeom（add_shape 默认行为）。
#   - 已经能完美对应 prstGeom 的简单形状（rect / ellipse / 箭头 / 星形 等）
#     继续走 prstGeom，保留 PowerPoint 原生可编辑性；
#     只有"客户端没指定 pptxShapeType"或"形状不在 _PRESET_GEOM_PREFERRED_TYPES
#     里"的复杂自定义路径才走 custGeom。

# SVG path 命令对应的参数个数（每组）
_SVG_PATH_PARAM_COUNT: Dict[str, int] = {
    "M": 2, "m": 2,
    "L": 2, "l": 2,
    "H": 1, "h": 1,
    "V": 1, "v": 1,
    "C": 6, "c": 6,
    "S": 4, "s": 4,
    "Q": 4, "q": 4,
    "T": 2, "t": 2,
    "A": 7, "a": 7,
    "Z": 0, "z": 0,
}

# 命令字母或数字（支持负号、小数、科学计数法）
_SVG_PATH_TOKEN_RE = re.compile(
    r"[MmLlHhVvCcSsQqTtAaZz]|[+-]?\d+\.?\d*(?:[eE][+-]?\d+)?"
)

# 已能用 OOXML prstGeom 完美表达的形状类型；存在 pptxShapeType 且命中时
# 继续走 prstGeom（保留 PowerPoint 内置形状的可编辑控点）。
_PRESET_GEOM_PREFERRED_TYPES: Set[str] = {
    # 基础矩形/椭圆/三角
    "rect", "roundRect", "round1Rect", "round2SameRect", "round2DiagRect",
    "snipRoundRect", "snip2DiagRect",
    "ellipse", "oval",
    "triangle", "rtTriangle", "rightTriangle", "diamond",
    "parallelogram", "trapezoid", "pentagon", "hexagon", "octagon",
    # 装饰
    "star4", "star5", "star6", "star6Point",
    "rightArrow", "leftArrow", "upArrow", "downArrow",
    "leftRightArrow", "upDownArrow", "notchedRightArrow",
    "heart", "lightningBolt", "plus", "cross", "cloud", "chevron",
    "callout1", "callout2",
}


def _iter_svg_path_commands(path_str: str):
    """生成 (cmd_letter, [floats]) 序列。

    支持 SVG 1.1 的隐式命令复用（多组参数共用一个字母），
    M 之后隐式重复变成 L、m 之后变成 l。
    """
    tokens = _SVG_PATH_TOKEN_RE.findall(path_str)
    n = len(tokens)
    i = 0
    while i < n:
        tok = tokens[i]
        if tok not in _SVG_PATH_PARAM_COUNT:
            # 异常 token（路径头尾的空白经常出现脏数据），跳过
            i += 1
            continue
        cmd = tok
        i += 1
        count = _SVG_PATH_PARAM_COUNT[cmd]
        if count == 0:
            yield (cmd, [])
            continue
        first = True
        while i + count <= n:
            chunk = tokens[i:i + count]
            if any(c in _SVG_PATH_PARAM_COUNT for c in chunk):
                break
            try:
                params = [float(c) for c in chunk]
            except ValueError:
                break
            yield (cmd, params)
            i += count
            if first:
                first = False
                # M 之后的隐式命令是 L、m 之后是 l
                if cmd == "M":
                    cmd = "L"
                elif cmd == "m":
                    cmd = "l"


def _apply_svg_path_custom_geometry(
    shape,
    path_str: str,
    view_box: List[float],
) -> bool:
    """将 SVG path 转为 OOXML <a:custGeom> 写入 shape。

    Args:
        shape: python-pptx Shape 对象（已通过 add_shape 创建）。
        path_str: SVG path d 属性，如 "M 0 0 L 100 0 ... Z"。
        view_box: [width, height]，SVG path 的坐标系尺寸。

    Returns:
        True 表示已写入 custGeom；False 表示参数无效或命令不支持，
        调用方应保留 add_shape 创建的 prstGeom 作为兜底。
    """
    if not path_str or not isinstance(path_str, str):
        return False
    if not view_box or len(view_box) < 2:
        return False

    try:
        vb_w = float(view_box[0])
        vb_h = float(view_box[1])
    except (TypeError, ValueError):
        return False
    if vb_w <= 0 or vb_h <= 0:
        return False

    try:
        from lxml import etree
        from pptx.oxml.ns import qn
    except ImportError:
        return False

    nsmap_a = "http://schemas.openxmlformats.org/drawingml/2006/main"
    sp_pr = _find_sp_pr(shape._element)
    if sp_pr is None:
        return False

    # 先把 SVG 命令解析出来；任何一条不支持的命令都让整体回退到 prstGeom。
    try:
        commands = list(_iter_svg_path_commands(path_str))
    except Exception as e:
        logger.warning(f"_apply_svg_path_custom_geometry: parse failed: {e}")
        return False
    if not commands:
        return False

    supported = {"M", "m", "L", "l", "H", "h", "V", "v", "Q", "q", "C", "c", "Z", "z"}
    for cmd, _ in commands:
        if cmd not in supported:
            logger.debug(
                f"_apply_svg_path_custom_geometry: command '{cmd}' not supported, "
                f"falling back to prstGeom"
            )
            return False

    # shape 的 EMU 尺寸；path w/h 直接用 shape 的 EMU 宽高，
    # 然后把 SVG 坐标按比例缩放到 EMU。
    try:
        shape_w = max(1, int(shape.width))
        shape_h = max(1, int(shape.height))
    except Exception:
        return False

    sx = shape_w / vb_w
    sy = shape_h / vb_h

    def _to_emu_x(x: float) -> int:
        return int(round(x * sx))

    def _to_emu_y(y: float) -> int:
        return int(round(y * sy))

    # 先收集所有要写入的子节点（先解析、后落盘），避免半成品 custGeom 污染 spPr。
    cur_x = 0.0
    cur_y = 0.0
    start_x = 0.0
    start_y = 0.0
    has_move = False
    pending: List[Tuple[str, List[Tuple[float, float]]]] = []

    for cmd, params in commands:
        upper = cmd.upper()
        is_rel = cmd.islower()

        if upper == "M":
            x, y = params[0], params[1]
            if is_rel and has_move:
                x += cur_x
                y += cur_y
            pending.append(("moveTo", [(x, y)]))
            cur_x, cur_y = x, y
            start_x, start_y = x, y
            has_move = True

        elif upper == "L":
            x, y = params[0], params[1]
            if is_rel:
                x += cur_x
                y += cur_y
            pending.append(("lnTo", [(x, y)]))
            cur_x, cur_y = x, y

        elif upper == "H":
            x = params[0]
            if is_rel:
                x += cur_x
            pending.append(("lnTo", [(x, cur_y)]))
            cur_x = x

        elif upper == "V":
            y = params[0]
            if is_rel:
                y += cur_y
            pending.append(("lnTo", [(cur_x, y)]))
            cur_y = y

        elif upper == "Q":
            cx, cy, x, y = params
            if is_rel:
                cx += cur_x
                cy += cur_y
                x += cur_x
                y += cur_y
            pending.append(("quadBezTo", [(cx, cy), (x, y)]))
            cur_x, cur_y = x, y

        elif upper == "C":
            c1x, c1y, c2x, c2y, x, y = params
            if is_rel:
                c1x += cur_x; c1y += cur_y
                c2x += cur_x; c2y += cur_y
                x += cur_x; y += cur_y
            pending.append(("cubicBezTo", [(c1x, c1y), (c2x, c2y), (x, y)]))
            cur_x, cur_y = x, y

        elif upper == "Z":
            pending.append(("close", []))
            cur_x, cur_y = start_x, start_y

    if not pending:
        return False
    # 必须至少有一个 moveTo
    if pending[0][0] != "moveTo":
        return False

    # 清理 add_shape 自动加的 prstGeom / 之前残留的 custGeom
    for old_geom in sp_pr.findall(f"{{{nsmap_a}}}prstGeom"):
        sp_pr.remove(old_geom)
    for old_geom in sp_pr.findall(f"{{{nsmap_a}}}custGeom"):
        sp_pr.remove(old_geom)

    cust_geom = etree.SubElement(sp_pr, qn("a:custGeom"))
    etree.SubElement(cust_geom, qn("a:avLst"))
    etree.SubElement(cust_geom, qn("a:gdLst"))
    etree.SubElement(cust_geom, qn("a:ahLst"))
    etree.SubElement(cust_geom, qn("a:cxnLst"))
    rect = etree.SubElement(cust_geom, qn("a:rect"))
    rect.set("l", "l")
    rect.set("t", "t")
    rect.set("r", "r")
    rect.set("b", "b")

    path_lst = etree.SubElement(cust_geom, qn("a:pathLst"))
    path = etree.SubElement(path_lst, qn("a:path"))
    path.set("w", str(shape_w))
    path.set("h", str(shape_h))

    def _append_pt(parent, x: float, y: float) -> None:
        pt = etree.SubElement(parent, qn("a:pt"))
        pt.set("x", str(_to_emu_x(x)))
        pt.set("y", str(_to_emu_y(y)))

    # spPr 子元素顺序要求：xfrm 之后才是 prstGeom/custGeom，
    # 这里用 SubElement 追加（与 _apply_round_rect_custom_geometry 一致）。
    # PowerPoint 容忍 spPr 末尾的 custGeom 与其他子元素的相对位置。

    for node_name, pts in pending:
        node = etree.SubElement(path, qn(f"a:{node_name}"))
        for x, y in pts:
            _append_pt(node, x, y)

    return True


def _write_shape_element(slide, props: dict, left, top, width, height):
    """写入形状元素（还原形状类型、填充、轮廓、文本）。返回 shape 对象。"""
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.dml.color import RGBColor
    from pptx.util import Pt

    # ── 1. 还原形状类型（不再全部退化为矩形）──
    shape_type_str = props.get("pptxShapeType") or "rect"
    mso_shape = _map_pptx_shape_to_mso(shape_type_str)
    shape = slide.shapes.add_shape(mso_shape, left, top, width, height)
    _remove_shape_theme_style(shape)

    # ── 1a. 可调关键点（keypoints）──
    keypoints = props.get("keypoints")
    path_formula = props.get("pathFormula")
    consumed_round_rect_keypoints = False
    is_round_rect_formula = path_formula == "roundRect" or (
        not path_formula and shape_type_str in _ROUND_RECT_PRESET_TYPES
    )
    if is_round_rect_formula and isinstance(keypoints, list) and keypoints:
        rr_keypoints = _normalize_round_rect_keypoints(keypoints, 0.1)
        if rr_keypoints:
            consumed_round_rect_keypoints = True
            if _round_rect_keypoints_uniform(rr_keypoints):
                _apply_round_rect_preset_geometry(shape, rr_keypoints[0])
            else:
                _apply_round_rect_custom_geometry(shape, rr_keypoints)

    if isinstance(keypoints, list) and keypoints and not consumed_round_rect_keypoints:
        _apply_shape_keypoints(
            shape=shape,
            shape_type=shape_type_str,
            keypoints=keypoints,
            path_formula=path_formula,
        )

    # ── 1b. 复杂自定义 SVG path → custGeom ──
    # 让有显式 path + viewBox 的形状真按路径渲染（波纹/L 形角标/连接线等），
    # 不再降级到 prstGeom rect 丢失细节。
    # 跳过条件：
    #   1) 圆角矩形已通过 _apply_round_rect_custom_geometry 写过 custGeom；
    #   2) 形状有 pathFormula（说明 keypoints 已写到 prstGeom/avLst）；
    #   3) 客户端显式指定 pptxShapeType 且在「简单 prstGeom 完美映射」清单内。
    explicit_pptx_shape_type = props.get("pptxShapeType")
    if isinstance(explicit_pptx_shape_type, str):
        explicit_pptx_shape_type = explicit_pptx_shape_type.strip() or None
    else:
        explicit_pptx_shape_type = None
    svg_path = props.get("path")
    view_box = props.get("viewBox")
    if (
        not consumed_round_rect_keypoints
        and not path_formula
        and isinstance(svg_path, str)
        and svg_path.strip()
        and isinstance(view_box, (list, tuple))
        and len(view_box) >= 2
        and (
            explicit_pptx_shape_type is None
            or explicit_pptx_shape_type not in _PRESET_GEOM_PREFERRED_TYPES
        )
    ):
        _apply_svg_path_custom_geometry(shape, svg_path, list(view_box))

    # ── 2. 填充：图片/图案 > 渐变 > 纯色 > 无 ──
    pattern = props.get("pattern")
    gradient = props.get("gradient")
    fill_color = props.get("fill")
    fill_theme_key = props.get("fillThemeKey")
    fill_transforms = props.get("fillThemeTransforms")

    pattern_applied = False
    if pattern and isinstance(pattern, str):
        pattern_applied = _apply_shape_pattern_fill(shape, pattern)

    if not pattern_applied:
        if gradient and isinstance(gradient, dict):
            _apply_gradient_fill(shape, gradient)
        elif (fill_color and isinstance(fill_color, str)) or fill_theme_key:
            try:
                shape.fill.solid()
                applied_theme_fill = False
                if fill_theme_key:
                    sp_pr = _find_sp_pr(shape._element)
                    if sp_pr is not None:
                        applied_theme_fill = _apply_theme_color_to_solid_parent(
                            sp_pr,
                            fill_theme_key,
                            fill_color,
                            transforms=fill_transforms,
                        )

                if not applied_theme_fill and fill_color and isinstance(fill_color, str):
                    hex_color, fill_alpha = _parse_css_color(fill_color)
                    hex_6 = hex_color.lstrip("#")
                    if len(hex_6) == 6:
                        shape.fill.fore_color.rgb = RGBColor.from_string(hex_6)
                        # 写回 fill alpha
                        if fill_alpha is not None and fill_alpha < 1.0:
                            _apply_color_alpha(shape, "fill", fill_alpha)
            except Exception:
                pass
        else:
            shape.fill.background()
    else:
        # 图片填充已通过 blipFill 写入，不再覆盖 fill
        pass

    # ── 3. 轮廓/边框 ──
    outline = props.get("outline")
    if outline and isinstance(outline, dict):
        try:
            ln = shape.line
            if outline.get("width"):
                ln.width = Pt(float(outline["width"]))
            outline_theme_key = outline.get("themeKey")
            applied_theme_outline = False
            if outline_theme_key:
                nsmap_a_outline = {"a": "http://schemas.openxmlformats.org/drawingml/2006/main"}
                sp_pr_outline = _find_sp_pr(shape._element)
                ln_elem_outline = sp_pr_outline.find("a:ln", nsmap_a_outline) if sp_pr_outline is not None else None
                if ln_elem_outline is not None:
                    applied_theme_outline = _apply_theme_color_to_solid_parent(
                        ln_elem_outline,
                        outline_theme_key,
                        outline.get("color"),
                    )

            if not applied_theme_outline and outline.get("color"):
                hex_c, ln_alpha = _parse_css_color(outline["color"])
                hex_6 = hex_c.lstrip("#")
                if len(hex_6) == 6:
                    ln.color.rgb = RGBColor.from_string(hex_6)
                if ln_alpha is not None and ln_alpha < 1.0:
                    _apply_color_alpha(shape, "outline", ln_alpha)
            dash_style = outline.get("style", "solid")
            from lxml import etree
            from pptx.oxml.ns import qn

            nsmap_a = {"a": "http://schemas.openxmlformats.org/drawingml/2006/main"}
            sp_pr = _find_sp_pr(shape._element)
            ln_elem = sp_pr.find("a:ln", nsmap_a) if sp_pr is not None else None
            if ln_elem is not None:
                # 清理历史 dash，避免重复节点和旧值残留导致读回误判
                ln_elem.attrib.pop("prstDash", None)
                for old_dash in ln_elem.findall("a:prstDash", nsmap_a):
                    ln_elem.remove(old_dash)

                _shape_dash_map = {
                    "dashed": "dash", "dotted": "dot",
                    "dashDot": "dashDot", "longDash": "lgDash", "longDashDot": "lgDashDot",
                }
                _ooxml_val = _shape_dash_map.get(dash_style)
                if _ooxml_val:
                    prst_dash = etree.SubElement(ln_elem, qn("a:prstDash"))
                    prst_dash.set("val", _ooxml_val)
        except Exception as e:
            logger.debug(f"Failed to apply outline: {e}")
    else:
        _apply_no_line_fill(shape)

    # ── 4. 内部文本（保留 HTML 格式）──
    text_props = props.get("text")
    if text_props and isinstance(text_props, dict) and text_props.get("content"):
        _write_shape_text(shape, text_props)

    return shape


def _map_pptx_shape_to_mso(shape_type_str: str):
    """将前端 pptxShapeType 字符串映射回 MSO_SHAPE 枚举"""
    from pptx.enum.shapes import MSO_SHAPE

    # 静态安全映射（只包含 python-pptx 确定支持的形状类型）
    mapping = {
        "rect": MSO_SHAPE.RECTANGLE,
        "roundRect": MSO_SHAPE.ROUNDED_RECTANGLE,
        "round1Rect": getattr(MSO_SHAPE, "ROUND_1_RECTANGLE", MSO_SHAPE.ROUNDED_RECTANGLE),
        "round2SameRect": getattr(MSO_SHAPE, "ROUND_2_SAME_RECTANGLE", MSO_SHAPE.ROUNDED_RECTANGLE),
        "round2DiagRect": getattr(MSO_SHAPE, "ROUND_2_DIAG_RECTANGLE", MSO_SHAPE.ROUNDED_RECTANGLE),
        "snip2DiagRect": getattr(MSO_SHAPE, "SNIP_2_DIAG_RECTANGLE", MSO_SHAPE.RECTANGLE),
        "snipRoundRect": getattr(MSO_SHAPE, "SNIP_ROUND_RECTANGLE", MSO_SHAPE.ROUNDED_RECTANGLE),
        "ellipse": MSO_SHAPE.OVAL,
        "diamond": MSO_SHAPE.DIAMOND,
        "triangle": MSO_SHAPE.ISOSCELES_TRIANGLE,
        "rtTriangle": MSO_SHAPE.RIGHT_TRIANGLE,
        "rightTriangle": MSO_SHAPE.RIGHT_TRIANGLE,
        "parallelogram": MSO_SHAPE.PARALLELOGRAM,
        "trapezoid": MSO_SHAPE.TRAPEZOID,
        "pentagon": MSO_SHAPE.PENTAGON,
        "hexagon": MSO_SHAPE.HEXAGON,
        "octagon": MSO_SHAPE.OCTAGON,
        "star5": MSO_SHAPE.STAR_5_POINT,
        "star4": MSO_SHAPE.STAR_4_POINT,
        "star6": MSO_SHAPE.STAR_6_POINT,
        "star6Point": MSO_SHAPE.STAR_6_POINT,
        "rightArrow": MSO_SHAPE.RIGHT_ARROW,
        "leftArrow": MSO_SHAPE.LEFT_ARROW,
        "upArrow": MSO_SHAPE.UP_ARROW,
        "downArrow": MSO_SHAPE.DOWN_ARROW,
        "leftRightArrow": MSO_SHAPE.LEFT_RIGHT_ARROW,
        "upDownArrow": MSO_SHAPE.UP_DOWN_ARROW,
        "heart": MSO_SHAPE.HEART,
        "lightningBolt": MSO_SHAPE.LIGHTNING_BOLT,
        "plus": MSO_SHAPE.CROSS,
        "cross": MSO_SHAPE.CROSS,
        "cloud": MSO_SHAPE.CLOUD,
        "chevron": MSO_SHAPE.CHEVRON,
        "notchedRightArrow": MSO_SHAPE.NOTCHED_RIGHT_ARROW,
        "callout1": getattr(MSO_SHAPE, "RECTANGULAR_CALLOUT", MSO_SHAPE.RECTANGLE),
        "callout2": getattr(MSO_SHAPE, "ROUNDED_RECTANGULAR_CALLOUT", MSO_SHAPE.RECTANGLE),
    }

    if shape_type_str in mapping:
        return mapping[shape_type_str]

    # 动态查找（兼容 python-pptx 版本差异）
    def _camel_to_upper_snake(name: str) -> str:
        chars = []
        for i, ch in enumerate(name):
            prev = name[i - 1] if i > 0 else ""
            if ch.isupper() and i > 0 and (prev.islower() or prev.isdigit()):
                chars.append("_")
            elif ch.isdigit() and i > 0 and prev.isalpha() and prev.islower():
                chars.append("_")
            chars.append(ch.upper())
        return "".join(chars)

    upper_name = shape_type_str.upper()
    snake_name = _camel_to_upper_snake(shape_type_str)
    candidates = [upper_name, snake_name]

    # 部分枚举名以 RECTANGLE 结尾
    if snake_name.endswith("_RECT"):
        candidates.append(snake_name + "ANGLE")

    for candidate in candidates:
        try:
            return MSO_SHAPE[candidate]
        except (KeyError, AttributeError):
            pass
        try:
            return getattr(MSO_SHAPE, candidate)
        except AttributeError:
            pass

    logger.debug(f"Unsupported shape type '{shape_type_str}', falling back to RECTANGLE")
    return MSO_SHAPE.RECTANGLE


def _apply_gradient_fill(shape, gradient: dict) -> None:
    """
    应用渐变填充到形状（写出 OOXML <a:gradFill>）。

    采用纯 lxml 路径，直接在 spPr 上挂 <a:gradFill>，并清理任何已存在的填充节点
    （noFill/solidFill/gradFill/blipFill/pattFill/grpFill），保证生成稳定可控的 XML。

    Wave 2 改造点：以前依赖 python-pptx 的 fill.gradient() 高级 API，对部分 shape 类型
    （含主题继承 / placeholder 等）不稳定，且 alpha 写入要从私有成员 _gs 取节点，
    维护成本高。现在直接由本函数控制 OOXML 结构，便于单元测试与排查。

    Gradient 字段（PPTElement.props.gradient，源头：dom_extractor 解析 CSS linear-gradient）：
        {
          "type": "linear" | "radial",
          "rotate": deg (linear 用，PPTX 语义：0=左→右，90=上→下),
          "colors": [{"pos": 0..1, "color": "#RRGGBB" | "#RRGGBBAA"}, ...],
          "center": {"x": 0..1, "y": 0..1}  # 仅 radial
        }
    """
    from lxml import etree
    from pptx.oxml.ns import qn

    try:
        colors = gradient.get("colors") or []
        if not colors:
            return

        nsmap_a = "http://schemas.openxmlformats.org/drawingml/2006/main"
        sp_pr = _find_sp_pr(shape._element)
        if sp_pr is None:
            return

        # 清理已有填充节点，避免 solidFill / gradFill 共存导致 PowerPoint 取错
        for fill_tag in ("noFill", "solidFill", "gradFill", "blipFill", "pattFill", "grpFill"):
            for old in sp_pr.findall(f"{{{nsmap_a}}}{fill_tag}"):
                sp_pr.remove(old)

        # 创建 gradFill；写入 flip/rotWithShape
        grad_fill = etree.SubElement(sp_pr, qn("a:gradFill"))
        grad_fill.set("flip", "none")
        grad_fill.set("rotWithShape", "1")

        # gsLst（gradient stops list）
        gs_lst = etree.SubElement(grad_fill, qn("a:gsLst"))
        n = len(colors)
        for i, stop in enumerate(colors):
            raw_color = stop.get("color", "#000000")
            hex_c, stop_alpha = _parse_css_color(raw_color)
            color_hex = hex_c.lstrip("#")[:6].upper() or "000000"

            if "pos" in stop and stop["pos"] is not None:
                try:
                    pos_norm = float(stop["pos"])
                except (TypeError, ValueError):
                    pos_norm = i / max(1, n - 1) if n > 1 else 0.0
            else:
                pos_norm = i / max(1, n - 1) if n > 1 else 0.0
            pos_norm = max(0.0, min(1.0, pos_norm))
            pos_val = int(round(pos_norm * 100000))

            gs = etree.SubElement(gs_lst, qn("a:gs"))
            gs.set("pos", str(pos_val))
            srgb = etree.SubElement(gs, qn("a:srgbClr"))
            srgb.set("val", color_hex)
            if stop_alpha is not None and 0 <= stop_alpha < 1.0:
                alpha_el = etree.SubElement(srgb, qn("a:alpha"))
                alpha_el.set("val", str(int(round(stop_alpha * 100000))))

        # 单 stop 兜底：OOXML 至少要 2 个 gs，缺则克隆首个到 100%
        if n < 2:
            first_gs = list(gs_lst)[0]
            extra = etree.fromstring(etree.tostring(first_gs))
            extra.set("pos", "100000")
            gs_lst.append(extra)

        # 渐变方向
        grad_type = (gradient.get("type") or "linear").lower()
        if grad_type == "radial":
            path_el = etree.SubElement(grad_fill, qn("a:path"))
            path_el.set("path", "circle")
            fill_to_rect = etree.SubElement(path_el, qn("a:fillToRect"))
            _center = gradient.get("center")
            if isinstance(_center, dict):
                _cx = int(max(0.0, min(1.0, float(_center.get("x", 0.5)))) * 100000)
                _cy = int(max(0.0, min(1.0, float(_center.get("y", 0.5)))) * 100000)
            else:
                _cx, _cy = 50000, 50000
            fill_to_rect.set("l", str(_cx))
            fill_to_rect.set("t", str(_cy))
            fill_to_rect.set("r", str(100000 - _cx))
            fill_to_rect.set("b", str(100000 - _cy))
        else:
            try:
                rotate_deg = float(gradient.get("rotate", 0) or 0)
            except (TypeError, ValueError):
                rotate_deg = 0.0
            # OOXML <a:lin ang> 单位 1/60000 度，0..21600000 区间
            ang_val = int(round(rotate_deg * 60000)) % (360 * 60000)
            lin = etree.SubElement(grad_fill, qn("a:lin"))
            lin.set("ang", str(ang_val))
            lin.set("scaled", "0")
    except Exception as e:
        logger.debug(f"Failed to apply gradient: {e}")


def _apply_shape_pattern_fill(shape, pattern_src: str) -> bool:
    """应用图片/图案填充到 shape（写入 blipFill）。"""
    import base64 as b64
    import io
    import os
    import urllib.request
    from urllib.parse import unquote_to_bytes
    from lxml import etree
    from pptx.oxml.ns import qn

    try:
        image_bytes = None
        if pattern_src.startswith("data:"):
            header, payload = pattern_src.split(",", 1)
            if ";base64" in header:
                image_bytes = b64.b64decode(payload)
            else:
                image_bytes = unquote_to_bytes(payload)
        elif pattern_src.startswith("http://") or pattern_src.startswith("https://"):
            try:
                # OSS 走 SDK，外部 URL 走 SSRF 安全 HTTPS（统一入口）
                image_bytes = _download_image_smart(pattern_src, max_bytes=MAX_EXPORT_IMAGE_BYTES)
            except ValueError:
                logger.warning("SSRF blocked pattern fill URL: %s", pattern_src[:120])
                return
            if len(image_bytes) > MAX_EXPORT_IMAGE_BYTES:
                logger.warning("Pattern fill image too large (%d bytes), skipping: %s…", len(image_bytes), pattern_src[:120])
                return
        else:
            logger.warning("Unsupported pattern source (local paths disallowed): %s…", pattern_src[:120])
            return False

        if not image_bytes:
            return False

        image_stream = io.BytesIO(image_bytes)
        _image_part, r_id = shape.part.get_or_add_image_part(image_stream)

        nsmap_a = "http://schemas.openxmlformats.org/drawingml/2006/main"
        nsmap_p = "http://schemas.openxmlformats.org/presentationml/2006/main"
        sp_elem = shape._element
        spPr = sp_elem.find(f"{{{nsmap_p}}}spPr")
        if spPr is None:
            spPr = sp_elem.find(f".//{{{nsmap_a}}}spPr")
        if spPr is None:
            return False

        # 移除已有填充，避免 fill 类型冲突
        for fill_tag in ("noFill", "solidFill", "gradFill", "blipFill", "pattFill", "grpFill"):
            for node in spPr.findall(f"{{{nsmap_a}}}{fill_tag}"):
                spPr.remove(node)

        blip_fill = etree.SubElement(spPr, qn("a:blipFill"))
        blip_fill.set("rotWithShape", "1")
        blip = etree.SubElement(blip_fill, qn("a:blip"))
        blip.set(qn("r:embed"), r_id)
        stretch = etree.SubElement(blip_fill, qn("a:stretch"))
        etree.SubElement(stretch, qn("a:fillRect"))
        return True
    except Exception as e:
        logger.debug(f"Failed to apply shape pattern fill: {e}")
        return False


def _write_shape_text(shape, text_props: dict) -> None:
    """写入形状内部文本（保留 HTML 富文本格式）"""
    import re
    from lxml import etree
    from pptx.util import Pt
    from pptx.enum.text import PP_ALIGN

    content = text_props.get("content", "")
    # 检查是否有任何文本内容
    clean_text = re.sub(r"<[^>]+>", "", content).strip()
    if not clean_text:
        return

    tf = shape.text_frame
    _apply_text_word_wrap(tf, text_props)
    tf.text = ""  # 清除默认文本

    default_font = text_props.get("defaultFontFamily") or text_props.get("defaultFontName")
    raw_default_size = text_props.get("defaultFontSize")
    default_size = f"{raw_default_size}px" if isinstance(raw_default_size, (int, float)) and raw_default_size > 0 else raw_default_size
    default_color = text_props.get("defaultColor")
    default_color_theme_key = text_props.get("defaultColorThemeKey")

    # 设置垂直对齐
    v_align = text_props.get("verticalAlign", "middle")
    try:
        nsmap = {"a": "http://schemas.openxmlformats.org/drawingml/2006/main"}
        body_pr = tf._txBody.find("a:bodyPr", nsmap)
        if body_pr is not None:
            anchor_xml = {"top": "t", "middle": "ctr", "bottom": "b"}
            body_pr.set("anchor", anchor_xml.get(v_align, "ctr"))
    except Exception:
        pass

    # 内边距：始终显式写出，避免 PowerPoint 套用默认幽灵内边距致形状内文字换行漂移。
    _apply_text_frame_insets(tf, text_props.get("margin"))

    # 溢出缩放：形状内文字默认 shrink-to-fit，缺字机器上回退字体更宽时自动缩小而非撑破形状。
    _apply_text_autofit(tf, text_props.get("autoFit"))

    # 解析 HTML 为段落结构（复用文本框的 HTML 解析器）
    paragraphs = _parse_html_to_paragraphs(content)
    if not paragraphs:
        paragraphs = [{"runs": [{"text": clean_text}], "align": None, "lineHeight": None}]

    align_map = {
        "left": PP_ALIGN.LEFT,
        "center": PP_ALIGN.CENTER,
        "right": PP_ALIGN.RIGHT,
        "justify": PP_ALIGN.JUSTIFY,
    }
    # 默认水平对齐（形状文本通常居中）
    h_align = text_props.get("align", "center")

    for p_idx, para_data in enumerate(paragraphs):
        if p_idx == 0:
            p = tf.paragraphs[0] if tf.paragraphs else tf.add_paragraph()
        else:
            p = tf.add_paragraph()

        # 段落对齐：优先使用段落自身的对齐，否则使用形状默认对齐
        p_align = para_data.get("align") or h_align
        if p_align and p_align in align_map:
            p.alignment = align_map[p_align]

        # 段落行距
        p_line_height = para_data.get("lineHeight")
        if p_line_height:
            lh_val = _parse_css_line_height_ratio(p_line_height)
            if lh_val is not None and lh_val > 0:
                p.line_spacing = lh_val

        # 段前间距
        p_margin_top = para_data.get("marginTop")
        if p_margin_top:
            mt_pt = _css_length_to_pt(p_margin_top)
            if mt_pt is not None and mt_pt > 0:
                p.space_before = Pt(round(mt_pt, 3))

        # 段后间距
        p_margin_bottom = para_data.get("marginBottom")
        if p_margin_bottom:
            mb_pt = _css_length_to_pt(p_margin_bottom)
            if mb_pt is not None and mb_pt > 0:
                p.space_after = Pt(round(mb_pt, 3))

        # 段落缩进（padding-left → marL, text-indent → indent）
        p_padding_left = para_data.get("paddingLeft")
        p_text_indent = para_data.get("textIndent")
        if p_padding_left or p_text_indent:
            try:
                nsmap_a_pi = "http://schemas.openxmlformats.org/drawingml/2006/main"
                p_pr = p._p.find(f"{{{nsmap_a_pi}}}pPr")
                if p_pr is None:
                    p_pr = etree.SubElement(p._p, f"{{{nsmap_a_pi}}}pPr")
                    p._p.insert(0, p_pr)
                if p_padding_left:
                    pl_px = _css_length_to_px(p_padding_left)
                    if pl_px is not None:
                        mar_l_emu = int(round(pl_px * PT_PER_PX * 12700))
                        p_pr.set("marL", str(mar_l_emu))
                if p_text_indent:
                    ti_pt = _css_length_to_pt(p_text_indent)
                    if ti_pt is not None:
                        indent_emu = int(round(ti_pt * 12700))
                        p_pr.set("indent", str(indent_emu))
            except Exception:
                pass

        runs_data = para_data.get("runs", [])
        if not runs_data:
            run = p.add_run()
            run.text = ""
            continue

        # 列表项目符号
        bullet_type = para_data.get("bullet")
        if bullet_type:
            try:
                nsmap_a = "http://schemas.openxmlformats.org/drawingml/2006/main"
                p_pr = p._p.find(f"{{{nsmap_a}}}pPr")
                if p_pr is None:
                    p_pr = etree.SubElement(p._p, f"{{{nsmap_a}}}pPr")
                    p._p.insert(0, p_pr)
                level = para_data.get("level", 0)
                if level > 0:
                    p_pr.set("lvl", str(level))
                    p_pr.set("marL", str(level * 457200))
                    p_pr.set("indent", str(-228600))
                # 项目符号样式（颜色/大小/字体）
                _bu_sty2 = para_data.get("bulletStyle")
                if _bu_sty2:
                    _bu_clr2 = _bu_sty2.get("color")
                    if _bu_clr2:
                        _buClrE2 = etree.SubElement(p_pr, f"{{{nsmap_a}}}buClr")
                        _srgbE2 = etree.SubElement(_buClrE2, f"{{{nsmap_a}}}srgbClr")
                        _srgbE2.set("val", _bu_clr2.lstrip("#")[:6].upper())
                    _bu_fsz2 = _bu_sty2.get("fontSize")
                    if _bu_fsz2:
                        if _bu_fsz2.endswith("%"):
                            _pct2 = int(_bu_fsz2.rstrip("%")) * 1000
                            _szE2 = etree.SubElement(p_pr, f"{{{nsmap_a}}}buSzPct")
                            _szE2.set("val", str(_pct2))
                        elif _bu_fsz2.endswith("pt"):
                            _pts2 = int(float(_bu_fsz2.rstrip("pt")) * 100)
                            _szE2 = etree.SubElement(p_pr, f"{{{nsmap_a}}}buSzPts")
                            _szE2.set("val", str(_pts2))
                    _bu_fnt2 = _bu_sty2.get("fontFamily")
                    if _bu_fnt2:
                        _buFntE2 = etree.SubElement(p_pr, f"{{{nsmap_a}}}buFont")
                        _buFntE2.set("typeface", _bu_fnt2)
                if bullet_type == "bullet":
                    bu_char = etree.SubElement(p_pr, f"{{{nsmap_a}}}buChar")
                    bu_char.set("char", para_data.get("bulletChar") or "\u2022")
                elif bullet_type == "number":
                    bu_auto = etree.SubElement(p_pr, f"{{{nsmap_a}}}buAutoNum")
                    num_fmt = para_data.get("numberFormat") or "arabicPeriod"
                    bu_auto.set("type", num_fmt)
            except Exception:
                pass

        for run_data in runs_data:
            run = p.add_run()
            run.text = run_data.get("text", "")
            _apply_run_style(
                run,
                run_data,
                default_font=default_font,
                default_size=default_size,
                default_color=default_color,
                default_color_theme_key=default_color_theme_key,
            )


def _write_line_element(
    slide,
    element: Dict[str, Any],
    slide_width_emu: int,
    slide_height_emu: int,
    canvas_width: int,
    canvas_height: int,
) -> None:
    """写入线条/连接器元素"""
    from pptx.oxml.ns import qn
    from lxml import etree

    props = element.get("props", {})
    start = props.get("start", [0, 0])
    end = props.get("end", [100, 0])
    line_color = props.get("color", "#333333")
    line_theme_key = props.get("colorThemeKey")
    line_width = props.get("lineWidth", 2)
    line_style = props.get("style", "solid")
    points_arr = props.get("points", ["", ""])
    try:
        line_width_pt = max(0.1, float(line_width) * _export_dpi_scale(slide_width_emu, canvas_width))
    except (TypeError, ValueError):
        line_width_pt = max(0.1, 2.0 * _export_dpi_scale(slide_width_emu, canvas_width))

    def _pair(raw: Any, fallback: List[float]) -> List[float]:
        try:
            if isinstance(raw, (list, tuple)) and len(raw) >= 2:
                return [float(raw[0]), float(raw[1])]
        except (TypeError, ValueError):
            pass
        return [float(fallback[0]), float(fallback[1])]

    start = _pair(start, [0.0, 0.0])
    end = _pair(end, [100.0, 0.0])
    broken = _pair(props.get("broken"), [0.0, 0.0]) if props.get("broken") is not None else None
    broken2 = _pair(props.get("broken2"), [0.0, 0.0]) if props.get("broken2") is not None else None
    curve = _pair(props.get("curve"), [0.0, 0.0]) if props.get("curve") is not None else None

    cubic_raw = props.get("cubic")
    cubic = None
    if isinstance(cubic_raw, (list, tuple)) and len(cubic_raw) >= 2:
        cp1 = _pair(cubic_raw[0], start)
        cp2 = _pair(cubic_raw[1], end)
        cubic = [cp1, cp2]

    # 计算真实几何包围盒（包含折点/控制点），避免仅按 start/end 导致曲线写回塌陷。
    all_pts: List[List[float]] = [start, end]
    if broken is not None:
        all_pts.append(broken)
    if broken2 is not None:
        all_pts.append(broken2)
        all_pts.append([(start[0] + end[0]) / 2.0, (broken2[1] + end[1]) / 2.0])
    if curve is not None:
        all_pts.append(curve)
    if cubic is not None:
        all_pts.extend(cubic)

    min_x = min(pt[0] for pt in all_pts)
    max_x = max(pt[0] for pt in all_pts)
    min_y = min(pt[1] for pt in all_pts)
    max_y = max(pt[1] for pt in all_pts)

    def _norm(pt: List[float]) -> List[float]:
        return [pt[0] - min_x, pt[1] - min_y]

    start_n = _norm(start)
    end_n = _norm(end)
    broken_n = _norm(broken) if broken is not None else None
    broken2_n = _norm(broken2) if broken2 is not None else None
    curve_n = _norm(curve) if curve is not None else None
    cubic_n = [_norm(cubic[0]), _norm(cubic[1])] if cubic is not None else None

    try:
        x = float(element.get("x", 0) or 0) + min_x
    except (TypeError, ValueError):
        x = min_x
    try:
        y = float(element.get("y", 0) or 0) + min_y
    except (TypeError, ValueError):
        y = min_y

    w = max_x - min_x
    h = max_y - min_y
    if w <= 0:
        w = max(abs(end_n[0] - start_n[0]), 1.0)
    if h < 0:
        h = 0.0

    left_emu = px_to_emu(x, canvas_width, slide_width_emu)
    top_emu = px_to_emu(y, canvas_height, slide_height_emu)
    width_emu = px_to_emu(w, canvas_width, slide_width_emu)
    height_emu = px_to_emu(max(h, 1), canvas_height, slide_height_emu)

    try:
        # 使用直线连接器（通过 shapes 的低级 API）
        cxnSp = slide.shapes._spTree.makeelement(
            qn("p:cxnSp"), {}
        )

        nsmap_a = "http://schemas.openxmlformats.org/drawingml/2006/main"

        # nvCxnSpPr
        nvCxnSpPr = etree.SubElement(cxnSp, qn("p:nvCxnSpPr"))
        cNvPr = etree.SubElement(nvCxnSpPr, qn("p:cNvPr"))
        cNvPr.set("id", str(_next_shape_id(slide.shapes._spTree)))
        cNvPr.set("name", element.get("name", "Connector"))
        if element.get("visible") is False:
            cNvPr.set("hidden", "1")
        cNvCxnSpPr = etree.SubElement(nvCxnSpPr, qn("p:cNvCxnSpPr"))
        line_locked = element.get("locked")
        if line_locked is None and isinstance(props, dict):
            line_locked = props.get("locked")
        if line_locked is True:
            cxn_sp_locks = etree.SubElement(cNvCxnSpPr, qn("a:cxnSpLocks"))
            cxn_sp_locks.set("noMove", "1")
            cxn_sp_locks.set("noResize", "1")
            cxn_sp_locks.set("noRot", "1")
        etree.SubElement(nvCxnSpPr, qn("p:nvPr"))

        # spPr
        spPr = etree.SubElement(cxnSp, qn("p:spPr"))

        # xfrm（变换）
        xfrm = etree.SubElement(spPr, f"{{{nsmap_a}}}xfrm")
        # 如果终点在起点左边或上面，需要 flip
        base_flip_h = end_n[0] < start_n[0]
        base_flip_v = end_n[1] < start_n[1]
        extra_flip_h = bool(_resolve_flip_flag(props, element, "flipH"))
        extra_flip_v = bool(_resolve_flip_flag(props, element, "flipV"))
        flip_h = base_flip_h ^ extra_flip_h
        flip_v = base_flip_v ^ extra_flip_v
        if flip_h:
            xfrm.set("flipH", "1")
        if flip_v:
            xfrm.set("flipV", "1")

        off = etree.SubElement(xfrm, f"{{{nsmap_a}}}off")
        off.set("x", str(left_emu))
        off.set("y", str(top_emu))
        ext = etree.SubElement(xfrm, f"{{{nsmap_a}}}ext")
        ext.set("cx", str(max(width_emu, 1)))
        ext.set("cy", str(max(height_emu, 1)))
        # 旋转（OOXML: 60000 = 1deg）
        try:
            rotate = float(element.get("rotate", 0) or 0)
            if abs(rotate) > 0.01:
                xfrm.set("rot", str(int((rotate % 360) * 60000)))
        except (TypeError, ValueError):
            pass

        # prstGeom — 根据线条类型选择连接器几何
        prstGeom = etree.SubElement(spPr, f"{{{nsmap_a}}}prstGeom")
        avLst = etree.SubElement(prstGeom, f"{{{nsmap_a}}}avLst")

        def _to_adj(raw_val: Any, denom: float, fallback: int = 50000) -> int:
            try:
                if not denom:
                    return fallback
                ratio = float(raw_val) / float(denom)
                return max(0, min(100000, int(round(ratio * 100000))))
            except (TypeError, ValueError):
                return fallback

        def _append_adj(name: str, val: int) -> None:
            gd = etree.SubElement(avLst, f"{{{nsmap_a}}}gd")
            gd.set("name", name)
            gd.set("fmla", f"val {val}")

        if broken2_n is not None:
            # 双折线：使用 bentConnector4 + 双调整值，保留 X/Y 折点语义。
            prstGeom.set("prst", "bentConnector4")
            _append_adj("adj1", _to_adj(broken2_n[0], w))
            _append_adj("adj2", _to_adj(broken2_n[1], h if h > 0 else 1.0))
        elif broken_n is not None:
            # 单折线：bentConnector3；额外写入 adj2 让回读能还原 Y 方向折点。
            prstGeom.set("prst", "bentConnector3")
            _append_adj("adj1", _to_adj(broken_n[0], w))
            _append_adj("adj2", _to_adj(broken_n[1], h if h > 0 else 1.0))
        elif cubic_n is not None:
            # 三次曲线使用 curvedConnector4，adj1/adj2 存 X 比率，adj3/adj4 存 Y 比率，
            # 保留完整的两个控制点信息，避免 roundtrip 时 Y 轴数据丢失。
            prstGeom.set("prst", "curvedConnector4")
            cp1 = cubic_n[0]
            cp2 = cubic_n[1]
            _append_adj("adj1", _to_adj(cp1[0], w))
            _append_adj("adj2", _to_adj(cp2[0], w))
            _append_adj("adj3", _to_adj(cp1[1], h if h > 0 else 1.0))
            _append_adj("adj4", _to_adj(cp2[1], h if h > 0 else 1.0))
        elif curve_n is not None:
            prstGeom.set("prst", "curvedConnector3")
            _append_adj("adj1", _to_adj(curve_n[0], w))
            _append_adj("adj2", _to_adj(curve_n[1], h if h > 0 else 1.0))
        else:
            prstGeom.set("prst", "line")

        # ln（线条样式）
        ln = etree.SubElement(spPr, f"{{{nsmap_a}}}ln")
        ln.set("w", str(int(line_width_pt * EMU_PER_PT)))

        hex_lc, lc_alpha = _parse_css_color(line_color)
        # 线条最终透明度 = 颜色 alpha × 元素 opacity
        try:
            element_opacity = float(element.get("opacity", 1.0))
        except (TypeError, ValueError):
            element_opacity = 1.0
        element_opacity = max(0.0, min(1.0, element_opacity))
        color_alpha = lc_alpha if lc_alpha is not None else 1.0
        final_alpha = max(0.0, min(1.0, round(color_alpha * element_opacity, 4)))

        applied_theme_line = False
        if line_theme_key:
            if final_alpha < 1.0:
                alpha_color = _hex_to_rgba(hex_lc, final_alpha)
            else:
                alpha_color = hex_lc
            applied_theme_line = _apply_theme_color_to_solid_parent(ln, line_theme_key, alpha_color)

        if not applied_theme_line:
            solidFill = etree.SubElement(ln, f"{{{nsmap_a}}}solidFill")
            srgbClr = etree.SubElement(solidFill, f"{{{nsmap_a}}}srgbClr")
            srgbClr.set("val", hex_lc.lstrip("#")[:6])
            if final_alpha < 1.0:
                alpha_el = etree.SubElement(srgbClr, qn("a:alpha"))
                alpha_el.set("val", str(int(final_alpha * 100000)))

        # dash style
        _line_style_to_ooxml = {
            "dashed": "dash",
            "dotted": "dot",
            "dashDot": "dashDot",
            "longDash": "lgDash",
            "longDashDot": "lgDashDot",
        }
        _ooxml_dash = _line_style_to_ooxml.get(line_style)
        if _ooxml_dash:
            prstDash = etree.SubElement(ln, f"{{{nsmap_a}}}prstDash")
            prstDash.set("val", _ooxml_dash)

        # 箭头/圆点端点（含尺寸 w/len）
        point_sizes = props.get("pointSizes", [{}, {}])
        if not isinstance(point_sizes, (list, tuple)) or len(point_sizes) < 2:
            point_sizes = [{}, {}]
        _valid_arrow_sizes = {"sm", "med", "lg"}
        if len(points_arr) >= 1 and points_arr[0]:
            head_type = _map_line_point_to_ooxml_arrow(str(points_arr[0]))
            if head_type != "none":
                headEnd = etree.SubElement(ln, f"{{{nsmap_a}}}headEnd")
                headEnd.set("type", head_type)
                hs = point_sizes[0] if isinstance(point_sizes[0], dict) else {}
                if hs.get("w") in _valid_arrow_sizes:
                    headEnd.set("w", hs["w"])
                if hs.get("len") in _valid_arrow_sizes:
                    headEnd.set("len", hs["len"])
        if len(points_arr) >= 2 and points_arr[1]:
            tail_type = _map_line_point_to_ooxml_arrow(str(points_arr[1]))
            if tail_type != "none":
                tailEnd = etree.SubElement(ln, f"{{{nsmap_a}}}tailEnd")
                tailEnd.set("type", tail_type)
                ts = point_sizes[1] if isinstance(point_sizes[1], dict) else {}
                if ts.get("w") in _valid_arrow_sizes:
                    tailEnd.set("w", ts["w"])
                if ts.get("len") in _valid_arrow_sizes:
                    tailEnd.set("len", ts["len"])

        # shadow（线条路径不走 _apply_common_write_props，需要在此单独处理）
        shadow = element.get("shadow") or props.get("shadow")
        if shadow and isinstance(shadow, dict):
            class _LineShapeProxy:
                def __init__(self, elem):
                    self._element = elem
            _apply_shadow_write(_LineShapeProxy(cxnSp), shadow)

        slide.shapes._spTree.append(cxnSp)

        # hyperlink（线条路径不走 _apply_common_write_props，这里单独处理）
        link_raw = element.get("link")
        if link_raw is None and isinstance(props, dict):
            link_raw = props.get("link")
        if link_raw is not None:
            try:
                line_shape = slide.shapes[-1]
                _apply_shape_hyperlink_write(line_shape, link_raw)
            except Exception:
                pass
    except Exception as e:
        logger.warning(f"Failed to write line element, falling back: {e}")


# ═══════════════════════════════════════════════
# 通用写回属性（rotate, opacity, flip, shadow）
# ═══════════════════════════════════════════════


def _apply_color_alpha(shape, target: str, alpha: float) -> None:
    """
    在写回 PPTX 时为颜色设置独立 alpha。

    target: "fill" → spPr/solidFill/颜色元素
            "outline" → spPr/ln/solidFill/颜色元素
    """
    from pptx.oxml.ns import qn
    from lxml import etree

    nsmap_a = "http://schemas.openxmlformats.org/drawingml/2006/main"
    try:
        spPr = _find_sp_pr(shape._element)
        if spPr is None:
            return

        if target == "fill":
            container = spPr.find(f"{{{nsmap_a}}}solidFill")
        elif target == "outline":
            ln = spPr.find(f"{{{nsmap_a}}}ln")
            container = ln.find(f"{{{nsmap_a}}}solidFill") if ln is not None else None
        else:
            return

        if container is not None and len(container) > 0:
            color_el = container[0]  # srgbClr or schemeClr
            # 移除已有 alpha
            for old in color_el.findall(f"{{{nsmap_a}}}alpha"):
                color_el.remove(old)
            alpha_el = etree.SubElement(color_el, qn("a:alpha"))
            alpha_el.set("val", str(int(alpha * 100000)))
    except Exception:
        pass


def _apply_common_write_props(shape, element: Dict[str, Any]) -> None:
    """应用通用写回属性：rotate, opacity, flipH/flipV, shadow, link"""
    from pptx.oxml.ns import qn
    from lxml import etree

    props = element.get("props", {})
    nsmap_a = "http://schemas.openxmlformats.org/drawingml/2006/main"
    nsmap_p = "http://schemas.openxmlformats.org/presentationml/2006/main"

    # ── 1. rotate（标准化到 [0, 360)，PPTX 只接受正角度） ──
    rotate = element.get("rotate", 0)
    if rotate:
        try:
            normalized = float(rotate) % 360
            shape.rotation = normalized
        except Exception:
            pass

    # ── 2. opacity ──
    # 语义：element.opacity 仅代表“元素整体透明度”。
    # - 对图片：写入 blip.alphaModFix（沿用历史兼容）。
    # - 对其它元素：写入 spPr.alphaModFix，不破坏 fill/outline 自身 alpha。
    opacity = element.get("opacity")
    if opacity is not None:
        try:
            opacity_val = max(0.0, min(1.0, float(opacity)))
        except (TypeError, ValueError):
            opacity_val = 1.0

        if opacity_val < 1.0:
            try:
                sp_elem = shape._element
                spPr = _find_sp_pr(sp_elem)
                if spPr is not None:
                    # 图片：优先写 blip 透明度
                    blip_fill = sp_elem.find(f".//{{{nsmap_p}}}blipFill")
                    if blip_fill is None:
                        blip_fill = sp_elem.find(f".//{{{nsmap_a}}}blipFill")
                    if blip_fill is not None:
                        blip = blip_fill.find(f"{{{nsmap_a}}}blip")
                        if blip is not None:
                            existing = 1.0
                            old_mod = blip.find(f"{{{nsmap_a}}}alphaModFix")
                            if old_mod is not None:
                                raw = old_mod.get("amt") or old_mod.get("val")
                                if raw:
                                    existing = int(raw) / 100000
                            final_alpha = max(0.0, min(1.0, round(existing * opacity_val, 6)))
                            for old in blip.findall(f"{{{nsmap_a}}}alphaModFix"):
                                blip.remove(old)
                            if final_alpha < 1.0:
                                alpha_mod = etree.SubElement(blip, qn("a:alphaModFix"))
                                alpha_mod.set("amt", str(int(final_alpha * 100000)))
                    else:
                        # 其它元素：写入 spPr 直属 alphaModFix
                        alpha_mod = spPr.find(f"{{{nsmap_a}}}alphaModFix")
                        existing = 1.0
                        if alpha_mod is not None:
                            raw = alpha_mod.get("val") or alpha_mod.get("amt")
                            if raw:
                                existing = int(raw) / 100000
                        final_alpha = max(0.0, min(1.0, round(existing * opacity_val, 6)))

                        if alpha_mod is None and final_alpha < 1.0:
                            alpha_mod = etree.SubElement(spPr, qn("a:alphaModFix"))
                        if alpha_mod is not None:
                            if final_alpha < 1.0:
                                alpha_mod.set("val", str(int(final_alpha * 100000)))
                                alpha_mod.attrib.pop("amt", None)
                            else:
                                spPr.remove(alpha_mod)
            except Exception:
                pass

    # ── 3. flipH / flipV ──
    flip_h = _resolve_flip_flag(props, element, "flipH")
    flip_v = _resolve_flip_flag(props, element, "flipV")
    if flip_h is not None or flip_v is not None:
        try:
            sp_elem = shape._element
            xfrm = _find_transform_xfrm(sp_elem)
            if xfrm is not None:
                if flip_h is True:
                    xfrm.set("flipH", "1")
                elif flip_h is False:
                    xfrm.attrib.pop("flipH", None)

                if flip_v is True:
                    xfrm.set("flipV", "1")
                elif flip_v is False:
                    xfrm.attrib.pop("flipV", None)
        except Exception:
            pass

    # ── 4. shadow ──
    shadow = element.get("shadow") or props.get("shadow")
    if shadow and isinstance(shadow, dict):
        _apply_shadow_write(shape, shadow)

    # ── 5. locked（写入/清理 *Locks） ──
    locked_raw = element.get("locked")
    if locked_raw is None:
        locked_raw = props.get("locked")
    if isinstance(locked_raw, bool):
        try:
            sp_elem = shape._element
            lock_targets = [
                (sp_elem.find(f".//{{{nsmap_p}}}cNvSpPr"), "spLocks"),
                (sp_elem.find(f".//{{{nsmap_p}}}cNvPicPr"), "picLocks"),
                (sp_elem.find(f".//{{{nsmap_p}}}cNvCxnSpPr"), "cxnSpLocks"),
                (sp_elem.find(f".//{{{nsmap_p}}}cNvGrpSpPr"), "grpSpLocks"),
            ]
            for c_nv, lock_tag in lock_targets:
                if c_nv is None:
                    continue
                locks = c_nv.find(f"{{{nsmap_a}}}{lock_tag}")
                if locked_raw:
                    if locks is None:
                        locks = etree.SubElement(c_nv, qn(f"a:{lock_tag}"))
                    locks.set("noMove", "1")
                    locks.set("noResize", "1")
                    locks.set("noRot", "1")
                else:
                    if locks is not None:
                        locks.attrib.pop("noMove", None)
                        locks.attrib.pop("noResize", None)
                        locks.attrib.pop("noRot", None)
                        if len(locks.attrib) == 0 and len(locks) == 0:
                            c_nv.remove(locks)
                break
        except Exception:
            pass

    # ── 6. visible（写入/清理 cNvPr@hidden） ──
    visible_raw = element.get("visible")
    if visible_raw is None:
        visible_raw = props.get("visible")
    if isinstance(visible_raw, bool):
        try:
            sp_elem = shape._element
            c_nv_pr = sp_elem.find(f".//{{{nsmap_p}}}cNvPr")
            if c_nv_pr is not None:
                if visible_raw is False:
                    c_nv_pr.set("hidden", "1")
                else:
                    c_nv_pr.attrib.pop("hidden", None)
        except Exception:
            pass

    # ── 7. link（元素整体超链接） ──
    link_raw: Any = None
    link_explicit = False
    if isinstance(element, dict) and "link" in element:
        link_raw = element.get("link")
        link_explicit = True
    elif isinstance(props, dict) and "link" in props:
        link_raw = props.get("link")
        link_explicit = True
    if link_explicit:
        _apply_shape_hyperlink_write(shape, link_raw)


def _apply_shadow_write(shape, shadow: dict) -> None:
    """写入元素阴影到 PPTX (outerShdw)"""
    import math
    from pptx.oxml.ns import qn
    from lxml import etree

    try:
        nsmap_a = "http://schemas.openxmlformats.org/drawingml/2006/main"
        nsmap_p = "http://schemas.openxmlformats.org/presentationml/2006/main"
        sp_elem = shape._element
        spPr = sp_elem.find(f"{{{nsmap_p}}}spPr")
        if spPr is None:
            spPr = sp_elem.find(f"{{{nsmap_a}}}spPr")
        if spPr is None:
            spPr = sp_elem.find(f".//{{{nsmap_p}}}spPr")
        if spPr is None:
            spPr = sp_elem.find(f".//{{{nsmap_a}}}spPr")
        if spPr is None:
            return

        effect_lst = spPr.find(f"{{{nsmap_a}}}effectLst")
        if effect_lst is None:
            effect_lst = etree.SubElement(spPr, qn("a:effectLst"))

        outer_shdw = etree.SubElement(effect_lst, qn("a:outerShdw"))

        h = float(shadow.get("h", 2))
        v = float(shadow.get("v", 2))
        blur = float(shadow.get("blur", 4))
        color = shadow.get("color", "#000000")

        outer_shdw.set("blurRad", str(int(blur * 12700)))

        dist = (h ** 2 + v ** 2) ** 0.5
        outer_shdw.set("dist", str(int(dist * 12700)))

        angle_rad = math.atan2(v, h)
        angle_deg = math.degrees(angle_rad)
        outer_shdw.set("dir", str(int(angle_deg * 60000)))

        # 解析颜色和透明度
        opacity_val = 1.0
        hex_color = color
        if color.startswith("rgba"):
            import re as _re
            m = _re.match(r"rgba?\((\d+),\s*(\d+),\s*(\d+),?\s*([\d.]*)\)", color)
            if m:
                r, g, b = int(m.group(1)), int(m.group(2)), int(m.group(3))
                opacity_val = float(m.group(4)) if m.group(4) else 1.0
                hex_color = f"#{r:02x}{g:02x}{b:02x}"

        # 独立 opacity 字段优先（提取阶段写入的标准格式）
        if "opacity" in shadow:
            try:
                opacity_val = float(shadow["opacity"])
            except (TypeError, ValueError):
                pass

        srgb = etree.SubElement(outer_shdw, qn("a:srgbClr"))
        srgb.set("val", hex_color.lstrip("#")[:6])

        if opacity_val < 1:
            alpha_el = etree.SubElement(srgb, qn("a:alpha"))
            alpha_el.set("val", str(int(opacity_val * 100000)))
    except Exception as e:
        logger.debug(f"Failed to apply shadow write: {e}")


def _set_slide_background_color(slide, color_str: str) -> None:
    """设置幻灯片纯色背景（支持 #RRGGBB / rgb() / rgba()）"""
    try:
        from pptx.dml.color import RGBColor
        from lxml import etree
        from pptx.oxml.ns import qn

        nsmap_a = "http://schemas.openxmlformats.org/drawingml/2006/main"
        hex_c, alpha = _parse_css_color(color_str)
        color_hex = hex_c.lstrip("#")[:6]

        bg = slide.background
        fill = bg.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor.from_string(color_hex)

        # 背景颜色透明度（写到 solidFill/srgbClr/alpha）
        # 限定搜索范围到 bg/bgPr 下，避免误命中 spTree 内形状的 solidFill
        nsmap_p = "http://schemas.openxmlformats.org/presentationml/2006/main"
        cSld = slide.background._element
        bg_node = cSld.find(f"{{{nsmap_p}}}bg")
        search_root = bg_node if bg_node is not None else cSld
        solid_fill = search_root.find(f".//{{{nsmap_a}}}solidFill")
        if solid_fill is not None and len(solid_fill) > 0:
            clr_el = solid_fill[0]
            for old in list(clr_el.findall(f"{{{nsmap_a}}}alpha")):
                clr_el.remove(old)
            if alpha is not None and alpha < 1.0:
                alpha_el = etree.SubElement(clr_el, qn("a:alpha"))
                alpha_el.set("val", str(int(alpha * 100000)))
    except Exception as e:
        logger.warning(f"Failed to set background color: {e}")


def _set_slide_background_theme(slide, theme_data: dict, fallback_color: Optional[str] = None) -> None:
    """设置幻灯片主题背景色（保留 theme 语义，失败时降级为纯色）"""
    try:
        from lxml import etree
        from pptx.enum.dml import MSO_THEME_COLOR
        from pptx.oxml.ns import qn

        raw_key = str(theme_data.get("key", "")).strip().lower()
        key_map = {
            "1": "DARK_1",
            "2": "LIGHT_1",
            "3": "DARK_2",
            "4": "LIGHT_2",
            "5": "ACCENT_1",
            "6": "ACCENT_2",
            "7": "ACCENT_3",
            "8": "ACCENT_4",
            "9": "ACCENT_5",
            "10": "ACCENT_6",
            "11": "HYPERLINK",
            "12": "FOLLOWED_HYPERLINK",
            "lt1": "LIGHT_1",
            "bg1": "BACKGROUND_1",
            "lt2": "LIGHT_2",
            "bg2": "BACKGROUND_2",
            "dk1": "DARK_1",
            "tx1": "TEXT_1",
            "dk2": "DARK_2",
            "tx2": "TEXT_2",
            "accent1": "ACCENT_1",
            "accent2": "ACCENT_2",
            "accent3": "ACCENT_3",
            "accent4": "ACCENT_4",
            "accent5": "ACCENT_5",
            "accent6": "ACCENT_6",
            "hlink": "HYPERLINK",
            "folhlink": "FOLLOWED_HYPERLINK",
        }
        enum_name = key_map.get(raw_key)
        if not enum_name:
            if fallback_color:
                _set_slide_background_color(slide, fallback_color)
            return

        bg = slide.background
        fill = bg.fill
        fill.solid()
        fill.fore_color.theme_color = getattr(MSO_THEME_COLOR, enum_name)

        nsmap_a = "http://schemas.openxmlformats.org/drawingml/2006/main"
        nsmap_p_bg = "http://schemas.openxmlformats.org/presentationml/2006/main"
        cSld_bg = slide.background._element
        bg_node_t = cSld_bg.find(f"{{{nsmap_p_bg}}}bg")
        search_root_t = bg_node_t if bg_node_t is not None else cSld_bg
        solid_fill = search_root_t.find(f".//{{{nsmap_a}}}solidFill")
        clr_el = None
        if solid_fill is not None and len(solid_fill) > 0:
            clr_el = solid_fill[0]  # schemeClr

        # 如果给了 rgba 颜色，附加 alpha（schemeClr/alpha）
        alpha_src = theme_data.get("color") or fallback_color
        if isinstance(alpha_src, str):
            _, alpha = _parse_css_color(alpha_src)
            if alpha is not None and alpha < 1.0 and clr_el is not None:
                for old in list(clr_el.findall(f"{{{nsmap_a}}}alpha")):
                    clr_el.remove(old)
                alpha_el = etree.SubElement(clr_el, qn("a:alpha"))
                alpha_el.set("val", str(int(alpha * 100000)))

        # 回写 lumMod/lumOff/tint/shade/satMod 变换参数（保持导入时的主题色修饰）
        transforms = theme_data.get("transforms")
        if isinstance(transforms, dict) and clr_el is not None:
            for tag_name in ("lumMod", "lumOff", "tint", "shade", "satMod"):
                val = transforms.get(tag_name)
                if val is not None:
                    try:
                        tr_el = etree.SubElement(clr_el, qn(f"a:{tag_name}"))
                        tr_el.set("val", str(int(val)))
                    except Exception:
                        pass
    except Exception as e:
        logger.warning(f"Failed to set background theme, fallback to color: {e}")
        if fallback_color:
            _set_slide_background_color(slide, fallback_color)


def _set_slide_background_gradient(slide, gradient: dict) -> None:
    """设置幻灯片渐变背景（支持 linear / radial 与 stop alpha）"""
    from pptx.dml.color import RGBColor

    try:
        colors = gradient.get("colors", [])
        if not colors or len(colors) < 2:
            return

        grad_type = str(gradient.get("type", "linear")).lower()
        nsmap_a = "http://schemas.openxmlformats.org/drawingml/2006/main"

        bg = slide.background
        fill = bg.fill
        fill.gradient()

        # 线性渐变角度
        rotate = gradient.get("rotate", 0)
        if grad_type != "radial":
            try:
                fill.gradient_angle = rotate
            except Exception:
                pass

        # 设置停靠点
        for i, stop in enumerate(colors):
            raw_c = stop.get("color", "#000000")
            hex_c, stop_alpha = _parse_css_color(raw_c)
            color_hex = hex_c.lstrip("#")[:6]
            pos_raw = stop.get("pos", i / (len(colors) - 1))
            try:
                pos = max(0.0, min(1.0, float(pos_raw)))
            except (TypeError, ValueError):
                pos = i / (len(colors) - 1)

            try:
                if i < len(fill.gradient_stops):
                    gs = fill.gradient_stops[i]
                else:
                    # python-pptx 某些版本没有 .add()，直接通过 XML 追加 gs 节点
                    try:
                        gs = fill.gradient_stops.add()
                    except (AttributeError, TypeError):
                        from lxml import etree as _et
                        from pptx.oxml.ns import qn as _qn
                        bg_elem = slide.background._element
                        grad_fill = bg_elem.find(f".//{{{nsmap_a}}}gradFill")
                        if grad_fill is not None:
                            gs_lst = grad_fill.find(f"{{{nsmap_a}}}gsLst")
                            if gs_lst is None:
                                gs_lst = _et.SubElement(grad_fill, _qn("a:gsLst"))
                            gs_node = _et.SubElement(gs_lst, _qn("a:gs"))
                            gs_node.set("pos", str(int(pos * 100000)))
                            srgb_node = _et.SubElement(gs_node, _qn("a:srgbClr"))
                            srgb_node.set("val", color_hex)
                            if stop_alpha is not None and stop_alpha < 1.0:
                                alpha_node = _et.SubElement(srgb_node, _qn("a:alpha"))
                                alpha_node.set("val", str(int(stop_alpha * 100000)))
                        continue
                gs.color.rgb = RGBColor.from_string(color_hex)
                gs.position = pos

                # per-stop alpha
                if stop_alpha is not None and stop_alpha < 1.0:
                    try:
                        from lxml import etree
                        from pptx.oxml.ns import qn
                        gs_el = gs._gs if hasattr(gs, "_gs") else None
                        if gs_el is not None:
                            srgb = gs_el.find(f"{{{nsmap_a}}}srgbClr")
                            if srgb is not None:
                                for old in list(srgb.findall(f"{{{nsmap_a}}}alpha")):
                                    srgb.remove(old)
                                alpha_el = etree.SubElement(srgb, qn("a:alpha"))
                                alpha_el.set("val", str(int(stop_alpha * 100000)))
                    except Exception:
                        pass
            except Exception:
                pass

        # 径向渐变：将 gradFill/lin 改为 gradFill/path
        if grad_type == "radial":
            try:
                from lxml import etree
                from pptx.oxml.ns import qn
                nsmap_p_g = "http://schemas.openxmlformats.org/presentationml/2006/main"
                cSld_g = slide.background._element
                bg_node_g = cSld_g.find(f"{{{nsmap_p_g}}}bg")
                search_root_g = bg_node_g if bg_node_g is not None else cSld_g
                grad_fill = search_root_g.find(f".//{{{nsmap_a}}}gradFill")
                if grad_fill is not None:
                    for lin in list(grad_fill.findall(f"{{{nsmap_a}}}lin")):
                        grad_fill.remove(lin)
                    path_el = grad_fill.find(f"{{{nsmap_a}}}path")
                    if path_el is None:
                        path_el = etree.SubElement(grad_fill, qn("a:path"))
                    path_el.set("path", "circle")
                    fill_to_rect = path_el.find(f"{{{nsmap_a}}}fillToRect")
                    if fill_to_rect is None:
                        fill_to_rect = etree.SubElement(path_el, qn("a:fillToRect"))
                    # 使用保存的中心点（如果有），否则默认居中
                    _center = gradient.get("center")
                    if isinstance(_center, dict):
                        _cx = int(max(0, min(1, float(_center.get("x", 0.5)))) * 100000)
                        _cy = int(max(0, min(1, float(_center.get("y", 0.5)))) * 100000)
                    else:
                        _cx, _cy = 50000, 50000
                    fill_to_rect.set("l", str(_cx))
                    fill_to_rect.set("t", str(_cy))
                    fill_to_rect.set("r", str(100000 - _cx))
                    fill_to_rect.set("b", str(100000 - _cy))
            except Exception:
                pass
    except Exception as e:
        logger.warning(f"Failed to set background gradient: {e}")


def _set_slide_background_image(slide, image_data: dict) -> None:
    """设置幻灯片图片背景（支持 cover / contain / repeat）"""
    import base64 as b64
    import io

    try:
        src = image_data.get("src", "")
        if not src:
            return

        size = str(image_data.get("size", "cover")).lower()
        if size not in ("cover", "contain", "repeat"):
            size = "cover"

        # 通过 XML 直接设置图片背景
        from pptx.oxml.ns import qn
        from lxml import etree

        # OSS 走 SDK，其他 URL 走 SSRF 安全 HTTPS；统一通过 _download_image_smart
        image_bytes = _download_image_smart(src, max_bytes=MAX_EXPORT_IMAGE_BYTES)
        if not image_bytes:
            logger.warning(f"Background image unavailable, skipping: {src[:120]}")
            return

        # 与普通图片元素保持一致：对 webp/svg 等格式做兼容转换（优先转 PNG）
        image_bytes = _normalize_image_bytes_for_pptx(image_bytes, src_hint=src)

        # 添加图片 part 到 slide
        image_stream = io.BytesIO(image_bytes)
        image_part, rId = slide.part.get_or_add_image_part(image_stream)
        _ = image_part  # keep part alive

        # 构建背景 XML
        # slide.background._element → cSld（CT_CommonSlideData）
        # 防御性检查：确保拿到的是 cSld 而非 bg 本身
        nsmap_p = "http://schemas.openxmlformats.org/presentationml/2006/main"
        _bg_elem = slide.background._element
        tag_local = etree.QName(_bg_elem.tag).localname if hasattr(_bg_elem, 'tag') else ""
        if tag_local == "bg":
            # 极端情况：_element 直接指向 bg 节点
            bg = _bg_elem
            cSld = bg.getparent()
        else:
            cSld = _bg_elem
            bg = cSld.find(f"{{{nsmap_p}}}bg")
        if bg is None:
            bg = etree.Element(qn("p:bg"))
            sp_tree = cSld.find(f"{{{nsmap_p}}}spTree") if cSld is not None else None
            if sp_tree is not None:
                cSld.insert(cSld.index(sp_tree), bg)
            elif cSld is not None:
                cSld.append(bg)
            else:
                logger.warning("Cannot find cSld to set background image")
                return

        # 清理原有 p:bg 子节点（避免破坏 spTree）
        for child in list(bg):
            bg.remove(child)

        bgPr = etree.SubElement(bg, qn("p:bgPr"))
        blipFill = etree.SubElement(bgPr, qn("a:blipFill"))
        blipFill.set("dpi", "0")
        blipFill.set("rotWithShape", "1")

        blip = etree.SubElement(blipFill, qn("a:blip"))
        blip.set(qn("r:embed"), rId)

        if size == "repeat":
            etree.SubElement(blipFill, qn("a:tile"))
        else:
            stretch = etree.SubElement(blipFill, qn("a:stretch"))
            fillRect = etree.SubElement(stretch, qn("a:fillRect"))
            if size == "contain":
                # 根据图片和幻灯片宽高比动态计算边距，实现 contain 效果
                try:
                    from PIL import Image as _PILImage
                    _img = _PILImage.open(io.BytesIO(image_bytes))
                    _iw, _ih = _img.size
                    _sw = slide.slide_layout.slide_master.slide_width if hasattr(slide, 'slide_layout') else 9144000
                    _sh = slide.slide_layout.slide_master.slide_height if hasattr(slide, 'slide_layout') else 6858000
                    _ir = _iw / _ih if _ih > 0 else 1.0
                    _sr = _sw / _sh if _sh > 0 else 1.0
                    if _ir > _sr:
                        _s = _sr / _ir
                        _m = int((1 - _s) / 2 * 100000)
                        fillRect.set("t", str(_m))
                        fillRect.set("b", str(_m))
                    else:
                        _s = _ir / _sr
                        _m = int((1 - _s) / 2 * 100000)
                        fillRect.set("l", str(_m))
                        fillRect.set("r", str(_m))
                except Exception:
                    fillRect.set("t", "5000")
                    fillRect.set("r", "5000")
                    fillRect.set("b", "5000")
                    fillRect.set("l", "5000")

        etree.SubElement(bgPr, qn("a:effectLst"))

    except Exception as e:
        logger.warning(f"Failed to set background image: {e}")
